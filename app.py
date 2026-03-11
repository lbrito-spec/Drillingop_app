# TNPIV29_app_fixed_v2 (con Viajes + Avance de profundidad)
# tnpiv3_app.py
# ------------------------------------------------------------
# Requisitos:
#   pip install streamlit pandas plotly reportlab python-pptx pillow kaleido
#
# Ejecutar:
#   streamlit run tnpiv3_app.py
# ------------------------------------------------------------

# --- FIX: alias seguro para evitar NameError (compatibilidad) ---
tipo_tiempo = None
# --- FIX: alias seguro para operaciones_seleccionadas ---
operaciones_seleccionadas = None
try:
    operaciones_seleccionadas = st.session_state.get('operacion_sel', None)
except Exception:
    pass

try:
    tipo_tiempo = st.session_state.get('tipo_time_general', None)
except Exception:
    pass

import os
import re
import base64
import json
import textwrap
from io import BytesIO
from datetime import datetime, date
import uuid

import streamlit as st

# Legacy safety default
legacy_calc_value = 0.0

mr_etapa_legacy = 0.0  # legacy var (kept to avoid NameError)
import pandas as pd

# --- Google OAuth + Drive (local) ---
# Requisitos:
#   pip install google-auth google-auth-oauthlib google-api-python-client requests
# Nota: para pruebas locales puede requerir:
#   export OAUTHLIB_INSECURE_TRANSPORT=1
try:
    import requests
    from google_auth_oauthlib.flow import Flow
    from google.oauth2.credentials import Credentials
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaInMemoryUpload
except Exception:
    # Permitimos que el script cargue aunque no estén instaladas las deps;
    # la UI mostrará un mensaje al usuario cuando intente usar OAuth/Drive.
    requests = None
    Flow = None
    Credentials = None
    build = None
    MediaInMemoryUpload = None

GOOGLE_SCOPES = [
    "openid",
    "https://www.googleapis.com/auth/userinfo.email",
    "https://www.googleapis.com/auth/userinfo.profile",
    "https://www.googleapis.com/auth/drive.file",
]

def _google_is_available() -> bool:
    return all([requests, Flow, Credentials, build, MediaInMemoryUpload])

def _drive_service():
    if not _google_is_available():
        return None
    gc = st.session_state.get("google_creds")
    if not gc:
        return None
    creds = Credentials(
        token=gc.get("token"),
        refresh_token=gc.get("refresh_token"),
        token_uri=gc.get("token_uri"),
        client_id=gc.get("client_id"),
        client_secret=gc.get("client_secret"),
        scopes=gc.get("scopes"),
    )
    return build("drive", "v3", credentials=creds)

def _ensure_drive_folder(drive, folder_name: str) -> str:
    q = f"mimeType='application/vnd.google-apps.folder' and name='{folder_name}' and trashed=false"
    res = drive.files().list(q=q, fields="files(id,name)").execute()
    files = res.get("files", [])
    if files:
        return files[0]["id"]
    meta = {"name": folder_name, "mimeType": "application/vnd.google-apps.folder"}
    folder = drive.files().create(body=meta, fields="id").execute()
    return folder["id"]

def _drive_upsert_json(drive, folder_id: str, filename: str, payload: dict) -> str:
    content = json.dumps(payload, ensure_ascii=False, indent=2, default=_json_default).encode("utf-8")
    media = MediaInMemoryUpload(content, mimetype="application/json", resumable=False)

    q = f"'{folder_id}' in parents and name='{filename}' and trashed=false"
    res = drive.files().list(q=q, fields="files(id,name)").execute()
    existing = (res.get("files") or [])

    if existing:
        file_id = existing[0]["id"]
        drive.files().update(fileId=file_id, media_body=media).execute()
        return file_id

    meta = {"name": filename, "parents": [folder_id]}
    created = drive.files().create(body=meta, media_body=media, fields="id").execute()
    return created["id"]

def _drive_list_json(drive, folder_id: str, limit: int = 100):
    q = f"'{folder_id}' in parents and mimeType='application/json' and trashed=false"
    res = drive.files().list(
        q=q,
        fields="files(id,name,modifiedTime)",
        orderBy="modifiedTime desc",
        pageSize=limit,
    ).execute()
    return res.get("files", [])

def _drive_download_json(drive, file_id: str) -> dict:
    data = drive.files().get_media(fileId=file_id).execute()
    return json.loads(data.decode("utf-8"))



def _calc_eff(prog: float, real: float) -> float:
    """Eficiencia (%): 100 si real <= programado; si real>prog => (prog/real)*100."""
    try:
        prog = float(prog)
        real = float(real)
    except Exception:
        return 0.0
    if real <= 0 or prog <= 0:
        return 0.0
    return 100.0 if real <= prog else (prog / real) * 100.0


# --- Helper seguro para float (evita None/valores inválidos) ---
def _safe_float(v, default=0.0) -> float:
    try:
        if v is None:
            return float(default)
        return float(v)
    except Exception:
        return float(default)


# --- FIX: asegurar RowID por registro (para edición en Detalle) ---
def _ensure_rowid(df_in: pd.DataFrame) -> pd.DataFrame:
    df = df_in.copy()
    if "RowID" not in df.columns:
        df.insert(0, "RowID", "")
    missing = df["RowID"].isna() | (df["RowID"].astype(str).str.strip() == "")
    if missing.any():
        df.loc[missing, "RowID"] = [str(uuid.uuid4()) for _ in range(int(missing.sum()))]
    return df


def _normalize_time_cause_columns(df_in: pd.DataFrame) -> pd.DataFrame:
    """Normaliza columnas de causas TNPI/TNP.

    - Asegura que existan Categoria/Detalle para TNPI y TNP.
    - Evita NaN/None en tablas y gráficas (usa "-").
    - Backward compatible: si hay registros antiguos con TNP guardado en
      Categoria_TNPI/Detalle_TNPI, los copia a Categoria_TNP/Detalle_TNP.
    """
    if df_in is None or df_in.empty:
        return df_in

    df = df_in.copy()

    for col in [
        "Categoria_TNPI",
        "Detalle_TNPI",
        "Categoria_TNP",
        "Detalle_TNP",
    ]:
        if col not in df.columns:
            df[col] = "-"

    # Asegura columnas de hora (opcional)
    for col in ["Hora_Inicio", "Hora_Fin"]:
        if col not in df.columns:
            df[col] = ""

    # Clean null-like values
    for col in ["Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP", "Detalle_TNP"]:
        df[col] = df[col].replace({None: "-"}).fillna("-")

    # Backfill TNP causes from TNPI columns if older data was stored there
    m_tnp = df.get("Tipo", "") == "TNP"
    if m_tnp.any():
        m_missing_cat = m_tnp & (df["Categoria_TNP"].isin(["-", "", "nan"]))
        m_missing_det = m_tnp & (df["Detalle_TNP"].isin(["-", "", "nan"]))
        if "Categoria_TNPI" in df.columns:
            df.loc[m_missing_cat, "Categoria_TNP"] = df.loc[m_missing_cat, "Categoria_TNPI"].replace({"nan": "-"})
        if "Detalle_TNPI" in df.columns:
            df.loc[m_missing_det, "Detalle_TNP"] = df.loc[m_missing_det, "Detalle_TNPI"].replace({"nan": "-"})

        # Final safety: don't leave NaNs
        df.loc[m_tnp, ["Categoria_TNP", "Detalle_TNP"]] = df.loc[m_tnp, ["Categoria_TNP", "Detalle_TNP"]].replace({None: "-"}).fillna("-")

    return df

import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import streamlit.components.v1 as components

# ------------------------------
# CSV robust loader (para CSVs con separador desconocido)
# ------------------------------
def robust_read_csv(src, encodings=("utf-8-sig","utf-8","latin-1"), seps=(",", ";", "\t", "|")) -> pd.DataFrame:
    """
    Lee CSV desde ruta o stream (UploadedFile) intentando:
      1) sep=None (Sniffer) con engine='python'
      2) una lista de separadores comunes
    También prueba encodings típicos (utf-8-sig/utf-8/latin-1).
    """
    # Obtener bytes si es stream (p.ej. st.file_uploader)
    data_bytes = None
    if hasattr(src, "getvalue"):
        try:
            data_bytes = src.getvalue()
        except Exception:
            data_bytes = None
    if data_bytes is None and hasattr(src, "read") and not isinstance(src, (str, bytes, os.PathLike)):
        try:
            # Ojo: read() consume el stream; por eso guardamos bytes
            data_bytes = src.read()
        except Exception:
            data_bytes = None

    def _try_read(buf, enc, sep):
        return pd.read_csv(buf, sep=sep, engine="python", encoding=enc)

    # Caso path en disco
    if isinstance(src, (str, os.PathLike)):
        last_err = None
        for enc in encodings:
            try:
                return pd.read_csv(src, sep=None, engine="python", encoding=enc)
            except Exception as e:
                last_err = e
            for s in seps:
                try:
                    return pd.read_csv(src, sep=s, engine="python", encoding=enc)
                except Exception as e:
                    last_err = e
        raise last_err if last_err else ValueError("No se pudo leer el CSV")

    # Caso stream / bytes
    if data_bytes is None:
        raise ValueError("No se pudo leer el archivo (stream vacío).")

    from io import BytesIO
    last_err = None
    for enc in encodings:
        try:
            return _try_read(BytesIO(data_bytes), enc, None)
        except Exception as e:
            last_err = e
        for s in seps:
            try:
                return _try_read(BytesIO(data_bytes), enc, s)
            except Exception as e:
                last_err = e
    raise last_err if last_err else ValueError("No se pudo leer el CSV")


from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from reportlab.lib.utils import ImageReader

from pptx import Presentation
from pptx.util import Inches, Pt

from PIL import Image

# ------------------------------
# PLOTLY EXPORT (kaleido)
# ------------------------------
PLOTLY_IMG_OK = True
try:
    import plotly.io as pio
    import plotly.graph_objects as go
    import importlib.util
    if importlib.util.find_spec("kaleido") is None:
        PLOTLY_IMG_OK = False
except Exception:
    PLOTLY_IMG_OK = False

# ------------------------------
# CONFIG STREAMLIT
# ------------------------------

def _semaforo_from_eff(eff):
    """Devuelve un semáforo (emoji) a partir de eficiencia en % (0-100)."""
    try:
        if eff is None:
            return "⚪"
        if isinstance(eff, str) and eff.strip()== "":
            return "⚪"
        val = float(eff)
    except Exception:
        return "⚪"
    if val >= 85:
        return "🟢"
    if val >= 75:
        return "🟡"
    return "🔴"

# Backward-compat alias used in some blocks
def _semaforo_text(eff):
    return _semaforo_from_eff(eff)

def semaforo_dot(eff):
    """Compat: devuelve bolita semáforo según eficiencia (%)."""
    return _semaforo_from_eff(eff)

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# FUNCIÓN AUXILIAR PARA RENDERIZAR HTML
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
def render_html(html_content: str, height: int = None):
    """
    Renderiza contenido HTML en Streamlit de manera robusta.
    Intenta usar st.html() primero, luego components.html(), y finalmente st.markdown().
    """
    # Intentar con st.html() (Streamlit 1.23.0+)
    if hasattr(st, 'html'):
        try:
            # Intentar con height si está disponible
            if height:
                st.html(html_content, height=height)
            else:
                st.html(html_content)
        except TypeError:
            # Si height no es soportado, usar sin parámetros
            st.html(html_content)
    # Fallback a components.html()
    elif hasattr(components, 'html'):
        components.html(html_content, height=height or 400, scrolling=True)
    # Último recurso: st.markdown con unsafe_allow_html
    else:
        st.markdown(html_content, unsafe_allow_html=True)

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# CHIPS UI (pro badges)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
def render_chip_row(items: list[dict], use_iframe: bool = False, height: int = 120) -> None:
    """
    Renderiza chips elegantes en una fila.
    items: [{"label": "...", "value": "...", "tone": "blue|green|amber|red|violet|gray"}]
    """
    light_mode = _is_light_theme()
    tones = {
        "blue": ("#38bdf8", "rgba(56, 189, 248, 0.28)", "rgba(56, 189, 248, 0.35)"),
        "green": ("#22c55e", "rgba(34, 197, 94, 0.28)", "rgba(34, 197, 94, 0.35)"),
        "amber": ("#f59e0b", "rgba(245, 158, 11, 0.28)", "rgba(245, 158, 11, 0.35)"),
        "red": ("#ef4444", "rgba(239, 68, 68, 0.28)", "rgba(239, 68, 68, 0.35)"),
        "violet": ("#8b5cf6", "rgba(139, 92, 246, 0.28)", "rgba(139, 92, 246, 0.35)"),
        "gray": ("#e2e8f0", "rgba(148, 163, 184, 0.20)", "rgba(148, 163, 184, 0.22)"),
    }
    if light_mode:
        tones["gray"] = ("#334155", "rgba(148, 163, 184, 0.18)", "rgba(148, 163, 184, 0.25)")
        tones["blue"] = ("#2563eb", "rgba(37, 99, 235, 0.16)", "rgba(37, 99, 235, 0.25)")
        tones["green"] = ("#16a34a", "rgba(22, 163, 74, 0.16)", "rgba(22, 163, 74, 0.25)")
        tones["amber"] = ("#d97706", "rgba(217, 119, 6, 0.16)", "rgba(217, 119, 6, 0.25)")
        tones["red"] = ("#dc2626", "rgba(220, 38, 38, 0.16)", "rgba(220, 38, 38, 0.25)")
        tones["violet"] = ("#7c3aed", "rgba(124, 58, 237, 0.16)", "rgba(124, 58, 237, 0.25)")
    chips_html = []
    for it in items:
        label = str(it.get("label", "")).strip()
        value = str(it.get("value", "")).strip()
        tone = it.get("tone", "gray")
        fg, bg, glow = tones.get(tone, tones["gray"])
        chips_html.append(
            f"""
            <div class="ds-chip" style="--chip-fg:{fg}; --chip-bg:{bg}; --chip-glow:{glow};">
              <span class="ds-chip-label">{label}</span>
              <span class="ds-chip-value">{value}</span>
            </div>
            """
        )

    label_color = "#475569" if light_mode else "rgba(255,255,255,0.70)"
    html = textwrap.dedent(
        f"""
        <style>
          .ds-chip-row {{
            display:flex; flex-wrap:wrap; gap:8px;
            padding: 2px 0 6px 0;
          }}
          .ds-chip {{
            display:inline-flex; align-items:center; gap:8px;
            padding: 6px 12px;
            border-radius: 999px;
            border: 1px solid rgba(255,255,255,0.16);
            background: linear-gradient(180deg, var(--chip-bg), rgba(255,255,255,0.02));
            box-shadow:
              inset 0 0 0 1px rgba(255,255,255,0.04),
              0 8px 20px rgba(0,0,0,0.28),
              0 0 16px var(--chip-glow);
            backdrop-filter: blur(6px);
            color: var(--chip-fg);
            font-size: 12px;
            font-weight: 800;
            letter-spacing: 0.2px;
          }}
          .ds-chip-label {{
            color: {label_color};
            font-weight: 700;
          }}
          .ds-chip-value {{
            color: var(--chip-fg);
          }}
        </style>
        <div class="ds-chip-row">
          {''.join(chips_html)}
        </div>
        """
    ).strip()
    if use_iframe and hasattr(components, "html"):
        components.html(html, height=height, scrolling=False)
    else:
        st.markdown(html, unsafe_allow_html=True)

def build_delta_chip_item(
    label: str,
    real: float,
    prog: float,
    unit: str = "",
    higher_is_better: bool = True,
    precision: int = 2,
) -> dict:
    """Devuelve un chip Δ vs prog con flecha y color."""
    try:
        real_v = float(real)
        prog_v = float(prog)
    except Exception:
        real_v = 0.0
        prog_v = 0.0
    delta = real_v - prog_v
    arrow = "↑" if delta >= 0 else "↓"
    good = (delta >= 0 and higher_is_better) or (delta <= 0 and not higher_is_better)
    tone = "green" if good else "red"
    fmt = f"{{delta:+.{precision}f}}"
    val = fmt.format(delta=delta)
    unit_txt = f" {unit}" if unit else ""
    return {
        "label": label,
        "value": f"{arrow} {val}{unit_txt} vs prog",
        "tone": tone,
    }

def _conn_exceso_suggestions(total_real_min: float, total_std_min: float, top_comp: str | None) -> list[dict]:
    """Chips pro con sugerencias cuando conexión supera estándar."""
    if total_std_min <= 0:
        return []
    over_min = max(0.0, float(total_real_min) - float(total_std_min))
    if over_min <= 0:
        return []
    over_pct = (over_min / float(total_std_min)) * 100.0 if total_std_min > 0 else 0.0
    tone = "red" if over_pct >= 20.0 else "amber"
    comp_txt = f"· foco en {top_comp}" if top_comp else ""
    return [
        {"label": "Exceso conexión", "value": f"{over_min:.1f} min", "tone": tone},
        {"label": "Sobre estándar", "value": f"{over_pct:.0f}% {comp_txt}".strip(), "tone": tone},
        {"label": "Sugerencia", "value": "Revisar pre/post y roles", "tone": "blue"},
        {"label": "Sugerencia", "value": "Checklist y materiales listos", "tone": "violet"},
    ]

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# MISSION CONTROL DASHBOARD (NASA Style)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
def mission_control_dashboard(etapa, eficiencia, tp_h, tnpi_h, tnp_h, total_real):
    """
    Panel de control tipo NASA con KPIs críticos
    """
    light_mode = _is_light_theme()
    # Determinar color del status basado en eficiencia
    if eficiencia >= 85:
        status_color = "#00ff88"
        status_text = "ÓPTIMO"
    elif eficiencia >= 75:
        status_color = "#ffaa00"
        status_text = "ATENCIÓN"
    else:
        status_color = "#ff4444"
        status_text = "CRÍTICO"
    if light_mode:
        panel_bg = "linear-gradient(180deg, #ffffff 0%, #f4f7fb 100%)"
        panel_border = "rgba(15,23,42,0.10)"
        panel_shadow = "0 8px 24px rgba(15,23,42,0.10)"
        card_bg = "rgba(15,23,42,0.03)"
        card_border = "rgba(15,23,42,0.08)"
        card_hover_bg = "rgba(15,23,42,0.06)"
        hover_border = "rgba(0, 136, 255, 0.35)"
        text_main = "#0f172a"
        text_muted = "#475569"
        badge_bg = "rgba(15,23,42,0.04)"
        progress_bg = "rgba(15,23,42,0.10)"
        divider = "rgba(15,23,42,0.10)"
    else:
        panel_bg = "linear-gradient(180deg, #0f1620 0%, #0a0e14 100%)"
        panel_border = "rgba(255,255,255,0.1)"
        panel_shadow = "0 8px 32px rgba(0,0,0,0.4)"
        card_bg = "rgba(255,255,255,0.05)"
        card_border = "rgba(255,255,255,0.08)"
        card_hover_bg = "rgba(255,255,255,0.1)"
        hover_border = "rgba(0, 136, 255, 0.3)"
        text_main = "#ffffff"
        text_muted = "rgba(255,255,255,0.7)"
        badge_bg = "rgba(255,255,255,0.05)"
        progress_bg = "rgba(255,255,255,0.1)"
        divider = "rgba(255,255,255,0.08)"

    # Calcular porcentajes
    tp_percent = (tp_h / total_real * 100) if total_real > 0 else 0
    tnpi_percent = (tnpi_h / total_real * 100) if total_real > 0 else 0
    tnp_percent = (tnp_h / total_real * 100) if total_real > 0 else 0
    
    return f"""
    <style>
    .mission-panel {{
        background: {panel_bg};
        border: 1px solid {panel_border};
        border-radius: 16px;
        padding: 20px;
        margin: 10px 0 20px 0;
        box-shadow: {panel_shadow};
        position: relative;
        overflow: hidden;
    }}
    .mission-panel::before {{
        content: '';
        position: absolute;
        top: 0;
        left: 0;
        right: 0;
        height: 2px;
        background: linear-gradient(90deg, #00ff88 0%, #0088ff 100%);
    }}
    .kpi-grid {{
        display: grid;
        grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
        gap: 16px;
        margin-top: 20px;
    }}
    .kpi-card {{
        background: {card_bg};
        border: 1px solid {card_border};
        border-radius: 12px;
        padding: 16px;
        text-align: center;
        transition: all 0.3s ease;
    }}
    .kpi-card:hover {{
        background: {card_hover_bg};
        border-color: {hover_border};
        transform: translateY(-2px);
    }}
    .kpi-value {{
        font-size: 28px;
        font-weight: 800;
        margin: 8px 0;
        color: {text_main};
        font-family: 'Courier New', monospace;
    }}
    .kpi-label {{
        font-size: 12px;
        color: {text_muted};
        text-transform: uppercase;
        letter-spacing: 1px;
    }}
    .status-indicator {{
        display: inline-block;
        width: 10px;
        height: 10px;
        border-radius: 50%;
        margin-right: 8px;
        animation: pulse 2s infinite;
    }}
    @keyframes pulse {{
        0% {{ opacity: 1; }}
        50% {{ opacity: 0.3; }}
        100% {{ opacity: 1; }}
    }}
    .mission-header {{
        display: flex;
        justify-content: space-between;
        align-items: center;
        margin-bottom: 15px;
    }}
    .mission-title {{
        font-size: 12px;
        color: {text_muted};
        text-transform: uppercase;
        letter-spacing: 1px;
    }}
    .mission-stage {{
        font-size: 20px;
        font-weight: 800;
        margin-top: 4px;
        background: linear-gradient(90deg, #00ff88, #0088ff);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }}
    .efficiency-badge {{
        background: {badge_bg};
        border: 1px solid {status_color};
        border-radius: 20px;
        padding: 8px 16px;
        text-align: center;
        min-width: 120px;
    }}
    .efficiency-value {{
        font-size: 32px;
        font-weight: 800;
        color: {status_color};
        line-height: 1;
    }}
    .efficiency-label {{
        font-size: 12px;
        color: {text_muted};
        margin-top: 2px;
    }}
    .progress-bar {{
        height: 6px;
        background: {progress_bg};
        border-radius: 3px;
        margin-top: 8px;
        overflow: hidden;
    }}
    .progress-fill {{
        height: 100%;
        border-radius: 3px;
    }}
    </style>
    
    <div class="mission-panel">
        <div class="mission-header">
            <div>
                <div class="mission-title">
                    <span class="status-indicator" style="background: {status_color};"></span>
                    MISSION CONTROL • DRILLSPOT
                </div>
                <div class="mission-stage">{etapa}</div>
            </div>
            <div class="efficiency-badge">
                <div class="efficiency-value">{eficiencia:.1f}%</div>
                <div class="efficiency-label">{status_text}</div>
            </div>
        </div>
        
        <div class="kpi-grid">
            <div class="kpi-card">
                <div class="kpi-label">TP Productivo</div>
                <div class="kpi-value" style="color: #00ff88;">{tp_h:.1f}h</div>
                <div class="progress-bar">
                    <div class="progress-fill" style="width: {tp_percent:.1f}%; background: #00ff88;"></div>
                </div>
                <div style="font-size: 11px; color: rgba(255,255,255,0.6);">{tp_percent:.1f}% del tiempo</div>
            </div>
            
            <div class="kpi-card">
                <div class="kpi-label">TNPI</div>
                <div class="kpi-value" style="color: #ffaa00;">{tnpi_h:.1f}h</div>
                <div class="progress-bar">
                    <div class="progress-fill" style="width: {tnpi_percent:.1f}%; background: #ffaa00;"></div>
                </div>
                <div style="font-size: 11px; color: rgba(255,255,255,0.6);">{tnpi_percent:.1f}% del tiempo</div>
            </div>
            
            <div class="kpi-card">
                <div class="kpi-label">TNP</div>
                <div class="kpi-value" style="color: #ff4444;">{tnp_h:.1f}h</div>
                <div class="progress-bar">
                    <div class="progress-fill" style="width: {tnp_percent:.1f}%; background: #ff4444;"></div>
                </div>
                <div style="font-size: 11px; color: rgba(255,255,255,0.6);">{tnp_percent:.1f}% del tiempo</div>
            </div>
            
            <div class="kpi-card">
                <div class="kpi-label">Tiempo Total</div>
                <div class="kpi-value" style="color: #0088ff;">{total_real:.1f}h</div>
                <div class="progress-bar">
                    <div class="progress-fill" style="width: 100%; background: linear-gradient(90deg, #00ff88, #0088ff);"></div>
                </div>
                <div style="font-size: 11px; color: rgba(255,255,255,0.6);">Operación actual</div>
            </div>
        </div>
        
        <div style="margin-top: 20px; padding-top: 15px; border-top: 1px solid {divider};">
            <div style="display: flex; justify-content: space-between; font-size: 11px; color: {text_muted};">
                <div>📊 <b>Estado:</b> {status_text}</div>
                <div>⏱️ <b>Actualizado:</b> {datetime.now().strftime('%H:%M:%S')}</div>
                <div>📍 <b>Etapa:</b> {etapa}</div>
            </div>
        </div>
    </div>
    """

# --- Helpers: coalesce duplicate columns (avoid losing data when columns are repeated) ---
def _coalesce_duplicate_columns(df: pd.DataFrame) -> pd.DataFrame:
    # If df has duplicate column names, coalesce them row-wise (first non-null/non-empty) into one.
    if df is None or df.empty or not df.columns.duplicated().any():
        return df
    out = df.copy()
    dup_names = [c for c in out.columns[out.columns.duplicated()].unique()]
    for name in dup_names:
        cols = [c for c in out.columns if c == name]
        base = out[cols[0]].copy()
        for c in cols[1:]:
            s = out[c]
            mask = base.isna() | (base.astype(str).str.strip() == '') | (base.astype(str).str.lower() == 'nan')
            base = base.where(~mask, s)
        out[name] = base
        # drop extra duplicate columns, keep the first occurrence only
        keep = []
        seen_first = False
        for c in out.columns:
            if c != name:
                keep.append(c)
            else:
                if not seen_first:
                    keep.append(c)
                    seen_first = True
        out = out.loc[:, keep]
    return out


def _decorate_turno_df(df: pd.DataFrame | None) -> pd.DataFrame | None:
    """Añade ☀️/🌙 en Turno para tablas sin tocar data base."""
    if df is None:
        return df
    _df = df.copy()
    if "Turno" in _df.columns:
        _df["Turno"] = _df["Turno"].replace({
            "Día": "Día ☀️",
            "Noche": "Noche 🌙",
            "Diurno": "Día ☀️",
            "Nocturno": "Noche 🌙",
        })
    return _df


def add_semaforo_column(df, eff_col="Eficiencia_pct"):
    """Agrega columna 'Semáforo' sin alterar estilos (solo texto)."""
    if df is None:
        return df
    df = _coalesce_duplicate_columns(df)
    if df is None or df.empty:
        return df
    if eff_col not in df.columns:
        return df
    _df = _decorate_turno_df(df)
    if _df is None:
        return _df
    _df["Semáforo"] = _df[eff_col].apply(_semaforo_from_eff)
    return _df

st.set_page_config(page_title="Dashboard Operativo DrillSpot", layout="wide")

# == == == == == == == == == == == == =
# Auth (solo Google OAuth)
# == == == == == == == == == == == == =

def _get_user_role(email: str) -> str:
    """Obtiene el rol desde secrets (por email)."""
    try:
        roles_map = dict(st.secrets.get("roles", {}))
        return str(roles_map.get(email, "") or "").strip() or "user"
    except Exception:
        return "user"

def _active_users_path() -> str:
    base_dir = os.path.dirname(os.path.abspath(__file__)) if "__file__" in globals() else os.getcwd()
    return os.path.join(base_dir, "active_users.json")

def _load_active_users() -> dict:
    try:
        p = _active_users_path()
        if os.path.exists(p):
            with open(p, "r", encoding="utf-8") as f:
                data = json.load(f)
            return data if isinstance(data, dict) else {}
    except Exception:
        pass
    return {}

def _save_active_users(data: dict) -> None:
    try:
        with open(_active_users_path(), "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception:
        pass

def _cleanup_active_users(data: dict, ttl_minutes: int = 30) -> dict:
    if not isinstance(data, dict):
        return {}
    out = {}
    now = datetime.now()
    for k, v in data.items():
        try:
            last_seen = datetime.fromisoformat(str(v.get("last_seen", "")))
        except Exception:
            last_seen = None
        if last_seen is None:
            continue
        if (now - last_seen).total_seconds() <= ttl_minutes * 60:
            out[k] = v
    return out

def _touch_active_user(user_meta: dict) -> None:
    if not isinstance(user_meta, dict):
        return
    key = (user_meta.get("email") or user_meta.get("username") or "").strip().lower()
    if not key:
        return
    data = _load_active_users()
    data = _cleanup_active_users(data, ttl_minutes=45)
    data[key] = {
        "name": user_meta.get("name") or user_meta.get("username") or key,
        "email": user_meta.get("email") or key,
        "photo_url": user_meta.get("photo_url") or "",
        "role": user_meta.get("role") or "user",
        "last_seen": datetime.now().isoformat(timespec="seconds"),
    }
    _save_active_users(data)

def _render_user_badge(user_meta: dict) -> str:
    name = (user_meta.get("name") or user_meta.get("username") or "").strip()
    photo = (user_meta.get("photo_url") or "").strip()
    role = (user_meta.get("role") or "").strip()
    initials = "".join([p[:1] for p in name.split() if p]).upper()[:2]
    # Badge fijo arriba derecha (no depende del header)
    if photo:
        img_html = f'<img src="{photo}" style="width:32px;height:32px;border-radius:999px;object-fit:cover;border:1px solid rgba(255,255,255,.25);" />'
    else:
        img_html = f'<div style="width:32px;height:32px;border-radius:999px;display:flex;align-items:center;justify-content:center;background:rgba(255,255,255,.12);border:1px solid rgba(255,255,255,.25);font-weight:800;font-size:12px;">{initials or "U"}</div>'
    return f"""
    <style>
      .user-badge {{
        position: fixed;
        top: 12px;
        right: 18px;
        z-index: 999999;
        display:flex;
        align-items:center;
        gap:10px;
        padding:6px 10px;
        border-radius:999px;
        background: rgba(15, 17, 22, 0.55);
        backdrop-filter: blur(8px);
        border: 1px solid rgba(255,255,255,.08);
        color: rgba(255,255,255,.92);
        font-size: 13px;
        line-height: 1;
      }}
      .user-badge .name {{
        max-width: 180px;
        overflow:hidden;
        text-overflow: ellipsis;
        white-space: nowrap;
      }}
    </style>
    <div class="user-badge">
      {img_html}
      <div class="name">{name}{f" · {role}" if role else ""}</div>
    </div>
    """


def _google_oauth_login_sidebar():
    # Estado base
    if "auth_ok" not in st.session_state:
        st.session_state["auth_ok"] = False
    if "auth_user" not in st.session_state:
        st.session_state["auth_user"] = None
    if "google_creds" not in st.session_state:
        st.session_state["google_creds"] = None

    with st.sidebar.expander("🔐 Acceso con Google", expanded=not st.session_state["auth_ok"]):
        if not _google_is_available():
            st.warning("Faltan dependencias para Google OAuth/Drive. Instala: google-auth, google-auth-oauthlib, google-api-python-client, requests.")
            return

        if st.session_state["auth_ok"] and st.session_state["auth_user"]:
            u = st.session_state["auth_user"]
            st.success(f"Sesión activa: {u.get('name', u.get('email',''))}")
            if u.get("photo_url"):
                st.sidebar.image(u["photo_url"], width=48)
            if st.button("Cerrar sesión", key="logout_btn_google"):
                st.session_state["auth_ok"] = False
                st.session_state["auth_user"] = None
                st.session_state["google_creds"] = None
                st.rerun()
            return

        # Validación de secrets (local)
        if "google_oauth" not in st.secrets:
            st.error("Falta configuración en .streamlit/secrets.toml: [google_oauth].")
            st.code("""[google_oauth]
client_id = "..."
client_secret = "..."
project_id = "..."
redirect_uri = "http://localhost:8501"
allowed_domain = ""  # opcional: "rogii.com"
""")
            return

        client_config = {
            "web": {
                "client_id": st.secrets["google_oauth"]["client_id"],
                "project_id": st.secrets["google_oauth"].get("project_id", ""),
                "auth_uri": "https://accounts.google.com/o/oauth2/auth",
                "token_uri": "https://oauth2.googleapis.com/token",
                "client_secret": st.secrets["google_oauth"]["client_secret"],
                "redirect_uris": [st.secrets["google_oauth"]["redirect_uri"]],
            }
        }

        flow = Flow.from_client_config(
            client_config,
            scopes=GOOGLE_SCOPES,
            redirect_uri=st.secrets["google_oauth"]["redirect_uri"],
        )

        auth_url, _ = flow.authorization_url(
            access_type="offline",
            include_granted_scopes="true",
            prompt="consent",
        )

        st.markdown(f"[➡️ Iniciar sesión con Google]({auth_url})")

        q = st.query_params
        code = q.get("code", None)
        if code:
            try:
                flow.fetch_token(code=code)
                creds = flow.credentials

                st.session_state["google_creds"] = {
                    "token": creds.token,
                    "refresh_token": creds.refresh_token,
                    "token_uri": creds.token_uri,
                    "client_id": creds.client_id,
                    "client_secret": creds.client_secret,
                    "scopes": creds.scopes,
                }

                r = requests.get(
                    "https://www.googleapis.com/oauth2/v2/userinfo",
                    headers={"Authorization": f"Bearer {creds.token}"},
                    timeout=20,
                )
                info = r.json() if r is not None else {}
                if not info.get("picture"):
                    r2 = requests.get(
                        "https://www.googleapis.com/oauth2/v3/userinfo",
                        headers={"Authorization": f"Bearer {creds.token}"},
                        timeout=20,
                    )
                    info2 = r2.json() if r2 is not None else {}
                    info = {**info2, **info}
                email = (info.get("email") or "").lower()

                allowed_domain = st.secrets["google_oauth"].get("allowed_domain", "")
                if allowed_domain:
                    allowed_domain = allowed_domain.lstrip("@").lower()
                    if not email.endswith("@" + allowed_domain):
                        st.error(f"Solo se permite acceso con dominio @{allowed_domain}")
                        st.session_state["auth_ok"] = False
                        st.session_state["auth_user"] = None
                        st.session_state["google_creds"] = None
                        return

                st.session_state["auth_ok"] = True
                pic = info.get("picture") or ""
                if pic:
                    if "googleusercontent.com" in pic and "=" not in pic:
                        pic = f"{pic}=s96-c"
                    if pic.startswith("http://"):
                        pic = "https://" + pic[len("http://"):]
                role = _get_user_role(email)
                st.session_state["auth_user"] = {
                    "name": info.get("name") or email,
                    "email": email,
                    "photo_url": pic,
                    "role": role,
                    "username": email,
                }

                st.query_params.clear()
                st.rerun()

            except Exception as e:
                st.error(f"No se pudo completar login: {e}")


# (El login legacy por usuario/contraseña fue removido: acceso solo por Google OAuth)

# ---------- TEMPORAL: bypass Google para pruebas en local ----------
# Poner en True para trabajar en local sin login Google. False = exigir login con Google.
BYPASS_GOOGLE_FOR_LOCAL = True

# ---------- Gate de acceso ----------
_google_oauth_login_sidebar()

if BYPASS_GOOGLE_FOR_LOCAL and not st.session_state.get("auth_ok"):
    st.session_state["auth_ok"] = True
    st.session_state["auth_user"] = {
        "name": "Usuario Local",
        "email": "local@local",
        "username": "local@local",
        "role": "user",
        "photo_url": "",
    }

if not st.session_state.get("auth_ok"):
    st.title("Dashboard Diario Operativo – DrillSpot / ROGII")
    st.info("Inicia sesión en el panel lateral para continuar.")
    st.stop()

# Mensaje de carga inicial (la app es pesada; así se ve que está respondiendo)
_loading_placeholder = None
if BYPASS_GOOGLE_FOR_LOCAL and st.session_state.get("auth_ok"):
    _loading_placeholder = st.empty()
    with _loading_placeholder.container():
        st.info("Cargando dashboard… (primera vez puede tardar 15–30 s)")

# Registrar usuario activo (latido simple)
try:
    if st.session_state.get("auth_user"):
        _touch_active_user(st.session_state["auth_user"])
except Exception:
    pass

# Badge usuario (foto + nombre)
try:
    if st.session_state.get("auth_user"):
        st.markdown(_render_user_badge(st.session_state["auth_user"]), unsafe_allow_html=True)
except Exception:
    pass

if _loading_placeholder is not None:
    _loading_placeholder.empty()

# --- Modo visual (forzar claro/oscuro independiente del theme de Streamlit) ---
# Esto controla los "cards" (HTML/iframes) y algunos estilos pro. No afecta cálculos.
if "ui_mode" not in st.session_state:
    # Si ya existe un turno (p.ej. BHA), úsalo como default. Si no, Diurno.
    st.session_state["ui_mode"] = st.session_state.get("turno", "Diurno")

with st.sidebar:
    st.radio("Modo visual", ["Diurno", "Nocturno"], key="ui_mode", horizontal=True)

with st.sidebar.expander("🟢 Usuarios activos", expanded=False):
    au = _cleanup_active_users(_load_active_users(), ttl_minutes=45)
    if not au:
        st.caption("Sin usuarios activos detectados.")
    else:
        users_sorted = sorted(au.values(), key=lambda x: str(x.get("last_seen", "")), reverse=True)
        for u in users_sorted[:12]:
            name = str(u.get("name") or u.get("email") or "").strip()
            role = str(u.get("role") or "").strip()
            last_seen = str(u.get("last_seen", "")).replace("T", " ")
            photo = str(u.get("photo_url") or "").strip()
            row = st.container()
            cols = row.columns([1, 5])
            if photo:
                cols[0].image(photo, width=28)
            else:
                cols[0].markdown("🧑")
            cols[1].markdown(f"**{name}**{f' · {role}' if role else ''}  \n`{last_seen}`")

# ------------------------------
# RUTAS (PC LOCAL)  ✅ AJUSTA ESTO
# ------------------------------
LOGO_PATH = r"C:\Users\l.brito_rogii\Downloads\DrillingOP_APP\ROGII_DINAMIC.gif"
TNPI_CSV_PATH = r"C:\Users\l.brito_rogii\Downloads\DrillingOP_APP\Detalles causas de TNPI.csv"

# ------------------------------
# ESTILO GLOBAL (HEADER PRO + UTILIDADES)
# ------------------------------
_light_mode_hdr = False
try:
    _ui_mode_hdr = st.session_state.get("ui_mode")
    if _ui_mode_hdr in ("Diurno", "Nocturno"):
        _light_mode_hdr = _ui_mode_hdr == "Diurno"
except Exception:
    _light_mode_hdr = False
if not _light_mode_hdr:
    try:
        _base_hdr = st.get_option("theme.base")
        _light_mode_hdr = str(_base_hdr).lower() == "light"
    except Exception:
        _light_mode_hdr = False

if _light_mode_hdr:
    _hdr_bg = "linear-gradient(180deg, rgba(255,255,255,0.96), rgba(244,247,251,0.98))"
    _hdr_border = "rgba(15,23,42,0.10)"
    _hdr_shadow = "0 18px 40px rgba(15,23,42,0.12)"
    _hdr_title = "#0f172a"
    _hdr_sub = "#475569"
    _hdr_status_bg = "rgba(15,23,42,0.04)"
    _hdr_status_border = "rgba(15,23,42,0.12)"
    _hdr_status_text = "#0f172a"
    _hdr_logo_bg = "rgba(15,23,42,0.04)"
    _hdr_logo_border = "rgba(15,23,42,0.10)"
    _hdr_logo_shadow = "inset 0 0 0 1px rgba(15,23,42,0.04)"
    _hdr_chip_border = "rgba(15,23,42,0.12)"
    _hdr_chip_shadow = "0 8px 16px rgba(15,23,42,0.12)"
else:
    _hdr_bg = "radial-gradient(1200px 240px at 20% -20%, rgba(40,180,99,0.22), transparent 60%), radial-gradient(1200px 240px at 80% 0%, rgba(46,134,193,0.22), transparent 55%), linear-gradient(180deg, rgba(18,18,20,0.95), rgba(8,8,10,0.96))"
    _hdr_border = "rgba(255,255,255,0.08)"
    _hdr_shadow = "0 18px 50px rgba(0,0,0,0.40)"
    _hdr_title = "rgba(255,255,255,0.95)"
    _hdr_sub = "rgba(255,255,255,0.72)"
    _hdr_status_bg = "rgba(255,255,255,0.06)"
    _hdr_status_border = "rgba(255,255,255,0.10)"
    _hdr_status_text = "rgba(255,255,255,0.88)"
    _hdr_logo_bg = "rgba(255,255,255,0.04)"
    _hdr_logo_border = "rgba(255,255,255,0.08)"
    _hdr_logo_shadow = "inset 0 0 0 1px rgba(255,255,255,0.02)"
    _hdr_chip_border = "rgba(255,255,255,0.10)"
    _hdr_chip_shadow = "0 8px 16px rgba(0,0,0,0.35)"

st.markdown(
    """
    <style>
      /* Quita margen arriba del main */
      .block-container {{ padding-top: 1.1rem; }}

      /* Header card */
      .ds-header {{
        border-radius: 22px;
        padding: 18px 20px;
        background: {hdr_bg};
        border: 1px solid {hdr_border};
        box-shadow: {hdr_shadow};
        display:flex;
        gap: 16px;
        align-items:center;
      }}
      .ds-logo-wrap{{
        width:64px;height:64px;border-radius:18px;
        background: {hdr_logo_bg};
        border: 1px solid {hdr_logo_border};
        display:flex;align-items:center;justify-content:center;
        box-shadow: {hdr_logo_shadow};
        overflow:hidden;
      }}
      .ds-logo {{
  width: 90px;
  height: auto;
  max-height: 70px;
  margin-right: 16px;
}}

.ds-logo.no-float {{
  animation: none !important;
}}

      @keyframes dsFloat{{
        0%{{ transform: translateY(0px) scale(1.00); }}
        50%{{ transform: translateY(-3px) scale(1.03); }}
        100%{{ transform: translateY(0px) scale(1.00); }}
      }}
      .ds-title{{
        font-size: 34px;
        font-weight: 900;
        line-height: 1.05;
        margin: 0;
        color: {hdr_title};
        letter-spacing: 0.2px;
      }}
      .ds-sub{{
        margin-top: 6px;
        color: {hdr_sub};
        font-size: 14px;
        font-weight: 600;
      }}

      /* Estado del día (pill) + glow dinámico por eficiencia */
      .ds-header {{ position: relative; overflow: hidden; }}
      .ds-header::after{{
        content:"";
        position:absolute; inset:-2px;
        background: radial-gradient(700px 260px at 12% 0%, var(--ds-glow, rgba(46,134,193,0.18)), transparent 60%),
                    radial-gradient(900px 260px at 88% 10%, var(--ds-glow2, rgba(40,180,99,0.18)), transparent 55%);
        pointer-events:none;
      }}
      .ds-header[data-status="ok"]{{ --ds-glow: rgba(40,180,99,0.22); --ds-glow2: rgba(46,134,193,0.18); }}
      .ds-header[data-status="warn"]{{ --ds-glow: rgba(241,196,15,0.22); --ds-glow2: rgba(46,134,193,0.14); }}
      .ds-header[data-status="crit"]{{ --ds-glow: rgba(231,76,60,0.28); --ds-glow2: rgba(241,196,15,0.12); }}

      .ds-status{{
        display:inline-flex; align-items:center; gap:8px;
        padding: 6px 10px;
        border-radius: 999px;
        border: 1px solid {hdr_status_border};
        background: {hdr_status_bg};
        color: {hdr_status_text};
        font-weight: 800;
        font-size: 12px;
        letter-spacing: 0.2px;
      }}
      .ds-status b{{ font-weight: 950; }}
      .ds-status .chip{{
        width:10px;height:10px;border-radius:999px;
        border: 2px solid {hdr_chip_border};
        box-shadow: {hdr_chip_shadow};
      }}
    </style>
    """.format(
        hdr_bg=_hdr_bg,
        hdr_border=_hdr_border,
        hdr_shadow=_hdr_shadow,
        hdr_title=_hdr_title,
        hdr_sub=_hdr_sub,
        hdr_status_bg=_hdr_status_bg,
        hdr_status_border=_hdr_status_border,
        hdr_status_text=_hdr_status_text,
        hdr_logo_bg=_hdr_logo_bg,
        hdr_logo_border=_hdr_logo_border,
        hdr_logo_shadow=_hdr_logo_shadow,
        hdr_chip_border=_hdr_chip_border,
        hdr_chip_shadow=_hdr_chip_shadow,
    ),
    unsafe_allow_html=True,
)

# ------------------------------
# HELPERS: base64 (PNG/GIF) para HTML
# ------------------------------
def file_to_base64(path: str) -> str:
    with open(path, "rb") as f:
        return base64.b64encode(f.read()).decode()


def mime_from_path(path: str) -> str:
    ext = os.path.splitext(path.lower())[1]
    if ext == ".gif":
        return "image/gif"
    if ext in [".jpg", ".jpeg"]:
        return "image/jpeg"
    return "image/png"


logo_b64 = ""
logo_mime = "image/png"
if LOGO_PATH and os.path.exists(LOGO_PATH):
    try:
        logo_b64 = file_to_base64(LOGO_PATH)
        logo_mime = mime_from_path(LOGO_PATH)
    except Exception:
        logo_b64 = ""
        logo_mime = "image/png"

# ------------------------------
# CONSTANTES
# ------------------------------
EXPORT_COLORWAY = ["#2E86C1", "#28B463", "#E74C3C", "#F1C40F", "#8E44AD", "#16A085"]
EQUIPO_TIPO = ["3000HP / AE", "2000HP"]
MODO_REPORTE_OPTS = ["Perforación", "Cambio de etapa"]
TIPO_AGUJERO = ["Entubado", "Descubierto"]
# Listas separadas: BNA = barrenas, Casing = etapas
BARRERAS_DEFAULT = ['36"', '26"', '18 1/2"', '17 1/2"', '16"', '14 1/2"', '13 1/2"', '12 1/4"', '10 5/8"', '8 1/2"', '6 1/4"']
SECCIONES_DEFAULT = ['30"', '20"', '16"', '13 3/8"', '11 3/4"', '9 5/8"', '7"', '5"']
TURNOS = ["Diurno", "Nocturno"]
# Límite de horas del día (al llegar aquí se considera día lleno). Turnos: 12h cada uno.
DAY_LIMIT_HOURS = 24.0
TURNO_LIMIT_HOURS = 12.0

ACTIVIDADES = [
    "Perforación",
    "Circula",
    "Rebaja cemento",
    "Prueba hermeticidad TR",
    "Instala UAP",
    "Desplaza",
    "Mantenimiento",
    "Succiona contrapozos",
    "Instala brida en cabezal",
    "Cambio de bombas",
    "Verifica parámetros",
    "Comandos fuera de la conexión",
    "Repaso fuera de la conexión",
    "Fallas",
    "Arma/Desarma BHA",
    "Conexión perforando",
    "Mete/levanta TR 30",
    "Mete/levanta TR 20",
    "Mete/levanta TR 16",
    "Mete/levanta TR 13 3/8",
    "Mete/levanta LN / TR (Lingadas)",
    "Mete/levanta LN / TR (TxT)",

    # Viajes (TRIPS)
    "Viaje metiendo con Pistolas",
    "Viaje sacando con pistolas",
    "Viaje metiendo con pescante",
    "Viaje levantando con pescante (asumiendo que se realizó la operación de pesca)",
    "Viaje inspeccionando roscas",
    "Viaje procedimiento quemado roscas nuevas",
    "Viaje de TLC",
    "Viaje metiendo con cuchara",
    "Viaje levantando con cuchara",
    "Viaje levantando/Metiendo TP de suelo natural",
    "Viaje levantando núcleo",
    "Viaje metiendo retenedor/PBR",
    "Viaje metiendo/levantando aplicando contrapresión (MPD)",
    "Viaje metiendo/levantando Alineados a MPD sin aplicar contrapresión",
    "Viaje levantando con tubería llena",
    "Viaje metiendo y sacando con conexión a top Drive (rotación y bombeo)",
    "Viaje metiendo y sacando con conexión a top Drive (rotación y bombeo, MPD)",
    "Viaje con conexión reductores de fricción / removedores de recortes (cada dos lingadas)",
    "Viaje con conexión usando llaves de fuerza",
    "Viaje con Calibración interna de TP",
    "Viaje Tramos dobles",
    "Viaje levantando empacador",
    "Viaje metiendo / levantando Aparejo doble",
    "Viaje metiendo / levantando aparejo de producción",
    "Viaje metiendo TP lingadas",
    "Viaje metiendo TP TxT",
    "Viaje levantando TP lingadas",
    "Viaje levantando TP TxT",
    "Viaje metiendo / levantando TP de 3 1/2\" - 2 7/8\" por lingadas",
    "Viaje metiendo / levantando TP de 3 1/2\" - 2 7/8\" TxT",
]


# ----------------------------------------------------------------------
# Catálogo de actividades para "Cambio de etapa" (CE)
# - Si existe el archivo 'actividades CE.csv' junto al script, se carga de ahí.
# - Si no existe, se usa un fallback mínimo.
# ----------------------------------------------------------------------
def _load_actividades_ce():
    fallback = [
        "Circular", "Bombear Bache", "Sacar sarta", "Eliminar BHA",
        "Instalar equipos para Introduccion TR /LN", "Bajar TR /LN",
        "Desmantelar equipo para introducción de TR", "Instalacion de equipo de cementacion",
        "Cementar", "Esperar fraguado", "WOC / Fraguado", "Prueba de presión", "NPT / Espera"
    ]
    try:
        _csv_candidates = [
            os.path.join(os.path.dirname(__file__), "actividades CE.csv"),
            os.path.join(os.getcwd(), "actividades CE.csv"),
        ]
        for _p in _csv_candidates:
            if os.path.exists(_p):
                _df = pd.read_csv(_p, encoding="latin-1")
                col = _df.columns[0]
                vals = [str(x).strip() for x in _df[col].tolist() if str(x).strip() and str(x).strip().lower() != "nan"]
                # quitar duplicados preservando orden
                seen = set()
                out = []
                for v in vals:
                    k = v.lower()
                    if k not in seen:
                        seen.add(k)
                        out.append(v)
                return out if out else fallback
    except Exception:
        pass
    return fallback

ACTIVIDADES_CE = _load_actividades_ce()
# Catálogo de objetivos para Viajes (m/h y min por conexión)
# Nota: estos valores vienen de la tabla de objetivos (velocidad y tiempo de conexión)
VIAJE_CATALOG = {

    "Mete/levanta TR 30": {"vel_mh": 48.0, "tconn_min": 8.0},
    "Mete/levanta TR 20": {"vel_mh": 75.0, "tconn_min": 5.5},
    "Mete/levanta TR 16": {"vel_mh": 112.0, "tconn_min": 5.0},
    "Mete/levanta TR 13 3/8": {"vel_mh": 120.0, "tconn_min": 4.5},
    "Mete/levanta LN / TR (Lingadas)": {"vel_mh": 242.0, "tconn_min": 4.0},
    "Mete/levanta LN / TR (TxT)": {"vel_mh": 140.0, "tconn_min": 4.0},
    "Viaje metiendo con Pistolas": {"vel_mh": 476.0, "tconn_min": 2.0},
    "Viaje sacando con pistolas": {"vel_mh": 476.0, "tconn_min": 2.0},
    "Viaje metiendo con pescante": {"vel_mh": 336.0, "tconn_min": 2.0},
    "Viaje levantando con pescante (asumiendo que se realizó la operación de pesca)": {"vel_mh": 306.0, "tconn_min": 2.5},
    "Viaje inspeccionando roscas": {"vel_mh": 336.0, "tconn_min": 4.0},
    "Viaje procedimiento quemado roscas nuevas": {"vel_mh": 252.0, "tconn_min": 5.5},
    "Viaje de TLC": {"vel_mh": 308.0, "tconn_min": 3.5},
    "Viaje metiendo con cuchara": {"vel_mh": 224.0, "tconn_min": 2.3},
    "Viaje levantando con cuchara": {"vel_mh": 224.0, "tconn_min": 2.0},
    "Viaje levantando/Metiendo TP de suelo natural": {"vel_mh": 252.0, "tconn_min": 5.0},
    "Viaje levantando núcleo": {"vel_mh": 364.0, "tconn_min": 2.5},
    "Viaje metiendo retenedor/PBR": {"vel_mh": 364.0, "tconn_min": 2.0},
    "Viaje metiendo/levantando aplicando contrapresión (MPD)": {"vel_mh": 252.0, "tconn_min": 4.0},
    "Viaje metiendo/levantando Alineados a MPD sin aplicar contrapresión": {"vel_mh": 430.0, "tconn_min": 2.0},
    "Viaje levantando con tubería llena": {"vel_mh": 476.0, "tconn_min": 2.5},
    "Viaje metiendo y sacando con conexión a top Drive (rotación y bombeo)": {"vel_mh": 252.0, "tconn_min": 5.0},
    "Viaje metiendo y sacando con conexión a top Drive (rotación y bombeo, MPD)": {"vel_mh": 196.0, "tconn_min": 7.0},
    "Viaje con conexión reductores de fricción / removedores de recortes (cada dos lingadas)": {"vel_mh": 210.0, "tconn_min": 7.0},
    "Viaje con conexión usando llaves de fuerza": {"vel_mh": 430.0, "tconn_min": 2.9},
    "Viaje con Calibración interna de TP": {"vel_mh": 470.0, "tconn_min": 2.3},
    "Viaje Tramos dobles": {"vel_mh": 250.0, "tconn_min": 2.9},
    "Viaje levantando empacador": {"vel_mh": 364.0, "tconn_min": 2.0},
    "Viaje metiendo / levantando Aparejo doble": {"vel_mh": 75.0, "tconn_min": 3.8},
    "Viaje metiendo / levantando aparejo de producción": {"vel_mh": 124.0, "tconn_min": 3.8},
    "Viaje metiendo TP lingadas": {"vel_mh": 640.0, "tconn_min": 1.5},
    "Viaje metiendo TP TxT": {"vel_mh": 192.0, "tconn_min": 2.0},
    "Viaje levantando TP lingadas": {"vel_mh": 732.0, "tconn_min": 1.5},
    "Viaje levantando TP TxT": {"vel_mh": 219.0, "tconn_min": 2.0},
    "Viaje metiendo / levantando TP de 3 1/2\" - 2 7/8\" por lingadas": {"vel_mh": 458.0, "tconn_min": 2.9},
    "Viaje metiendo / levantando TP de 3 1/2\" - 2 7/8\" TxT": {"vel_mh": 156.0, "tconn_min": 2.9},
}

# Conexiones
CONN_COMPONENTS = [
    "Preconexión",
    "Conexión",
    "Postconexión",
    "Repaso",
    "Survey",
    "Comandos RSS",
    "Bache",
    "Presión reducida",
]
CONN_COLOR_MAP = {
    "Repaso": "#7F8C8D",
    "Preconexión": "#F9E79F",
    "Conexión": "#00A8E8",
    "Postconexión": "#D5DBDB",
    "Bache": "#48C9B0",
    "Survey": "#5B2C6F",
    "Comandos RSS": "#E67E22",
    "Presión reducida": "#85C1E9",
}
CONN_ORDER = [
    "Repaso", "Preconexión", "Survey", "Conexión",
    "Postconexión", "Bache", "Comandos RSS", "Presión reducida"
]

CONN_TYPE_OPTS = ["Fondo a fondo", "Fondo a fondo con MPD"]
ANGLE_BUCKETS = ["<30°", "30° - 60°", ">60°"]

CONN_STDS = {
    ("Fondo a fondo", "<30°"): {"Preconexión": 5, "Conexión": 5, "Postconexión": 5, "TOTAL": 15},
    ("Fondo a fondo", "30° - 60°"): {"Preconexión": 12, "Conexión": 5, "Postconexión": 5, "TOTAL": 22},
    ("Fondo a fondo", ">60°"): {"Preconexión": 25, "Conexión": 5, "Postconexión": 5, "TOTAL": 35},
    ("Fondo a fondo con MPD", "<30°"): {"Preconexión": 8, "Conexión": 7, "Postconexión": 5, "TOTAL": 20},
    ("Fondo a fondo con MPD", "30° - 60°"): {"Preconexión": 12, "Conexión": 7, "Postconexión": 8, "TOTAL": 27},
    ("Fondo a fondo con MPD", ">60°"): {"Preconexión": 25, "Conexión": 7, "Postconexión": 8, "TOTAL": 40},
}

# BHA estándares -> (objetivo arma, objetivo desarma)
BHA_TYPES = {
    1:  ("Sarta lisa y/o Empacada y/o Péndulo", 4.0, 3.0),
    2:  ("Motor - Fondo/ MLPWD", 6.5, 5.0),
    3:  ("Rotatorio / MLPWD", 6.0, 4.5),
    4:  ("Rotatorio - MLPWD - Ampliador", 7.0, 5.5),
    5:  ("Rotatorio - MWD/LWD/PWD - Densidad Neutron (fuente radioactiva)/Sónico", 8.5, 6.5),
    6:  ("Sarta de limpieza, coronas, molinos, empacador de prueba y pescante", 3.5, 2.5),
    7:  ("Cucharas (Armado/Desarmado)", 4.5, 2.5),
    8:  ("Motor o Rotatorio - MWD/LWD/PWD - Densidad Neutrón/Sónico - 1 o más ampliador", 10.5, 7.5),
    9:  ("Sartas de Jetteo para aguas profundas (Casing / liner Drilling)", 3.0, 3.5),
    10: ("Equipo de Producción/Disparos/Toma de Registros y Operaciones Terminación", 10.0, 10.0),
}

# ------------------------------
# ACRÓNIMOS y casing
# ------------------------------
ACRONYMS = {"TR", "TP", "TNPI", "TNP", "BHA", "VCP", "WITS", "MPD", "AE", "RT", "BOP", "ROP", "RSS", "CRT", "GWD", "MLPWD", "MWD", "LWD", "PWD"}

def smart_case(text: str) -> str:
    if text is None:
        return ""
    t = str(text).strip()
    if t == "":
        return ""
    base = t[:1].upper() + t[1:].lower()
    out = base
    for a in sorted(ACRONYMS, key=len, reverse=True):
        out = re.sub(rf"\b{re.escape(a.lower())}\b", a, out, flags=re.IGNORECASE)
    out = out.replace("Tnpi", "TNPI").replace("Tnp", "TNP").replace("Tp", "TP").replace("Rop", "ROP")
    return out

def clamp_0_100(x: float) -> float:
    try:
        return max(0.0, min(float(x), 100.0))
    except Exception:
        return 0.0

def safe_pct(num: float, den: float) -> float:
    return (num / den * 100.0) if den and den > 0 else 0.0

def semaforo_color(v_0_100: float) -> str:
    v = clamp_0_100(v_0_100)
    if v >= 85:
        return "#2ECC71"
    if v >= 75:
        return "#F1C40F"
    return "#E74C3C"

def status_from_eff(eff: float) -> tuple[str, str, str]:
    """returns (status_key, label, color_hex)"""
    e = clamp_0_100(eff)
    if e >= 85:
        return ("ok", "OK", "#2ECC71")
    if e >= 75:
        return ("warn", "ATENCIÓN", "#F1C40F")
    return ("crit", "CRÍTICO", "#E74C3C")


# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# HELPERS: FECHAS (histórico diario / comparativo)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
def _df_fecha_to_date(s: pd.Series) -> pd.Series:
    """Coerce Fecha to datetime.date (accepts date/datetime/str)."""
    if pd.api.types.is_datetime64_any_dtype(s):
        return s.dt.date
    return pd.to_datetime(s, errors="coerce").dt.date

def _available_days(df: pd.DataFrame) -> list[date]:
    if df is None or df.empty or "Fecha" not in df.columns:
        return []
    d = _df_fecha_to_date(df["Fecha"])
    return sorted([x for x in d.dropna().unique().tolist()])

def split_day(df: pd.DataFrame, day: date, date_col: str = "Fecha") -> pd.DataFrame:
    """Return rows of df that match the given date."""
    if df is None or df.empty or date_col not in df.columns:
        return pd.DataFrame(columns=df.columns if df is not None else [])
    tmp = df.copy()
    tmp["_Fecha_dt"] = pd.to_datetime(tmp[date_col], errors="coerce")
    return tmp[tmp["_Fecha_dt"].dt.date == day].drop(columns=["_Fecha_dt"], errors="ignore").copy()


# ------------------------------
# TNPI catálogo
# ------------------------------
@st.cache_data(show_spinner=False)
def load_tnpi_catalog(csv_path: str) -> pd.DataFrame:
    """Carga catálogo de causas.
    Soporta CSVs con:
      - 2 columnas (Categoria_TNPI, Detalle_TNPI)
      - 4 columnas (Categoria_TNPI, Detalle_TNPI, Categoria_TNP, Detalle_TNP)
    Si no vienen columnas de TNP, las crea vacías (o replica TNPI como fallback).
    """
    cols_4 = ["Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP", "Detalle_TNP"]
    cols_2 = ["Categoria_TNPI", "Detalle_TNPI"]

    if csv_path and os.path.exists(csv_path):
        try:
            df = robust_read_csv(csv_path)
        except Exception:
            df = pd.read_csv(csv_path, encoding="utf-8", errors="ignore")

        # Normaliza número de columnas
        if df.shape[1] >= 4:
            df = df.iloc[:, :4].copy()
            df.columns = cols_4
        elif df.shape[1] == 2:
            df = df.iloc[:, :2].copy()
            df.columns = cols_2
            # crea columnas TNP (fallback a TNPI para no romper UI)
            df["Categoria_TNP"] = df["Categoria_TNPI"]
            df["Detalle_TNP"] = df["Detalle_TNPI"]
        elif df.shape[1] == 3:
            df = df.iloc[:, :3].copy()
            df.columns = ["Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP"]
            df["Detalle_TNP"] = ""
        else:
            # vacío / inesperado
            df = pd.DataFrame(columns=cols_4)

        for c in cols_4:
            if c not in df.columns:
                df[c] = ""
            df[c] = df[c].fillna("").astype(str).str.strip()

        return df

    # fallback mínimo
    return pd.DataFrame(columns=cols_4)

def render_export_diario_calendario():
    # -----------------------------------------------------------------
    # EXPORT AUTOMÁTICO DIARIO (por calendario) - PDF / PPTX / CSV
    # -----------------------------------------------------------------
    def _safe_float(v, default=0.0) -> float:
        try:
            if v is None:
                return float(default)
            return float(v)
        except Exception:
            return float(default)

    with st.expander("Export automático diario (calendario)", expanded=False):
        df_base = st.session_state.get("df", pd.DataFrame()).copy()
        days_all = _available_days(df_base)
        if len(days_all) == 0:
            st.info("Aún no hay datos con Fecha para exportar reportes diarios.")
        else:
            dia_exp = st.date_input("Día a exportar", value=days_all[-1], min_value=days_all[0], max_value=days_all[-1], key="exp_day_pick")
            # Alcance del reporte diario
            scope_rep = st.radio(
                "Alcance del reporte",
                ["Por pozo (todas las etapas del día)", "Por etapa (solo etapa seleccionada)"],
                index=0,
                horizontal=True,
                key="exp_scope_pick",
            )
            df_day = split_day(df_base, dia_exp, date_col="Fecha")
            if (not df_day.empty) and (scope_rep.startswith("Por etapa")):
                etapas_dia = [e for e in df_day.get("Etapa", pd.Series(dtype=str)).fillna("").astype(str).unique().tolist() if e != ""]
                etapa_default = st.session_state.get("etapa_val", "")
                if etapa_default in etapas_dia:
                    idx_def = etapas_dia.index(etapa_default)
                else:
                    idx_def = 0
                etapa_pick = st.selectbox("Etapa a exportar", options=etapas_dia if etapas_dia else [""], index=idx_def if etapas_dia else 0, key="exp_etapa_pick")
                if etapa_pick:
                    df_day = df_day[df_day["Etapa"].astype(str) == str(etapa_pick)].copy()

            if df_day.empty:
                st.warning("No hay registros para ese día.")
            else:
                # Meta (toma la primera fila que encuentre)
                r0 = df_day.iloc[0]
                meta_d = {
                    "equipo": str(r0.get("Equipo", st.session_state.get("equipo_val", "")) or ""),
                    "pozo": str(r0.get("Pozo", st.session_state.get("pozo_val", "")) or ""),
                    "etapa": str(r0.get("Etapa", "") or ""),
                    "fecha": dia_exp.isoformat(),
                }

                total = float(df_day.get("Horas_Reales", pd.Series(dtype=float)).fillna(0).sum())
                tp = float(df_day[df_day.get("Tipo", "") == "TP"]["Horas_Reales"].sum()) if "Tipo" in df_day.columns else total
                tnpi = float(df_day[df_day.get("Tipo", "") == "TNPI"]["Horas_Reales"].sum()) if "Tipo" in df_day.columns else 0.0
                tnp = float(df_day[df_day.get("Tipo", "") == "TNP"]["Horas_Reales"].sum()) if "Tipo" in df_day.columns else 0.0
                eff = clamp_0_100(safe_pct(tp, total)) if total > 0 else 0.0

                kpis_d = {
                    "TP (h)": f"{tp:.2f}",
                    "TNPI (h)": f"{tnpi:.2f}",
                    "TNP (h)": f"{tnp:.2f}",
                    "Horas total (h)": f"{total:.2f}",
                    "Eficiencia del día": f"{eff:.0f}%",
                }

                sig_day = f"{dia_exp.isoformat()}|{scope_rep}|{meta_d.get('etapa','')}|{len(df_day)}|{total:.2f}|{tp:.2f}|{tnpi:.2f}|{tnp:.2f}"
                if st.session_state.get("exp_day_sig") != sig_day:
                    st.session_state["exp_day_sig"] = sig_day
                    st.session_state.pop("exp_day_pdf", None)
                    st.session_state.pop("exp_day_ppt", None)
                    st.session_state.pop("exp_day_csv", None)

                if st.session_state.get("exp_day_pdf") is None or st.session_state.get("exp_day_ppt") is None:
                    st.caption("Para acelerar la interfaz, genera los archivos bajo demanda.")
                    if st.button("Preparar exportables (Día)", use_container_width=True, key="exp_day_prepare"):
                        with st.spinner("Generando exportables del día..."):
                            prog = st.progress(0)
                            prog_msg = st.empty()
                            prog_msg.caption("Iniciando...")
                            charts_d = {}
                            # Pie tiempos
                            if "Tipo" in df_day.columns and "Horas_Reales" in df_day.columns:
                                df_t = df_day.groupby("Tipo", as_index=False)["Horas_Reales"].sum()
                                if not df_t.empty:
                                    charts_d["TP vs TNPI vs TNP (Diario)"] = px.pie(
                                        df_t,
                                        names="Tipo",
                                        values="Horas_Reales",
                                        hole=0.55,
                                        title=f"TP vs TNPI vs TNP - {dia_exp.isoformat()}",
                                    )
                            prog.progress(12)
                            prog_msg.caption("Graficas de tiempos listas.")
                            # Pie actividades
                            if "Actividad" in df_day.columns and "Horas_Reales" in df_day.columns:
                                df_a = df_day.groupby("Actividad", as_index=False)["Horas_Reales"].sum().sort_values("Horas_Reales", ascending=False).head(10)
                                if not df_a.empty:
                                    charts_d["Top actividades (Diario)"] = px.pie(
                                        df_a,
                                        names="Actividad",
                                        values="Horas_Reales",
                                        hole=0.35,
                                        title=f"Top actividades - {dia_exp.isoformat()}",
                                    )
                            prog.progress(24)
                            prog_msg.caption("Graficas de actividades listas.")
                            # BHA (Arma/Desarma) - diario
                            df_bha_all = st.session_state.get("df_bha", pd.DataFrame()).copy()
                            if not df_bha_all.empty and "Fecha" in df_bha_all.columns:
                                df_bha_all["Fecha"] = df_bha_all["Fecha"].astype(str)
                                df_bha_d = df_bha_all[df_bha_all["Fecha"] == str(dia_exp)].copy()
                                if scope_rep.startswith("Por etapa") and meta_d.get("etapa"):
                                    if "Etapa" in df_bha_d.columns:
                                        df_bha_d = df_bha_d[df_bha_d["Etapa"] == str(meta_d.get("etapa"))].copy()
                                if not df_bha_d.empty:
                                    df_long_bha = df_bha_d.melt(
                                        id_vars=[c for c in ["BHA_Tipo", "Accion"] if c in df_bha_d.columns],
                                        value_vars=[c for c in ["Estandar_h", "Real_h"] if c in df_bha_d.columns],
                                        var_name="Serie",
                                        value_name="Horas",
                                    )
                                    if not df_long_bha.empty:
                                        fig_bha_d = px.bar(
                                            df_long_bha,
                                            x="BHA_Tipo" if "BHA_Tipo" in df_long_bha.columns else "Accion",
                                            y="Horas",
                                            color="Serie",
                                            barmode="group",
                                            title=f"BHA - {dia_exp.isoformat()}",
                                            color_discrete_sequence=EXPORT_COLORWAY,
                                        )
                                        charts_d["BHA (Estándar vs Real)"] = fig_bha_d
                            prog.progress(38)
                            prog_msg.caption("Graficas BHA listas.")

                            # ROP diario (Día vs Noche): por etapa o consolidado por pozo
                            if modo_reporte == "Perforación":
                                rop_prog_d = 0.0
                                rop_rd = 0.0
                                rop_rn = 0.0
                                por_etapa = st.session_state.drill_day.get("por_etapa", {})
                                if scope_rep.startswith("Por etapa") and meta_d.get("etapa"):
                                    etapa_key = str(meta_d.get("etapa"))
                                    etapa_data_rop_d = por_etapa.get(etapa_key, {})
                                    _prog_map = etapa_data_rop_d.get("rop_prog_by_date", {}) or {}
                                    _rd_map = etapa_data_rop_d.get("rop_real_dia_by_date", {}) or {}
                                    _rn_map = etapa_data_rop_d.get("rop_real_noche_by_date", {}) or {}
                                    _p_entry = _prog_map.get(str(dia_exp), {})
                                    rop_prog_d = _safe_float(_p_entry.get("rop_prog") if isinstance(_p_entry, dict) else (_p_entry or 0.0))
                                    rop_rd = _safe_float(_rd_map.get(str(dia_exp), 0.0) or 0.0)
                                    rop_rn = _safe_float(_rn_map.get(str(dia_exp), 0.0) or 0.0)
                                else:
                                    # Consolidado por pozo: sumar por etapa si hay datos diarios
                                    for _, etapa_data_rop_d in (por_etapa or {}).items():
                                        _prog_map = etapa_data_rop_d.get("rop_prog_by_date", {}) or {}
                                        _rd_map = etapa_data_rop_d.get("rop_real_dia_by_date", {}) or {}
                                        _rn_map = etapa_data_rop_d.get("rop_real_noche_by_date", {}) or {}
                                        _p_entry = _prog_map.get(str(dia_exp), {})
                                        rop_prog_d += _safe_float(_p_entry.get("rop_prog") if isinstance(_p_entry, dict) else (_p_entry or 0.0))
                                        rop_rd += _safe_float(_rd_map.get(str(dia_exp), 0.0) or 0.0)
                                        rop_rn += _safe_float(_rn_map.get(str(dia_exp), 0.0) or 0.0)
                                if (rop_prog_d + rop_rd + rop_rn) > 0:
                                    df_rop_d = pd.DataFrame(
                                        [
                                            {"Turno": "Día ☀️", "Programado (m/h)": rop_prog_d, "Real (m/h)": rop_rd},
                                            {"Turno": "Noche 🌙", "Programado (m/h)": rop_prog_d, "Real (m/h)": rop_rn},
                                        ]
                                    )
                                    fig_rop_d = px.bar(
                                        df_rop_d,
                                        x="Turno",
                                        y=["Programado (m/h)", "Real (m/h)"],
                                        barmode="group",
                                        text_auto=True,
                                        title=f"ROP - {dia_exp.isoformat()}",
                                        color_discrete_sequence=EXPORT_COLORWAY,
                                    )
                                    charts_d["ROP (Diario)"] = fig_rop_d
                            prog.progress(52)
                            prog_msg.caption("Graficas ROP listas.")

                            # Metros perforados diarios (Real vs Programado): por etapa o consolidado
                            if modo_reporte == "Perforación":
                                mp_d = 0.0
                                mr_d = 0.0
                                mr_n = 0.0
                                por_etapa = st.session_state.drill_day.get("por_etapa", {})
                                if scope_rep.startswith("Por etapa") and meta_d.get("etapa"):
                                    etapa_key = str(meta_d.get("etapa"))
                                    etapa_data_m = por_etapa.get(etapa_key, {})
                                    _mp_map = etapa_data_m.get("metros_prog_by_date", {}) or {}
                                    _md_map = etapa_data_m.get("metros_real_dia_by_date", {}) or {}
                                    _mn_map = etapa_data_m.get("metros_real_noche_by_date", {}) or {}
                                    _mp_entry = _mp_map.get(str(dia_exp), {})
                                    mp_d = _safe_float(_mp_entry.get("metros_prog") if isinstance(_mp_entry, dict) else (_mp_entry or 0.0))
                                    mr_d = _safe_float(_md_map.get(str(dia_exp), 0.0) or 0.0)
                                    mr_n = _safe_float(_mn_map.get(str(dia_exp), 0.0) or 0.0)
                                else:
                                    for _, etapa_data_m in (por_etapa or {}).items():
                                        _mp_map = etapa_data_m.get("metros_prog_by_date", {}) or {}
                                        _md_map = etapa_data_m.get("metros_real_dia_by_date", {}) or {}
                                        _mn_map = etapa_data_m.get("metros_real_noche_by_date", {}) or {}
                                        _mp_entry = _mp_map.get(str(dia_exp), {})
                                        mp_d += _safe_float(_mp_entry.get("metros_prog") if isinstance(_mp_entry, dict) else (_mp_entry or 0.0))
                                        mr_d += _safe_float(_md_map.get(str(dia_exp), 0.0) or 0.0)
                                        mr_n += _safe_float(_mn_map.get(str(dia_exp), 0.0) or 0.0)
                                mr_t = mr_d + mr_n
                                if (mp_d + mr_d + mr_n) > 0:
                                    df_m_d = pd.DataFrame(
                                        [
                                            {"Tipo": "Programado (total)", "Metros (m)": mp_d},
                                            {"Tipo": "Real Día ☀️", "Metros (m)": mr_d},
                                            {"Tipo": "Real Noche 🌙", "Metros (m)": mr_n},
                                            {"Tipo": "Real Total", "Metros (m)": mr_t},
                                        ]
                                    )
                                    fig_m_d = px.bar(
                                        df_m_d,
                                        x="Tipo",
                                        y="Metros (m)",
                                        text_auto=True,
                                        title=f"Metros - {dia_exp.isoformat()}",
                                        color="Tipo",
                                        color_discrete_map={
                                            "Programado (total)": "#6B7280",
                                            "Real Día ☀️": "#F59E0B",
                                            "Real Noche 🌙": "#1D4ED8",
                                            "Real Total": "#22C55E",
                                        },
                                    )
                                    charts_d["Metros perforados (Diario)"] = fig_m_d
                            prog.progress(66)
                            prog_msg.caption("Graficas de metros listas.")

                            # Conexiones perforando (diario)
                            df_conn_all = st.session_state.get("df_conn", pd.DataFrame()).copy()
                            if not df_conn_all.empty and "Fecha" in df_conn_all.columns:
                                df_conn_all["Fecha"] = df_conn_all["Fecha"].astype(str)
                                df_conn_d = df_conn_all[df_conn_all["Fecha"] == str(dia_exp)].copy()
                                if scope_rep.startswith("Por etapa") and meta_d.get("etapa") and "Etapa" in df_conn_d.columns:
                                    df_conn_d = df_conn_d[df_conn_d["Etapa"] == str(meta_d.get("etapa"))].copy()
                                if not df_conn_d.empty and {"Componente", "Minutos_Reales"}.issubset(df_conn_d.columns):
                                    df_conn_sum = df_conn_d.groupby("Componente", as_index=False)["Minutos_Reales"].sum()
                                    df_conn_sum["Componente"] = pd.Categorical(df_conn_sum["Componente"], categories=CONN_ORDER, ordered=True)
                                    df_conn_sum = df_conn_sum.sort_values("Componente")
                                    charts_d["Conexiones (Distribución)"] = px.pie(
                                        df_conn_sum,
                                        names="Componente",
                                        values="Minutos_Reales",
                                        hole=0.35,
                                        title=f"Conexiones - {dia_exp.isoformat()}",
                                        color="Componente",
                                        color_discrete_map=CONN_COLOR_MAP,
                                    )

                                    df_stack = df_conn_d.copy()
                                    df_stack["Conn_Label"] = df_stack["Profundidad_m"].fillna(df_stack["Conn_No"]).astype(float).astype(int).astype(str)
                                    df_stack["Componente"] = pd.Categorical(df_stack["Componente"], categories=CONN_ORDER, ordered=True)
                                    df_stack_g = df_stack.groupby(["Conn_Label", "Componente"], as_index=False)["Minutos_Reales"].sum().sort_values(["Conn_Label", "Componente"])
                                    fig_conn_stack = px.bar(
                                        df_stack_g,
                                        x="Conn_Label",
                                        y="Minutos_Reales",
                                        color="Componente",
                                        category_orders={"Componente": CONN_ORDER},
                                        color_discrete_map=CONN_COLOR_MAP,
                                        barmode="stack",
                                        title=f"Conexiones perforando - {dia_exp.isoformat()}",
                                        labels={"Conn_Label": "Profundidad (m)", "Minutos_Reales": "Tiempo (min)"},
                                    )
                                    charts_d["Conexiones perforando (Stack)"] = fig_conn_stack
                            prog.progress(78)
                            prog_msg.caption("Graficas de conexiones listas.")

                            # Viajes (si existen datos por hora)
                            viajes_store = st.session_state.get("viajes_hourly_store", {})
                            if isinstance(viajes_store, dict) and len(viajes_store) > 0:
                                for v_name, v_obj in viajes_store.items():
                                    hourly_df = v_obj.get("hourly") if isinstance(v_obj, dict) else None
                                    if isinstance(hourly_df, pd.DataFrame) and not hourly_df.empty:
                                        df_plot = hourly_df.copy().sort_values("hour").reset_index(drop=True)
                                        df_plot["hour_str"] = df_plot["hour"].astype(int)
                                        fig_v = px.bar(
                                            df_plot,
                                            x="hour_str",
                                            y="speed_mh",
                                            labels={"hour_str": "Hora", "speed_mh": "m/h"},
                                            title=f"Viaje – {v_name}",
                                        )
                                        fig_v.update_traces(marker_color=EXPORT_COLORWAY[0])
                                        charts_d[f"Viaje – Velocidad ({v_name})"] = fig_v
                                        fig_c = px.bar(
                                            df_plot,
                                            x="hour_str",
                                            y="conn_min",
                                            labels={"hour_str": "Hora", "conn_min": "min"},
                                            title=f"Viaje – Conexiones ({v_name})",
                                        )
                                        fig_c.update_traces(marker_color=EXPORT_COLORWAY[1] if len(EXPORT_COLORWAY) > 1 else EXPORT_COLORWAY[0])
                                        charts_d[f"Viaje – Conexiones ({v_name})"] = fig_c
                            prog.progress(86)
                            prog_msg.caption("Graficas de viajes listas.")

                            st.session_state["exp_day_pdf"] = build_pdf(meta_d, kpis_d, charts=charts_d)
                            prog.progress(93)
                            prog_msg.caption("PDF listo.")
                            st.session_state["exp_day_ppt"] = build_pptx(meta_d, kpis_d, charts=charts_d)
                            prog.progress(98)
                            prog_msg.caption("PowerPoint listo.")
                            st.session_state["exp_day_csv"] = df_day.to_csv(index=False).encode("utf-8")
                            prog.progress(100)
                            prog_msg.caption("CSV listo.")

                if st.session_state.get("exp_day_pdf") is not None:
                    colx1, colx2, colx3 = st.columns(3)
                    with colx1:
                        st.download_button(
                            "Descargar PDF (Día)",
                            data=st.session_state.get("exp_day_pdf"),
                            file_name=f"Reporte_Diario_{meta_d['pozo']}_{dia_exp.isoformat()}.pdf",
                            mime="application/pdf",
                            use_container_width=True,
                            key="dl_pdf_day",
                        )
                    with colx2:
                        st.download_button(
                            "Descargar PPTX (Día)",
                            data=st.session_state.get("exp_day_ppt"),
                            file_name=f"Reporte_Diario_{meta_d['pozo']}_{dia_exp.isoformat()}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                            key="dl_ppt_day",
                        )
                    with colx3:
                        st.download_button(
                            "Descargar CSV (Día)",
                            data=st.session_state.get("exp_day_csv"),
                            file_name=f"Datos_Diarios_{meta_d['pozo']}_{dia_exp.isoformat()}.csv",
                            mime="text/csv",
                            use_container_width=True,
                            key="dl_csv_day",
                        )

                with st.expander("Vista previa (tabla del día)", expanded=False):
                    st.dataframe(df_day, use_container_width=True, height=260)


def style_for_export(fig):
    if not PLOTLY_IMG_OK:
        return fig
    f = go.Figure(fig.to_dict())
    f.update_layout(
        template="plotly_white",
        paper_bgcolor="white",
        plot_bgcolor="white",
        font=dict(color="black", size=18),
        margin=dict(l=40, r=40, t=70, b=40),
        legend=dict(bgcolor="rgba(255,255,255,0.85)", borderwidth=0, font=dict(size=24)),
        title=dict(x=0.02),
        colorway=EXPORT_COLORWAY,
        uniformtext=dict(minsize=16, mode="show"),
    )
    f.update_xaxes(tickfont=dict(size=24), title_font=dict(size=24), automargin=True)
    f.update_yaxes(tickfont=dict(size=24), title_font=dict(size=24), automargin=True)
    f.update_traces(
        textfont=dict(size=24),
        insidetextfont=dict(size=26),
        outsidetextfont=dict(size=24),
        selector=dict(type="pie"),
    )
    legend_items = {t.name for t in f.data if getattr(t, "name", None)}
    if len(legend_items) >= 6:
        f.update_layout(
            legend=dict(orientation="h", yanchor="bottom", y=-0.25, xanchor="center", x=0.5, font=dict(size=24)),
            margin=dict(l=40, r=40, t=70, b=130),
        )
    return f

def plotly_to_png_bytes(fig) -> bytes | None:
    if not PLOTLY_IMG_OK:
        return None
    try:
        fig_export = style_for_export(fig)
        png = pio.to_image(fig_export, format="png", width=1800, height=1000, scale=2)
        im = Image.open(BytesIO(png)).convert("RGBA")
        bg = Image.new("RGBA", im.size, (255, 255, 255, 255))
        bg.paste(im, (0, 0), im)
        out = bg.convert("RGB")
        b = BytesIO()
        out.save(b, format="PNG", optimize=True)
        return b.getvalue()
    except Exception:
        return None

def build_pdf(meta: dict, kpis: dict, charts: dict) -> bytes:
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter

    def write_text(txt, y, size=12, bold=False):
        c.setFont("Helvetica-Bold" if bold else "Helvetica", size)
        c.drawString(0.75 * inch, y, txt)
        return y - 0.26 * inch

    def write_chart(fig, y, title):
        img_bytes = plotly_to_png_bytes(fig)
        if img_bytes is None:
            y = write_text(f"{title} (gráfica no disponible: instala kaleido)", y, size=10, bold=False)
            return y
        y = write_text(title, y, size=12, bold=True)
        img_h = 3.1 * inch
        img_w = width - 1.5 * inch
        y_img = y - img_h

        if y_img < 0.75 * inch:
            c.showPage()
            y = height - 0.75 * inch
            y = write_text(title, y, bold=True)
            y_img = y - img_h

        img_reader = ImageReader(BytesIO(img_bytes))
        c.drawImage(
            img_reader,
            0.75 * inch,
            y_img,
            width=img_w,
            height=img_h,
            preserveAspectRatio=True,
            mask=None,
        )
        return y_img - 0.25 * inch

    y = height - 0.75 * inch
    y = write_text("Reporte DrillSpot / ROGII", y, size=18, bold=True)
    y = write_text(f"Equipo: {meta.get('equipo','')}", y)
    y = write_text(f"Pozo: {meta.get('pozo','')}", y)
    y = write_text(f"Etapa: {meta.get('etapa','')}", y)
    y = write_text(f"Fecha: {meta.get('fecha','')}", y)
    y -= 0.1 * inch

    y = write_text("KPIs", y, size=14, bold=True)
    for k, v in kpis.items():
        y = write_text(f"- {k}: {v}", y, size=11)
        if y < 1.0 * inch:
            c.showPage()
            y = height - 0.75 * inch

    if charts:
        c.showPage()
        y = height - 0.75 * inch
        y = write_text("Gráficas", y, size=14, bold=True)
        for name, fig in charts.items():
            y = write_chart(fig, y, name)
            if y < 1.0 * inch:
                c.showPage()
                y = height - 0.75 * inch

    c.save()
    buffer.seek(0)
    return buffer.getvalue()

def build_pptx(meta: dict, kpis: dict, charts: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Reporte DrillSpot / ROGII"
    slide.placeholders[1].text = (
        f"Equipo: {meta.get('equipo','')} | Pozo: {meta.get('pozo','')} | "
        f"Etapa: {meta.get('etapa','')} | Fecha: {meta.get('fecha','')}"
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "KPIs"
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.2))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, (k, v) in enumerate(kpis.items()):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"{k}: {v}"
        p.font.size = Pt(18)

    for title, fig in (charts or {}).items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        img_bytes = plotly_to_png_bytes(fig)
        if img_bytes is None:
            slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.0)).text_frame.text = (
                "No se pudo embebir imagen (instala kaleido)."
            )
        else:
            slide.shapes.add_picture(BytesIO(img_bytes), Inches(0.8), Inches(1.4), width=Inches(11.6))

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()

# ------------------------------
# Gauge principal
# ------------------------------
def build_gauge(title: str, value_0_100: float):
    if not PLOTLY_IMG_OK:
        return None
    light_mode = _is_light_theme()
    text_color = "#0f172a" if light_mode else "white"
    tick_color = "rgba(15,23,42,0.35)" if light_mode else "rgba(255,255,255,0.35)"
    v = clamp_0_100(value_0_100)
    # Color dinámico por rango (pro look)
    if v >= 85:
        bar_color = "#22c55e"
        delta_color = "#22c55e"
    elif v >= 75:
        bar_color = "#f59e0b"
        delta_color = "#f59e0b"
    else:
        bar_color = "#ef4444"
        delta_color = "#ef4444"
    _sk, status_label, status_color = status_from_eff(v)
    fig = go.Figure(
        go.Indicator(
            mode="gauge+number+delta",
            value=v,
            number={"suffix": "%", "font": {"size": 58, "family": "Arial Black", "color": text_color}},
            delta={
                "reference": 85,
                "increasing": {"color": delta_color},
                "decreasing": {"color": delta_color},
                "position": "bottom",
                "valueformat": ".0f",
                "prefix": "Δ ",
                "suffix": " vs 85%",
            },
            title={"text": title.upper(), "font": {"size": 20, "family": "Arial Black", "color": text_color}},
            gauge={
                "axis": {
                    "range": [0, 100],
                    "tickwidth": 1,
                    "tickcolor": tick_color,
                    "tickvals": [0, 50, 100],
                    "ticktext": ["0", "50", "100"],
                },
                "bar": {"thickness": 0.34, "color": bar_color},
                "steps": [
                    {"range": [0, 75], "color": "rgba(239,68,68,0.28)"},
                    {"range": [75, 85], "color": "rgba(245,158,11,0.28)"},
                    {"range": [85, 100], "color": "rgba(34,197,94,0.28)"},
                ],
                "threshold": {"line": {"color": "#8b5cf6", "width": 4}, "thickness": 0.78, "value": 85},
            },
        )
    )
    fig.update_layout(
        height=360,
        margin=dict(l=20, r=20, t=60, b=10),
        paper_bgcolor="rgba(0,0,0,0)",
        font=dict(color=text_color),
    )
    fig.add_annotation(
        text=f"<b>{status_label}</b>",
        x=0.5,
        y=0.05,
        xref="paper",
        yref="paper",
        showarrow=False,
        font=dict(size=16, color=status_color, family="Arial Black"),
    )
    return fig

# ------------------------------
# HTML PRO: CSS embebido
# ------------------------------

def _is_light_theme() -> bool:
    """Determina si debemos renderizar en modo claro.

    Prioridad:
    1) st.session_state['ui_mode'] (Diurno/Nocturno) - controla el look de los cards pro.
    2) theme.base de Streamlit.
    """
    try:
        ui_mode = st.session_state.get("ui_mode")
        if ui_mode in ("Diurno", "Nocturno"):
            return ui_mode == "Diurno"
    except Exception:
        pass

    try:
        base = st.get_option("theme.base")
        return str(base).lower() == "light"
    except Exception:
        return False


def _pro_iframe_css(light: bool = False) -> str:
    """CSS base for 'pro' embedded tables/cards (used in HTML iframes)."""
    if light:
        bg = "#ffffff"
        card = "#ffffff"
        border = "#e5e7eb"
        text = "#111827"
        muted = "#6b7280"
        row_hover = "#f3f4f6"
        header = "#f9fafb"
        shadow = "0 8px 18px rgba(0,0,0,.10)"
        track = "#e5e7eb"
    else:
        bg = "#0b0f14"
        card = "#0f1620"
        border = "#223043"
        text = "#e6edf3"
        muted = "#9aa7b2"
        row_hover = "#132033"
        header = "#0c121b"
        shadow = "0 8px 22px rgba(0,0,0,.35)"
        track = "#223043"

    return f"""
    <style>
      :root {{
        --bg: {bg};
        --card: {card};
        --border: {border};
        --text: {text};
        --muted: {muted};
        --row-hover: {row_hover};
        --header: {header};
        --shadow: {shadow};
        --track: {track};
      }}
      html, body {{
        margin:0; padding:0;
        font-family: ui-sans-serif, system-ui, -apple-system, Segoe UI, Roboto, Arial;
        background: transparent;
        color: var(--text);
      }}
      .wrap {{
        background: var(--bg);
        border: 1px solid var(--border);
        border-radius: 18px;
        padding: 18px 18px 14px 18px;
        box-shadow: var(--shadow);
      }}
      .ds-card {{
        background: var(--card);
        border: 1px solid var(--border);
        border-radius: 18px;
        padding: 18px 18px 14px 18px;
        box-shadow: var(--shadow);
      }}
      .title {{
        font-size: 28px;
        font-weight: 800;
        letter-spacing: .2px;
        margin: 4px 0 14px;
      }}
      .sub {{
        color: var(--muted);
        margin-top: -10px;
        margin-bottom: 14px;
        font-size: 14px;
      }}
      table {{
        width: 100%;
        border-collapse: collapse;
        border-spacing: 0;
        overflow: hidden;
        border-radius: 14px;
      }}
      thead th {{
        text-align: left;
        font-size: 14px;
        color: var(--muted);
        font-weight: 700;
        padding: 12px 14px;
        background: var(--header);
        border-bottom: 1px solid var(--border);
      }}
      tbody td {{
        padding: 12px 14px;
        border-bottom: 1px solid var(--border);
        font-size: 16px;
      }}
      /* Numeric columns alignment (used by the activity indicator table) */
      .ds-name {{
        text-align: left;
      }}
      .ds-num {{
        text-align: right;
        font-variant-numeric: tabular-nums;
      }}
      tbody tr:hover td {{
        background: var(--row-hover);
      }}
      .kpi {{
        font-weight: 800;
        font-size: 26px;
      }}
      .pill {{
        display:inline-flex;
        align-items:center;
        gap:10px;
      }}
      .dot {{
        width: 14px; height: 14px;
        border-radius: 50%;
        display:inline-block;
      }}
      .dot.red {{ background: #ef4444; box-shadow: 0 0 0 4px rgba(239,68,68,.18); }}
      .dot.ylw {{ background: #f59e0b; box-shadow: 0 0 0 4px rgba(245,158,11,.18); }}
      .dot.grn {{ background: #22c55e; box-shadow: 0 0 0 4px rgba(34,197,94,.18); }}
      .bar {{
        width: 290px;
        height: 14px;
        border-radius: 999px;
        background: var(--track);
        overflow: hidden;
      }}
      .bar > span {{
        display:block;
        height: 100%;
        border-radius: 999px;
      }}
      .barwrap {{
        display: flex;
        align-items: center;
        justify-content: flex-end;
        gap: 10px;
      }}
      .pct {{
        min-width: 40px;
        text-align: right;
        color: var(--muted);
        font-weight: 700;
      }}
    </style>
    """
def kpi_table_html(rows: list[dict]) -> str:
    def dot(color, pulse=False, tooltip=""):
        cls = "dot pulse" if pulse else "dot"
        tt = f' title="{tooltip}"' if tooltip else ""
        return f'<span class="{cls}" style="background:{color};"{tt}></span>'

    tr = ""
    for r in rows:
        eff = clamp_0_100(r.get("eff", 0))
        color = semaforo_color(eff)
        pulse = eff < 75
        tooltip = "Eficiencia < 75% (revisar TNPI / causas)" if pulse else ""
        tr += f"""
        <tr>
          <td class="ds-name">{r.get('kpi','')}</td>
          <td class="ds-num">{r.get('real','')}</td>
          <td class="ds-num">{r.get('tnpi','')}</td>
          <td class="ds-num">{r.get('tnp','')}</td>
          <td class="ds-num">{eff:.0f} {dot(color, pulse=pulse, tooltip=tooltip)}</td>
        </tr>
        """

    return f"""
    {_pro_iframe_css(light=_is_light_theme())}
    <div class="ds-card">
      <div style="font-size:26px;font-weight:900;color:var(--text);margin:2px 0 10px 0;">
        Indicador de desempeño
      </div>
      <table class="ds-t">
        <thead>
          <tr>
            <th>KPI</th>
            <th style="text-align:right;">Real</th>
            <th style="text-align:right;">TNPI</th>
            <th style="text-align:right;">TNP</th>
            <th style="text-align:right;">Eficiencia (%)</th>
          </tr>
        </thead>
        <tbody>{tr}</tbody>
      </table>
      <div style="margin-top:10px;color:var(--muted);font-size:13px;font-weight:700;display:flex;gap:18px;align-items:center;">
        <span><span class="dot" style="background:#E74C3C;"></span> &nbsp;&lt; 75%</span>
        <span><span class="dot" style="background:#F1C40F;"></span> &nbsp;75–85%</span>
        <span><span class="dot" style="background:#2ECC71;"></span> &nbsp;&ge; 85%</span>
      </div>
    </div>
    """

def indicators_table_html(title: str, rows: list[dict], kind: str = "actividad") -> str:
    def dot(color, pulse=False, tooltip=""):
        cls = "dot pulse" if pulse else "dot"
        tt = f' title="{tooltip}"' if tooltip else ""
        return f'<span class="{cls}" style="background:{color};"{tt}></span>'

    th_name = "Actividad" if kind == "actividad" else "Conexión"
    th_real = "Real (h)" if kind == "actividad" else "Real (min)"
    th_tnpi = "TNPI (h)" if kind == "actividad" else "TNPI (min)"
    th_tnp = "TNP (h)" if kind == "actividad" else "TNP (min)"

    tr = ""
    for r in rows:
        eff = clamp_0_100(r.get("eff", 0))
        color = semaforo_color(eff)
        pulse = eff < 75
        tooltip = "Eficiencia < 75% (revisar TNPI / causas)" if pulse else ""
        width = max(0, min(int(round(eff)), 100))

        tr += f"""
        <tr>
          <td class="ds-name">{r.get("name","")}</td>
          <td class="ds-num">{r.get("real","")}</td>
          <td class="ds-num">{r.get("tnpi","")}</td>
          <td class="ds-num">{r.get("tnp","")}</td>
          <td class="ds-num">
            <div class="barwrap">
              <div class="bar"><span style="width:{width}%; background:{color};"></span></div>
              <div class="pct">{eff:.0f}%</div>
            </div>
          </td>
          <td>{dot(color, pulse=pulse, tooltip=tooltip)}</td>
        </tr>
        """

    return f"""
    {_pro_iframe_css(light=_is_light_theme())}
    <div class="ds-card">
      <div style="font-size:34px;font-weight:950;color:var(--text);margin:4px 0 12px 0;">
        {title}
      </div>
      <table class="ds-t">
        <thead>
          <tr>
            <th>{th_name}</th>
            <th style="text-align:right;">{th_real}</th>
            <th style="text-align:right;">{th_tnpi}</th>
            <th style="text-align:right;">{th_tnp}</th>
            <th style="text-align:right;">Eficiencia (%)</th>
            <th>Semáforo</th>
          </tr>
        </thead>
        <tbody>{tr}</tbody>
      </table>
    </div>
    """

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# SESSION STATE INIT (ANTES del header preview!)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# --- FIX: inicializar depth_rows para Viajes ---
if "depth_rows" not in st.session_state:
    st.session_state.depth_rows = pd.DataFrame(
        columns=["Etapa", "PT_programada_m", "PT_actual_m"]
    )

if "df" not in st.session_state:
    st.session_state.df = pd.DataFrame(
        columns=["RowID",
        "Equipo", "Pozo", "Etapa", "Fecha", "Equipo_Tipo", "Modo_Reporte",
            "Seccion", "Corrida", "Tipo_Agujero", "Operacion", "Actividad", "Turno",
            "Tipo", "Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP", "Detalle_TNP",
            "Hora_Inicio", "Hora_Fin",
            "Horas_Prog", "Horas_Reales",
            "ROP_Prog_mh", "ROP_Real_mh",
            "Comentario", "Origen", "BHA_ID", "CONN_ID", "VIAJE_TIPO",
        ]
    )

if "df_conn" not in st.session_state:
    st.session_state.df_conn = pd.DataFrame(
        columns=[
            "Equipo", "Pozo", "Etapa", "Fecha", "Equipo_Tipo", "Seccion", "Corrida",
            "Tipo_Agujero", "Turno", "Conn_No", "Profundidad_m",
            "Conn_Tipo", "Angulo_Bucket",
            "Componente", "Minutos_Reales", "Minutos_Estandar", "Minutos_TNPI",
            "Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP", "Detalle_TNP", "Comentario", "CONN_ID",
        ]
    )

if "df_bha" not in st.session_state:
    st.session_state.df_bha = pd.DataFrame(
        columns=[
            "Equipo", "Pozo", "Etapa", "Fecha", "Turno",
            "Barrena", "BHA_Tipo", "BHA_Componentes", "Accion", "BHA_ID",
            "Estandar_h", "Real_h", "TNPI_h", "TNP_h", "Eficiencia_pct"
        ]
    )

if "drill_day" not in st.session_state:
    st.session_state.drill_day = {
        # Datos globales (para compatibilidad)
        "metros_prog_total": 0.0,
        "metros_real_dia": 0.0,
        "metros_real_noche": 0.0,
        "rop_prog_total": 0.0,
        "rop_real_dia": 0.0,
        "rop_real_noche": 0.0,
        "tnpi_metros_h": 0.0,
        "pt_programada_m": 0.0,
        "prof_actual_m": 0.0,

        # Nuevo: Datos por etapa
        "por_etapa": {
            # Ejemplo:
            # "36'": {"pt_prog": 1000, "prof_actual": 500, "metros_prog": 200, ...},
            # "26'": {"pt_prog": 800, "prof_actual": 300, "metros_prog": 150, ...},
        }
    }

if "custom_actividades" not in st.session_state:
    st.session_state.custom_actividades = []


# --- FIX: desglose de TNPI para BHA (múltiples causas) ---
if "bha_tnpi_breakdown" not in st.session_state:
    st.session_state.bha_tnpi_breakdown = []  # lista de dicts: {Categoria_TNPI, Detalle_TNPI, Horas_TNPI_h, Comentario}

if "bha_tnpi_breakdown_draft" not in st.session_state:
    st.session_state.bha_tnpi_breakdown_draft = []  # borrador antes de "Guardar cambios"
if "bha_tnpi_breakdown_saved" not in st.session_state:
    st.session_state.bha_tnpi_breakdown_saved = False

# --- FIX: desglose de TNPI para actividades generales (múltiples causas) ---
if "act_tnpi_breakdown" not in st.session_state:
    st.session_state.act_tnpi_breakdown = []
if "act_tnpi_breakdown_draft" not in st.session_state:
    st.session_state.act_tnpi_breakdown_draft = []
if "act_tnpi_breakdown_saved" not in st.session_state:
    st.session_state.act_tnpi_breakdown_saved = False

# FUNCIÓN PARA OBTENER/ACTUALIZAR DATOS POR ETAPA (PONER JUSTO DESPUÉS)
def get_etapa_data(etapa_nombre):
    """Obtiene o crea los datos de una etapa específica"""
    if "por_etapa" not in st.session_state.drill_day:
        st.session_state.drill_day["por_etapa"] = {}

    if etapa_nombre not in st.session_state.drill_day["por_etapa"]:
        # Crear estructura inicial para la etapa
        st.session_state.drill_day["por_etapa"][etapa_nombre] = {
            "pt_programada_m": 0.0,
            "prof_actual_m": 0.0,

            # Metros / ROP diarios (último valor capturado)
            "metros_prog_total": 0.0,
            "metros_real_dia": 0.0,
            "metros_real_noche": 0.0,
            "rop_prog_total": 0.0,
            "rop_real_dia": 0.0,
            "rop_real_noche": 0.0,

            # Metas por etapa
            "rop_prog_etapa": 0.0,

            # Históricos por fecha (para acumulados / promedios por etapa)
            "metros_real_dia_by_date": {},
            "metros_real_noche_by_date": {},
            "rop_real_dia_by_date": {},
            "rop_real_noche_by_date": {},

            # ROP programada por corrida (maestro) y por fecha (registro diario)
            "rop_prog_by_corrida": {},
            "rop_prog_by_corrida_meta": {},
            "rop_prog_by_date": {},

            "tnpi_metros_h": 0.0,
        }

    return st.session_state.drill_day["por_etapa"][etapa_nombre]

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# HEADER PRO (preview eficiencia para glow/estado)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
_df_prev = st.session_state.df
_total_prev = float(_df_prev["Horas_Reales"].sum()) if not _df_prev.empty else 0.0
_tp_prev = float(_df_prev[_df_prev["Tipo"] == "TP"]["Horas_Reales"].sum()) if not _df_prev.empty else 0.0
_eff_prev = clamp_0_100(safe_pct(_tp_prev, _total_prev)) if _total_prev > 0 else 0.0
_status_key, _status_label, _status_color = status_from_eff(_eff_prev)

if logo_b64:
    st.markdown(
        f"""
        <div class="ds-header" data-status="{_status_key}">
          <div class="ds-logo-wrap">
            <img class="ds-logo" src="data:{logo_mime};base64,{logo_b64}" />
          </div>
          <div style="flex:1; position:relative; z-index:1;">
            <div class="ds-title">Dashboard Diario Operativo – DrillSpot / ROGII</div>
            <div class="ds-sub">Operational Report</div>
          </div>
          <div style="display:flex; flex-direction:column; gap:8px; align-items:flex-end; position:relative; z-index:1;">
            <div class="ds-status">
              <span class="chip" style="background:{_status_color};"></span>
              Estado del día: <b>{_status_label}</b>
            </div>
            <div class="ds-status" title="Eficiencia del día (TP / Real total)">
              <span class="chip" style="background:rgba(255,255,255,0.20);"></span>
              Eficiencia: <b>{_eff_prev:.0f}%</b>
            </div>
          </div>
        </div>
        """,
        unsafe_allow_html=True
    )
else:
    st.markdown(
        f"""
        <div class="ds-header" data-status="{_status_key}">
          <div style="flex:1; position:relative; z-index:1;">
            <div class="ds-title">Dashboard Diario Operativo – DrillSpot / ROGII</div>
            <div class="ds-sub">Operational Report</div>
          </div>
          <div style="display:flex; flex-direction:column; gap:8px; align-items:flex-end; position:relative; z-index:1;">
            <div class="ds-status">
              <span class="chip" style="background:{_status_color};"></span>
              Estado del día: <b>{_status_label}</b>
            </div>
            <div class="ds-status" title="Eficiencia del día (TP / Real total)">
              <span class="chip" style="background:rgba(255,255,255,0.20);"></span>
              Eficiencia: <b>{_eff_prev:.0f}%</b>
            </div>
          </div>
        </div>
        """,
        unsafe_allow_html=True
    )



# ------------------------------
# Toggle global (defínelo ANTES de usarlo en gráficos previos al sidebar)
# ------------------------------
show_charts = bool(st.session_state.get("show_charts", True))

# --- ROP Programado vs Real ---
# (Movido a la pestaña dedicada "ROP" para evitar duplicidad/confusión)

st.divider()


# ------------------------------
# MODO DE REPORTE (DEFAULT SEGURO)
# ------------------------------
# Se usa antes del sidebar (por el bloque Avance de profundidad).
modo_reporte = st.session_state.get("modo_reporte", MODO_REPORTE_OPTS[0])


# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# GUARDAR / CARGAR JORNADA (JSON local)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
def _default_jornada_path(equipo: str, pozo: str, fecha_str: str) -> str:
    safe = lambda s: re.sub(r"[^A-Za-z0-9_-]+", "_", str(s)).strip("_")
    script_dir = os.path.dirname(os.path.abspath(__file__)) if "__file__" in globals() else os.getcwd()
    return os.path.join(script_dir, f"jornada_{safe(equipo)}_{safe(pozo)}_{safe(fecha_str)}.json")

def _list_local_jornadas(limit: int = 60) -> list[tuple[str, str]]:
    script_dir = os.path.dirname(os.path.abspath(__file__)) if "__file__" in globals() else os.getcwd()
    try:
        files = []
        for f in os.listdir(script_dir):
            if f.lower().startswith("jornada_") and f.lower().endswith(".json"):
                full = os.path.join(script_dir, f)
                try:
                    mtime = datetime.fromtimestamp(os.path.getmtime(full))
                except Exception:
                    mtime = None
                files.append((full, mtime))
        files.sort(key=lambda x: x[1] or datetime.min, reverse=True)
        out = []
        for full, mtime in files[:limit]:
            label = os.path.basename(full)
            if mtime:
                label = f"{label} · {mtime.strftime('%Y-%m-%d %H:%M')}"
            out.append((label, full))
        return out
    except Exception:
        return []

def _kpi_summary_from_payload(payload: dict) -> dict:
    """Calcula KPIs básicos desde un payload de jornada."""
    try:
        meta = {}
        if isinstance(payload, dict):
            meta = payload.get("meta") or {}
            if not meta:
                drill_day = payload.get("drill_day") or {}
                meta = drill_day.get("meta") or {}
        rows = payload.get("df", []) if isinstance(payload, dict) else []
        df_k = pd.DataFrame(rows)
        if df_k.empty or "Horas_Reales" not in df_k.columns:
            return {}
        df_k["Horas_Reales"] = pd.to_numeric(df_k["Horas_Reales"], errors="coerce").fillna(0.0)
        total_h = float(df_k["Horas_Reales"].sum())
        tp_h = float(df_k[df_k.get("Tipo", "") == "TP"]["Horas_Reales"].sum()) if "Tipo" in df_k.columns else total_h
        tnpi_h = float(df_k[df_k.get("Tipo", "") == "TNPI"]["Horas_Reales"].sum()) if "Tipo" in df_k.columns else 0.0
        tnp_h = float(df_k[df_k.get("Tipo", "") == "TNP"]["Horas_Reales"].sum()) if "Tipo" in df_k.columns else 0.0
        eff = clamp_0_100(safe_pct(tp_h, total_h)) if total_h > 0 else 0.0
        return {
            "pozo": str(meta.get("pozo", "") or "").strip(),
            "etapa": str(meta.get("etapa", "") or meta.get("etapa_manual_val", "") or "").strip(),
            "fecha": str(meta.get("fecha", "") or "").strip(),
            "total_h": total_h,
            "tp_h": tp_h,
            "tnpi_h": tnpi_h,
            "tnp_h": tnp_h,
            "eff": eff,
        }
    except Exception:
        return {}

def _render_kpi_summary(k: dict, title: str = "Resumen KPI") -> None:
    if not k:
        st.sidebar.caption("Sin datos para resumen KPI.")
        return
    st.sidebar.markdown(f"**{title}**")
    meta_parts = []
    if k.get("pozo"):
        meta_parts.append(f"Pozo: {k['pozo']}")
    if k.get("etapa"):
        meta_parts.append(f"Etapa: {k['etapa']}")
    if k.get("fecha"):
        meta_parts.append(f"Fecha: {k['fecha']}")
    if meta_parts:
        st.sidebar.caption(" · ".join(meta_parts))
    def _bar(label: str, value: float, max_value: float, color: str) -> str:
        pct = 0.0 if max_value <= 0 else max(0.0, min(100.0, (value / max_value) * 100.0))
        return (
            "<div style='margin:6px 0 8px;'>"
            f"<div style='font-size:11px;opacity:.75;margin-bottom:4px;'>{label}: {value:.1f} h</div>"
            "<div style='height:6px;border-radius:999px;background:rgba(255,255,255,.08);overflow:hidden;'>"
            f"<div style='height:100%;width:{pct:.1f}%;background:{color};'></div>"
            "</div>"
            "</div>"
        )

    eff = float(k.get("eff", 0.0) or 0.0)
    eff_color = "#16a34a" if eff >= 70 else "#f59e0b" if eff >= 50 else "#ef4444"
    total_h = float(k.get("total_h", 0.0) or 0.0)
    tp_h = float(k.get("tp_h", 0.0) or 0.0)
    tnpi_h = float(k.get("tnpi_h", 0.0) or 0.0)
    tnp_h = float(k.get("tnp_h", 0.0) or 0.0)

    st.sidebar.markdown(
        "<div style='font-size:11px;opacity:.75;margin-top:4px;'>"
        f"Total: <b>{total_h:.1f} h</b> · Eficiencia: "
        f"<span style='color:{eff_color};font-weight:700;'>{eff:.0f}%</span>"
        "</div>",
        unsafe_allow_html=True,
    )
    st.sidebar.markdown(
        _bar("TP", tp_h, total_h, "#22c55e")
        + _bar("TNPI", tnpi_h, total_h, "#f59e0b")
        + _bar("TNP", tnp_h, total_h, "#ef4444"),
        unsafe_allow_html=True,
    )

def save_jornada_json(path_out: str) -> None:
    # Meta/contexto del sidebar (para reconstrucción confiable al cargar)
    meta = {
        "equipo": st.session_state.get("equipo_val", ""),
        "pozo": st.session_state.get("pozo_val", ""),
        "fecha": str(st.session_state.get("fecha_val", "")),
        "equipo_tipo": st.session_state.get("equipo_tipo_val", ""),
        "etapa_manual": bool(st.session_state.get("etapa_manual_chk", False)),
        "etapa": st.session_state.get("etapa_sel", ""),
        "etapa_manual_val": st.session_state.get("etapa_manual_val", ""),
        "modo_reporte": st.session_state.get("modo_reporte", ""),
        "show_charts": bool(st.session_state.get("show_charts", True)),
    }

    # También guardamos el meta dentro de drill_day para que quede autocontenido
    st.session_state.drill_day["meta"] = meta

    payload = {
        "version": "1.1",
        "saved_at": datetime.now().isoformat(timespec="seconds"),
        "meta": meta,
        "df": st.session_state.df.to_dict(orient="records"),
        "df_conn": st.session_state.df_conn.to_dict(orient="records"),
        "df_bha": st.session_state.df_bha.to_dict(orient="records"),
        "drill_day": st.session_state.drill_day,
        "custom_actividades": st.session_state.get("custom_actividades", []),
    }
    with open(path_out, "w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False, indent=2)

def _normalize_df_for_hash(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame()
    out = df.copy()
    for c in out.columns:
        if pd.api.types.is_numeric_dtype(out[c]):
            out[c] = pd.to_numeric(out[c], errors="coerce").round(6)
        out[c] = out[c].fillna("")
        out[c] = out[c].astype(str)
    return out

def _merge_df_rows(base: pd.DataFrame, incoming: pd.DataFrame) -> tuple[pd.DataFrame, int]:
    if incoming is None or incoming.empty:
        return base, 0
    if base is None or base.empty:
        return incoming.copy(), len(incoming)
    incoming = incoming.reindex(columns=base.columns)
    base_norm = _normalize_df_for_hash(base)
    inc_norm = _normalize_df_for_hash(incoming)
    base_hash = pd.util.hash_pandas_object(base_norm, index=False)
    inc_hash = pd.util.hash_pandas_object(inc_norm, index=False)
    keep_mask = ~inc_hash.isin(set(base_hash))
    added = int(keep_mask.sum())
    if added > 0:
        base = pd.concat([base, incoming.loc[keep_mask].copy()], ignore_index=True)
    return base, added

def _merge_dict_no_overwrite(base: dict, incoming: dict) -> dict:
    if not isinstance(base, dict):
        base = {}
    if not isinstance(incoming, dict):
        return base
    for k, v in incoming.items():
        if k not in base:
            base[k] = v
        elif isinstance(base.get(k), dict) and isinstance(v, dict):
            base[k] = _merge_dict_no_overwrite(base.get(k, {}), v)
    return base

def _filter_df_by_date(df: pd.DataFrame, fecha_sel) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame()
    if "Fecha" not in df.columns:
        return pd.DataFrame()
    fecha_str = str(fecha_sel)
    df_local = df.copy()
    df_local["Fecha"] = df_local["Fecha"].astype(str)
    return df_local[df_local["Fecha"] == fecha_str].copy()

def _day_used_hours(df: pd.DataFrame, fecha_sel) -> float:
    """Suma horas reales registradas en un día (según columna Fecha)."""
    if df is None or df.empty or "Fecha" not in df.columns:
        return 0.0
    try:
        fecha_str = str(fecha_sel)
        df_local = df.copy()
        df_local["Fecha"] = df_local["Fecha"].astype(str)
        df_local = df_local[df_local["Fecha"] == fecha_str]
        return float(pd.to_numeric(df_local.get("Horas_Reales", 0.0), errors="coerce").fillna(0.0).sum())
    except Exception:
        return 0.0

def _remaining_day_hours(df: pd.DataFrame, fecha_sel, day_limit: float = None) -> float:
    """Horas restantes disponibles para el día (cap en DAY_LIMIT_HOURS)."""
    if day_limit is None:
        day_limit = DAY_LIMIT_HOURS
    used = _day_used_hours(df, fecha_sel)
    try:
        return max(0.0, float(day_limit) - float(used))
    except Exception:
        return 0.0

def _day_used_hours_by_turno(df: pd.DataFrame, fecha_sel, turno_nombre: str) -> float:
    """Suma horas ya cargadas en el día para un turno ('Diurno' o 'Nocturno')."""
    if df is None or df.empty or "Fecha" not in df.columns or "Turno" not in df.columns:
        return 0.0
    try:
        fecha_str = str(fecha_sel)
        df_local = df.copy()
        df_local["Fecha"] = df_local["Fecha"].astype(str)
        df_local = df_local[df_local["Fecha"] == fecha_str]
        turno_col = df_local.get("Turno", pd.Series(dtype=str)).fillna("").astype(str).str.lower()
        is_diurno = (
            turno_col.str.contains("diurno", na=False) | turno_col.str.contains("dia", na=False)
            | turno_col.str.contains("día", na=False) | turno_col.str.contains("day", na=False)
            | turno_col.str.contains("☀", na=False)
        )
        is_nocturno = (
            turno_col.str.contains("nocturno", na=False) | turno_col.str.contains("noche", na=False)
            | turno_col.str.contains("night", na=False) | turno_col.str.contains("🌙", na=False)
        )
        tn = str(turno_nombre or "").strip().lower()
        if "diurno" in tn or "dia" in tn or "día" in tn:
            mask = is_diurno
        elif "nocturno" in tn or "noche" in tn:
            mask = is_nocturno
        else:
            return 0.0
        return float(pd.to_numeric(df_local.loc[mask, "Horas_Reales"], errors="coerce").fillna(0).sum())
    except Exception:
        return 0.0

def _sync_bha_from_df(df_main: pd.DataFrame, df_bha_in: pd.DataFrame) -> pd.DataFrame:
    """Sincroniza df_bha con las actividades BHA presentes en df principal."""
    if df_bha_in is None:
        return df_bha_in
    if df_main is None or df_main.empty:
        return df_bha_in.iloc[0:0].copy()

    out = df_bha_in.copy()

    # 1) Si existen IDs de BHA en el df principal, usarlo como fuente de verdad
    bha_ids = set()
    if "BHA_ID" in df_main.columns:
        bha_ids = set(
            df_main["BHA_ID"]
            .dropna()
            .astype(str)
            .str.strip()
            .tolist()
        )
        bha_ids = {x for x in bha_ids if x}
    if bha_ids and "BHA_ID" in out.columns:
        out["BHA_ID"] = out["BHA_ID"].astype(str)
        out = out[out["BHA_ID"].isin(bha_ids)]
        return out

    # 2) Fallback: match por actividad BHA en df principal
    if "Actividad" not in df_main.columns:
        return out
    df_bha_act = df_main[df_main["Actividad"].astype(str).str.contains("Arma/Desarma BHA", case=False, na=False)].copy()
    if df_bha_act.empty:
        return out.iloc[0:0].copy()

    def _extract_bha_tipo(act: str) -> str:
        try:
            m = re.search(r"Tipo\s*(\d+)", str(act))
            return m.group(1) if m else ""
        except Exception:
            return ""

    df_bha_act["_BHA_Tipo"] = df_bha_act["Actividad"].apply(_extract_bha_tipo)
    df_bha_act["_Fecha"] = df_bha_act.get("Fecha", "").astype(str)
    df_bha_act["_Etapa"] = df_bha_act.get("Etapa", "").astype(str)
    df_bha_act["_Turno"] = df_bha_act.get("Turno", "").astype(str)
    keys = set(
        tuple(x)
        for x in df_bha_act[["_Fecha", "_Etapa", "_Turno", "_BHA_Tipo"]]
        .dropna()
        .values.tolist()
    )

    if not keys:
        return out

    out["_Fecha"] = out.get("Fecha", "").astype(str)
    out["_Etapa"] = out.get("Etapa", "").astype(str)
    out["_Turno"] = out.get("Turno", "").astype(str)
    out["_BHA_Tipo"] = out.get("BHA_Tipo", "").astype(str)
    out["_key"] = list(zip(out["_Fecha"], out["_Etapa"], out["_Turno"], out["_BHA_Tipo"]))
    out = out[out["_key"].isin(keys)].copy()
    out.drop(columns=["_Fecha", "_Etapa", "_Turno", "_BHA_Tipo", "_key"], inplace=True, errors="ignore")
    return out

def _sync_conn_from_df(df_main: pd.DataFrame, df_conn_in: pd.DataFrame) -> pd.DataFrame:
    """Sincroniza df_conn con las conexiones presentes en df principal."""
    if df_conn_in is None:
        return df_conn_in
    if df_main is None or df_main.empty:
        return df_conn_in.iloc[0:0].copy()

    if "CONN_ID" not in df_main.columns or "CONN_ID" not in df_conn_in.columns:
        return df_conn_in

    conn_ids = set(
        df_main["CONN_ID"]
        .dropna()
        .astype(str)
        .str.strip()
        .tolist()
    )
    conn_ids = {x for x in conn_ids if x}
    if not conn_ids:
        return df_conn_in.iloc[0:0].copy()

    out = df_conn_in.copy()
    out["CONN_ID"] = out["CONN_ID"].astype(str)
    out = out[out["CONN_ID"].isin(conn_ids)]
    return out

def _build_day_payload(fecha_sel, autor: str = "") -> dict:
    meta = {
        "equipo": st.session_state.get("equipo_val", ""),
        "pozo": st.session_state.get("pozo_val", ""),
        "fecha": str(fecha_sel),
        "equipo_tipo": st.session_state.get("equipo_tipo_val", ""),
        "etapa": st.session_state.get("etapa_sel", ""),
        "autor": autor or "",
    }
    df_day = _filter_df_by_date(st.session_state.df, fecha_sel)
    df_conn_day = _filter_df_by_date(st.session_state.df_conn, fecha_sel)
    df_bha_day = _filter_df_by_date(st.session_state.df_bha, fecha_sel)

    drill_day_in = st.session_state.get("drill_day", {}) or {}
    drill_day_out = {"meta": drill_day_in.get("meta", {})}
    por_etapa_out = {}
    fecha_str = str(fecha_sel)
    for etapa_k, etapa_data in (drill_day_in.get("por_etapa", {}) or {}).items():
        if not isinstance(etapa_data, dict):
            continue
        data_out = {}
        for key in [
            "rop_prog_by_date", "rop_real_dia_by_date", "rop_real_noche_by_date",
            "metros_prog_by_date", "metros_real_dia_by_date", "metros_real_noche_by_date",
        ]:
            m = etapa_data.get(key, {})
            if isinstance(m, dict) and fecha_str in m:
                data_out[key] = {fecha_str: m.get(fecha_str)}
        for key in [
            "pt_programada_m", "prof_actual_m", "metros_prog_total",
            "metros_real_dia", "metros_real_noche", "rop_prog_total",
            "rop_real_dia", "rop_real_noche",
        ]:
            if key in etapa_data:
                data_out[key] = etapa_data.get(key)
        if data_out:
            por_etapa_out[str(etapa_k)] = data_out
    if por_etapa_out:
        drill_day_out["por_etapa"] = por_etapa_out

    return {
        "version": "1.0-day",
        "saved_at": datetime.now().isoformat(timespec="seconds"),
        "meta": meta,
        "df": df_day.to_dict(orient="records"),
        "df_conn": df_conn_day.to_dict(orient="records"),
        "df_bha": df_bha_day.to_dict(orient="records"),
        "drill_day": drill_day_out,
    }

def load_jornada_json(path_in: str) -> bool:
    if not path_in or not os.path.exists(path_in):
        return False
    with open(path_in, "r", encoding="utf-8") as f:
        payload = json.load(f)

    return _apply_jornada_payload(payload)

def _apply_jornada_payload(payload: dict) -> bool:
    try:
        meta = payload.get("meta") or {}
        # Fecha de la jornada que estamos cargando (para merge por día, no reemplazo total)
        _fecha_raw = str(meta.get("fecha", ""))
        payload_fecha_str = None
        for fmt in ("%Y-%m-%d", "%Y/%m/%d"):
            try:
                dt = datetime.strptime(_fecha_raw, fmt)
                payload_fecha_str = dt.strftime("%Y-%m-%d")
                break
            except Exception:
                pass
        if not payload_fecha_str:
            payload_fecha_str = _fecha_raw

        def _merge_jornada_table(current: pd.DataFrame, payload_rows: list, cols) -> pd.DataFrame:
            """Mantiene todos los días actuales; solo reemplaza el día de la jornada si el payload tiene más (o igual) datos; si en sesión hay más registros (ej. viajes recién agregados), se conservan."""
            new_day = pd.DataFrame(payload_rows, columns=cols) if payload_rows else pd.DataFrame(columns=cols)
            if current is None or current.empty:
                return new_day if not new_day.empty else pd.DataFrame(columns=cols)
            if "Fecha" not in current.columns:
                return new_day if not new_day.empty else current
            cur = current.copy()
            cur["Fecha"] = cur["Fecha"].astype(str).str.replace("/", "-")
            other = cur[cur["Fecha"] != payload_fecha_str]
            existing_for_date = cur[cur["Fecha"] == payload_fecha_str]
            # Si en sesión ya hay datos de ese día y tiene más (o igual) filas que el JSON, no reemplazar (evita perder viajes/actividades recién cargados al cambiar de día)
            if len(existing_for_date) > 0 and len(existing_for_date) >= len(new_day) and not new_day.empty:
                return current.reset_index(drop=True)
            # Si el payload viene vacío para ese día, conservar lo que hay en sesión
            if new_day.empty:
                return cur.reset_index(drop=True)
            new_day = new_day.copy()
            new_day["Fecha"] = payload_fecha_str
            return pd.concat([other, new_day], ignore_index=True).reset_index(drop=True)

        # Merge por fecha: se actualiza solo el día de la jornada, se conservan los demás días
        st.session_state.df = _merge_jornada_table(
            st.session_state.df, payload.get("df", []), st.session_state.df.columns
        )
        st.session_state.df_conn = _merge_jornada_table(
            st.session_state.df_conn, payload.get("df_conn", []), st.session_state.df_conn.columns
        )
        st.session_state.df_bha = _merge_jornada_table(
            st.session_state.df_bha, payload.get("df_bha", []), st.session_state.df_bha.columns
        )

        # drill_day + meta
        st.session_state.drill_day = payload.get("drill_day", st.session_state.drill_day) or st.session_state.drill_day
        meta = meta or st.session_state.drill_day.get("meta") or {}

        # Actividades personalizadas
        st.session_state.custom_actividades = payload.get("custom_actividades", []) or []
        # Restauración segura del sidebar: NO modificar keys de widgets después de instanciados.
        # Guardamos valores para aplicarlos al inicio del script (antes de render del sidebar) y forzamos rerun.
        if meta:
            pending = {
                'equipo_val': meta.get('equipo', ''),
                'pozo_val': meta.get('pozo', ''),
            }

            # fecha viene como string "YYYY-MM-DD" o "YYYY/MM/DD"
            _fecha_raw = str(meta.get('fecha', ''))
            _fecha = None
            for fmt in ('%Y-%m-%d', '%Y/%m/%d'):
                try:
                    _fecha = datetime.strptime(_fecha_raw, fmt).date()
                    break
                except Exception:
                    pass
            if _fecha is not None:
                pending['fecha_val'] = _fecha

            pending['equipo_tipo_val'] = meta.get('equipo_tipo', '')
            pending['etapa_manual_chk'] = bool(meta.get('etapa_manual', False))
            pending['etapa_sel'] = meta.get('etapa', meta.get('etapa_manual_val', ''))
            pending['etapa_manual_val'] = meta.get('etapa_manual_val', meta.get('etapa', ''))

            if 'modo_reporte' in meta:
                pending['modo_reporte'] = meta.get('modo_reporte', st.session_state.get('modo_reporte', ''))
            if 'show_charts' in meta:
                pending['show_charts'] = bool(meta.get('show_charts', True))

            st.session_state['_pending_sidebar_restore'] = pending
            # Mantener meta también dentro de drill_day
            st.session_state.drill_day['meta'] = meta

        return True
    except Exception as e:
        st.sidebar.error(f"No se pudo aplicar la jornada: {e}")
        return False

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# SIDEBAR (con modo presentación)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# --- RESTORE SEGURO: aplicar valores cargados ANTES de instanciar widgets ---
_pending = st.session_state.pop('_pending_sidebar_restore', None)
if isinstance(_pending, dict) and _pending:
    for _k, _v in _pending.items():
        try:
            st.session_state[_k] = _v
        except Exception:
            pass

st.sidebar.title("Panel de Control")
presentacion = st.sidebar.toggle("Modo presentación (ocultar sidebar)", value=False)

if presentacion:
    st.markdown(
        """
        <style>
          [data-testid="stSidebar"] { display: none; }
          .block-container { padding-top: 1.0rem; }
        </style>
        """,
        unsafe_allow_html=True
    )

with st.sidebar.container(border=True):
    st.sidebar.markdown("### Reporte")
    equipo = st.sidebar.text_input("Equipo", value=st.session_state.get("equipo_val","PM 2402"), key="equipo_val")
    equipo_tipo = st.sidebar.selectbox("Tipo de equipo", options=EQUIPO_TIPO, index=EQUIPO_TIPO.index(st.session_state.get("equipo_tipo_val", EQUIPO_TIPO[0])) if st.session_state.get("equipo_tipo_val", EQUIPO_TIPO[0]) in EQUIPO_TIPO else 0, key="equipo_tipo_val")
    pozo = st.sidebar.text_input("Pozo", value=st.session_state.get("pozo_val","OME 1 EXP"), key="pozo_val")
    # Etapa (sección) - lista + opción manual
    _default_etapa = st.session_state.get("etapa_sel", SECCIONES_DEFAULT[2])
    _opts_etapa = SECCIONES_DEFAULT + ["Otra (manual)"]
    _idx = _opts_etapa.index(_default_etapa) if _default_etapa in _opts_etapa else 2
    etapa_pick = st.sidebar.selectbox("Etapa", _opts_etapa, index=_idx, key="etapa_select")
    if etapa_pick == "Otra (manual)":
        etapa = st.sidebar.text_input("Etapa (manual)", value=st.session_state.get("etapa_manual_val", ""), key="etapa_manual_input")
        st.session_state["etapa_manual_val"] = etapa
        st.session_state["etapa_sel"] = etapa
    else:
        etapa = etapa_pick
        st.session_state["etapa_sel"] = etapa
    fecha = st.sidebar.date_input("Fecha", value=st.session_state.get("fecha_val", datetime.today().date()), key="fecha_val")

    # Progreso de carga del dia (horas reales vs DAY_LIMIT_HOURS, p.ej. 25h)
    try:
        _df_day = st.session_state.get("df", pd.DataFrame())
        if isinstance(_df_day, pd.DataFrame) and (not _df_day.empty) and ("Fecha" in _df_day.columns):
            _df_day_local = _df_day.copy()
            _df_day_local["Fecha"] = _df_day_local["Fecha"].astype(str)
            _df_day_local = _df_day_local[_df_day_local["Fecha"] == str(fecha)]
        else:
            _df_day_local = pd.DataFrame()
        _hrs_day = float(pd.to_numeric(_df_day_local.get("Horas_Reales", pd.Series(dtype=float)), errors="coerce").fillna(0).sum())
    except Exception:
        _hrs_day = 0.0
    _pct_day = clamp_0_100(safe_pct(_hrs_day, DAY_LIMIT_HOURS)) if _hrs_day >= 0 else 0.0
    _rest_day = max(0.0, DAY_LIMIT_HOURS - _hrs_day)
    st.sidebar.markdown(f"**Avance del dia ({DAY_LIMIT_HOURS:.0f}h)**")
    st.sidebar.progress(_pct_day / 100.0)
    st.sidebar.caption(f"{_hrs_day:.2f} h cargadas · faltan {_rest_day:.2f} h")

    # Avance por turno (Diurno/Nocturno) - TURNO_LIMIT_HOURS cada uno; la barra muestra máx 100% y "12 h / 12 h" aunque se pasen
    _turno_col = _df_day_local.get("Turno", pd.Series(dtype=str)).fillna("").astype(str)
    _turno_norm = _turno_col.str.lower()
    _is_day_turno = _turno_norm.str.contains("diurno") | _turno_norm.str.contains("dia") | _turno_norm.str.contains("día") | _turno_norm.str.contains("day") | _turno_norm.str.contains("☀")
    _is_night_turno = _turno_norm.str.contains("nocturno") | _turno_norm.str.contains("noche") | _turno_norm.str.contains("night") | _turno_norm.str.contains("🌙")
    _hrs_day_turno = float(pd.to_numeric(_df_day_local.loc[_is_day_turno, "Horas_Reales"], errors="coerce").fillna(0).sum()) if not _df_day_local.empty else 0.0
    _hrs_night_turno = float(pd.to_numeric(_df_day_local.loc[_is_night_turno, "Horas_Reales"], errors="coerce").fillna(0).sum()) if not _df_day_local.empty else 0.0
    _pct_day_turno = clamp_0_100(safe_pct(_hrs_day_turno, TURNO_LIMIT_HOURS)) if _hrs_day_turno >= 0 else 0.0
    _pct_night_turno = clamp_0_100(safe_pct(_hrs_night_turno, TURNO_LIMIT_HOURS)) if _hrs_night_turno >= 0 else 0.0
    # Mostrar horas cap en 12 para que la barra no muestre "18.5 h / 12 h"; a 12h la barra queda al 100%
    _hrs_day_display = min(_hrs_day_turno, TURNO_LIMIT_HOURS)
    _hrs_night_display = min(_hrs_night_turno, TURNO_LIMIT_HOURS)

    st.sidebar.markdown("**Avance por turno**")
    _bar_tpl = """
    <div style="margin: 6px 0 4px 0;">
      <div style="font-size: 0.88rem; font-weight: 600; color: #111827; display:flex; align-items:center; gap:6px;">
        <span>{icon}</span><span>{label}</span><span style="color:#6b7280;">{pct:.0f}%</span>
      </div>
      <div style="height:10px; background:#e5e7eb; border-radius:999px; overflow:hidden; border:1px solid #e5e7eb;">
        <div style="height:100%; width:{pct:.2f}%; background:{bar_color}; border-radius:999px;"></div>
      </div>
      <div style="font-size: 0.8rem; color:#6b7280; margin-top:2px;">{hrs:.2f} h / {limit:.0f} h</div>
    </div>
    """
    st.sidebar.markdown(
        _bar_tpl.format(icon="☀️", label="Diurno", pct=_pct_day_turno, bar_color="#F59E0B", hrs=_hrs_day_display, limit=TURNO_LIMIT_HOURS),
        unsafe_allow_html=True,
    )
    st.sidebar.markdown(
        _bar_tpl.format(icon="🌙", label="Nocturno", pct=_pct_night_turno, bar_color="#2563EB", hrs=_hrs_night_display, limit=TURNO_LIMIT_HOURS),
        unsafe_allow_html=True,
    )
    
# --- Sync contexto actual a drill_day/meta (para que el JSON siempre quede completo) ---
_meta_now = {
    "equipo": st.session_state.get("equipo_val", ""),
    "pozo": st.session_state.get("pozo_val", ""),
    "fecha": str(fecha),
    "equipo_tipo": st.session_state.get("equipo_tipo_val", ""),
    "etapa_manual": bool(st.session_state.get("etapa_manual_val", "")),
    "etapa": etapa,
    "etapa_manual_val": st.session_state.get("etapa_manual_val", ""),
    "modo_reporte": st.session_state.get("modo_reporte", ""),
    "show_charts": bool(st.session_state.get("show_charts", True)),
}
st.session_state.drill_day["meta"] = _meta_now

def _json_default(obj):
    if isinstance(obj, (datetime, date)):
        return obj.isoformat()
    if isinstance(obj, (np.integer, np.floating)):
        return obj.item()
    if isinstance(obj, np.ndarray):
        return obj.tolist()
    try:
        if isinstance(obj, pd.Timestamp):
            return obj.isoformat()
        if pd.isna(obj):
            return None
    except Exception:
        pass
    raise TypeError(f"Object of type {obj.__class__.__name__} is not JSON serializable")

def _render_jornada_avanzado(container) -> None:
    c = container or st.sidebar
    c.markdown("#### Guardar / cargar (avanzado)")

    # Nombre sugerido del archivo (solo nombre, sin ruta)
    _fname_full = _default_jornada_path(equipo, pozo, str(fecha))
    _fname = os.path.basename(_fname_full)

    # Construir payload completo (incluye meta, df, conexiones, BHA, drill_day, etc.)
    def _build_jornada_payload() -> dict:
        meta = {
            "equipo": st.session_state.get("equipo_val", ""),
            "pozo": st.session_state.get("pozo_val", ""),
            "fecha": str(st.session_state.get("fecha_val", "")),
            "equipo_tipo": st.session_state.get("equipo_tipo_val", ""),
            "etapa_manual": bool(st.session_state.get("etapa_manual_chk", False)),
            "etapa": st.session_state.get("etapa_sel", ""),
            "etapa_manual_val": st.session_state.get("etapa_manual_val", ""),
            "modo_reporte": st.session_state.get("modo_reporte", ""),
            "show_charts": bool(st.session_state.get("show_charts", True)),
        }
        # Autocontenible
        st.session_state.drill_day["meta"] = meta

        return {
            "version": "1.2",
            "saved_at": datetime.now().isoformat(timespec="seconds"),
            "meta": meta,
            "df": st.session_state.df.to_dict(orient="records"),
            "df_conn": st.session_state.df_conn.to_dict(orient="records"),
            "df_bha": st.session_state.df_bha.to_dict(orient="records"),
            "drill_day": st.session_state.drill_day,
            "custom_actividades": st.session_state.get("custom_actividades", []),
        }

    _payload = _build_jornada_payload()
    _payload_str = json.dumps(_payload, ensure_ascii=False, indent=2, default=_json_default)

    if use_drive:
        folder_name = st.secrets.get("drive", {}).get("jornadas_folder", "DrillSpot_Jornadas")
        folder_id = _ensure_drive_folder(drive, folder_name)

        # Guardar en Drive
        if c.button("💾 Guardar jornada en Drive", use_container_width=True):
            try:
                _drive_upsert_json(drive, folder_id, _fname, _payload)
                c.success("Guardado en Drive ✅")
            except Exception as e:
                c.error(f"No se pudo guardar en Drive: {e}")

        c.divider()

        # Cargar desde Drive
        c.caption("Cargar desde Drive")
        _files = _drive_list_json(drive, folder_id, limit=100)
        if not _files:
            c.info("No hay jornadas en Drive todavía.")
            up_jornada = None
        else:
            _options = {f'{f["name"]} · {f.get("modifiedTime","")}': f["id"] for f in _files}
            _pick = c.selectbox("Selecciona jornada", list(_options.keys()), key="drive_jornada_pick")
            if c.button("📥 Descargar selección a memoria", use_container_width=True):
                try:
                    payload = _drive_download_json(drive, _options[_pick])
                    # Guardamos en memoria como si fuera un upload
                    st.session_state["_drive_payload_cache"] = payload
                    c.success("Lista para aplicar ✅ (pulsa 'Aplicar jornada')")
                except Exception as e:
                    c.error(f"No se pudo descargar: {e}")
            up_jornada = None  # el apply tomará del cache
    else:
        c.download_button(
            label="Guardar jornada (.json)",
            data=_payload_str,
            file_name=_fname,
            mime="application/json",
            use_container_width=True,
        )

        c.divider()

        up_jornada = c.file_uploader(
            "Cargar jornada (.json)",
            type=["json"],
            accept_multiple_files=False,
            key="jornada_uploader",
            help="Carga un .json previamente guardado para continuar donde se dejó (incluye etapa, estadísticas, etc.).",
        )

    if c.button("Aplicar jornada", use_container_width=True, disabled=(up_jornada is None and st.session_state.get("_drive_payload_cache") is None)):
        payload = None

        # 1) Si viene de Drive (cache en memoria)
        if st.session_state.get("_drive_payload_cache") is not None:
            payload = st.session_state.get("_drive_payload_cache")
        # 2) Si viene de upload local
        elif up_jornada is not None:
            try:
                payload = json.loads(up_jornada.getvalue().decode("utf-8"))
            except Exception as e:
                c.error(f"No se pudo leer el JSON: {e}")
                payload = None

        if isinstance(payload, dict):
            if _apply_jornada_payload(payload):
                c.success("Jornada cargada ✅")
                st.session_state["_drive_payload_cache"] = None
                st.rerun()

# --- Guardar / Cargar jornada ---
# Preferencia: Google Drive (si hay sesión OAuth). Fallback: descarga/subida local.
drive = _drive_service()
use_drive = drive is not None

st.sidebar.markdown("**Selector rápido de jornadas**")
if use_drive:
    folder_name = st.secrets.get("drive", {}).get("jornadas_folder", "DrillSpot_Jornadas")
    folder_id = _ensure_drive_folder(drive, folder_name)
    _files_quick = _drive_list_json(drive, folder_id, limit=50)
    if _files_quick:
        _options_quick = {f'{f["name"]} · {f.get("modifiedTime","")}': f["id"] for f in _files_quick}
        _pick_quick = st.sidebar.selectbox("Jornadas en Drive", list(_options_quick.keys()), key="drive_jornada_quick")
        try:
            _payload_preview = _drive_download_json(drive, _options_quick[_pick_quick])
            _kpi_prev = _kpi_summary_from_payload(_payload_preview)
            _render_kpi_summary(_kpi_prev, title="Resumen KPI (vista previa)")
        except Exception:
            st.sidebar.caption("No se pudo generar vista previa.")
        if st.sidebar.button("Cargar jornada (Drive)", use_container_width=True, key="drive_jornada_quick_btn"):
            try:
                payload = _drive_download_json(drive, _options_quick[_pick_quick])
                if _apply_jornada_payload(payload):
                    st.sidebar.success("Jornada cargada ✅")
                    st.rerun()
            except Exception as e:
                st.sidebar.error(f"No se pudo cargar: {e}")
    else:
        st.sidebar.caption("No hay jornadas en Drive todavía.")
else:
    _local_list = _list_local_jornadas()
    if _local_list:
        _labels = [x[0] for x in _local_list]
        _map = {x[0]: x[1] for x in _local_list}
        _pick_local = st.sidebar.selectbox("Jornadas locales", _labels, key="local_jornada_quick")
        try:
            with open(_map[_pick_local], "r", encoding="utf-8") as f:
                _payload_preview = json.load(f)
            _kpi_prev = _kpi_summary_from_payload(_payload_preview)
            _render_kpi_summary(_kpi_prev, title="Resumen KPI (vista previa)")
        except Exception:
            st.sidebar.caption("No se pudo generar vista previa.")
        if st.sidebar.button("Cargar jornada local", use_container_width=True, key="local_jornada_quick_btn"):
            if load_jornada_json(_map[_pick_local]):
                st.sidebar.success("Jornada cargada ✅")
                st.rerun()
            else:
                st.sidebar.error("No se pudo cargar la jornada local.")
    else:
        st.sidebar.caption("No hay jornadas locales guardadas.")

st.sidebar.markdown("**Subir y aplicar .json**")
up_quick = st.sidebar.file_uploader(
    "Subir jornada (.json)",
    type=["json"],
    accept_multiple_files=False,
    key="quick_jornada_uploader",
)
if st.sidebar.button("Aplicar jornada (subida)", use_container_width=True, disabled=up_quick is None, key="quick_jornada_apply_btn"):
    try:
        payload = json.loads(up_quick.getvalue().decode("utf-8"))
        if _apply_jornada_payload(payload):
            st.sidebar.success("Jornada cargada ✅")
            st.rerun()
    except Exception as e:
        st.sidebar.error(f"No se pudo aplicar la jornada: {e}")

_exp_adv = st.sidebar.expander("Jornada (avanzado)", expanded=False)
_render_jornada_avanzado(_exp_adv)

with st.sidebar.container(border=True):
    st.sidebar.markdown("### Carga colaborativa (por dia)")
    st.sidebar.caption("Cada persona exporta su dia y luego se hace merge sin sobreescribir.")
    colab_name = st.sidebar.text_input("Colaborador", value=st.session_state.get("colab_name", ""), key="colab_name")

    _safe = lambda s: re.sub(r"[^A-Za-z0-9_-]+", "_", str(s)).strip("_")
    _colab_tag = f"_{_safe(colab_name)}" if colab_name else ""
    _day_fname = f"dia_{_safe(pozo)}_{_safe(str(fecha))}{_colab_tag}.json"
    _day_payload = _build_day_payload(fecha, colab_name)
    _day_payload_str = json.dumps(_day_payload, ensure_ascii=False, indent=2, default=_json_default)

    st.sidebar.download_button(
        label="Exportar dia (colaborativo)",
        data=_day_payload_str,
        file_name=_day_fname,
        mime="application/json",
        use_container_width=True,
    )

    st.sidebar.divider()

    up_days = st.sidebar.file_uploader(
        "Importar dias (merge)",
        type=["json"],
        accept_multiple_files=True,
        key="merge_days_uploader",
        help="Sube uno o varios JSON diarios para unir sin sobreescribir registros existentes.",
    )

    if st.sidebar.button("Aplicar merge (dias)", use_container_width=True, disabled=not up_days):
        added_df = 0
        added_conn = 0
        added_bha = 0
        merged_files = 0
        for f in up_days or []:
            try:
                payload = json.loads(f.getvalue().decode("utf-8"))
            except Exception:
                continue
            if not isinstance(payload, dict):
                continue
            inc_df = pd.DataFrame(payload.get("df", []), columns=st.session_state.df.columns)
            inc_df_conn = pd.DataFrame(payload.get("df_conn", []), columns=st.session_state.df_conn.columns)
            inc_df_bha = pd.DataFrame(payload.get("df_bha", []), columns=st.session_state.df_bha.columns)

            st.session_state.df, _a = _merge_df_rows(st.session_state.df, inc_df)
            added_df += _a
            st.session_state.df_conn, _b = _merge_df_rows(st.session_state.df_conn, inc_df_conn)
            added_conn += _b
            st.session_state.df_bha, _c = _merge_df_rows(st.session_state.df_bha, inc_df_bha)
            added_bha += _c

            st.session_state.drill_day = _merge_dict_no_overwrite(
                st.session_state.drill_day, payload.get("drill_day", {})
            )
            merged_files += 1

        st.sidebar.success(
            f"Merge listo: {merged_files} archivos. "
            f"Filas nuevas -> actividades: {added_df}, conexiones: {added_conn}, BHA: {added_bha}"
        )
        st.rerun()

with st.sidebar.container(border=True):
    st.sidebar.markdown("### Modo")
    modo_reporte = st.sidebar.radio(
        "Tipo",
        MODO_REPORTE_OPTS,
        index=MODO_REPORTE_OPTS.index(modo_reporte) if modo_reporte in MODO_REPORTE_OPTS else 0
    )
    st.session_state["modo_reporte"] = modo_reporte

# Toggle liviano para evitar render pesado
show_charts = st.sidebar.toggle(
    "Mostrar gráficas (mejor rendimiento)",
    value=bool(st.session_state.get("show_charts", True)),
    key="show_charts",
)

with st.sidebar.container(border=True):
    st.sidebar.markdown("### Catálogo TNPI (CSV)")
    script_dir = os.path.dirname(os.path.abspath(__file__)) if "__file__" in globals() else os.getcwd()
    candidate_local = os.path.join(script_dir, "Detalles causas de TNPI.csv")
    csv_path_use = TNPI_CSV_PATH if (TNPI_CSV_PATH and os.path.exists(TNPI_CSV_PATH)) else (candidate_local if os.path.exists(candidate_local) else "")
    up = st.sidebar.file_uploader("Cargar CSV", type=["csv"], accept_multiple_files=False)

    if up is not None:
        try:
            df_tnpi_cat = robust_read_csv(up)
        except Exception as e:
            df_tnpi_cat = None
            st.sidebar.error(f"Error leyendo CSV TNPI: {e}")

        if df_tnpi_cat is not None and not df_tnpi_cat.empty:
            det_col = None
            cat_col = None
            for c in df_tnpi_cat.columns:
                cl = str(c).lower()
                if det_col is None and ("detalle" in cl or "causa" in cl):
                    det_col = c
                if cat_col is None and ("categoria" in cl or "categoría" in cl):
                    cat_col = c

            # Fallbacks comunes
            if det_col is None:
                for cand in ["Detalle de causa de TNPI", "Detalle", "Causa", "Detalle_TNPI"]:
                    if cand in df_tnpi_cat.columns:
                        det_col = cand
                        break
            if cat_col is None:
                for cand in ["Categoria", "Categoría", "Categoria_TNPI"]:
                    if cand in df_tnpi_cat.columns:
                        cat_col = cand
                        break

            if det_col is None or cat_col is None:
                st.sidebar.error("No pude identificar columnas de Categoria/Detalle en el CSV TNPI.")
            else:
                df_tnpi_cat = df_tnpi_cat[[cat_col, det_col]].copy()
                df_tnpi_cat.columns = ["Categoria_TNPI", "Detalle_TNPI"]
                df_tnpi_cat["Categoria_TNPI"] = df_tnpi_cat["Categoria_TNPI"].apply(smart_case)
                df_tnpi_cat["Detalle_TNPI"] = df_tnpi_cat["Detalle_TNPI"].apply(smart_case)
                # Mantener columnas TNP para compatibilidad con el resto de la app
                df_tnpi_cat["Categoria_TNP"] = df_tnpi_cat["Categoria_TNPI"]
                df_tnpi_cat["Detalle_TNP"] = df_tnpi_cat["Detalle_TNPI"]
                df_tnpi_cat = df_tnpi_cat.dropna().drop_duplicates().reset_index(drop=True)
                st.sidebar.success("CSV TNPI cargado")
    else:
        df_tnpi_cat = load_tnpi_catalog(csv_path_use)
        if not csv_path_use:
            st.sidebar.warning("No se encontró el CSV. Usando catálogo mínimo.")


# ------------------------------
# Catálogo de causas TNP (Tiempo No Productivo) - similar a TNPI
# ------------------------------
def load_tnp_catalog(path_csv: str) -> pd.DataFrame:
    """Carga catálogo TNP desde CSV. Soporta utf-8 / latin-1."""
    if not path_csv or not os.path.exists(path_csv):
        return pd.DataFrame({"Categoria_TNP": ["-"], "Detalle_TNP": ["-"]})
    try:
        df0 = robust_read_csv(path_csv)
    except Exception:
        df0 = None
    if df0 is None or df0.empty:
        return pd.DataFrame({"Categoria_TNP": ["-"], "Detalle_TNP": ["-"]})

    det_col = None
    cat_col = None
    for c in df0.columns:
        cl = str(c).lower()
        if det_col is None and ("detalle" in cl or "causa" in cl):
            det_col = c
        if cat_col is None and ("categoria" in cl or "categoría" in cl):
            cat_col = c
    # Fallback por nombres esperados del archivo que nos compartiste
    if cat_col is None and "Categoria" in df0.columns:
        cat_col = "Categoria"
    if det_col is None and "Detalle de causa de TNP" in df0.columns:
        det_col = "Detalle de causa de TNP"

    if cat_col is None or det_col is None:
        # Intento: tomar primeras 2 cols
        cols = list(df0.columns)[:2]
        if len(cols) >= 2:
            cat_col, det_col = cols[1], cols[0]

    df0 = df0[[cat_col, det_col]].copy()
    df0.columns = ["Categoria_TNP", "Detalle_TNP"]
    df0["Categoria_TNP"] = df0["Categoria_TNP"].astype(str).apply(smart_case)
    df0["Detalle_TNP"] = df0["Detalle_TNP"].astype(str).apply(smart_case)
    df0 = df0.dropna().drop_duplicates().reset_index(drop=True)
    return df0

# Default: buscamos el CSV TNP junto al script o en rutas conocidas
TNP_CSV_PATH = ""
candidate_tnp_local = os.path.join(os.getcwd(), "Detalles causas de TNP2.csv")
candidate_tnp_alt = os.path.join(os.path.dirname(__file__) if "__file__" in globals() else os.getcwd(), "Detalles causas de TNP2.csv")
if os.path.exists(candidate_tnp_local):
    TNP_CSV_PATH = candidate_tnp_local
elif os.path.exists(candidate_tnp_alt):
    TNP_CSV_PATH = candidate_tnp_alt

st.sidebar.markdown("---")
st.sidebar.subheader("Catálogo TNP (causas)")
up_tnp = st.sidebar.file_uploader("Cargar CSV TNP", type=["csv"], accept_multiple_files=False, key="up_tnp_cat")

if up_tnp is not None:
    try:
        # DrillSpot exports a veces vienen con delimitador/encoding irregular
        df_tnp_cat = robust_read_csv(up_tnp)
    except Exception as e:
        df_tnp_cat = None
        st.sidebar.error(f"Error leyendo CSV TNP: {e}")

    if df_tnp_cat is not None and not df_tnp_cat.empty:
        det_col = None
        cat_col = None
        for c in df_tnp_cat.columns:
            cl = str(c).lower()
            if det_col is None and ("detalle" in cl or "causa" in cl):
                det_col = c
            if cat_col is None and ("categoria" in cl or "categoría" in cl):
                cat_col = c

        # Fallbacks comunes
        if cat_col is None:
            for cand in ["Categoria", "Categoría", "Categoria_TNP"]:
                if cand in df_tnp_cat.columns:
                    cat_col = cand
                    break
        if det_col is None:
            for cand in ["Detalle de causa de TNP", "Detalle", "Causa", "Detalle_TNP"]:
                if cand in df_tnp_cat.columns:
                    det_col = cand
                    break

        if det_col is None or cat_col is None:
            st.sidebar.error("No pude identificar columnas de Categoria/Detalle en el CSV TNP.")
        else:
            df_tnp_cat = df_tnp_cat[[cat_col, det_col]].copy()
            df_tnp_cat.columns = ["Categoria_TNP", "Detalle_TNP"]
            df_tnp_cat["Categoria_TNP"] = df_tnp_cat["Categoria_TNP"].astype(str).apply(smart_case)
            df_tnp_cat["Detalle_TNP"] = df_tnp_cat["Detalle_TNP"].astype(str).apply(smart_case)
            df_tnp_cat = df_tnp_cat.dropna().drop_duplicates().reset_index(drop=True)
            st.sidebar.success("CSV TNP cargado")
else:
    df_tnp_cat = load_tnp_catalog(TNP_CSV_PATH)
    if not TNP_CSV_PATH:
        st.sidebar.warning("No se encontró CSV TNP. Usando catálogo mínimo.")

tnp_cat_list = sorted(df_tnp_cat["Categoria_TNP"].dropna().unique().tolist())

cat_list = sorted(df_tnpi_cat["Categoria_TNPI"].dropna().unique().tolist()) or ["Proceso"]

# Inputs perforación (metros/ROP) + PT/Prof actual
if modo_reporte == "Perforación":
    with st.sidebar.container(border=True):
        st.sidebar.markdown("### Profundidad (avance) - Por Etapa")
        
        # Obtener datos específicos de esta etapa
        etapa_data = get_etapa_data(etapa)
        
        etapa_data["pt_programada_m"] = st.sidebar.number_input(
            f"PT programada (m) - {etapa}",
            0.0, step=1.0, 
            value=float(etapa_data["pt_programada_m"])
        )
        
        etapa_data["rop_prog_etapa"] = st.sidebar.number_input(
            f"ROP programada (m/h) - {etapa}",
            0.0, step=0.1,
            value=float(etapa_data.get("rop_prog_etapa", 0.0))
        )

        etapa_data["prof_actual_m"] = st.sidebar.number_input(
            f"Profundidad actual (m) - {etapa}",
            0.0, step=1.0, 
            value=float(etapa_data["prof_actual_m"])
        )
        
        # Mantener compatibilidad con datos globales (opcional)
        st.session_state.drill_day["pt_programada_m"] = etapa_data["pt_programada_m"]
        st.session_state.drill_day["prof_actual_m"] = etapa_data["prof_actual_m"]

    # (Metros perforados (día) movido a la pestaña ROP)
# CONTEXTO ACTUAL (PONER DESPUÉS DE LOS INPUTS DE PROFUNDIDAD)
with st.sidebar.container(border=True):
    st.sidebar.markdown("### Contexto Actual")
    
    # Mostrar claramente qué etapa estamos trabajando
    st.sidebar.markdown(f"""
        <div style='background: rgba(40, 180, 99, 0.1); padding: 8px; border-radius: 8px; border-left: 3px solid #28B463; margin-bottom: 10px;'>
            <div style='font-size: 12px; color: #28B463;'>Etapa actual:</div>
            <div style='font-size: 16px; color: white; font-weight: bold;'>{etapa}</div>
        </div>
    """, unsafe_allow_html=True)

    # Barrena (BNA) global
    _bna_default = st.session_state.get("barrena_global", BARRERAS_DEFAULT[0])
    _bna_opts = BARRERAS_DEFAULT + ["Otra (manual)"]
    _bna_idx = _bna_opts.index(_bna_default) if _bna_default in _bna_opts else 0
    barrena_pick = st.sidebar.selectbox("Barrena (BNA)", options=_bna_opts, index=_bna_idx, key="barrena_global_sel")
    if barrena_pick == "Otra (manual)":
        barrena_global = st.sidebar.text_input("Barrena (manual)", value=st.session_state.get("barrena_manual", ""), key="barrena_manual").strip()
    else:
        barrena_global = barrena_pick
    st.session_state["barrena_global"] = barrena_global
    
    # Indicador de qué datos se están capturando
    if modo_reporte == "Perforación":
        # Contar actividades en esta etapa
                # FIX: usar siempre los DataFrames del session_state (df aún no está definido aquí)
        _df_loc = st.session_state.df
        _dfc_loc = st.session_state.df_conn
        actividades_etapa = len(_df_loc[_df_loc["Etapa"] == etapa]) if not _df_loc.empty else 0
        conexiones_etapa = len(_dfc_loc[_dfc_loc["Etapa"] == etapa]) if not _dfc_loc.empty else 0
        
        st.sidebar.markdown(f"""
            <div style='font-size: 12px; color: rgba(255,255,255,0.7);'>
                📊 <b>Actividades:</b> {actividades_etapa}<br>
                🔗 <b>Conexiones:</b> {conexiones_etapa}
            </div>
        """, unsafe_allow_html=True)

    st.sidebar.markdown("### Captura actividad (general)")
    # Corrida activa (Run): se usa como contexto global y para ROP programada por corrida
    _corrida_prev = st.session_state.get("corrida_activa", None)
    if _corrida_prev is None and "Corrida (Run)" in st.session_state:
        # compatibilidad con versiones previas (cuando no había key explícita)
        _corrida_prev = st.session_state.get("Corrida (Run)")
    corrida = st.sidebar.text_input("Corrida (Run)", _corrida_prev or "Run 1", key="corrida_activa")
    st.session_state.drill_day["corrida_activa"] = corrida
    tipo_agujero = st.sidebar.radio("Tipo de agujero", TIPO_AGUJERO, horizontal=True)
    turno = st.sidebar.radio("Turno", TURNOS, horizontal=True)

    operacion = "Perforación" if modo_reporte == "Perforación" else st.sidebar.selectbox(
        "Operación", ["Superficie", "TR", "Otra"], index=0
    )

    # --- Actividad (catálogo + personalizadas + otra) ---
    if modo_reporte == "Cambio de etapa":
        actividades_base = ACTIVIDADES_CE
        actividades_opts = actividades_base
    else:
        actividades_base = ACTIVIDADES
        actividades_opts = actividades_base + sorted(st.session_state.get("custom_actividades", []))

    actividad_sel = st.sidebar.selectbox(
        "Actividad",
        actividades_opts + ["Otra (especificar)"],
        key="actividad_select_sidebar",
        help="Catálogo según el modo de reporte (Perforación / Cambio de etapa)."
    )

    actividad = actividad_sel
    if actividad_sel == "Otra (especificar)":
        actividad = st.sidebar.text_input("Especifica actividad", "", key="actividad_otro").strip()

    # Tipo de tiempo (SIEMPRE visible)
    tipo = st.sidebar.radio("Tipo de tiempo", ["TP", "TNPI", "TNP"], horizontal=True, key="tipo_time_general")

    # Hora (opcional) para discretizar por horas
    registrar_hora = st.sidebar.checkbox("Registrar hora", value=False, key="act_use_time")
    hora_ini = None
    hora_fin = None
    bitacora_enabled = False
    bitacora_entries = st.session_state.get("act_bitacora_entries", [])
    bitacora_total_h = float(st.session_state.get("act_bitacora_total_h", 0.0) or 0.0)
    if registrar_hora:
        hora_ini = st.sidebar.time_input(
            "Hora inicio",
            value=st.session_state.get("act_hora_ini", datetime.now().time()),
            key="act_hora_ini",
        )
        hora_fin = st.sidebar.time_input(
            "Hora fin",
            value=st.session_state.get("act_hora_fin", datetime.now().time()),
            key="act_hora_fin",
        )
        bitacora_enabled = st.sidebar.toggle("Bitácora por horas", value=False, key="act_use_bitacora")
    hora_ini_txt = hora_ini.strftime("%H:%M") if (registrar_hora and hora_ini) else ""
    hora_fin_txt = hora_fin.strftime("%H:%M") if (registrar_hora and hora_fin) else ""
    if registrar_hora and bitacora_enabled:
        st.sidebar.caption("Completa la bitácora en la pestaña **Bitácora por horas**.")

    render_chip_row([
        {"label": "Modo", "value": modo_reporte, "tone": "blue"},
        {"label": "Turno", "value": turno, "tone": "violet"},
        {"label": "Tipo", "value": tipo, "tone": "amber" if tipo == "TNPI" else ("red" if tipo == "TNP" else "green")},
        {"label": "Operación", "value": operacion, "tone": "gray"},
        {"label": "Actividad", "value": actividad or "-", "tone": "blue"},
        {"label": "Corrida", "value": corrida, "tone": "gray"},
    ], use_iframe=True, height=120)

    # -------------------------------------------------
    # Helper: Viajes (calcular estándar sugerido)
    # Estándar (h) = distancia(m)/velocidad(m/h) + conexiones * tconn(min)/60
    # -------------------------------------------------
    if actividad in VIAJE_CATALOG:
        with st.sidebar.expander("Viaje – calculadora estándar (TNPI)", expanded=False):
            v = float(VIAJE_CATALOG[actividad].get("vel_mh", 0.0) or 0.0)
            tc = float(VIAJE_CATALOG[actividad].get("tconn_min", 0.0) or 0.0)

            etapa_viajes_sel = st.selectbox("Etapa base para viaje", options=st.session_state.depth_rows["Etapa"].tolist(), index=0)
            _drow = st.session_state.depth_rows[st.session_state.depth_rows["Etapa"] == etapa_viajes_sel].iloc[0]
            pt_prog_v = float(_drow["PT_programada_m"] or 0.0)
            pt_act_v = float(_drow["PT_actual_m"] or 0.0)

            dist = st.number_input("Distancia (m)", min_value=0.0, value=max(pt_prog_v - pt_act_v, 0.0), step=10.0, key="dist_viaje")
            conexiones_etapa = len(st.session_state.df_conn[st.session_state.df_conn["Etapa"] == etapa_viajes_sel]) if "df_conn" in st.session_state else 0
            nconn = st.number_input("Conexiones (#)", min_value=0, value=int(conexiones_etapa), step=1, key="nconn_viaje")

            if v > 0:
                est = dist / v + (nconn * tc / 60.0)
                st.caption(f"Estándar sugerido: {est:.2f} h (v={v:.0f} m/h, tconn={tc:.1f} min)")
            else:
                st.caption("Configura vel_mh > 0 en VIAJE_CATALOG para cálculo automático.")

    # Detalles TNPI/TNP (SIEMPRE disponibles cuando aplique)
    categoria_tnpi = "-"
    detalle_tnpi = "-"
    categoria_tnp = "-"
    detalle_tnp = "-"

    if tipo == "TNPI":
        # Usa el catálogo TNPI cargado (df_tnpi_cat) y su lista de categorías (cat_list)
        categoria_tnpi = st.sidebar.selectbox(
            "Categoría TNPI",
            options=cat_list if "cat_list" in globals() else ["-"],
            key="cat_tnpi_general",
        )
        det_all = (
            df_tnpi_cat[df_tnpi_cat["Categoria_TNPI"] == categoria_tnpi]["Detalle_TNPI"].tolist()
            if "df_tnpi_cat" in globals()
            else ["-"]
        )
        q = (st.sidebar.text_input("Buscar detalle TNPI", "", key="q_tnpi_general") or "").strip().lower()
        det_filtered = [d for d in det_all if q in str(d).lower()] if q else det_all
        detalle_tnpi = st.sidebar.selectbox(
            "Detalle TNPI",
            options=det_filtered if det_filtered else det_all,
            key="det_tnpi_general",
        )

    elif tipo == "TNP":
        categoria_tnp = st.sidebar.selectbox(
            "Categoría TNP",
            options=tnp_cat_list if "tnp_cat_list" in globals() else ["-"],
            key="cat_tnp_general",
        )
        det_all_tnp = (
            df_tnp_cat[df_tnp_cat["Categoria_TNP"] == categoria_tnp]["Detalle_TNP"].tolist()
            if "df_tnp_cat" in globals()
            else ["-"]
        )
        q2 = (st.sidebar.text_input("Buscar detalle TNP", "", key="q_tnp_general") or "").strip().lower()
        det_filtered_tnp = [d for d in det_all_tnp if q2 in str(d).lower()] if q2 else det_all_tnp
        detalle_tnp = st.sidebar.selectbox(
            "Detalle TNP",
            options=det_filtered_tnp if det_filtered_tnp else det_all_tnp,
            key="det_tnp_general",
        )

    horas_prog = st.sidebar.number_input("Horas estándar / programadas (h)", 0.0, step=0.25, key="hp_general")
    if registrar_hora and bitacora_enabled:
        horas_real = float(bitacora_total_h or 0.0)
        if horas_real > 0:
            st.sidebar.caption(f"Horas reales (bitácora): {horas_real:.2f} h")
        else:
            st.sidebar.caption("Horas reales (bitácora): 0.00 h")
    elif registrar_hora and hora_ini and hora_fin:
        _dt_ini = datetime.combine(datetime.today().date(), hora_ini)
        _dt_fin = datetime.combine(datetime.today().date(), hora_fin)
        horas_real = (_dt_fin - _dt_ini).total_seconds() / 3600.0
        if horas_real < 0:
            horas_real += 24.0
        st.sidebar.caption(f"Horas reales calculadas: {horas_real:.2f} h")
    else:
        horas_real = st.sidebar.number_input("Horas reales (h)", 0.0, step=0.25, key="hr_general")
    rop_prog = 0.0
    rop_real = 0.0
    # ROP por actividad (opcional) se centraliza en la pestaña "ROP" para evitar confusión.

    

    # --- Auto TNPI por exceso (solo cuando capturas como TP y hay estándar) ---

    tnpi_exceso_h = 0.0

    if tipo == "TP" and float(horas_prog) > 0:
        tnpi_exceso_h = max(0.0, float(horas_real) - float(horas_prog))

    if tnpi_exceso_h > 0:
        st.sidebar.markdown(f"**TNPI por exceso detectado:** {tnpi_exceso_h:.2f} h")
        st.sidebar.caption(f"(Real {float(horas_real):.2f} h − Estándar {float(horas_prog):.2f} h)")

        # --- Exceso (Real > Estándar) ---
        st.sidebar.markdown("**Asignación del exceso (Real > Estándar)**")

        exceso_tipo = st.sidebar.radio(
            "¿Cómo registrar el exceso?",
            options=["TNPI", "TNP"],
            horizontal=True,
            key="exceso_tipo_general_choice",
        )
        st.session_state["exceso_tipo_general"] = exceso_tipo

        # Helpers de catálogo
        _has_cat = "df_tnpi_cat" in globals() and isinstance(df_tnpi_cat, pd.DataFrame)

        # --- Asignación simple (una sola causa) ---
        if exceso_tipo == "TNPI":
            st.sidebar.markdown("**Asignación simple (TNPI)**")
            act_cat_simple = st.sidebar.selectbox(
                "Categoría TNPI (exceso)",
                options=cat_list if "cat_list" in globals() else ["-"],
                index=0,
                key="act_cat_simple",
            )

            if _has_cat and "Categoria_TNPI" in df_tnpi_cat.columns and "Detalle_TNPI" in df_tnpi_cat.columns:
                _det_opts_act = (
                    df_tnpi_cat[df_tnpi_cat["Categoria_TNPI"] == act_cat_simple]["Detalle_TNPI"]
                    .dropna()
                    .unique()
                    .tolist()
                )
            else:
                _det_opts_act = ["-"]

            act_det_simple = st.sidebar.selectbox(
                "Detalle TNPI (exceso)",
                options=_det_opts_act if len(_det_opts_act) else ["-"],
                index=0,
                key="act_det_simple",
            )

            # placeholders para no romper lógica abajo
            act_cat_simple_tnp = "-"
            act_det_simple_tnp = "-"

        else:
            st.sidebar.markdown("**Asignación simple (TNP)**")

            # Lista categorías TNP (fallback: TNPI si no existe catálogo TNP)
            if _has_cat and "Categoria_TNP" in df_tnpi_cat.columns:
                _cat_tnp = df_tnpi_cat["Categoria_TNP"].dropna().unique().tolist()
                _cat_tnp = [c for c in _cat_tnp if str(c).strip() not in ("", "-", "nan", "None")]
            else:
                _cat_tnp = []

            if not _cat_tnp and "cat_list" in globals():
                _cat_tnp = cat_list

            act_cat_simple_tnp = st.sidebar.selectbox(
                "Categoría TNP (exceso)",
                options=_cat_tnp if _cat_tnp else ["-"],
                index=0,
                key="act_cat_simple_tnp",
            )

            if _has_cat and "Categoria_TNP" in df_tnpi_cat.columns and "Detalle_TNP" in df_tnpi_cat.columns:
                _det_opts_tnp = (
                    df_tnpi_cat[df_tnpi_cat["Categoria_TNP"] == act_cat_simple_tnp]["Detalle_TNP"]
                    .dropna()
                    .unique()
                    .tolist()
                )
            else:
                _det_opts_tnp = ["-"]

            act_det_simple_tnp = st.sidebar.selectbox(
                "Detalle TNP (exceso)",
                options=_det_opts_tnp if len(_det_opts_tnp) else ["-"],
                index=0,
                key="act_det_simple_tnp",
            )

            # placeholders para no romper lógica TNPI abajo
            act_cat_simple = "-"
            act_det_simple = "-"
# --- Desglose opcional (múltiples causas) ---
        with st.sidebar.expander("Detalle TNPI (opcional) - desglose por múltiples causas", expanded=False):
            st.caption("Agrega varias causas y presiona **Guardar**. La suma debe ser igual al TNPI por exceso detectado.")

            if (not st.session_state.get("act_tnpi_breakdown_draft")) and (not st.session_state.get("act_tnpi_breakdown")):
                st.session_state.act_tnpi_breakdown_draft = [{
                    "Categoria_TNPI": (cat_list[0] if "cat_list" in globals() and isinstance(cat_list, list) and len(cat_list) else "-"),
                    "Detalle_TNPI": "-",
                    "Horas_TNPI_h": float(tnpi_exceso_h),
                    "Comentario": "",
                }]

            _cat_new_act = st.selectbox(
                "Categoría TNPI (nuevo renglón)",
                options=cat_list if "cat_list" in globals() else ["-"],
                index=0,
                key="act_bd_cat_new",
            )

            if "df_tnpi_cat" in globals() and "Categoria_TNPI" in df_tnpi_cat.columns and "Detalle_TNPI" in df_tnpi_cat.columns:
                _det_new_opts_act = (
                    df_tnpi_cat[df_tnpi_cat["Categoria_TNPI"] == _cat_new_act]["Detalle_TNPI"]
                    .dropna()
                    .unique()
                    .tolist()
                )
            else:
                _det_new_opts_act = ["-"]

            _det_new_act = st.selectbox(
                "Detalle TNPI (nuevo renglón)",
                options=_det_new_opts_act if len(_det_new_opts_act) else ["-"],
                index=0,
                key="act_bd_det_new",
            )

            _h_new_act = st.number_input(
                "Horas TNPI (h) (nuevo renglón)",
                min_value=0.0,
                value=0.0,
                step=0.25,
                format="%.2f",
                key="act_bd_h_new",
            )

            _com_new_act = st.text_input("Comentario (nuevo renglón)", "", key="act_bd_com_new")

            # --- Botones del desglose ---
            if st.button("Agregar", use_container_width=True, key="act_bd_add"):
                if _h_new_act <= 0:
                    st.warning("Las horas deben ser > 0 para agregar la causa.")
                else:
                    st.session_state.act_tnpi_breakdown_draft.append({
                        "Categoria_TNPI": _cat_new_act,
                        "Detalle_TNPI": _det_new_act,
                        "Horas_TNPI_h": float(_h_new_act),
                        "Comentario": _com_new_act,
                    })
                    st.rerun()

            if st.button("Limpiar", use_container_width=True, key="act_bd_clear"):
                st.session_state.act_tnpi_breakdown_draft = []
                st.session_state.act_tnpi_breakdown_saved = False
                st.rerun()

            # Tabla borrador
            bd_act = pd.DataFrame(
                st.session_state.get("act_tnpi_breakdown_draft", []),
                columns=["Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP", "Detalle_TNP", "Horas_TNPI_h", "Comentario"],
            )
            bd_act["Horas_TNPI_h"] = pd.to_numeric(bd_act["Horas_TNPI_h"], errors="coerce").fillna(0.0)

            st.markdown("**Causas en borrador**")
            if bd_act.empty:
                st.info("Aún no hay causas en el borrador.")
            else:
                st.dataframe(bd_act, use_container_width=True, hide_index=True)

            sum_bd_act = float(bd_act["Horas_TNPI_h"].sum()) if not bd_act.empty else 0.0
            st.caption(
                f"Suma borrador: **{sum_bd_act:.2f} h**  |  TNPI por exceso: **{float(tnpi_exceso_h):.2f} h**"
            )

            if st.button("Guardar", use_container_width=True, key="act_bd_save"):
                if bd_act.empty:
                        st.error("No hay renglones en el borrador para guardar.")
                else:
                    bd2 = bd_act[bd_act["Horas_TNPI_h"] > 0].copy()
                    sum2 = float(bd2["Horas_TNPI_h"].sum())
                    if abs(sum2 - float(tnpi_exceso_h)) > 1e-6:
                        st.error(
                            f"La suma del desglose ({sum2:.2f} h) debe ser igual al TNPI por exceso ({float(tnpi_exceso_h):.2f} h)."
                        )
                    else:
                        st.session_state.act_tnpi_breakdown = bd2.to_dict(orient="records")
                        st.session_state.act_tnpi_breakdown_saved = True
                        st.success(
                            "Desglose guardado. Al agregar la actividad se registrarán varias causas TNPI."
                        )
    comentario = st.sidebar.text_input("Comentario", "", key="com_general")
    # Aviso preventivo: menos de 1 hora disponible en el día
    _remaining_day = _remaining_day_hours(st.session_state.df, fecha)
    if 0 < _remaining_day <= 1.0:
        st.sidebar.warning(f"Queda {float(_remaining_day):.2f} h disponible en el día.")
    disable_general_add = (actividad in ["Conexión perforando", "Arma/Desarma BHA"]) or (actividad_sel == "Otra (especificar)" and not actividad)
    if st.sidebar.button("Agregar actividad", use_container_width=True, disabled=disable_general_add):
        if actividad_sel == "Otra (especificar)" and not actividad:
            st.warning("Especifica el nombre de la actividad antes de agregarla.")
            st.stop()

        # Memorizar actividad nueva en esta sesión (y se persiste si guardas la jornada)
        if actividad_sel == "Otra (especificar)" and actividad:
            base_lower = {a.strip().lower() for a in ACTIVIDADES}
            custom_lower = {a.strip().lower() for a in st.session_state.get("custom_actividades", [])}
            if actividad.lower() not in base_lower and actividad.lower() not in custom_lower:
                st.session_state.custom_actividades.append(actividad)

        etapa_use = ((etapa_viajes_sel or etapa) if "etapa_viajes_sel" in globals() else etapa)

        # --- Lógica general ---
        # 1) Si capturas como TP y el Real supera el Estándar, se divide en:
        #    - TP = Estándar
        #    - TNPI = exceso
        # 2) Para ese TNPI por exceso: puedes asignarlo a 1 causa (simple) o desglosarlo (detalle opcional).

        add_rows = []

        # Caso: bitácora por horas -> múltiples filas por tramo
        if registrar_hora and bitacora_enabled and len(bitacora_entries) > 0:
            total_bit = float(sum(r["Horas_Reales"] for r in bitacora_entries))
            exceso_total_h = 0.0
            if tipo == "TP" and float(horas_prog) > 0 and total_bit > float(horas_prog):
                exceso_total_h = float(total_bit) - float(horas_prog)
            for r in bitacora_entries:
                _factor = float(r["Horas_Reales"]) / total_bit if total_bit > 0 else 0.0
                _prog = 0.0
                if float(horas_prog) > 0 and total_bit > 0:
                    _prog = float(horas_prog) * _factor
                base_row = {
                    "Equipo": equipo,
                    "Pozo": pozo,
                    "Etapa": etapa_use,
                    "Fecha": str(fecha),
                    "Equipo_Tipo": st.session_state.get("equipo_tipo_val", ""),
                    "Modo_Reporte": modo_reporte,
                    "Seccion": etapa,
                    "Corrida": corrida,
                    "Tipo_Agujero": tipo_agujero,
                    "Operacion": operacion,
                    "Actividad": actividad,
                    "Turno": turno,
                    "Hora_Inicio": str(r.get("Hora_Inicio", "")),
                    "Hora_Fin": str(r.get("Hora_Fin", "")),
                    "ROP_Prog_mh": float(rop_prog),
                    "ROP_Real_mh": float(rop_real),
                    "Comentario": str(r.get("Comentario", "") or comentario),
                    "Origen": "Manual",
                }

                if tipo == "TP" and exceso_total_h > 0:
                    # TP hasta el estándar proporcional + exceso como TNPI/TNP
                    tp_h = float(_prog)
                    exceso_h = float(r.get("Horas_Reales", 0.0)) - tp_h
                    add_rows.append({
                        **base_row,
                        "Tipo": "TP",
                        "Categoria_TNPI": "-",
                        "Detalle_TNPI": "-",
                        "Categoria_TNP": "-",
                        "Detalle_TNP": "-",
                        "Horas_Prog": float(_prog),
                        "Horas_Reales": float(tp_h),
                    })

                    exceso_tipo = st.session_state.get("exceso_tipo_general", "TNPI")
                    exceso_tipo = "TNP" if str(exceso_tipo).upper() == "TNP" else "TNPI"
                    if exceso_h > 0:
                        if exceso_tipo == "TNP":
                            add_rows.append({
                                **base_row,
                                "Tipo": "TNP",
                                "Categoria_TNPI": "-",
                                "Detalle_TNPI": "-",
                                "Categoria_TNP": (act_cat_simple_tnp if "act_cat_simple_tnp" in locals() else categoria_tnp),
                                "Detalle_TNP": (act_det_simple_tnp if "act_det_simple_tnp" in locals() else detalle_tnp),
                                "Horas_Prog": 0.0,
                                "Horas_Reales": float(exceso_h),
                                "Comentario": f"Exceso (Real {total_bit:.2f}h > Estándar {horas_prog:.2f}h) registrado como TNP.",
                                "Origen": "EXCESO",
                            })
                        else:
                            # TNPI: si hay desglose guardado, distribuir proporcionalmente por tramo
                            act_tnpi_breakdown = st.session_state.get("act_tnpi_breakdown", None)
                            bd_saved = pd.DataFrame()
                            if act_tnpi_breakdown is not None and len(act_tnpi_breakdown) > 0 and bool(st.session_state.get("act_tnpi_breakdown_saved", False)):
                                try:
                                    bd_saved = pd.DataFrame(act_tnpi_breakdown)
                                    if "Categoria_TNPI" not in bd_saved.columns and "Categoria" in bd_saved.columns:
                                        bd_saved["Categoria_TNPI"] = bd_saved["Categoria"]
                                    if "Detalle_TNPI" not in bd_saved.columns and "Detalle" in bd_saved.columns:
                                        bd_saved["Detalle_TNPI"] = bd_saved["Detalle"]
                                    if "Horas_Reales" not in bd_saved.columns and "Horas" in bd_saved.columns:
                                        bd_saved["Horas_Reales"] = bd_saved["Horas"]
                                    bd_saved = bd_saved[["Categoria_TNPI", "Detalle_TNPI", "Horas_Reales"]].copy()
                                    bd_saved["Horas_Reales"] = pd.to_numeric(bd_saved["Horas_Reales"], errors="coerce").fillna(0.0)
                                except Exception:
                                    bd_saved = pd.DataFrame()

                            if not bd_saved.empty:
                                share = float(exceso_h) / float(exceso_total_h) if exceso_total_h > 0 else 0.0
                                for _, br in bd_saved.iterrows():
                                    _h = float(br.get("Horas_Reales", 0.0) or 0.0) * share
                                    if _h <= 0:
                                        continue
                                    add_rows.append({
                                        **base_row,
                                        "Tipo": "TNPI",
                                        "Categoria_TNPI": str(br.get("Categoria_TNPI", "-") or "-"),
                                        "Detalle_TNPI": str(br.get("Detalle_TNPI", "-") or "-"),
                                        "Categoria_TNP": "-",
                                        "Detalle_TNP": "-",
                                        "Horas_Prog": 0.0,
                                        "Horas_Reales": float(_h),
                                        "Comentario": f"Exceso (Real {total_bit:.2f}h > Estándar {horas_prog:.2f}h) registrado como TNPI.",
                                        "Origen": "EXCESO",
                                    })
                            else:
                                add_rows.append({
                                    **base_row,
                                    "Tipo": "TNPI",
                                    "Categoria_TNPI": (act_cat_simple if "act_cat_simple" in locals() else categoria_tnpi),
                                    "Detalle_TNPI": (act_det_simple if "act_det_simple" in locals() else detalle_tnpi),
                                    "Categoria_TNP": "-",
                                    "Detalle_TNP": "-",
                                    "Horas_Prog": 0.0,
                                    "Horas_Reales": float(exceso_h),
                                    "Comentario": f"Exceso (Real {total_bit:.2f}h > Estándar {horas_prog:.2f}h) registrado como TNPI.",
                                    "Origen": "EXCESO",
                                })
                else:
                    add_rows.append({
                        **base_row,
                        "Tipo": tipo,
                        "Categoria_TNPI": categoria_tnpi if tipo == "TNPI" else "-",
                        "Detalle_TNPI": detalle_tnpi if tipo == "TNPI" else "-",
                        "Categoria_TNP": categoria_tnp if tipo == "TNP" else "-",
                        "Detalle_TNP": detalle_tnp if tipo == "TNP" else "-",
                        "Horas_Prog": float(_prog),
                        "Horas_Reales": float(r.get("Horas_Reales", 0.0)),
                    })

        # Caso: TP con exceso -> split TP + TNPI
        elif tipo == "TP" and float(horas_prog) > 0 and float(horas_real) > float(horas_prog):
            exceso_h = max(0.0, float(horas_real) - float(horas_prog))
            tipo_exceso = st.session_state.get("exceso_tipo_general", "TNPI")  # Obtener el tipo seleccionado

            base = {
                "Equipo": equipo,
                "Pozo": pozo,
                "Etapa": etapa_use,
                "Fecha": str(fecha),
                "Equipo_Tipo": st.session_state.get("equipo_tipo_val", ""),
                "Modo_Reporte": modo_reporte,
                "Seccion": etapa,
                "Corrida": corrida,
                "Tipo_Agujero": tipo_agujero,
                "Operacion": operacion,
                "Actividad": actividad,
                "Turno": turno,
                "Hora_Inicio": hora_ini_txt,
                "Hora_Fin": hora_fin_txt,
                "ROP_Prog_mh": float(rop_prog),
                "ROP_Real_mh": float(rop_real),
                "Comentario": comentario,
                "Origen": "Manual",
            }

            # TP (hasta el estándar)
            add_rows.append({
                        **base,
                        "Tipo": "TP",
                        "Categoria_TNPI": "-",
                        "Detalle_TNPI": "-",
                        "Categoria_TNP": "-",
                        "Detalle_TNP": "-",
                        "Horas_Prog": float(horas_prog),
                        "Horas_Reales": float(horas_prog),
                    })

                        # TNPI/TNP por exceso: desglose guardado o asignación simple
            exceso_tipo = st.session_state.get("exceso_tipo_general", "TNPI")
            exceso_tipo = "TNP" if str(exceso_tipo).upper() == "TNP" else "TNPI"

            bd_saved = pd.DataFrame()
            if exceso_tipo == "TNP":
                # Registrar exceso como TNP (usa los selectores de Categoría/Detalle TNP)
                add_rows.append({
                    **base,
                    "Tipo": "TNP",
                    "Categoria_TNPI": "-",
                    "Detalle_TNPI": "-",
                    "Categoria_TNP": (act_cat_simple_tnp if "act_cat_simple_tnp" in locals() else categoria_tnp),
                    "Detalle_TNP": (act_det_simple_tnp if "act_det_simple_tnp" in locals() else detalle_tnp),
                    "Horas_Prog": 0.0,
                    "Horas_Reales": float(tnpi_exceso_h),
                    "Comentario": f"Exceso (Real {horas_real:.2f}h > Estándar {horas_prog:.2f}h) registrado como TNP.",
                    "Origen": "EXCESO",
                })
            else:
                # Registrar exceso como TNPI. Si hay desglose guardado, lo aplicamos.
                act_tnpi_breakdown = st.session_state.get("act_tnpi_breakdown", None)
                if act_tnpi_breakdown is not None and len(act_tnpi_breakdown) > 0:
                    try:
                        bd_saved = pd.DataFrame(act_tnpi_breakdown)
                        # Normaliza columnas esperadas
                        if "Categoria_TNPI" not in bd_saved.columns and "Categoria" in bd_saved.columns:
                            bd_saved["Categoria_TNPI"] = bd_saved["Categoria"]
                        if "Detalle_TNPI" not in bd_saved.columns and "Detalle" in bd_saved.columns:
                            bd_saved["Detalle_TNPI"] = bd_saved["Detalle"]
                        if "Horas_Reales" not in bd_saved.columns and "Horas" in bd_saved.columns:
                            bd_saved["Horas_Reales"] = bd_saved["Horas"]
                        bd_saved = bd_saved[["Categoria_TNPI", "Detalle_TNPI", "Horas_Reales"]].copy()
                        bd_saved["Horas_Reales"] = pd.to_numeric(bd_saved["Horas_Reales"], errors="coerce").fillna(0.0)
                    except Exception:
                        bd_saved = pd.DataFrame()
            
                if bd_saved.empty:
                    add_rows.append({
                        **base,
                        "Tipo": "TNPI",
                        "Categoria_TNPI": (act_cat_simple if "act_cat_simple" in locals() else categoria_tnpi),
                        "Detalle_TNPI": (act_det_simple if "act_det_simple" in locals() else detalle_tnpi),
                        "Categoria_TNP": "-",
                        "Detalle_TNP": "-",
                        "Horas_Prog": 0.0,
                        "Horas_Reales": float(tnpi_exceso_h),
                        "Comentario": f"Exceso (Real {horas_real:.2f}h > Estándar {horas_prog:.2f}h) registrado como TNPI.",
                        "Origen": "EXCESO",
                    })
            # Si existe desglose TNPI guardado, lo aplicamos como múltiples filas TNPI
            if (not bd_saved.empty) and (exceso_tipo == "TNPI"):
                sum_bd = float(bd_saved["Horas_Reales"].sum()) if "Horas_Reales" in bd_saved.columns else 0.0
                if sum_bd > 0 and abs(sum_bd - float(tnpi_exceso_h)) > 1e-6:
                    bd_saved["Horas_Reales"] = bd_saved["Horas_Reales"] * (float(tnpi_exceso_h) / sum_bd)
                for _, r in bd_saved.iterrows():
                    add_rows.append({
                        **base,
                        "Tipo": "TNPI",
                        "Categoria_TNPI": str(r.get("Categoria_TNPI", "-")),
                        "Detalle_TNPI": str(r.get("Detalle_TNPI", "-")),
                        "Categoria_TNP": "-",
                        "Detalle_TNP": "-",
                        "Horas_Prog": 0.0,
                        "Horas_Reales": float(r.get("Horas_Reales", 0.0)),
                        "Comentario": f"Exceso (Real {horas_real:.2f}h > Estándar {horas_prog:.2f}h) registrado como TNPI (desglose).",
                        "Origen": "EXCESO",
                    })

        else:
            # Sin split: respeta el tipo seleccionado (TP/TNPI/TNP)
            add_rows.append({
                "Equipo": equipo,
                "Pozo": pozo,
                "Etapa": etapa_use,
                "Fecha": str(fecha),
                "Equipo_Tipo": st.session_state.get("equipo_tipo_val", ""),
                "Modo_Reporte": modo_reporte,
                "Seccion": etapa,
                "Corrida": corrida,
                "Tipo_Agujero": tipo_agujero,
                "Operacion": operacion,
                "Actividad": actividad,
                "Turno": turno,
                "Hora_Inicio": hora_ini_txt,
                "Hora_Fin": hora_fin_txt,
                "Tipo": tipo,
                "Categoria_TNPI": categoria_tnpi if tipo == "TNPI" else "-",
                "Detalle_TNPI": detalle_tnpi if tipo == "TNPI" else "-",
                "Categoria_TNP": categoria_tnp if tipo == "TNP" else "-",
                "Detalle_TNP": detalle_tnp if tipo == "TNP" else "-",
                "Horas_Prog": float(horas_prog),
                "Horas_Reales": float(horas_real),
                "ROP_Prog_mh": float(rop_prog),
                "ROP_Real_mh": float(rop_real),
                "Comentario": comentario,
                "Origen": "Manual",
            })

        # Validación: no permitir que el día supere DAY_LIMIT_HOURS
        new_hours = float(sum([_safe_float(r.get("Horas_Reales", 0.0)) for r in add_rows]))
        remaining = _remaining_day_hours(st.session_state.df, fecha)
        if remaining <= 0:
            st.error(f"El día ya completó {DAY_LIMIT_HOURS:.0f}h. No se pueden agregar más actividades.")
            st.stop()
        if new_hours > remaining + 1e-6:
            st.error(f"No se puede agregar: quedan {remaining:.2f} h disponibles en el día.")
            st.stop()
        # Validación: no permitir que el turno supere TURNO_LIMIT_HOURS (12h)
        hrs_turno_actual = _day_used_hours_by_turno(st.session_state.df, fecha, turno)
        restante_turno = max(0.0, TURNO_LIMIT_HOURS - hrs_turno_actual)
        if new_hours > restante_turno + 1e-6:
            st.error(
                f"El turno **{turno}** ya tiene {hrs_turno_actual:.2f} h cargadas. "
                f"No se pueden cargar más de {TURNO_LIMIT_HOURS:.0f} h por turno (quedan {restante_turno:.2f} h)."
            )
            st.stop()

        st.session_state.df = pd.concat([st.session_state.df, pd.DataFrame(add_rows)], ignore_index=True)
        
        st.session_state.df = _ensure_rowid(st.session_state.df)
        st.session_state.df = _normalize_time_cause_columns(st.session_state.df)

        # Eliminar columnas duplicadas (puede ocurrir por compatibilidad / merges)
        if st.session_state.df.columns.duplicated().any():
            st.session_state.df = _coalesce_duplicate_columns(st.session_state.df)

        st.sidebar.success("Actividad agregada")
        st.rerun()


# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# CAPTURA ESPECIAL: CONEXIÓN PERFORANDO (MEJORADO - CON ETAPA ESPECÍFICA)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
if modo_reporte == "Perforación" and actividad == "Conexión perforando":
    with st.sidebar.expander("Conexión perforando (captura)", expanded=True):
        st.caption(f"Fecha en trabajo: {str(fecha)}")
        # Usa los inputs globales (etapa, corrida, tipo de agujero, turno)
        profundidad_m = st.number_input("Profundidad (m)", 0.0, step=1.0, key="prof_conn")
        
        conn_tipo = st.selectbox("Tipo de conexión", CONN_TYPE_OPTS, key="conn_tipo")
        ang_bucket = st.selectbox("Rango de ángulo", ANGLE_BUCKETS, key="ang_bucket")
        
        st.markdown("**Componentes (min reales)**")
        mins_real = {}
        for comp in CONN_COMPONENTS:
            mins_real[comp] = st.number_input(comp, min_value=0.0, step=0.1, value=0.0, key=f"min_{comp}")

        # Estándares (min) según tipo de conexión / ángulo
        std_map = CONN_STDS.get((conn_tipo, ang_bucket), {})
        std_pre = float(std_map.get("Preconexión", 0))
        std_conn = float(std_map.get("Conexión", 0))
        std_post = float(std_map.get("Postconexión", 0))
        std_total_line = float(std_map.get("TOTAL", std_pre + std_conn + std_post))
        total_real_min_ui = float(sum(mins_real.values())) if mins_real else 0.0
        st.caption(f"Total real: {total_real_min_ui:.1f} min")

        tipo_tiempo_conn = st.radio("Tipo de tiempo (Conexión)", options=["TP", "TNP"], horizontal=True, key="tipo_tiempo_conn")

        # Catálogo TNP (mismo archivo que TNPI)
        cat_list_tnp = sorted([c for c in df_tnpi_cat.get("Categoria_TNP", pd.Series(dtype=str)).dropna().astype(str).unique().tolist() if c.strip() != ""])
        if not cat_list_tnp:
            cat_list_tnp = ["-"]

        cat_tnpi_conn, det_tnpi_conn = "-", "-"
        cat_tnp_conn, det_tnp_conn = "-", "-"
        if tipo_tiempo_conn == "TP":
            if total_real_min_ui > std_total_line:
                st.markdown("**Exceso (Real > Estándar)**")
                exceso_policy_conn = st.radio(
                    "¿Cómo registrar el exceso?",
                    options=["TNPI", "TNP"],
                    horizontal=True,
                    key="conn_exceso_policy",
                    help="Si Real supera el estándar, el exceso puede registrarse como TNPI (no productivo/improductivo) o como TNP (no planeado).",
                )

                if exceso_policy_conn == "TNPI":
                    st.markdown("**Causa TNPI (solo para el exceso)**")
                    cat_tnpi_conn = st.selectbox("Categoría TNPI (exceso)", options=cat_list, key="conn_cat_tnpi")
                    det_all = df_tnpi_cat[df_tnpi_cat["Categoria_TNPI"] == cat_tnpi_conn]["Detalle_TNPI"].tolist()
                    q2 = (st.text_input("Buscar detalle TNPI (exceso)", value="", key="q_conn_tnpi") or "").strip().lower()
                    det_filtered = [d for d in det_all if q2 in str(d).lower()] if q2 else det_all
                    det_tnpi_conn = st.selectbox(
                        "Detalle TNPI (exceso)",
                        options=det_filtered if det_filtered else det_all,
                        key="det_conn_tnpi",
                    )
                else:
                    st.markdown("**Causa TNP (solo para el exceso)**")
                    cat_tnp_conn = st.selectbox("Categoría TNP (exceso)", options=cat_list_tnp, key="conn_cat_tnp")
                    det_all = df_tnpi_cat[df_tnpi_cat.get("Categoria_TNP", "") == cat_tnp_conn].get("Detalle_TNP", pd.Series(dtype=str)).tolist()
                    q2 = (st.text_input("Buscar detalle TNP (exceso)", value="", key="q_conn_tnp") or "").strip().lower()
                    det_filtered = [d for d in det_all if q2 in str(d).lower()] if q2 else det_all
                    det_tnp_conn = st.selectbox(
                        "Detalle TNP (exceso)",
                        options=det_filtered if det_filtered else det_all if det_all else ["-"],
                        key="det_conn_tnp",
                    )
            else:
                exceso_policy_conn = "TNPI"
                st.caption("No hay exceso: Real <= Estándar.")
        else:
            # Toda la conexión se registra como TNP (no hay TNPI automático aquí)
            exceso_policy_conn = "TNP"
            st.markdown("**Causa TNP (toda la conexión)**")
            cat_tnp_conn = st.selectbox("Categoría TNP", options=cat_list_tnp, key="conn_cat_tnp_full")
            det_all = df_tnpi_cat[df_tnpi_cat.get("Categoria_TNP", "") == cat_tnp_conn].get("Detalle_TNP", pd.Series(dtype=str)).tolist()
            q2 = (st.text_input("Buscar detalle TNP", value="", key="q_conn_tnp_full") or "").strip().lower()
            det_filtered = [d for d in det_all if q2 in str(d).lower()] if q2 else det_all
            det_tnp_conn = st.selectbox(
                "Detalle TNP",
                options=det_filtered if det_filtered else det_all if det_all else ["-"],
                key="det_conn_tnp_full",
            )
            cat_tnpi_conn, det_tnpi_conn = "-", "-"

        # Hora (opcional) para conexión
        conn_use_time = st.checkbox("Registrar hora (conexión)", value=False, key="conn_use_time")
        conn_hora_ini = None
        conn_hora_fin = None
        conn_bitacora_enabled = False
        conn_bitacora_entries = st.session_state.get("act_bitacora_entries", [])
        conn_bitacora_total_h = float(st.session_state.get("act_bitacora_total_h", 0.0) or 0.0)
        if conn_use_time:
            conn_hora_ini = st.time_input(
                "Hora inicio (conexión)",
                value=st.session_state.get("conn_hora_ini", datetime.now().time()),
                key="conn_hora_ini",
            )
            conn_hora_fin = st.time_input(
                "Hora fin (conexión)",
                value=st.session_state.get("conn_hora_fin", datetime.now().time()),
                key="conn_hora_fin",
            )
            conn_bitacora_enabled = st.toggle("Bitácora por horas (conexión)", value=False, key="conn_use_bitacora")
        conn_hora_ini_txt = conn_hora_ini.strftime("%H:%M") if (conn_use_time and conn_hora_ini) else ""
        conn_hora_fin_txt = conn_hora_fin.strftime("%H:%M") if (conn_use_time and conn_hora_fin) else ""
        if conn_use_time and conn_bitacora_enabled:
            st.caption("Completa la bitácora en la pestaña **Bitácora por horas**.")
            conn_bitacora_mode = st.radio(
                "Uso de bitácora (conexión)",
                options=["Usar bitácora como total de conexión", "Seguir registrando por componentes"],
                horizontal=True,
                key="conn_bitacora_mode",
                help="Define si la bitácora reemplaza el total real de conexión o solo sirve como referencia.",
            )
        else:
            conn_bitacora_mode = "Seguir registrando por componentes"

        conn_comment = st.text_input("Comentario conexión", "", key="conn_comment")

        if st.button("Agregar conexión", use_container_width=True):
            conn_id = str(uuid.uuid4())
            conn_no = int(st.session_state.df_conn["Conn_No"].max()) + 1 if not st.session_state.df_conn.empty else 1

            rows = []
            for comp in CONN_COMPONENTS:
                real = float(mins_real.get(comp, 0.0))
                if comp == "Preconexión":
                    std_use = std_pre
                elif comp == "Conexión":
                    std_use = std_conn
                elif comp == "Postconexión":
                    std_use = std_post
                else:
                    std_use = 0.0
                
                tnpi_exceso = 0.0
                minutos_tnp = 0.0
                if tipo_tiempo_conn == "TP":
                    exceso = max(0.0, float(real) - float(std_use))
                    if exceso_policy_conn == "TNP":
                        # El exceso se registra como TNP (no TNPI)
                        tnpi_exceso = 0.0
                        minutos_tnp = float(exceso)
                        cat_tnpi_use, det_tnpi_use = "-", "-"
                        cat_tnp_use, det_tnp_use = cat_tnp_conn, det_tnp_conn
                    else:
                        # El exceso se registra como TNPI
                        tnpi_exceso = float(exceso)
                        minutos_tnp = 0.0
                        cat_tnpi_use, det_tnpi_use = cat_tnpi_conn, det_tnpi_conn
                        cat_tnp_use, det_tnp_use = "-", "-"
                else:
                    # Toda la conexión como TNP
                    tnpi_exceso = 0.0
                    minutos_tnp = float(real)
                    cat_tnpi_use, det_tnpi_use = "-", "-"
                    cat_tnp_use, det_tnp_use = cat_tnp_conn, det_tnp_conn
                rows.append(
                    {
                        "Equipo": equipo,
                        "Pozo": pozo,
                        "Etapa": etapa,
                        "Fecha": str(fecha),
                        "Equipo_Tipo": st.session_state.get("equipo_tipo_val", ""),
                        "Seccion": etapa,  # También en Seccion
                        "Corrida": corrida,
                        "Tipo_Agujero": tipo_agujero,
                        "Turno": turno,
                        "Conn_No": conn_no,
                        "Profundidad_m": float(profundidad_m),
                        "Conn_Tipo": conn_tipo,
                        "Angulo_Bucket": ang_bucket,
                        "Componente": comp,
                        "Minutos_Reales": real,
                        "Minutos_Estandar": float(std_use),
                        "Minutos_TNPI": float(tnpi_exceso),
                        "Minutos_TNP": float(minutos_tnp),
                        "Categoria_TNPI": (cat_tnpi_use if float(tnpi_exceso) > 0 else "-"),
                        "Detalle_TNPI": (det_tnpi_use if float(tnpi_exceso) > 0 else "-"),
                        "Categoria_TNP": (cat_tnp_use if float(minutos_tnp) > 0 else "-"),
                        "Detalle_TNP": (det_tnp_use if float(minutos_tnp) > 0 else "-"),
                        "Comentario": conn_comment,
                        "CONN_ID": conn_id,
                    }
                )
            
            df_new = pd.DataFrame(rows)
            st.session_state.df_conn = pd.concat([st.session_state.df_conn, df_new], ignore_index=True)
            st.session_state["_toast_conn"] = True
            

            total_real_min = float(df_new["Minutos_Reales"].sum())
            if conn_use_time and conn_bitacora_enabled and conn_bitacora_total_h > 0 and conn_bitacora_mode == "Usar bitácora como total de conexión":
                total_real_min = float(conn_bitacora_total_h) * 60.0
                st.caption(f"Bitácora aplicada como total: {conn_bitacora_total_h:.2f} h")
                if float(df_new["Minutos_Reales"].sum() or 0.0) <= 0.0:
                    st.warning("Bitácora aplicada como total, pero los componentes están en 0. Esto afecta KPIs de conexiones.")
            std_total_line = float(std_map.get("TOTAL", std_pre + std_conn + std_post))
            exceso_total_min = max(0.0, total_real_min - std_total_line)

            # Parte base (hasta el estándar) siempre conserva el tipo seleccionado (TP o TNP).
            base_min = min(total_real_min, std_total_line)

            # ¿Cómo registrar el exceso?
            # - Si el usuario eligió registrar el exceso como TNPI -> lo mandamos a Minutos_TNPI
            # - Si eligió TNP -> lo mandamos a Minutos_TNP y guardamos categoría/detalle TNP
            if tipo_tiempo_conn == "TP":
                if exceso_policy_conn == "TNPI":
                    tnpi_min = exceso_total_min
                    tnp_min = 0.0
                else:  # "TNP"
                    tnpi_min = 0.0
                    tnp_min = exceso_total_min
            else:
                # Si la conexión completa se está registrando como TNP, no hay desglose de exceso.
                tnpi_min = 0.0
                tnp_min = float(total_real_min)

            base = dict(
                Equipo=equipo,
                Pozo=pozo,
                Etapa=etapa,
                Fecha=str(fecha),
                Equipo_Tipo=st.session_state.get("equipo_tipo_val", ""),
                Modo_Reporte="Perforación",
                Seccion=etapa,
                Corrida=corrida,
                Tipo_Agujero=tipo_agujero,
                Operacion="Perforación",
                Actividad=f"Conexión perforando ({ang_bucket})",
                Turno=turno,
                Hora_Inicio=conn_hora_ini_txt,
                Hora_Fin=conn_hora_fin_txt,
                Tipo=tipo_tiempo_conn,
                Categoria_TNPI="-",
                Detalle_TNPI="-",
                Categoria_TNP=(cat_tnp_conn if tipo_tiempo_conn == "TNP" else "-"),
                Detalle_TNP=(det_tnp_conn if tipo_tiempo_conn == "TNP" else "-"),
                Horas_Prog=float(std_total_line / 60.0) if std_total_line else 0.0,
                Horas_Reales=float(base_min / 60.0),
                ROP_Prog_mh=0.0,
                ROP_Real_mh=0.0,
                Comentario=st.session_state.get("comentario_conn", "") or "",
                Origen="Manual",
                CONN_ID=conn_id,
                Eficiencia_pct=float(_calc_eff(std_total_line / 60.0, base_min / 60.0)),
                Semáforo=_semaforo_text(float(_calc_eff(std_total_line / 60.0, base_min / 60.0))),
            )

            add_rows = [base]

            # Fila de exceso (solo si la base era TP y hubo exceso)
            if tipo_tiempo_conn == "TP" and exceso_total_min > 0:
                if exceso_policy_conn == "TNPI" and tnpi_min > 0:
                    add_rows.append({**base,
                        "Tipo": "TNPI",
                        "Categoria_TNPI": (cat_tnpi_use or "-"),
                        "Detalle_TNPI": (det_tnpi_use or "-"),
                        "Categoria_TNP": "-",
                        "Detalle_TNP": "-",
                        "Horas_Prog": 0.0,
                        "Horas_Reales": float(tnpi_min / 60.0),
                        "Eficiencia_pct": 0.0,
                        "Semáforo": _semaforo_text(0.0),
                    })
                if exceso_policy_conn == "TNP" and tnp_min > 0:
                    add_rows.append({**base,
                        "Tipo": "TNP",
                        "Categoria_TNPI": "-",
                        "Detalle_TNPI": "-",
                        "Categoria_TNP": (cat_tnp_conn or "-"),
                        "Detalle_TNP": (det_tnp_conn or "-"),
                        "Horas_Prog": 0.0,
                        "Horas_Reales": float(tnp_min / 60.0),
                        "Eficiencia_pct": 0.0,
                        "Semáforo": _semaforo_text(0.0),
                    })

            # Validación: no permitir que el día supere DAY_LIMIT_HOURS
            new_hours = float(sum([_safe_float(r.get("Horas_Reales", 0.0)) for r in add_rows]))
            remaining = _remaining_day_hours(st.session_state.df, fecha)
            if remaining <= 0:
                st.error(f"El día ya completó {DAY_LIMIT_HOURS:.0f}h. No se pueden agregar más conexiones.")
                st.stop()
            if new_hours > remaining + 1e-6:
                st.error(f"No se puede agregar: quedan {remaining:.2f} h disponibles en el día.")
                st.stop()
            # Validación: no permitir que el turno supere TURNO_LIMIT_HOURS (12h)
            hrs_turno_conn = _day_used_hours_by_turno(st.session_state.df, fecha, turno)
            restante_turno_conn = max(0.0, TURNO_LIMIT_HOURS - hrs_turno_conn)
            if new_hours > restante_turno_conn + 1e-6:
                st.error(
                    f"El turno **{turno}** ya tiene {hrs_turno_conn:.2f} h cargadas. "
                    f"No se pueden cargar más de {TURNO_LIMIT_HOURS:.0f} h por turno (quedan {restante_turno_conn:.2f} h)."
                )
                st.stop()

            st.session_state.df = pd.concat([st.session_state.df, pd.DataFrame(add_rows)], ignore_index=True)
            st.sidebar.success("Conexión agregada")
            st.rerun()
            
        st.session_state.df = _ensure_rowid(st.session_state.df)

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# CAPTURA ESPECIAL: ARMA/DESARMA BHA
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
if actividad == "Arma/Desarma BHA":
    with st.sidebar.expander("Arma/Desarma BHA (captura)", expanded=True):
        st.caption(f"Fecha en trabajo: {str(fecha)}")
        bha_tipo_tiempo = st.session_state.get("tipo_time_general", "TP")
        st.session_state["bha_tipo_tiempo"] = bha_tipo_tiempo

        bha_turno = st.radio("Turno (BHA)", TURNOS, horizontal=True, key="bha_turno")
        barrena = st.session_state.get("barrena_global", "")
        bha_tipo = st.selectbox("Tipo (1–10)", options=list(BHA_TYPES.keys()), index=0, key="bha_tipo")

        desc, std_arma, std_desarma = BHA_TYPES[int(bha_tipo)]
        accion = st.radio("Acción", ["Arma", "Desarma"], horizontal=True, key="bha_accion")

        std_default = float(std_arma if accion == "Arma" else std_desarma)
        override = st.checkbox("Editar estándar manualmente", value=False, key="bha_override")
        if override:
            estandar_h = st.number_input("Estándar (h)", min_value=0.0, step=0.25, value=float(std_default), key="bha_std_manual")
        else:
            estandar_h = float(std_default)
            st.caption(f"Estándar automático: **{estandar_h:.2f} h**")

        label_real_bha = "Real (h)"
        if st.session_state.get("bha_tipo_tiempo", "TP") == "TNP":
            label_real_bha = "TNP (h) - tiempo real"
        elif st.session_state.get("bha_tipo_tiempo", "TP") == "TNPI":
            label_real_bha = "TNPI (h) - tiempo real"
        elif st.session_state.get("bha_tipo_tiempo", "TP") == "TP":
            label_real_bha = "TP (h) - tiempo real"

        # Hora (opcional) para BHA + bitácora
        bha_use_time = st.checkbox("Registrar hora (BHA)", value=False, key="bha_use_time")
        bha_hora_ini = None
        bha_hora_fin = None
        bha_bitacora_enabled = False
        bha_bitacora_entries = st.session_state.get("act_bitacora_entries", [])
        bha_bitacora_total_h = float(st.session_state.get("act_bitacora_total_h", 0.0) or 0.0)
        if bha_use_time:
            bha_hora_ini = st.time_input(
                "Hora inicio (BHA)",
                value=st.session_state.get("bha_hora_ini", datetime.now().time()),
                key="bha_hora_ini",
            )
            bha_hora_fin = st.time_input(
                "Hora fin (BHA)",
                value=st.session_state.get("bha_hora_fin", datetime.now().time()),
                key="bha_hora_fin",
            )
            bha_bitacora_enabled = st.toggle("Bitácora por horas (BHA)", value=False, key="bha_use_bitacora")
        bha_hora_ini_txt = bha_hora_ini.strftime("%H:%M") if (bha_use_time and bha_hora_ini) else ""
        bha_hora_fin_txt = bha_hora_fin.strftime("%H:%M") if (bha_use_time and bha_hora_fin) else ""
        if bha_use_time and bha_bitacora_enabled:
            st.caption("Completa la bitácora en la pestaña **Bitácora por horas**.")

        if bha_use_time:
            if bha_bitacora_enabled and bha_bitacora_total_h > 0:
                real_h = float(bha_bitacora_total_h)
                st.caption(f"Horas reales (bitácora): {real_h:.2f} h")
            elif bha_hora_ini and bha_hora_fin:
                _dt_ini = datetime.combine(datetime.today().date(), bha_hora_ini)
                _dt_fin = datetime.combine(datetime.today().date(), bha_hora_fin)
                real_h = (_dt_fin - _dt_ini).total_seconds() / 3600.0
                if real_h < 0:
                    real_h += 24.0
                st.caption(f"Horas reales calculadas: {real_h:.2f} h")
            else:
                real_h = 0.0
                st.caption("Horas reales calculadas: 0.00 h")
        else:
            real_h = st.number_input(label_real_bha, min_value=0.0, step=0.25, value=0.0, key="bha_real_h")

        tnpi_h = max(0.0, float(real_h) - float(estandar_h))
        tnp_h = 0.0
        if (st.session_state.get("tipo_time_bha") or st.session_state.get("tipo_time_general") or "TP") == "TNP":
            tnp_h = float(real_h)
            tnpi_h = 0.0

        if bha_use_time and bha_bitacora_enabled and bha_bitacora_total_h > 0:
            st.caption(
                f"Bitácora aplicada: {bha_bitacora_total_h:.2f} h "
                f"({'total' if bha_bitacora_enabled else 'referencia'})"
            )

            tnp_h = float(real_h)
            tnpi_h = 0.0
        tp_h_local = max(0.0, float(real_h) - float(tnpi_h))
        eff_bha = clamp_0_100(safe_pct(tp_h_local, float(real_h))) if real_h > 0 else 0.0

        st.caption(f"Componentes: **{desc}**")
        st.caption(f"TNPI por exceso: **{tnpi_h:.2f} h** | Eficiencia: **{eff_bha:.0f}%**")

        bha_cat = "-"
        bha_det = "-"

        # --- NUEVO: Seleccionar tipo de exceso para BHA ---

        if tnpi_h > 0 and bha_tipo_tiempo == "TP":

            st.markdown("**Exceso (BHA)**")

            

            bha_exceso_tipo = st.radio(

                "¿Cómo registrar el exceso?",

                options=["TNPI", "TNP"],

                horizontal=True,

                key="bha_exceso_tipo",

                help="El exceso de horas puede registrarse como TNPI (no productivo/improductivo) o como TNP (no planeado)."

            )

            

            # Actualizar variable para uso posterior
        # session_state already managed by widget key
            # Mostrar configuración según el tipo seleccionado

            if bha_exceso_tipo == "TNPI":

                st.markdown("**Asignación simple (una sola causa)**")

                bha_cat = st.selectbox(

                    "Categoría TNPI (BHA)",

                    options=cat_list if "cat_list" in globals() else ["-"],

                    index=0,

                    key="bha_cat_simple_1",

                )



                if "df_tnpi_cat" in globals() and "Categoria_TNPI" in df_tnpi_cat.columns and "Detalle_TNPI" in df_tnpi_cat.columns:

                    _det_opts = (

                        df_tnpi_cat[df_tnpi_cat["Categoria_TNPI"] == bha_cat]["Detalle_TNPI"]

                        .dropna()

                        .unique()

                        .tolist()

                    )

                else:

                    _det_opts = ["-"]



                bha_det = st.selectbox(

                    "Detalle TNPI (BHA)",

                    options=_det_opts if len(_det_opts) else ["-"],

                    index=0,

                    key="bha_det_simple_1",

                )

            else:  # TNP

                # Configuración para TNP

                st.markdown("**Asignación simple para TNP (una sola causa)**")

                bha_cat = st.selectbox(

                    "Categoría TNP (BHA)",

                    options=tnp_cat_list if "tnp_cat_list" in globals() else ["-"],

                    index=0,

                    key="bha_cat_simple_2",

                )

                

                if "df_tnp_cat" in globals() and "Categoria_TNP" in df_tnp_cat.columns and "Detalle_TNP" in df_tnp_cat.columns:

                    _det_opts = (

                        df_tnp_cat[df_tnp_cat["Categoria_TNP"] == bha_cat]["Detalle_TNP"]

                        .dropna()

                        .unique()

                        .tolist()

                    )

                else:

                    _det_opts = ["-"]

                

                bha_det = st.selectbox(

                    "Detalle TNP (BHA)",

                    options=_det_opts if len(_det_opts) else ["-"],

                    index=0,

                    key="bha_det_simple_2",

                )

            # Nota: no sobrescribir la selección TNP con un selector TNPI.

            # --- Desglose opcional (múltiples causas) ---
            with st.expander("Detalle TNPI (opcional) - desglose por múltiples causas", expanded=False):
                st.caption("Agrega varias causas y presiona **Guardar cambios**. La suma debe ser igual al TNPI calculado.")

                # Inicializar borrador si está vacío y no hay nada guardado
                if (not st.session_state.get("bha_tnpi_breakdown_draft")) and (not st.session_state.get("bha_tnpi_breakdown")):
                    st.session_state.bha_tnpi_breakdown_draft = [{
                        "Categoria_TNPI": (cat_list[0] if "cat_list" in globals() and isinstance(cat_list, list) and len(cat_list) else "-"),
                        "Detalle_TNPI": "-",
                        "Horas_TNPI_h": float(tnpi_h),
                        "Comentario": "",
                    }]

                _cat_new = st.selectbox(
                    "Categoría TNPI (nuevo renglón)",
                    options=cat_list if "cat_list" in globals() else ["-"],
                    index=0,
                    key="bha_bd_cat_new",
                )

                if "df_tnpi_cat" in globals() and "Categoria_TNPI" in df_tnpi_cat.columns and "Detalle_TNPI" in df_tnpi_cat.columns:
                    _det_new_opts = (
                        df_tnpi_cat[df_tnpi_cat["Categoria_TNPI"] == _cat_new]["Detalle_TNPI"]
                        .dropna()
                        .unique()
                        .tolist()
                    )
                else:
                    _det_new_opts = ["-"]

                _det_new = st.selectbox(
                    "Detalle TNPI (nuevo renglón)",
                    options=_det_new_opts if len(_det_new_opts) else ["-"],
                    index=0,
                    key="bha_bd_det_new",
                )

                _h_new = st.number_input(
                    "Horas TNPI (h) (nuevo renglón)",
                    min_value=0.0,
                    value=0.0,
                    step=0.25,
                    format="%.2f",
                    key="bha_bd_h_new",
                )

                _com_new = st.text_input("Comentario (nuevo renglón)", "", key="bha_bd_com_new")

                if st.button("Agregar", use_container_width=True, key="bha_bd_add"):
                        if _h_new <= 0:
                            st.warning("Las horas deben ser > 0 para agregar la causa.")
                        else:
                            st.session_state.bha_tnpi_breakdown_draft.append({
                                "Categoria_TNPI": _cat_new,
                                "Detalle_TNPI": _det_new,
                                "Horas_TNPI_h": float(_h_new),
                                "Comentario": _com_new,
                            })
                            st.rerun()

                if st.button("Limpiar", use_container_width=True, key="bha_bd_clear"):
                    st.session_state.bha_tnpi_breakdown_draft = []
                    st.session_state.bha_tnpi_breakdown_saved = False
                    st.rerun()

                bd = pd.DataFrame(
                    st.session_state.get("bha_tnpi_breakdown_draft", []),
                    columns=["Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP", "Detalle_TNP", "Horas_TNPI_h", "Comentario"]
                )
                bd["Horas_TNPI_h"] = pd.to_numeric(bd["Horas_TNPI_h"], errors="coerce").fillna(0.0)

                st.markdown("**Causas en borrador**")
                if bd.empty:
                    st.info("Aún no hay causas en el borrador.")
                else:
                    for i, r in bd.reset_index(drop=True).iterrows():
                        cols = st.columns([2, 3, 1, 3, 1])
                        cols[0].write(str(r.get("Categoria_TNPI", "-")))
                        cols[1].write(str(r.get("Detalle_TNPI", "-")))
                        cols[2].write(f'{float(r.get("Horas_TNPI_h", 0.0)):.2f}')
                        cols[3].write(str(r.get("Comentario", "")))
                        if cols[4].button("Eliminar", key=f"bha_bd_del_{i}"):
                            st.session_state.bha_tnpi_breakdown_draft.pop(i)
                            st.rerun()

                sum_bd = float(bd["Horas_TNPI_h"].sum()) if not bd.empty else 0.0
                st.caption(f"Suma borrador: **{sum_bd:.2f} h**  |  TNPI calculado: **{float(tnpi_h):.2f} h**")

                if st.button("Guardar", use_container_width=True, key="bha_bd_save"):
                    if bd.empty:
                        st.error("No hay renglones en el borrador para guardar.")
                    else:
                        bd2 = bd[bd["Horas_TNPI_h"] > 0].copy()
                        sum2 = float(bd2["Horas_TNPI_h"].sum())
                        if abs(sum2 - float(tnpi_h)) > 1e-6:
                            st.error(f"La suma del desglose ({sum2:.2f} h) debe ser igual al TNPI calculado ({float(tnpi_h):.2f} h).")
                        else:
                            st.session_state.bha_tnpi_breakdown = bd2.to_dict(orient="records")
                            st.session_state.bha_tnpi_breakdown_saved = True
                            st.success("Desglose guardado. Al agregar el BHA se registrarán varias causas TNPI.")

        bha_comment = st.text_input("Comentario BHA", "", key="bha_comment")

        if st.button("Agregar BHA", use_container_width=True):
            bha_id = str(uuid.uuid4())
            row_bha = {
                "Equipo": equipo,
                "Pozo": pozo,
                "Etapa": ((etapa_viajes_sel or etapa) if "etapa_viajes_sel" in globals() else etapa),
                "Fecha": str(fecha),
                "Turno": bha_turno,
                "Barrena": barrena,
                "BHA_Tipo": int(bha_tipo),
                "BHA_Componentes": desc,
                "Accion": accion,
                "BHA_ID": bha_id,
                "Estandar_h": float(estandar_h),
                "Real_h": float(real_h),
                "TNPI_h": float(tnpi_h),
                "TNP_h": float(tnp_h),
                "Eficiencia_pct": float(_calc_eff(estandar_h, real_h)),
            }
            st.session_state.df_bha = pd.concat([st.session_state.df_bha, pd.DataFrame([row_bha])], ignore_index=True)

            base = dict(
                Equipo=equipo,
                Pozo=pozo,
                Etapa=etapa,
                Fecha=str(fecha),
                Equipo_Tipo=st.session_state.get("equipo_tipo_val", ""),
                Modo_Reporte=modo_reporte,
                Seccion=etapa,
                Corrida=corrida,
                Tipo_Agujero=tipo_agujero,
                Operacion=operacion,
                Actividad=f"Arma/Desarma BHA (Tipo {int(bha_tipo)})",
                Turno=bha_turno,
                Hora_Inicio=bha_hora_ini_txt,
                Hora_Fin=bha_hora_fin_txt,
                ROP_Prog_mh=0.0,
                ROP_Real_mh=0.0,
                Comentario=bha_comment.strip(),
                Origen="BHA",
                BHA_ID=bha_id,
            )

            add = [
                {
                    **base,
                    "Tipo": "TP",
                    "Categoria_TNPI": "-",
                    "Detalle_TNPI": "-",
                    "Horas_Prog": float(estandar_h),
                    "Horas_Reales": float(max(0.0, float(real_h) - float(tnpi_h))),
                }
            ]
            if tnpi_h > 0 and bha_tipo_tiempo == "TP":
                # Si hay desglose guardado, úsalo. Si no, usa asignación simple (una sola causa).
                bd_saved = pd.DataFrame(
                    st.session_state.get("bha_tnpi_breakdown", []),
                    columns=["Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP", "Detalle_TNP", "Horas_TNPI_h", "Comentario"]
                )
                if (not bd_saved.empty) and bool(st.session_state.get("bha_tnpi_breakdown_saved", False)):
                    bd_saved["Horas_TNPI_h"] = pd.to_numeric(bd_saved["Horas_TNPI_h"], errors="coerce").fillna(0.0)
                    bd_saved = bd_saved[bd_saved["Horas_TNPI_h"] > 0].copy()

                    sum_bd = float(bd_saved["Horas_TNPI_h"].sum())
                    if abs(sum_bd - float(tnpi_h)) > 1e-6:
                        st.error(f"La suma del desglose TNPI guardado ({sum_bd:.2f} h) debe ser igual al TNPI calculado ({float(tnpi_h):.2f} h).")
                        st.stop()

                    for _, r in bd_saved.iterrows():
                        add.append(
                            {
                                **base,
                                "Tipo": "TNPI",
                                "Categoria_TNPI": str(r.get("Categoria_TNPI", "-") or "-"),
                                "Detalle_TNPI": str(r.get("Detalle_TNPI", "-") or "-"),
                                "Horas_Prog": 0.0,
                                "Horas_Reales": float(r.get("Horas_TNPI_h", 0.0) or 0.0),
                                "Comentario": (base.get("Comentario","") + (f" | {r.get('Comentario')}" if r.get("Comentario") else "")).strip(" |"),
                            }
                        )
                else:
                    # Asignación simple (una sola causa)
                    add.append(
                        {
                            **base,
                            "Tipo": "TNPI",
                            "Categoria_TNPI": bha_cat,
                            "Detalle_TNPI": bha_det,
                            "Horas_Prog": 0.0,
                            "Horas_Reales": float(tnpi_h),
                        }
                    )

                # Limpia desglose para el siguiente BHA
                st.session_state.bha_tnpi_breakdown = []
                st.session_state.bha_tnpi_breakdown_draft = []
                st.session_state.bha_tnpi_breakdown_saved = False

            # Si hay bitácora por horas, dividir en tramos y repartir horas proporcionalmente
            if bha_use_time and bha_bitacora_enabled and bha_bitacora_total_h > 0 and len(bha_bitacora_entries) > 0:
                total_bit = float(sum(r["Horas_Reales"] for r in bha_bitacora_entries))
                new_rows = []
                for r in bha_bitacora_entries:
                    _factor = float(r.get("Horas_Reales", 0.0)) / total_bit if total_bit > 0 else 0.0
                    for row in add:
                        _hr = float(row.get("Horas_Reales", 0.0))
                        if _hr <= 0 or _factor <= 0:
                            continue
                        _row = dict(row)
                        _row["Hora_Inicio"] = str(r.get("Hora_Inicio", ""))
                        _row["Hora_Fin"] = str(r.get("Hora_Fin", ""))
                        _row["Horas_Reales"] = _hr * _factor
                        if "Horas_Prog" in _row:
                            _row["Horas_Prog"] = float(_row.get("Horas_Prog", 0.0)) * _factor
                        _com = str(r.get("Comentario", "") or "").strip()
                        if _com:
                            _row["Comentario"] = (str(_row.get("Comentario", "") or "") + f" | {_com}").strip(" |")
                        new_rows.append(_row)
                add = new_rows

            # Validación: no permitir que el día supere DAY_LIMIT_HOURS
            new_hours = float(sum([_safe_float(r.get("Horas_Reales", 0.0)) for r in add]))
            remaining = _remaining_day_hours(st.session_state.df, fecha)
            if remaining <= 0:
                st.error(f"El día ya completó {DAY_LIMIT_HOURS:.0f}h. No se pueden agregar más actividades.")
                st.stop()
            if new_hours > remaining + 1e-6:
                st.error(f"No se puede agregar: quedan {remaining:.2f} h disponibles en el día.")
                st.stop()
            # Validación: no permitir que el turno BHA supere TURNO_LIMIT_HOURS (12h)
            hrs_bha_turno = _day_used_hours_by_turno(st.session_state.df, fecha, bha_turno)
            restante_bha_turno = max(0.0, TURNO_LIMIT_HOURS - hrs_bha_turno)
            if new_hours > restante_bha_turno + 1e-6:
                st.error(
                    f"El turno **{bha_turno}** ya tiene {hrs_bha_turno:.2f} h cargadas. "
                    f"No se pueden cargar más de {TURNO_LIMIT_HOURS:.0f} h por turno (quedan {restante_bha_turno:.2f} h)."
                )
                st.stop()

            st.session_state.df = pd.concat([st.session_state.df, pd.DataFrame(add)], ignore_index=True)
            st.success("BHA agregado")
            st.rerun()

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# MAIN DATA
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
df = st.session_state.df.copy()
df_conn = st.session_state.df_conn.copy()
df_bha = st.session_state.df_bha.copy()

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# BHA: GRAFICA ESTÁNDAR VS REAL (cuando estás capturando Arma/Desarma)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# Nota: se muestra solo cuando en el sidebar eliges la actividad "Arma/Desarma BHA"


# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# KPIs base
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
total_prog = float(df["Horas_Prog"].sum()) if not df.empty else 0.0
total_real = float(df["Horas_Reales"].sum()) if not df.empty else 0.0
tp_h = float(df[df["Tipo"] == "TP"]["Horas_Reales"].sum()) if not df.empty else 0.0
tnpi_h = float(df[df["Tipo"] == "TNPI"]["Horas_Reales"].sum()) if not df.empty else 0.0
tnp_h = float(df[df["Tipo"] == "TNP"]["Horas_Reales"].sum()) if not df.empty else 0.0
eficiencia_dia = clamp_0_100(safe_pct(tp_h, total_real)) if total_real > 0 else 0.0

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# METROS / ROP (IMPORTANTE: define variables SIEMPRE)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
mr_total = 0.0
tnpi_m_h = 0.0
eff_m = 0.0
rr = 0.0
eff_rop = 0.0

if modo_reporte == "Perforación":
    mp = _safe_float(st.session_state.drill_day.get("metros_prog_total", 0.0))
    mr_d = _safe_float(st.session_state.drill_day.get("metros_real_dia", 0.0))
    mr_n = _safe_float(st.session_state.drill_day.get("metros_real_noche", 0.0))
    mr_total = mr_d + mr_n

    tnpi_m_h = _safe_float(st.session_state.drill_day.get("tnpi_metros_h", 0.0))
    eff_m = clamp_0_100(safe_pct(mr_total, mp)) if mp > 0 else 0.0

    rp = _safe_float(st.session_state.drill_day.get("rop_prog_total", 0.0))
    rr_d = _safe_float(st.session_state.drill_day.get("rop_real_dia", 0.0))
    rr_n = _safe_float(st.session_state.drill_day.get("rop_real_noche", 0.0))
    rr = (rr_d + rr_n) / (2 if (rr_d > 0 and rr_n > 0) else 1) if (rr_d > 0 or rr_n > 0) else 0.0
    eff_rop = clamp_0_100(safe_pct(rr, rp)) if rp > 0 else 0.0

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# KPI CONEXIONES (IMPORTANTE: define variables SIEMPRE)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
conn_real_min = 0.0
conn_std_min = 0.0
conn_tnpi_min = 0.0
eff_conn = 0.0

if modo_reporte == "Perforación" and not df_conn.empty:
    conn_real_min = float(df_conn.groupby(["Conn_No"])["Minutos_Reales"].sum().sum())
    per_conn2 = df_conn.groupby("Conn_No", as_index=False).first()[["Conn_No", "Conn_Tipo", "Angulo_Bucket"]]
    per_conn2["Std_Total"] = per_conn2.apply(
        lambda r: float(CONN_STDS.get((r["Conn_Tipo"], r["Angulo_Bucket"]), {}).get("TOTAL", 0.0)),
        axis=1
    )
    conn_std_min = float(per_conn2["Std_Total"].sum())

    conn_tp_min = min(conn_real_min, conn_std_min) if conn_std_min > 0 else conn_real_min
    conn_tnpi_min = max(0.0, conn_real_min - conn_std_min) if conn_std_min > 0 else 0.0
    eff_conn = clamp_0_100(safe_pct(conn_tp_min, conn_real_min)) if conn_real_min > 0 else 0.0


# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# DrillSpot KPI Export (XLSX) -> Viajes & Conexiones (por hora)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
def _clean_drillspot_kpi_df(df_raw: pd.DataFrame) -> pd.DataFrame:
    """
    Espera el formato típico del export 'KPI Report' de DrillSpot:
    columnas: Start Time, End Time, Start Bit Depth, End Bit Depth, KPI, Duration, ...
    Nota: la primera fila suele traer unidades ('date','m','name','min', etc). Se elimina.
    """
    if df_raw is None or df_raw.empty:
        return pd.DataFrame()

    df = df_raw.copy()

    # Normaliza nombres (por si vienen con espacios raros)
    df.columns = [str(c).strip() for c in df.columns]

    # Quita primera fila de unidades si aplica
    if "Start Time" in df.columns:
        first = str(df.iloc[0]["Start Time"]).strip().lower()
        if first in {"date", "datetime"}:
            df = df.iloc[1:].reset_index(drop=True)

    # Tipos
    for c in ["Start Time", "End Time"]:
        if c in df.columns:
            df[c] = pd.to_datetime(df[c], errors="coerce")

    for c in ["Start Bit Depth", "End Bit Depth", "Duration"]:
        if c in df.columns:
            df[c] = pd.to_numeric(df[c], errors="coerce")

    if "KPI" in df.columns:
        df["KPI"] = df["KPI"].astype(str)

    df = df.dropna(subset=["Start Time", "End Time", "KPI", "Duration"]).reset_index(drop=True)
    return df


def load_drillspot_kpi_xlsx(uploaded_file) -> pd.DataFrame:
    """Lee el XLSX exportado por DrillSpot y devuelve un DataFrame limpio."""
    if uploaded_file is None:
        return pd.DataFrame()
    try:
        xls = pd.ExcelFile(uploaded_file)
        sheet = "KPI Report" if "KPI Report" in xls.sheet_names else xls.sheet_names[0]
        df_raw = xls.parse(sheet)
        return _clean_drillspot_kpi_df(df_raw)
    except Exception:
        return pd.DataFrame()


def compute_viaje_conn_hourly_from_kpi(
    df_kpi: pd.DataFrame,
    direction: str,
) -> tuple[pd.DataFrame, dict]:
    """
    direction: 'Trip In' o 'Trip Out'
    Retorna:
      - hourly_df con columnas: hour (0-23), speed_mh (real), conn_min (real)
      - meta dict: distance_m_total, running_minutes_total, conn_events_total
    """
    if df_kpi is None or df_kpi.empty:
        hourly = pd.DataFrame({"hour": list(range(24)), "speed_mh": [0.0]*24, "conn_min": [0.0]*24})
        return hourly, {"distance_m_total": 0.0, "running_minutes_total": 0.0, "conn_events_total": 0}

    df = df_kpi.copy()

    # Filtra KPIs
    run_key = f"{direction}: Running"
    conn_key = f"{direction}: Connection"

    df_run = df[df["KPI"].str.contains(run_key, case=False, na=False)].copy()
    df_conn = df[df["KPI"].str.contains(conn_key, case=False, na=False)].copy()

    # Running -> velocidad (m/h) por hora (ponderado por tiempo)
    if not df_run.empty:
        df_run["hour"] = df_run["Start Time"].dt.hour.astype(int)
        df_run["dist_m"] = (df_run["End Bit Depth"] - df_run["Start Bit Depth"]).abs()
        df_run["dur_h"] = (df_run["Duration"] / 60.0).replace(0, np.nan)
        df_run["speed_mh"] = (df_run["dist_m"] / df_run["dur_h"]).replace([np.inf, -np.inf], np.nan).fillna(0.0)

        g = df_run.groupby("hour", as_index=False).agg(
            dist_m=("dist_m", "sum"),
            dur_h=("dur_h", "sum"),
        )
        g["speed_mh"] = g.apply(lambda r: (r["dist_m"] / r["dur_h"]) if r["dur_h"] and r["dur_h"] > 0 else 0.0, axis=1)
        speed_map = dict(zip(g["hour"].astype(int), g["speed_mh"].astype(float)))
        dist_total = float(df_run["dist_m"].sum())
        run_min_total = float(df_run["Duration"].sum())
    else:
        speed_map = {}
        dist_total = 0.0
        run_min_total = 0.0

    # Connection -> minutos promedio por hora (promedio por evento)
    if not df_conn.empty:
        df_conn["hour"] = df_conn["Start Time"].dt.hour.astype(int)
        g2 = df_conn.groupby("hour", as_index=False).agg(
            conn_min=("Duration", "mean"),
            n=("Duration", "count"),
        )
        conn_map = dict(zip(g2["hour"].astype(int), g2["conn_min"].astype(float)))
        conn_events_total = int(df_conn.shape[0])
    else:
        conn_map = {}
        conn_events_total = 0

    rows = []
    for h in range(24):
        rows.append(
            {
                "hour": h,
                "speed_mh": float(speed_map.get(h, 0.0)),
                "conn_min": float(conn_map.get(h, 0.0)),
            }
        )
    hourly = pd.DataFrame(rows)
    return hourly, {
        "distance_m_total": dist_total,
        "running_minutes_total": run_min_total,
        "conn_events_total": conn_events_total,
    }


def default_trip_direction_from_activity(activity_name: str) -> str:
    """Heurística simple para mapear tu 'Viaje ...' a Trip In / Trip Out del export de KPIs."""
    a = (activity_name or "").lower()
    if any(k in a for k in ["metiendo", "bajando", "entrando"]):
        return "Trip In"
    if any(k in a for k in ["sacando", "levantando", "subiendo", "saliendo"]):
        return "Trip Out"
    # fallback
    return "Trip In"

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# CACHE: generar figuras (reduce lentitud)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
@st.cache_data(show_spinner=False)
def _make_figs(df_json: str, df_conn_json: str, modo_reporte: str):
    df_local = pd.read_json(df_json, orient="split") if df_json else pd.DataFrame()
    dfc_local = pd.read_json(df_conn_json, orient="split") if df_conn_json else pd.DataFrame()

    figs = {"tiempos": None, "act_pie": None, "act_bar": None, "conn_pie": None, "conn_stack": None}

    # tiempos
    if not df_local.empty and {"Tipo", "Horas_Reales"}.issubset(df_local.columns):
        figs["tiempos"] = px.pie(df_local, names="Tipo", values="Horas_Reales", hole=0.55, title="TP vs TNPI vs TNP")

    # actividades
    if not df_local.empty and {"Actividad", "Horas_Reales"}.issubset(df_local.columns):
        df_act = df_local.groupby("Actividad", as_index=False)["Horas_Reales"].sum().sort_values("Horas_Reales", ascending=False)
        figs["act_pie"] = px.pie(df_act, names="Actividad", values="Horas_Reales", hole=0.35, title="Horas por actividad")

        palette = px.colors.qualitative.Set3 + px.colors.qualitative.Pastel + px.colors.qualitative.Bold
        act_names = df_act["Actividad"].tolist()
        act_color_map = {a: palette[i % len(palette)] for i, a in enumerate(act_names)}

        figs["act_bar"] = px.bar(
            df_act, x="Actividad", y="Horas_Reales", color="Actividad",
            title="Distribución de actividades (24 h)",
            color_discrete_map=act_color_map,
            text="Horas_Reales",
        )
        figs["act_bar"].update_layout(showlegend=False)

    # conexiones
    if modo_reporte == "Perforación" and not dfc_local.empty and {"Componente", "Minutos_Reales"}.issubset(dfc_local.columns):
        df_conn_sum = dfc_local.groupby("Componente", as_index=False)["Minutos_Reales"].sum()
        df_conn_sum["Componente"] = pd.Categorical(df_conn_sum["Componente"], categories=CONN_ORDER, ordered=True)
        df_conn_sum = df_conn_sum.sort_values("Componente")

        figs["conn_pie"] = px.pie(
            df_conn_sum, names="Componente", values="Minutos_Reales", hole=0.35,
            title="Distribución de tiempo en conexión (min/% )",
            color="Componente", color_discrete_map=CONN_COLOR_MAP
        )

        df_stack = dfc_local.copy()
        df_stack["Conn_Label"] = df_stack["Profundidad_m"].fillna(df_stack["Conn_No"]).astype(float).astype(int).astype(str)
        df_stack["Componente"] = pd.Categorical(df_stack["Componente"], categories=CONN_ORDER, ordered=True)

        df_stack_g = df_stack.groupby(["Conn_Label", "Componente"], as_index=False)["Minutos_Reales"].sum().sort_values(["Conn_Label", "Componente"])

        per_conn = df_stack.groupby("Conn_Label", as_index=False).first()[["Conn_Label", "Conn_Tipo", "Angulo_Bucket"]]
        per_conn["Std_Total"] = per_conn.apply(
            lambda r: float(CONN_STDS.get((r["Conn_Tipo"], r["Angulo_Bucket"]), {}).get("TOTAL", 0.0)),
            axis=1,
        )
        std_line = float(per_conn["Std_Total"].mean()) if not per_conn.empty else 0.0

        fig_conn_stack = px.bar(
            df_stack_g,
            x="Conn_Label",
            y="Minutos_Reales",
            color="Componente",
            category_orders={"Componente": CONN_ORDER},
            color_discrete_map=CONN_COLOR_MAP,
            barmode="stack",
            title="Conexiones perforando",
            labels={"Conn_Label": "Profundidad (m)", "Minutos_Reales": "Tiempo (min)"},
        )

        if std_line > 0:
            fig_conn_stack.add_hline(
                y=std_line,
                line_dash="dash",
                line_color="#9C640C",
                annotation_text=f"{std_line:.1f}",
                annotation_position="top left",
                annotation_font_color="#9C640C",
            )

        df_tot = df_stack.groupby("Conn_Label", as_index=False)["Minutos_Reales"].sum().rename(columns={"Minutos_Reales": "Real_Total"})
        tot_map = dict(zip(df_tot["Conn_Label"].astype(str), df_tot["Real_Total"]))
        for x in sorted(df_tot["Conn_Label"].astype(str).unique(), key=lambda v: float(v) if v.replace(".", "", 1).isdigit() else v):
            y = float(tot_map.get(x, 0))
            fig_conn_stack.add_annotation(x=x, y=y, text=f"<b>{y:.0f}</b>", showarrow=False, yshift=10)

        fig_conn_stack.update_layout(legend_title_text="", xaxis_tickangle=0)
        figs["conn_stack"] = fig_conn_stack

    return figs

df_json = df.to_json(orient="split") if not df.empty else ""
df_conn_json = df_conn.to_json(orient="split") if not df_conn.empty else ""
figs = _make_figs(df_json, df_conn_json, modo_reporte) if show_charts else {"tiempos": None, "act_pie": None, "act_bar": None, "conn_pie": None, "conn_stack": None}

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# NAV PRO: TABS
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
st.session_state.df = _ensure_rowid(st.session_state.df)

def _any_state_key(prefix: str) -> bool:
    return any(k.startswith(prefix) and bool(st.session_state.get(k)) for k in st.session_state.keys())

any_viaje_time = _any_state_key("viaje_use_time_")
any_viaje_bitacora = _any_state_key("viaje_use_bitacora_")

show_bitacora_tab = bool(st.session_state.get("act_use_time", False)) or bool(st.session_state.get("bha_use_time", False)) or bool(st.session_state.get("conn_use_time", False)) or any_viaje_time
tab_labels = []
if show_bitacora_tab:
    tab_labels.append("Bitácora por horas")
tab_labels += [
    "Resumen", "Indicadores (Actividades)", "Top TNPI/TNP", "Conexiones", "Viajes y conexiones",
    "BHA (Arma/Desarma)", "ROP", "Metros", "Detalle", "Comparativa de Etapas",
    "Estadísticas CE", "Estadísticas por Etapa", "Estadísticas por Corrida", "Estadísticas DrillSpot",
    "Reporte General del Pozo", "Ejecutivo", "Exportar"
]

_tabs = st.tabs(tab_labels)
if show_bitacora_tab:
    tab_bitacora = _tabs[0]
    (tab_resumen, tab_act, tab_top, tab_conn, tab_viajes, tab_bha, tab_rop, tab_metros, tab_detalle, tab_comp,
     tab_ce, tab_estadisticas, tab_corridas, tab_drillspot, tab_general, tab_ejecutivo, tab_export) = _tabs[1:]
else:
    tab_bitacora = None
    (tab_resumen, tab_act, tab_top, tab_conn, tab_viajes, tab_bha, tab_rop, tab_metros, tab_detalle, tab_comp,
     tab_ce, tab_estadisticas, tab_corridas, tab_drillspot, tab_general, tab_ejecutivo, tab_export) = _tabs
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: BITÁCORA (ACTIVIDADES)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
if tab_bitacora is not None:
    with tab_bitacora:
        st.subheader("Bitácora por horas (todas las actividades)")
        bitacora_enabled_any = bool(st.session_state.get("act_use_bitacora", False)) or bool(st.session_state.get("bha_use_bitacora", False)) or bool(st.session_state.get("conn_use_bitacora", False)) or any_viaje_bitacora
        _act_label = st.session_state.get("actividad_select_sidebar", "")
        _act_label = _act_label or st.session_state.get("actividad_otro", "") or "-"
        st.caption(f"Actividad activa: **{_act_label}**")
        if not show_bitacora_tab:
            st.info("Activa **Registrar hora** en alguna actividad para habilitar esta bitácora.")
        elif not bitacora_enabled_any:
            st.info("Activa **Bitácora por horas** en la actividad que estás registrando.")
        else:
            if st.button("Limpiar bitácora", use_container_width=False):
                st.session_state["act_bitacora_rows"] = pd.DataFrame([{"Hora inicio": "", "Hora fin": "", "Comentario": ""}])
                st.session_state["act_bitacora_entries"] = []
                st.session_state["act_bitacora_total_h"] = 0.0
                st.success("Bitácora limpia.")
                st.rerun()
            st.caption("Cada tramo usa el tipo seleccionado en el sidebar.")
            _default_bit = st.session_state.get("act_bitacora_rows", None)
            if _default_bit is None or not isinstance(_default_bit, pd.DataFrame) or _default_bit.empty:
                _default_bit = pd.DataFrame([{"Hora inicio": "", "Hora fin": "", "Comentario": ""}])
            bit_df = st.data_editor(
                _default_bit,
                num_rows="dynamic",
                use_container_width=True,
                key="act_bitacora_rows",
                column_config={
                    "Hora inicio": st.column_config.TextColumn("Hora inicio", help="Formato HH:MM"),
                    "Hora fin": st.column_config.TextColumn("Hora fin", help="Formato HH:MM"),
                    "Comentario": st.column_config.TextColumn("Comentario"),
                },
            )
            invalid_rows = 0
            bitacora_entries = []
            if isinstance(bit_df, pd.DataFrame) and not bit_df.empty:
                for _, r in bit_df.iterrows():
                    _ini_txt = str(r.get("Hora inicio", "") or "").strip()
                    _fin_txt = str(r.get("Hora fin", "") or "").strip()
                    if not _ini_txt or not _fin_txt:
                        continue
                    _ini_dt = pd.to_datetime(_ini_txt, format="%H:%M", errors="coerce")
                    _fin_dt = pd.to_datetime(_fin_txt, format="%H:%M", errors="coerce")
                    if pd.isna(_ini_dt) or pd.isna(_fin_dt):
                        invalid_rows += 1
                        continue
                    _dt_ini = datetime.combine(datetime.today().date(), _ini_dt.time())
                    _dt_fin = datetime.combine(datetime.today().date(), _fin_dt.time())
                    _h = (_dt_fin - _dt_ini).total_seconds() / 3600.0
                    if _h < 0:
                        _h += 24.0
                    if _h <= 0:
                        continue
                    bitacora_entries.append({
                        "Hora_Inicio": _ini_dt.strftime("%H:%M"),
                        "Hora_Fin": _fin_dt.strftime("%H:%M"),
                        "Horas_Reales": float(_h),
                        "Comentario": str(r.get("Comentario", "") or "").strip(),
                    })
            bitacora_total_h = float(sum(x["Horas_Reales"] for x in bitacora_entries))
            st.session_state["act_bitacora_entries"] = bitacora_entries
            st.session_state["act_bitacora_total_h"] = float(bitacora_total_h)
            st.caption(f"Total bitácora: {bitacora_total_h:.2f} h")
            if invalid_rows > 0:
                st.warning(f"Hay {invalid_rows} filas con hora inválida (usa HH:MM).")

# TAB: RESUMEN
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: RESUMEN (MODIFICADO CON FILTRO DE ETAPA)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
with tab_resumen:

    # --- MISSION CONTROL DASHBOARD ---
    st.markdown("### 🧭 Centro de Control de Misión")

    # Calcular KPIs generales (todas las etapas)
    total_prog = float(df["Horas_Prog"].sum()) if not df.empty else 0.0
    total_real = float(df["Horas_Reales"].sum()) if not df.empty else 0.0
    tp_h = float(df[df["Tipo"] == "TP"]["Horas_Reales"].sum()) if not df.empty else 0.0
    tnpi_h = float(df[df["Tipo"] == "TNPI"]["Horas_Reales"].sum()) if not df.empty else 0.0
    tnp_h = float(df[df["Tipo"] == "TNP"]["Horas_Reales"].sum()) if not df.empty else 0.0
    eficiencia_dia = clamp_0_100(safe_pct(tp_h, total_real)) if total_real > 0 else 0.0

    # Mostrar el dashboard NASA (vista general)
    dashboard_html = mission_control_dashboard(
        etapa="VISIÓN GENERAL",
        eficiencia=eficiencia_dia,
        tp_h=tp_h,
        tnpi_h=tnpi_h,
        tnp_h=tnp_h,
        total_real=total_real
    )

    render_html(dashboard_html, height=450)

    # --- FILTRO DE ETAPA EN EL RESUMEN ---
    col_filtro1, col_filtro2 = st.columns([1, 3])

    with col_filtro1:
        # Obtener todas las etapas disponibles
        etapas_set = set()
        if not df.empty and "Etapa" in df.columns:
            etapas_set.update([str(x).strip() for x in df["Etapa"].unique().tolist()])
        etapas_set.update([str(x).strip() for x in (st.session_state.drill_day.get("por_etapa", {}) or {}).keys()])
        if isinstance(etapa, str) and etapa.strip():
            etapas_set.add(etapa.strip())
        etapas_disponibles = sorted([e for e in etapas_set if e and e.lower() != "nan"])
        if not etapas_disponibles:
            etapas_disponibles = ["Sin datos"]

        # Selector de etapa para el resumen
        etapa_resumen = st.selectbox(
            "Etapa para resumen",
            options=etapas_disponibles,
            index=etapas_disponibles.index(etapa) if etapa in etapas_disponibles else 0,
            key="etapa_resumen"
        )

    # Filtrar datos por la etapa seleccionada
    if etapa_resumen != "Sin datos":
        df_resumen_filtrado = df[df["Etapa"] == etapa_resumen].copy()
        df_conn_filtrado = df_conn[df_conn["Etapa"] == etapa_resumen].copy()
    else:
        df_resumen_filtrado = pd.DataFrame()
        df_conn_filtrado = pd.DataFrame()

    # Recalcular KPIs con datos filtrados
    total_prog_filtrado = float(df_resumen_filtrado["Horas_Prog"].sum()) if not df_resumen_filtrado.empty else 0.0
    total_real_filtrado = float(df_resumen_filtrado["Horas_Reales"].sum()) if not df_resumen_filtrado.empty else 0.0
    tp_h_filtrado = float(df_resumen_filtrado[df_resumen_filtrado["Tipo"] == "TP"]["Horas_Reales"].sum()) if not df_resumen_filtrado.empty else 0.0
    tnpi_h_filtrado = float(df_resumen_filtrado[df_resumen_filtrado["Tipo"] == "TNPI"]["Horas_Reales"].sum()) if not df_resumen_filtrado.empty else 0.0
    tnp_h_filtrado = float(df_resumen_filtrado[df_resumen_filtrado["Tipo"] == "TNP"]["Horas_Reales"].sum()) if not df_resumen_filtrado.empty else 0.0
    eficiencia_dia_filtrado = clamp_0_100(safe_pct(tp_h_filtrado, total_real_filtrado)) if total_real_filtrado > 0 else 0.0

    # Recalcular KPIs de conexiones filtradas
    conn_real_min_filtrado = 0.0
    conn_std_min_filtrado = 0.0
    conn_tnpi_min_filtrado = 0.0
    conn_tnp_min_filtrado = 0.0
    eff_conn_filtrado = 0.0

    if not df_conn_filtrado.empty:
        conn_real_min_filtrado = float(df_conn_filtrado.groupby(["Conn_No"])["Minutos_Reales"].sum().sum())
        per_conn2_filtrado = df_conn_filtrado.groupby("Conn_No", as_index=False).first()[["Conn_No", "Conn_Tipo", "Angulo_Bucket"]]
        per_conn2_filtrado["Std_Total"] = per_conn2_filtrado.apply(
            lambda r: float(CONN_STDS.get((r["Conn_Tipo"], r["Angulo_Bucket"]), {}).get("TOTAL", 0.0)),
            axis=1
        )
        conn_std_min_filtrado = float(per_conn2_filtrado["Std_Total"].sum())

        conn_tp_min_filtrado = min(conn_real_min_filtrado, conn_std_min_filtrado) if conn_std_min_filtrado > 0 else conn_real_min_filtrado
        conn_tnpi_min_filtrado = max(0.0, conn_real_min_filtrado - conn_std_min_filtrado) if conn_std_min_filtrado > 0 else 0.0
        eff_conn_filtrado = clamp_0_100(safe_pct(conn_tp_min_filtrado, conn_real_min_filtrado)) if conn_real_min_filtrado > 0 else 0.0

    # --- MISSION CONTROL PARA ETAPA ESPECÍFICA ---
    if etapa_resumen != "Sin datos" and not df_resumen_filtrado.empty:
        st.markdown(f"### 📊 Control de Misión - {etapa_resumen}")

        etapa_dashboard_html = mission_control_dashboard(
            etapa=etapa_resumen,
            eficiencia=eficiencia_dia_filtrado,
            tp_h=tp_h_filtrado,
            tnpi_h=tnpi_h_filtrado,
            tnp_h=tnp_h_filtrado,
            total_real=total_real_filtrado
        )

        render_html(etapa_dashboard_html, height=450)

    st.subheader("Indicador de desempeño (principal)")

    # Gauge con eficiencia filtrada
    fig_g = build_gauge(f"Eficiencia - {etapa_resumen}", eficiencia_dia_filtrado) if PLOTLY_IMG_OK else None
    col_g, col_t = st.columns([1.05, 2.0])

    with col_g:
        if fig_g is not None:
            # Usar etapa_resumen para hacer la clave única
            etapa_key = etapa_resumen.replace(" ", "_").replace("/", "_").replace('"', "") if etapa_resumen != "Sin datos" else "general"
            st.plotly_chart(fig_g, use_container_width=True, key=f"gauge_resumen_{etapa_key}")
        else:
            st.info("Para gauge instala kaleido: pip install -U kaleido")

    with col_t:
        # KPIs específicos de la etapa seleccionada
        kpi_rows = [
            {"kpi": "Horas Totales", "real": f"{total_real_filtrado:.1f} h", "tnpi": f"{tnpi_h_filtrado:.1f} h", "eff": eficiencia_dia_filtrado},
            {"kpi": "Conexión perforando", "real": f"{(conn_real_min_filtrado/60.0):.2f} h", "tnpi": f"{(conn_tnpi_min_filtrado/60.0):.2f} h", "eff": eff_conn_filtrado},
        ]

        # Agregar metros y ROP si tenemos datos por etapa
        if modo_reporte == "Perforación" and etapa_resumen != "Sin datos":
            # Obtener datos específicos de esta etapa
            etapa_data_resumen = get_etapa_data(etapa_resumen)

            mr_total = float(etapa_data_resumen.get("metros_real_dia", 0.0)) + float(etapa_data_resumen.get("metros_real_noche", 0.0))
            tnpi_m_h = float(etapa_data_resumen.get("tnpi_metros_h", 0.0))
            mp_total = float(etapa_data_resumen.get("metros_prog_total", 0.0))

            eff_m = clamp_0_100(safe_pct(mr_total, mp_total)) if mp_total > 0 else 0.0

            kpi_rows.insert(0, {"kpi": "Metros perforados", "real": f"{mr_total:.0f} m", "tnpi": f"{tnpi_m_h:.2f} h", "eff": eff_m})

        components.html(kpi_table_html(kpi_rows), height=300, scrolling=False)

    
    # Mostrar indicador claro de qué etapa estamos viendo
    with col_filtro2:
        light_mode = _is_light_theme()
        if light_mode:
            etapa_bg = "rgba(37, 99, 235, 0.08)"
            etapa_border = "4px solid #1d4ed8"
            etapa_label = "#1d4ed8"
            etapa_value = "#0f172a"
            etapa_sub = "#475569"
        else:
            etapa_bg = "rgba(46, 134, 193, 0.1)"
            etapa_border = "4px solid #2E86C1"
            etapa_label = "#2E86C1"
            etapa_value = "white"
            etapa_sub = "rgba(255,255,255,0.7)"
        st.markdown(f"""
            <div style='background: {etapa_bg}; padding: 10px; border-radius: 10px; border-left: {etapa_border}; margin-top: 10px;'>
                <div style='font-size: 14px; color: {etapa_label}; font-weight: bold;'>Etapa seleccionada:</div>
                <div style='font-size: 18px; color: {etapa_value}; font-weight: bold;'>{etapa_resumen}</div>
                <div style='font-size: 12px; color: {etapa_sub}; margin-top: 5px;'>
                    {len(df_resumen_filtrado)} actividades | {len(df_conn_filtrado)} conexiones
                </div>
            </div>
        """, unsafe_allow_html=True)
    
    # ------------------------------
    # Avance de profundidad (solo Perforación)
    # ------------------------------
    if modo_reporte == "Perforación" and etapa_resumen != "Sin datos":
        # Obtener datos específicos de esta etapa
        etapa_data_resumen = get_etapa_data(etapa_resumen)
        
        pt_prog = float(etapa_data_resumen.get("pt_programada_m", 0.0) or 0.0)
        prof_act = float(etapa_data_resumen.get("prof_actual_m", 0.0) or 0.0)
        
        avance = (prof_act / pt_prog) if pt_prog > 0 else 0.0
        avance = max(0.0, min(1.0, avance))

        restante = max(0.0, pt_prog - prof_act)
        tone_av = "green" if avance >= 0.85 else ("amber" if avance >= 0.70 else "red")
        tone_rest = "green" if restante <= 0 else "gray"
        riesgo_txt = "BAJO" if avance >= 0.85 else ("MEDIO" if avance >= 0.70 else "ALTO")

        st.markdown("### Avance de profundidad")

        # Chips pro + barra futurista
        render_chip_row([
            {"label": "Etapa", "value": etapa_resumen, "tone": "blue"},
            {"label": "PT programada", "value": f"{pt_prog:,.0f} m", "tone": "violet"},
            {"label": "Prof. actual", "value": f"{prof_act:,.0f} m", "tone": "blue"},
            {"label": "Restante", "value": f"{restante:,.0f} m", "tone": tone_rest},
            {"label": "Avance", "value": f"{avance*100:.1f}%", "tone": tone_av},
            {"label": "Riesgo", "value": riesgo_txt, "tone": tone_av},
        ], use_iframe=True, height=120)

        if avance >= 0.85:
            bar_grad = "linear-gradient(90deg, #00ff88, #00c3ff 60%, #8b5cf6)"
            glow_color = "rgba(0, 255, 136, 0.35)"
            chip_bg = "rgba(10, 35, 24, 0.9)"
            chip_border = "rgba(0, 255, 136, 0.65)"
            chip_fg = "#c7ffe6"
            if _is_light_theme():
                chip_bg = "rgba(0, 255, 136, 0.12)"
                chip_border = "rgba(0, 255, 136, 0.55)"
                chip_fg = "#065f46"
        elif avance >= 0.70:
            bar_grad = "linear-gradient(90deg, #f59e0b, #fbbf24 60%, #fb7185)"
            glow_color = "rgba(245, 158, 11, 0.35)"
            chip_bg = "rgba(35, 25, 8, 0.9)"
            chip_border = "rgba(245, 158, 11, 0.75)"
            chip_fg = "#ffe6b0"
            if _is_light_theme():
                chip_bg = "rgba(245, 158, 11, 0.12)"
                chip_border = "rgba(245, 158, 11, 0.55)"
                chip_fg = "#92400e"
        else:
            bar_grad = "linear-gradient(90deg, #ef4444, #f97316 60%, #f59e0b)"
            glow_color = "rgba(239, 68, 68, 0.35)"
            chip_bg = "rgba(40, 10, 10, 0.9)"
            chip_border = "rgba(239, 68, 68, 0.75)"
            chip_fg = "#ffd3d3"
            if _is_light_theme():
                chip_bg = "rgba(239, 68, 68, 0.12)"
                chip_border = "rgba(239, 68, 68, 0.55)"
                chip_fg = "#991b1b"

        light_mode = _is_light_theme()
        if light_mode:
            prog_wrap_bg = "linear-gradient(180deg, rgba(255,255,255,0.98), rgba(244,247,251,0.98))"
            prog_border = "rgba(15,23,42,0.10)"
            prog_shadow = "0 10px 24px rgba(15,23,42,0.10)"
            prog_head = "#475569"
            prog_bar_bg = "rgba(15,23,42,0.08)"
            prog_ticks = "rgba(15,23,42,0.45)"
            prog_glow = "radial-gradient(120px 20px at 20% 50%, rgba(0,255,136,0.18), transparent 60%)"
            chip_shadow = "0 8px 18px rgba(15,23,42,0.12), 0 0 12px {glow_color}"
            chip_arrow_shadow = "drop-shadow(0 2px 4px rgba(15,23,42,0.15))"
        else:
            prog_wrap_bg = "linear-gradient(180deg, rgba(18,18,22,0.92), rgba(10,10,14,0.95))"
            prog_border = "rgba(255,255,255,0.08)"
            prog_shadow = "0 10px 28px rgba(0,0,0,0.45)"
            prog_head = "rgba(255,255,255,0.8)"
            prog_bar_bg = "rgba(255,255,255,0.08)"
            prog_ticks = "rgba(255,255,255,0.55)"
            prog_glow = "radial-gradient(120px 20px at 20% 50%, rgba(0,255,136,0.25), transparent 60%)"
            chip_shadow = "0 8px 18px rgba(0,0,0,0.45), 0 0 12px {glow_color}"
            chip_arrow_shadow = "drop-shadow(0 2px 4px rgba(0,0,0,0.35))"

        progress_html = f"""
        <style>
          .ds-progress-wrap {{
            border-radius: 16px;
            padding: 14px 16px;
            background: {prog_wrap_bg};
            border: 1px solid {prog_border};
            box-shadow: {prog_shadow};
          }}
          .ds-progress-head {{
            display:flex; align-items:center; justify-content:space-between;
            color: {prog_head}; font-size:12px; font-weight:700;
            letter-spacing:0.4px; text-transform:uppercase;
          }}
          .ds-progress-bar {{
            position: relative; height: 14px; border-radius: 999px;
            background: {prog_bar_bg};
            overflow: hidden; margin-top: 10px;
          }}
          .ds-progress-fill {{
            height: 100%;
            width: {avance*100:.2f}%;
            border-radius: 999px;
            background: {bar_grad};
            box-shadow: 0 0 18px {glow_color};
          }}
          .ds-progress-sheen {{
            position:absolute; inset:0;
            background: linear-gradient(120deg, rgba(255,255,255,0.0), rgba(255,255,255,0.25), rgba(255,255,255,0.0));
            transform: translateX(-120%);
            animation: sheen 3.2s ease-in-out infinite;
            mix-blend-mode: screen;
            pointer-events:none;
          }}
          .ds-progress-glow {{
            position:absolute; inset:0;
            background: {prog_glow};
            mix-blend-mode: screen;
          }}
          .ds-progress-ticks {{
            display:flex; justify-content:space-between; margin-top:8px;
            font-size:11px; color: {prog_ticks};
          }}
          .ds-progress-chip {{
            position:absolute; top:-28px;
            left: calc({avance*100:.2f}%);
            transform: translateX(-50%);
            padding: 4px 8px;
            border-radius: 999px;
            background: {chip_bg};
            border: 1px solid {chip_border};
            color: {chip_fg};
            font-size: 11px;
            font-weight: 800;
            letter-spacing: 0.2px;
            box-shadow: {chip_shadow};
            backdrop-filter: blur(6px);
            white-space: nowrap;
            animation: chipPulse 2.4s ease-in-out infinite;
          }}
          .ds-progress-chip::after {{
            content:"";
            position:absolute; left: 50%; bottom: -5px;
            transform: translateX(-50%);
            width: 0; height: 0;
            border-left: 5px solid transparent;
            border-right: 5px solid transparent;
            border-top: 6px solid {chip_border};
            filter: {chip_arrow_shadow};
          }}
          @keyframes chipPulse {{
            0% {{ transform: translateX(-50%) scale(1); opacity: 1; }}
            50% {{ transform: translateX(-50%) scale(1.03); opacity: 0.88; }}
            100% {{ transform: translateX(-50%) scale(1); opacity: 1; }}
          }}
          @keyframes sheen {{
            0% {{ transform: translateX(-120%); }}
            50% {{ transform: translateX(20%); }}
            100% {{ transform: translateX(120%); }}
          }}
        </style>
        <div class="ds-progress-wrap">
          <div class="ds-progress-head">
            <span>PROFUNDIDAD</span>
            <span>{prof_act:,.0f} / {pt_prog:,.0f} m</span>
          </div>
          <div class="ds-progress-bar">
            <div class="ds-progress-fill"></div>
            <div class="ds-progress-glow"></div>
            <div class="ds-progress-sheen"></div>
            <div class="ds-progress-chip">Avance {avance*100:.1f}%</div>
          </div>
          <div class="ds-progress-ticks">
            <span>0%</span><span>25%</span><span>50%</span><span>75%</span><span>100%</span>
          </div>
        </div>
        """
        st.markdown(progress_html, unsafe_allow_html=True)

        c1, c2, c3 = st.columns(3)
        c1.metric("PT programada (m)", f"{pt_prog:,.0f}")
        c2.metric("Profundidad actual (m)", f"{prof_act:,.0f}")
        c3.metric("Avance", f"{avance*100:.1f}%")
        
        st.divider()

    if show_charts and etapa_resumen != "Sin datos":
        st.divider()
        st.subheader(f"Gráficas - {etapa_resumen}")
        
        # Generar figuras específicas para esta etapa
        if not df_resumen_filtrado.empty:
            # Tiempos (TP vs TNPI vs TNP)
            df_tiempos = df_resumen_filtrado.groupby("Tipo")["Horas_Reales"].sum().reset_index()
            if not df_tiempos.empty:
                fig_tiempos = px.pie(df_tiempos, names="Tipo", values="Horas_Reales", 
                                     hole=0.55, title=f"TP vs TNPI vs TNP - {etapa_resumen}")
                st.plotly_chart(fig_tiempos, use_container_width=True, key="pie_tiempos_resumen")
            
            # Actividades
            df_act = df_resumen_filtrado.groupby("Actividad", as_index=False)["Horas_Reales"].sum().sort_values("Horas_Reales", ascending=False)
            if not df_act.empty:
                fig_act_pie = px.pie(df_act.head(8), names="Actividad", values="Horas_Reales", 
                                     hole=0.35, title=f"Top Actividades - {etapa_resumen}")
                st.plotly_chart(fig_act_pie, use_container_width=True, key="pie_actividades_resumen")

    # -----------------------------------------------------------------
    # RESUMEN DIARIO (mismas gráficas pero por Fecha)
    # -----------------------------------------------------------------
    with st.expander("Resumen diario (por fecha)", expanded=False):
        if not df_resumen_filtrado.empty and "Fecha" in df_resumen_filtrado.columns:
            fechas_disp = (
                sorted(df_resumen_filtrado["Fecha"].astype(str).dropna().unique().tolist())
                if not df_resumen_filtrado["Fecha"].isna().all()
                else []
            )
        else:
            fechas_disp = []
        if len(fechas_disp) == 0:
            st.info("No hay fechas disponibles para generar el resumen diario.")
        else:
            fecha_resumen = st.selectbox(
                "Fecha para resumen diario",
                options=fechas_disp,
                index=len(fechas_disp) - 1,
                key="fecha_resumen_diario"
            )

            df_diario = df_resumen_filtrado[df_resumen_filtrado["Fecha"].astype(str) == str(fecha_resumen)].copy()

            if df_diario.empty:
                st.info(f"No hay datos para la fecha {fecha_resumen} (etapa {etapa_resumen}).")
            else:
                # KPIs diarios
                total_real_d = float(df_diario["Horas_Reales"].sum()) if "Horas_Reales" in df_diario.columns else 0.0
                tp_h_d = float(df_diario[df_diario["Tipo"] == "TP"]["Horas_Reales"].sum()) if "Tipo" in df_diario.columns else 0.0
                tnpi_h_d = float(df_diario[df_diario["Tipo"] == "TNPI"]["Horas_Reales"].sum()) if "Tipo" in df_diario.columns else 0.0
                tnp_h_d = float(df_diario[df_diario["Tipo"] == "TNP"]["Horas_Reales"].sum()) if "Tipo" in df_diario.columns else 0.0
                eff_d = clamp_0_100(safe_pct(tp_h_d, total_real_d)) if total_real_d > 0 else 0.0

                # Mission Control diario (cap DAY_LIMIT_HOURS)
                total_cap = min(total_real_d, DAY_LIMIT_HOURS)
                if total_real_d > DAY_LIMIT_HOURS:
                    scale = total_cap / total_real_d if total_real_d > 0 else 0.0
                    tp_cap = tp_h_d * scale
                    tnpi_cap = tnpi_h_d * scale
                    tnp_cap = tnp_h_d * scale
                else:
                    tp_cap, tnpi_cap, tnp_cap = tp_h_d, tnpi_h_d, tnp_h_d
                eff_cap = clamp_0_100(safe_pct(tp_cap, total_cap)) if total_cap > 0 else 0.0

                st.markdown(f"### 🧭 Mission Control Diario (cap {DAY_LIMIT_HOURS:.0f}h)")
                render_html(
                    mission_control_dashboard(
                        etapa=f"{etapa_resumen} / {fecha_resumen}",
                        eficiencia=eff_cap,
                        tp_h=tp_cap,
                        tnpi_h=tnpi_cap,
                        tnp_h=tnp_cap,
                        total_real=total_cap,
                    ),
                    height=450,
                )

                # Perforación + conexiones (diario)
                st.markdown("### Perforación + conexiones (diario)")
                include_conn_d = st.toggle("Incluir conexiones perforando", value=True, key="resumen_diario_include_conn")
                perf_h_d = 0.0
                if not df_diario.empty and "Actividad" in df_diario.columns:
                    _act_d = df_diario["Actividad"].astype(str).str.strip().str.lower()
                    perf_h_d = float(df_diario[_act_d.isin(["perforación", "perforacion"])]["Horas_Reales"].sum())
                conn_h_d = 0.0
                conn_tp_h_d = 0.0
                conn_tnpi_h_d = 0.0
                conn_tnp_h_d = 0.0
                if not df_conn_filtrado.empty and "Fecha" in df_conn_filtrado.columns:
                    _df_conn_d = df_conn_filtrado[df_conn_filtrado["Fecha"].astype(str) == str(fecha_resumen)].copy()
                    if not _df_conn_d.empty:
                        conn_real_min = float(_df_conn_d.groupby("Conn_No")["Minutos_Reales"].sum().sum())
                        conn_tnpi_min = float(_df_conn_d["Minutos_TNPI"].sum()) if "Minutos_TNPI" in _df_conn_d.columns else 0.0
                        conn_tnp_min = float(_df_conn_d["Minutos_TNP"].sum()) if "Minutos_TNP" in _df_conn_d.columns else 0.0
                        conn_tp_min = max(0.0, conn_real_min - conn_tnpi_min - conn_tnp_min)
                        conn_h_d = conn_real_min / 60.0
                        conn_tp_h_d = conn_tp_min / 60.0
                        conn_tnpi_h_d = conn_tnpi_min / 60.0
                        conn_tnp_h_d = conn_tnp_min / 60.0

                col_m1, col_m2, col_m3 = st.columns(3)
                col_m1.metric("Horas perforación", f"{perf_h_d:.2f} h")
                col_m2.metric("Horas conexión", f"{conn_h_d:.2f} h")
                col_m3.metric("Total", f"{(perf_h_d + (conn_h_d if include_conn_d else 0.0)):.2f} h")

                render_chip_row(
                    [
                        {"label": "Perforación", "value": f"{perf_h_d:.2f} h", "tone": "blue"},
                        {"label": "Conexión", "value": f"{conn_h_d:.2f} h", "tone": "amber"},
                        {"label": "Total", "value": f"{(perf_h_d + (conn_h_d if include_conn_d else 0.0)):.2f} h", "tone": "green"},
                        {"label": "TNPI", "value": f"{tnpi_h_d:.2f} h", "tone": "orange"},
                        {"label": "TNP", "value": f"{tnp_h_d:.2f} h", "tone": "red"},
                    ],
                    use_iframe=True,
                    height=90,
                )

                rows_d = [{"Tipo": "Perforación", "Segmento": "Perforación", "Horas": perf_h_d}]
                if include_conn_d:
                    rows_d.extend([
                        {"Tipo": "Conexión perforando", "Segmento": "TP", "Horas": conn_tp_h_d},
                        {"Tipo": "Conexión perforando", "Segmento": "TNPI", "Horas": conn_tnpi_h_d},
                        {"Tipo": "Conexión perforando", "Segmento": "TNP", "Horas": conn_tnp_h_d},
                    ])
                df_plot_d = pd.DataFrame(rows_d)
                if not df_plot_d.empty:
                    fig_pc_d = px.bar(
                        df_plot_d,
                        x="Tipo",
                        y="Horas",
                        title="Horas de perforación (con opción de conexiones)",
                        color="Segmento",
                        barmode="stack",
                        color_discrete_map={
                            "Perforación": "#2563EB",
                            "TP": "#10B981",
                            "TNPI": "#F59E0B",
                            "TNP": "#EF4444",
                        },
                    )
                    fig_pc_d.update_traces(text=None)
                    totals_d = df_plot_d.groupby("Tipo", as_index=False)["Horas"].sum()
                    fig_pc_d.add_trace(
                        go.Scatter(
                            x=totals_d["Tipo"],
                            y=totals_d["Horas"],
                            text=totals_d["Horas"].map(lambda v: f"{v:.2f} h"),
                            mode="text",
                            textposition="top center",
                            showlegend=False,
                        )
                    )
                    fig_pc_d.update_layout(yaxis_title="Horas", xaxis_title="")
                    st.plotly_chart(fig_pc_d, use_container_width=True)

                c1, c2, c3, c4, c5 = st.columns(5)
                c1.metric("Horas (Real)", f"{total_real_d:.2f}")
                c2.metric("TP (h)", f"{tp_h_d:.2f}")
                c3.metric("TNPI (h)", f"{tnpi_h_d:.2f}")
                c4.metric("TNP (h)", f"{tnp_h_d:.2f}")
                c5.metric("Eficiencia", f"{eff_d:.1f}%")

                # Gráfica KPI diaria
                df_tiempos_d = df_diario.groupby("Tipo")["Horas_Reales"].sum().reset_index()
                if not df_tiempos_d.empty:
                    fig_tiempos_d = px.pie(
                        df_tiempos_d, names="Tipo", values="Horas_Reales",
                        hole=0.55, title=f"TP vs TNPI vs TNP (Diario) - {fecha_resumen} / {etapa_resumen}"
                    )
                    st.plotly_chart(fig_tiempos_d, use_container_width=True)

                # Actividades diarias
                df_act_d = (
                    df_diario.groupby("Actividad", as_index=False)["Horas_Reales"]
                    .sum()
                    .sort_values("Horas_Reales", ascending=False)
                )
                if not df_act_d.empty:
                    fig_act_pie_d = px.pie(
                        df_act_d.head(10), names="Actividad", values="Horas_Reales",
                        hole=0.35, title=f"Top Actividades (Diario) - {fecha_resumen} / {etapa_resumen}"
                    )
                    st.plotly_chart(fig_act_pie_d, use_container_width=True)

                # ------------------------------
                # KPI diario pro (semáforos)
                # ------------------------------
                st.markdown("### KPIs diarios (pro)")
                sem_tp = semaforo_dot(safe_pct(tp_h_d, total_real_d)) if total_real_d > 0 else "⚪"
                sem_eff = semaforo_dot(eff_d)
                kpi_rows_d = [
                    {"Métrica": "Horas reales", "Valor": f"{total_real_d:.2f} h", "Semáforo": "⚪"},
                    {"Métrica": "TP", "Valor": f"{tp_h_d:.2f} h", "Semáforo": sem_tp},
                    {"Métrica": "TNPI", "Valor": f"{tnpi_h_d:.2f} h", "Semáforo": semaforo_dot(safe_pct(tp_h_d, total_real_d)) if total_real_d > 0 else "⚪"},
                    {"Métrica": "TNP", "Valor": f"{tnp_h_d:.2f} h", "Semáforo": "⚪"},
                    {"Métrica": "Eficiencia", "Valor": f"{eff_d:.1f}%", "Semáforo": sem_eff},
                ]
                st.dataframe(pd.DataFrame(kpi_rows_d), use_container_width=True, hide_index=True)

                # ------------------------------
                # Top 5 causas TNPI/TNP (diario)
                # ------------------------------
                st.markdown("### Top 5 causas TNPI/TNP del día")
                col_t1, col_t2 = st.columns(2)

                with col_t1:
                    df_tnpi_d = df_diario[df_diario.get("Tipo", "") == "TNPI"].copy()
                    if df_tnpi_d.empty:
                        st.info("No hay TNPI para este día.")
                    else:
                        df_tnpi_d["Detalle_TNPI"] = df_tnpi_d.get("Detalle_TNPI", "-").replace({"-": "Sin detalle"}).astype(str)
                        g_tnpi = (
                            df_tnpi_d.groupby("Detalle_TNPI", as_index=False)["Horas_Reales"]
                            .sum()
                            .sort_values("Horas_Reales", ascending=False)
                            .head(5)
                        )
                        total_tnpi = float(df_tnpi_d["Horas_Reales"].sum())
                        g_tnpi["%"] = g_tnpi["Horas_Reales"].apply(lambda v: (float(v) / total_tnpi * 100.0) if total_tnpi > 0 else 0.0)
                        fig_tnpi = px.bar(
                            g_tnpi.sort_values("Horas_Reales"),
                            x="Horas_Reales",
                            y="Detalle_TNPI",
                            orientation="h",
                            title="TNPI - Top 5 causas",
                            text=g_tnpi["%"].map(lambda v: f"{v:.0f}%"),
                        )
                        fig_tnpi.update_layout(xaxis_title="Horas", yaxis_title="Detalle")
                        fig_tnpi.update_traces(marker_color="#EF4444", textposition="outside")
                        st.plotly_chart(fig_tnpi, use_container_width=True)

                with col_t2:
                    df_tnp_d = df_diario[df_diario.get("Tipo", "") == "TNP"].copy()
                    if df_tnp_d.empty:
                        st.info("No hay TNP para este día.")
                    else:
                        df_tnp_d["Detalle_TNP"] = df_tnp_d.get("Detalle_TNP", "-").replace({"-": "Sin detalle"}).astype(str)
                        g_tnp = (
                            df_tnp_d.groupby("Detalle_TNP", as_index=False)["Horas_Reales"]
                            .sum()
                            .sort_values("Horas_Reales", ascending=False)
                            .head(5)
                        )
                        total_tnp = float(df_tnp_d["Horas_Reales"].sum())
                        g_tnp["%"] = g_tnp["Horas_Reales"].apply(lambda v: (float(v) / total_tnp * 100.0) if total_tnp > 0 else 0.0)
                        fig_tnp = px.bar(
                            g_tnp.sort_values("Horas_Reales"),
                            x="Horas_Reales",
                            y="Detalle_TNP",
                            orientation="h",
                            title="TNP - Top 5 causas",
                            text=g_tnp["%"].map(lambda v: f"{v:.0f}%"),
                        )
                        fig_tnp.update_layout(xaxis_title="Horas", yaxis_title="Detalle")
                        fig_tnp.update_traces(marker_color="#3B82F6", textposition="outside")
                        st.plotly_chart(fig_tnp, use_container_width=True)

                # ------------------------------
                # ROP diario (Día vs Noche)
                # ------------------------------
                st.markdown("### ROP real vs programado (Día/Noche)")
                etapa_data_rop_d = get_etapa_data(etapa_resumen) if etapa_resumen != "Sin datos" else {}
                _prog_map = etapa_data_rop_d.get("rop_prog_by_date", {}) or {}
                _rd_map = etapa_data_rop_d.get("rop_real_dia_by_date", {}) or {}
                _rn_map = etapa_data_rop_d.get("rop_real_noche_by_date", {}) or {}
                _p_entry = _prog_map.get(str(fecha_resumen), {})
                rop_prog_d = _safe_float(_p_entry.get("rop_prog") if isinstance(_p_entry, dict) else (_p_entry or 0.0))
                rop_rd = _safe_float(_rd_map.get(str(fecha_resumen), 0.0) or 0.0)
                rop_rn = _safe_float(_rn_map.get(str(fecha_resumen), 0.0) or 0.0)
                df_rop_d = pd.DataFrame(
                    [
                        {"Turno": "Día ☀️", "Programado (m/h)": rop_prog_d, "Real (m/h)": rop_rd},
                        {"Turno": "Noche 🌙", "Programado (m/h)": rop_prog_d, "Real (m/h)": rop_rn},
                    ]
                )
                if (rop_prog_d + rop_rd + rop_rn) > 0:
                    fig_rop_d = px.bar(
                        df_rop_d,
                        x="Turno",
                        y=["Programado (m/h)", "Real (m/h)"],
                        barmode="group",
                        text_auto=True,
                    )
                    fig_rop_d.update_layout(margin=dict(l=10, r=10, t=30, b=10), height=320, legend_title_text="Serie")
                    st.plotly_chart(fig_rop_d, use_container_width=True)
                else:
                    st.info("No hay datos de ROP para este día.")

                # ------------------------------
                # Metros perforados diario (Día/Noche)
                # ------------------------------
                st.markdown("### Metros perforados (Real vs Programado)")
                _mp_map = etapa_data_rop_d.get("metros_prog_by_date", {}) or {}
                _md_map = etapa_data_rop_d.get("metros_real_dia_by_date", {}) or {}
                _mn_map = etapa_data_rop_d.get("metros_real_noche_by_date", {}) or {}
                _mp_entry = _mp_map.get(str(fecha_resumen), {})
                mp_d = _safe_float(_mp_entry.get("metros_prog") if isinstance(_mp_entry, dict) else (_mp_entry or 0.0))
                mr_d = _safe_float(_md_map.get(str(fecha_resumen), 0.0) or 0.0)
                mr_n = _safe_float(_mn_map.get(str(fecha_resumen), 0.0) or 0.0)
                mr_t = mr_d + mr_n
                df_m_d = pd.DataFrame(
                    [
                        {"Tipo": "Programado (total)", "Metros (m)": mp_d},
                        {"Tipo": "Real Día ☀️", "Metros (m)": mr_d},
                        {"Tipo": "Real Noche 🌙", "Metros (m)": mr_n},
                        {"Tipo": "Real Total", "Metros (m)": mr_t},
                    ]
                )
                if (mp_d + mr_d + mr_n) > 0:
                    fig_m_d = px.bar(
                        df_m_d,
                        x="Tipo",
                        y="Metros (m)",
                        text_auto=True,
                        color="Tipo",
                        color_discrete_map={
                            "Programado (total)": "#6B7280",
                            "Real Día ☀️": "#F59E0B",
                            "Real Noche 🌙": "#1D4ED8",
                            "Real Total": "#22C55E",
                        },
                    )
                    fig_m_d.update_layout(margin=dict(l=10, r=10, t=30, b=10), height=320)
                    st.plotly_chart(fig_m_d, use_container_width=True)
                else:
                    st.info("No hay datos de metros para este día.")

                # ------------------------------
                # BHA diario (Real vs Estándar)
                # ------------------------------
                st.markdown("### BHA (Real vs Estándar)")
                df_bha_d = st.session_state.get("df_bha", pd.DataFrame()).copy()
                if not df_bha_d.empty and "Fecha" in df_bha_d.columns:
                    df_bha_d["Fecha"] = df_bha_d["Fecha"].astype(str)
                    df_bha_d = df_bha_d[df_bha_d["Fecha"] == str(fecha_resumen)].copy()
                    if "Etapa" in df_bha_d.columns and etapa_resumen != "Sin datos":
                        df_bha_d = df_bha_d[df_bha_d["Etapa"] == etapa_resumen].copy()
                if not df_bha_d.empty:
                    df_long_bha = df_bha_d.melt(
                        id_vars=[c for c in ["Accion", "BHA_Tipo"] if c in df_bha_d.columns],
                        value_vars=[c for c in ["Estandar_h", "Real_h"] if c in df_bha_d.columns],
                        var_name="Serie",
                        value_name="Horas",
                    )
                    fig_bha_d = px.bar(
                        df_long_bha,
                        x="BHA_Tipo" if "BHA_Tipo" in df_long_bha.columns else "Accion",
                        y="Horas",
                        color="Serie",
                        barmode="group",
                        title=f"BHA - {fecha_resumen} / {etapa_resumen}",
                    )
                    fig_bha_d.update_layout(margin=dict(l=10, r=10, t=30, b=10), height=320)
                    st.plotly_chart(fig_bha_d, use_container_width=True)
                else:
                    st.info("No hay registros BHA para este día.")

                # ------------------------------
                # Conexiones perforando (diario)
                # ------------------------------
                st.markdown("### Conexiones perforando (diario)")
                df_conn_d = df_conn_filtrado.copy()
                if not df_conn_d.empty and "Fecha" in df_conn_d.columns:
                    df_conn_d["Fecha"] = df_conn_d["Fecha"].astype(str)
                    df_conn_d = df_conn_d[df_conn_d["Fecha"] == str(fecha_resumen)].copy()
                if not df_conn_d.empty:
                    df_conn_sum = df_conn_d.groupby("Componente", as_index=False)["Minutos_Reales"].sum()
                    df_conn_sum["Componente"] = pd.Categorical(df_conn_sum["Componente"], categories=CONN_ORDER, ordered=True)
                    df_conn_sum = df_conn_sum.sort_values("Componente")
                    fig_conn_pie_d = px.pie(
                        df_conn_sum,
                        names="Componente",
                        values="Minutos_Reales",
                        hole=0.35,
                        title=f"Distribución tiempo en conexión - {fecha_resumen}",
                        color="Componente",
                        color_discrete_map=CONN_COLOR_MAP,
                    )
                    st.plotly_chart(fig_conn_pie_d, use_container_width=True)

                    df_stack = df_conn_d.copy()
                    df_stack["Conn_Label"] = df_stack["Profundidad_m"].fillna(df_stack["Conn_No"]).astype(float).astype(int).astype(str)
                    df_stack["Componente"] = pd.Categorical(df_stack["Componente"], categories=CONN_ORDER, ordered=True)
                    df_stack_g = (
                        df_stack.groupby(["Conn_Label", "Componente"], as_index=False)["Minutos_Reales"]
                        .sum()
                        .sort_values(["Conn_Label", "Componente"])
                    )
                    per_conn = df_stack.groupby("Conn_Label", as_index=False).first()[["Conn_Label", "Conn_Tipo", "Angulo_Bucket"]]
                    per_conn["Std_Total"] = per_conn.apply(
                        lambda r: float(CONN_STDS.get((r["Conn_Tipo"], r["Angulo_Bucket"]), {}).get("TOTAL", 0.0)),
                        axis=1,
                    )
                    std_line = float(per_conn["Std_Total"].mean()) if not per_conn.empty else 0.0
                    fig_conn_stack_d = px.bar(
                        df_stack_g,
                        x="Conn_Label",
                        y="Minutos_Reales",
                        color="Componente",
                        category_orders={"Componente": CONN_ORDER},
                        color_discrete_map=CONN_COLOR_MAP,
                        barmode="stack",
                        title=f"Conexiones perforando - {fecha_resumen}",
                        labels={"Conn_Label": "Profundidad (m)", "Minutos_Reales": "Tiempo (min)"},
                    )
                    if std_line > 0:
                        fig_conn_stack_d.add_hline(
                            y=std_line,
                            line_dash="dash",
                            line_color="#9C640C",
                            annotation_text=f"{std_line:.1f}",
                            annotation_position="top left",
                            annotation_font_color="#9C640C",
                        )
                    fig_conn_stack_d.update_layout(legend_title_text="", xaxis_tickangle=0, height=320)
                    st.plotly_chart(fig_conn_stack_d, use_container_width=True)
                else:
                    st.info("No hay conexiones para este día.")

                # ------------------------------
                # Viajes (si aplica)
                # ------------------------------
                st.markdown("### Viajes (si aplica)")
                viajes_store = st.session_state.get("viajes_hourly_store", {})
                if isinstance(viajes_store, dict) and len(viajes_store) > 0:
                    viaje_tipo_sel = st.selectbox("Tipo de viaje (resumen diario)", options=sorted(list(viajes_store.keys())))
                    store_v = viajes_store.get(viaje_tipo_sel, {})
                    vel_std = float(VIAJE_CATALOG.get(viaje_tipo_sel, {}).get("vel_mh", 0.0)) if viaje_tipo_sel else 0.0
                    tconn_std = float(VIAJE_CATALOG.get(viaje_tipo_sel, {}).get("tconn_min", 0.0)) if viaje_tipo_sel else 0.0
                    usar_std_variable = bool(st.session_state.get(f"viaje_std_var_{viaje_tipo_sel}", False))
                    std_hourly_df = store_v.get("std_hourly")
                    hourly_df = store_v.get("hourly")
                    if isinstance(hourly_df, pd.DataFrame) and not hourly_df.empty:
                        df_plot = hourly_df.copy().sort_values("hour").reset_index(drop=True)
                        df_plot["hour_str"] = df_plot["hour"].astype(int)
                        day_start = int(st.session_state.get("day_start", 6))
                        day_end = int(st.session_state.get("day_end", 18))
                        def _is_day(h: int) -> bool:
                            if day_start == day_end:
                                return True
                            if day_start < day_end:
                                return day_start <= h < day_end
                            return (h >= day_start) or (h < day_end)
                        df_plot["Turno"] = df_plot["hour"].astype(int).apply(lambda h: "Día ☀️" if _is_day(h) else "Noche 🌙")
                        fig_v = px.bar(
                            df_plot,
                            x="hour_str",
                            y="speed_mh",
                            color="Turno",
                            color_discrete_map={"Día ☀️": "#F59E0B", "Noche 🌙": "#1D4ED8"},
                            labels={"hour_str": "Hora", "speed_mh": "m/h", "Turno": "Turno"},
                            title=f"Viaje – {viaje_tipo_sel}",
                        )
                        if usar_std_variable and isinstance(std_hourly_df, pd.DataFrame) and not std_hourly_df.empty:
                            _s = std_hourly_df.copy()
                            _s["hour_str"] = _s["hour"].astype(int)
                            fig_v.add_scatter(
                                x=_s["hour_str"],
                                y=_s["std_speed_mh"],
                                mode="lines",
                                name="Estándar",
                                line=dict(dash="dash", color="red"),
                            )
                        elif vel_std > 0:
                            fig_v.add_hline(
                                y=vel_std,
                                line_dash="dash",
                                line_color="red",
                                annotation_text=f"Estándar {vel_std:.0f}",
                                annotation_position="top left",
                            )
                        fig_v.update_layout(xaxis=dict(dtick=1))
                        st.plotly_chart(fig_v, use_container_width=True)

                        fig_c = px.bar(
                            df_plot,
                            x="hour_str",
                            y="conn_min",
                            color="Turno",
                            color_discrete_map={"Día ☀️": "#F59E0B", "Noche 🌙": "#1D4ED8"},
                            labels={"hour_str": "Hora", "conn_min": "min", "Turno": "Turno"},
                            title=f"Conexiones – {viaje_tipo_sel}",
                        )
                        if usar_std_variable and isinstance(std_hourly_df, pd.DataFrame) and not std_hourly_df.empty:
                            _s = std_hourly_df.copy()
                            _s["hour_str"] = _s["hour"].astype(int)
                            fig_c.add_scatter(
                                x=_s["hour_str"],
                                y=_s["std_conn_min"],
                                mode="lines",
                                name="Estándar",
                                line=dict(dash="dash", color="red"),
                            )
                        elif tconn_std > 0:
                            fig_c.add_hline(
                                y=tconn_std,
                                line_dash="dash",
                                line_color="red",
                                annotation_text=f"Estándar {tconn_std:.1f}",
                                annotation_position="top left",
                            )
                        fig_c.update_layout(xaxis=dict(dtick=1))
                        st.plotly_chart(fig_c, use_container_width=True)
                    else:
                        st.info("No hay datos de viajes para mostrar.")
                else:
                    st.info("No hay viajes registrados para el resumen diario.")

                # Tabla resumen diario
                with st.expander("Ver tabla diaria (etapa + fecha)", expanded=False):
                    cols_show = [c for c in ["Fecha","Etapa","Actividad","Tipo","Horas_Prog","Horas_Reales","Categoria_TNPI","Detalle_TNPI","Categoria_TNP","Detalle_TNP","Comentario"] if c in df_diario.columns]
                    st.dataframe(df_diario[cols_show], use_container_width=True, height=260)

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: INDICADORES ACTIVIDADES
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
with tab_act:
    # --- NUEVO: Vista de indicadores (diario vs acumulado) ---
    vista_ind = st.radio(
        "Vista de indicadores",
        ["Día seleccionado", "Acumulado (toda la jornada)"],
        index=0,
        horizontal=True,
        key="vista_indicadores",
    )

    # Base dataframe para indicadores
    df_ind_base = st.session_state.get("df", pd.DataFrame()).copy()

    # Filtrar por fecha seleccionada (puede incluir varias etapas)
    if vista_ind == "Día seleccionado":
        fecha_sel = st.session_state.get("fecha_val", None)
        if fecha_sel is not None and "Fecha" in df_ind_base.columns:
            df_ind_base["_Fecha_dt"] = pd.to_datetime(df_ind_base["Fecha"], errors="coerce")
            try:
                fecha_date = fecha_sel if hasattr(fecha_sel, "year") else pd.to_datetime(fecha_sel).date()
            except Exception:
                fecha_date = pd.to_datetime(fecha_sel, errors="coerce").date()
            df_ind_base = df_ind_base[df_ind_base["_Fecha_dt"].dt.date == fecha_date].copy()
            df_ind_base.drop(columns=["_Fecha_dt"], inplace=True, errors="ignore")

    st.subheader("Indicador de desempeño por actividades")
    rows_act = []
    if not df_ind_base.empty:
        g = df_ind_base.groupby(["Actividad", "Tipo"], as_index=False)["Horas_Reales"].sum()
        piv = g.pivot_table(index="Actividad", columns="Tipo", values="Horas_Reales", aggfunc="sum", fill_value=0.0).reset_index()
        for col in ["TP", "TNPI", "TNP"]:
            if col not in piv.columns:
                piv[col] = 0.0
        piv["Real"] = piv["TP"] + piv["TNPI"] + piv["TNP"]
        def _safe_pct(n, d):
            return (float(n) / float(d) * 100.0) if float(d) > 0 else 0.0
        piv["Eficiencia"] = piv.apply(lambda r: max(0.0, min(100.0, _safe_pct(r["TP"], r["Real"]))), axis=1)
        piv = piv.sort_values("Real", ascending=False)
        for _, r in piv.iterrows():
            rows_act.append({
                "name": str(r["Actividad"]),
                "real": f"{float(r['Real']):.2f}",
                "tnpi": f"{float(r['TNPI']):.2f}",
                "tnp": f"{float(r['TNP']):.2f}",
                "eff": float(r["Eficiencia"]),
                "semaforo": _semaforo_text(float(r["Eficiencia"]))
            })
    if rows_act:
        components.html(indicators_table_html("Indicador de desempeño por actividades", rows_act, kind="actividad"), height=520, scrolling=True)
    else:
        st.info("Aún no hay datos suficientes para indicador por actividades.")

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: TOP TNPI/TNP
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
with tab_top:
    st.subheader("Top 5 categorías – TNPI / TNP")

    df_top = st.session_state.get("df", pd.DataFrame()).copy()
    if df_top.empty:
        st.info("Aún no hay datos para calcular el top de TNPI/TNP.")
    else:
        df_top["Horas_Reales"] = pd.to_numeric(df_top.get("Horas_Reales", 0.0), errors="coerce").fillna(0.0)

        col_f1, col_f2, col_f3 = st.columns(3)
        with col_f1:
            days = _available_days(df_top)
            fecha_opts = ["Todas"] + [d.isoformat() for d in days]
            fecha_sel = st.selectbox(
                "Fecha",
                options=fecha_opts,
                index=(len(fecha_opts) - 1 if len(fecha_opts) > 1 else 0),
                key="top_tnpi_fecha",
            )
        with col_f2:
            etapas = sorted([str(x) for x in df_top.get("Etapa", pd.Series(dtype=str)).fillna("").unique().tolist() if str(x).strip() != ""])
            etapa_sel = st.selectbox(
                "Etapa",
                options=["Todas"] + etapas,
                index=0,
                key="top_tnpi_etapa",
            )
        with col_f3:
            modo_opts = ["Todos"]
            if "Modo_Reporte" in df_top.columns:
                modos = sorted([str(x) for x in df_top["Modo_Reporte"].dropna().unique().tolist() if str(x).strip() != ""])
                modo_opts += modos
            modo_sel = st.selectbox("Modo de reporte", options=modo_opts, index=0, key="top_tnpi_modo")

        # Aplicar filtros
        df_f = df_top.copy()
        if fecha_sel != "Todas" and "Fecha" in df_f.columns:
            try:
                fecha_dt = datetime.strptime(str(fecha_sel), "%Y-%m-%d").date()
            except Exception:
                fecha_dt = pd.to_datetime(str(fecha_sel), errors="coerce").date()
            df_f["_Fecha_dt"] = pd.to_datetime(df_f["Fecha"], errors="coerce").dt.date
            df_f = df_f[df_f["_Fecha_dt"] == fecha_dt].copy()
            df_f.drop(columns=["_Fecha_dt"], inplace=True, errors="ignore")
        if etapa_sel != "Todas" and "Etapa" in df_f.columns:
            df_f = df_f[df_f["Etapa"].astype(str) == str(etapa_sel)].copy()
        if modo_sel != "Todos" and "Modo_Reporte" in df_f.columns:
            df_f = df_f[df_f["Modo_Reporte"].astype(str) == str(modo_sel)].copy()

        # KPIs rápidos (chips)
        total_h_f = float(df_f["Horas_Reales"].sum()) if not df_f.empty and "Horas_Reales" in df_f.columns else 0.0
        tp_h_f = float(df_f[df_f["Tipo"] == "TP"]["Horas_Reales"].sum()) if "Tipo" in df_f.columns else 0.0
        tnpi_h_f = float(df_f[df_f["Tipo"] == "TNPI"]["Horas_Reales"].sum()) if "Tipo" in df_f.columns else 0.0
        tnp_h_f = float(df_f[df_f["Tipo"] == "TNP"]["Horas_Reales"].sum()) if "Tipo" in df_f.columns else 0.0
        eff_f = clamp_0_100(safe_pct(tp_h_f, total_h_f)) if total_h_f > 0 else 0.0
        tone_eff = "green" if eff_f >= 85 else ("amber" if eff_f >= 70 else "red")
        render_chip_row([
            {"label": "Horas total", "value": f"{total_h_f:.1f} h", "tone": "gray"},
            {"label": "TP", "value": f"{tp_h_f:.1f} h", "tone": "green"},
            {"label": "TNPI", "value": f"{tnpi_h_f:.1f} h", "tone": "amber"},
            {"label": "TNP", "value": f"{tnp_h_f:.1f} h", "tone": "red"},
            {"label": "Eficiencia", "value": f"{eff_f:.0f}%", "tone": tone_eff},
        ], use_iframe=True, height=110)

        # Chips con flechas: comparativo vs dia anterior disponible
        prev_day = None
        if fecha_sel != "Todas":
            try:
                fecha_dt_sel = datetime.strptime(str(fecha_sel), "%Y-%m-%d").date()
            except Exception:
                fecha_dt_sel = pd.to_datetime(str(fecha_sel), errors="coerce").date()
            days_all = _available_days(df_top)
            prevs = [d for d in days_all if d < fecha_dt_sel]
            prev_day = prevs[-1] if prevs else None

        if prev_day:
            df_prev = df_top.copy()
            if etapa_sel != "Todas" and "Etapa" in df_prev.columns:
                df_prev = df_prev[df_prev["Etapa"].astype(str) == str(etapa_sel)].copy()
            if modo_sel != "Todos" and "Modo_Reporte" in df_prev.columns:
                df_prev = df_prev[df_prev["Modo_Reporte"].astype(str) == str(modo_sel)].copy()
            if "Fecha" in df_prev.columns:
                df_prev["_Fecha_dt"] = pd.to_datetime(df_prev["Fecha"], errors="coerce").dt.date
                df_prev = df_prev[df_prev["_Fecha_dt"] == prev_day].copy()
                df_prev.drop(columns=["_Fecha_dt"], inplace=True, errors="ignore")

            total_prev = float(df_prev["Horas_Reales"].sum()) if not df_prev.empty and "Horas_Reales" in df_prev.columns else 0.0
            tp_prev = float(df_prev[df_prev["Tipo"] == "TP"]["Horas_Reales"].sum()) if "Tipo" in df_prev.columns else 0.0
            tnpi_prev = float(df_prev[df_prev["Tipo"] == "TNPI"]["Horas_Reales"].sum()) if "Tipo" in df_prev.columns else 0.0
            tnp_prev = float(df_prev[df_prev["Tipo"] == "TNP"]["Horas_Reales"].sum()) if "Tipo" in df_prev.columns else 0.0
            eff_prev = clamp_0_100(safe_pct(tp_prev, total_prev)) if total_prev > 0 else 0.0

            st.caption(f"Comparativo vs {prev_day.isoformat()}")
            render_chip_row([
                build_delta_chip_item("Δ Eficiencia", real=eff_f, prog=eff_prev, unit="%", higher_is_better=True, precision=0),
                build_delta_chip_item("Δ TNPI", real=tnpi_h_f, prog=tnpi_prev, unit="h", higher_is_better=False, precision=2),
                build_delta_chip_item("Δ TNP", real=tnp_h_f, prog=tnp_prev, unit="h", higher_is_better=False, precision=2),
                build_delta_chip_item("Δ Horas", real=total_h_f, prog=total_prev, unit="h", higher_is_better=False, precision=2),
            ], use_iframe=True, height=110)

        def _top_categorias(df_in: pd.DataFrame, tipo: str) -> pd.DataFrame:
            if df_in.empty or "Tipo" not in df_in.columns:
                return pd.DataFrame()
            d = df_in[df_in["Tipo"].astype(str).str.strip() == tipo].copy()
            if d.empty:
                return pd.DataFrame()

            def _pick_categoria_col(df_tipo: pd.DataFrame, tipo_local: str) -> str | None:
                if tipo_local == "TNPI":
                    if "Categoria_TNPI" in df_tipo.columns:
                        return "Categoria_TNPI"
                    if "Categoria" in df_tipo.columns:
                        return "Categoria"
                if tipo_local == "TNP":
                    if "Categoria_TNP" in df_tipo.columns:
                        return "Categoria_TNP"
                    if "Categoria_TNPI" in df_tipo.columns:
                        return "Categoria_TNPI"
                    if "Categoria" in df_tipo.columns:
                        return "Categoria"
                return None

            cat_col = _pick_categoria_col(d, tipo)
            if not cat_col:
                return pd.DataFrame()

            d["_Categoria_view"] = d[cat_col].astype(str)
            d["_Categoria_view"] = d["_Categoria_view"].replace(
                {"": "Sin categoría", "nan": "Sin categoría", "None": "Sin categoría", "-": "Sin categoría"}
            ).fillna("Sin categoría")

            if (
                tipo == "TNP"
                and cat_col == "Categoria_TNP"
                and d["_Categoria_view"].astype(str).str.strip().eq("Sin categoría").all()
                and "Categoria_TNPI" in d.columns
            ):
                d["_Categoria_view"] = d["Categoria_TNPI"].astype(str)
                d["_Categoria_view"] = d["_Categoria_view"].replace(
                    {"": "Sin categoría", "nan": "Sin categoría", "None": "Sin categoría", "-": "Sin categoría"}
                ).fillna("Sin categoría")

            g = d.groupby("_Categoria_view", as_index=False)["Horas_Reales"].sum()
            g = g.sort_values("Horas_Reales", ascending=False).head(5)
            g = g.rename(columns={"_Categoria_view": "Categoria"})
            return g

        top_tnpi = _top_categorias(df_f, "TNPI")
        top_tnp = _top_categorias(df_f, "TNP")

        c1, c2 = st.columns(2)
        with c1:
            st.markdown("#### TNPI")
            if top_tnpi.empty:
                st.info("No hay TNPI para los filtros seleccionados.")
            else:
                total_tnpi = float(top_tnpi["Horas_Reales"].sum())
                top_tnpi = top_tnpi.copy()
                top_tnpi["%"] = top_tnpi["Horas_Reales"].apply(lambda v: (float(v) / total_tnpi * 100.0) if total_tnpi > 0 else 0.0)
                top_tnpi["Semáforo"] = top_tnpi["%"].apply(semaforo_dot)
                fig_tnpi = px.bar(
                    top_tnpi.sort_values("Horas_Reales"),
                    x="Horas_Reales",
                    y="Categoria",
                    orientation="h",
                    title="Top 5 TNPI por categoría (h)",
                )
                fig_tnpi.update_traces(marker_color="#EF4444")
                st.plotly_chart(fig_tnpi, use_container_width=True)
                st.dataframe(top_tnpi, use_container_width=True, hide_index=True)
        with c2:
            st.markdown("#### TNP")
            if top_tnp.empty:
                st.info("No hay TNP para los filtros seleccionados.")
            else:
                total_tnp = float(top_tnp["Horas_Reales"].sum())
                top_tnp = top_tnp.copy()
                top_tnp["%"] = top_tnp["Horas_Reales"].apply(lambda v: (float(v) / total_tnp * 100.0) if total_tnp > 0 else 0.0)
                top_tnp["Semáforo"] = top_tnp["%"].apply(semaforo_dot)
                fig_tnp = px.bar(
                    top_tnp.sort_values("Horas_Reales"),
                    x="Horas_Reales",
                    y="Categoria",
                    orientation="h",
                    title="Top 5 TNP por categoría (h)",
                )
                fig_tnp.update_traces(marker_color="#3B82F6")
                st.plotly_chart(fig_tnp, use_container_width=True)
                st.dataframe(top_tnp, use_container_width=True, hide_index=True)

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: CONEXIONES
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
with tab_conn:
    if st.session_state.get("_toast_conn", False):
        st.success("Conexión agregada correctamente")
        st.session_state["_toast_conn"] = False

    st.subheader("Conexiones perforando")
    st.caption(f"Fecha en trabajo: {str(st.session_state.get('fecha_val', ''))}")

    # Usar siempre la copia más reciente de session_state para que al guardar en Detalle se vea el cambio sin añadir otra fila
    _df_conn_tab = st.session_state.get("df_conn", pd.DataFrame()).copy()
    if _df_conn_tab.empty and not df_conn.empty:
        _df_conn_tab = df_conn.copy()

    if modo_reporte != "Perforación":
        st.info("Cambia a modo **Perforación** para ver conexiones.")
    else:
        # ------------------------------
        # Selector de etapa (para separar gráficas por etapa)
        # ------------------------------
        etapas_conn = sorted(_df_conn_tab["Etapa"].dropna().unique().tolist()) if not _df_conn_tab.empty else []
        etapa_conn_view = st.selectbox(
            "Etapa para conexiones",
            options=etapas_conn if etapas_conn else ["Sin datos"],
            index=(etapas_conn.index(etapa) if etapas_conn and etapa in etapas_conn else 0),
            key="etapa_conn_view",
            help="Filtra las conexiones y sus gráficas por etapa (evita mezclar varias etapas en la misma gráfica).",
        )

        df_conn_view = _df_conn_tab[_df_conn_tab["Etapa"] == etapa_conn_view].copy() if (etapa_conn_view != "Sin datos" and not _df_conn_tab.empty) else pd.DataFrame()

        # ------------------------------
        # Chips pro: exceso vs estándar (con sugerencias)
        # ------------------------------
        if not df_conn_view.empty:
            try:
                per_conn = df_conn_view.groupby("Conn_No", as_index=False).first()[["Conn_No", "Conn_Tipo", "Angulo_Bucket"]]
                per_conn["Std_Total"] = per_conn.apply(
                    lambda r: float(CONN_STDS.get((r["Conn_Tipo"], r["Angulo_Bucket"]), {}).get("TOTAL", 0.0)),
                    axis=1,
                )
                total_std_min = float(per_conn["Std_Total"].sum())
                total_real_min = float(df_conn_view.groupby(["Conn_No"])["Minutos_Reales"].sum().sum())

                # Componente con mayor exceso
                comp_over = None
                if {"Componente", "Minutos_Reales", "Minutos_Estandar"}.issubset(df_conn_view.columns):
                    comp_sum = df_conn_view.groupby("Componente", as_index=False).agg(
                        real=("Minutos_Reales", "sum"),
                        std=("Minutos_Estandar", "sum"),
                    )
                    comp_sum["over"] = comp_sum["real"] - comp_sum["std"]
                    comp_sum = comp_sum.sort_values("over", ascending=False)
                    if not comp_sum.empty and float(comp_sum.iloc[0]["over"]) > 0:
                        comp_over = str(comp_sum.iloc[0]["Componente"])

                chips_exceso = _conn_exceso_suggestions(total_real_min, total_std_min, comp_over)
                if chips_exceso:
                    render_chip_row(chips_exceso, use_iframe=True, height=110)
            except Exception:
                pass

        # ------------------------------
        # Gráficas (pie + stacked) por etapa
        # ------------------------------
        if show_charts:
            if df_conn_view.empty:
                st.info("Aún no hay datos de conexiones para la etapa seleccionada.")
            else:
                # Pie por componentes
                if {"Componente", "Minutos_Reales"}.issubset(df_conn_view.columns):
                    df_conn_sum = df_conn_view.groupby("Componente", as_index=False)["Minutos_Reales"].sum()
                    df_conn_sum["Componente"] = pd.Categorical(df_conn_sum["Componente"], categories=CONN_ORDER, ordered=True)
                    df_conn_sum = df_conn_sum.sort_values("Componente")

                    fig_conn_pie = px.pie(
                        df_conn_sum,
                        names="Componente",
                        values="Minutos_Reales",
                        hole=0.35,
                        title=f"Distribución de tiempo en conexión - {etapa_conn_view}",
                        color="Componente",
                        color_discrete_map=CONN_COLOR_MAP,
                    )
                    st.plotly_chart(fig_conn_pie, use_container_width=True, key="pie_conexiones")

                # Stacked por conexión/profundidad
                df_stack = df_conn_view.copy()
                df_stack["Conn_Label"] = df_stack["Profundidad_m"].fillna(df_stack["Conn_No"]).astype(float).astype(int).astype(str)
                df_stack["Componente"] = pd.Categorical(df_stack["Componente"], categories=CONN_ORDER, ordered=True)
                df_stack_g = (
                    df_stack.groupby(["Conn_Label", "Componente"], as_index=False)["Minutos_Reales"]
                    .sum()
                    .sort_values(["Conn_Label", "Componente"])
                )

                per_conn = df_stack.groupby("Conn_Label", as_index=False).first()[["Conn_Label", "Conn_Tipo", "Angulo_Bucket"]]
                per_conn["Std_Total"] = per_conn.apply(
                    lambda r: float(CONN_STDS.get((r["Conn_Tipo"], r["Angulo_Bucket"]), {}).get("TOTAL", 0.0)),
                    axis=1,
                )
                std_line = float(per_conn["Std_Total"].mean()) if not per_conn.empty else 0.0

                fig_conn_stack = px.bar(
                    df_stack_g,
                    x="Conn_Label",
                    y="Minutos_Reales",
                    color="Componente",
                    category_orders={"Componente": CONN_ORDER},
                    color_discrete_map=CONN_COLOR_MAP,
                    barmode="stack",
                    title=f"Conexiones perforando - {etapa_conn_view}",
                    labels={"Conn_Label": "Profundidad (m)", "Minutos_Reales": "Tiempo (min)"},
                )

                if std_line > 0:
                    fig_conn_stack.add_hline(
                        y=std_line,
                        line_dash="dash",
                        line_color="#9C640C",
                        annotation_text=f"{std_line:.1f}",
                        annotation_position="top left",
                        annotation_font_color="#9C640C",
                    )

                df_tot = (
                    df_stack.groupby("Conn_Label", as_index=False)["Minutos_Reales"]
                    .sum()
                    .rename(columns={"Minutos_Reales": "Real_Total"})
                )
                tot_map = dict(zip(df_tot["Conn_Label"].astype(str), df_tot["Real_Total"]))
                for x in sorted(df_tot["Conn_Label"].astype(str).unique(), key=lambda v: float(v) if v.replace(".", "", 1).isdigit() else v):
                    y = float(tot_map.get(x, 0))
                    fig_conn_stack.add_annotation(x=x, y=y, text=f"<b>{y:.0f}</b>", showarrow=False, yshift=10)

                fig_conn_stack.update_layout(legend_title_text="", xaxis_tickangle=0)
                st.plotly_chart(fig_conn_stack, use_container_width=True, key="stack_conexiones")

        st.subheader("Indicador de desempeño por conexiones")
        rows_conn = []
        if not df_conn_view.empty:
            per = df_conn_view.groupby(["Conn_No", "Profundidad_m"], as_index=False).agg(
                real_min=("Minutos_Reales", "sum"),
                tnpi_min=("Minutos_TNPI", "sum"),
                tnp_min=("Minutos_TNP", "sum") if "Minutos_TNP" in df_conn_view.columns else ("Minutos_TNPI", "sum"),
            )
            per["eff"] = per.apply(
                lambda r: clamp_0_100(safe_pct(r["real_min"] - r["tnpi_min"] - r.get("tnp_min", 0.0), r["real_min"])) if r["real_min"] > 0 else 0.0,
                axis=1,
            )
            per = per.sort_values("Conn_No", ascending=True)

            for _, r in per.iterrows():
                name = f"#{int(r['Conn_No'])}  (Prof {float(r['Profundidad_m']):.0f} m)"
                rows_conn.append(
                    {
                        "name": name,
                        "real": f"{float(r['real_min']):.0f}",
                        "tnpi": f"{float(r['tnpi_min']):.0f}",
                        "tnp": f"{float(r.get('tnp_min', 0.0)):.0f}",
                        "eff": float(r["eff"]),
                    }
                )

        if rows_conn:
            components.html(indicators_table_html(f"Indicador de desempeño por conexiones - {etapa_conn_view}", rows_conn, kind="conexion"), height=420, scrolling=True)
        else:
            st.info("Aún no hay conexiones para indicador en la etapa seleccionada.")


# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: ROP (REAL VS PROGRAMADO)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =


with tab_viajes:
    st.subheader("Viajes y conexiones de TP")
    st.caption(f"Fecha en trabajo: {str(st.session_state.get('fecha_val', ''))}")

    _df_main = st.session_state.df

    # --- Selector por corrida (Run) y día grabado en esa corrida ---
    with st.expander("📅 Ver por corrida (Run) y día", expanded=True):
        corridas_en_df = []
        if not _df_main.empty and "Corrida" in _df_main.columns:
            corridas_en_df = sorted(_df_main["Corrida"].dropna().astype(str).str.strip().unique().tolist())
            corridas_en_df = [c for c in corridas_en_df if c]
        opts_corrida = ["(Usar fecha del sidebar)"] + corridas_en_df
        corrida_viajes_sel = st.selectbox(
            "Corrida (Run)",
            options=opts_corrida,
            index=0,
            key="viaje_corrida_sel",
            help="Elige una corrida para ver los días grabados en ella, o usa la fecha del sidebar.",
        )
        fecha_viajes_desde_selector = None
        if corrida_viajes_sel and corrida_viajes_sel != "(Usar fecha del sidebar)":
            _mask = _df_main["Corrida"].astype(str).str.strip() == str(corrida_viajes_sel).strip()
            _df_corrida = _df_main.loc[_mask]
            if not _df_corrida.empty and "Fecha" in _df_corrida.columns:
                _fechas_raw = _df_corrida["Fecha"].dropna().astype(str).str.strip().unique().tolist()
                _fechas_ordenadas = sorted(set(_fechas_raw))
                if _fechas_ordenadas:
                    dia_en_corrida_sel = st.selectbox(
                        "Día (grabado en esta corrida)",
                        options=_fechas_ordenadas,
                        index=len(_fechas_ordenadas) - 1,
                        key="viaje_dia_corrida_sel",
                        help="Días con registros en la corrida seleccionada.",
                    )
                    try:
                        for fmt in ("%Y-%m-%d", "%Y/%m/%d", "%d/%m/%Y"):
                            try:
                                fecha_viajes_desde_selector = datetime.strptime(str(dia_en_corrida_sel), fmt).date()
                                break
                            except ValueError:
                                continue
                    except Exception:
                        pass
        if fecha_viajes_desde_selector is not None:
            fecha_viajes = fecha_viajes_desde_selector
            st.caption(f"Mostrando datos del día **{str(fecha_viajes)}** (corrida *{corrida_viajes_sel}*).")
        else:
            fecha_viajes = st.session_state.get("fecha_val", datetime.today().date())

    # --- FILTRO DE ETAPA (Viajes y conexiones) ---
    _etapas_v = sorted(_df_main["Etapa"].dropna().unique().tolist()) if (not _df_main.empty and "Etapa" in _df_main.columns) else []
    etapa_viajes_sel = st.selectbox(
        "Etapa para viajes",
        options=_etapas_v,
        index=0 if _etapas_v else None,
        help="Filtra la vista/registro de viajes por etapa.",
        key="etapa_viajes_sel",
    ) if _etapas_v else None


    if "viajes_hourly_store" not in st.session_state:
        # Store por tipo de viaje y fecha (cada día su propia gráfica/datos)
        st.session_state["viajes_hourly_store"] = {}

    # fecha_viajes ya quedó definida arriba (selector por corrida/día o fecha del sidebar)
    colA, colB, colC = st.columns([1.4, 1.0, 1.0])

    with colA:
        viaje_tipo = st.selectbox(
            "Tipo de viaje",
            options=sorted(list(VIAJE_CATALOG.keys())) if "VIAJE_CATALOG" in globals() else [],
            help="Selecciona el tipo de viaje (catálogo de objetivos).",
            key="viaje_tipo_sel",
        )

    # Dirección (Trip In / Trip Out): mismo día puede tener ambos, se identifican con chips
    _base_key = f"{viaje_tipo}|{str(fecha_viajes)}" if viaje_tipo else ""
    _store = st.session_state.get("viajes_hourly_store", {})
    _has_trip_in = bool(_base_key and _store.get(f"{_base_key}|Trip In"))
    _has_trip_out = bool(_base_key and _store.get(f"{_base_key}|Trip Out"))
    direction_viajes = st.radio(
        "Dirección del viaje",
        options=["Trip In", "Trip Out"],
        index=0,
        key="viaje_direction_sel",
        horizontal=True,
        help="En un mismo día puedes tener Trip In y Trip Out; los chips abajo indican cuáles tienen datos.",
    )
    _bg_in = "#22c55e" if _has_trip_in else "rgba(255,255,255,0.12)"
    _fg_in = "#fff" if _has_trip_in else "rgba(255,255,255,0.6)"
    _bg_out = "#3b82f6" if _has_trip_out else "rgba(255,255,255,0.12)"
    _fg_out = "#fff" if _has_trip_out else "rgba(255,255,255,0.6)"
    _label_in = "Trip In ✓" if _has_trip_in else "Trip In (sin datos)"
    _label_out = "Trip Out ✓" if _has_trip_out else "Trip Out (sin datos)"
    st.markdown(
        f'<div style="display:flex; gap:8px; align-items:center; flex-wrap:wrap;">'
        f'<span style="background:{_bg_in}; color:{_fg_in}; padding:4px 10px; border-radius:999px; font-size:0.85rem;">{_label_in}</span>'
        f'<span style="background:{_bg_out}; color:{_fg_out}; padding:4px 10px; border-radius:999px; font-size:0.85rem;">{_label_out}</span>'
        f'</div>',
        unsafe_allow_html=True,
    )
    st.caption("Chips: verde/azul = hay datos grabados para ese sentido en este día.")

    _viaje_store_key = f"{_base_key}|{direction_viajes}" if _base_key else ""
    _viaje_ui_suffix = (_viaje_store_key or "").replace("|", "_").replace(" ", "_") or str(fecha_viajes)

    # Standards por catálogo
    vel_std = float(VIAJE_CATALOG.get(viaje_tipo, {}).get("vel_mh", 0.0)) if viaje_tipo else 0.0
    tconn_std = float(VIAJE_CATALOG.get(viaje_tipo, {}).get("tconn_min", 0.0)) if viaje_tipo else 0.0

    
    # Aliases (compatibilidad con bloques de cálculo/registro)
    v_std_mh = vel_std
    conn_std_min = tconn_std
    with colB:
        considerar_conexion = st.toggle(
            "Considerar tiempo de conexión",
            value=True,
            help="Si lo apagas, se omite el KPI de conexiones (solo viaje)."
        )

    with colC:
        distancia_manual = st.number_input(
            "Longitud (m) (opcional)",
            min_value=0.0,
            step=1.0,
            value=float(st.session_state.get("viaje_distancia_m", 0.0) or 0.0),
            help="Si importas KPIs, la longitud se calcula automáticamente; aquí puedes ajustar manual."
        )
        st.session_state["viaje_distancia_m"] = float(distancia_manual)

    st.caption(f"**Estándar:** {vel_std:.0f} m/h | **Conexión estándar:** {tconn_std:.1f} min")


    # ------------------------------
    # CORTE DE TURNOS (editable)
    # ------------------------------
    with st.expander("Corte de turnos (para colorear Día/Noche)", expanded=False):
        cts1, cts2 = st.columns(2)
        with cts1:
            day_start = st.number_input(
                "Inicio turno Día (hora 0–23)",
                min_value=0, max_value=23,
                value=int(st.session_state.get("day_start", 6)),
                step=1,
                key="viajes_day_start",
            )
        with cts2:
            day_end = st.number_input(
                "Fin turno Día (hora 0–23)",
                min_value=0, max_value=23,
                value=int(st.session_state.get("day_end", 18)),
                step=1,
                key="viajes_day_end",
            )
        st.session_state["day_start"] = int(day_start)
        st.session_state["day_end"] = int(day_end)
        st.caption(
            "Regla: Día si la hora está entre Inicio (incl.) y Fin (excl.). "
            "Si Inicio > Fin, se asume que el turno Día cruza medianoche."
        )

    # ------------------------------
    # IMPORTAR KPIs DrillSpot
    # ------------------------------
    with st.expander("Importar KPIs de DrillSpot (XLSX) para autocalcular por hora", expanded=False):
        up_kpi = st.file_uploader("Sube el export de KPIs (XLSX)", type=["xlsx"], key="kpi_xlsx_viajes")

        direction_default = default_trip_direction_from_activity(viaje_tipo) if viaje_tipo else "Trip In"
        direction = st.selectbox("Dirección para el cálculo", options=["Trip In", "Trip Out"], index=0 if direction_default == "Trip In" else 1)

        if st.button("Calcular automáticamente desde el XLSX", use_container_width=True, disabled=(up_kpi is None or not viaje_tipo)):
            df_kpi = load_drillspot_kpi_xlsx(up_kpi)
            hourly_df, meta = compute_viaje_conn_hourly_from_kpi(df_kpi, direction=direction)

            # Guarda en session por tipo, fecha y dirección (Trip In / Trip Out)
            _key_xlsx = f"{_base_key}|{direction}" if _base_key else ""
            if _key_xlsx:
                st.session_state["viajes_hourly_store"][_key_xlsx] = {
                    "hourly": hourly_df,
                    "meta": meta,
                    "direction": direction,
                    "considerar_conexion": considerar_conexion,
                }
            # Si hay longitud del KPI, úsala (pero permite ajustar)
            if meta.get("distance_m_total", 0.0) > 0:
                st.session_state["viaje_distancia_m"] = float(meta["distance_m_total"])
            st.success("KPIs importados y calculados ✅ (puedes editar manualmente abajo)")

    # ------------------------------
    # DATA MANUAL / EDITABLE
    # ------------------------------
    store = st.session_state["viajes_hourly_store"].get(_viaje_store_key, {})
    hourly_df = store.get("hourly")
    meta = store.get("meta", {}) if isinstance(store, dict) else {}

    if hourly_df is None or not isinstance(hourly_df, pd.DataFrame) or hourly_df.empty:
        hourly_df = pd.DataFrame({"hour": list(range(24)), "speed_mh": [0.0]*24, "conn_min": [0.0]*24})

    st.markdown("### Carga manual (por hora)")
    st.caption("Ingresa la **velocidad promedio (m/h)** por hora y (opcional) el **tiempo de conexión promedio (min)** por hora. "
               "Si importaste el XLSX, aquí podrás ajustar valores puntuales.")

    editable = hourly_df.copy()
    editable = editable.sort_values("hour").reset_index(drop=True)
    editable.rename(columns={"hour": "Hora", "speed_mh": "Velocidad real (m/h)", "conn_min": "Conexión real (min)"}, inplace=True)

    edited = st.data_editor(
        editable,
        use_container_width=True,
        hide_index=True,
        column_config={
            "Hora": st.column_config.NumberColumn("Hora", min_value=0, max_value=23, step=1, disabled=True),
            "Velocidad real (m/h)": st.column_config.NumberColumn("Velocidad real (m/h)", min_value=0.0, step=1.0),
            "Conexión real (min)": st.column_config.NumberColumn("Conexión real (min)", min_value=0.0, step=0.1),
        },
        num_rows="fixed",
        key=f"viajes_hourly_editor_{_viaje_ui_suffix}",
    )

    csave1, csave2 = st.columns([1, 1])
    with csave1:
        if st.button("Guardar ajustes manuales", use_container_width=True, disabled=(not viaje_tipo)):
            h2 = edited.rename(columns={"Hora": "hour", "Velocidad real (m/h)": "speed_mh", "Conexión real (min)": "conn_min"}).copy()
            h2["hour"] = h2["hour"].astype(int)
            for c in ["speed_mh", "conn_min"]:
                h2[c] = pd.to_numeric(h2[c], errors="coerce").fillna(0.0)

            if _viaje_store_key:
                st.session_state["viajes_hourly_store"][_viaje_store_key] = {
                    "hourly": h2,
                    "meta": meta,
                    "direction": store.get("direction", default_trip_direction_from_activity(viaje_tipo)),
                    "considerar_conexion": considerar_conexion,
                }
            st.success("Ajustes guardados ✅")

    with csave2:
        if st.button("Limpiar (poner en cero)", use_container_width=True, disabled=(not viaje_tipo)):
            h2 = pd.DataFrame({"hour": list(range(24)), "speed_mh": [0.0]*24, "conn_min": [0.0]*24})
            if _viaje_store_key:
                st.session_state["viajes_hourly_store"][_viaje_store_key] = {
                    "hourly": h2,
                    "meta": {},
                    "direction": store.get("direction", default_trip_direction_from_activity(viaje_tipo)),
                    "considerar_conexion": considerar_conexion,
                }
            st.success("Valores reiniciados ✅")
            st.rerun()

    # Recupera la versión guardada (después de edición)
    store = st.session_state["viajes_hourly_store"].get(_viaje_store_key, {})
    hourly_df = store.get("hourly", pd.DataFrame({"hour": list(range(24)), "speed_mh": [0.0]*24, "conn_min": [0.0]*24}))
    hourly_df = hourly_df.sort_values("hour").reset_index(drop=True)

    # ------------------------------
    # ESTÁNDAR VARIABLE POR HORA (OPCIONAL)
    # ------------------------------
    usar_std_variable = st.checkbox(
        "Estándar variable por hora (opcional)",
        value=bool(st.session_state.get(f"viaje_std_var_{viaje_tipo}", False)),
        key=f"viaje_std_var_{viaje_tipo}",
        help="Actívalo solo cuando el estándar cambie durante el viaje (por tramo / lingadas vs TxT, etc.). "
             "Si está apagado, se usa el estándar general (línea roja fija) como está hoy."
    )

    std_hourly_df = None
    if usar_std_variable and viaje_tipo:
        st.caption("Edita el estándar por hora. Esto NO reemplaza tu estándar general; solo se usa si activas este modo.")
        std_store = store.get("std_hourly")
        if std_store is None or not isinstance(std_store, pd.DataFrame) or std_store.empty:
            std_store = pd.DataFrame({
                "hour": list(range(24)),
                "std_speed_mh": [float(v_std_mh or 0.0)] * 24,
                "std_conn_min": [float(tconn_std or 0.0)] * 24,
                "conn_count": [0] * 24,
            })

        std_edit = std_store.copy().sort_values("hour").reset_index(drop=True)
        std_edit.rename(columns={
            "hour": "Hora",
            "std_speed_mh": "Estándar velocidad (m/h)",
            "std_conn_min": "Estándar conexión (min)",
            "conn_count": "Conexiones (#) en la hora",
        }, inplace=True)

        std_edited = st.data_editor(
            std_edit,
            use_container_width=True,
            hide_index=True,
            column_config={
                "Hora": st.column_config.NumberColumn("Hora", min_value=0, max_value=23, step=1, disabled=True),
                "Estándar velocidad (m/h)": st.column_config.NumberColumn("Estándar velocidad (m/h)", min_value=0.0, step=1.0),
                "Estándar conexión (min)": st.column_config.NumberColumn("Estándar conexión (min)", min_value=0.0, step=0.1),
                "Conexiones (#) en la hora": st.column_config.NumberColumn("Conexiones (#) en la hora", min_value=0, step=1),
            },
        )

        cstd1, cstd2 = st.columns(2)
        with cstd1:
            if st.button("Guardar estándar por hora", use_container_width=True, disabled=(not viaje_tipo)):
                s2 = std_edited.copy()
                s2.rename(columns={
                    "Hora": "hour",
                    "Estándar velocidad (m/h)": "std_speed_mh",
                    "Estándar conexión (min)": "std_conn_min",
                    "Conexiones (#) en la hora": "conn_count",
                }, inplace=True)
                s2["hour"] = s2["hour"].astype(int)
                for c in ["std_speed_mh", "std_conn_min"]:
                    s2[c] = pd.to_numeric(s2[c], errors="coerce").fillna(0.0)
                s2["conn_count"] = pd.to_numeric(s2["conn_count"], errors="coerce").fillna(0).astype(int)

                # Persistimos junto con el store del viaje (por fecha)
                if _viaje_store_key:
                    st.session_state["viajes_hourly_store"][_viaje_store_key] = {
                        "hourly": hourly_df,
                        "std_hourly": s2,
                        "meta": meta,
                        "direction": store.get("direction", default_trip_direction_from_activity(viaje_tipo)),
                        "considerar_conexion": considerar_conexion,
                    }
                st.success("Estándar por hora guardado ✅")
                st.rerun()

        with cstd2:
            if st.button("Reset estándar por hora", use_container_width=True, disabled=(not viaje_tipo)):
                s2 = pd.DataFrame({
                    "hour": list(range(24)),
                    "std_speed_mh": [float(v_std_mh or 0.0)] * 24,
                    "std_conn_min": [float(tconn_std or 0.0)] * 24,
                    "conn_count": [0] * 24,
                })
                if _viaje_store_key:
                    st.session_state["viajes_hourly_store"][_viaje_store_key] = {
                        "hourly": hourly_df,
                        "std_hourly": s2,
                        "meta": meta,
                        "direction": store.get("direction", default_trip_direction_from_activity(viaje_tipo)),
                        "considerar_conexion": considerar_conexion,
                    }
                st.success("Estándar por hora reiniciado ✅")
                st.rerun()

        # Recarga (después de guardar/reset)
        store = st.session_state["viajes_hourly_store"].get(_viaje_store_key, {})
        std_hourly_df = store.get("std_hourly")
        if std_hourly_df is not None and isinstance(std_hourly_df, pd.DataFrame) and not std_hourly_df.empty:
            std_hourly_df = std_hourly_df.sort_values("hour").reset_index(drop=True)


    # ------------------------------
    # GRÁFICAS
    # ------------------------------
    st.divider()
    st.markdown("### Gráficas")

    df_plot = hourly_df.copy()
    df_plot["hour_str"] = df_plot["hour"].astype(int)

    # Turno por hora (para colores Día/Noche)
    day_start = int(st.session_state.get("day_start", 6))
    day_end = int(st.session_state.get("day_end", 18))

    def _is_day(h: int) -> bool:
        if day_start == day_end:
            return True  # todo el día (caso extremo)
        if day_start < day_end:
            return day_start <= h < day_end
        # Cruza medianoche
        return (h >= day_start) or (h < day_end)

    df_plot["Turno"] = df_plot["hour"].astype(int).apply(lambda h: "Día ☀️" if _is_day(h) else "Noche 🌙")


    fig_v = px.bar(
        df_plot,
        x="hour_str",
        y="speed_mh",
        color="Turno",
        color_discrete_map={"Día ☀️": "#F59E0B", "Noche 🌙": "#1D4ED8"},
        labels={"hour_str": "Hora", "speed_mh": "m/h", "Turno": "Turno"},
        title=f"Viaje – {viaje_tipo}"
    )
    if usar_std_variable and std_hourly_df is not None and isinstance(std_hourly_df, pd.DataFrame) and not std_hourly_df.empty:
        # Línea estándar variable (por hora)
        _s = std_hourly_df.copy()
        _s["hour_str"] = _s["hour"].astype(int)
        fig_v.add_scatter(
            x=_s["hour_str"],
            y=_s["std_speed_mh"],
            mode="lines",
            name="Estándar",
            line=dict(dash="dash", color="red"),
        )
    elif vel_std > 0:
        fig_v.add_hline(
            y=vel_std,
            line_dash="dash",
            line_color="red",
            annotation_text=f"Estándar {vel_std:.0f}",
            annotation_position="top left",
        )
    fig_v.update_layout(showlegend=True, legend_title_text='', xaxis=dict(dtick=1))
    st.plotly_chart(fig_v, use_container_width=True, key=f"bar_viajes_velocidad_{_viaje_ui_suffix}")

    if considerar_conexion:
        fig_c = px.bar(
            df_plot,
            x="hour_str",
            y="conn_min",
            color="Turno",
            color_discrete_map={"Día ☀️": "#F59E0B", "Noche 🌙": "#1D4ED8"},
            labels={"hour_str": "Hora", "conn_min": "min", "Turno": "Turno"},
            title=f"Conexiones – {viaje_tipo}"
        )
        if usar_std_variable and std_hourly_df is not None and isinstance(std_hourly_df, pd.DataFrame) and not std_hourly_df.empty:
            _s = std_hourly_df.copy()
            _s["hour_str"] = _s["hour"].astype(int)
            fig_c.add_scatter(
                x=_s["hour_str"],
                y=_s["std_conn_min"],
                mode="lines",
                name="Estándar",
                line=dict(dash="dash", color="red"),
            )
        elif tconn_std > 0:
            fig_c.add_hline(
                y=tconn_std,
                line_dash="dash",
                line_color="red",
                annotation_text=f"Estándar {tconn_std:.1f}",
                annotation_position="top left",
            )

        fig_c.update_layout(showlegend=True, legend_title_text='', xaxis=dict(dtick=1))
        st.plotly_chart(fig_c, use_container_width=True, key=f"bar_viajes_conexiones_{_viaje_ui_suffix}")

    # ------------------------------
    # AVISO: botón de guardar más abajo
    # ------------------------------
    st.info(f"**Para guardar este viaje** en las actividades del día **{str(fecha_viajes)}**: baja a la sección **Registro en actividades** y pulsa el botón **Registrar este viaje en actividades** (ahí se guardan los datos de las gráficas de arriba).")

    # ------------------------------
    # RESUMEN (TABLA)
    # ------------------------------
    st.markdown("### Resumen")
    dist = float(st.session_state.get("viaje_distancia_m", 0.0) or 0.0)
    if isinstance(meta, dict) and meta.get("distance_m_total", 0.0) and dist <= 0:
        dist = float(meta.get("distance_m_total", 0.0))

    # Real promedio (sobre horas con dato > 0)
    speed_vals = hourly_df["speed_mh"].astype(float)
    speed_real = float(speed_vals[speed_vals > 0].mean()) if (speed_vals > 0).any() else 0.0

    conn_vals = hourly_df["conn_min"].astype(float)
    conn_real = float(conn_vals[conn_vals > 0].mean()) if (conn_vals > 0).any() else 0.0

    

    # Aliases (compatibilidad con lógica TNPI/registro)
    v_real_mh = speed_real
    conn_real_min = float(conn_real or 0.0)
    sum_df = pd.DataFrame([{
        "Tipo de viaje": viaje_tipo or "-",
        "Longitud (m)": dist if dist > 0 else "-",
        "Estándar (m/h)": vel_std if vel_std > 0 else "-",
        "Real (m/h)": round(speed_real, 1) if speed_real > 0 else "-",
        "Estándar (min)": tconn_std if (considerar_conexion and tconn_std > 0) else "-",
        "Real (min)": round(conn_real, 2) if (considerar_conexion and conn_real > 0) else "-",
    }])

    st.dataframe(sum_df, use_container_width=True, hide_index=True)


    # ------------------------------
    # REGISTRO EN ACTIVIDADES (para que cuente en TNPI / distribución / detalle)
    # ------------------------------
    st.markdown("### Registro en actividades")

    # Conexiones totales (para convertir min/conn a horas)
    n_conn_total_default = int(st.session_state.get("viaje_n_conn", 0) or 0)

    with st.expander("Configurar registro (opcional)", expanded=False):
        col_r1, col_r2, col_r3 = st.columns([1, 1, 1])
        with col_r1:
            turno_viaje = st.radio("Turno del viaje", options=["Diurno", "Nocturno"], horizontal=True, key=f"viaje_turno_{viaje_tipo}")
        with col_r2:
            n_conn_total = st.number_input(
                "Conexiones totales (#)",
                min_value=0,
                step=1,
                value=n_conn_total_default,
                key=f"viaje_nconn_total_{viaje_tipo}"
            )
        
        conexiones_total = int(n_conn_total or 0)
        with col_r3:
            actividad_nombre = st.text_input(
                "Nombre en actividades",
                value=(viaje_tipo if (viaje_tipo or '').lower().startswith('viaje') else f"Viaje {viaje_tipo}") if viaje_tipo else "Viaje",
                key=f"viaje_actname_{viaje_tipo}"
            )

        st.caption("El cálculo de horas usa: Horas = Distancia/Velocidad + (#Conexiones × min/conexión)/60 (si está habilitado).")

        # Categoría/detalle para TNPI si aplica
        cat_opts = (cat_list if 'cat_list' in globals() else ["-"])
        categoria_viaje = st.selectbox(
            "Categoría TNPI (si aplica)",
            options=cat_opts,
            index=0,
            key=f"viaje_cat_{viaje_tipo}"
        )
        detalle_viaje = st.text_input(
            "Detalle TNPI (si aplica)",
            value="",
            key=f"viaje_det_{viaje_tipo}"
        )
        comentario_viaje = st.text_input(
            "Comentario (opcional)",
            value="",
            key=f"viaje_com_{viaje_tipo}"
        )

        # Hora (opcional) para viaje
        viaje_use_time = st.checkbox(
            "Registrar hora (viaje)",
            value=False,
            key=f"viaje_use_time_{viaje_tipo}",
        )
        viaje_hora_ini = None
        viaje_hora_fin = None
        viaje_bitacora_enabled = False
        viaje_bitacora_entries = st.session_state.get("act_bitacora_entries", [])
        viaje_bitacora_total_h = float(st.session_state.get("act_bitacora_total_h", 0.0) or 0.0)
        if viaje_use_time:
            viaje_hora_ini = st.time_input(
                "Hora inicio (viaje)",
                value=st.session_state.get(f"viaje_hora_ini_{viaje_tipo}", datetime.now().time()),
                key=f"viaje_hora_ini_{viaje_tipo}",
            )
            viaje_hora_fin = st.time_input(
                "Hora fin (viaje)",
                value=st.session_state.get(f"viaje_hora_fin_{viaje_tipo}", datetime.now().time()),
                key=f"viaje_hora_fin_{viaje_tipo}",
            )
            viaje_bitacora_enabled = st.toggle("Bitácora por horas (viaje)", value=False, key=f"viaje_use_bitacora_{viaje_tipo}")
        viaje_hora_ini_txt = viaje_hora_ini.strftime("%H:%M") if (viaje_use_time and viaje_hora_ini) else ""
        viaje_hora_fin_txt = viaje_hora_fin.strftime("%H:%M") if (viaje_use_time and viaje_hora_fin) else ""
        if viaje_use_time and viaje_bitacora_enabled:
            st.caption("Completa la bitácora en la pestaña **Bitácora por horas**.")
            viaje_bitacora_mode = st.radio(
                "Uso de bitácora (viaje)",
                options=["Usar bitácora como total del viaje", "Seguir con cálculo estándar/KPI"],
                horizontal=True,
                key=f"viaje_bitacora_mode_{viaje_tipo}",
                help="Define si la bitácora reemplaza las horas reales del viaje o solo sirve como referencia.",
            )
        else:
            viaje_bitacora_mode = "Seguir con cálculo estándar/KPI"

        # Permite override de horas reales si no hay suficientes datos
        auto_real_h = 0.0
        if dist > 0 and speed_real > 0:
            auto_real_h = dist / speed_real
            if considerar_conexion and n_conn_total and conn_real > 0:
                auto_real_h += (float(n_conn_total) * float(conn_real) / 60.0)

        if viaje_use_time:
            if viaje_bitacora_enabled and viaje_bitacora_total_h > 0 and viaje_bitacora_mode == "Usar bitácora como total del viaje":
                horas_reales_override = float(viaje_bitacora_total_h)
                st.caption(f"Bitácora aplicada como total: {horas_reales_override:.2f} h")
            elif viaje_hora_ini and viaje_hora_fin:
                _dt_ini = datetime.combine(datetime.today().date(), viaje_hora_ini)
                _dt_fin = datetime.combine(datetime.today().date(), viaje_hora_fin)
                horas_reales_override = (_dt_fin - _dt_ini).total_seconds() / 3600.0
                if horas_reales_override < 0:
                    horas_reales_override += 24.0
                st.caption(f"Horas reales calculadas: {horas_reales_override:.2f} h")
            else:
                horas_reales_override = 0.0
                st.caption("Horas reales calculadas: 0.00 h")
            st.session_state[f"viaje_realh_override_{viaje_tipo}"] = float(horas_reales_override)
        else:
            horas_reales_override = st.number_input(
                "Horas reales (override, opcional)",
                min_value=0.0,
                step=0.1,
                value=float(auto_real_h) if auto_real_h > 0 else 0.0,
                key=f"viaje_realh_override_{viaje_tipo}",
                help="Si no quieres usar el cálculo automático (por velocidad), escribe aquí las horas reales totales del viaje."
            )

    # Horas estándar (desde catálogo) y reales (auto/override)
    n_conn_used = int(st.session_state.get(f"viaje_nconn_total_{viaje_tipo}", n_conn_total_default) or 0)

    # ------------------------------
    # Cálculos (estándar/real/TNPI) para registro
    # ------------------------------
    # Si NO está activado estándar variable por hora: usamos el estándar general (línea roja fija) como hasta ahora.
    # Si SÍ está activado: usamos std_hourly_df (por hora) para calcular estándar/real y TNPI por velocidad + conexiones.

    tnpi_vel_h = 0.0
    tnpi_conn_h = 0.0
    std_h_viaje = 0.0
    std_h_conn = 0.0
    real_h_viaje = 0.0
    real_h_conn = 0.0

    if usar_std_variable and std_hourly_df is not None and isinstance(std_hourly_df, pd.DataFrame) and not std_hourly_df.empty:
        # --- Viaje por horas (velocidad) ---
        _h = hourly_df.copy().sort_values("hour").reset_index(drop=True)
        _s = std_hourly_df.copy().sort_values("hour").reset_index(drop=True)

        # merge por hora
        _m = pd.merge(_h, _s, on="hour", how="left")
        _m["speed_mh"] = pd.to_numeric(_m["speed_mh"], errors="coerce").fillna(0.0)
        _m["std_speed_mh"] = pd.to_numeric(_m["std_speed_mh"], errors="coerce").fillna(0.0)

        # distancia por hora inferida (m). Si no cuadra con dist, escalamos para que sume "dist".
        _m["dist_h"] = _m["speed_mh"].clip(lower=0.0)  # m/h * 1h
        dist_infer = float(_m["dist_h"].sum() or 0.0)
        dist_obj = float(dist or 0.0)

        if dist_obj > 0 and dist_infer > 0:
            factor = dist_obj / dist_infer
            _m["dist_h"] = _m["dist_h"] * factor
        elif dist_obj > 0 and dist_infer == 0:
            # Sin distribución por hora: cae al método global
            _m["dist_h"] = 0.0

        # tiempo real por hora = dist_h / v_real (si v_real>0)
        _m["t_real_h"] = 0.0
        mask_vr = _m["speed_mh"] > 0
        _m.loc[mask_vr, "t_real_h"] = _m.loc[mask_vr, "dist_h"] / _m.loc[mask_vr, "speed_mh"]

        # tiempo estándar por hora = dist_h / v_std (si v_std>0)
        _m["t_std_h"] = 0.0
        mask_vs = _m["std_speed_mh"] > 0
        _m.loc[mask_vs, "t_std_h"] = _m.loc[mask_vs, "dist_h"] / _m.loc[mask_vs, "std_speed_mh"]

        # TNPI por velocidad (solo si real > std)
        tnpi_vel_h = float((_m["t_real_h"] - _m["t_std_h"]).clip(lower=0.0).sum() or 0.0)

        std_h_viaje = float(_m["t_std_h"].sum() or 0.0)
        real_h_viaje = float(_m["t_real_h"].sum() or 0.0)

        # --- Conexiones por horas ---
        _m["conn_min"] = pd.to_numeric(_m["conn_min"], errors="coerce").fillna(0.0)
        _m["std_conn_min"] = pd.to_numeric(_m["std_conn_min"], errors="coerce").fillna(0.0)
        _m["conn_count"] = pd.to_numeric(_m["conn_count"], errors="coerce").fillna(0).astype(int)

        real_h_conn = float((_m["conn_min"] * _m["conn_count"]).sum() / 60.0) if considerar_conexion else 0.0
        std_h_conn = float((_m["std_conn_min"] * _m["conn_count"]).sum() / 60.0) if considerar_conexion else 0.0

        tnpi_conn_h = max(0.0, real_h_conn - std_h_conn) if considerar_conexion else 0.0

        std_h = std_h_viaje + (std_h_conn if considerar_conexion else 0.0)
        real_h = real_h_viaje + (real_h_conn if considerar_conexion else 0.0)
        tnpi_h = tnpi_vel_h + tnpi_conn_h
        tp_h = max(0.0, real_h - tnpi_h)

    else:
        # --- Estándar global (como estaba) ---
        std_h = 0.0
        if dist > 0 and vel_std > 0:
            std_h = dist / vel_std
            if considerar_conexion and n_conn_used and tconn_std > 0:
                std_h += (float(n_conn_used) * float(tconn_std) / 60.0)

        real_h = float(st.session_state.get(f"viaje_realh_override_{viaje_tipo}", 0.0) or 0.0)
        if real_h <= 0 and dist > 0 and speed_real > 0:
            real_h = dist / speed_real
            if considerar_conexion and n_conn_used and conn_real > 0:
                real_h += (float(n_conn_used) * float(conn_real) / 60.0)


        # Componentes (global)
        std_h_viaje = (dist / vel_std) if (dist > 0 and vel_std > 0) else 0.0
        real_h_viaje = (dist / speed_real) if (dist > 0 and speed_real > 0) else 0.0
        std_h_conn = (float(n_conn_used) * float(tconn_std) / 60.0) if (considerar_conexion and n_conn_used and tconn_std > 0) else 0.0
        real_h_conn = (float(n_conn_used) * float(conn_real) / 60.0) if (considerar_conexion and n_conn_used and conn_real > 0) else 0.0

        # TNPI por velocidad (solo si v_real < v_std) y TNPI por conexiones (si aplica)
        tnpi_vel_h = max(0.0, (dist / speed_real) - (dist / vel_std)) if (dist > 0 and speed_real > 0 and vel_std > 0) else 0.0
        tnpi_conn_h = max(0.0, ((float(n_conn_used) * float(conn_real) / 60.0) - (float(n_conn_used) * float(tconn_std) / 60.0))) if (considerar_conexion and n_conn_used and conn_real > 0 and tconn_std > 0) else 0.0
        tnpi_h = tnpi_vel_h + (tnpi_conn_h if considerar_conexion else 0.0)

        tp_h = max(0.0, real_h - tnpi_h)



    cM1, cM2, cM3 = st.columns(3)
    cM1.metric("Estándar (h)", f"{std_h:.2f}")
    cM2.metric("Real (h)", f"{real_h:.2f}")
    cM3.metric("TNPI por exceso (h)", f"{tnpi_h:.2f}")

    chip_items_viaje = [
        build_delta_chip_item(
            "Δ Tiempo real",
            real=real_h,
            prog=std_h,
            unit="h",
            higher_is_better=False,
            precision=2,
        ),
        build_delta_chip_item(
            "Δ Velocidad",
            real=speed_real,
            prog=vel_std,
            unit="m/h",
            higher_is_better=True,
            precision=1,
        ),
    ]
    if considerar_conexion:
        chip_items_viaje.append(
            build_delta_chip_item(
                "Δ Conexión",
                real=conn_real,
                prog=tconn_std,
                unit="min",
                higher_is_better=False,
                precision=1,
            )
        )
    render_chip_row(chip_items_viaje, use_iframe=True, height=140)

    # Botón para registrar en el DataFrame principal (st.session_state.df)
    # Decide si al registrar quieres separar automáticamente el exceso como TNPI (sin perder el estándar general).
    auto_tnpi_por_desempeno = st.toggle(
        "Registrar TNPI automáticamente (exceso vs estándar)",
        value=True,
        help="Si está activo: se registra TP hasta el estándar y el exceso como TNPI. Si está apagado: se registra una sola fila con el tipo seleccionado (TP/TNPI/TNP)."
    )

    # Cuando el usuario quiere capturar TNPI de viajes de forma manual (p. ej. causas exógenas)
    # puede que no haya estándar calculable (std_h = 0) o que el TNPI por desempeño resulte 0.
    # En esos casos, este selector asegura que el registro se guarde como TNPI/TNP según corresponda.
    tipo_time_viaje = st.selectbox(
        "Tipo de tiempo a registrar (si no hay TNPI automático)",
        ["TP", "TNPI", "TNP"],
        index=1,
        key="tipo_time_viaje",
        help="Si el TNPI por desempeño sale 0 (o no hay estándar), selecciona TNPI para contabilizarlo en causa–raíz."
    )

    comp_tnpi_viaje = st.selectbox(
        "Componente TNPI (Viajes)",
        ["Velocidad", "Conexiones", "Otro"],
        index=0,
        key="comp_tnpi_viaje",
        help="Usado para graficar/desglosar TNPI de viajes en el tab Ejecutivo cuando el registro es manual (sin TNPI automático)."
    )

    if st.button("Registrar este viaje en actividades", use_container_width=True):
        # Validaciones básicas
        if float(real_h or 0.0) <= 0.0:
            st.warning("No hay horas para registrar (revisa longitud, velocidades y/o conexiones).")
        else:
            # TNPI calculado por desempeño (exceso en tiempo por velocidad + exceso en tiempo por conexiones)
            _tnpi_total_h = float(max(0.0, (tnpi_vel_h or 0.0) + (tnpi_conn_h or 0.0)))
            _std_h = float(std_h or 0.0)
            _real_h = float(real_h or 0.0)
            _tp_h = float(max(0.0, _real_h - _tnpi_total_h))

            # Estándares (para trazabilidad en el registro)
            # Aliases para compatibilidad (evita NameError si cambian nombres)
            conn_std = tconn_std
            conn_real = conn_real_min
            _std_speed_mh = float(v_std_mh or 0.0)
            _real_speed_mh = float(v_real_mh or 0.0)
            try:
                _std_conn_min = float(conn_std or 0.0)
            except Exception:
                _std_conn_min = 0.0
            try:
                _real_conn_min = float(conn_real or 0.0)
            except Exception:
                _real_conn_min = 0.0

            # Asegura valor por defecto si no existe en el scope local
            if "turno_registro" not in locals():
                turno_registro = turno
            if "actividad_registro" not in locals():
                actividad_registro = "Viaje"
            if "detalle_registro" not in locals():
                detalle_registro = ""
            if "categoria_tnpi_registro" not in locals():
                categoria_tnpi_registro = ""

            # Base común del registro (mismo esquema que el registro general)
            #
            # IMPORTANTE:
            # - El tab "Ejecutivo" filtra TNPI de viajes por Origen.
            # - Para evitar ceros por mismatch de etiquetas, usamos siempre:
            #     Origen = "Viajes y conexiones"
            # - Además, cuando se calcula TNPI automático por desempeño,
            #   registramos TNPI separado por componente (Velocidad / Conexiones)
            #   para que el executive pueda desglosarlo.
            _base = {
                "Equipo": equipo,
                "Pozo": pozo,
                "Etapa": ((etapa_viajes_sel or etapa) if "etapa_viajes_sel" in globals() else etapa),
                "Fecha": fecha,
                "Equipo_Tipo": st.session_state.get("equipo_tipo_val", ""),
                "Seccion": etapa,
                "Corrida": corrida,
                "Tipo_Agujero": tipo_agujero,
                "Operacion": operacion,
                "Turno": turno_registro if "turno_registro" in locals() else turno,
                "Actividad": actividad_registro if "actividad_registro" in locals() else "Viaje",
                "Detalle_TNPI": detalle_registro if "detalle_registro" in locals() else "",
                "Categoria_TNPI": categoria_tnpi_registro if "categoria_tnpi_registro" in locals() else "",
                "Origen": "Viajes y conexiones",
                "Hora_Inicio": viaje_hora_ini_txt,
                "Hora_Fin": viaje_hora_fin_txt,
                "Longitud_m": float(dist or 0.0),
                "std_speed_mh": _std_speed_mh,
                "real_speed_mh": _real_speed_mh,
                "std_conn_min": _std_conn_min,
                "real_conn_min": _real_conn_min,
                "VIAJE_TIPO": str(viaje_tipo or "").strip(),
            }

            _rows = []
            if auto_tnpi_por_desempeno and _tnpi_total_h > 0.0:
                # 1) Parte productiva (TP) hasta el estándar
                _rows.append({
                    **_base,
                    "Tipo": "TP",
                    "Horas_Prog": _std_h,
                    "Horas_Reales": _tp_h,
                    "TP_h": _tp_h,
                    "TNPI_h": 0.0,
                    "TNP_h": 0.0,
                    "ROP_mh": _real_speed_mh,
                })
                # 2) Exceso como TNPI, separado por componente
                #    (esto habilita el desglose Velocidad vs Conexiones en el tab Ejecutivo)
                _detalle_user = (detalle_registro if "detalle_registro" in locals() else "").strip()
                _cat_user = (categoria_tnpi_registro if "categoria_tnpi_registro" in locals() else "").strip()

                if float(tnpi_vel_h or 0.0) > 0.0:
                    _rows.append({
                        **_base,
                        "Tipo": "TNPI",
                        "Categoria_TNPI": _cat_user,
                        "Detalle_TNPI": f"Velocidad - {_detalle_user}" if _detalle_user else "Velocidad",
                        "Horas_Prog": 0.0,
                        "Horas_Reales": float(tnpi_vel_h),
                        "TP_h": 0.0,
                        "TNPI_h": float(tnpi_vel_h),
                        "TNP_h": 0.0,
                        "ROP_mh": 0.0,
                    })
                if float(tnpi_conn_h or 0.0) > 0.0:
                    _rows.append({
                        **_base,
                        "Tipo": "TNPI",
                        "Categoria_TNPI": _cat_user,
                        "Detalle_TNPI": f"Conexiones - {_detalle_user}" if _detalle_user else "Conexiones",
                        "Horas_Prog": 0.0,
                        "Horas_Reales": float(tnpi_conn_h),
                        "TP_h": 0.0,
                        "TNPI_h": float(tnpi_conn_h),
                        "TNP_h": 0.0,
                        "ROP_mh": 0.0,
                    })
            else:
                # Registro tradicional: una sola fila con el tipo elegido
                _tipo = st.session_state.get("tipo_time_viaje", "TNPI")
                
                # Registro tradicional: una sola fila con el tipo elegido
                _tipo = st.session_state.get("tipo_time_viaje", "TNPI")

                # Para registros manuales de TNPI (sin TNPI automático), prefijamos el detalle con el componente
                # para que el tab Ejecutivo pueda desglosar y graficar (Velocidad vs Conexiones).
                _base_row = dict(_base)
                if _tipo == "TNPI":
                    _comp = st.session_state.get("comp_tnpi_viaje", "Otro") or "Otro"
                    _det = str(_base_row.get("Detalle_TNPI", "") or "").strip()
                    if _comp in ("Velocidad", "Conexiones"):
                        if _det:
                            if not _det.lower().startswith(_comp.lower()):
                                _det = f"{_comp} - {_det}"
                        else:
                            _det = _comp
                    _base_row["Detalle_TNPI"] = _det

                _rows.append({
                    **_base_row,
                    "Tipo": _tipo,
                    "Horas_Prog": _std_h,
                    "Horas_Reales": _real_h,
                    "TP_h": _real_h if _tipo == "TP" else 0.0,
                    "TNPI_h": _real_h if _tipo == "TNPI" else 0.0,
                    "TNP_h": _real_h if _tipo == "TNP" else 0.0,
                    "ROP_mh": _real_speed_mh if _tipo == "TP" else 0.0,
                })

            if not _rows:
                st.warning("No hay horas para registrar (revisa longitud, velocidades y/o conexiones).")
            else:
                # Validación: no permitir que el día supere DAY_LIMIT_HOURS
                new_hours = float(sum([_safe_float(r.get("Horas_Reales", 0.0)) for r in _rows]))
                remaining = _remaining_day_hours(st.session_state.df, fecha)
                if remaining <= 0:
                    st.error(f"El día ya completó {DAY_LIMIT_HOURS:.0f}h. No se pueden agregar más actividades.")
                    st.stop()
                if new_hours > remaining + 1e-6:
                    st.error(f"No se puede agregar: quedan {remaining:.2f} h disponibles en el día.")
                    st.stop()
                # Validación: no permitir que el turno supere TURNO_LIMIT_HOURS (12h)
                _turno_viaje = locals().get("turno_registro", turno)
                hrs_turno_viaje = _day_used_hours_by_turno(st.session_state.df, fecha, _turno_viaje)
                restante_turno_viaje = max(0.0, TURNO_LIMIT_HOURS - hrs_turno_viaje)
                if new_hours > restante_turno_viaje + 1e-6:
                    st.error(
                        f"El turno **{_turno_viaje}** ya tiene {hrs_turno_viaje:.2f} h cargadas. "
                        f"No se pueden cargar más de {TURNO_LIMIT_HOURS:.0f} h por turno (quedan {restante_turno_viaje:.2f} h)."
                    )
                    st.stop()

                nueva = pd.DataFrame(_rows)
                st.session_state.df = pd.concat([st.session_state.df, nueva], ignore_index=True)
                st.success(f"Registro agregado: {len(_rows)} fila(s).")
                st.rerun()


with tab_bha:
    st.subheader("BHA (Arma/Desarma)")
    st.caption(f"Fecha en trabajo: {str(st.session_state.get('fecha_val', ''))}")

    df_bha = st.session_state.df_bha
    # --- FILTRO DE ETAPA (BHA) ---
    if (not df_bha.empty) and ("Etapa" in df_bha.columns):
        _etapas_bha = sorted(df_bha["Etapa"].dropna().unique().tolist())
        etapa_bha_sel = st.selectbox(
            "Etapa para BHA",
            options=_etapas_bha,
            index=0 if _etapas_bha else None,
            help="Filtra los registros BHA por etapa."
        ) if _etapas_bha else None
        df_bha = df_bha[df_bha["Etapa"] == etapa_bha_sel] if etapa_bha_sel else df_bha

    if df_bha.empty:
        st.info("Aún no hay registros BHA para graficar.")
    else:
        std_sum = float(pd.to_numeric(df_bha.get("Estandar_h", 0.0), errors="coerce").fillna(0.0).sum())
        real_sum = float(pd.to_numeric(df_bha.get("Real_h", 0.0), errors="coerce").fillna(0.0).sum())
        eff_avg = clamp_0_100(safe_pct(std_sum, real_sum)) if real_sum > 0 else 0.0
        eff_tone = "green" if eff_avg >= 85 else ("amber" if eff_avg >= 75 else "red")

        render_chip_row([
            build_delta_chip_item(
                "Δ Tiempo real BHA",
                real=real_sum,
                prog=std_sum,
                unit="h",
                higher_is_better=False,
                precision=2,
            ),
            {"label": "Eficiencia prom.", "value": f"{eff_avg:.0f}%", "tone": eff_tone},
            {"label": "Registros", "value": f"{len(df_bha)}", "tone": "gray"},
        ], use_iframe=True, height=120)

        n_bha = n_max_bha = min(50, len(df_bha))
        if n_max_bha <= 1:
            n_bha = n_max_bha
            st.caption("Mostrando el único registro disponible." if n_bha == 1 else "Sin registros para graficar.")
        else:
            n_bha = st.slider("Últimos registros a graficar", min_value=1, max_value=n_max_bha, value=min(12, n_max_bha))
        df_bha_last = df_bha.tail(n_bha).copy()

        # Eficiencia y semáforo (igual que en otras vistas)
        df_bha_last["Eficiencia_pct"] = df_bha_last.apply(
            lambda r: (float(r.get("Estandar_h", 0.0)) / float(r.get("Real_h", 0.0)) * 100.0) if float(r.get("Real_h", 0.0) or 0.0) > 0 else 0.0,
            axis=1
        )
        df_bha_last["Semáforo"] = df_bha_last["Eficiencia_pct"].apply(semaforo_dot)

        def _bha_label(row):
            try:
                t = int(row.get("BHA_Tipo", 0))
            except Exception:
                t = row.get("BHA_Tipo", "")
            return f"T{t} - {row.get('Accion','')}".strip(" -")

        df_bha_last["Etiqueta"] = df_bha_last.apply(_bha_label, axis=1)

        df_long = df_bha_last.melt(
            id_vars=["Etiqueta"],
            value_vars=["Estandar_h", "Real_h"],
            var_name="Serie",
            value_name="Horas"
        )

        fig_bha = px.bar(
            df_long,
            x="Etiqueta",
            y="Horas",
            color="Serie",
            barmode="group",
            title="BHA: Estándar vs Real (últimos registros)"
        )
        fig_bha.update_layout(xaxis_title="Etiqueta", yaxis_title="Horas", legend_title="Serie")
        fig_bha.update_traces(texttemplate="%{y:.0f}", textposition="inside")
        st.plotly_chart(fig_bha, use_container_width=True, key="bar_bha")

        st.dataframe(df_bha_last, use_container_width=True, hide_index=True)



with tab_rop:
    st.subheader("ROP – Registro diario (Real) + Programado por corrida")

    if modo_reporte != "Perforación":
        st.info("Esta pestaña aplica para modo **Perforación**.")
    else:
        # --- captura por fecha (evita que se 'arrastre' al cambiar de día) ---
        fecha_key = str(st.session_state.get("fecha_val", datetime.today().date()))

        def _get_by_date(etapa_data: dict, k: str, default: float = 0.0) -> float:
            try:
                return float((etapa_data.get(k, {}) or {}).get(fecha_key, default))
            except Exception:
                return float(default)

        # --- aplicar reseteos pendientes ANTES de instanciar widgets (evita StreamlitAPIException) ---
        if st.session_state.get("_pending_widget_resets"):
            for _k, _v in list(st.session_state["_pending_widget_resets"].items()):
                st.session_state[_k] = _v
            st.session_state["_pending_widget_resets"].clear()

        # asegurar que etapa_data exista antes de usarse
        etapa_data_rop = get_etapa_data(etapa)

        # Corrida activa (viene del sidebar: 'Corrida (Run)')
        corrida_activa = str(st.session_state.get("corrida_activa") or st.session_state.drill_day.get("corrida_activa") or "Run 1")

        st.caption(f"📌 Corrida activa: **{corrida_activa}** · Fecha seleccionada: **{fecha_key}**")

        sub_diario, sub_corrida = st.tabs(["📅 Diario", "🏷️ ROP programado por corrida"])

        # == == == == == == == == == == == == =
        # 📅 Diario: ROP real día/noche + ROP programado del día
        # == == == == == == == == == == == == =
        with sub_diario:
            c1, c2, c3 = st.columns(3)

            # --- ROP programada del día (se guarda por fecha, con corrida asociada) ---
            with c1:
                etapa_data_rop.setdefault("rop_prog_by_date", {})
                etapa_data_rop.setdefault("rop_prog_by_corrida", {})

                _rp_entry = (etapa_data_rop.get("rop_prog_by_date") or {}).get(fecha_key)
                if isinstance(_rp_entry, dict) and "rop_prog" in _rp_entry:
                    rp_default = float(_rp_entry.get("rop_prog") or 0.0)
                else:
                    # default sugerido: maestro por corrida -> (fallback) último valor por etapa -> rop_prog_etapa
                    rp_default = float((etapa_data_rop.get("rop_prog_by_corrida") or {}).get(corrida_activa, 0.0) or 0.0)
                    if rp_default <= 0:
                        rp_default = float(etapa_data_rop.get("rop_prog_total", 0.0) or 0.0)
                    if rp_default <= 0:
                        rp_default = float(etapa_data_rop.get("rop_prog_etapa", 0.0) or 0.0)

                rp = st.number_input(
                    f"ROP programada (m/h) - {fecha_key}",
                    min_value=0.0, step=0.1,
                    value=float(rp_default),
                    key=f"rop_prog_diaria_{etapa}_{fecha_key}",
                    help="Registro diario. Por defecto toma el valor del maestro de la corrida activa (si existe).",
                )

                # Guardar "foto del día" (no se recalcula si mañana cambias de corrida/plan)
                etapa_data_rop["rop_prog_by_date"][fecha_key] = {"corrida_id": corrida_activa, "rop_prog": float(rp)}

                # Compatibilidad con otros bloques que aún lean rop_prog_total
                etapa_data_rop["rop_prog_total"] = float(rp)
                st.session_state.drill_day["rop_prog_total"] = float(rp)

            # --- ROP real Día ---
            with c2:
                rop_dia_val = _get_by_date(etapa_data_rop, "rop_real_dia_by_date", 0.0)
                rop_dia_val = st.number_input(
                    f"ROP real Día ☀️ - {etapa} (m/h)",
                    min_value=0.0, step=0.1,
                    value=float(rop_dia_val),
                    key=f"rop_real_dia_{etapa}_{fecha_key}",
                )
                etapa_data_rop["rop_real_dia"] = float(rop_dia_val)
                st.session_state.drill_day["rop_real_dia"] = float(rop_dia_val)

                etapa_data_rop.setdefault("rop_real_dia_by_date", {})
                if float(rop_dia_val) > 0:
                    etapa_data_rop["rop_real_dia_by_date"][fecha_key] = float(rop_dia_val)

            # --- ROP real Noche ---
            with c3:
                rop_noche_val = _get_by_date(etapa_data_rop, "rop_real_noche_by_date", 0.0)
                rop_noche_val = st.number_input(
                    f"ROP real Noche 🌙 - {etapa} (m/h)",
                    min_value=0.0, step=0.1,
                    value=float(rop_noche_val),
                    key=f"rop_real_noche_{etapa}_{fecha_key}",
                )
                etapa_data_rop["rop_real_noche"] = float(rop_noche_val)
                st.session_state.drill_day["rop_real_noche"] = float(rop_noche_val)

                etapa_data_rop.setdefault("rop_real_noche_by_date", {})
                if float(rop_noche_val) > 0:
                    etapa_data_rop["rop_real_noche_by_date"][fecha_key] = float(rop_noche_val)

            # (Metros perforados se registran en la pestaña 'Metros')


            # Sincroniza (compatibilidad con otros bloques que lean claves sueltas)
            st.session_state["rop_prog_total"] = float(st.session_state.drill_day["rop_prog_total"])
            st.session_state["rop_real_diurno"] = float(st.session_state.drill_day["rop_real_dia"])
            st.session_state["rop_real_nocturno"] = float(st.session_state.drill_day["rop_real_noche"])

            rp = float(st.session_state.drill_day.get("rop_prog_total", 0.0) or 0.0)
            rd = float(st.session_state.drill_day.get("rop_real_dia", 0.0) or 0.0)
            rn = float(st.session_state.drill_day.get("rop_real_noche", 0.0) or 0.0)

            # Promedio ponderado por metros (si están capturados), si no, promedio simple de turnos no-cero
            md = float(st.session_state.drill_day.get("metros_real_dia", 0.0) or 0.0)
            mn = float(st.session_state.drill_day.get("metros_real_noche", 0.0) or 0.0)
            if (md + mn) > 0:
                rr_avg = ((rd * md) + (rn * mn)) / (md + mn)
            else:
                vals = [v for v in [rd, rn] if v > 0]
                rr_avg = sum(vals) / len(vals) if vals else 0.0

            eff_rop_day = clamp_0_100(safe_pct(rr_avg, rp)) if rp > 0 else 0.0
            _, sl, sc = status_from_eff(eff_rop_day)

            k1, k2, k3, k4 = st.columns([1.2, 1.2, 1.2, 1.0])
            k1.metric("ROP real promedio (m/h)", f"{rr_avg:.2f}")
            k2.metric("ROP programada del día (m/h)", f"{rp:.2f}")
            k3.metric("Eficiencia ROP (%)", f"{eff_rop_day:.0f}%")
            with k4:
                st.markdown(
                    f"""<div style="display:flex;align-items:center;gap:10px;margin-top:28px;">
                        <span style="display:inline-block;width:12px;height:12px;border-radius:50%;background:{sc};box-shadow:0 0 0 2px rgba(255,255,255,0.08);"></span>
                        <div style="font-weight:800;font-size:22px;letter-spacing:0.5px;">{sl}</div>
                    </div>""",
                    unsafe_allow_html=True,
                )

            render_chip_row([
                build_delta_chip_item(
                    "Δ ROP real",
                    real=rr_avg,
                    prog=rp,
                    unit="m/h",
                    higher_is_better=True,
                    precision=2,
                )
            ], use_iframe=True, height=90)

            # Gráfica
            df_rop = pd.DataFrame(
                [
                    {"Turno": "Día ☀️", "Programado (m/h)": rp, "Real (m/h)": rd},
                    {"Turno": "Noche 🌙", "Programado (m/h)": rp, "Real (m/h)": rn},
                ]
            )
            fig_rop = px.bar(df_rop, x="Turno", y=["Programado (m/h)", "Real (m/h)"], barmode="group", text_auto=True)
            fig_rop.update_layout(margin=dict(l=10, r=10, t=30, b=10), height=340, legend_title_text="Serie")
            st.plotly_chart(fig_rop, use_container_width=True, key="bar_rop")


            # == == == == == == == == == == == == =
            # Tendencia por fecha (ROP)
            # == == == == == == == == == == == == =
            st.markdown("### Tendencia por fecha")

            # Construir serie por fecha desde registros diarios
            _prog_map = etapa_data_rop.get("rop_prog_by_date") or {}
            _rd_map = etapa_data_rop.get("rop_real_dia_by_date") or {}
            _rn_map = etapa_data_rop.get("rop_real_noche_by_date") or {}
            _md_map = etapa_data_rop.get("metros_real_dia_by_date") or {}
            _mn_map = etapa_data_rop.get("metros_real_noche_by_date") or {}

            _dates = sorted({*list(_prog_map.keys()), *list(_rd_map.keys()), *list(_rn_map.keys())})
            trend_rows = []
            if _dates:
                for _d in _dates:
                    _p_entry = _prog_map.get(_d, {})
                    _prog = _safe_float(_p_entry.get("rop_prog") if isinstance(_p_entry, dict) else (_p_entry or 0.0))
                    _rd = _safe_float(_rd_map.get(_d, 0.0) or 0.0)
                    _rn = _safe_float(_rn_map.get(_d, 0.0) or 0.0)

                    _md = _safe_float(_md_map.get(_d, 0.0) or 0.0)
                    _mn = _safe_float(_mn_map.get(_d, 0.0) or 0.0)

                    # mismo criterio que el KPI diario: ponderado por metros si existen
                    if (_md + _mn) > 0:
                        _real_avg = ((_rd * _md) + (_rn * _mn)) / (_md + _mn)
                    else:
                        _vals = [v for v in [_rd, _rn] if v > 0]
                        _real_avg = sum(_vals) / len(_vals) if _vals else 0.0

                    _eff = clamp_0_100(safe_pct(_real_avg, _prog)) if _prog > 0 else 0.0
                    _sem = "🟢" if _eff >= 85 else ("🟡" if _eff >= 70 else "🔴")

                    trend_rows.append({
                        "Fecha": _d,
                        "Programado": round(_prog, 2),
                        "Real Día ☀️": round(_rd, 2),
                        "Real Noche 🌙": round(_rn, 2),
                        "Real Promedio": round(_real_avg, 4),
                        "Eficiencia_pct": round(_eff, 4),
                        "Semáforo": _sem,
                    })

                df_tr = pd.DataFrame(trend_rows)

                fig_tr = go.Figure()
                fig_tr.add_bar(x=df_tr["Fecha"], y=df_tr["Real Día ☀️"], name="Real Día ☀️", marker_color="rgba(245,158,11,0.90)")
                fig_tr.add_bar(x=df_tr["Fecha"], y=df_tr["Real Noche 🌙"], name="Real Noche 🌙", marker_color="rgba(29,78,216,0.85)")
                fig_tr.add_trace(go.Scatter(x=df_tr["Fecha"], y=df_tr["Programado"], mode="lines+markers", name="Programado", line=dict(color="rgba(96,165,250,0.95)", width=3), marker=dict(size=8)))
                fig_tr.add_trace(go.Scatter(x=df_tr["Fecha"], y=df_tr["Real Promedio"], mode="lines+markers", name="Real Promedio", line=dict(color="rgba(249,115,22,0.95)", width=3), marker=dict(size=8)))

                fig_tr.update_layout(
                    barmode="stack",
                    height=420,
                    margin=dict(l=10, r=10, t=20, b=10),
                    legend=dict(orientation="v", x=1.02, y=1.0),
                    xaxis_title="Fecha",
                    yaxis_title="ROP (m/h)",
                )
                st.plotly_chart(fig_tr, use_container_width=True, key=f"trend_rop_{etapa}")

                st.dataframe(df_tr, use_container_width=True, hide_index=True)
            else:
                st.info("Aún no hay suficientes registros diarios para mostrar la tendencia.")

            # Detalle + semáforo por turno
            def _eff_turno(real_v: float, prog_v: float) -> float:
                return clamp_0_100(safe_pct(real_v, prog_v)) if prog_v > 0 else 0.0

            rows = []
            for turno_lbl, real_v in [("Día ☀️", rd), ("Noche 🌙", rn)]:
                e = _eff_turno(real_v, rp)
                rows.append(
                    {
                        "Turno": turno_lbl,
                        "ROP Programado (m/h)": round(rp, 2),
                        "ROP Real (m/h)": round(real_v, 2),
                        "Eficiencia (%)": round(e, 0),
                        "Semáforo": "🟢" if e >= 85 else ("🟡" if e >= 70 else "🔴"),
                    }
                )
            st.markdown("### Detalle")
            st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True)

        # == == == == == == == == == == == == =
        # 🏷️ Maestro: ROP programado por corrida (editable)
        # == == == == == == == == == == == == =
        with sub_corrida:
            st.markdown("### Maestro: ROP programado por corrida")
            st.caption("Guarda aquí el **ROP programado base** de la corrida. Este valor se propone como *default* en el registro diario si ese día aún no tiene ROP programada.")

            etapa_data_rop.setdefault("rop_prog_by_corrida", {})
            etapa_data_rop.setdefault("rop_prog_by_corrida_meta", {})
            etapa_data_rop.setdefault("rop_prog_by_date", {})

            # Valor guardado actualmente (si existe)
            rp_saved = float((etapa_data_rop["rop_prog_by_corrida"].get(corrida_activa, 0.0)) or 0.0)
            rp_master = st.number_input(
                f"ROP programada de la corrida **{corrida_activa}** (m/h)",
                min_value=0.0, step=0.1,
                value=float(rp_saved),
                key=f"rop_prog_corrida_{etapa}_{corrida_activa}",
                help="⚠️ Nota: escribir el valor NO lo guarda automáticamente; debes presionar **Guardar** para que quede registrado en el maestro.",
            )

            # Estado visual: guardado vs pendiente
            is_pending = abs(float(rp_master) - float(rp_saved)) > 1e-9
            last_meta = (etapa_data_rop["rop_prog_by_corrida_meta"].get(corrida_activa) or {})
            last_ts = str(last_meta.get("updated_at") or "").strip()
            last_by = str(last_meta.get("updated_by") or "").strip()

            top = st.columns([1.1, 1.1, 1.2])
            with top[0]:
                st.metric("Guardado (maestro)", f"{rp_saved:.2f} m/h" if rp_saved > 0 else "-")
            with top[1]:
                st.metric("Entrada actual", f"{float(rp_master):.2f} m/h" if float(rp_master) > 0 else "-")
            with top[2]:
                badge = "🟡 Pendiente de guardar" if is_pending else ("🟢 Guardado" if rp_saved > 0 else "⚪ Sin registrar")
                extra = (f" · {last_ts}" if last_ts else "")
                st.markdown(f"**Estado:** {badge}{extra}")

            btn_a, btn_b, btn_c = st.columns([1.1, 1.1, 0.9])
            with btn_a:
                if st.button("💾 Guardar en maestro", use_container_width=True, key=f"save_master_{etapa}_{corrida_activa}"):
                    etapa_data_rop["rop_prog_by_corrida"][corrida_activa] = float(rp_master)
                    etapa_data_rop["rop_prog_by_corrida_meta"][corrida_activa] = {
                        "updated_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        "updated_by": (st.session_state.get("username") or ""),
                    }
                    st.success("Maestro actualizado (ROP programado por corrida).")
                    st.rerun()
            with btn_b:
                if st.button("📅 Copiar a registro diario (fecha)", use_container_width=True, key=f"copy_master_to_day_{etapa}_{corrida_activa}"):
                    # asegura que el maestro quede guardado también
                    etapa_data_rop["rop_prog_by_corrida"][corrida_activa] = float(rp_master)
                    etapa_data_rop["rop_prog_by_corrida_meta"][corrida_activa] = {
                        "updated_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        "updated_by": (st.session_state.get("username") or ""),
                    }
                    # "foto del día"
                    etapa_data_rop["rop_prog_by_date"][fecha_key] = {"corrida_id": corrida_activa, "rop_prog": float(rp_master)}
                    etapa_data_rop["rop_prog_total"] = float(rp_master)
                    st.session_state.drill_day["rop_prog_total"] = float(rp_master)
                    st.success(f"Registro diario actualizado para {fecha_key}.")
                    st.rerun()
            with btn_c:
                _show_master_key = f"_show_master_{etapa}"
                _is_showing = bool(st.session_state.get(_show_master_key, False))
                if not _is_showing:
                    if st.button("🧾 Ver maestro", use_container_width=True, key=f"show_master_{etapa}_{corrida_activa}"):
                        st.session_state[_show_master_key] = True
                        st.rerun()
                else:
                    if st.button("🙈 Ocultar maestro", use_container_width=True, key=f"hide_master_{etapa}_{corrida_activa}"):
                        st.session_state[_show_master_key] = False
                        st.rerun()

            show_master = bool(st.session_state.get(f"_show_master_{etapa}", False))

            # Tabla del maestro (con metadatos)
            st.markdown("---")
            st.markdown("#### Maestro registrado (por corrida)")
            if etapa_data_rop["rop_prog_by_corrida"]:
                rows = []
                for k, v in (etapa_data_rop["rop_prog_by_corrida"].items() or []):
                    meta = (etapa_data_rop.get("rop_prog_by_corrida_meta", {}) or {}).get(k, {}) or {}
                    rows.append({
                        "Corrida": str(k),
                        "ROP_Prog_mh": float(v),
                        "Actualizado": str(meta.get("updated_at") or ""),
                        "Usuario": str(meta.get("updated_by") or ""),
                    })
                df_master = pd.DataFrame(rows).sort_values("Corrida").reset_index(drop=True)

                if not show_master:
                    st.info("Maestro oculto (presiona **Ver maestro** para mostrarlo).")
                else:
                    st.dataframe(df_master, use_container_width=True, hide_index=True)

                # Ayuda: por qué 'no se ve guardado'
                st.caption("💡 Si escribes un valor y cambias de pestaña sin presionar **Guardar en maestro**, el valor queda solo como *entrada* y no se registra en el maestro.")
            else:
                st.info("Aún no hay ROP programado guardado por corrida para esta etapa.")

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
with tab_metros:
    st.subheader("Metros perforados - Registro diario")

    corrida_activa = str(st.session_state.get("corrida_activa", ""))
    fecha_key = str(st.session_state.get("fecha_val", datetime.today().date()))
    st.caption(f"📌 Corrida activa: **{corrida_activa or '-'}** · Fecha seleccionada: **{fecha_key}**")

    etapa_data_rop = get_etapa_data(etapa)

    def _get_by_date(d: dict, key: str, default=0.0):
        dd = d.get(key) or {}
        try:
            return float(dd.get(fecha_key, default))
        except Exception:
            return float(default)

    st.markdown("### Metros perforados (registro diario)")
    m1, m2, m3 = st.columns(3)

    # --- Metros programados del día (total) ---
    with m1:
        etapa_data_rop.setdefault("metros_prog_by_date", {})
        _mp_entry = (etapa_data_rop.get("metros_prog_by_date") or {}).get(fecha_key)
        if isinstance(_mp_entry, dict) and "metros_prog" in _mp_entry:
            mp_default = _safe_float(_mp_entry.get("metros_prog") or 0.0)
        else:
            mp_default = _safe_float(etapa_data_rop.get("metros_prog_total", 0.0) or 0.0)

        mp = st.number_input(
            f"Metros programados (m) - {fecha_key}",
            min_value=0.0, step=0.1,
            value=float(mp_default),
            key=f"metros_prog_diaria_{etapa}_{fecha_key}",
            help="Registro diario. Se guarda por fecha y se asocia a la corrida activa."
        )
        etapa_data_rop["metros_prog_total"] = float(mp)
        st.session_state.drill_day["metros_prog_total"] = float(mp)
        etapa_data_rop["metros_prog_by_date"][fecha_key] = {"corrida_id": corrida_activa, "metros_prog": float(mp)}

    # --- Metros reales Día ---
    with m2:
        mr_dia_val = _get_by_date(etapa_data_rop, "metros_real_dia_by_date", 0.0)
        mr_dia_val = st.number_input(
            f"Metros reales Día ☀️ - {etapa} (m)",
            min_value=0.0, step=0.1,
            value=float(mr_dia_val),
            key=f"metros_real_dia_{etapa}_{fecha_key}",
        )
        etapa_data_rop["metros_real_dia"] = float(mr_dia_val)
        st.session_state.drill_day["metros_real_dia"] = float(mr_dia_val)
        etapa_data_rop.setdefault("metros_real_dia_by_date", {})
        if float(mr_dia_val) > 0:
            etapa_data_rop["metros_real_dia_by_date"][fecha_key] = float(mr_dia_val)

    # --- Metros reales Noche ---
    with m3:
        mr_noche_val = _get_by_date(etapa_data_rop, "metros_real_noche_by_date", 0.0)
        mr_noche_val = st.number_input(
            f"Metros reales Noche 🌙 - {etapa} (m)",
            min_value=0.0, step=0.1,
            value=float(mr_noche_val),
            key=f"metros_real_noche_{etapa}_{fecha_key}",
        )
        etapa_data_rop["metros_real_noche"] = float(mr_noche_val)
        st.session_state.drill_day["metros_real_noche"] = float(mr_noche_val)
        etapa_data_rop.setdefault("metros_real_noche_by_date", {})
        if float(mr_noche_val) > 0:
            etapa_data_rop["metros_real_noche_by_date"][fecha_key] = float(mr_noche_val)

    # KPI + gráfica
    mr_total = float(st.session_state.drill_day.get("metros_real_dia", 0.0)) + float(st.session_state.drill_day.get("metros_real_noche", 0.0))
    eff_m = clamp_0_100(safe_pct(mr_total, float(mp))) if float(mp) > 0 else 0.0
    st.caption(f"📌 Metros reales total: **{mr_total:.2f} m** · Eficiencia metros: **{eff_m:.0f}%**")


    render_chip_row([
        build_delta_chip_item(
            "Δ Metros reales",
            real=mr_total,
            prog=float(mp),
            unit="m",
            higher_is_better=True,
            precision=1,
        )
    ], use_iframe=True, height=90)

    df_m = pd.DataFrame([
        {"Tipo": "Programado (total)", "Metros (m)": float(mp)},
        {"Tipo": "Real Día ☀️", "Metros (m)": float(st.session_state.drill_day.get("metros_real_dia", 0.0))},
        {"Tipo": "Real Noche 🌙", "Metros (m)": float(st.session_state.drill_day.get("metros_real_noche", 0.0))},
        {"Tipo": "Real Total", "Metros (m)": mr_total},
    ])
    fig_m = px.bar(
        df_m,
        x="Tipo",
        y="Metros (m)",
        text_auto=True,
        color="Tipo",
        color_discrete_map={
            "Programado (total)": "#6B7280",  # gris (plan)
            "Real Día ☀️": "#F59E0B",  # ámbar (día)
            "Real Noche 🌙": "#1D4ED8",  # azul (noche)
            "Real Total": "#22C55E",  # verde (total)
        },
    )
    fig_m.update_layout(
        title="Metros perforados — Programado vs Real (☀️/🌙)",
        margin=dict(l=10, r=10, t=40, b=10),
        height=340,
    )
    st.plotly_chart(fig_m, use_container_width=True, key=f"bar_metros_{etapa}_{fecha_key}")


    # == == == == == == == == == == == == =
    # Tendencia por fecha (Metros)
    # == == == == == == == == == == == == =
    st.markdown("### Tendencia por fecha")

    _prog_map = etapa_data_rop.get("metros_prog_by_date") or {}
    _rd_map = etapa_data_rop.get("metros_real_dia_by_date") or {}
    _rn_map = etapa_data_rop.get("metros_real_noche_by_date") or {}

    _dates = sorted({*list(_prog_map.keys()), *list(_rd_map.keys()), *list(_rn_map.keys())})
    trend_rows = []
    if _dates:
        for _d in _dates:
            _p_entry = _prog_map.get(_d, {})
            _prog = _safe_float(_p_entry.get("metros_prog") if isinstance(_p_entry, dict) else (_p_entry or 0.0))
            _rd = _safe_float(_rd_map.get(_d, 0.0) or 0.0)
            _rn = _safe_float(_rn_map.get(_d, 0.0) or 0.0)
            _rt = _rd + _rn
            _eff = clamp_0_100(safe_pct(_rt, _prog)) if _prog > 0 else 0.0
            _sem = "🟢" if _eff >= 85 else ("🟡" if _eff >= 70 else "🔴")
            trend_rows.append({
                "Fecha": _d,
                "Programado": round(_prog, 2),
                "Real Día ☀️": round(_rd, 2),
                "Real Noche 🌙": round(_rn, 2),
                "Real Total": round(_rt, 2),
                "Eficiencia_pct": round(_eff, 2),
                "Semáforo": _sem,
            })

        df_tr = pd.DataFrame(trend_rows)

        fig_tr = go.Figure()
        fig_tr.add_bar(x=df_tr["Fecha"], y=df_tr["Real Día ☀️"], name="Real Día ☀️", marker_color="rgba(245,158,11,0.90)")
        fig_tr.add_bar(x=df_tr["Fecha"], y=df_tr["Real Noche 🌙"], name="Real Noche 🌙", marker_color="rgba(29,78,216,0.85)")
        fig_tr.add_trace(go.Scatter(x=df_tr["Fecha"], y=df_tr["Programado"], mode="lines+markers", name="Programado", line=dict(color="rgba(96,165,250,0.95)", width=3), marker=dict(size=8)))
        fig_tr.add_trace(go.Scatter(x=df_tr["Fecha"], y=df_tr["Real Total"], mode="lines+markers", name="Real Total", line=dict(color="rgba(249,115,22,0.95)", width=3), marker=dict(size=8)))

        fig_tr.update_layout(
            barmode="stack",
            height=420,
            margin=dict(l=10, r=10, t=20, b=10),
            legend=dict(orientation="v", x=1.02, y=1.0),
            xaxis_title="Fecha",
            yaxis_title="Metros (m)",
        )
        st.plotly_chart(fig_tr, use_container_width=True, key=f"trend_metros_{etapa}")

        st.dataframe(df_tr, use_container_width=True, hide_index=True)
    else:
        st.info("Aún no hay suficientes registros diarios para mostrar la tendencia.")


with tab_detalle:

    st.markdown("## Edición manual (TNPI/TNP/TP)")
    st.caption("Puedes editar tipo, horas y causas. Al guardar, se actualizan las gráficas automáticamente.")

    df_det = st.session_state.get("df", pd.DataFrame()).copy()
    df_det = _ensure_rowid(df_det)

    if df_det.empty:
        st.info("No hay registros para editar.")
    else:
        with st.expander("Editar registros en tabla (guardar cambios)", expanded=True):
            editable_cols = [
                "RowID",
                "Fecha",
                "Etapa",
                "Operacion",
                "Actividad",
                "Turno",
                "Corrida",
                "Tipo",
                "Categoria_TNPI",
                "Detalle_TNPI",
                "Categoria_TNP",
                "Detalle_TNP",
                "Horas_Prog",
                "Horas_Reales",
                "Comentario",
            ]
            show_cols = [c for c in editable_cols if c in df_det.columns]
            if "Eliminar" not in df_det.columns:
                df_det["Eliminar"] = False
            show_cols = ["Eliminar"] + show_cols

            # Opciones de catálogos
            cat_tnpi_opts = ["-"]
            det_tnpi_opts = ["-"]
            if "df_tnpi_cat" in globals():
                if "Categoria_TNPI" in df_tnpi_cat.columns:
                    cat_tnpi_opts = sorted(df_tnpi_cat["Categoria_TNPI"].dropna().unique().tolist())
                if "Detalle_TNPI" in df_tnpi_cat.columns:
                    det_tnpi_opts = sorted(df_tnpi_cat["Detalle_TNPI"].dropna().unique().tolist())

            cat_tnp_opts = ["-"]
            det_tnp_opts = ["-"]
            if "df_tnp_cat" in globals():
                if "Categoria_TNP" in df_tnp_cat.columns:
                    cat_tnp_opts = sorted(df_tnp_cat["Categoria_TNP"].dropna().unique().tolist())
                if "Detalle_TNP" in df_tnp_cat.columns:
                    det_tnp_opts = sorted(df_tnp_cat["Detalle_TNP"].dropna().unique().tolist())

            actividades_opts = sorted(list(set(ACTIVIDADES + ACTIVIDADES_CE + st.session_state.get("custom_actividades", []))))
            etapas_opts = sorted(list(set(SECCIONES_DEFAULT + df_det.get("Etapa", pd.Series(dtype=str)).dropna().astype(str).tolist())))
            corridas_opts = sorted(list(set(df_det.get("Corrida", pd.Series(dtype=str)).dropna().astype(str).tolist())))

            edited = st.data_editor(
                df_det[show_cols],
                use_container_width=True,
                hide_index=True,
                num_rows="dynamic",
                column_config={
                    "Eliminar": st.column_config.CheckboxColumn("Eliminar", help="Marca para borrar el registro"),
                    "Tipo": st.column_config.SelectboxColumn("Tipo", options=["TP", "TNPI", "TNP"]),
                    "Operacion": st.column_config.SelectboxColumn("Operación", options=["Perforación", "Superficie", "TR", "Otra"]),
                    "Actividad": st.column_config.SelectboxColumn("Actividad", options=actividades_opts),
                    "Turno": st.column_config.SelectboxColumn("Turno", options=TURNOS),
                    "Etapa": st.column_config.SelectboxColumn("Etapa", options=etapas_opts),
                    "Corrida": st.column_config.SelectboxColumn("Corrida", options=corridas_opts) if corridas_opts else st.column_config.TextColumn("Corrida"),
                    "Categoria_TNPI": st.column_config.SelectboxColumn("Categoría TNPI", options=cat_tnpi_opts),
                    "Detalle_TNPI": st.column_config.SelectboxColumn("Detalle TNPI", options=det_tnpi_opts),
                    "Categoria_TNP": st.column_config.SelectboxColumn("Categoría TNP", options=cat_tnp_opts),
                    "Detalle_TNP": st.column_config.SelectboxColumn("Detalle TNP", options=det_tnp_opts),
                    "Horas_Prog": st.column_config.NumberColumn("Horas Prog", min_value=0.0, step=0.25, format="%.2f"),
                    "Horas_Reales": st.column_config.NumberColumn("Horas Reales", min_value=0.0, step=0.25, format="%.2f"),
                },
                key="detalle_editor_df_" + str(st.session_state.get("_detalle_editor_version", 0)),
            )

            if st.button("Guardar cambios (Detalle)", use_container_width=True):
                ed = edited.copy()

                # Eliminar filas marcadas
                if "Eliminar" in ed.columns:
                    ed = ed[~ed["Eliminar"].astype(bool)].copy()
                if "Eliminar" in ed.columns:
                    ed.drop(columns=["Eliminar"], inplace=True, errors="ignore")

                # Asegurar RowID en nuevos registros
                if "RowID" in ed.columns:
                    ed["RowID"] = ed["RowID"].astype(str)
                    missing = ed["RowID"].isna() | (ed["RowID"].astype(str).str.strip() == "")
                    if missing.any():
                        ed.loc[missing, "RowID"] = [str(uuid.uuid4()) for _ in range(int(missing.sum()))]

                if "Horas_Prog" in ed.columns:
                    ed["Horas_Prog"] = pd.to_numeric(ed["Horas_Prog"], errors="coerce").fillna(0.0)
                if "Horas_Reales" in ed.columns:
                    ed["Horas_Reales"] = pd.to_numeric(ed["Horas_Reales"], errors="coerce").fillna(0.0)

                # Limpieza de categorías según tipo
                if "Tipo" in ed.columns:
                    mask_not_tnpi = ed["Tipo"].astype(str).str.upper() != "TNPI"
                    mask_not_tnp = ed["Tipo"].astype(str).str.upper() != "TNP"
                    if "Categoria_TNPI" in ed.columns:
                        ed.loc[mask_not_tnpi, "Categoria_TNPI"] = "-"
                    if "Detalle_TNPI" in ed.columns:
                        ed.loc[mask_not_tnpi, "Detalle_TNPI"] = "-"
                    if "Categoria_TNP" in ed.columns:
                        ed.loc[mask_not_tnp, "Categoria_TNP"] = "-"
                    if "Detalle_TNP" in ed.columns:
                        ed.loc[mask_not_tnp, "Detalle_TNP"] = "-"

                # Merge seguro por RowID
                master = st.session_state.get("df", pd.DataFrame()).copy()
                master = _ensure_rowid(master)
                master = master.set_index("RowID")
                ed2 = ed.set_index("RowID")
                common = [c for c in ed2.columns if c in master.columns]
                master.update(ed2[common])

                # Agregar nuevos registros (RowID no existente)
                new_rows = ed2.loc[~ed2.index.isin(master.index)]
                if not new_rows.empty:
                    master = pd.concat([master, new_rows[common]], axis=0)

                master = master.reset_index()
                st.session_state.df = _ensure_rowid(master)
                # Sincroniza BHA con lo que queda en actividades
                try:
                    st.session_state.df_bha = _sync_bha_from_df(st.session_state.df, st.session_state.df_bha)
                except Exception:
                    pass
                # Sincroniza Conexiones con lo que queda en actividades
                try:
                    st.session_state.df_conn = _sync_conn_from_df(st.session_state.df, st.session_state.df_conn)
                except Exception:
                    pass
                # Sincroniza viajes (limpia store si ya no existen en actividades)
                try:
                    if "viajes_hourly_store" in st.session_state and "VIAJE_TIPO" in st.session_state.df.columns:
                        valid = set(
                            st.session_state.df["VIAJE_TIPO"]
                            .dropna()
                            .astype(str)
                            .str.strip()
                            .tolist()
                        )
                        valid = {v for v in valid if v}
                        if valid:
                            def _viaje_key_tipo(k):
                                return (str(k).split("|")[0].strip() if "|" in str(k) else str(k).strip())
                            st.session_state["viajes_hourly_store"] = {
                                k: v for k, v in st.session_state["viajes_hourly_store"].items()
                                if _viaje_key_tipo(k) in valid
                            }
                        else:
                            st.session_state["viajes_hourly_store"] = {}
                except Exception:
                    pass
                # Limpia caches para que estadísticas/figuras se recalculen
                try:
                    _make_figs.clear()
                except Exception:
                    pass
                # Invalidar exportables para evitar desalineación con cambios
                for k in [
                    "exp_main_sig", "exp_main_pdf", "exp_main_ppt",
                    "exp_day_sig", "exp_day_pdf", "exp_day_ppt", "exp_day_csv",
                ]:
                    try:
                        st.session_state.pop(k, None)
                    except Exception:
                        pass
                st.session_state["_detalle_editor_version"] = st.session_state.get("_detalle_editor_version", 0) + 1
                st.success("Cambios guardados. Las gráficas y pestañas (p. ej. Conexiones) se actualizaron.")
                st.rerun()

    st.subheader("Detalle de actividades")
    # Eficiencia por fila (si hay estándar): Horas_Prog / Horas_Reales
    df_disp = df.copy()
    if "Horas_Prog" in df_disp.columns and "Horas_Reales" in df_disp.columns:
        hr = pd.to_numeric(df_disp["Horas_Reales"], errors="coerce").fillna(0.0)
        hp = pd.to_numeric(df_disp["Horas_Prog"], errors="coerce").fillna(0.0)
        df_disp["Eficiencia_pct"] = np.where(hr > 0, (hp / hr) * 100.0, 0.0)
        df_disp["Eficiencia_pct"] = df_disp["Eficiencia_pct"].clip(lower=0, upper=100)
    df_disp = _coalesce_duplicate_columns(df_disp)
    st.dataframe(add_semaforo_column(df_disp), use_container_width=True, height=340)

    if modo_reporte == "Perforación":
        st.subheader("Detalle de conexiones")
        dfc = df_conn.copy()
        if "Minutos_Estandar" in dfc.columns and "Minutos_Reales" in dfc.columns:
            mr = pd.to_numeric(dfc["Minutos_Reales"], errors="coerce").fillna(0.0)
            ms = pd.to_numeric(dfc["Minutos_Estandar"], errors="coerce").fillna(0.0)
            dfc["Eficiencia_pct"] = np.where(mr > 0, (ms / mr) * 100.0, 0.0)
            dfc["Eficiencia_pct"] = dfc["Eficiencia_pct"].clip(lower=0, upper=100)
        st.dataframe(add_semaforo_column(dfc), use_container_width=True, height=280)

    if not df_bha.empty:
        st.subheader("Detalle BHA")
        st.dataframe(add_semaforo_column(df_bha), use_container_width=True, height=280)

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: COMPARATIVA DE ETAPAS
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
with tab_comp:
    st.subheader("Comparativa de Etapas (Pro)")

    if df.empty or "Etapa" not in df.columns:
        st.info("No hay datos suficientes para comparar etapas.")
    else:
        # Estilo neutro para chips/tags del multiselect (evita rojo)
        st.markdown(
            """
            <style>
            div[data-baseweb="tag"]{
                background-color: rgba(255,255,255,0.10) !important;
                border: 1px solid rgba(255,255,255,0.18) !important;
            }
            div[data-baseweb="tag"] span{
                color: rgba(255,255,255,0.90) !important;
            }
            div[data-baseweb="tag"] svg{
                fill: rgba(255,255,255,0.70) !important;
            }
            </style>
            """,
            unsafe_allow_html=True,
        )

        with st.expander("Comparativo por días (calendario)", expanded=False):
            days_all = _available_days(df)
            if len(days_all) < 1:
                st.info("No hay fechas disponibles para comparar.")
            else:
                col_d1, col_d2 = st.columns(2)
                with col_d1:
                    day_a = st.date_input("Día A", value=days_all[-1], min_value=days_all[0], max_value=days_all[-1], key="cmp_day_a")
                with col_d2:
                    day_b_default = days_all[-2] if len(days_all) >= 2 else days_all[-1]
                    day_b = st.date_input("Día B", value=day_b_default, min_value=days_all[0], max_value=days_all[-1], key="cmp_day_b")

                df_a = split_day(df, day_a, date_col="Fecha")
                df_b = split_day(df, day_b, date_col="Fecha")

                def _kpis_day(dfin: pd.DataFrame) -> dict:
                    total = float(dfin.get("Horas_Reales", pd.Series(dtype=float)).fillna(0).sum()) if not dfin.empty else 0.0
                    tp = float(dfin[dfin.get("Tipo", "") == "TP"]["Horas_Reales"].sum()) if "Tipo" in dfin.columns else total
                    tnpi = float(dfin[dfin.get("Tipo", "") == "TNPI"]["Horas_Reales"].sum()) if "Tipo" in dfin.columns else 0.0
                    tnp = float(dfin[dfin.get("Tipo", "") == "TNP"]["Horas_Reales"].sum()) if "Tipo" in dfin.columns else 0.0
                    eff = clamp_0_100(safe_pct(tp, total)) if total > 0 else 0.0
                    return {"TP": tp, "TNPI": tnpi, "TNP": tnp, "Total": total, "Eficiencia": eff}

                k_a = _kpis_day(df_a)
                k_b = _kpis_day(df_b)

                c1, c2, c3, c4, c5 = st.columns(5)
                c1.metric(f"Total A ({day_a})", f"{k_a['Total']:.1f} h")
                c2.metric(f"TP A", f"{k_a['TP']:.1f} h")
                c3.metric(f"TNPI A", f"{k_a['TNPI']:.1f} h")
                c4.metric(f"TNP A", f"{k_a['TNP']:.1f} h")
                c5.metric(f"Eficiencia A", f"{k_a['Eficiencia']:.0f}%")

                c1b, c2b, c3b, c4b, c5b = st.columns(5)
                c1b.metric(f"Total B ({day_b})", f"{k_b['Total']:.1f} h")
                c2b.metric("TP B", f"{k_b['TP']:.1f} h")
                c3b.metric("TNPI B", f"{k_b['TNPI']:.1f} h")
                c4b.metric("TNP B", f"{k_b['TNP']:.1f} h")
                c5b.metric("Eficiencia B", f"{k_b['Eficiencia']:.0f}%")

                render_chip_row([
                    {"label": f"Día A {day_a}", "value": f"{k_a['Total']:.1f} h", "tone": "blue"},
                    {"label": "TP A", "value": f"{k_a['TP']:.1f} h", "tone": "green"},
                    {"label": "TNPI A", "value": f"{k_a['TNPI']:.1f} h", "tone": "amber"},
                    {"label": "TNP A", "value": f"{k_a['TNP']:.1f} h", "tone": "red"},
                    {"label": "Eficiencia A", "value": f"{k_a['Eficiencia']:.0f}%", "tone": "blue"},
                ], use_iframe=True, height=110)

                render_chip_row([
                    {"label": f"Día B {day_b}", "value": f"{k_b['Total']:.1f} h", "tone": "blue"},
                    {"label": "TP B", "value": f"{k_b['TP']:.1f} h", "tone": "green"},
                    {"label": "TNPI B", "value": f"{k_b['TNPI']:.1f} h", "tone": "amber"},
                    {"label": "TNP B", "value": f"{k_b['TNP']:.1f} h", "tone": "red"},
                    {"label": "Eficiencia B", "value": f"{k_b['Eficiencia']:.0f}%", "tone": "blue"},
                ], use_iframe=True, height=110)

                if show_charts:
                    df_cmp_days = pd.DataFrame(
                        [
                            {"Día": str(day_a), "Total": k_a["Total"], "TP": k_a["TP"], "TNPI": k_a["TNPI"], "TNP": k_a["TNP"]},
                            {"Día": str(day_b), "Total": k_b["Total"], "TP": k_b["TP"], "TNPI": k_b["TNPI"], "TNP": k_b["TNP"]},
                        ]
                    )
                    fig_days = px.bar(
                        df_cmp_days,
                        x="Día",
                        y=["TP", "TNPI", "TNP"],
                        barmode="stack",
                        title="Comparativo de tiempos (Día A vs Día B)",
                    )
                    st.plotly_chart(fig_days, use_container_width=True)

        st.markdown("## Comparativa de Etapas")
        col_a, col_b = st.columns(2)
        etapas_all = sorted([e for e in df["Etapa"].dropna().unique().tolist() if str(e).strip() != ""])
        with col_a:
            etapa_a = st.selectbox("Etapa A", options=etapas_all, index=0, key="cmp_etapa_a")
        with col_b:
            etapa_b = st.selectbox("Etapa B", options=etapas_all, index=1 if len(etapas_all) > 1 else 0, key="cmp_etapa_b")

        def _kpis_etapa(etp: str) -> dict:
            dfx = df[df["Etapa"] == etp].copy()
            total = float(dfx.get("Horas_Reales", pd.Series(dtype=float)).fillna(0).sum()) if not dfx.empty else 0.0
            tp = float(dfx[dfx.get("Tipo", "") == "TP"]["Horas_Reales"].sum()) if "Tipo" in dfx.columns else total
            tnpi = float(dfx[dfx.get("Tipo", "") == "TNPI"]["Horas_Reales"].sum()) if "Tipo" in dfx.columns else 0.0
            tnp = float(dfx[dfx.get("Tipo", "") == "TNP"]["Horas_Reales"].sum()) if "Tipo" in dfx.columns else 0.0
            eff = clamp_0_100(safe_pct(tp, total)) if total > 0 else 0.0
            return {"TP": tp, "TNPI": tnpi, "TNP": tnp, "Total": total, "Eficiencia": eff}

        k_a = _kpis_etapa(etapa_a)
        k_b = _kpis_etapa(etapa_b)

        st.markdown("### Comparativa de Eficiencia por Etapa")
        if show_charts:
            col_m1, col_m2, col_m3 = st.columns([1, 1, 2])
            delta_ab = k_a["Eficiencia"] - k_b["Eficiencia"]
            meta_obj = 85.0
            with col_m1:
                st.metric("Eficiencia A", f"{k_a['Eficiencia']:.1f}%", delta=f"{delta_ab:+.1f}% vs B")
            with col_m2:
                st.metric("Eficiencia B", f"{k_b['Eficiencia']:.1f}%")
            with col_m3:
                st.metric("Meta objetivo", f"{meta_obj:.0f}%")

            df_eff2 = pd.DataFrame(
                [
                    {"Etiqueta": f"Etapa A ({etapa_a})", "Eficiencia": k_a["Eficiencia"], "Etapa": etapa_a},
                    {"Etiqueta": f"Etapa B ({etapa_b})", "Eficiencia": k_b["Eficiencia"], "Etapa": etapa_b},
                ]
            )
            fig_eff2 = px.bar(
                df_eff2,
                x="Eficiencia",
                y="Etiqueta",
                orientation="h",
                text_auto=True,
                title="Eficiencia (%)",
            )
            fig_eff2.update_traces(marker_color=["#22c55e", "#3b82f6"], textposition="outside")
            fig_eff2.update_layout(
                xaxis_title="Eficiencia (%)",
                yaxis_title="",
                xaxis_range=[0, 100],
                margin=dict(l=10, r=10, t=40, b=10),
            )
            fig_eff2.add_vline(x=meta_obj, line_dash="dash", line_color="#16a34a", annotation_text="Meta 85%", annotation_position="top")
            fig_eff2.add_vline(x=70, line_dash="dot", line_color="#f59e0b", annotation_text="Alerta 70%", annotation_position="top")
            st.plotly_chart(fig_eff2, use_container_width=True)

        render_chip_row([
            {"label": f"Etapa A {etapa_a}", "value": f"{k_a['Total']:.1f} h", "tone": "blue"},
            {"label": "TP A", "value": f"{k_a['TP']:.1f} h", "tone": "green"},
            {"label": "TNPI A", "value": f"{k_a['TNPI']:.1f} h", "tone": "amber"},
            {"label": "TNP A", "value": f"{k_a['TNP']:.1f} h", "tone": "red"},
            {"label": "Eficiencia A", "value": f"{k_a['Eficiencia']:.0f}%", "tone": "blue"},
        ], use_iframe=True, height=110)

        render_chip_row([
            {"label": f"Etapa B {etapa_b}", "value": f"{k_b['Total']:.1f} h", "tone": "blue"},
            {"label": "TP B", "value": f"{k_b['TP']:.1f} h", "tone": "green"},
            {"label": "TNPI B", "value": f"{k_b['TNPI']:.1f} h", "tone": "amber"},
            {"label": "TNP B", "value": f"{k_b['TNP']:.1f} h", "tone": "red"},
            {"label": "Eficiencia B", "value": f"{k_b['Eficiencia']:.0f}%", "tone": "blue"},
        ], use_iframe=True, height=110)

        st.divider()
        st.markdown("### Detalle A vs B")
        def _semaforo_emoji(v):
            v = clamp_0_100(v)
            return "🟢" if v >= 85 else ("🟡" if v >= 70 else "🔴")

        df_cmp_etapas = pd.DataFrame(
            [
                {
                    "Etapa": etapa_a,
                    "Horas Totales": k_a["Total"],
                    "TP (h)": k_a["TP"],
                    "TNPI (h)": k_a["TNPI"],
                    "TNP (h)": k_a["TNP"],
                    "Eficiencia %": k_a["Eficiencia"],
                    "Semáforo": _semaforo_emoji(k_a["Eficiencia"]),
                    "Gap vs Meta": k_a["Eficiencia"] - meta_obj,
                },
                {
                    "Etapa": etapa_b,
                    "Horas Totales": k_b["Total"],
                    "TP (h)": k_b["TP"],
                    "TNPI (h)": k_b["TNPI"],
                    "TNP (h)": k_b["TNP"],
                    "Eficiencia %": k_b["Eficiencia"],
                    "Semáforo": _semaforo_emoji(k_b["Eficiencia"]),
                    "Gap vs Meta": k_b["Eficiencia"] - meta_obj,
                },
            ]
        )
        df_cmp_etapas["Gap vs Meta"] = df_cmp_etapas["Gap vs Meta"].map(lambda v: f"{v:+.1f}%")
        st.dataframe(df_cmp_etapas, use_container_width=True, hide_index=True)

        st.divider()
        # --- Comparativo multi-etapas (radar/heatmap + TNP + resumen) ---
        etapas_all = sorted([e for e in df["Etapa"].dropna().unique().tolist() if str(e).strip() != ""])
        etapas_sel = st.multiselect(
            "Etapas a comparar",
            options=etapas_all,
            default=etapas_all[:2] if len(etapas_all) >= 2 else etapas_all,
            key="cmp_etapas_sel",
        )
        if not etapas_sel:
            st.info("Selecciona al menos una etapa.")
        else:
            df_cmp = df[df["Etapa"].isin(etapas_sel)].copy()

            g = (
                df_cmp.groupby(["Etapa", "Tipo"], dropna=False)["Horas_Reales"]
                .sum()
                .reset_index()
            )
            piv = (
                g.pivot_table(index="Etapa", columns="Tipo", values="Horas_Reales", fill_value=0.0)
                .reset_index()
            )
            for col in ["TP", "TNPI", "TNP"]:
                if col not in piv.columns:
                    piv[col] = 0.0

            piv["Total_h"] = piv["TP"] + piv["TNPI"] + piv["TNP"]
            piv["Eficiencia_pct"] = piv.apply(
                lambda r: clamp_0_100(safe_pct(r["TP"], r["Total_h"])) if r["Total_h"] > 0 else 0.0,
                axis=1,
            )
            piv["Semáforo"] = piv["Eficiencia_pct"].apply(semaforo_dot)

            # Conexiones por etapa
            conn_map = {}
            if not df_conn.empty and "Etapa" in df_conn.columns:
                dfc = df_conn[df_conn["Etapa"].isin(etapas_sel)].copy()
                conn_map = dfc.groupby("Etapa")["Conn_No"].nunique().to_dict()
            piv["Conexiones"] = piv["Etapa"].map(lambda e: int(conn_map.get(e, 0)))

            # Normalización 0-100
            def _norm_series(s: pd.Series) -> pd.Series:
                try:
                    s = pd.to_numeric(s, errors="coerce").fillna(0.0)
                    mn, mx = float(s.min()), float(s.max())
                    if mx == mn:
                        return s.apply(lambda v: 100.0 if v > 0 else 0.0)
                    return (s - mn) / (mx - mn) * 100.0
                except Exception:
                    return pd.Series([0.0] * len(s))

            radar_cols = ["Total_h", "TP", "TNPI", "TNP", "Eficiencia_pct", "Conexiones"]
            radar_labels = ["Horas Totales", "TP (h)", "TNPI (h)", "TNP (h)", "Eficiencia %", "Conexiones"]
            norm_df = piv.copy()
            for c in radar_cols:
                norm_df[c] = _norm_series(norm_df[c])

            if show_charts:
                st.markdown("### Radar comparativo (normalizado 0–100)")
                fig_r = go.Figure()
                palette = px.colors.qualitative.Vivid
                for _, r in norm_df.iterrows():
                    color = palette[int(_ % len(palette))]
                    fig_r.add_trace(
                        go.Scatterpolar(
                            r=[float(r[c]) for c in radar_cols],
                            theta=radar_labels,
                            fill="toself",
                            name=str(r["Etapa"]),
                            opacity=0.35,
                            line=dict(color=color, width=2),
                            fillcolor=color,
                        )
                    )
                fig_r.update_layout(
                    polar=dict(radialaxis=dict(visible=True, range=[0, 100])),
                    showlegend=True,
                    margin=dict(l=20, r=20, t=20, b=20),
                    height=380,
                )
                st.plotly_chart(fig_r, use_container_width=True)

                st.markdown("### Comparativo normalizado (heatmap 0–100)")
                hm = norm_df.set_index("Etapa")[radar_cols]
                hm.columns = radar_labels
                fig_hm = px.imshow(
                    hm,
                    color_continuous_scale="Turbo",
                    range_color=[0, 100],
                    aspect="auto",
                )
                fig_hm.update_layout(margin=dict(l=20, r=20, t=30, b=20), height=320)
                st.plotly_chart(fig_hm, use_container_width=True)

            st.markdown("### Análisis de TNP (comparativo)")
            df_tnp = df_cmp[df_cmp["Tipo"] == "TNP"].copy()
            if df_tnp.empty:
                st.info("No hay registros TNP en las etapas seleccionadas.")
            else:
                for c, fb in [("Categoria_TNP", "Sin categoría"), ("Detalle_TNP", "Sin detalle")]:
                    if c not in df_tnp.columns:
                        df_tnp[c] = fb
                    df_tnp[c] = df_tnp[c].fillna(fb).replace({"-": fb, "None": fb, "nan": fb})

                tnp_et = (
                    df_tnp.groupby("Etapa", as_index=False)["Horas_Reales"]
                    .sum()
                    .sort_values("Horas_Reales", ascending=False)
                )
                if show_charts:
                    fig_tnp = px.bar(
                        tnp_et,
                        x="Etapa",
                        y="Horas_Reales",
                        title="TNP por etapa (h)",
                        text_auto=True,
                    )
                    fig_tnp.update_traces(marker_color="#f59e0b")
                    st.plotly_chart(fig_tnp, use_container_width=True)

                tnp_det = (
                    df_tnp.groupby(["Etapa", "Categoria_TNP", "Detalle_TNP"], as_index=False)["Horas_Reales"]
                    .sum()
                    .sort_values("Horas_Reales", ascending=False)
                )
                st.dataframe(tnp_det, use_container_width=True, hide_index=True)

            st.markdown("### Resumen comparativo")
            show_cols = ["Etapa", "Total_h", "TP", "TNPI", "TNP", "Eficiencia_pct", "Conexiones", "Semáforo"]
            st.dataframe(piv[show_cols].sort_values("Etapa"), use_container_width=True, hide_index=True)

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: ESTADÍSTICAS POR ETAPA
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: ESTADÍSTICAS POR ETAPA
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: ESTADÍSTICAS CAMBIO DE ETAPA (CE)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
with tab_ce:
    st.markdown("### 🔁 ")
    # --- CE metrics safety defaults ---
    horas_total_ce = 0.0
    tp_ce = 0.0
    tnpi_ce = 0.0
    tnp_ce = 0.0
    eficiencia_ce = 0.0

    st.markdown("### Estadísticas - Cambio de etapa (CE)")
    df_all = st.session_state.df.copy()

    if df_all.empty:
        st.info("Aún no hay actividades registradas. Agrega actividades y vuelve aquí para ver estadísticas.")
    else:
        # Filtrar CE
        if "Modo_Reporte" in df_all.columns:
            df_ce = df_all[df_all["Modo_Reporte"].astype(str) == "Cambio de etapa"].copy()
        else:
            df_ce = df_all.copy()

        if df_ce.empty:
            st.warning("No hay actividades registradas con Modo de reporte = 'Cambio de etapa'.")
            st.caption("Tip: cambia el 'Modo reporte' en el panel lateral antes de agregar actividades de CE.")
        else:
            # Normalizar columnas mínimas
            if "Fecha" in df_ce.columns:
                df_ce["Fecha"] = pd.to_datetime(df_ce["Fecha"], errors="coerce").dt.date
            else:
                df_ce["Fecha"] = pd.NaT

            df_ce["Horas_Reales"] = pd.to_numeric(df_ce.get("Horas_Reales", 0), errors="coerce").fillna(0.0)

            # Filtros
            c1, c2, c3 = st.columns([1, 1, 1])
            with c1:
                fechas = [d for d in df_ce["Fecha"].dropna().unique().tolist() if d]
                if fechas:
                    fmin, fmax = min(fechas), max(fechas)
                else:
                    fmin = fmax = datetime.today().date()
                rango = st.date_input("Rango de fechas", value=(fmin, fmax), key="ce_rango")
                if isinstance(rango, tuple) and len(rango) == 2:
                    f_ini, f_fin = rango
                else:
                    f_ini, f_fin = fmin, fmax
            with c2:
                etapas = sorted([str(x) for x in df_ce.get("Etapa", pd.Series(dtype=str)).fillna("").unique().tolist() if str(x).strip() != ""])
                etapa_f = st.selectbox("Etapa", options=["(Todas)"] + etapas, index=0, key="ce_etapa")
            with c3:
                mostrar_detalle = st.toggle("Ver tabla detalle", value=False, key="ce_det_toggle")

            if "Fecha" in df_ce.columns and f_ini and f_fin:
                df_ce = df_ce[(df_ce["Fecha"] >= f_ini) & (df_ce["Fecha"] <= f_fin)].copy()
            if etapa_f != "(Todas)" and "Etapa" in df_ce.columns:
                df_ce = df_ce[df_ce["Etapa"].astype(str) == str(etapa_f)].copy()

            if df_ce.empty:
                st.warning("No hay datos CE para ese filtro.")
                st.stop()

            # KPIs
            total_h = float(df_ce["Horas_Reales"].sum())
            tp_h = float(df_ce[df_ce.get("Tipo", "") == "TP"]["Horas_Reales"].sum()) if "Tipo" in df_ce.columns else total_h
            tnpi_h = float(df_ce[df_ce.get("Tipo", "") == "TNPI"]["Horas_Reales"].sum()) if "Tipo" in df_ce.columns else 0.0
            tnp_h = float(df_ce[df_ce.get("Tipo", "") == "TNP"]["Horas_Reales"].sum()) if "Tipo" in df_ce.columns else 0.0
            eff = (tp_h / total_h * 100.0) if total_h > 0 else 0.0

            # Semáforo (ajustable)
            warn_below = 75.0
            crit_below = 60.0
            if eff >= warn_below:
                tone = "green"
                sem_txt = "OK"
            elif eff >= crit_below:
                tone = "amber"
                sem_txt = "ATENCIÓN"
            else:
                tone = "red"
                sem_txt = "CRÍTICO"

            # Persistir KPIs CE globales (para otros bloques)
            horas_total_ce = total_h
            tp_ce = tp_h
            tnpi_ce = tnpi_h
            tnp_ce = tnp_h
            eficiencia_ce = eff

            # Chips pro (KPIs CE)
            render_chip_row([
                {"label": "Horas total (CE)", "value": f"{total_h:.2f} h", "tone": "blue"},
                {"label": "TP", "value": f"{tp_h:.2f} h", "tone": "green"},
                {"label": "TNPI", "value": f"{tnpi_h:.2f} h", "tone": "amber"},
                {"label": "TNP", "value": f"{tnp_h:.2f} h", "tone": "red"},
                {"label": "Eficiencia", "value": f"{eff:.0f}% · {sem_txt}", "tone": tone},
            ], use_iframe=True, height=120)

            # Chips adicionales (pro)
            n_act = int(df_ce["Actividad"].nunique()) if "Actividad" in df_ce.columns else 0
            n_days = int(df_ce["Fecha"].nunique()) if "Fecha" in df_ce.columns else 0
            avg_day = (total_h / n_days) if n_days > 0 else 0.0
            render_chip_row([
                {"label": "Días con CE", "value": f"{n_days}", "tone": "gray"},
                {"label": "Actividades", "value": f"{n_act}", "tone": "violet"},
                {"label": "Promedio diario", "value": f"{avg_day:.2f} h", "tone": "blue"},
            ], use_iframe=True, height=110)

            st.divider()

            # Gráficas principales
            g1, g2 = st.columns([1, 1])
            with g1:
                if "Tipo" in df_ce.columns:
                    df_tipo = df_ce.groupby("Tipo", as_index=False)["Horas_Reales"].sum()
                    if not df_tipo.empty:
                        fig = px.pie(df_tipo, names="Tipo", values="Horas_Reales", hole=0.55, title="Distribución TP / TNPI / TNP (CE)")
                        st.plotly_chart(fig, use_container_width=True)
            with g2:
                if "Actividad" in df_ce.columns:
                    df_a = df_ce.groupby("Actividad", as_index=False)["Horas_Reales"].sum().sort_values("Horas_Reales", ascending=False).head(15)
                    if not df_a.empty:
                        palette = px.colors.qualitative.Set3 + px.colors.qualitative.Pastel + px.colors.qualitative.Bold
                        act_names = df_a["Actividad"].tolist()
                        act_color_map = {a: palette[i % len(palette)] for i, a in enumerate(act_names)}
                        fig = px.bar(
                            df_a,
                            x="Actividad",
                            y="Horas_Reales",
                            color="Actividad",
                            title="Top actividades por horas (CE)",
                            color_discrete_map=act_color_map,
                        )
                        fig.update_layout(xaxis_title="", yaxis_title="Horas", xaxis_tickangle=-35, showlegend=False)
                        st.plotly_chart(fig, use_container_width=True)

            # Tendencia por fecha / hora
            if "Fecha" in df_ce.columns and df_ce["Fecha"].notna().any():
                st.markdown("### Tendencia (CE)")
                has_time = ("Hora_Inicio" in df_ce.columns) and df_ce["Hora_Inicio"].astype(str).str.strip().ne("").any()
                tendencia_mode = st.radio(
                    "Vista",
                    ["Por día", "Por hora"],
                    index=0,
                    horizontal=True,
                    key="ce_tendencia_mode",
                )
                if tendencia_mode == "Por hora" and not has_time:
                    st.info("No hay horas registradas. Activa 'Registrar hora' al capturar CE.")
                    tendencia_mode = "Por día"

                if tendencia_mode == "Por hora":
                    df_tmp = df_ce.copy()
                    df_tmp["_Hora"] = pd.to_datetime(df_tmp["Hora_Inicio"], format="%H:%M", errors="coerce").dt.hour
                    df_tmp = df_tmp.dropna(subset=["_Hora"])
                    if df_tmp.empty:
                        st.info("No hay horas válidas para generar tendencia por hora.")
                    else:
                        g = df_tmp.groupby(["_Hora", "Tipo"], as_index=False)["Horas_Reales"].sum()
                        piv = g.pivot_table(index="_Hora", columns="Tipo", values="Horas_Reales", fill_value=0.0).reset_index()
                        for c in ["TP", "TNPI", "TNP"]:
                            if c not in piv.columns:
                                piv[c] = 0.0
                        piv["Total_h"] = piv["TP"] + piv["TNPI"] + piv["TNP"]
                        piv = piv.sort_values("_Hora")
                        piv["Eficiencia_pct"] = piv.apply(lambda r: (r["TP"] / r["Total_h"] * 100.0) if r["Total_h"] > 0 else 0.0, axis=1)
                        piv["Semáforo"] = piv["Eficiencia_pct"].apply(semaforo_dot)

                        # Chips pro
                        best_row = piv.sort_values("Eficiencia_pct", ascending=False).iloc[0]
                        worst_row = piv.sort_values("Eficiencia_pct", ascending=True).iloc[0]
                        avg_eff = float(piv["Eficiencia_pct"].mean()) if len(piv) > 0 else 0.0
                        render_chip_row([
                            {"label": "Mejor hora", "value": f"{int(best_row['_Hora']):02d}:00 · {best_row['Eficiencia_pct']:.0f}%", "tone": "green"},
                            {"label": "Peor hora", "value": f"{int(worst_row['_Hora']):02d}:00 · {worst_row['Eficiencia_pct']:.0f}%", "tone": "red"},
                            {"label": "Eficiencia promedio", "value": f"{avg_eff:.0f}%", "tone": "blue"},
                            {"label": "Horas con CE", "value": f"{len(piv)}", "tone": "gray"},
                        ], use_iframe=True, height=120)

                        df_long = piv.melt(
                            id_vars=["_Hora"],
                            value_vars=["Total_h", "TP", "TNPI", "TNP"],
                            var_name="Serie",
                            value_name="Horas",
                        )
                        fig = px.line(
                            df_long,
                            x="_Hora",
                            y="Horas",
                            color="Serie",
                            markers=True,
                            title="Tendencia por hora (CE)",
                        )
                        fig.update_layout(xaxis_title="Hora", yaxis_title="Horas")
                        fig.update_xaxes(dtick=1)
                        st.plotly_chart(fig, use_container_width=True)

                        st.markdown("#### Semáforo por hora (CE)")
                        st.dataframe(
                            piv[["_Hora", "Total_h", "TP", "TNPI", "TNP", "Eficiencia_pct", "Semáforo"]],
                            use_container_width=True,
                            hide_index=True,
                        )
                else:
                    df_d = df_ce.groupby("Fecha", as_index=False).agg(
                        Total_h=("Horas_Reales", "sum"),
                        TP_h=("Horas_Reales", lambda s: float(s[df_ce.loc[s.index, "Tipo"].astype(str) == "TP"].sum()) if "Tipo" in df_ce.columns else float(s.sum())),
                        TNPI_h=("Horas_Reales", lambda s: float(s[df_ce.loc[s.index, "Tipo"].astype(str) == "TNPI"].sum()) if "Tipo" in df_ce.columns else 0.0),
                        TNP_h=("Horas_Reales", lambda s: float(s[df_ce.loc[s.index, "Tipo"].astype(str) == "TNP"].sum()) if "Tipo" in df_ce.columns else 0.0),
                    ).sort_values("Fecha")

                    df_d["Eficiencia_pct"] = df_d.apply(lambda r: (r["TP_h"]/r["Total_h"]*100.0) if r["Total_h"]>0 else 0.0, axis=1)
                    df_d["Semáforo"] = df_d["Eficiencia_pct"].apply(semaforo_dot)
                    df_d["Fecha"] = pd.to_datetime(df_d["Fecha"], errors="coerce")

                    # Chips pro arriba de la tendencia
                    if not df_d.empty:
                        best_row = df_d.sort_values("Eficiencia_pct", ascending=False).iloc[0]
                        worst_row = df_d.sort_values("Eficiencia_pct", ascending=True).iloc[0]
                        avg_eff = float(df_d["Eficiencia_pct"].mean()) if len(df_d) > 0 else 0.0
                        best_day = best_row["Fecha"].date().isoformat() if pd.notna(best_row["Fecha"]) else "-"
                        worst_day = worst_row["Fecha"].date().isoformat() if pd.notna(worst_row["Fecha"]) else "-"
                        render_chip_row([
                            {"label": "Mejor día", "value": f"{best_day} · {best_row['Eficiencia_pct']:.0f}%", "tone": "green"},
                            {"label": "Peor día", "value": f"{worst_day} · {worst_row['Eficiencia_pct']:.0f}%", "tone": "red"},
                            {"label": "Eficiencia promedio", "value": f"{avg_eff:.0f}%", "tone": "blue"},
                            {"label": "Días", "value": f"{len(df_d)}", "tone": "gray"},
                        ], use_iframe=True, height=120)

                    df_long = df_d.melt(
                        id_vars=["Fecha"],
                        value_vars=["Total_h", "TP_h", "TNPI_h", "TNP_h"],
                        var_name="Serie",
                        value_name="Horas",
                    )
                    fig = px.line(
                        df_long,
                        x="Fecha",
                        y="Horas",
                        color="Serie",
                        markers=True,
                        title="Tendencia por fecha (CE)",
                    )
                    fig.update_layout(xaxis_title="", yaxis_title="Horas")
                    fig.update_xaxes(dtick="D1", tickformat="%Y-%m-%d")
                    st.plotly_chart(fig, use_container_width=True)

                    # Semáforos por fecha (tabla + chips)
                    st.markdown("#### Semáforo por fecha (CE)")
                    st.dataframe(
                        df_d[["Fecha", "Total_h", "TP_h", "TNPI_h", "TNP_h", "Eficiencia_pct", "Semáforo"]],
                        use_container_width=True,
                        hide_index=True,
                    )

            # Tabla resumen por actividad
            if "Actividad" in df_ce.columns:
                if "Tipo" in df_ce.columns:
                    piv = df_ce.pivot_table(index="Actividad", columns="Tipo", values="Horas_Reales", aggfunc="sum", fill_value=0.0)
                    for c in ["TP","TNPI","TNP"]:
                        if c not in piv.columns:
                            piv[c]=0.0
                    piv["Total"] = piv[["TP","TNPI","TNP"]].sum(axis=1)
                    piv["Eficiencia_%"] = piv.apply(lambda r: (r["TP"]/r["Total"]*100.0) if r["Total"]>0 else 0.0, axis=1)
                    piv["Semáforo"] = piv["Eficiencia_%"].apply(semaforo_dot)
                    piv = piv.sort_values("Total", ascending=False).reset_index()
                else:
                    piv = df_ce.groupby("Actividad", as_index=False)["Horas_Reales"].sum().rename(columns={"Horas_Reales":"Total"})
                    piv["TP"]=piv["Total"]; piv["TNPI"]=0.0; piv["TNP"]=0.0; piv["Eficiencia_%"]=100.0
                    piv["Semáforo"] = piv["Eficiencia_%"].apply(semaforo_dot)

                st.markdown("#### Resumen por actividad (CE)")
                st.dataframe(piv, use_container_width=True, hide_index=True)

            if mostrar_detalle:
                st.markdown("#### Detalle (CE)")
                st.dataframe(
                    _decorate_turno_df(df_ce.sort_values(["Fecha", "Turno"], ascending=[True, True])),
                    use_container_width=True,
                    hide_index=True
                )

            st.caption("Recomendación: usa CE para capturar tiempos de transición (cambio de herramienta/etapa, cementación, WOC, etc.). Esto permite separar desempeño de perforación vs tiempos de cambio de etapa.")


with tab_estadisticas:
    st.subheader("📊 Estadísticas por Etapa")
    
    # Selector de modo: Etapa actual vs Todas las etapas
    col_modo1, col_modo2 = st.columns([1, 3])
    
    with col_modo1:
        modo_estadisticas = st.radio(
            "Modo de análisis",
            options=["Etapa actual", "Todas las etapas"],
            horizontal=True,
            key="modo_estadisticas"
        )
    
    if modo_estadisticas == "Etapa actual":
        # Filtro para seleccionar etapa
        etapas_disponibles = df["Etapa"].unique().tolist() if not df.empty else []
        if not etapas_disponibles:
            st.info("No hay datos disponibles. Por favor, captura algunas actividades primero.")
        else:
            etapa_seleccionada = st.selectbox("Seleccionar etapa para análisis", etapas_disponibles)
            
            # Filtrar datos por etapa
            df_etapa = df[df["Etapa"] == etapa_seleccionada].copy()
            df_conn_etapa = df_conn[df_conn["Seccion"] == etapa_seleccionada].copy()
            df_bha_etapa = df_bha[df_bha["Etapa"] == etapa_seleccionada].copy()
            
            # ---- SECCIÓN 1: KPIs PRINCIPALES ----
            st.markdown("### 📈 KPIs Principales")
            
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                tp_h_etapa = float(df_etapa[df_etapa["Tipo"] == "TP"]["Horas_Reales"].sum()) if not df_etapa.empty else 0.0
                st.metric("TP (h)", f"{tp_h_etapa:.1f}")
            
            with col2:
                tnpi_h_etapa = float(df_etapa[df_etapa["Tipo"] == "TNPI"]["Horas_Reales"].sum()) if not df_etapa.empty else 0.0
                st.metric("TNPI (h)", f"{tnpi_h_etapa:.1f}")
            
            with col3:
                tnp_h_etapa = float(df_etapa[df_etapa["Tipo"] == "TNP"]["Horas_Reales"].sum()) if not df_etapa.empty else 0.0
                st.metric("TNP (h)", f"{tnp_h_etapa:.1f}")
            
            with col4:
                total_h_etapa = float(df_etapa["Horas_Reales"].sum()) if not df_etapa.empty else 0.0
                eficiencia_etapa = clamp_0_100(safe_pct(tp_h_etapa, total_h_etapa)) if total_h_etapa > 0 else 0.0
                sk, sl, sc = status_from_eff(eficiencia_etapa)
                st.markdown(f"""
                    <div style="text-align:center">
                        <div style="font-size:24px;font-weight:bold;color:{sc}">{eficiencia_etapa:.0f}%</div>
                        <div style="font-size:12px;color:#888">Eficiencia</div>
                    </div>
                """, unsafe_allow_html=True)
            
            # ---- SECCIÓN 2: GRÁFICAS ----
            st.markdown("### 📊 Distribuciones")
            
            col_chart1, col_chart2 = st.columns(2)
            
            with col_chart1:
                # Distribución de tiempos
                if not df_etapa.empty:
                    df_tiempos = df_etapa.groupby("Tipo")["Horas_Reales"].sum().reset_index()
                    fig_tiempos = px.pie(df_tiempos, names="Tipo", values="Horas_Reales", 
                                        title="Distribución de Tiempos (%)", hole=0.4,
                                        color="Tipo", color_discrete_map={"TP": "#2ECC71", "TNPI": "#E74C3C", "TNP": "#F1C40F"})
                    fig_tiempos.update_traces(textposition='inside', textinfo='percent+label')
                    st.plotly_chart(fig_tiempos, use_container_width=True)
                else:
                    st.info("No hay datos de tiempos")
            
            with col_chart2:
                # Distribución de operaciones
                if not df_etapa.empty:
                    df_operaciones = df_etapa.groupby("Operacion")["Horas_Reales"].sum().reset_index()
                    df_operaciones = df_operaciones.sort_values("Horas_Reales", ascending=False).head(5)
                    fig_operaciones = px.bar(df_operaciones, x="Operacion", y="Horas_Reales",
                                            title="Top 5 - Operaciones (h)", text_auto=True,
                                            color="Horas_Reales", color_continuous_scale="Viridis")
                    st.plotly_chart(fig_operaciones, use_container_width=True)
                else:
                    st.info("No hay datos de operaciones")
            
            # ---- SECCIÓN 3: TABLAS DETALLADAS ----
            st.markdown("### 📋 Detalles Específicos")
            
            # Inicializar variables fuera de los tabs
            df_conn_summary = pd.DataFrame()
            conexiones_count = 0
            
            tab1, tab2, tab3 = st.tabs(["📊 Metros y ROP", "🔧 BHA", "🔗 Conexiones"])
            
            with tab1:
                # Metros perforados y ROP
                if modo_reporte == "Perforación":
                    # Usar datos por etapa (no globales) para que Programado/Real correspondan a la etapa seleccionada
                    etapa_data = get_etapa_data(etapa_seleccionada)

                    # Metros programados por etapa: usamos PT programada (m)
                    mp_etapa = float(etapa_data.get("pt_programada_m", 0.0) or 0.0)

                    # Metros reales por etapa: acumulado de metros diarios capturados (día + noche)
                    _mr_d_map = etapa_data.get("metros_real_dia_by_date", {}) or {}
                    _mr_n_map = etapa_data.get("metros_real_noche_by_date", {}) or {}
                    mr_total_calc = float(sum(_mr_d_map.values()) + sum(_mr_n_map.values()))
                    if legacy_calc_value == 0.0:
                        legacy_calc_value = float(
                            (etapa_data.get("metros_real_dia", 0.0) or 0.0)
                            + (etapa_data.get("metros_real_noche", 0.0) or 0.0)
                        )
                    mr_etapa = mr_total_calc if mr_total_calc > 0 else float(legacy_calc_value or 0.0)

                    # ROP programada por etapa (meta)
                    rp_etapa = float(etapa_data.get("rop_prog_etapa", 0.0) or 0.0)

                    # ROP real promedio por etapa: promedio simple de los ROP diarios capturados (manual)
                    _rop_d_map = etapa_data.get("rop_real_dia_by_date", {}) or {}
                    _rop_n_map = etapa_data.get("rop_real_noche_by_date", {}) or {}
                    _rop_vals = [float(v) for v in list(_rop_d_map.values()) + list(_rop_n_map.values()) if float(v) > 0]
                    if _rop_vals:
                        rr_etapa = float(sum(_rop_vals) / len(_rop_vals))
                    else:
                        _tmp = []
                        if float(etapa_data.get("rop_real_dia", 0.0) or 0.0) > 0:
                            _tmp.append(float(etapa_data.get("rop_real_dia", 0.0) or 0.0))
                        if float(etapa_data.get("rop_real_noche", 0.0) or 0.0) > 0:
                            _tmp.append(float(etapa_data.get("rop_real_noche", 0.0) or 0.0))
                        rr_etapa = float(sum(_tmp) / len(_tmp)) if _tmp else 0.0

                    eficiencia_metros = (mr_etapa / mp_etapa * 100) if mp_etapa > 0 else 0.0
                    eficiencia_rop = (rr_etapa / rp_etapa * 100) if rp_etapa > 0 else 0.0

                    df_metros = pd.DataFrame({
                        "Concepto": ["Programado", "Real", "Eficiencia"],
                        "Metros (m)": [mp_etapa, mr_etapa, eficiencia_metros],
                        "ROP (m/h)": [rp_etapa, rr_etapa, eficiencia_rop],
                    })

                    # Semáforos SOLO en la fila de eficiencia
                    df_metros["Semáforo Metros"] = ""
                    df_metros["Semáforo ROP"] = ""
                    df_metros.loc[df_metros["Concepto"] == "Eficiencia", "Semáforo Metros"] = semaforo_dot(eficiencia_metros)
                    df_metros.loc[df_metros["Concepto"] == "Eficiencia", "Semáforo ROP"] = semaforo_dot(eficiencia_rop)

                    st.dataframe(df_metros, use_container_width=True, hide_index=True)
                else:
                    st.info("Esta sección aplica solo para el modo Perforación.")
            with tab2:
                # BHA
                if not df_bha_etapa.empty:
                    df_bha_display = df_bha_etapa.copy()
                    df_bha_display["Eficiencia_pct"] = df_bha_display.apply(
                        lambda r: (r["Estandar_h"] / r["Real_h"] * 100) if r["Real_h"] > 0 else 0,
                        axis=1
                    )
                    df_bha_display["Semáforo"] = df_bha_display["Eficiencia_pct"].apply(semaforo_dot)
                    
                    # Gráfica de BHA
                    fig_bha_etapa = px.bar(df_bha_display, x="BHA_Tipo", y=["Estandar_h", "Real_h"],
                                          title="BHA: Estándar vs Real por Tipo", barmode="group",
                                          labels={"value": "Horas", "variable": "Tipo"})
                    st.plotly_chart(fig_bha_etapa, use_container_width=True)
                    
                    st.dataframe(df_bha_display[["BHA_Tipo", "BHA_Componentes", "Accion", "Estandar_h", "Real_h", "TNPI_h", "Eficiencia_pct", "Semáforo"]], 
                               use_container_width=True, hide_index=True)
                else:
                    st.info("No hay datos BHA para esta etapa")
            
            with tab3:
                # Conexiones
                if not df_conn_etapa.empty:
                    # Resumen por conexión
                    df_conn_summary = df_conn_etapa.groupby("Conn_No").agg({
                        "Minutos_Reales": "sum",
                        "Minutos_TNPI": "sum"
                    }).reset_index()
                    df_conn_summary["TP_min"] = df_conn_summary["Minutos_Reales"] - df_conn_summary["Minutos_TNPI"]
                    df_conn_summary["Eficiencia_pct"] = df_conn_summary.apply(
                        lambda r: (r["TP_min"] / r["Minutos_Reales"] * 100) if r["Minutos_Reales"] > 0 else 0,
                        axis=1
                    )
                    df_conn_summary["Semáforo"] = df_conn_summary["Eficiencia_pct"].apply(semaforo_dot)
                    conexiones_count = len(df_conn_summary)
                    
                    # Gráfica de conexiones
                    fig_conn_etapa = px.bar(df_conn_summary, x="Conn_No", y=["TP_min", "Minutos_TNPI"],
                                           title="Conexiones: TP vs TNPI", barmode="stack",
                                           labels={"value": "Minutos", "variable": "Tipo"})
                    st.plotly_chart(fig_conn_etapa, use_container_width=True)
                    
                    st.dataframe(df_conn_summary[["Conn_No", "Minutos_Reales", "TP_min", "Minutos_TNPI", "Eficiencia_pct", "Semáforo"]],
                               use_container_width=True, hide_index=True)
                else:
                    st.info("No hay datos de conexiones para esta etapa")
                    conexiones_count = 0
            
            # ---- SECCIÓN 4: ANÁLISIS TNPI ----
            st.markdown("### 🔍 Análisis de TNPI")
            
            if tnpi_h_etapa > 0:
                # Top causas de TNPI
                df_tnpi_causas = df_etapa[df_etapa["Tipo"] == "TNPI"].groupby(["Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP", "Detalle_TNP"])["Horas_Reales"].sum().reset_index()
                df_tnpi_causas = df_tnpi_causas.sort_values("Horas_Reales", ascending=False).head(10)
                
                col_causas1, col_causas2 = st.columns(2)
                
                with col_causas1:
                    # Gráfica de causas
                    if not df_tnpi_causas.empty:
                        fig_causas = px.bar(df_tnpi_causas, x="Detalle_TNPI", y="Horas_Reales",
                                           title="Top 10 - Causas de TNPI (h)",
                                           color="Horas_Reales", color_continuous_scale="Reds")
                        fig_causas.update_layout(xaxis_tickangle=45)
                        st.plotly_chart(fig_causas, use_container_width=True)
                
                with col_causas2:
                    # Tabla de causas
                    if not df_tnpi_causas.empty:
                        st.dataframe(df_tnpi_causas[["Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP", "Detalle_TNP", "Horas_Reales"]],
                                   use_container_width=True, hide_index=True)
                
                # Distribución por categoría
                df_tnpi_cat = df_etapa[df_etapa["Tipo"] == "TNPI"].groupby("Categoria_TNPI")["Horas_Reales"].sum().reset_index()
                if not df_tnpi_cat.empty:
                    fig_tnpi_cat = px.pie(df_tnpi_cat, names="Categoria_TNPI", values="Horas_Reales",
                                         title="TNPI por Categoría (%)", hole=0.3)
                    st.plotly_chart(fig_tnpi_cat, use_container_width=True)
            else:
                st.success("🎉 No hay TNPI registrado para esta etapa")
            

            # ---- SECCIÓN 4B: ANÁLISIS TNP ----
            st.markdown("### 🔵 Análisis de TNP")

            df_tnp_etapa = df_etapa[df_etapa["Tipo"] == "TNP"].copy()
            if not df_tnp_etapa.empty:
                # Normalizar nulos/guiones para evitar 'nan'
                for col in ["Categoria_TNP", "Detalle_TNP"]:
                    if col not in df_tnp_etapa.columns:
                        df_tnp_etapa[col] = ""
                    df_tnp_etapa[col] = (
                        df_tnp_etapa[col]
                        .astype(str)
                        .replace({"nan": "", "None": "", "-": ""})
                        .fillna("")
                        .str.strip()
                    )

                df_tnp_etapa["Categoria_TNP"] = df_tnp_etapa["Categoria_TNP"].replace({"": "Sin categoría"})
                df_tnp_etapa["Detalle_TNP"] = df_tnp_etapa["Detalle_TNP"].replace({"": "Sin detalle"})

                # Top detalles
                df_tnp_top = (
                    df_tnp_etapa.groupby(["Categoria_TNP", "Detalle_TNP"])["Horas_Reales"]
                    .sum()
                    .reset_index()
                    .sort_values("Horas_Reales", ascending=False)
                    .head(10)
                )

                col_tnp1, col_tnp2 = st.columns(2)
                with col_tnp1:
                    if not df_tnp_top.empty:
                        fig_tnp_top = px.bar(
                            df_tnp_top,
                            x="Detalle_TNP",
                            y="Horas_Reales",
                            title="Top 10 - Causas de TNP (h)",
                            color="Horas_Reales",
                            color_continuous_scale="Blues",
                        )
                        fig_tnp_top.update_layout(xaxis_tickangle=45)
                        st.plotly_chart(fig_tnp_top, use_container_width=True)

                with col_tnp2:
                    st.dataframe(
                        df_tnp_top[["Categoria_TNP", "Detalle_TNP", "Horas_Reales"]],
                        use_container_width=True,
                        hide_index=True,
                    )

                # Distribución por categoría
                df_tnp_cat = (
                    df_tnp_etapa.groupby("Categoria_TNP")["Horas_Reales"]
                    .sum()
                    .reset_index()
                    .sort_values("Horas_Reales", ascending=False)
                )
                if not df_tnp_cat.empty:
                    fig_tnp_cat = px.pie(
                        df_tnp_cat,
                        names="Categoria_TNP",
                        values="Horas_Reales",
                        title="TNP por Categoría (%)",
                        hole=0.3,
                    )
                    st.plotly_chart(fig_tnp_cat, use_container_width=True)
            else:
                st.success("🎉 No hay TNP registrado para esta etapa")

            # ---- SECCIÓN 5: RESUMEN EJECUTIVO ----
            st.markdown("### 📋 Resumen Ejecutivo")
            
            # Asegurar valores por defecto si no se calcularon en esta ruta
            rr_etapa = float(rr_etapa) if "rr_etapa" in locals() else 0.0

            # Crear resumen ejecutivo
            resumen_data = {
                "Métrica": ["Horas Totales", "TP (Horas Productivas)", "TNPI (Horas No Productivas)", 
                           "TNP (Tiempo No Productivo)", "Eficiencia General", "Metros Perforados", 
                           "ROP Promedio", "Conexiones Realizadas", "Operaciones BHA"],
                "Valor": [
                    f"{total_h_etapa:.1f} h",
                    f"{tp_h_etapa:.1f} h",
                    f"{tnpi_h_etapa:.1f} h",
                    f"{tnp_h_etapa:.1f} h",
                    f"{eficiencia_etapa:.0f}%",
                    f"{0.0:.0f} m" if modo_reporte == "Perforación" else "N/A",
                    f"{rr_etapa:.1f} m/h" if modo_reporte == "Perforación" else "N/A",
                    f"{conexiones_count}",
                    f"{len(df_bha_etapa)}" if not df_bha_etapa.empty else "0"
                ],
                "Estado": [
                    "🟢" if total_h_etapa > 0 else "⚪",
                    "🟢" if tp_h_etapa > 0 else "⚪",
                    "🟡" if 0 < tnpi_h_etapa < 5 else ("🔴" if tnpi_h_etapa >= 5 else "🟢"),
                    "🟡" if 0 < tnp_h_etapa < 3 else ("🔴" if tnp_h_etapa >= 3 else "🟢"),
                    semaforo_dot(eficiencia_etapa),
                    "🟢" if horas_total_ce > 0 else "⚪",
                    "🟢" if rr_etapa > 0 else "⚪",
                    "🟢" if conexiones_count > 0 else "⚪",
                    "🟢" if len(df_bha_etapa) > 0 else "⚪"
                ]
            }
            
            df_resumen = pd.DataFrame(resumen_data)
            st.dataframe(df_resumen, use_container_width=True, hide_index=True)
            
            # Botón para exportar reporte de etapa
            col_exp1, col_exp2 = st.columns(2)
            with col_exp1:
                if st.button("📥 Exportar Reporte de Etapa (PDF)", use_container_width=True):
                    # Aquí iría la lógica para exportar el reporte de etapa
                    st.success("Funcionalidad de exportación en desarrollo")
            
            with col_exp2:
                if st.button("📊 Generar Dashboard Ejecutivo", use_container_width=True):
                    st.success("Dashboard generado para revisión ejecutiva")
    
    else:  # Modo "Todas las etapas"
        st.info("Mostrando estadísticas consolidadas de todas las etapas")
        
        # Mostrar un mensaje y botón para ir al reporte general
        st.markdown("""
        **Para ver el reporte general completo con todas las etapas, por favor ve a la pestaña:**
        ### 📊 **"Reporte General del Pozo"**
        
        Allí encontrarás:
        - KPIs consolidados de todas las etapas
        - Gráficas de distribución general
        - Análisis de TNPI por categoría y etapa
        - Tablas resumen detalladas
        - Opciones de exportación
        """)
        
        # Botón para ir directamente al tab general
        if st.button("Ir a Reporte General del Pozo", use_container_width=True):
            # No hay forma directa de cambiar tabs en Streamlit, pero podemos usar session state
            st.session_state["active_tab"] = "Reporte General del Pozo"
            st.rerun()

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =

# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: ESTADÍSTICAS POR CORRIDA
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
with tab_corridas:
    st.subheader("Estadísticas por corrida")

    if "corridas_manual" not in st.session_state:
        st.session_state.corridas_manual = []  # lista de dicts: {Nombre, Etapa, Prof_ini, Prof_fin, Fecha_ini, Fecha_fin}

    df = st.session_state.get("df", pd.DataFrame()).copy()
    if df.empty:
        st.info("Aún no hay actividades cargadas para calcular estadísticas.")
    else:
        metodo = st.radio(
            "Método de definición de corrida",
            options=["Por Corrida (campo Corrida)", "Auto por BHA (Arma/Desarma)", "Por profundidad (manual)"],
            horizontal=True,
        )

        # -----------------------------
        # 1) Por columna Corrida
        # -----------------------------
        if metodo == "Por Corrida (campo Corrida)":
            if "Corrida" not in df.columns:
                st.warning('No existe la columna "Corrida" en la tabla de actividades.')
            else:
                corridas = [c for c in df["Corrida"].dropna().unique().tolist() if str(c).strip() != ""]
                if not corridas:
                    st.info('No hay valores en "Corrida". Puedes capturar el campo Corrida en el sidebar o usar los otros métodos.')
                else:
                    corrida_sel = st.selectbox("Selecciona corrida", options=sorted(corridas, key=lambda x: str(x)))
                    d = df[df["Corrida"] == corrida_sel].copy()

                    # KPIs
                    d["Horas_Reales"] = pd.to_numeric(d["Horas_Reales"], errors="coerce").fillna(0.0)
                    d["Horas_Prog"] = pd.to_numeric(d["Horas_Prog"], errors="coerce").fillna(0.0)
                    total_h = float(d["Horas_Reales"].sum())
                    tp_h = float(d.loc[d["Tipo"] == "TP", "Horas_Reales"].sum())
                    tnpi_h = float(d.loc[d["Tipo"] == "TNPI", "Horas_Reales"].sum())
                    tnp_h = float(d.loc[d["Tipo"] == "TNP", "Horas_Reales"].sum())
                    eff = (tp_h / total_h * 100.0) if total_h > 0 else 0.0

                    c1, c2, c3, c4, c5 = st.columns(5)
                    c1.metric("Total (h)", f"{total_h:.2f}")
                    c2.metric("TP (h)", f"{tp_h:.2f}")
                    c3.metric("TNPI (h)", f"{tnpi_h:.2f}")
                    c4.metric("TNP (h)", f"{tnp_h:.2f}")
                    c5.metric("Eficiencia (%)", f"{eff:.1f}")

                    st.markdown("### Composición de tiempos (TP/TNPI/TNP)")
                    tipo_color_map = {"TP": "#22C55E", "TNPI": "#F59E0B", "TNP": "#EF4444"}
                    try:
                        _comp = (
                            d.groupby(["Actividad", "Tipo"], dropna=False)["Horas_Reales"]
                            .sum()
                            .reset_index()
                        )
                        _comp = _comp[_comp["Horas_Reales"] > 0]
                        if _comp.empty:
                            st.info("No hay horas para graficar en la corrida seleccionada.")
                        else:
                            fig_stack = px.bar(
                                _comp,
                                x="Horas_Reales",
                                y="Actividad",
                                color="Tipo",
                                color_discrete_map=tipo_color_map,
                                orientation="h",
                                title="Horas por actividad (apilado por tipo)",
                            )
                            fig_stack.update_layout(xaxis_title="Horas", yaxis_title="Actividad", barmode="stack")
                            st.plotly_chart(fig_stack, use_container_width=True)

                            _tot = (
                                d.groupby("Tipo", dropna=False)["Horas_Reales"]
                                .sum()
                                .reset_index()
                            )
                            _tot = _tot[_tot["Horas_Reales"] > 0]
                            if not _tot.empty:
                                fig_donut = px.pie(
                                    _tot,
                                    names="Tipo",
                                    values="Horas_Reales",
                                    title="Composición total de tiempos",
                                    hole=0.35,
                                    color="Tipo",
                                    color_discrete_map=tipo_color_map,
                                )
                                st.plotly_chart(fig_donut, use_container_width=True)
                    except Exception as _e:
                        st.warning(f"No pude generar gráficas combinadas: {_e}")

                    st.markdown("### Distribución TNPI por categoría")
                    d_tnpi = d[d["Tipo"] == "TNPI"].copy()
                    for col, fb in [("Categoria_TNPI", "Sin categoría"), ("Detalle_TNPI", "Sin detalle")]:
                        if col not in d_tnpi.columns:
                            d_tnpi[col] = fb
                        d_tnpi[col] = (
                            d_tnpi[col]
                            .astype(str)
                            .replace({"-": fb, "": fb, "None": fb, "nan": fb})
                            .fillna(fb)
                        )
                    pareto = (
                        d_tnpi
                        .groupby(["Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP", "Detalle_TNP"], dropna=False)["Horas_Reales"]
                        .sum()
                        .sort_values(ascending=False)
                        .reset_index()
                    )
                    st.dataframe(pareto, use_container_width=True, hide_index=True)

                    # --- Gráficas (pro) TNPI por corrida ---
                    try:
                        if not pareto.empty:
                            _cat = (
                                pareto.groupby("Categoria_TNPI", dropna=False)["Horas_Reales"]
                                .sum()
                                .sort_values(ascending=False)
                                .reset_index()
                            )
                            _cat = _cat[_cat["Horas_Reales"] > 0]
                            if not _cat.empty:
                                total_tnpi_cat = float(_cat["Horas_Reales"].sum())
                                top_cat_name = str(_cat.iloc[0]["Categoria_TNPI"])
                                top_cat_val = float(_cat.iloc[0]["Horas_Reales"])
                                top_cat_pct = (top_cat_val / total_tnpi_cat * 100.0) if total_tnpi_cat > 0 else 0.0
                                render_chip_row(
                                    [
                                        {"label": "TNPI total", "value": f"{total_tnpi_cat:.2f} h", "tone": "orange"},
                                        {"label": "Categorías", "value": f"{len(_cat)}", "tone": "gray"},
                                        {"label": "Top categoría", "value": f"{top_cat_name} · {top_cat_pct:.0f}%", "tone": "blue"},
                                    ],
                                    use_iframe=True,
                                    height=100,
                                )
                                _top_cat = _cat.head(5)
                                fig_bar = px.bar(
                                    _top_cat.sort_values("Horas_Reales", ascending=True),
                                    x="Horas_Reales",
                                    y="Categoria_TNPI",
                                    color="Categoria_TNPI",
                                    color_discrete_sequence=px.colors.qualitative.Bold,
                                    orientation="h",
                                    title="Top 5 TNPI por categoría (h)",
                                )
                                fig_bar.update_layout(xaxis_title="Horas", yaxis_title="Categoría TNPI")
                                st.plotly_chart(fig_bar, use_container_width=True)

                                fig_pie = px.pie(
                                    _cat,
                                    names="Categoria_TNPI",
                                    values="Horas_Reales",
                                    title="Distribución TNPI por categoría",
                                    hole=0.35,
                                    color="Categoria_TNPI",
                                    color_discrete_sequence=px.colors.qualitative.Bold,
                                )
                                st.plotly_chart(fig_pie, use_container_width=True)
                    except Exception as _e:
                        st.warning(f"No pude generar gráficas por corrida: {_e}")


                st.markdown("### Distribución TNP por categoría")
                try:
                    df_tnp = d[d["Tipo"] == "TNP"].copy() if "d" in locals() else pd.DataFrame()
                    if df_tnp.empty and "d" in locals():
                        df_tnp = d[d["Tipo"] == "TNP"].copy()

                    if df_tnp.empty:
                        st.info("No hay registros TNP para la corrida seleccionada.")
                    else:
                        df_tnp["_Categoria_TNP_view"] = "-"
                        if "Categoria_TNP" in df_tnp.columns:
                            df_tnp["_Categoria_TNP_view"] = df_tnp["Categoria_TNP"].astype(str)
                        if (df_tnp["_Categoria_TNP_view"].astype(str).str.strip().eq("-").all()
                            and "Categoria_TNPI" in df_tnp.columns):
                            df_tnp["_Categoria_TNP_view"] = df_tnp["Categoria_TNPI"].astype(str)
                        df_tnp["_Categoria_TNP_view"] = (
                            df_tnp["_Categoria_TNP_view"]
                            .astype(str)
                            .replace({"-": "Sin categoría", "": "Sin categoría", "None": "Sin categoría", "nan": "Sin categoría"})
                            .fillna("Sin categoría")
                        )

                        tnp_cat = (
                            df_tnp.groupby("_Categoria_TNP_view", dropna=False)["Horas_Reales"]
                            .sum()
                            .sort_values(ascending=False)
                            .reset_index()
                        )
                        
                        tnp_cat = tnp_cat.rename(columns={"_Categoria_TNP_view": "Categoria_TNP"})
                        total_tnp_cat = float(tnp_cat["Horas_Reales"].sum()) if not tnp_cat.empty else 0.0
                        top_cat_name = str(tnp_cat.iloc[0]["Categoria_TNP"]) if not tnp_cat.empty else "-"
                        top_cat_val = float(tnp_cat.iloc[0]["Horas_Reales"]) if not tnp_cat.empty else 0.0
                        top_cat_pct = (top_cat_val / total_tnp_cat * 100.0) if total_tnp_cat > 0 else 0.0
                        render_chip_row([
                            {"label": "TNP total", "value": f"{total_tnp_cat:.2f} h", "tone": "red"},
                            {"label": "Categorías", "value": f"{len(tnp_cat)}", "tone": "gray"},
                            {"label": "Top categoría", "value": f"{top_cat_name} · {top_cat_pct:.0f}%", "tone": "blue"},
                        ], use_iframe=True, height=100)
                        st.dataframe(tnp_cat, use_container_width=True, hide_index=True)

                        # Barras: Top 5 por categoría
                        top_tnp = tnp_cat.head(5).copy()
                        fig_tnp_bar = px.bar(
                            top_tnp.sort_values("Horas_Reales", ascending=True),
                            x="Horas_Reales",
                            y="Categoria_TNP",
                            color="Categoria_TNP",
                            orientation="h",
                            title="Top 5 TNP por categoría (h)",
                            color_discrete_sequence=px.colors.qualitative.Bold,
                        )
                        fig_tnp_bar.update_layout(xaxis_title="Horas", yaxis_title="Categoría TNP")
                        st.plotly_chart(fig_tnp_bar, use_container_width=True)

                        # Donut: distribución por categoría
                        cat_tbl = tnp_cat[tnp_cat["Horas_Reales"] > 0]
                        if not cat_tbl.empty:
                            fig_tnp_pie = px.pie(
                                cat_tbl,
                                names="Categoria_TNP",
                                values="Horas_Reales",
                                title="Distribución TNP por categoría",
                                hole=0.35,
                                color="Categoria_TNP",
                                color_discrete_sequence=px.colors.qualitative.Bold,
                            )
                            st.plotly_chart(fig_tnp_pie, use_container_width=True)
                except Exception as _e:
                    st.warning(f"No pude generar gráficas TNP por corrida: {_e}")


                    st.markdown("### Detalle de actividades")
                    st.dataframe(d, use_container_width=True, hide_index=True)

        # -----------------------------
        # 2) Auto por BHA
        # -----------------------------
        elif metodo == "Auto por BHA (Arma/Desarma)":
            df_bha = st.session_state.get("df_bha", pd.DataFrame()).copy()
            if df_bha.empty or "Accion" not in df_bha.columns:
                st.info("No hay registros de BHA para calcular corridas automáticamente.")
            else:
                # Normalizar fecha
                if "Fecha" in df_bha.columns:
                    df_bha["Fecha"] = pd.to_datetime(df_bha["Fecha"], errors="coerce")
                df_bha = df_bha.sort_values(["Fecha"], na_position="last")

                # Crear RunId incremental: inicia con acciones tipo 'Arma'
                def _is_start(a):
                    a = str(a).strip().lower()
                    return a.startswith("arma") or a.startswith("arm") or a.startswith("a")

                run_id = 0
                run_ids = []
                for a in df_bha["Accion"].fillna("").tolist():
                    if _is_start(a):
                        run_id += 1
                    run_ids.append(run_id if run_id > 0 else 1)
                df_bha["RunId_Auto"] = run_ids

                run_sel = st.selectbox("Selecciona corrida (auto por BHA)", options=sorted(df_bha["RunId_Auto"].unique().tolist()))
                # Filtrar actividades que pertenecen a esa corrida usando fecha ventana de BHA
                win = df_bha[df_bha["RunId_Auto"] == run_sel].copy()
                tmin = win["Fecha"].min()
                tmax = win["Fecha"].max()

                d = df.copy()
                if "Fecha" in d.columns:
                    d["Fecha_dt"] = pd.to_datetime(d["Fecha"], errors="coerce")
                    d = d[(d["Fecha_dt"] >= tmin) & (d["Fecha_dt"] <= tmax)].copy()

                d["Horas_Reales"] = pd.to_numeric(d["Horas_Reales"], errors="coerce").fillna(0.0)
                total_h = float(d["Horas_Reales"].sum())
                tp_h = float(d.loc[d["Tipo"] == "TP", "Horas_Reales"].sum())
                tnpi_h = float(d.loc[d["Tipo"] == "TNPI", "Horas_Reales"].sum())
                eff = (tp_h / total_h * 100.0) if total_h > 0 else 0.0

                c1, c2, c3, c4 = st.columns(4)
                c1.metric("Ventana", f"{tmin.date()} → {tmax.date()}" if pd.notna(tmin) and pd.notna(tmax) else "-")
                c2.metric("Total (h)", f"{total_h:.2f}")
                c3.metric("TNPI (h)", f"{tnpi_h:.2f}")
                c4.metric("Eficiencia (%)", f"{eff:.1f}")

                st.markdown("### Eventos BHA de la corrida")
                st.dataframe(win, use_container_width=True, hide_index=True)

                st.markdown("### Actividades dentro de la ventana")
                st.dataframe(d.drop(columns=["Fecha_dt"], errors="ignore"), use_container_width=True, hide_index=True)

        # -----------------------------
        # 3) Manual por profundidad
        # -----------------------------
        else:
            st.markdown("Define corridas manuales por profundidad (y opcionalmente por fechas) para calcular estadísticas.")
            with st.expander("Crear / editar corrida (manual)", expanded=False):
                nombre = st.text_input("Nombre corrida", "")
                etapa_sel = st.selectbox("Etapa", options=sorted(df["Etapa"].dropna().unique().tolist(), key=lambda x: str(x)))
                prof_ini = st.number_input("Profundidad inicio (m)", min_value=0.0, value=0.0, step=1.0)
                prof_fin = st.number_input("Profundidad fin (m)", min_value=0.0, value=0.0, step=1.0)

                # Fecha opcional para filtrar actividades (más confiable que inferir por profundidad sin log continuo)
                usar_fechas = st.checkbox("Filtrar actividades por rango de fechas (opcional)", value=True)
                fecha_ini = st.date_input("Fecha inicio", value=pd.Timestamp.today().date())
                fecha_fin = st.date_input("Fecha fin", value=pd.Timestamp.today().date())

                if st.button("Guardar corrida manual", use_container_width=True):
                    if not nombre.strip():
                        st.warning("Escribe un nombre para la corrida.")
                    elif prof_fin <= prof_ini:
                        st.warning("Profundidad fin debe ser mayor que profundidad inicio.")
                    else:
                        st.session_state.corridas_manual.append({
                            "Nombre": nombre.strip(),
                            "Etapa": etapa_sel,
                            "Prof_ini": float(prof_ini),
                            "Prof_fin": float(prof_fin),
                            "Usar_fechas": bool(usar_fechas),
                            "Fecha_ini": str(fecha_ini),
                            "Fecha_fin": str(fecha_fin),
                        })
                        st.success("Corrida guardada.")
                        st.rerun()

            if not st.session_state.corridas_manual:
                st.info("Aún no has creado corridas manuales.")
            else:
                nombres = [c["Nombre"] for c in st.session_state.corridas_manual]
                corrida_sel = st.selectbox("Selecciona corrida manual", options=nombres)
                corrida = next(c for c in st.session_state.corridas_manual if c["Nombre"] == corrida_sel)

                d = df[df["Etapa"] == corrida["Etapa"]].copy()
                if corrida.get("Usar_fechas", True):
                    d["Fecha_dt"] = pd.to_datetime(d["Fecha"], errors="coerce")
                    d = d[
                        (d["Fecha_dt"] >= pd.to_datetime(corrida["Fecha_ini"])) &
                        (d["Fecha_dt"] <= pd.to_datetime(corrida["Fecha_fin"]) + pd.Timedelta(days=1) - pd.Timedelta(seconds=1))
                    ].copy()

                d["Horas_Reales"] = pd.to_numeric(d["Horas_Reales"], errors="coerce").fillna(0.0)
                total_h = float(d["Horas_Reales"].sum())
                tp_h = float(d.loc[d["Tipo"] == "TP", "Horas_Reales"].sum())
                tnpi_h = float(d.loc[d["Tipo"] == "TNPI", "Horas_Reales"].sum())
                eff = (tp_h / total_h * 100.0) if total_h > 0 else 0.0
                metros = float(corrida["Prof_fin"] - corrida["Prof_ini"])
                rop = (metros / tp_h) if tp_h > 0 else 0.0

                c1, c2, c3, c4 = st.columns(4)
                c1.metric("Metros (ΔProf)", f"{metros:.1f} m")
                c2.metric("TP (h)", f"{tp_h:.2f}")
                c3.metric("TNPI (h)", f"{tnpi_h:.2f}")
                c4.metric("ROP (m/h)", f"{rop:.2f}")

                st.markdown("### Distribución TNPI (por horas)")
                pareto = (
                    d[d["Tipo"] == "TNPI"]
                    .groupby(["Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP", "Detalle_TNP"], dropna=False)["Horas_Reales"]
                    .sum()
                    .sort_values(ascending=False)
                    .reset_index()
                )
                st.dataframe(pareto, use_container_width=True, hide_index=True)

                st.markdown("### Detalle de actividades (filtrado)")
                st.dataframe(d.drop(columns=["Fecha_dt"], errors="ignore"), use_container_width=True, hide_index=True)


# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: ESTADÍSTICAS DRILLSPOT (KPI EXPORT)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
with tab_drillspot:
    st.subheader("Estadísticas DrillSpot (KPI Export)")
    st.caption("Estas estadísticas se calculan aparte y no modifican tus actividades. Carga el KPI Export (CSV o XLSX) para activar la vista.")

    up = st.file_uploader("Cargar KPI Export de DrillSpot", type=["csv", "xlsx"])
    if up is None:
        st.info("Carga un archivo para ver las estadísticas DrillSpot.")
    else:
        try:
            if up.name.lower().endswith(".csv"):
                kpi_raw = pd.read_csv(up)
            else:
                # intenta hoja por defecto; si falla, lee la primera
                xls = pd.ExcelFile(up)
                sheet = "KPI Report" if "KPI Report" in xls.sheet_names else xls.sheet_names[0]
                kpi_raw = pd.read_excel(up, sheet_name=sheet)

            # Limpieza robusta: eliminar fila de unidades típica (Start Time == 'date')
            if "Start Time" in kpi_raw.columns:
                kpi = kpi_raw.copy()
                kpi = kpi[kpi["Start Time"].astype(str).str.lower() != "date"].copy()
            else:
                kpi = kpi_raw.copy()

            # Tipos
            if "Start Time" in kpi.columns:
                kpi["Start Time"] = pd.to_datetime(kpi["Start Time"], errors="coerce")
            if "End Time" in kpi.columns:
                kpi["End Time"] = pd.to_datetime(kpi["End Time"], errors="coerce")
            for c in ["Start Bit Depth", "End Bit Depth", "Duration"]:
                if c in kpi.columns:
                    kpi[c] = pd.to_numeric(kpi[c], errors="coerce")

            # Duración en horas (Duration viene en minutos en el export)
            if "Duration" in kpi.columns:
                kpi["Duration_h"] = kpi["Duration"] / 60.0

            st.success(f"Archivo cargado: {up.name} ({len(kpi):,} filas)")

            # Resumen general
            total_h = float(kpi["Duration_h"].sum()) if "Duration_h" in kpi.columns else 0.0
            prof_ini = float(kpi["Start Bit Depth"].min()) if "Start Bit Depth" in kpi.columns else 0.0
            prof_fin = float(kpi["End Bit Depth"].max()) if "End Bit Depth" in kpi.columns else 0.0
            net_m = prof_fin - prof_ini
            c1, c2, c3, c4 = st.columns(4)
            c1.metric("Total (h)", f"{total_h:.2f}")
            c2.metric("Prof. inicio (m)", f"{prof_ini:.2f}")
            c3.metric("Prof. fin (m)", f"{prof_fin:.2f}")
            c4.metric("ΔProf neta (m)", f"{net_m:.2f}")

            # Chips pro acumulativos con semáforos (ROP: verde ≥6, ámbar ≥4, rojo <4 m/h)
            rop_acum = (net_m / total_h) if total_h and total_h > 0 else 0.0
            n_events = len(kpi)
            _st = kpi.get("Start Time")
            n_days = _st.dropna().dt.date.nunique() if _st is not None and hasattr(_st.dt, "date") else 0
            _rop_tone = "green" if rop_acum >= 6 else ("amber" if rop_acum >= 4 else "red") if total_h and total_h > 0 else "gray"
            render_chip_row([
                {"label": "Total (h)", "value": f"{total_h:.2f}", "tone": "blue"},
                {"label": "ΔProf neta (m)", "value": f"{net_m:.2f}", "tone": "violet"},
                {"label": "ROP (m/h)", "value": f"{rop_acum:.2f}", "tone": _rop_tone},
                {"label": "Días cubiertos", "value": str(n_days), "tone": "gray"},
                {"label": "Eventos", "value": f"{n_events:,}", "tone": "gray"},
            ], use_iframe=True, height=100)

            # Stats por Run
            if "Run" in kpi.columns:
                st.markdown("### Resumen por Run")
                run = (
                    kpi.groupby("Run", dropna=False)
                    .agg(
                        Inicio=("Start Bit Depth", "min"),
                        Fin=("End Bit Depth", "max"),
                        Delta=("End Bit Depth", lambda s: float(s.max()) - float(kpi.loc[s.index, "Start Bit Depth"].min()) if "Start Bit Depth" in kpi.columns else 0.0),
                        Horas=("Duration_h", "sum"),
                        Eventos=("KPI", "count"),
                    )
                    .reset_index()
                )
                # Semáforo por avance (Delta): verde ≥100 m, ámbar 0–100, rojo ≤0
                def _sem_run(d):
                    if d is None or (isinstance(d, float) and pd.isna(d)):
                        return "⚪"
                    try:
                        v = float(d)
                        if v >= 100:
                            return "🟢"
                        if v > 0:
                            return "🟡"
                        return "🔴"
                    except Exception:
                        return "⚪"
                run["Semáforo"] = run["Delta"].apply(_sem_run)
                run = run[["Run", "Semáforo", "Inicio", "Fin", "Delta", "Horas", "Eventos"]]
                st.dataframe(run, use_container_width=True, hide_index=True)
                try:
                    r = run.copy()
                    try:
                        r["Run_num"] = pd.to_numeric(r["Run"], errors="coerce")
                        r = r.sort_values("Run_num")
                    except Exception:
                        pass

                    fig_run = px.bar(
                        r,
                        x="Run",
                        y="Horas",
                        title="Horas por Run",
                    )
                    fig_run.update_layout(
                        xaxis_title="Run",
                        yaxis_title="Horas",
                        legend_title_text="",
                    )
                    st.plotly_chart(fig_run, use_container_width=True)
                except Exception as _e:
                    st.warning(f"No pude generar gráfica por Run: {_e}")
            # Pareto KPI
            if "KPI" in kpi.columns and "Duration_h" in kpi.columns:
                st.markdown("### Pareto (KPI por horas)")
                pareto = (
                    kpi.groupby("KPI", dropna=False)["Duration_h"]
                    .sum()
                    .sort_values(ascending=False)
                    .reset_index()
                )
                total_pareto = pareto["Duration_h"].sum()
                pareto["%"] = (pareto["Duration_h"] / total_pareto * 100).round(1) if total_pareto else 0
                st.dataframe(pareto, use_container_width=True, hide_index=True)
                # ------------------------------
                # Gráficas (Plotly)
                # ------------------------------
                try:
                    # Barras: Top KPIs por horas
                    top = pareto.head(12).copy()
                    fig_bar = px.bar(
                        top.sort_values("Duration_h", ascending=True),
                        x="Duration_h",
                        y="KPI",
                        orientation="h",
                        title="Top KPIs por horas",
                    )
                    fig_bar.update_layout(xaxis_title="Horas", yaxis_title="KPI")
                    st.plotly_chart(fig_bar, use_container_width=True)

                    # Pastel: distribución (Top 8 + Otros)
                    pie_df = pareto.head(8).copy()
                    otros_h = float(pareto["Duration_h"].sum() - pie_df["Duration_h"].sum())
                    if otros_h > 0:
                        pie_df = pd.concat(
                            [pie_df, pd.DataFrame([{"KPI": "Otros", "Duration_h": otros_h}])],
                            ignore_index=True,
                        )

                    fig_pie = px.pie(
                        pie_df,
                        names="KPI",
                        values="Duration_h",
                        title="Distribución de horas por KPI",
                        hole=0.35,
                    )
                    st.plotly_chart(fig_pie, use_container_width=True)
                except Exception as _e:
                    st.warning(f"No pude generar gráficas KPI: {_e}")

            # ------------------------------ Estadísticas discretizadas por día
            st.markdown("### Estadísticas por día")
            _st_col = kpi.get("Start Time")
            if _st_col is not None and _st_col.notna().any():
                fechas_disp = sorted(_st_col.dropna().dt.date.unique().tolist())
                if fechas_disp:
                    sel_fecha = st.date_input(
                        "Seleccionar día",
                        value=min(fechas_disp) if fechas_disp else date.today(),
                        min_value=min(fechas_disp),
                        max_value=max(fechas_disp),
                        key="drillspot_dia",
                    )
                    if sel_fecha in fechas_disp:
                        kpi_dia = kpi[_st_col.dt.date == sel_fecha].copy()
                        total_h_d = float(kpi_dia["Duration_h"].sum()) if "Duration_h" in kpi_dia.columns else 0.0
                        prof_ini_d = float(kpi_dia["Start Bit Depth"].min()) if "Start Bit Depth" in kpi_dia.columns else 0.0
                        prof_fin_d = float(kpi_dia["End Bit Depth"].max()) if "End Bit Depth" in kpi_dia.columns else 0.0
                        net_m_d = prof_fin_d - prof_ini_d
                        rop_d = (net_m_d / total_h_d) if total_h_d and total_h_d > 0 else 0.0
                        n_ev_d = len(kpi_dia)
                        _rop_tone_d = "green" if rop_d >= 6 else ("amber" if rop_d >= 4 else "red") if total_h_d else "gray"
                        render_chip_row([
                            {"label": "Total (h)", "value": f"{total_h_d:.2f}", "tone": "blue"},
                            {"label": "ΔProf (m)", "value": f"{net_m_d:.2f}", "tone": "violet"},
                            {"label": "ROP (m/h)", "value": f"{rop_d:.2f}", "tone": _rop_tone_d},
                            {"label": "Eventos", "value": f"{n_ev_d:,}", "tone": "gray"},
                        ], use_iframe=True, height=100)
                        # Gráfica pro: horas del día (0–23) por KPI (barras apiladas) + línea de profundidad
                        kpi_dia["hora"] = kpi_dia["Start Time"].dt.hour
                        por_hora_kpi = kpi_dia.groupby(["hora", "KPI"], dropna=False)["Duration_h"].sum().reset_index()
                        pivot_h = por_hora_kpi.pivot(index="hora", columns="KPI", values="Duration_h").fillna(0)
                        pivot_h = pivot_h.reindex(range(24), fill_value=0).fillna(0)
                        fig_dia = go.Figure()
                        colors = px.colors.qualitative.Set2 + px.colors.qualitative.Pastel
                        for i, col in enumerate(pivot_h.columns):
                            fig_dia.add_trace(go.Bar(
                                x=pivot_h.index.tolist(),
                                y=pivot_h[col].tolist(),
                                name=str(col) if col else "N/A",
                                marker_color=colors[i % len(colors)],
                            ))
                        has_depth = "End Bit Depth" in kpi_dia.columns and kpi_dia["End Bit Depth"].notna().any()
                        added_depth_line = False
                        if has_depth and "End Time" in kpi_dia.columns:
                            et = pd.to_datetime(kpi_dia["End Time"], errors="coerce")
                            depth_at_h = []
                            for h in range(24):
                                t_end = pd.Timestamp(sel_fecha) + pd.Timedelta(hours=h + 1)
                                mask = et <= t_end
                                val = kpi_dia.loc[mask, "End Bit Depth"].max() if mask.any() else None
                                depth_at_h.append(None if val is None or (isinstance(val, float) and pd.isna(val)) else float(val))
                            valid = [x is not None for x in depth_at_h]
                            if any(valid):
                                y2 = [d if v else None for d, v in zip(depth_at_h, valid)]
                                fig_dia.add_trace(
                                    go.Scatter(
                                        x=list(range(24)),
                                        y=y2,
                                        name="Profundidad (m)",
                                        line=dict(color="#22c55e", width=2),
                                        yaxis="y2",
                                    )
                                )
                                added_depth_line = True
                        layout_kw = dict(
                            barmode="stack",
                            title=f"Actividad por hora del día — {sel_fecha}",
                            xaxis_title="Hora del día",
                            yaxis_title="Horas",
                            xaxis=dict(dtick=1, range=[-0.5, 23.5]),
                            legend=dict(orientation="h", yanchor="bottom", y=1.02),
                            height=420,
                        )
                        if added_depth_line:
                            layout_kw["yaxis2"] = dict(title="Profundidad (m)", overlaying="y", side="right", showgrid=False)
                        fig_dia.update_layout(**layout_kw)
                        st.plotly_chart(fig_dia, use_container_width=True)

                        # Segunda gráfica: actividad discretizada por evento (no acumulada), minuto dentro de la hora
                        st.markdown("#### Timeline por evento (minutos dentro de cada hora)")
                        st.caption("Cada barra es un evento en su minuto de inicio; la altura es la duración. No se acumulan.")
                        kpi_dia_st = kpi_dia["Start Time"].dropna()
                        kpi_dia_et = kpi_dia["End Time"].dropna() if "End Time" in kpi_dia.columns else None
                        kpi_dia_dur = kpi_dia["Duration"] if "Duration" in kpi_dia.columns else (kpi_dia["Duration_h"] * 60.0)
                        kpi_dia_kpi = kpi_dia["KPI"] if "KPI" in kpi_dia.columns else pd.Series(["Evento"] * len(kpi_dia))
                        segments = []
                        for i in kpi_dia.index:
                            start = kpi_dia_st.get(i)
                            if start is None or (hasattr(start, "tz") and pd.isna(start)):
                                continue
                            dur_min = float(kpi_dia_dur.iloc[kpi_dia.index.get_loc(i)] or 0)
                            if dur_min <= 0:
                                continue
                            k = kpi_dia_kpi.iloc[kpi_dia.index.get_loc(i)]
                            if pd.isna(k):
                                k = "N/A"
                            start = pd.Timestamp(start)
                            hour0 = start.hour
                            remaining = dur_min
                            min0 = start.minute + start.second / 60.0 + start.microsecond / 60000000.0
                            while remaining > 1e-6 and hour0 < 24:
                                seg_dur = min(remaining, 60.0 - min0)
                                if seg_dur > 0:
                                    segments.append((hour0, min0, seg_dur, str(k)))
                                remaining -= seg_dur
                                hour0 += 1
                                min0 = 0
                        if segments:
                            seg_df = pd.DataFrame(segments, columns=["hora", "min_inicio", "dur_min", "KPI"])
                            kpi_orden = seg_df["KPI"].unique().tolist()
                            colors_d = {k: px.colors.qualitative.Set2[i % len(px.colors.qualitative.Set2)] for i, k in enumerate(kpi_orden)}
                            if len(kpi_orden) > len(px.colors.qualitative.Set2):
                                extra = px.colors.qualitative.Pastel
                                for i, k in enumerate(kpi_orden[len(px.colors.qualitative.Set2):]):
                                    colors_d[k] = extra[i % len(extra)]
                            fig_tl = go.Figure()
                            for k in kpi_orden:
                                sub = seg_df[seg_df["KPI"] == k]
                                fig_tl.add_trace(go.Bar(
                                    x=sub["hora"].tolist(),
                                    y=sub["dur_min"].tolist(),
                                    base=sub["min_inicio"].tolist(),
                                    name=str(k),
                                    marker_color=colors_d.get(k, "gray"),
                                    width=0.85,
                                ))
                            layout_tl = dict(
                                barmode="overlay",
                                title=f"Actividad por evento (minuto dentro de la hora) — {sel_fecha}",
                                xaxis_title="Hora del día",
                                yaxis_title="Minutos dentro de la hora",
                                xaxis=dict(dtick=1, range=[-0.5, 23.5]),
                                yaxis=dict(range=[0, 60], dtick=5),
                                legend=dict(orientation="h", yanchor="bottom", y=1.02),
                                height=420,
                            )
                            if added_depth_line:
                                fig_tl.add_trace(
                                    go.Scatter(
                                        x=list(range(24)),
                                        y=y2,
                                        name="Profundidad (m)",
                                        line=dict(color="#22c55e", width=2),
                                        yaxis="y2",
                                    )
                                )
                                layout_tl["yaxis2"] = dict(title="Profundidad (m)", overlaying="y", side="right", showgrid=False)
                            fig_tl.update_layout(**layout_tl)
                            st.plotly_chart(fig_tl, use_container_width=True)
                        else:
                            st.info("No hay eventos con hora de inicio para mostrar el timeline por evento.")
                    else:
                        st.info("No hay datos para la fecha seleccionada.")
                else:
                    st.caption("No se detectaron fechas en el archivo para estadísticas por día.")
            else:
                st.caption("Se requiere columna 'Start Time' con fechas para estadísticas por día.")

            st.markdown("### Datos crudos (preview)")
            st.dataframe(kpi.head(200), use_container_width=True, hide_index=True)

        except Exception as e:
            st.error(f"No pude leer el archivo. Error: {e}")

    # ------------------------------ Rig Activities (timeline por evento, mismo estilo que KPI)
    st.markdown("---")
    st.subheader("Rig Activities")
    st.caption("Carga un CSV de actividades del equipo: columna de timestamp (ej. YYYY-MM-DDTHH:MM:SS), «Rig Activity», «Bit depth (m)», «Hole Depth (m)». Se muestran actividades completas en timeline por evento.")
    up_rig = st.file_uploader("Cargar CSV Rig Activities", type=["csv"], key="rig_activities_upload")
    if up_rig is not None:
        try:
            rig_raw = pd.read_csv(up_rig)
            # Detectar columna de tiempo (primera columna o la que contenga 'time'/'date' o formato ISO)
            time_col = None
            for c in rig_raw.columns:
                c_low = str(c).strip().lower()
                if c_low in ("time", "timestamp", "date", "datetime") or "time" in c_low or "date" in c_low:
                    time_col = c
                    break
            if time_col is None and len(rig_raw.columns) > 0:
                sample = str(rig_raw.iloc[0, 0]) if len(rig_raw) > 0 else ""
                if "T" in sample and "-" in sample and ":" in sample:
                    time_col = rig_raw.columns[0]
            if time_col is None and len(rig_raw.columns) > 0:
                time_col = rig_raw.columns[0]
            # Actividad y profundidad
            activity_col = None
            for c in rig_raw.columns:
                if "rig" in str(c).lower() and "activ" in str(c).lower():
                    activity_col = c
                    break
            if activity_col is None:
                for c in rig_raw.columns:
                    if "activ" in str(c).lower():
                        activity_col = c
                        break
            if activity_col is None:
                st.error("No se encontró columna de actividad (ej. «Rig Activity»).")
            else:
                rig = rig_raw.copy()
                rig[time_col] = pd.to_datetime(rig[time_col], errors="coerce")
                rig = rig.dropna(subset=[time_col]).sort_values(time_col).reset_index(drop=True)
                if rig.empty:
                    st.warning("No hay filas con timestamp válido.")
                else:
                    # Segmentos: agrupar filas consecutivas con la misma actividad (start = primera fila, end = timestamp de la siguiente)
                    rig["_act"] = rig[activity_col].astype(str).fillna("")
                    segs = []
                    i = 0
                    while i < len(rig):
                        act = rig["_act"].iloc[i]
                        t_start = rig[time_col].iloc[i]
                        j = i + 1
                        while j < len(rig) and rig["_act"].iloc[j] == act:
                            j += 1
                        if j < len(rig):
                            t_end = rig[time_col].iloc[j]
                        else:
                            # Último bloque: duración = diferencia con el siguiente registro o 1 registro = 0 min → usar intervalo medio
                            if len(rig) > 1:
                                dt_med = (rig[time_col].diff().dropna()).median()
                                t_end = t_start + (dt_med if pd.notna(dt_med) else pd.Timedelta(seconds=10))
                            else:
                                t_end = t_start + pd.Timedelta(seconds=10)
                        dur_min = (t_end - t_start).total_seconds() / 60.0
                        if dur_min > 0:
                            segs.append({"Start Time": t_start, "End Time": t_end, "Duration": dur_min, "Rig Activity": act})
                        i = j
                    if not segs:
                        st.info("No se generaron segmentos de actividad (revisa el CSV).")
                    else:
                        rig_seg = pd.DataFrame(segs)
                        st.success(f"Rig Activities cargado: {up_rig.name} → {len(rig_seg):,} segmentos.")

                        # Profundidad: Bit depth o Hole Depth
                        depth_col = None
                        for c in rig_raw.columns:
                            if "bit" in str(c).lower() and "depth" in str(c).lower():
                                depth_col = c
                                break
                        if depth_col is None:
                            for c in rig_raw.columns:
                                if "hole" in str(c).lower() and "depth" in str(c).lower():
                                    depth_col = c
                                    break
                        if depth_col:
                            rig[depth_col] = pd.to_numeric(rig[depth_col], errors="coerce")
                            rig_has_depth = rig[depth_col].notna().any() and (rig[depth_col] != 0).any()
                        else:
                            rig_has_depth = False

                        fechas_rig = sorted(pd.to_datetime(rig_seg["Start Time"]).dt.date.unique().tolist())
                        if fechas_rig:
                            sel_fecha_rig = st.date_input(
                                "Día (Rig Activities)",
                                value=min(fechas_rig),
                                min_value=min(fechas_rig),
                                max_value=max(fechas_rig),
                                key="rig_dia",
                            )
                            if sel_fecha_rig in fechas_rig:
                                rig_dia = rig_seg[pd.to_datetime(rig_seg["Start Time"]).dt.date == sel_fecha_rig].copy()
                                if rig_dia.empty:
                                    st.info("No hay segmentos para la fecha seleccionada.")
                                else:
                                    st.markdown("#### Timeline Rig Activities (por evento)")
                                    st.caption("Cada barra es un segmento de actividad en su minuto de inicio; la altura es la duración. No se acumulan.")
                                    # Chips pro: métricas del día (antes del gráfico para evitar superposición con leyenda)
                                    total_min_rig = float(rig_dia["Duration"].sum())
                                    total_h_rig = total_min_rig / 60.0
                                    top_act = rig_dia.groupby("Rig Activity", dropna=False)["Duration"].sum().sort_values(ascending=False)
                                    top_act_name = str(top_act.index[0]) if not top_act.empty else "-"
                                    top_act_min = float(top_act.iloc[0]) if not top_act.empty else 0.0
                                    n_act = int(rig_dia["Rig Activity"].nunique()) if "Rig Activity" in rig_dia.columns else 0
                                    # Semáforo: cobertura del día (horas con actividad / 24)
                                    horas_con_act = rig_dia["Start Time"].apply(lambda t: pd.Timestamp(t).hour).nunique() if "Start Time" in rig_dia.columns else 0
                                    cobertura_pct = (horas_con_act / 24.0 * 100.0) if horas_con_act else 0.0
                                    _tone_cob = "green" if cobertura_pct >= 80 else ("amber" if cobertura_pct >= 50 else "red")
                                    _tone_top = "green" if "drill" in top_act_name.lower() or "perfor" in top_act_name.lower() else "blue"
                                    render_chip_row([
                                        {"label": "Total día", "value": f"{total_h_rig:.2f} h", "tone": "blue"},
                                        {"label": "Actividad principal", "value": f"{top_act_name[:20]}{'…' if len(top_act_name) > 20 else ''}", "tone": _tone_top},
                                        {"label": "Principal (min)", "value": f"{top_act_min:.0f} min", "tone": "violet"},
                                        {"label": "Cobertura día", "value": f"{cobertura_pct:.0f}% ({horas_con_act}h)", "tone": _tone_cob},
                                        {"label": "Actividades distintas", "value": str(n_act), "tone": "gray"},
                                    ], use_iframe=True, height=100)
                                    segments_rig = []
                                    for _, r in rig_dia.iterrows():
                                        start = pd.Timestamp(r["Start Time"])
                                        dur_min = float(r["Duration"])
                                        if dur_min <= 0:
                                            continue
                                        k = str(r["Rig Activity"]) if pd.notna(r["Rig Activity"]) else "N/A"
                                        hour0 = start.hour
                                        remaining = dur_min
                                        min0 = start.minute + start.second / 60.0 + start.microsecond / 60000000.0
                                        while remaining > 1e-6 and hour0 < 24:
                                            seg_dur = min(remaining, 60.0 - min0)
                                            if seg_dur > 0:
                                                segments_rig.append((hour0, min0, seg_dur, k))
                                            remaining -= seg_dur
                                            hour0 += 1
                                            min0 = 0
                                    if segments_rig:
                                        seg_rig_df = pd.DataFrame(segments_rig, columns=["hora", "min_inicio", "dur_min", "Rig Activity"])
                                        act_orden = seg_rig_df["Rig Activity"].unique().tolist()
                                        colors_rig = {a: px.colors.qualitative.Set2[i % len(px.colors.qualitative.Set2)] for i, a in enumerate(act_orden)}
                                        if len(act_orden) > len(px.colors.qualitative.Set2):
                                            for i, a in enumerate(act_orden[len(px.colors.qualitative.Set2):]):
                                                colors_rig[a] = px.colors.qualitative.Pastel[i % len(px.colors.qualitative.Pastel)]
                                        fig_rig = go.Figure()
                                        for a in act_orden:
                                            sub = seg_rig_df[seg_rig_df["Rig Activity"] == a]
                                            fig_rig.add_trace(go.Bar(
                                                x=sub["hora"].tolist(),
                                                y=sub["dur_min"].tolist(),
                                                base=sub["min_inicio"].tolist(),
                                                name=str(a),
                                                marker_color=colors_rig.get(a, "gray"),
                                                width=0.85,
                                            ))
                                        layout_rig = dict(
                                            barmode="overlay",
                                            title=dict(text=f"Rig Activities por evento — {sel_fecha_rig}", x=0.5, xanchor="center", font=dict(size=16)),
                                            margin=dict(t=56, b=80, l=50, r=50),
                                            xaxis_title="Hora del día",
                                            yaxis_title="Minutos dentro de la hora",
                                            xaxis=dict(dtick=1, range=[-0.5, 23.5]),
                                            yaxis=dict(range=[0, 60], dtick=5),
                                            legend=dict(orientation="h", yanchor="top", y=-0.14, xanchor="center", x=0.5),
                                            height=420,
                                        )
                                        y2_rig = None
                                        if depth_col and rig_has_depth and time_col:
                                            # Profundidad al cierre de cada hora (desde datos crudos del día)
                                            rig_dia_dates = rig[pd.to_datetime(rig[time_col]).dt.date == sel_fecha_rig]
                                            if not rig_dia_dates.empty:
                                                ts = pd.to_datetime(rig_dia_dates[time_col])
                                                depth_vals = rig_dia_dates[depth_col]
                                                depth_at_h = []
                                                for h in range(24):
                                                    t_end = pd.Timestamp(sel_fecha_rig) + pd.Timedelta(hours=h + 1)
                                                    mask = ts <= t_end
                                                    if mask.any():
                                                        val = depth_vals.loc[mask].max()
                                                        depth_at_h.append(None if pd.isna(val) else float(val))
                                                    else:
                                                        depth_at_h.append(None)
                                                if any(x is not None for x in depth_at_h):
                                                    y2_rig = [x if x is not None else None for x in depth_at_h]
                                                    fig_rig.add_trace(go.Scatter(
                                                        x=list(range(24)),
                                                        y=y2_rig,
                                                        name="Profundidad (m)",
                                                        line=dict(color="#22c55e", width=2),
                                                        yaxis="y2",
                                                    ))
                                                    layout_rig["yaxis2"] = dict(title="Profundidad (m)", overlaying="y", side="right", showgrid=False)
                                        fig_rig.update_layout(**layout_rig)
                                        st.plotly_chart(fig_rig, use_container_width=True)
                                    else:
                                        st.info("No hay segmentos con duración para esta fecha.")
                            else:
                                st.info("Selecciona una fecha con datos.")
                        else:
                            st.caption("No se detectaron fechas en los segmentos de Rig Activities.")
        except Exception as e_rig:
            st.error(f"No se pudo procesar el CSV de Rig Activities. Error: {e_rig}")

# NUEVA TAB: REPORTE GENERAL DEL POZO (TODAS LAS ETAPAS)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
with tab_general:
    st.subheader("📊 Reporte General del Pozo - Todas las Etapas")
    
    # Verificar si hay datos
    if df.empty:
        st.info("No hay datos disponibles. Por favor, captura algunas actividades primero.")
    else:
        # ---- ESTILO: quitar chips rojos de multiselect (usar tonos neutros) ----
        st.markdown(
            """
            <style>
            /* Multiselect tags/chips */
            div[data-baseweb="tag"]{
                background-color: rgba(255,255,255,0.10) !important;
                border: 1px solid rgba(255,255,255,0.18) !important;
            }
            div[data-baseweb="tag"] span{
                color: rgba(255,255,255,0.90) !important;
            }
            /* 'x' button inside tag */
            div[data-baseweb="tag"] svg{
                fill: rgba(255,255,255,0.70) !important;
            }
            /* Select/multiselect control border */
            div[data-baseweb="select"] > div{
                border-color: rgba(255,255,255,0.18) !important;
                box-shadow: none !important;
            }
            div[data-baseweb="select"] > div:focus-within{
                border-color: rgba(255,255,255,0.35) !important;
                box-shadow: none !important;
            }
            </style>
            """,
            unsafe_allow_html=True,
        )

        # ---- FILTROS ----
        col_filt1, col_filt2, col_filt3 = st.columns(3)

        with col_filt1:
            if "Fecha" in df.columns:
                _fechas_dt = pd.to_datetime(df["Fecha"], errors="coerce")
                fechas_disponibles = sorted(_fechas_dt.dt.date.dropna().unique().tolist())
            else:
                fechas_disponibles = []
            fecha_seleccionada = st.selectbox(
                "Filtrar por fecha",
                options=["Todas las fechas"] + fechas_disponibles,
                index=0,
                key="filtro_fecha_general",
            )

        with col_filt2:
            # Selector de Tipo de tiempo (sin chips rojos): selectbox
            _opciones_tt = ['Todos', 'TP', 'TNPI', 'TNP']
            tipo_tiempo_sel = st.selectbox('Tipo de tiempo', options=_opciones_tt, index=0, key='tipo_tiempo_sel')
            tipos_tiempo_sel = ['TP','TNPI','TNP'] if tipo_tiempo_sel == 'Todos' else [tipo_tiempo_sel]

        with col_filt3:
            operaciones_disponibles = sorted(df["Operacion"].dropna().unique().tolist()) if "Operacion" in df.columns else []
            # Selector de Operación (sin chips rojos): selectbox
            _ops_operacion = ['Todas', 'Perforación', 'Viaje', 'Conexión', 'BHA', 'NPT', 'Otro']
            operacion_sel = st.selectbox('Filtrar por operación', options=_ops_operacion, index=0, key='operacion_sel')
            operaciones_sel = None if operacion_sel == 'Todas' else [operacion_sel]

        # Aplicar filtros
        df_filtrado = df.copy()

        if fecha_seleccionada != "Todas las fechas":
            df_filtrado = df_filtrado[df_filtrado["Fecha"] == fecha_seleccionada]

        # aplicar filtro de tipo de tiempo
        df_filtrado = df_filtrado[df_filtrado["Tipo"].isin(tipos_tiempo_sel)]

        # aplicar filtro de operación
        if operaciones_sel is not None:
            df_filtrado = df_filtrado[df_filtrado["Operacion"].isin(operaciones_sel)]

        # ---- KPIs GENERALES ----
        st.markdown("### 📈 KPIs Generales del Pozo")
        
        col_kpi1, col_kpi2, col_kpi3, col_kpi4, col_kpi5 = st.columns(5)
        
        with col_kpi1:
            total_horas = float(df_filtrado["Horas_Reales"].sum()) if not df_filtrado.empty else 0.0
            st.metric("Horas Totales", f"{total_horas:.1f} h")
        
        with col_kpi2:
            tp_horas = float(df_filtrado[df_filtrado["Tipo"] == "TP"]["Horas_Reales"].sum()) if not df_filtrado.empty else 0.0
            st.metric("TP (Horas Productivas)", f"{tp_horas:.1f} h")
        
        with col_kpi3:
            tnpi_horas = float(df_filtrado[df_filtrado["Tipo"] == "TNPI"]["Horas_Reales"].sum()) if not df_filtrado.empty else 0.0
            st.metric("TNPI (Horas)", f"{tnpi_horas:.1f} h")
        
        with col_kpi4:
            tnp_horas = float(df_filtrado[df_filtrado["Tipo"] == "TNP"]["Horas_Reales"].sum()) if not df_filtrado.empty else 0.0
            st.metric("TNP (Horas)", f"{tnp_horas:.1f} h")

        with col_kpi5:
            eficiencia_general = clamp_0_100(safe_pct(tp_horas, total_horas)) if total_horas > 0 else 0.0
            sk, sl, sc = status_from_eff(eficiencia_general)
            st.markdown(f"""
                <div style="text-align:center">
                    <div style="font-size:24px;font-weight:bold;color:{sc}">{eficiencia_general:.0f}%</div>
                    <div style="font-size:12px;color:#888">Eficiencia General</div>
                </div>
            """, unsafe_allow_html=True)
        
        # ---- GRÁFICAS GENERALES ----
        st.markdown("### 📊 Distribución General")
        
        # Gráfica 1: Horas por Etapa (Stacked)
        if not df_filtrado.empty:
            # Preparar datos para gráfica de etapas
            df_etapas = df_filtrado.groupby(["Etapa", "Tipo"])["Horas_Reales"].sum().reset_index()
            
            # Pivot table para stacked bar
            df_pivot = df_etapas.pivot_table(index="Etapa", columns="Tipo", values="Horas_Reales", fill_value=0).reset_index()
            
            # Ordenar por total de horas
            df_pivot["Total"] = df_pivot.sum(axis=1, numeric_only=True)
            df_pivot = df_pivot.sort_values("Total", ascending=True)
            
            fig_etapas = go.Figure()
            
            # Colores para los tipos
            colores = {"TP": "#2ECC71", "TNPI": "#E74C3C", "TNP": "#F1C40F"}
            
            for tipo in ["TNP", "TNPI", "TP"]:  # Orden inverso para mejor visualización
                if tipo in df_pivot.columns:
                    fig_etapas.add_trace(go.Bar(
                        name=tipo,
                        y=df_pivot["Etapa"],
                        x=df_pivot[tipo],
                        orientation='h',
                        marker_color=colores.get(tipo, "#3498DB"),
                        text=df_pivot[tipo].round(1),
                        textposition='auto',
                    ))
            
            fig_etapas.update_layout(
                title="Horas por Etapa - Todas las Etapas",
                barmode='stack',
                height=400,
                xaxis_title="Horas",
                yaxis_title="Etapa",
                showlegend=True,
                legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1)
            )
            
            st.plotly_chart(fig_etapas, use_container_width=True)
        
        # Gráfica 2: Distribución de actividades principales
        if not df_filtrado.empty:
            df_actividades = df_filtrado.groupby("Actividad")["Horas_Reales"].sum().reset_index()
            df_actividades = df_actividades.sort_values("Horas_Reales", ascending=False).head(10)
            
            fig_actividades = px.bar(
                df_actividades, 
                x="Horas_Reales", 
                y="Actividad", 
                orientation='h',
                title="Top 10 Actividades (Horas)",
                color="Horas_Reales",
                color_continuous_scale="Viridis"
            )
            fig_actividades.update_layout(height=400)
            st.plotly_chart(fig_actividades, use_container_width=True)
        
        # ---- ANÁLISIS DE CAUSAS (TNPI / TNP) ----
        st.markdown("### 🔎 Análisis de causas")

        col_a1, col_a2 = st.columns(2)

        with col_a1:
            st.markdown("#### 🔴 TNPI")
            df_tnpi_rg = df_filtrado[df_filtrado["Tipo"] == "TNPI"].copy()
            if df_tnpi_rg.empty:
                st.info("No hay registros TNPI para los filtros seleccionados.")
            else:
                for c, fb in [("Categoria_TNPI","Sin categoría"),("Detalle_TNPI","Sin detalle")]:
                    if c not in df_tnpi_rg.columns:
                        df_tnpi_rg[c]=fb
                    df_tnpi_rg[c]=df_tnpi_rg[c].fillna(fb).replace({"-":fb,"None":fb})
                df_cat = df_tnpi_rg.groupby("Categoria_TNPI", as_index=False)["Horas_Reales"].sum().sort_values("Horas_Reales", ascending=False)
                fig = px.pie(df_cat, names="Categoria_TNPI", values="Horas_Reales", hole=0.55, title="TNPI por categoría (h)")
                st.plotly_chart(fig, use_container_width=True)
                df_det = df_tnpi_rg.groupby("Detalle_TNPI", as_index=False)["Horas_Reales"].sum().sort_values("Horas_Reales", ascending=False).head(10)
                fig2 = px.bar(df_det, x="Horas_Reales", y="Detalle_TNPI", orientation='h', title="Top 10 - Detalle TNPI (h)")
                st.plotly_chart(fig2, use_container_width=True)

        with col_a2:
            st.markdown("#### 🟡 TNP")
            df_tnp_rg = df_filtrado[df_filtrado["Tipo"] == "TNP"].copy()
            if df_tnp_rg.empty:
                st.info("No hay registros TNP para los filtros seleccionados.")
            else:
                for c, fb in [("Categoria_TNP","Sin categoría"),("Detalle_TNP","Sin detalle")]:
                    if c not in df_tnp_rg.columns:
                        df_tnp_rg[c]=fb
                    df_tnp_rg[c]=df_tnp_rg[c].fillna(fb).replace({"-":fb,"None":fb})
                df_cat = df_tnp_rg.groupby("Categoria_TNP", as_index=False)["Horas_Reales"].sum().sort_values("Horas_Reales", ascending=False)
                fig = px.pie(df_cat, names="Categoria_TNP", values="Horas_Reales", hole=0.55, title="TNP por categoría (h)")
                st.plotly_chart(fig, use_container_width=True)
                df_det = df_tnp_rg.groupby("Detalle_TNP", as_index=False)["Horas_Reales"].sum().sort_values("Horas_Reales", ascending=False).head(10)
                fig2 = px.bar(df_det, x="Horas_Reales", y="Detalle_TNP", orientation='h', title="Top 10 - Detalle TNP (h)")
                st.plotly_chart(fig2, use_container_width=True)

        # ---- TABLAS DETALLADAS ----
        st.markdown("### 📋 Resumen por Etapa")
        
        # Crear resumen por etapa
        if not df_filtrado.empty:
            resumen_etapas = []
            etapas_unicas = sorted(df_filtrado["Etapa"].unique())
            
            for etapa_actual in etapas_unicas:
                df_etapa_actual = df_filtrado[df_filtrado["Etapa"] == etapa_actual]
                
                # Calcular KPIs para esta etapa
                total_etapa = float(df_etapa_actual["Horas_Reales"].sum())
                tp_etapa = float(df_etapa_actual[df_etapa_actual["Tipo"] == "TP"]["Horas_Reales"].sum())
                tnpi_etapa = float(df_etapa_actual[df_etapa_actual["Tipo"] == "TNPI"]["Horas_Reales"].sum())
                tnp_etapa = float(df_etapa_actual[df_etapa_actual["Tipo"] == "TNP"]["Horas_Reales"].sum())
                
                eficiencia_etapa = clamp_0_100(safe_pct(tp_etapa, total_etapa)) if total_etapa > 0 else 0.0
                
                # Contar conexiones para esta etapa
                conexiones_etapa = 0
                if not df_conn.empty and "Seccion" in df_conn.columns:
                    conexiones_etapa = len(df_conn[df_conn["Seccion"] == etapa_actual]["Conn_No"].unique())
                
                # Contar BHA para esta etapa
                bha_etapa = 0
                if not df_bha.empty and "Etapa" in df_bha.columns:
                    bha_etapa = len(df_bha[df_bha["Etapa"] == etapa_actual])
                
                resumen_etapas.append({
                    "Etapa": etapa_actual,
                    "Horas Totales": f"{total_etapa:.1f}",
                    "TP (h)": f"{tp_etapa:.1f}",
                    "TNPI (h)": f"{tnpi_etapa:.1f}",
                    "TNP (h)": f"{tnp_etapa:.1f}",
                    "Eficiencia (%)": f"{eficiencia_etapa:.0f}%",
                    "Conexiones": f"{conexiones_etapa}",
                    "Operaciones BHA": f"{bha_etapa}",
                    "Semáforo": semaforo_dot(eficiencia_etapa)
                })
            
            # Crear DataFrame y mostrar
            df_resumen_etapas = pd.DataFrame(resumen_etapas)
            st.dataframe(df_resumen_etapas, use_container_width=True, hide_index=True)
        
        # ---- ANÁLISIS DE TNPI GENERAL ----
        st.markdown("### 🔍 Análisis de TNPI - Todas las Etapas")
        
        if tnpi_horas > 0:
            # Top causas de TNPI en todas las etapas
            df_tnpi_general = df_filtrado[df_filtrado["Tipo"] == "TNPI"].copy()
            
            col_tnpi1, col_tnpi2 = st.columns(2)
            
            with col_tnpi1:
                # Por categoría
                df_tnpi_cat = df_tnpi_general.groupby("Categoria_TNPI")["Horas_Reales"].sum().reset_index()
                df_tnpi_cat = df_tnpi_cat.sort_values("Horas_Reales", ascending=False)
                
                if not df_tnpi_cat.empty:
                    fig_tnpi_cat = px.bar(
                        df_tnpi_cat, 
                        x="Horas_Reales", 
                        y="Categoria_TNPI", 
                        orientation='h',
                        title="TNPI por Categoría (h)",
                        color="Horas_Reales",
                        color_continuous_scale="Reds"
                    )
                    fig_tnpi_cat.update_layout(height=300)
                    st.plotly_chart(fig_tnpi_cat, use_container_width=True)
            
            with col_tnpi2:
                # Por etapa
                df_tnpi_etapa = df_tnpi_general.groupby("Etapa")["Horas_Reales"].sum().reset_index()
                df_tnpi_etapa = df_tnpi_etapa.sort_values("Horas_Reales", ascending=True)
                
                if not df_tnpi_etapa.empty:
                    fig_tnpi_etapa = px.bar(
                        df_tnpi_etapa, 
                        x="Horas_Reales", 
                        y="Etapa", 
                        orientation='h',
                        title="TNPI por Etapa (h)",
                        color="Horas_Reales",
                        color_continuous_scale="Oranges"
                    )
                    fig_tnpi_etapa.update_layout(height=300)
                    st.plotly_chart(fig_tnpi_etapa, use_container_width=True)
            
            # Tabla detallada de TNPI
            st.markdown("**Detalle de TNPI por etapa y categoría**")
            df_tnpi_detalle = df_tnpi_general.groupby(["Etapa", "Categoria_TNPI", "Detalle_TNPI", "Categoria_TNP", "Detalle_TNP"])["Horas_Reales"].sum().reset_index()
            df_tnpi_detalle = df_tnpi_detalle.sort_values(["Etapa", "Horas_Reales"], ascending=[True, False])
            
            if not df_tnpi_detalle.empty:
                st.dataframe(df_tnpi_detalle, use_container_width=True, height=300)
        else:
            st.success("🎉 No hay TNPI registrado en el período seleccionado")
        


        # ---- ANÁLISIS DE TNP GENERAL ----
        st.markdown("### 🔍 Análisis de TNP - Todas las Etapas")

        if tnp_horas > 0:
            df_tnp_general = df_filtrado[df_filtrado["Tipo"] == "TNP"].copy()

            # Normalizar (evitar NaN / '-')
            for col, fallback in [("Categoria_TNP", "Sin categoría"), ("Detalle_TNP", "Sin detalle")]:
                if col not in df_tnp_general.columns:
                    df_tnp_general[col] = fallback
                df_tnp_general[col] = (
                    df_tnp_general[col]
                    .fillna(fallback)
                    .astype(str)
                    .replace({"-": fallback, "None": fallback, "nan": fallback})
                )

            col_tnp1, col_tnp2 = st.columns(2)

            with col_tnp1:
                df_tnp_cat = df_tnp_general.groupby("Categoria_TNP")["Horas_Reales"].sum().reset_index()
                df_tnp_cat = df_tnp_cat.sort_values("Horas_Reales", ascending=True)
                if not df_tnp_cat.empty:
                    fig_tnp_cat = px.bar(
                        df_tnp_cat,
                        x="Horas_Reales",
                        y="Categoria_TNP",
                        orientation='h',
                        title="TNP por Categoría (h)",
                        color="Horas_Reales",
                        color_continuous_scale="Blues"
                    )
                    fig_tnp_cat.update_layout(height=300)
                    st.plotly_chart(fig_tnp_cat, use_container_width=True)

            with col_tnp2:
                df_tnp_etapa = df_tnp_general.groupby("Etapa")["Horas_Reales"].sum().reset_index()
                df_tnp_etapa = df_tnp_etapa.sort_values("Horas_Reales", ascending=True)
                if not df_tnp_etapa.empty:
                    fig_tnp_etapa = px.bar(
                        df_tnp_etapa,
                        x="Horas_Reales",
                        y="Etapa",
                        orientation='h',
                        title="TNP por Etapa (h)",
                        color="Horas_Reales",
                        color_continuous_scale="Teal"
                    )
                    fig_tnp_etapa.update_layout(height=300)
                    st.plotly_chart(fig_tnp_etapa, use_container_width=True)

            st.markdown("**Detalle de TNP por etapa, categoría y detalle**")
            df_tnp_detalle = df_tnp_general.groupby(["Etapa", "Categoria_TNP", "Detalle_TNP"])["Horas_Reales"].sum().reset_index()
            df_tnp_detalle = df_tnp_detalle.sort_values(["Etapa", "Horas_Reales"], ascending=[True, False])
            if not df_tnp_detalle.empty:
                st.dataframe(df_tnp_detalle, use_container_width=True, height=300)
        else:
            st.success("🎉 No hay TNP registrado en el período seleccionado")

        # ---- EXPORTAR REPORTE GENERAL ----
        st.markdown("### 📥 Exportar Reporte General")
        
        col_exp1, col_exp2 = st.columns(2)
        
        with col_exp1:
            if st.button("📊 Generar Reporte PDF", use_container_width=True):
                # Aquí iría la lógica para generar PDF del reporte general
                st.success("Reporte general generado (funcionalidad en desarrollo)")
        
        with col_exp2:
            if st.button("📈 Exportar Datos a Excel", use_container_width=True):
                # Crear Excel con datos generales
                output = BytesIO()
                with pd.ExcelWriter(output, engine='openpyxl') as writer:
                    df_filtrado.to_excel(writer, sheet_name='Datos_Completos', index=False)
                    if not df_conn.empty:
                        df_conn.to_excel(writer, sheet_name='Conexiones', index=False)
                    if not df_bha.empty:
                        df_bha.to_excel(writer, sheet_name='BHA', index=False)
                
                output.seek(0)
                st.download_button(
                    label="Descargar Excel",
                    data=output,
                    file_name=f"Reporte_General_{pozo}_{datetime.now().strftime('%Y%m%d')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
# TAB: EJECUTIVO (Causa–raíz + Recomendaciones + PDF)
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
with tab_ejecutivo:
    st.subheader("Análisis causa–raíz (Viajes)")
    df_main = st.session_state.df.copy()

    # --- TNPI Viajes: velocidad vs conexiones ---
    df_vtnpi = df_main[(df_main["Tipo"] == "TNPI") & (df_main["Origen"].fillna("") == "Viajes y conexiones")].copy()

    vel_mask = df_vtnpi["Detalle_TNPI"].fillna("").str.contains("Velocidad", case=False)
    conn_mask = df_vtnpi["Detalle_TNPI"].fillna("").str.contains("Conexi", case=False)

    tnpi_total_h = float(df_vtnpi["Horas_Reales"].sum()) if not df_vtnpi.empty else 0.0
    tnpi_vel_h = float(df_vtnpi.loc[vel_mask, "Horas_Reales"].sum()) if (not df_vtnpi.empty) else 0.0
    tnpi_conn_h = float(df_vtnpi.loc[conn_mask, "Horas_Reales"].sum()) if (not df_vtnpi.empty) else 0.0
    tnpi_otros_h = max(0.0, tnpi_total_h - tnpi_vel_h - tnpi_conn_h)


    c1, c2, c3 = st.columns(3)
    c1.metric("TNPI Viajes – Velocidad (h)", f"{tnpi_vel_h:.2f}")
    c2.metric("TNPI Viajes – Conexiones (h)", f"{tnpi_conn_h:.2f}")
    c3.metric("TNPI Viajes – Total (h)", f"{tnpi_total_h:.2f}")

    # Donut % (si hay datos)
    fig_donut = None
    if tnpi_total_h > 0 and PLOTLY_IMG_OK:
        ddf = pd.DataFrame(
            {"Causa": ["Velocidad", "Conexiones", "Otros"], "Horas": [tnpi_vel_h, tnpi_conn_h, tnpi_otros_h]}
        )
        fig_donut = px.pie(ddf, names="Causa", values="Horas", hole=0.55, title="TNPI Viajes – Distribución (%)")
        fig_donut.update_traces(textinfo="percent+label")
        st.plotly_chart(fig_donut, use_container_width=True)
    elif tnpi_total_h == 0:
        st.info("Aún no hay TNPI de viajes registrado para el día.")

    st.divider()

    # --- Recomendaciones automáticas ---
    st.subheader("Recomendaciones automáticas")
    recos = []
    razones = []

    if tnpi_total_h == 0:
        recos.append("Sin TNPI en viajes registrado: mantener parámetros y disciplina operativa.")
    else:
        p_vel = tnpi_vel_h / tnpi_total_h if tnpi_total_h > 0 else 0.0
        p_conn = tnpi_conn_h / tnpi_total_h if tnpi_total_h > 0 else 0.0

        if p_conn >= 0.60:
            recos += [
                "Priorizar mejora de conexiones: checklist, roles claros y preparación previa (preconexión).",
                "Revisar herramientas/llave/MPD y tiempos muertos recurrentes durante conexiones.",
                "Validar handover turno a turno y asegurar que materiales/herramientas estén listos antes del pico de conexiones."
            ]
            razones.append(f"Conexiones representan {p_conn*100:.0f}% del TNPI de viajes.")
        if p_vel >= 0.60:
            recos += [
                "Priorizar mejora de velocidad de viaje: revisar arrastre/fricción y condiciones del hoyo.",
                "Ajustar prácticas (barrido/limpieza) y revisar límites operativos que reduzcan velocidad.",
                "Evaluar si el método (Lingadas/TxT) está siendo aplicado correctamente por tramo."
            ]
            razones.append(f"Velocidad representa {p_vel*100:.0f}% del TNPI de viajes.")

        if not recos:
            recos.append("TNPI distribuido entre velocidad y conexiones: atacar las 2 principales horas críticas y estandarizar el método por tramo.")

    # Horas críticas (top 3)
    try:
        df_viajes_h = st.session_state.get("viajes_por_hora_df", None)
    except Exception:
        df_viajes_h = None

    if df_viajes_h is not None and isinstance(df_viajes_h, pd.DataFrame) and not df_viajes_h.empty:
        # Espera columnas: hour, tnpi_vel_h, tnpi_conn_h o tnpi_total_h
        cand_cols = [c for c in df_viajes_h.columns if "tnpi" in c.lower()]
        if cand_cols:
            df_tmp = df_viajes_h.copy()
            if "tnpi_total_h" not in df_tmp.columns:
                # crea total si hay vel/conn
                vcol = "tnpi_vel_h" if "tnpi_vel_h" in df_tmp.columns else None
                ccol = "tnpi_conn_h" if "tnpi_conn_h" in df_tmp.columns else None
                if vcol or ccol:
                    df_tmp["tnpi_total_h"] = (df_tmp[vcol].fillna(0) if vcol else 0) + (df_tmp[ccol].fillna(0) if ccol else 0)
            top = df_tmp.sort_values("tnpi_total_h", ascending=False).head(3)
            horas = [f"{int(h):02d}:00" for h in top["hour"].tolist()] if "hour" in top.columns else []
            if horas and float(top["tnpi_total_h"].sum()) > 0:
                razones.append("Horas críticas (mayor TNPI): " + ", ".join(horas))

    if razones:
        st.caption(" • " + " • ".join(razones))

    for r in recos[:6]:
        st.write("• " + r)

    st.divider()

    # --- Export PDF Ejecutivo ---
    st.subheader("Export ejecutivo (PDF)")
    st.caption("Genera un PDF en tamaño Carta con KPIs + gráficas clave (Viajes/Conexiones) + recomendaciones.")

    # Tomamos las figuras de la pestaña de viajes si existen en session_state (si no, no falla)
    fig_speed = st.session_state.get("fig_viaje_speed", None)
    fig_conn = st.session_state.get("fig_viaje_conn", None)

    meta_pdf = {"equipo": st.session_state.get("equipo_val", ""), "pozo": st.session_state.get("pozo_val", ""), "etapa": etapa, "fecha": str(fecha)}
    kpis_pdf = {
        "Eficiencia global (%)": f"{_eff_prev:.0f}%",
        "TNPI Viajes (h)": f"{tnpi_total_h:.2f}",
        "TNPI Velocidad (h)": f"{tnpi_vel_h:.2f}",
        "TNPI Conexiones (h)": f"{tnpi_conn_h:.2f}",
    }

    charts_pdf = {}
    if fig_speed is not None:
        charts_pdf["Viaje – Velocidad por hora"] = fig_speed
    if fig_conn is not None:
        charts_pdf["Conexiones – Min/conn por hora"] = fig_conn
    if fig_donut is not None:
        charts_pdf["TNPI Viajes – Distribución (%)"] = fig_donut

    # Adjuntamos recomendaciones como string en KPIs (para que salgan en el PDF actual sin remaquetar)
    if recos:
        kpis_pdf["Recomendaciones"] = " | ".join(recos[:4])

    pdf_bytes = build_pdf(meta_pdf, kpis_pdf, charts_pdf)
    st.download_button(
        "Descargar PDF diario (Carta)",
        data=pdf_bytes,
        file_name=f"Reporte_DrillSpot_{pozo}_{str(fecha)}.pdf",
        mime="application/pdf",
        use_container_width=True,
    )


# TAB: EXPORTAR
# == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == == =
with tab_export:
    st.subheader("Exportar (PDF / PowerPoint)")
    render_export_diario_calendario()

    meta = {"equipo": st.session_state.get("equipo_val", ""), "pozo": st.session_state.get("pozo_val", ""), "etapa": etapa, "fecha": str(fecha)}
    kpis_export = {
        "Modo": modo_reporte,
        "TP (h)": f"{tp_h:.2f}",
        "TNPI (h)": f"{tnpi_h:.2f}",
        "TNP (h)": f"{tnp_h:.2f}",
        "Eficiencia del día": f"{eficiencia_dia:.0f}%",
    }

    if modo_reporte == "Perforación":
        kpis_export["PT programada (m)"] = f"{float(st.session_state.drill_day['pt_programada_m']):.0f}"
        kpis_export["Profundidad actual (m)"] = f"{float(st.session_state.drill_day['prof_actual_m']):.0f}"
        kpis_export["Metros programa (m)"] = f"{float(st.session_state.drill_day['metros_prog_total']):.0f}"
        kpis_export["Metros real (m)"] = f"{float(st.session_state.drill_day['metros_real_dia'] + st.session_state.drill_day['metros_real_noche']):.0f}"
        kpis_export["ROP programa (m/h)"] = f"{float(st.session_state.drill_day['rop_prog_total']):.2f}"
        rr_d = float(st.session_state.drill_day["rop_real_dia"])
        rr_n = float(st.session_state.drill_day["rop_real_noche"])
        rr_local = (rr_d + rr_n) / (2 if (rr_d > 0 and rr_n > 0) else 1) if (rr_d > 0 or rr_n > 0) else 0.0
        kpis_export["ROP real (m/h)"] = f"{rr_local:.2f}"

    charts_export = {}
    if show_charts:
        for key, label in [
            ("tiempos", "Distribución de tiempos"),
            ("act_pie", "Distribución actividades (pie)"),
            ("act_bar", "Distribución actividades (bar)"),
            ("conn_pie", "Distribución tiempo en conexión (pie)"),
            ("conn_stack", "Conexiones perforando (stack)"),
        ]:
            if figs.get(key) is not None:
                charts_export[label] = figs[key]

    sig_main = f"{pozo}|{etapa}|{fecha}|{modo_reporte}|{repr(kpis_export)}|{list(charts_export.keys())}"
    if st.session_state.get("exp_main_sig") != sig_main:
        st.session_state["exp_main_sig"] = sig_main
        st.session_state.pop("exp_main_pdf", None)
        st.session_state.pop("exp_main_ppt", None)

    if st.session_state.get("exp_main_pdf") is None or st.session_state.get("exp_main_ppt") is None:
        st.caption("Para acelerar la interfaz, genera el PDF/PPTX bajo demanda.")
        if st.button("Preparar exportables", use_container_width=True, key="exp_main_prepare"):
            with st.spinner("Generando exportables..."):
                prog_main = st.progress(0)
                prog_main_msg = st.empty()
                prog_main_msg.caption("Iniciando...")
                st.session_state["exp_main_pdf"] = build_pdf(meta, kpis_export, charts=charts_export)
                prog_main.progress(55)
                prog_main_msg.caption("PDF listo.")
                st.session_state["exp_main_ppt"] = build_pptx(meta, kpis_export, charts_export)
                prog_main.progress(100)
                prog_main_msg.caption("PowerPoint listo.")

    if st.session_state.get("exp_main_pdf") is not None:
        col_pdf, col_ppt = st.columns(2)
        with col_pdf:
            fname_pdf = f"Reporte_DrillSpot_{pozo}_{etapa}_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf"
            st.download_button(
                "Descargar PDF",
                data=st.session_state.get("exp_main_pdf"),
                file_name=fname_pdf,
                mime="application/pdf",
                use_container_width=True,
            )

        with col_ppt:
            fname_pptx = f"Reporte_DrillSpot_{pozo}_{etapa}_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
            st.download_button(
                "Descargar PowerPoint",
                data=st.session_state.get("exp_main_ppt"),
                file_name=fname_pptx,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )

    if not PLOTLY_IMG_OK:
        st.caption("Para exportar gráficas como imágenes instala: `pip install -U kaleido`.")

# NOTE: Added corrected Captura actividad block (see above).

# --- FIX: sincronizar depth_rows desde drill_day (por etapa) ---
rows = []
for _etapa_k, _data in st.session_state.drill_day.get("por_etapa", {}).items():
    rows.append({
        "Etapa": _etapa_k,
        "PT_programada_m": _data.get("pt_programada_m", 0.0),
        "PT_actual_m": _data.get("prof_actual_m", 0.0),
    })
st.session_state.depth_rows = pd.DataFrame(rows)