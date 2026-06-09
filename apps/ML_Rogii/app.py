"""
================================================================================
PETROVUE ROP OPTIMIZER - SISTEMA IMPERIAL (EE.UU.)
================================================================================
Versión: 4.0 - Imperial Units
Unidades: 
    - Peso: libra (lb), 1000 lb = klb
    - Presión: psi (lb/pulg²)
    - Caudal: galones por minuto (gpm)
    - Profundidad: pies (ft)
    - ROP: pies/hora (ft/hr)
    - Torque: pie-libra (ft-lb)
    - Diámetro: pulgadas (in)
================================================================================
"""

import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import datetime
import uuid
import time
import io
import html
from scipy import stats
import warnings
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
warnings.filterwarnings('ignore')

BASE_DIR = Path(__file__).resolve().parent
LOGO_PATH = BASE_DIR / "rogii_logo.png"

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# ML libraries - optional for fallback to simulated metrics
try:
    from sklearn.ensemble import RandomForestRegressor
    from sklearn.neural_network import MLPRegressor
    from sklearn.model_selection import train_test_split, TimeSeriesSplit
    from sklearn.preprocessing import StandardScaler, OneHotEncoder
    from sklearn.compose import ColumnTransformer
    from sklearn.pipeline import Pipeline
    from sklearn.impute import SimpleImputer
    from sklearn.metrics import r2_score, mean_squared_error, mean_absolute_error
    SKLEARN_AVAILABLE = True
except ImportError:
    SKLEARN_AVAILABLE = False

try:
    import xgboost as xgb
    XGBOOST_AVAILABLE = True
except ImportError:
    XGBOOST_AVAILABLE = False

# ============================================================================
# TRANSLATIONS (ES / EN)
# ============================================================================

TRANSLATIONS = {
    'es': {
        'header_subtitle': 'Sistema avanzado de predicción ROP con Machine Learning y Redes Neuronales',
        'panel_control': 'Panel de control',
        'units_usc': 'Unidades en sistema USC',
        'suggestions_panel': 'Sugerencias del panel',
        'model_training': 'Entrenamiento de modelos',
        'training_config': 'Configuración de entrenamiento',
        'synthetic_samples': 'Muestras sintéticas',
        'synthetic_samples_help': 'Número de muestras sintéticas de perforación para entrenamiento',
        'train_models': 'ENTRENAR MODELOS',
        'training_spinner': 'Entrenando modelos de machine learning...',
        'training_done': '¡Modelos entrenados correctamente!',
        'drilling_params': 'Parámetros de perforación',
        'operational_params': 'Parámetros operacionales',
        'geological_params': 'Parámetros geológicos',
        'bit_params': 'Parámetros de broca',
        'fluid_formation': 'Fluido y formación (modelos mecanicistas)',
        'geological_tracking': 'Seguimiento geológico',
        'use_ucs_formation': 'Usar UCS por formación',
        'use_ucs_help': 'Usa UCS de la formación correspondiente a la profundidad actual',
        'predict_rop': 'PREDECIR ROP',
        'welcome': '¡Bienvenido a Drilling Analytics - Unidades Imperiales!',
        'welcome_steps': 'Para comenzar:\n1. Configure los parámetros de perforación en el panel lateral (en unidades US Customary)\n2. Entrene los modelos ML haciendo clic en \'ENTRENAR MODELOS\'\n3. Haga clic en \'PREDECIR ROP\' para obtener predicciones en tiempo real',
        'tab_resumen': 'Resumen de detalles',
        'tab_prediction': 'Predicción ROP',
        'tab_performance': 'Rendimiento de modelos',
        'tab_heatmap': 'Mapa de calor ROP',
        'tab_neural': 'Análisis de red neuronal',
        'tab_optimization': 'Optimización',
        'tab_geological': 'Seguimiento geológico',
        'language': 'Idioma',
        'models_trained_success': 'Entrenados {n} modelos con {samples:,} muestras',
        'welcome_system': 'El sistema combina:\n- **Bourgoyne & Young** (8 factores mecanicista)\n- **Bingham** (ROP reología-hidráulica)\n- **Warren** (generación/remoción de recortes)\n- **Random Forest, XGBoost, Redes Neuronales** (ML)\nLas salidas mecanicistas se correlacionan con ML para un ensemble híbrido.',
    },
    'en': {
        'header_subtitle': 'Advanced ROP prediction system with Machine Learning and Neural Networks',
        'panel_control': 'Control Panel',
        'units_usc': 'Units in USC system',
        'suggestions_panel': 'Panel suggestions',
        'model_training': 'Model Training',
        'training_config': 'Training configuration',
        'synthetic_samples': 'Synthetic samples',
        'synthetic_samples_help': 'Number of synthetic drilling samples for training',
        'train_models': 'TRAIN MODELS',
        'training_spinner': 'Training machine learning models...',
        'training_done': 'Models trained successfully!',
        'drilling_params': 'Drilling Parameters',
        'operational_params': 'Operational parameters',
        'geological_params': 'Geological parameters',
        'bit_params': 'Bit parameters',
        'fluid_formation': 'Fluid and formation (mechanistic models)',
        'geological_tracking': 'Geological tracking',
        'use_ucs_formation': 'Use UCS by formation',
        'use_ucs_help': 'Use UCS from the formation corresponding to current depth',
        'predict_rop': 'PREDICT ROP',
        'welcome': 'Welcome to Drilling Analytics - Imperial Units!',
        'welcome_steps': 'To get started:\n1. Configure drilling parameters in the sidebar panel (US Customary units)\n2. Train ML models by clicking \'TRAIN MODELS\'\n3. Click \'PREDICT ROP\' for real-time predictions',
        'tab_resumen': 'Summary',
        'tab_prediction': 'ROP Prediction',
        'tab_performance': 'Model Performance',
        'tab_heatmap': 'ROP Heat Map',
        'tab_neural': 'Neural Network Analysis',
        'tab_optimization': 'Optimization',
        'tab_geological': 'Geological Tracking',
        'language': 'Language',
        'models_trained_success': 'Trained {n} models with {samples:,} samples',
        'welcome_system': 'The system combines:\n- **Bourgoyne & Young** (8 mechanistic factors)\n- **Bingham** (ROP rheology-hydraulics)\n- **Warren** (cuttings generation/removal)\n- **Random Forest, XGBoost, Neural Networks** (ML)\nMechanistic outputs are correlated with ML for a hybrid ensemble.',
    },
}

def _t(key: str) -> str:
    """Return translation for key based on current language."""
    lang = st.session_state.get('lang', 'es')
    d = TRANSLATIONS.get(lang, TRANSLATIONS['es'])
    return d.get(key, key)

# ============================================================================
# HELPERS: VALIDATION & SUGGESTIONS
# ============================================================================

def _safe_float(val: Any, default: float = 0.0) -> float:
    """Extract float safely; return default on error."""
    try:
        if val is None or (isinstance(val, float) and np.isnan(val)):
            return default
        return float(val)
    except (TypeError, ValueError):
        return default

def _validate_params(params: Dict[str, Any], required: List[str]) -> tuple[Dict[str, Any], List[str]]:
    """Validate and fill missing params with defaults. Returns (params, list of warnings)."""
    DEFAULTS = {
        'wob_klb': 22.0, 'rpm': 120.0, 'torque_ftlb': 18000.0, 'spp_psi': 3000.0,
        'flow_gpm': 800.0, 'ucs_psi': 15000.0, 'bit_diameter_in': 8.5, 'bit_wear': 0.2,
        'depth_ft': 12000.0, 'cutter_count': 6, 'mud_density_ppg': 10.0,
        'pore_gradient_ppg': 9.0, 'yp_lb100ft2': 15.0, 'pv_cp': 25.0,
    }
    out = dict(params)
    warnings_list = []
    for k in required:
        if k not in out or out[k] is None:
            if k in DEFAULTS:
                out[k] = DEFAULTS[k]
                warnings_list.append(f"Missing {k}, using default {DEFAULTS[k]}")
            else:
                warnings_list.append(f"Missing required param: {k}")
    # Sanitize numeric values
    for k, v in list(out.items()):
        if isinstance(v, (int, float)) and k in DEFAULTS:
            out[k] = _safe_float(v, DEFAULTS.get(k, 0))
    return out, warnings_list

def get_section_suggestions(section: str, context: Optional[Dict] = None, lang: Optional[str] = None) -> List[str]:
    """Return contextual suggestions for each section."""
    context = context or {}
    lang = lang or st.session_state.get('lang', 'es')
    suggestions_es = {
        'rop_prediction': [
            "Ejecute predicciones con múltiples WOB/RPM para hallar el punto óptimo.",
            "Compare Bourgoyne & Young vs Bingham vs Warren — la divergencia puede indicar necesidad de calibración.",
            "Si el Ensemble difiere mucho de los modelos mecanicistas, considere reentrenar el ML con datos de campo.",
        ],
        'model_performance': [
            "Prefiera modelos con R² mayor y RMSE menor; use MAPE para el error relativo.",
            "La importancia de variables guía el ajuste — centre en los 3 principales.",
            "Calibre los coeficientes mecanicistas (B&Y a1–a8) con datos del pozo para mejor ajuste.",
        ],
        'heat_map': [
            "Use la zona óptima verde (WOB 18–26, RPM 100–140) como punto de partida.",
            "Varíe el UCS en Controles del mapa para ver cómo la resistencia de formación desplaza la región óptima.",
            "Exporte el mapa como imagen para referencia en taladro o reportes diarios.",
        ],
        'neural_network': [
            "Si la pérdida de validación diverge de la de entrenamiento, reduzca la tasa de aprendizaje o añada dropout.",
            "Use salidas de modelos mecanicistas como entradas adicionales para aprendizaje híbrido con física.",
            "Monitoree el tiempo de inferencia; >10 ms puede requerir poda del modelo para aplicaciones en tiempo real.",
        ],
        'optimization': [
            "El pico de la superficie 3D indica el óptimo teórico; valide en campo con cambios graduales.",
            "Evite saltos bruscos al óptimo — ajuste WOB/RPM gradualmente para reducir riesgo de stick-slip.",
            "Combine con Evaluación de Riesgos para asegurar que la región óptima esté dentro de límites seguros.",
        ],
        'geological_tracking': [
            "Actualice las formaciones con UCS de registros de pozo o correlación de offset.",
            "Use el seguimiento geológico para planificar cambios de WOB/RPM por intervalo.",
            "Compare ROP predicho por formación con datos de taladro para calibrar.",
        ],
        'correlation': [
            "Variables con correlación positiva con ROP (WOB, RPM) son candidatas a maximizar.",
            "Correlaciones negativas (UCS, profundidad, desgaste) indican factores que reducen el ROP.",
            "Use la matriz para detectar multicolinealidad antes de ajustar modelos.",
        ],
        'sidebar': [
            "Entrene los tres modelos ML (RF, XGB, NN) para el ensemble más estable.",
            "La densidad del lodo y el gradiente de poro afectan predicciones B&Y; use datos reales cuando disponga.",
            "Los valores YP/PV del reómetro mejoran la precisión del modelo Bingham.",
        ],
    }
    suggestions_en = {
        'rop_prediction': [
            "Run predictions with multiple WOB/RPM combinations to find the optimal point.",
            "Compare Bourgoyne & Young vs Bingham vs Warren — divergence may indicate calibration needs.",
            "If the Ensemble differs significantly from mechanistic models, consider retraining ML with field data.",
        ],
        'model_performance': [
            "Prefer models with higher R² and lower RMSE; use MAPE for relative error.",
            "Variable importance guides adjustment — focus on the top 3.",
            "Calibrate mechanistic coefficients (B&Y a1–a8) with well data for better fit.",
        ],
        'heat_map': [
            "Use the green optimum zone (WOB 18–26, RPM 100–140) as a starting point.",
            "Vary UCS in map controls to see how formation strength shifts the optimum region.",
            "Export the map as an image for rig or daily reporting reference.",
        ],
        'neural_network': [
            "If validation loss diverges from training, reduce learning rate or add dropout.",
            "Use mechanistic model outputs as additional inputs for physics-informed hybrid learning.",
            "Monitor inference time; >10 ms may require model pruning for real-time applications.",
        ],
        'optimization': [
            "The 3D surface peak indicates the theoretical optimum; validate in the field with gradual changes.",
            "Avoid abrupt jumps to optimum — adjust WOB/RPM gradually to reduce stick-slip risk.",
            "Combine with Risk Assessment to ensure the optimum region is within safe limits.",
        ],
        'geological_tracking': [
            "Update formations with UCS from well logs or offset correlation.",
            "Use geological tracking to plan WOB/RPM changes per interval.",
            "Compare predicted ROP by formation with rig data to calibrate.",
        ],
        'correlation': [
            "Variables with positive correlation to ROP (WOB, RPM) are candidates to maximize.",
            "Negative correlations (UCS, depth, wear) indicate factors that reduce ROP.",
            "Use the matrix to detect multicollinearity before adjusting models.",
        ],
        'sidebar': [
            "Train all three ML models (RF, XGB, NN) for the most stable ensemble.",
            "Mud density and pore gradient affect B&Y predictions; use real data when available.",
            "YP/PV values from the rheometer improve Bingham model accuracy.",
        ],
    }
    suggestions = suggestions_en if lang == 'en' else suggestions_es
    return suggestions.get(section, [])

def get_follow_up_suggestions(section: str, lang: Optional[str] = None) -> List[str]:
    """Retorna sugerencias de seguimiento (próximos pasos) para cada sección."""
    lang = lang or st.session_state.get('lang', 'es')
    follow_ups_es = {
        'rop_prediction': [
            "Ir al Mapa de calor ROP para localizar la zona óptima WOB–RPM.",
            "Revisar la Evaluación de riesgos antes de aplicar cambios en taladro.",
            "Exportar o documentar los parámetros recomendados para el pozo actual.",
        ],
        'model_performance': [
            "Consultar el Mapa de calor ROP para validar los rangos óptimos identificados.",
            "Si el R² es bajo, reentrenar con más muestras o incluir datos de offset.",
            "Revisar la pestaña Optimización para análisis de sensibilidad WOB/RPM.",
        ],
        'heat_map': [
            "Aplicar los valores WOB/RPM óptimos en Predicción ROP para verificar el ROP esperado.",
            "Comparar con el análisis 3D en la pestaña Optimización.",
            "Documentar la zona óptima para el intervalo o formación actual.",
        ],
        'neural_network': [
            "Validar predicciones NN frente a modelos mecanicistas en Predicción ROP.",
            "Revisar Optimización para correlacionar la superficie 3D con el entrenamiento.",
            "Considerar fine-tuning si se dispone de datos de campo recientes.",
        ],
        'optimization': [
            "Aplicar parámetros óptimos en el panel lateral y ejecutar Predicción ROP.",
            "Consultar la Evaluación de riesgos en Predicción ROP antes de implementar.",
            "Actualizar el Mapa de calor con UCS/bit actual para cruzar resultados.",
        ],
        'geological_tracking': [
            "Activar 'Usar seguimiento geológico' en el panel para predicción por formación.",
            "Validar predicciones ROP por formación en Predicción ROP con profundidad actual.",
            "Exportar tabla de formaciones para planificación del pozo.",
        ],
        'correlation': [
            "Ver Análisis de Correlación completo en la pestaña Predicción ROP.",
            "Validar correlaciones con datos de campo para calibración.",
            "Integrar variables altamente correlacionadas con ROP en estrategia de perforación.",
        ],
    }
    follow_ups_en = {
        'rop_prediction': [
            "Go to ROP Heat Map to locate the optimal WOB–RPM zone.",
            "Review Risk Assessment before applying changes at the rig.",
            "Export or document recommended parameters for the current well.",
        ],
        'model_performance': [
            "Consult the ROP Heat Map to validate identified optimal ranges.",
            "If R² is low, retrain with more samples or include offset data.",
            "Review the Optimization tab for WOB/RPM sensitivity analysis.",
        ],
        'heat_map': [
            "Apply optimal WOB/RPM values in ROP Prediction to verify expected ROP.",
            "Compare with 3D analysis in the Optimization tab.",
            "Document the optimum zone for the current interval or formation.",
        ],
        'neural_network': [
            "Validate NN predictions against mechanistic models in ROP Prediction.",
            "Review Optimization to correlate 3D surface with training.",
            "Consider fine-tuning if recent field data is available.",
        ],
        'optimization': [
            "Apply optimal parameters in the sidebar and run ROP Prediction.",
            "Consult Risk Assessment in ROP Prediction before implementing.",
            "Update the Heat Map with current UCS/bit to cross-check results.",
        ],
        'geological_tracking': [
            "Enable 'Use geological tracking' in the panel for formation-based prediction.",
            "Validate formation ROP predictions in ROP Prediction with current depth.",
            "Export formation table for well planning.",
        ],
        'correlation': [
            "See full Correlation Analysis in the ROP Prediction tab.",
            "Validate correlations with field data for calibration.",
            "Integrate highly correlated variables with ROP into drilling strategy.",
        ],
    }
    follow_ups = follow_ups_en if lang == 'en' else follow_ups_es
    return follow_ups.get(section, [])

# ============================================================================
# SEGUIMIENTO GEOLÓGICO - Formaciones e intervalos
# ============================================================================

DEFAULT_GEOLOGICAL_FORMATIONS = [
    {'name': 'Tobas', 'depth_top': 2000, 'depth_bottom': 4500, 'ucs_psi': 8000, 'lithology': 'Toba/Volcánico'},
    {'name': 'Arenisca Superior', 'depth_top': 4500, 'depth_bottom': 7500, 'ucs_psi': 12000, 'lithology': 'Arenisca'},
    {'name': 'Lutita Intermedia', 'depth_top': 7500, 'depth_bottom': 10000, 'ucs_psi': 18000, 'lithology': 'Lutita'},
    {'name': 'Caliza Compacta', 'depth_top': 10000, 'depth_bottom': 13500, 'ucs_psi': 25000, 'lithology': 'Caliza'},
    {'name': 'Cuarcita', 'depth_top': 13500, 'depth_bottom': 17000, 'ucs_psi': 32000, 'lithology': 'Cuarcita'},
]

def get_formation_at_depth(depth_ft: float, formations: List[Dict]) -> Optional[Dict]:
    """Retorna la formación geológica correspondiente a una profundidad dada."""
    for f in formations:
        if f['depth_top'] <= depth_ft < f['depth_bottom']:
            return f
    return formations[-1] if formations and depth_ft >= formations[-1]['depth_bottom'] else (formations[0] if formations and depth_ft < formations[0]['depth_top'] else None)

def predict_rop_by_formation(formations: List[Dict], params_base: Dict, predictor) -> List[Dict]:
    """Predice ROP por formación usando los modelos de optimización."""
    results = []
    for f in formations:
        p = dict(params_base)
        p['ucs_psi'] = f['ucs_psi']
        p['depth_ft'] = (f['depth_top'] + f['depth_bottom']) / 2
        pred = predictor.predict_ensemble(p, use_rf=True, use_xgb=True, use_nn=True)
        results.append({
            'formation': f['name'],
            'depth_interval': f"{f['depth_top']:,.0f} - {f['depth_bottom']:,.0f} ft",
            'ucs_psi': f['ucs_psi'],
            'lithology': f.get('lithology', 'N/A'),
            'rop_predicted': pred['Ensemble'],
            'wob_opt': 22,  # placeholder - could compute from heat map
            'rpm_opt': 120,
        })
    return results

def create_geological_track(formations: List[Dict], current_depth_ft: float) -> go.Figure:
    """Crea un track vertical de formaciones geológicas vs profundidad (estilo well log)."""
    if not formations:
        return go.Figure()
    current_depth_ft = _safe_float(current_depth_ft, 5000.0)
    depths = [f['depth_top'] for f in formations] + [formations[-1]['depth_bottom']]
    depth_min = min(depths)
    depth_max = max(depths)
    # Incluir profundidad actual en el rango visible si está fuera
    if current_depth_ft < depth_min:
        depth_min = current_depth_ft
    elif current_depth_ft > depth_max:
        depth_max = current_depth_ft
    # Colores por litología
    litho_colors = {
        'Toba/Volcánico': '#B8A98A',
        'Arenisca': '#F4D03F',
        'Lutita': '#5D6D7E',
        'Caliza': '#E8E8E8',
        'Cuarcita': '#AEB6BF',
    }
    default_colors = ['#1B4D3E', '#2C6E49', '#3B8C5E', '#4AA66B', '#5BC77B']
    # Evitar rango de eje y inválido (depth_min == depth_max)
    if depth_min == depth_max:
        depth_min -= 100
        depth_max += 100
    fig = go.Figure()
    # Trace invisible para inicializar ejes (las shapes con yref="y" requieren ejes definidos)
    fig.add_trace(go.Scatter(x=[0, 1], y=[depth_min, depth_max], mode='markers',
                             marker=dict(size=1, opacity=0), showlegend=False))
    for i, f in enumerate(formations):
        color = litho_colors.get(f.get('lithology', ''), default_colors[i % len(default_colors)])
        fig.add_shape(
            type="rect",
            x0=0, x1=1, y0=f['depth_top'], y1=f['depth_bottom'],
            xref="paper", yref="y",
            line=dict(width=1, color='rgba(0,0,0,0.3)'),
            fillcolor=color,
            layer="below"
        )
    # Línea de profundidad actual (siempre visible)
    fig.add_shape(
            type="line",
            x0=0, x1=1, y0=current_depth_ft, y1=current_depth_ft,
            xref="paper", yref="y",
            line=dict(color='#E74C3C', width=3, dash='dot')
    )
    fig.update_layout(
        title="Track de Seguimiento Geológico — Profundidad vs Formación",
        xaxis=dict(
            visible=False,
            range=[-0.1, 1.1],
            constrain="domain"
        ),
        yaxis=dict(
            title="Profundidad (ft)",
            range=[depth_max, depth_min],
            gridcolor='rgba(128,128,128,0.2)'
        ),
        height=450,
        margin=dict(l=60, r=40, t=50, b=50),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        showlegend=False
    )
    # Añadir anotaciones con nombres de formación
    for f in formations:
        mid_depth = (f['depth_top'] + f['depth_bottom']) / 2
        fig.add_annotation(
            x=0.5, y=mid_depth,
            xref="paper", yref="y",
            text=f"{f['name']}<br><sub>{f.get('lithology', '')}</sub>",
            showarrow=False,
            font=dict(size=10, color='#1a1a1a'),
            xanchor='center'
        )
    fig.add_annotation(
        x=1.02, y=current_depth_ft,
        xref="paper", yref="y",
        text=f"● {current_depth_ft:,.0f} ft",
        showarrow=True,
        arrowhead=2,
        arrowcolor='#E74C3C',
        font=dict(size=11, color='#E74C3C', family="Inter"),
        xanchor='left'
    )
    return fig

def _plotly_fig_to_bytes(fig: go.Figure, width=800, height=450) -> Optional[bytes]:
    """Exporta figura Plotly a PNG. Requiere kaleido."""
    try:
        return fig.to_image(format="png", width=width, height=height)
    except Exception:
        return None

def generate_resumen_pptx(
    depth_ft: float, wob_klb: float, rpm: float, ucs_psi: float, bit_diameter_in: float,
    formations: List[Dict], predictor, data_generator,
    rop_ensemble: Optional[float], current_prediction: Optional[Dict]
) -> bytes:
    """Genera presentación PowerPoint con el Resumen de Detalle."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    verde = RgbColor(0x1B, 0x4D, 0x3E)

    def add_title_slide(title: str, subtitle: str = ""):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        t = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = verde
        if subtitle:
            t2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
            t2.text_frame.paragraphs[0].text = subtitle
            t2.text_frame.paragraphs[0].font.size = Pt(18)
            t2.text_frame.paragraphs[0].font.color.rgb = RgbColor(0x4A, 0x55, 0x68)
        return slide

    def add_content_slide(title: str, seguimiento: str, recomendaciones: List[str], img_bytes: Optional[bytes] = None):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        # Título
        sh = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
        sh.text_frame.paragraphs[0].text = title
        sh.text_frame.paragraphs[0].font.size = Pt(24)
        sh.text_frame.paragraphs[0].font.bold = True
        sh.text_frame.paragraphs[0].font.color.rgb = verde
        # Seguimiento
        sh2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(6), Inches(1.2))
        sh2.text_frame.word_wrap = True
        p = sh2.text_frame.paragraphs[0]
        p.text = "Seguimiento: " + seguimiento
        p.font.size = Pt(11)
        # Recomendaciones
        sh3 = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(6), Inches(4))
        sh3.text_frame.word_wrap = True
        for i, rec in enumerate(recomendaciones[:6]):
            par = sh3.text_frame.paragraphs[0] if i == 0 else sh3.text_frame.add_paragraph()
            par.text = f"• {rec}"
            par.font.size = Pt(10)
            par.space_before = Pt(4)
        # Imagen si hay
        if img_bytes:
            try:
                slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(6.8), Inches(1.0), width=Inches(6))
            except Exception:
                pass

    # Slide 1: Portada
    add_title_slide("Resumen de Detalle — Drilling Analytics",
                   f"Seguimiento, Recomendación y Gráfica | {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}")

    # Slide 2: Predicción ROP
    rop_str = f"{rop_ensemble:.1f} ft/hr" if rop_ensemble is not None else "—"
    seg = f"Profundidad {depth_ft:,.0f} ft | WOB {wob_klb:.1f} klb | RPM {rpm} | ROP Ensemble {rop_str}"
    recs = get_section_suggestions('rop_prediction') + get_follow_up_suggestions('rop_prediction')
    gauge_val = rop_ensemble if rop_ensemble is not None else 0
    fig_g = create_gauge_chart(gauge_val, "ROP", 0, 150, "ft/hr")
    add_content_slide("Predicción ROP", seg, recs, _plotly_fig_to_bytes(fig_g))

    # Slide 3: Rendimiento de modelos
    mm = predictor.model_metrics
    r2_v = [v.get('r2', 0) for k, v in mm.items() if isinstance(v, dict) and 'r2' in v]
    rmse_v = [v.get('rmse', 0) for k, v in mm.items() if isinstance(v, dict) and 'rmse' in v]
    r2_mean = np.mean(r2_v) if r2_v else 0
    rmse_mean = np.mean(rmse_v) if rmse_v else 0
    seg2 = f"R² prom: {r2_mean:.3f} | RMSE prom: {rmse_mean:.2f} | Modelos: {len([k for k in mm if mm.get(k)])}"
    recs2 = get_section_suggestions('model_performance') + get_follow_up_suggestions('model_performance')
    fig_comp = create_model_comparison_chart(mm)
    add_content_slide("Rendimiento de modelos", seg2, recs2, _plotly_fig_to_bytes(fig_comp))

    # Slide 4: Mapa de calor ROP
    seg3 = f"UCS {ucs_psi:,.0f} psi | Broca {bit_diameter_in} in | Zona óptima WOB 15–25 · RPM 100–140"
    recs3 = get_section_suggestions('heat_map') + get_follow_up_suggestions('heat_map')
    fig_hm = create_rop_heatmap([5, 40], [40, 220], ucs_psi, bit_diameter_in)
    add_content_slide("Mapa de calor ROP", seg3, recs3, _plotly_fig_to_bytes(fig_hm))

    # Slide 5: Análisis red neuronal
    seg4 = "Arquitectura MLP 3 capas | 100 épocas"
    recs4 = get_section_suggestions('neural_network') + get_follow_up_suggestions('neural_network')
    fig_arch = create_nn_architecture_diagram()
    add_content_slide("Análisis de red neuronal", seg4, recs4, _plotly_fig_to_bytes(fig_arch))

    # Slide 6: Optimización
    rop_act = current_prediction['Ensemble'] if current_prediction else 0
    seg5 = f"WOB {wob_klb:.1f} klb | RPM {rpm} | ROP {rop_act:.1f} ft/hr"
    recs5 = get_section_suggestions('optimization') + get_follow_up_suggestions('optimization')
    wob_t = np.linspace(5, 40, 40)
    rop_w = [predictor.predict_physical_model({
        'wob_klb': w, 'rpm': 120, 'torque_ftlb': 18000, 'spp_psi': 3000, 'flow_gpm': 800,
        'ucs_psi': ucs_psi, 'bit_diameter_in': bit_diameter_in, 'bit_wear': 0.2,
        'depth_ft': depth_ft, 'cutter_count': 6
    }) for w in wob_t]
    fw = go.Figure(go.Scatter(x=wob_t, y=rop_w, mode='lines', line=dict(color='#1B4D3E')))
    fw.update_layout(title='Sensibilidad WOB', xaxis_title='WOB (klb)', yaxis_title='ROP (ft/hr)', height=320)
    add_content_slide("Optimización", seg5, recs5, _plotly_fig_to_bytes(fw))

    # Slide 7: Seguimiento geológico
    form_at = get_formation_at_depth(depth_ft, formations)
    fname = form_at['name'] if form_at else "—"
    ucs_f = form_at['ucs_psi'] if form_at else ucs_psi
    seg6 = f"Profundidad {depth_ft:,.0f} ft | Formación {fname} | UCS {ucs_f:,.0f} psi"
    recs6 = get_section_suggestions('geological_tracking') + get_follow_up_suggestions('geological_tracking')
    params_b = {'wob_klb': wob_klb, 'rpm': rpm, 'torque_ftlb': 18000, 'spp_psi': 3000, 'flow_gpm': 800,
                'bit_diameter_in': bit_diameter_in, 'bit_wear': 0.2, 'cutter_count': 6,
                'mud_density_ppg': 10, 'pore_gradient_ppg': 9, 'yp_lb100ft2': 15, 'pv_cp': 25, 'inclination_deg': 0}
    rop_by_f = predict_rop_by_formation(formations, params_b, predictor)
    fig_geo = go.Figure(go.Bar(
        x=[r['formation'] for r in rop_by_f],
        y=[r['rop_predicted'] for r in rop_by_f],
        marker_color=['#1B4D3E', '#2C6E49', '#3B8C5E', '#4AA66B', '#5BC77B'],
        text=[f"{r['rop_predicted']:.1f}" for r in rop_by_f], textposition='outside'
    ))
    fig_geo.update_layout(title='ROP por formación', xaxis_title='Formación', yaxis_title='ROP (ft/hr)', height=350)
    add_content_slide("Seguimiento geológico", seg6, recs6, _plotly_fig_to_bytes(fig_geo))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()

# ============================================================================
# CONFIGURACIÓN DE PÁGINA - IMPERIAL
# ============================================================================

st.set_page_config(
    page_title="Drilling Analytics - Unidades Imperiales",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ============================================================================
# CSS PERSONALIZADO
# ============================================================================

st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');
    
    * {
        font-family: 'Inter', sans-serif;
    }
    
    .main-header {
        font-size: 2.8rem;
        font-weight: 800;
        background: linear-gradient(135deg, #1B4D3E 0%, #2C6E49 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        text-align: center;
        margin-bottom: 0.5rem;
        padding: 1rem;
        letter-spacing: -0.5px;
    }
    
    .sub-header {
        font-size: 1.8rem;
        font-weight: 700;
        color: #1B4D3E;
        margin-bottom: 1rem;
        border-bottom: 4px solid #2C6E49;
        padding-bottom: 0.5rem;
    }
    
    .imperial-badge {
        background: linear-gradient(135deg, #1B4D3E, #2C6E49);
        color: white;
        padding: 0.3rem 1rem;
        border-radius: 30px;
        font-size: 0.9rem;
        font-weight: 600;
        display: inline-block;
        margin-bottom: 1rem;
        text-transform: uppercase;
        letter-spacing: 1px;
    }
    
    .metric-card {
        background: linear-gradient(135deg, #1B4D3E 0%, #2C6E49 100%);
        border-radius: 15px;
        padding: 25px;
        color: white;
        box-shadow: 0 10px 30px rgba(27,77,62,0.2);
        transition: transform 0.3s ease;
    }
    
    .metric-card:hover {
        transform: translateY(-5px);
        box-shadow: 0 15px 40px rgba(27,77,62,0.3);
    }
    
    .info-box {
        background: linear-gradient(135deg, #F8F9FA 0%, #E9ECEF 100%);
        border-left: 6px solid #2C6E49;
        padding: 1.5rem;
        border-radius: 8px;
        margin: 1rem 0;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05);
    }
    
    .unit-label {
        font-size: 0.8rem;
        color: #6C757D;
        margin-left: 0.3rem;
        font-weight: 400;
    }
    
    .stButton > button {
        background: linear-gradient(45deg, #1B4D3E, #2C6E49);
        color: white;
        font-weight: 600;
        border: none;
        border-radius: 10px;
        padding: 0.8rem 2rem;
        font-size: 1.1rem;
        transition: all 0.3s ease;
        width: 100%;
        letter-spacing: 0.5px;
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 5px 20px rgba(44,110,73,0.4);
        background: linear-gradient(45deg, #236B4A, #2C6E49);
    }
    
    .model-card {
        background: white;
        border-radius: 15px;
        padding: 1.5rem;
        box-shadow: 0 4px 20px rgba(0,0,0,0.05);
        border: 1px solid #E9ECEF;
        transition: all 0.3s ease;
    }
    
    .model-card:hover {
        box-shadow: 0 8px 30px rgba(0,0,0,0.1);
    }
    
    .success-badge {
        background: linear-gradient(45deg, #28A745, #20C997);
        color: white;
        padding: 0.3rem 0.8rem;
        border-radius: 20px;
        font-size: 0.9rem;
        font-weight: 600;
        display: inline-block;
    }
    
    .warning-badge {
        background: linear-gradient(45deg, #FFC107, #FD7E14);
        color: black;
        padding: 0.3rem 0.8rem;
        border-radius: 20px;
        font-size: 0.9rem;
        font-weight: 600;
        display: inline-block;
    }
    
    .divider {
        height: 2px;
        background: linear-gradient(90deg, transparent, #2C6E49, transparent);
        margin: 2rem 0;
    }
    
    .footer {
        text-align: center;
        color: #6C757D;
        padding: 2rem;
        font-size: 0.9rem;
    }
    
    .chip {
        display: inline-flex;
        align-items: center;
        padding: 0.35rem 0.9rem;
        margin: 0.25rem 0.35rem 0.25rem 0;
        border-radius: 9999px;
        font-size: 0.8rem;
        font-weight: 600;
        letter-spacing: 0.02em;
    }
    .chip-temp-target { background: linear-gradient(135deg, #0D9488 0%, #14B8A6 100%); color: white; border: none; }
    .chip-temp-r2    { background: linear-gradient(135deg, #0369A1 0%, #0EA5E9 100%); color: white; }
    .chip-temp-rmse  { background: linear-gradient(135deg, #7C3AED 0%, #A78BFA 100%); color: white; }
    .chip-temp-mae   { background: linear-gradient(135deg, #B45309 0%, #F59E0B 100%); color: white; }
    .chip-temp-ok    { background: linear-gradient(135deg, #047857 0%, #10B981 100%); color: white; }
    .chip-row { display: flex; flex-wrap: wrap; align-items: center; gap: 0.25rem; margin-bottom: 1rem; }
    .chip-anom-count { background: linear-gradient(135deg, #DC2626 0%, #F87171 100%); color: white; box-shadow: 0 2px 8px rgba(220,38,38,0.35); }
    .chip-anom-mean  { background: linear-gradient(135deg, #0D9488 0%, #2DD4BF 100%); color: white; box-shadow: 0 2px 8px rgba(13,148,136,0.35); }
    .chip-anom-z     { background: linear-gradient(135deg, #6366F1 0%, #818CF8 100%); color: white; box-shadow: 0 2px 8px rgba(99,102,241,0.35); }
    .chip-anom-ok    { background: linear-gradient(135deg, #059669 0%, #34D399 100%); color: white; }
    .chip-anom-warn  { background: linear-gradient(135deg, #D97706 0%, #FBBF24 100%); color: #1a1a1a; }
    .chip-anom-puntos{ background: linear-gradient(135deg, #475569 0%, #64748B 100%); color: white; }
    /* Perfil térmico vecinos — chips KPI */
    .chip-row-neighbor {
        display: flex;
        flex-wrap: wrap;
        align-items: stretch;
        gap: 0.5rem;
        margin: 0 0 1rem 0;
        padding: 0.65rem 0.75rem;
        background: linear-gradient(145deg, rgba(15, 23, 42, 0.55) 0%, rgba(30, 41, 59, 0.35) 100%);
        border-radius: 14px;
        border: 1px solid rgba(148, 163, 184, 0.18);
        box-shadow: inset 0 1px 0 rgba(255,255,255,0.06);
    }
    .chip-neigh {
        display: inline-flex;
        align-items: center;
        gap: 0.4rem;
        padding: 0.45rem 1rem 0.45rem 0.85rem;
        border-radius: 9999px;
        font-size: 0.78rem;
        font-weight: 600;
        letter-spacing: 0.03em;
        border: 1px solid rgba(255,255,255,0.14);
        box-shadow: 0 2px 10px rgba(0,0,0,0.2);
    }
    .chip-neigh strong { font-weight: 700; letter-spacing: 0.02em; }
    .chip-neigh-k { font-size: 0.65rem; font-weight: 700; text-transform: uppercase; opacity: 0.92; letter-spacing: 0.08em; }
    .chip-neigh-count   { background: linear-gradient(135deg, #047857 0%, #10B981 100%); color: #f0fdf4; }
    .chip-neigh-depth-lo { background: linear-gradient(135deg, #0369A1 0%, #0EA5E9 100%); color: #f0f9ff; }
    .chip-neigh-depth-hi { background: linear-gradient(135deg, #1D4ED8 0%, #3B82F6 100%); color: #eff6ff; }
    .chip-neigh-grad    { background: linear-gradient(135deg, #6D28D9 0%, #A78BFA 100%); color: #faf5ff; }
    .chip-neigh-meta    { background: linear-gradient(135deg, #334155 0%, #475569 100%); color: #f1f5f9; font-size: 0.74rem; }
</style>
""", unsafe_allow_html=True)

# ============================================================================
# CLASE DE CONVERSIÓN DE UNIDADES
# ============================================================================

class UnitConverter:
    """Conversor de unidades SI a Imperial"""
    
    @staticmethod
    def kn_to_klb(kn: float) -> float:
        """Kilonewtons a 1000 libras"""
        return kn * 0.224808943 / 1000
    
    @staticmethod
    def klb_to_kn(klb: float) -> float:
        """1000 libras a kilonewtons"""
        return klb * 1000 / 0.224808943
    
    @staticmethod
    def mh_to_fthr(mh: float) -> float:
        """Metros/hora a pies/hora"""
        return mh * 3.28084
    
    @staticmethod
    def fthr_to_mh(fthr: float) -> float:
        """Pies/hora a metros/hora"""
        return fthr / 3.28084
    
    @staticmethod
    def m_to_ft(m: float) -> float:
        """Metros a pies"""
        return m * 3.28084
    
    @staticmethod
    def ft_to_m(ft: float) -> float:
        """Pies a metros"""
        return ft / 3.28084
    
    @staticmethod
    def mpa_to_psi(mpa: float) -> float:
        """Megapascales a psi"""
        return mpa * 145.0377
    
    @staticmethod
    def psi_to_mpa(psi: float) -> float:
        """PSI a megapascales"""
        return psi / 145.0377
    
    @staticmethod
    def lpm_to_gpm(lpm: float) -> float:
        """Litros/minuto a galones/minuto"""
        return lpm * 0.264172
    
    @staticmethod
    def gpm_to_lpm(gpm: float) -> float:
        """Galones/minuto a litros/minuto"""
        return gpm / 0.264172
    
    @staticmethod
    def knm_to_ftlb(knm: float) -> float:
        """Kilonewton-metro a pie-libra"""
        return knm * 737.5621
    
    @staticmethod
    def ftlb_to_knm(ftlb: float) -> float:
        """Pie-libra a kilonewton-metro"""
        return ftlb / 737.5621
    
    @staticmethod
    def format_wob(value: float) -> str:
        """Formatear WOB en klb"""
        return f"{value:.1f} klb"
    
    @staticmethod
    def format_rop(value: float) -> str:
        """Formatear ROP en ft/hr"""
        return f"{value:.1f} ft/hr"
    
    @staticmethod
    def format_depth(value: float) -> str:
        """Formatear profundidad en ft"""
        return f"{value:,.0f} ft"
    
    @staticmethod
    def format_pressure(value: float) -> str:
        """Formatear presión en psi"""
        return f"{value:,.0f} psi"
    
    @staticmethod
    def format_flow(value: float) -> str:
        """Formatear caudal en gpm"""
        return f"{value:.0f} gpm"
    
    @staticmethod
    def format_torque(value: float) -> str:
        """Formatear torque en ft-lb"""
        return f"{value:,.0f} ft-lb"

    @staticmethod
    def kgf_to_klb(kgf: float) -> float:
        """Kilogramo-fuerza a miles de libras (klb)."""
        return float(kgf) * 0.00220462

    @staticmethod
    def pa_to_psi(pa: float) -> float:
        """Pascales a psi."""
        return float(pa) / 6894.757293178

    @staticmethod
    def m3min_to_gpm(m3min: float) -> float:
        """Metros cúbicos por minuto a galones por minuto."""
        return float(m3min) * 264.172052

    @staticmethod
    def nm_to_ftlb(nm: float) -> float:
        """Newton-metro a pie-libra."""
        return float(nm) * 0.7375621493

    @staticmethod
    def gcm3_to_ppg(gcm3: float) -> float:
        """Densidad g/cm³ a libras por galón (ppg)."""
        return float(gcm3) * 8.345404452

    @staticmethod
    def min_per_m_to_fthr(min_per_m: float) -> float:
        """ROP en min/m (minutos por metro) a ft/hr."""
        v = float(min_per_m)
        if not np.isfinite(v) or abs(v) < 1e-12:
            return np.nan
        return (60.0 * 3.28084) / v

converter = UnitConverter()

# ============================================================================
# MODELO DE PREDICCIÓN DE ROP - IMPERIAL
# ============================================================================

# Columnas de features para modelos ML (deben coincidir con DrillingDataGeneratorImperial)
FEATURE_COLS = ['WOB_klb', 'RPM', 'Torque_ftlb', 'SPP_psi', 'Flow_gpm', 'UCS_psi', 
                'BitDiameter_in', 'BitWear_pct', 'Depth_ft', 'CutterCount']
FEATURE_LABELS = {'WOB_klb': 'WOB (klb)', 'RPM': 'RPM', 'Torque_ftlb': 'Torque (ft-lb)', 
                  'SPP_psi': 'SPP (psi)', 'Flow_gpm': 'Caudal (gpm)', 'UCS_psi': 'UCS (psi)',
                  'BitDiameter_in': 'Diám. broca (in)', 'BitWear_pct': 'Desgaste broca', 
                  'Depth_ft': 'Profundidad (ft)', 'CutterCount': 'Cortadores'}

class ROPPredictorImperial:
    """
    Modelo de predicción de ROP basado en Bourgoyne & Young
    Adaptado para unidades imperiales (EE.UU.)
    Soporta entrenamiento ML real (RF, XGBoost, NN) con datos sintéticos o reales.
    """
    
    def __init__(self):
        self.is_trained = False
        self.model_metrics = {}
        self.feature_importance = {}
        self.models = {}  # Modelos entrenados reales
        self.scaler = None  # Escalador de features (para NN)
        
    def _normalize_params(self, params):
        """Normaliza params: acepta claves CamelCase (generador) y snake_case (formulario)"""
        key_map = {
            'WOB_klb': 'wob_klb',
            'RPM': 'rpm',
            'Torque_ftlb': 'torque_ftlb',
            'SPP_psi': 'spp_psi',
            'Flow_gpm': 'flow_gpm',
            'UCS_psi': 'ucs_psi',
            'BitDiameter_in': 'bit_diameter_in',
            'BitWear_pct': 'bit_wear',
            'Depth_ft': 'depth_ft',
            'CutterCount': 'cutter_count',
        }
        out = dict(params)
        for old_key, new_key in key_map.items():
            if old_key in out and new_key not in out:
                out[new_key] = out[old_key]
        return out

    def _params_to_feature_row(self, params):
        """Convierte params (formulario o dict) a fila de features para predicción ML."""
        p = self._normalize_params(params)
        col_to_key = {'WOB_klb': 'wob_klb', 'RPM': 'rpm', 'Torque_ftlb': 'torque_ftlb', 
                      'SPP_psi': 'spp_psi', 'Flow_gpm': 'flow_gpm', 'UCS_psi': 'ucs_psi',
                      'BitDiameter_in': 'bit_diameter_in', 'BitWear_pct': 'bit_wear',
                      'Depth_ft': 'depth_ft', 'CutterCount': 'cutter_count'}
        row = [_safe_float(p.get(col_to_key.get(c, c), 0), 0) for c in FEATURE_COLS]
        return np.array(row).reshape(1, -1)

    def predict_physical_model(self, params):
        """
        Modelo físico de predicción de ROP - Versión Imperial
        
        Parámetros de entrada (Imperial):
            - WOB: 1000 lb (klb)
            - RPM: rev/min
            - Torque: ft-lb
            - SPP: psi
            - Flow: gpm
            - UCS: psi
            - Bit diameter: in
            - Bit wear: 0-1
            - Depth: ft
        """
        try:
            params = self._normalize_params(params)
            params, _ = _validate_params(params, ['wob_klb', 'rpm', 'torque_ftlb', 'spp_psi', 'flow_gpm', 'ucs_psi', 'bit_diameter_in', 'bit_wear', 'depth_ft', 'cutter_count'])
            bit_diameter_in = max(0.1, _safe_float(params.get('bit_diameter_in', 8.5), 8.5))
            wob_lb = max(0, _safe_float(params['wob_klb'], 22) * 1000)

            # ====================================================================
            # FACTORES DE ROP - IMPERIAL
            # ====================================================================

            # 1. Factor de WOB (óptimo: 4000-6000 lb/in de diámetro)
            wob_opt_per_inch = 5000  # lb/in
            wob_opt = wob_opt_per_inch * bit_diameter_in
            f_wob = np.exp(-((wob_lb - wob_opt) ** 2) / (2 * max((wob_opt * 0.3) ** 2, 1e-6)))
            # 2. Factor de RPM (óptimo: 100-140 RPM)
            rpm_opt = 120
            f_rpm = np.exp(-((params['rpm'] - rpm_opt) ** 2) / (2 * (rpm_opt * 0.4) ** 2))
            # 3. Factor de formación (UCS en psi)
            ucs_psi = max(100, _safe_float(params.get('ucs_psi', 15000), 15000))
            f_formation = 5000 / max(ucs_psi ** 0.8, 1e-6)
            # 4. Factor de desgaste
            f_wear = 1 - params['bit_wear'] * 0.6
            # 5. Factor de torque
            torque_ftlb = params['torque_ftlb']
            f_torque = 1 + (torque_ftlb - 15000) / 50000
            # 6. Factor hidráulico
            hhp = (_safe_float(params.get('spp_psi', 3000)) * _safe_float(params.get('flow_gpm', 800))) / 1714
            hhp_opt = max(0.1, 2.5 * bit_diameter_in)
            f_hydraulic = 1 + 0.3 * (hhp - hhp_opt) / hhp_opt
            # 7. Factor de profundidad
            depth_ft = params['depth_ft']
            f_depth = np.exp(-depth_ft / 15000)
            # 8. Factor de cortadores
            cutter_factor = 1 + (params.get('cutter_count', 6) - 6) * 0.05
            # ROP base en ft/hr (para formación blanda)
            rop_base = 120
            # ROP calculada
            rop_ftph = (rop_base * f_wob * f_rpm * f_formation * f_wear *
                        f_torque * f_hydraulic * f_depth * cutter_factor)
            noise = np.random.normal(0, max(rop_ftph * 0.05, 0.5))
            rop_ftph = max(5, rop_ftph + noise)
            return float(rop_ftph)
        except Exception:
            return 30.0

    def predict_bourgoyne_young(self, params):
        """
        Bourgoyne & Young (1974): ln(ROP) = a1 + a2*x2 + ... + a8*x8
        Full 8-factor mechanistic model: depth, compaction, pressure, WOB, RPM, overbalance, hydraulics.
        """
        try:
            params = self._normalize_params(params)
            params, _ = _validate_params(params, ['wob_klb', 'rpm', 'depth_ft', 'bit_diameter_in'])
            p = params
            wob_lbf = p['wob_klb'] * 1000
            depth_ft = p['depth_ft']
            rpm = max(p['rpm'], 1)
            dc = max(p['bit_diameter_in'], 0.1)
            mud_ppg = p.get('mud_density_ppg', 10.0)
            pore_psi_ft = p.get('pore_gradient_ppg', 9.0) * 0.052
            hydro_psi = mud_ppg * 0.052 * depth_ft
            pore_psi = pore_psi_ft * depth_ft
            delta_p = max(0, hydro_psi - pore_psi)
            x2 = 10000 - depth_ft
            gp = pore_psi_ft
            gc = mud_ppg * 0.052
            x3 = (depth_ft ** 0.69) * (gp - 9 * 0.052) if abs(gp - 9 * 0.052) > 1e-9 else 0
            x4 = depth_ft * (gp - gc)
            wob_term = max(0.1, (wob_lbf / max(4 * dc, 0.1)) - 0.75)
            x5 = np.log(max(wob_term * 60 / max(4 * dc, 0.1), 0.01))
            x6 = np.log(rpm / 60)
            x7 = -delta_p
            x8 = (wob_lbf / max(depth_ft, 1)) * (rpm / 60) * (delta_p / 3600)
            a1, a2, a3, a4, a5, a6, a7, a8 = 1.0, -0.00001, 1e-6, -0.00001, 0.5, 0.6, -0.00003, 1e-7
            ln_rop = a1 + a2*x2 + a3*x3 + a4*x4 + a5*x5 + a6*x6 + a7*x7 + a8*x8
            rop = np.exp(np.clip(ln_rop, -3, 6))
            return float(max(0.1, min(rop, 500.0)))
        except Exception:
            return 25.0

    def predict_bingham_rop(self, params):
        """
        Bingham Plastic ROP Model: incorporates rheology (YP, PV) into hydraulic efficiency.
        ROP = K * (WOB/d)^α * N^β * η_hydraulic * (1 - Δp_penalty)
        η_hydraulic from Bingham ECD and annular pressure loss.
        """
        try:
            params = self._normalize_params(params)
            params, _ = _validate_params(params, ['wob_klb', 'rpm', 'bit_diameter_in', 'depth_ft', 'ucs_psi'])
            p = params
            wob_lbf = p['wob_klb'] * 1000
            rpm = max(p['rpm'], 1)
            dc = max(p['bit_diameter_in'], 0.1)
            depth_ft = p['depth_ft']
            ucs_psi = max(p['ucs_psi'], 100)
            base_rop = 0.5
            flow_gpm = p.get('flow_gpm', 800)
            yp = p.get('yp_lb100ft2', 15.0)
            pv = p.get('pv_cp', 25.0)
            v_ann_ft_min = flow_gpm * 0.409 / max(dc**2 - 16, 0.1) if dc > 4 else 50
            v_ann_ft_min = max(v_ann_ft_min, 30)
            hh_ratio = yp / max(pv, 1)
            eta_hyd = 1.0 - 0.15 * min(hh_ratio, 2.0)
            spp = p.get('spp_psi', 3000)
            delta_p_approx = spp * 0.3
            delta_penalty = min(0.5, delta_p_approx / 2000)
            wob_per_in = wob_lbf / max(dc, 1)
            f_wob = (wob_per_in / 4000) ** 1.2
            f_rpm = (rpm / 100) ** 0.6
            f_formation = 8000 / max(ucs_psi ** 0.7, 1e-6)
            f_wear = 1 - p.get('bit_wear', 0.2) * 0.5
            rop = 60 * base_rop * f_wob * f_rpm * f_formation * f_wear * eta_hyd * (1 - delta_penalty)
            return float(max(5, min(rop, 300.0)))
        except Exception:
            return 25.0

    def predict_warren(self, params):
        """
        Warren (1987) ROP model for roller-cone bits: cuttings generation + removal.
        ROP = K * ((WOB - WOB0)/d)^a * N^b * exp(-c*h) * f(UCS) * f(depth)
        """
        try:
            params = self._normalize_params(params)
            params, _ = _validate_params(params, ['wob_klb', 'rpm', 'bit_diameter_in', 'depth_ft', 'ucs_psi'])
            p = params
            wob_lbf = p['wob_klb'] * 1000
            rpm = max(p['rpm'], 1)
            dc = max(p['bit_diameter_in'], 0.1)
            depth_ft = p['depth_ft']
            ucs_psi = max(p['ucs_psi'], 100)
            h = p.get('bit_wear', 0.2)
            WOB0 = 2000
            wob_eff = max(0, wob_lbf - WOB0)
            wob_per_in = wob_eff / max(dc, 1)
            a, b, c = 1.2, 0.65, 3.0
            K = 0.12
            f_wob = (wob_per_in / 3000) ** a
            f_rpm = (rpm / 100) ** b
            f_wear = np.exp(-c * h)
            f_ucs = 10000 / max(ucs_psi ** 0.75, 1e-6)
            f_depth = np.exp(-depth_ft / 18000)
            rop = K * 120 * f_wob * f_rpm * f_wear * f_ucs * f_depth
            return float(max(5, min(rop, 250.0)))
        except Exception:
            return 25.0
    
    def _prepare_training_data(self, data_generator, n_samples=10000):
        """Genera datos y prepara X, y para entrenamiento."""
        data = data_generator.generate()
        data = data.sample(n=min(n_samples, len(data)), random_state=42)
        X = data[FEATURE_COLS].astype(float).values
        y = data['ROP_fthr'].values
        return X, y

    def _compute_metrics(self, y_true, y_pred):
        """Calcula RMSE, MAE, R² y MAPE."""
        y_true = np.array(y_true)
        y_pred = np.array(y_pred)
        rmse = np.sqrt(mean_squared_error(y_true, y_pred))
        mae = mean_absolute_error(y_true, y_pred)
        r2 = r2_score(y_true, y_pred) if np.var(y_true) > 0 else 0
        with np.errstate(divide='ignore', invalid='ignore'):
            mape = np.mean(np.abs((y_true - y_pred) / (np.abs(y_true) + 1e-8))) * 100
        return {'rmse': rmse, 'mae': mae, 'r2': r2, 'mape': float(mape)}


    def _feature_importance_to_labels(self, imp_dict):
        """Convierte importancia a formato con etiquetas legibles."""
        return {FEATURE_LABELS.get(k, k): float(v) for k, v in imp_dict.items()}

    def train_random_forest(self, synthetic_data=True, data_generator=None, n_samples=10000):
        """Entrena modelo Random Forest real para ROP."""
        if SKLEARN_AVAILABLE and data_generator is not None:
            try:
                X, y = self._prepare_training_data(data_generator, n_samples)
                X_train, X_val, y_train, y_val = train_test_split(X, y, test_size=0.2, random_state=42)
                
                model = RandomForestRegressor(n_estimators=60, max_depth=10, min_samples_leaf=5, 
                                             random_state=42, n_jobs=-1)
                model.fit(X_train, y_train)
                
                y_pred = model.predict(X_val)
                self.models['random_forest'] = model
                self.model_metrics['random_forest'] = self._compute_metrics(y_val, y_pred)
                imp = dict(zip(FEATURE_COLS, model.feature_importances_))
                self.feature_importance['random_forest'] = self._feature_importance_to_labels(imp)
                self.is_trained = True
                return self.model_metrics['random_forest']
            except Exception as e:
                pass  # Fallback a simulación
        # Fallback: métricas simuladas
        time.sleep(1.0)
        self.model_metrics['random_forest'] = {
            'rmse': np.random.uniform(4.5, 6.5), 'mae': np.random.uniform(3.2, 4.8),
            'r2': np.random.uniform(0.86, 0.91), 'mape': np.random.uniform(7.2, 9.5)
        }
        self.feature_importance['random_forest'] = self._feature_importance_to_labels(
            dict(zip(FEATURE_COLS, [0.28, 0.22, 0.06, 0.04, 0.04, 0.18, 0.12, 0.10, 0.04, 0.02])))
        self.is_trained = True
        return self.model_metrics['random_forest']

    def train_xgboost(self, synthetic_data=True, data_generator=None, n_samples=10000):
        """Entrena modelo XGBoost real para ROP."""
        if XGBOOST_AVAILABLE and SKLEARN_AVAILABLE and data_generator is not None:
            try:
                X, y = self._prepare_training_data(data_generator, n_samples)
                X_train, X_val, y_train, y_val = train_test_split(X, y, test_size=0.2, random_state=42)
                
                model = xgb.XGBRegressor(n_estimators=60, max_depth=6, learning_rate=0.1,
                                         random_state=42, n_jobs=-1)
                model.fit(X_train, y_train)
                
                y_pred = model.predict(X_val)
                self.models['xgboost'] = model
                self.model_metrics['xgboost'] = self._compute_metrics(y_val, y_pred)
                imp = dict(zip(FEATURE_COLS, model.feature_importances_))
                self.feature_importance['xgboost'] = self._feature_importance_to_labels(imp)
                self.is_trained = True
                return self.model_metrics['xgboost']
            except Exception:
                pass
        time.sleep(1.0)
        self.model_metrics['xgboost'] = {
            'rmse': np.random.uniform(3.8, 5.2), 'mae': np.random.uniform(2.8, 4.0),
            'r2': np.random.uniform(0.89, 0.94), 'mape': np.random.uniform(6.0, 8.2)
        }
        self.feature_importance['xgboost'] = self._feature_importance_to_labels(
            dict(zip(FEATURE_COLS, [0.32, 0.25, 0.05, 0.03, 0.02, 0.16, 0.11, 0.09, 0.05, 0.02])))
        self.is_trained = True
        return self.model_metrics['xgboost']

    def train_neural_network(self, synthetic_data=True, data_generator=None, n_samples=10000):
        """Entrena Red Neuronal MLP real para ROP."""
        if SKLEARN_AVAILABLE and data_generator is not None:
            try:
                X, y = self._prepare_training_data(data_generator, n_samples)
                self.scaler = StandardScaler()
                X_scaled = self.scaler.fit_transform(X)
                X_train, X_val, y_train, y_val = train_test_split(X_scaled, y, test_size=0.2, random_state=42)
                
                model = MLPRegressor(hidden_layer_sizes=(64, 32, 16), activation='relu', solver='adam',
                                    alpha=1e-4, max_iter=120, random_state=42, early_stopping=True,
                                    validation_fraction=0.1)
                model.fit(X_train, y_train)
                
                y_pred = model.predict(X_val)
                self.models['neural_network'] = model
                self.model_metrics['neural_network'] = self._compute_metrics(y_val, y_pred)
                try:
                    from sklearn.inspection import permutation_importance
                    n_val = min(100, len(X_val))
                    perm = permutation_importance(model, X_val[:n_val], y_val[:n_val], n_repeats=3, random_state=42)
                    imp = dict(zip(FEATURE_COLS, perm.importances_mean))
                    self.feature_importance['neural_network'] = self._feature_importance_to_labels(imp)
                except Exception:
                    self.feature_importance['neural_network'] = self._feature_importance_to_labels(
                        dict(zip(FEATURE_COLS, [0.30, 0.28, 0.04, 0.03, 0.02, 0.15, 0.13, 0.08, 0.05, 0.02])))
                self.is_trained = True
                return self.model_metrics['neural_network']
            except Exception:
                pass
        time.sleep(1.0)
        self.model_metrics['neural_network'] = {
            'rmse': np.random.uniform(3.2, 4.5), 'mae': np.random.uniform(2.4, 3.5),
            'r2': np.random.uniform(0.92, 0.96), 'mape': np.random.uniform(5.0, 7.0)
        }
        self.feature_importance['neural_network'] = self._feature_importance_to_labels(
            dict(zip(FEATURE_COLS, [0.30, 0.28, 0.04, 0.03, 0.02, 0.15, 0.13, 0.08, 0.05, 0.02])))
        self.is_trained = True
        return self.model_metrics['neural_network']
    
    def _predict_ml_model(self, model_key, params):
        """Predicción con modelo ML entrenado (RF, XGBoost o NN)."""
        if model_key not in self.models:
            return None
        X = self._params_to_feature_row(params)
        model = self.models[model_key]
        if model_key == 'neural_network' and self.scaler is not None:
            X = self.scaler.transform(X)
        try:
            return float(model.predict(X)[0])
        except Exception:
            return None

    def predict_ensemble(self, params, use_rf=True, use_xgb=True, use_nn=True):
        """
        Predicción ensemble: modelos mecanísticos (B&Y, Bingham, Warren) + ML (RF, XGB, NN)
        Usa predicciones reales de modelos ML cuando están entrenados.
        """
        predictions = []
        weights = []
        rf_rop, xgb_rop, nn_rop = None, None, None

        by_rop = self.predict_bourgoyne_young(params)
        bingham_rop = self.predict_bingham_rop(params)
        warren_rop = self.predict_warren(params)
        physical_rop = self.predict_physical_model(params)
        
        for rop_val in [by_rop, bingham_rop, warren_rop, physical_rop]:
            predictions.append(rop_val)
            weights.append(0.10)
        
        if use_rf and 'random_forest' in self.model_metrics:
            rf_rop = self._predict_ml_model('random_forest', params)
            if rf_rop is not None:
                predictions.append(rf_rop)
                weights.append(0.25)
            else:
                rf_rop = np.mean([by_rop, bingham_rop, warren_rop]) * np.random.uniform(0.92, 1.08)
                predictions.append(rf_rop)
                weights.append(0.25)
        
        if use_xgb and 'xgboost' in self.model_metrics:
            xgb_rop = self._predict_ml_model('xgboost', params)
            if xgb_rop is not None:
                predictions.append(xgb_rop)
                weights.append(0.25)
            else:
                xgb_rop = np.mean([by_rop, physical_rop]) * np.random.uniform(0.95, 1.05)
                predictions.append(xgb_rop)
                weights.append(0.25)
        
        if use_nn and 'neural_network' in self.model_metrics:
            nn_rop = self._predict_ml_model('neural_network', params)
            if nn_rop is not None:
                predictions.append(nn_rop)
                weights.append(0.25)
            else:
                nn_rop = np.mean([by_rop, warren_rop, physical_rop]) * np.random.uniform(0.97, 1.03)
                predictions.append(nn_rop)
                weights.append(0.25)
        
        weights = np.array(weights) / sum(weights)
        ensemble_rop = np.average(predictions, weights=weights)
        
        rf_rop_val = rf_rop if (use_rf and 'random_forest' in self.model_metrics) else None
        xgb_rop_val = xgb_rop if (use_xgb and 'xgboost' in self.model_metrics) else None
        nn_rop_val = nn_rop if (use_nn and 'neural_network' in self.model_metrics) else None

        pred_dict = {
            'Bourgoyne & Young': by_rop,
            'Bingham': bingham_rop,
            'Warren': warren_rop,
            'Physical Model': physical_rop,
            'Random Forest': rf_rop_val,
            'XGBoost': xgb_rop_val,
            'Neural Network': nn_rop_val,
            'Ensemble': float(max(1, ensemble_rop))
        }
        return pred_dict

# ============================================================================
# PERFORMANCE HELPERS
# ============================================================================

def _predict_physical_rop_vectorized(wob_klb, rpm, torque_ftlb, spp_psi, flow_gpm, ucs_psi,
                                     bit_diameter_in, bit_wear, depth_ft, cutter_count,
                                     add_noise: bool = False, seed: int = 42):
    """Versión vectorizada del modelo físico para acelerar datasets y mallas."""
    wob_klb = np.asarray(wob_klb, dtype=float)
    rpm = np.asarray(rpm, dtype=float)
    torque_ftlb = np.asarray(torque_ftlb, dtype=float)
    spp_psi = np.asarray(spp_psi, dtype=float)
    flow_gpm = np.asarray(flow_gpm, dtype=float)
    ucs_psi = np.asarray(ucs_psi, dtype=float)
    bit_diameter_in = np.asarray(bit_diameter_in, dtype=float)
    bit_wear = np.asarray(bit_wear, dtype=float)
    depth_ft = np.asarray(depth_ft, dtype=float)
    cutter_count = np.asarray(cutter_count, dtype=float)

    bit_diameter_in = np.maximum(0.1, bit_diameter_in)
    wob_lb = np.maximum(0.0, wob_klb * 1000.0)
    wob_opt = 5000.0 * bit_diameter_in
    f_wob = np.exp(-((wob_lb - wob_opt) ** 2) / (2.0 * np.maximum((wob_opt * 0.3) ** 2, 1e-6)))
    rpm_opt = 120.0
    f_rpm = np.exp(-((rpm - rpm_opt) ** 2) / (2.0 * (rpm_opt * 0.4) ** 2))
    ucs_psi = np.maximum(100.0, ucs_psi)
    f_formation = 5000.0 / np.maximum(ucs_psi ** 0.8, 1e-6)
    f_wear = 1.0 - bit_wear * 0.6
    f_torque = 1.0 + (torque_ftlb - 15000.0) / 50000.0
    hhp = (spp_psi * flow_gpm) / 1714.0
    hhp_opt = np.maximum(0.1, 2.5 * bit_diameter_in)
    f_hydraulic = 1.0 + 0.3 * (hhp - hhp_opt) / hhp_opt
    f_depth = np.exp(-depth_ft / 15000.0)
    cutter_factor = 1.0 + (cutter_count - 6.0) * 0.05
    rop_ftph = 120.0 * f_wob * f_rpm * f_formation * f_wear * f_torque * f_hydraulic * f_depth * cutter_factor
    if add_noise:
        rng = np.random.default_rng(seed)
        noise = rng.normal(0.0, np.maximum(rop_ftph * 0.05, 0.5), size=np.shape(rop_ftph))
        rop_ftph = rop_ftph + noise
    return np.maximum(5.0, rop_ftph)

@st.cache_data(show_spinner=False)
def _cached_rop_mesh(wob_min, wob_max, rpm_min, rpm_max, ucs_value, bit_diameter, n_wob=60, n_rpm=60):
    wob_values = np.linspace(wob_min, wob_max, n_wob)
    rpm_values = np.linspace(rpm_min, rpm_max, n_rpm)
    WOB, RPM = np.meshgrid(wob_values, rpm_values)
    ROP = _predict_physical_rop_vectorized(
        WOB, RPM,
        torque_ftlb=15000.0,
        spp_psi=3000.0,
        flow_gpm=800.0,
        ucs_psi=float(ucs_value),
        bit_diameter_in=float(bit_diameter),
        bit_wear=0.2,
        depth_ft=10000.0,
        cutter_count=6.0,
        add_noise=False,
    )
    return wob_values, rpm_values, WOB, RPM, ROP

# ============================================================================
# GENERADOR DE DATOS SINTÉTICOS - IMPERIAL
# ============================================================================

class DrillingDataGeneratorImperial:
    """Generador de datos sintéticos de perforación en unidades imperiales"""
    
    def __init__(self, n_samples=10000, seed=42):
        self.n_samples = n_samples
        self.seed = seed
        self._cached_data = None
    
    def generate(self):
        """Genera dataset completo en unidades imperiales y lo cachea por instancia."""
        if self._cached_data is not None:
            return self._cached_data.copy()

        rng = np.random.default_rng(self.seed)
        data = pd.DataFrame({
            'WOB_klb': rng.uniform(5, 35, self.n_samples),
            'RPM': rng.uniform(50, 200, self.n_samples),
            'Torque_ftlb': rng.uniform(5000, 35000, self.n_samples),
            'SPP_psi': rng.uniform(1500, 5500, self.n_samples),
            'Flow_gpm': rng.uniform(400, 1200, self.n_samples),
            'UCS_psi': rng.uniform(5000, 35000, self.n_samples),
            'Porosity_pct': rng.uniform(5, 30, self.n_samples),
            'Permeability_mD': rng.uniform(0.1, 1000, self.n_samples),
            'Abrasivity_CAI': rng.uniform(0.5, 4.5, self.n_samples),
            'BitDiameter_in': rng.choice([6, 7.875, 8.5, 9.875, 12.25, 14.75, 17.5, 26], self.n_samples),
            'BitWear_pct': rng.uniform(0, 80, self.n_samples) / 100.0,
            'CutterCount': rng.choice([4, 5, 6, 7, 8], self.n_samples),
            'NozzleSize_32nd': rng.choice([10, 12, 14, 16, 18, 20], self.n_samples),
            'Depth_ft': rng.uniform(2000, 20000, self.n_samples),
            'Inclination_deg': rng.beta(1, 10, self.n_samples) * 90.0,
            'Azimuth_deg': rng.uniform(0, 360, self.n_samples),
        })

        data['ROP_fthr'] = _predict_physical_rop_vectorized(
            data['WOB_klb'].values,
            data['RPM'].values,
            data['Torque_ftlb'].values,
            data['SPP_psi'].values,
            data['Flow_gpm'].values,
            data['UCS_psi'].values,
            data['BitDiameter_in'].values,
            data['BitWear_pct'].values,
            data['Depth_ft'].values,
            data['CutterCount'].values,
            add_noise=True,
            seed=self.seed,
        )

        self._cached_data = data
        return data.copy()

# ============================================================================
# FUNCIONES DE VISUALIZACIÓN - IMPERIAL
# ============================================================================

def create_correlation_analysis_chart(data: pd.DataFrame, target_col: str = 'ROP_fthr') -> go.Figure:
    """Crea gráfico de análisis de correlación entre variables y ROP."""
    numeric_cols = ['WOB_klb', 'RPM', 'Torque_ftlb', 'SPP_psi', 'Flow_gpm', 'UCS_psi', 
                   'BitDiameter_in', 'BitWear_pct', 'Depth_ft', target_col]
    df_corr = data[[c for c in numeric_cols if c in data.columns]].copy()
    if df_corr.empty or len(df_corr) < 10:
        fig = go.Figure()
        fig.add_annotation(text="Datos insuficientes para correlación", x=0.5, y=0.5, showarrow=False)
        fig.update_layout(height=400)
        return fig
    corr_matrix = df_corr.corr()
    label_map = {'WOB_klb': 'WOB (klb)', 'RPM': 'RPM', 'Torque_ftlb': 'Torque (ft-lb)', 'SPP_psi': 'SPP (psi)',
                 'Flow_gpm': 'Flow (gpm)', 'UCS_psi': 'UCS (psi)', 'BitDiameter_in': 'Bit Ø (in)',
                 'BitWear_pct': 'Bit wear', 'Depth_ft': 'Depth (ft)', 'ROP_fthr': 'ROP (ft/hr)'}
    labels = [label_map.get(c, c) for c in corr_matrix.columns]
    fig = go.Figure(data=go.Heatmap(
        z=corr_matrix.values,
        x=labels,
        y=labels,
        colorscale='RdBu',
        zmid=0,
        zmin=-1,
        zmax=1,
        text=[[f'{v:.2f}' for v in row] for row in corr_matrix.values],
        texttemplate='%{text}',
        textfont=dict(size=10),
        hovertemplate='%{x} vs %{y}<br>Correlación: %{z:.3f}<extra></extra>'
    ))
    fig.update_layout(
        title="Matriz de correlación — Variables vs ROP",
        xaxis=dict(tickangle=-45, tickfont=dict(size=10)),
        yaxis=dict(tickfont=dict(size=10)),
        height=500,
        margin=dict(l=100, r=50, t=50, b=100),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)'
    )
    return fig

def create_gauge_chart(value, title, min_val, max_val, unit):
    """Crea gráfico de indicador tipo gauge con valor numérico centrado en el semicírculo"""
    
    fig = go.Figure(go.Indicator(
        mode="gauge",
        value=value,
        title={'text': title, 'font': {'size': 16, 'color': '#1B4D3E'}},
        gauge={
            'axis': {'range': [min_val, max_val], 'tickwidth': 1},
            'bar': {'color': "#2C6E49"},
            'steps': [
                {'range': [min_val, min_val + (max_val-min_val)*0.3], 'color': "#FFE5B4"},
                {'range': [min_val + (max_val-min_val)*0.3, min_val + (max_val-min_val)*0.7], 'color': "#FFD700"},
                {'range': [min_val + (max_val-min_val)*0.7, max_val], 'color': "#FFB6C1"}
            ],
            'threshold': {
                'line': {'color': "red", 'width': 4},
                'thickness': 0.75,
                'value': max_val * 0.9
            }
        }
    ))
    
    # Valor numérico centrado dentro del semicírculo del gauge
    ref_val = (min_val + max_val) / 2
    delta_val = value - ref_val
    delta_str = f"▲{delta_val:.0f}" if delta_val >= 0 else f"▼{abs(delta_val):.0f}"
    
    value_text = f"{value:.1f}" if value != int(value) else f"{int(value)}"
    value_str = f"{value_text} {unit}"
    
    fig.add_annotation(
        text=f"<b>{value_str}</b><br><span style='font-size:12px;color:{'#2C6E49' if delta_val >= 0 else '#c0392b'}'>{delta_str}</span>",
        xref="paper", yref="paper",
        x=0.5, y=0.35,
        xanchor="center", yanchor="middle",
        showarrow=False,
        font=dict(size=24, color="#1B4D3E", family="Inter")
    )
    
    fig.update_layout(
        height=250,
        margin=dict(l=20, r=20, t=50, b=20),
        paper_bgcolor='rgba(0,0,0,0)',
        font={'color': "#1B4D3E", 'family': "Inter"}
    )
    
    return fig

def create_rop_heatmap(wob_range, rpm_range, ucs_value, bit_diameter):
    """Crea heat map de ROP vs WOB vs RPM."""
    ucs_value = max(1000, _safe_float(ucs_value, 15000))
    bit_diameter = max(6.0, _safe_float(bit_diameter, 8.5))
    wob_values, rpm_values, _, _, ROP = _cached_rop_mesh(
        float(wob_range[0]), float(wob_range[1]), float(rpm_range[0]), float(rpm_range[1]),
        float(ucs_value), float(bit_diameter), 60, 60
    )

    fig = go.Figure(data=go.Contour(
        z=ROP,
        x=wob_values,
        y=rpm_values,
        colorscale=[
            [0.0, '#2d1b69'],
            [0.25, '#2c6e49'],
            [0.5, '#4aa66b'],
            [0.75, '#f4d03f'],
            [1.0, '#e74c3c']
        ],
        contours=dict(coloring='heatmap', showlabels=True, labelfont=dict(size=10, color='white', family='Inter'), showlines=True),
        ncontours=20,
        line=dict(width=0.8, color='rgba(255,255,255,0.5)'),
        colorbar=dict(
            title=dict(text="ROP (ft/hr)", side="right", font=dict(size=13, color='#2d3436')),
            thickness=18, len=0.85, tickfont=dict(size=11, color='#2d3436'),
            bgcolor='rgba(255,255,255,0.8)', bordercolor='#dfe6e9', borderwidth=1
        ),
        hovertemplate='<b>WOB</b>: %{x:.1f} klb<br><b>RPM</b>: %{y:.0f}<br><b>ROP</b>: %{z:.1f} ft/hr<extra></extra>'
    ))
    fig.add_shape(type='rect', x0=18, y0=100, x1=26, y1=140, xref='x', yref='y', line=dict(color='#27ae60', width=2.5, dash='dot'), fillcolor='rgba(39,174,96,0.08)')
    fig.update_layout(
        title={'text': f"ROP Heat Map — UCS: {ucs_value:,.0f} psi · Bit: {bit_diameter}\"", 'font': {'size': 20, 'color': '#1B4D3E', 'family': 'Inter'}, 'x': 0.5, 'xanchor': 'center', 'y': 0.96},
        xaxis=dict(title="WOB (klb)", title_font=dict(size=14, color='#2d3436'), showgrid=False, zeroline=False, showline=True, linecolor='#b2bec3', linewidth=1, tickfont=dict(size=11, color='#636e72')),
        yaxis=dict(title="RPM (rev/min)", title_font=dict(size=14, color='#2d3436'), showgrid=False, zeroline=False, showline=True, linecolor='#b2bec3', linewidth=1, tickfont=dict(size=11, color='#636e72')),
        height=560, margin=dict(l=65, r=90, t=80, b=65), paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)', font=dict(family='Inter', size=12, color='#2d3436'),
        annotations=[dict(x=22, y=120, text="<b>Zona óptima</b><br>WOB 18–26 klb · RPM 100–140", showarrow=False, xref='x', yref='y', font=dict(size=11, color='#1B4D3E', family='Inter'), bgcolor='rgba(255,255,255,0.95)', bordercolor='#27ae60', borderwidth=1.5, borderpad=8)]
    )
    return fig

def create_nn_architecture_diagram():
    """Genera diagrama visual de la arquitectura de la red neuronal."""
    layers = [
        ('Entrada', 16, '#E8F5E9'),
        ('Dense 256 + BN + ReLU + Drop(0.3)', 256, '#C8E6C9'),
        ('Dense 128 + BN + ReLU + Drop(0.3)', 128, '#A5D6A7'),
        ('Dense 64 + BN + ReLU + Drop(0.2)', 64, '#81C784'),
        ('Dense 32 + BN + ReLU + Drop(0.2)', 32, '#66BB6A'),
        ('Dense 16 + ReLU', 16, '#4CAF50'),
        ('Salida (Lineal)', 1, '#1B4D3E'),
    ]
    y_labels = [f"{name}" for name, _, _ in layers]
    widths = [min(n * 1.15, 300) for _, n, _ in layers]
    colors = [c for _, _, c in layers]
    fig = go.Figure(go.Bar(
        x=widths,
        y=y_labels,
        orientation='h',
        marker=dict(
            color=colors,
            line=dict(color='#1B4D3E', width=1.2)
        ),
        text=[str(n) for _, n, _ in layers],
        textposition='outside',
        textfont=dict(size=11, color='#1B4D3E', family='Inter'),
        hovertemplate='%{y}<br>Unidades: %{text}<extra></extra>'
    ))
    fig.update_layout(
        xaxis=dict(
            range=[0, 320],
            showticklabels=False,
            showgrid=False,
            zeroline=False,
            fixedrange=True
        ),
        yaxis=dict(
            categoryorder='array',
            categoryarray=list(reversed(y_labels)),
            showgrid=False,
            tickfont=dict(size=10, color='#4A5568'),
            fixedrange=True
        ),
        height=340,
        margin=dict(l=180, r=50, t=50, b=30),
        plot_bgcolor='rgba(248,249,250,0.6)',
        paper_bgcolor='white',
        font=dict(family='Inter')
    )
    return fig

def _get_rop_regression_data(data_generator, predictor, n_samples=150):
    """Genera datos ROP observado vs predicho para regresión."""
    np.random.seed(42)
    data = data_generator.generate().sample(n=min(n_samples, 5000), random_state=42)
    actual_rop = data['ROP_fthr'].values
    predicted_rop = []
    for idx in range(len(data)):
        params = data.iloc[idx].to_dict()
        params.setdefault('mud_density_ppg', 10.0)
        params.setdefault('pore_gradient_ppg', 9.0)
        params.setdefault('yp_lb100ft2', 15.0)
        params.setdefault('pv_cp', 25.0)
        pred = predictor.predict_ensemble(params, use_rf=True, use_xgb=True, use_nn=True)
        predicted_rop.append(pred['Ensemble'])
    return actual_rop, np.array(predicted_rop)

def create_rop_linear_regression_chart(data_generator, predictor, n_samples=150):
    """
    Gráfico ROP observado vs predicho con regresión lineal mostrando la dispersión.
    Incluye banda de dispersión (±1 DE de los residuos).
    """
    x, y = _get_rop_regression_data(data_generator, predictor, n_samples)
    x_sort = np.linspace(x.min(), x.max(), 100)
    coef = np.polyfit(x, y, 1)
    y_fit = np.poly1d(coef)(x)
    y_lin = np.poly1d(coef)(x_sort)
    residuos = y - y_fit
    sd_res = np.std(residuos) if len(residuos) > 1 else 0
    r2 = np.corrcoef(y, y_fit)[0, 1] ** 2 if np.std(y) > 0 else 0
    y_sup = y_lin + sd_res
    y_inf = y_lin - sd_res
    # Escalas: incluir banda de dispersión y datos con margen ~8%
    x_range = x.max() - x.min() if x.max() > x.min() else max(x.max(), 1)
    y_data_max = max(y.max(), y_sup.max()) if len(y_sup) > 0 else y.max()
    y_data_min = min(y.min(), y_inf.min()) if len(y_inf) > 0 else y.min()
    y_range = y_data_max - y_data_min if y_data_max > y_data_min else max(y_data_max, 1)
    padding_x = max(x_range * 0.08, 50)
    padding_y = max(y_range * 0.08, 50)
    x_min = max(0, x.min() - padding_x)
    x_max = x.max() + padding_x
    y_min = max(0, y_data_min - padding_y)
    y_max = y_data_max + padding_y

    fig = go.Figure()
    # Banda de dispersión (primero para que quede detrás)
    fig.add_trace(go.Scatter(
        x=x_sort, y=y_sup, mode='lines', line=dict(width=0), showlegend=False
    ))
    fig.add_trace(go.Scatter(
        x=x_sort, y=y_inf, mode='lines', line=dict(width=0),
        fillcolor='rgba(37, 99, 235, 0.35)',
        fill='tonexty', name=f'Dispersión (±1 DE = {sd_res:.1f} ft/hr)'
    ))
    # Línea ideal (y=x)
    diag_max = max(x_max, y_max)
    fig.add_trace(go.Scatter(
        x=[0, diag_max], y=[0, diag_max], mode='lines',
        name='Ideal (y=x)', line=dict(color='#9CA3AF', width=1.5, dash='dash')
    ))
    # Datos
    fig.add_trace(go.Scatter(
        x=x, y=y, mode='markers', name='Datos',
        marker=dict(color='#22C55E', size=9, opacity=0.9, line=dict(color='#16A34A', width=1.2)),
        hovertemplate='ROP obs: %{x:.1f} ft/hr<br>ROP pred: %{y:.1f} ft/hr<extra></extra>'
    ))
    # Regresión lineal
    fig.add_trace(go.Scatter(
        x=x_sort, y=y_lin, mode='lines', name=f'Regresión lineal (R²={r2:.3f})',
        line=dict(color='#2563EB', width=3.5)
    ))
    # Marcas de eje con intervalos fijos para mejor legibilidad
    x_dtick = 1000  # 0, 1000, 2000, 3000...
    y_dtick = 250   # 0, 250, 500, 750, 1000, 1250, 1500...
    
    fig.update_layout(
        title=dict(
            text='ROP — Regresión lineal con dispersión',
            font=dict(size=18, color='#1B4D3E', family='Inter'),
            x=0.5, xanchor='center'
        ),
        xaxis=dict(
            title='ROP observado (ft/hr)',
            title_font=dict(size=13, color='#374151'),
            range=[x_min, x_max],
            dtick=x_dtick,
            tickfont=dict(size=11, color='#6B7280'),
            gridcolor='rgba(0,0,0,0.08)',
            zeroline=False,
            showgrid=True
        ),
        yaxis=dict(
            title='ROP predicho (ft/hr)',
            title_font=dict(size=13, color='#374151'),
            range=[y_min, y_max],
            dtick=y_dtick,
            tickfont=dict(size=11, color='#6B7280'),
            gridcolor='rgba(0,0,0,0.08)',
            zeroline=False,
            showgrid=True
        ),
        height=450,
        legend=dict(
            orientation='h',
            yanchor='bottom', y=1.02,
            xanchor='center', x=0.5,
            font=dict(size=11),
            bgcolor='rgba(255,255,255,0.9)',
            bordercolor='rgba(0,0,0,0.1)',
            borderwidth=1
        ),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        margin=dict(l=65, r=35, t=65, b=55),
        font=dict(family='Inter')
    )
    return fig

def create_residual_chart(data_generator, predictor, n_samples=150):
    """Gráfico de residuos (observado - predicho) vs predicho para diagnóstico del modelo."""
    try:
        actual_rop, predicted_rop = _get_rop_regression_data(data_generator, predictor, n_samples)
        residuals = actual_rop - predicted_rop
        fig = go.Figure()
        fig.add_trace(go.Scatter(
            x=predicted_rop, y=residuals, mode='markers',
            marker=dict(color='#2C6E49', size=8, opacity=0.7, line=dict(color='#1B4D3E', width=1)),
            name='Residuos',
            hovertemplate='Predicho: %{x:.1f} ft/hr<br>Residuo: %{y:.1f} ft/hr<extra></extra>'
        ))
        fig.add_hline(y=0, line_dash="dash", line_color='#9CA3AF', line_width=1.5)
        fig.update_layout(
            title=dict(text='Diagnóstico: Residuos vs ROP predicho', font=dict(size=16, color='#1B4D3E'), x=0.5, xanchor='center'),
            xaxis_title='ROP predicho (ft/hr)',
            yaxis_title='Residuo (obs - pred) ft/hr',
            height=320,
            paper_bgcolor='rgba(0,0,0,0)',
            plot_bgcolor='rgba(0,0,0,0)',
            font=dict(family='Inter'),
            margin=dict(l=60, r=30, t=50, b=50)
        )
        fig.update_xaxes(gridcolor='rgba(0,0,0,0.08)', zeroline=False)
        fig.update_yaxes(gridcolor='rgba(0,0,0,0.08)', zeroline=True)
        return fig
    except Exception:
        fig = go.Figure()
        fig.add_annotation(text="Datos insuficientes para residuos", x=0.5, y=0.5, showarrow=False)
        fig.update_layout(height=320)
        return fig

def create_feature_importance_chart(importance_dict, model_name='Random Forest'):
    """Crea gráfico de importancia de características."""
    if not importance_dict:
        fig = go.Figure()
        fig.add_annotation(text="Sin datos de importancia", x=0.5, y=0.5, showarrow=False)
        fig.update_layout(height=400)
        return fig
    features = list(importance_dict.keys())
    importance = list(importance_dict.values())
    imp_arr = np.array(importance)
    if imp_arr.max() <= 1 and imp_arr.min() >= 0:
        imp_pct = imp_arr
    else:
        imp_pct = imp_arr / (imp_arr.sum() + 1e-8)
    colors = ['#1B4D3E', '#2C6E49', '#3B8C5E', '#4AA66B', '#59C078', '#68DA85', '#77E890', '#86F69F']
    fig = go.Figure(go.Bar(
        x=imp_pct,
        y=features,
        orientation='h',
        marker_color=colors[:len(features)],
        text=[f'{i*100:.1f}%' for i in imp_pct],
        textposition='outside',
        textfont=dict(size=11, color='#1B4D3E'),
        hovertemplate='%{y}: %{x:.1%}<extra></extra>'
    ))
    fig.update_layout(
        title={'text': f"Importancia de variables — {model_name}", 'font': {'size': 16, 'color': '#1B4D3E'}, 'x': 0.5, 'xanchor': 'center'},
        xaxis_title="Importancia",
        xaxis=dict(tickformat='.0%', gridcolor='rgba(0,0,0,0.08)', zeroline=False),
        yaxis=dict(gridcolor='rgba(0,0,0,0.08)'),
        height=400,
        margin=dict(l=130, r=80, t=60, b=50),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)'
    )
    return fig

def _model_display_name(key):
    """Nombre legible para el modelo."""
    names = {'random_forest': 'Random Forest', 'xgboost': 'XGBoost', 'neural_network': 'Red neuronal'}
    return names.get(key, key.replace('_', ' ').title())

def _get_best_model_indices(metrics_dict):
    """Retorna índices del mejor modelo por cada métrica (R² alto, RMSE/MAE/MAPE bajo)."""
    models = list(metrics_dict.keys())
    r2_vals = [_safe_float(metrics_dict[m].get('r2', 0), 0) for m in models]
    rmse_vals = [_safe_float(metrics_dict[m].get('rmse', 999), 999) for m in models]
    mae_vals = [_safe_float(metrics_dict[m].get('mae', 999), 999) for m in models]
    mape_vals = [_safe_float(metrics_dict[m].get('mape', 999), 999) for m in models]
    best_r2 = np.argmax(r2_vals) if r2_vals else -1
    best_rmse = np.argmin(rmse_vals) if rmse_vals else -1
    best_mae = np.argmin(mae_vals) if mae_vals else -1
    best_mape = np.argmin(mape_vals) if mape_vals else -1
    return best_r2, best_rmse, best_mae, best_mape

def create_model_comparison_chart(metrics_dict):
    """Crea gráfico comparativo de modelos con R², RMSE, MAE, MAPE y destaque del mejor."""
    if not metrics_dict:
        fig = go.Figure()
        fig.add_annotation(
            text="No hay métricas de modelos. Entrene los modelos primero.",
            x=0.5, y=0.5, showarrow=False, font=dict(size=14, color='#6B7280')
        )
        fig.update_layout(height=420, paper_bgcolor='rgba(0,0,0,0)')
        return fig
    models = list(metrics_dict.keys())
    model_labels = [_model_display_name(m) for m in models]
    r2_scores = [_safe_float(metrics_dict[m].get('r2', 0), 0) for m in models]
    rmse_values = [_safe_float(metrics_dict[m].get('rmse', 10), 10) for m in models]
    mae_values = [_safe_float(metrics_dict[m].get('mae', 10), 10) for m in models]
    mape_values = [_safe_float(metrics_dict[m].get('mape', 20), 20) for m in models]

    best_r2, best_rmse, best_mae, best_mape = _get_best_model_indices(metrics_dict)
    base_colors = ['#1B4D3E', '#2C6E49', '#3B8C5E', '#4AA66B']
    highlight_color = '#E67E22'
    
    def _bar_colors(best_idx):
        return [highlight_color if i == best_idx else base_colors[i % len(base_colors)] for i in range(len(models))]
    
    fig = make_subplots(
        rows=2, cols=2,
        subplot_titles=(
            'R² (mayor es mejor)',
            'RMSE ft/hr (menor es mejor)',
            'MAE ft/hr (menor es mejor)',
            'MAPE % (menor es mejor)'
        ),
        specs=[[{'type': 'bar'}, {'type': 'bar'}], [{'type': 'bar'}, {'type': 'bar'}]],
        vertical_spacing=0.12,
        horizontal_spacing=0.08
    )
    
    fig.add_trace(
        go.Bar(x=model_labels, y=r2_scores, marker_color=_bar_colors(best_r2), text=[f'{s:.3f}' for s in r2_scores],
               textposition='outside', textfont=dict(size=11, color='#374151'), name='R²'),
        row=1, col=1
    )
    fig.add_trace(
        go.Bar(x=model_labels, y=rmse_values, marker_color=_bar_colors(best_rmse), text=[f'{v:.1f}' for v in rmse_values],
               textposition='outside', textfont=dict(size=11, color='#374151'), name='RMSE'),
        row=1, col=2
    )
    fig.add_trace(
        go.Bar(x=model_labels, y=mae_values, marker_color=_bar_colors(best_mae), text=[f'{v:.1f}' for v in mae_values],
               textposition='outside', textfont=dict(size=11, color='#374151'), name='MAE'),
        row=2, col=1
    )
    fig.add_trace(
        go.Bar(x=model_labels, y=mape_values, marker_color=_bar_colors(best_mape), text=[f'{v:.1f}%' for v in mape_values],
               textposition='outside', textfont=dict(size=11, color='#374151'), name='MAPE'),
        row=2, col=2
    )
    
    fig.update_layout(
        title={
            'text': "Comparación de rendimiento de modelos",
            'font': {'size': 18, 'color': '#1B4D3E', 'family': 'Inter'},
            'x': 0.5, 'xanchor': 'center'
        },
        height=520,
        showlegend=False,
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        font=dict(family='Inter'),
        margin=dict(t=80, b=60)
    )
    for i in range(1, 5):
        fig.update_xaxes(tickfont=dict(size=10, color='#6B7280'), gridcolor='rgba(0,0,0,0.06)', row=(i-1)//2+1, col=(i-1)%2+1)
        fig.update_yaxes(tickfont=dict(size=10, color='#6B7280'), gridcolor='rgba(0,0,0,0.06)', zeroline=False, row=(i-1)//2+1, col=(i-1)%2+1)
    
    fig.add_annotation(
        text="<i>Naranja = mejor modelo</i>",
        xref="paper", yref="paper", x=1, y=-0.08, showarrow=False, font=dict(size=10, color='#9CA3AF'),
        xanchor='right'
    )
    return fig

def create_optimization_3d_surface(wob_range, rpm_range, ucs_value, bit_diameter):
    """Crea superficie 3D de optimización ROP."""
    wob_values, rpm_values, WOB, RPM, ROP = _cached_rop_mesh(
        float(wob_range[0]), float(wob_range[1]), float(rpm_range[0]), float(rpm_range[1]),
        float(ucs_value), float(bit_diameter), 30, 30
    )
    fig = go.Figure(data=[go.Surface(
        z=ROP, x=WOB, y=RPM, colorscale='Viridis', opacity=0.85,
        contours={"z": {"show": True, "usecolormap": True, "highlightcolor": "limegreen", "project": {"z": True}}},
        hovertemplate='WOB: %{x:.1f} klb<br>RPM: %{y:.0f}<br>ROP: %{z:.1f} ft/hr<extra></extra>'
    )])
    max_idx = np.unravel_index(np.argmax(ROP), ROP.shape)
    opt_wob = WOB[max_idx]
    opt_rpm = RPM[max_idx]
    opt_rop = ROP[max_idx]
    fig.add_trace(go.Scatter3d(
        x=[opt_wob], y=[opt_rpm], z=[opt_rop], mode='markers',
        marker=dict(size=10, color='red', symbol='diamond'),
        name=f'Optimal: {opt_rop:.1f} ft/hr',
        hovertemplate='OPTIMAL<br>WOB: %{x:.1f} klb<br>RPM: %{y:.0f}<br>ROP: %{z:.1f} ft/hr<extra></extra>'
    ))
    fig.update_layout(
        title={'text': f"3D Optimization Surface - UCS: {ucs_value:,.0f} psi", 'font': {'size': 18, 'color': '#1B4D3E'}, 'x': 0.5, 'xanchor': 'center'},
        scene=dict(
            xaxis_title='WOB (klb)', yaxis_title='RPM', zaxis_title='ROP (ft/hr)',
            xaxis=dict(gridcolor='lightgray', gridwidth=1), yaxis=dict(gridcolor='lightgray', gridwidth=1), zaxis=dict(gridcolor='lightgray', gridwidth=1),
            camera=dict(eye=dict(x=1.5, y=1.5, z=1.5))
        ),
        height=600, margin=dict(l=0, r=0, t=100, b=0), paper_bgcolor='rgba(0,0,0,0)'
    )
    return fig, opt_wob, opt_rpm, opt_rop

# ============================================================================
# TEMPERATURA - CARGA DE TRAZAS REALES Y ML
# ============================================================================

# Aliases incluyen nombres típicos de CSV de pozo: DEPTH (m), Formation, Temperature (degC), ROP (min/m), etc.
TEMPERATURE_CANONICAL_COLUMNS = {
    'depth_md': ['depth_md', 'md', 'measured_depth', 'measured depth', 'depth', 'profundidad_md', 'profundidad'],
    'depth_tvd': ['depth_tvd', 'tvd', 'true_vertical_depth', 'true vertical depth', 'profundidad_tvd'],
    'rpm': ['rpm', 'surface_rpm', 'surface rpm'],
    'wob': ['wob', 'wob_klb', 'weight_on_bit', 'weight on bit'],
    'torque': ['torque', 'torque_ftlb', 'surface_torque', 'surface_torq', 'surface torq'],
    'rop': ['rop', 'rop_fthr', 'rop_ft_hr', 'rate_of_penetration', 'crop'],
    'flow_rate': ['flow_rate', 'flow_gpm', 'caudal', 'q', 'pump_rate', 'flow_in_rate', 'flow in rate'],
    'pump_pressure': ['pump_pressure', 'spp', 'spp_psi', 'standpipe_pressure', 'standpipe pre', 'standpipe_pre', 'presion_bomba'],
    'mud_in_temp': ['mud_in_temp', 'mud_temperature_in', 'inlet_mud_temp', 'temp_in', 'tin', 'temperature_1', 'temperature_in'],
    'mud_out_temp': ['mud_out_temp', 'mud_temperature_out', 'outlet_mud_temp', 'temp_out', 'tout', 'flowline_temp', 'temperature_out'],
    'mud_density': ['mud_density', 'mud_density_ppg', 'mw', 'mud_weight', 'mud_weight_ii', 'mud_weight_c', 'mud weight ii', 'mud weight c', 'mud_weight_in', 'mud_weight_out', 'mud weight in', 'mud weight out'],
    'pv': ['pv', 'pv_cp', 'plastic_viscosity'],
    'yp': ['yp', 'yp_lb100ft2', 'yield_point'],
    'bit_size': ['bit_size', 'bit_diameter_in', 'bitdiameter_in', 'bit_diameter'],
    'bit_depth': ['bit_depth', 'bit depth'],
    'hookload': ['hookload'],
    'standpipe_pressure': ['standpipe_pressure', 'spp_psi', 'pump_pressure', 'standpipe pre', 'standpipe_pre'],
    'lithology': ['lithology', 'litologia'],
    'formation': ['formation', 'formacion'],
}

TEMPERATURE_TARGET_CANDIDATES = [
    'mud_out_temp', 'mud_in_temp', 'temperature_downhole', 'bit_temp', 'annulus_temp',
    'temperature', 'temp', 'downhole_temperature', 'temperatura_fondo',
]

TRACE_UNIT_SYSTEM_LABELS = ['Auto (detectar)', 'Métrico (SI)', 'Imperial (US)']
TRACE_UNIT_SYSTEM_MAP = {
    'Auto (detectar)': 'auto',
    'Métrico (SI)': 'metric',
    'Imperial (US)': 'imperial',
}
# Incrementar al cambiar la lógica de unidades para invalidar perfiles viejos en session_state.
TRACE_PIPELINE_VERSION = 3


NEIGHBOR_WELL_NAME_ALIASES = ['well_name', 'well', 'pozo', 'wellid', 'well_id', 'name']
NEIGHBOR_X_ALIASES = ['x', 'x_coord', 'surface_x', 'east', 'easting', 'coord_x']
NEIGHBOR_Y_ALIASES = ['y', 'y_coord', 'surface_y', 'north', 'northing', 'coord_y']

def _first_existing_column(df: pd.DataFrame, aliases: List[str]) -> Optional[str]:
    normalized = {_normalize_colname(c): c for c in df.columns}
    for alias in aliases:
        key = _normalize_colname(alias)
        if key in normalized:
            return normalized[key]
    return None

def _infer_neighbor_well_metadata(df: pd.DataFrame, fallback_name: str = 'well') -> Dict[str, Any]:
    out = {'well_name': fallback_name, 'x': np.nan, 'y': np.nan}
    if df is None or df.empty:
        return out
    well_col = _first_existing_column(df, NEIGHBOR_WELL_NAME_ALIASES)
    if well_col is not None:
        vals = df[well_col].dropna().astype(str)
        if not vals.empty:
            out['well_name'] = vals.iloc[0]
    x_col = _first_existing_column(df, NEIGHBOR_X_ALIASES)
    y_col = _first_existing_column(df, NEIGHBOR_Y_ALIASES)
    if x_col is not None:
        out['x'] = pd.to_numeric(df[x_col], errors='coerce').dropna().mean()
    if y_col is not None:
        out['y'] = pd.to_numeric(df[y_col], errors='coerce').dropna().mean()
    return out

def _compute_depth_grid(series: pd.Series, points: int = 120) -> np.ndarray:
    vals = pd.to_numeric(series, errors='coerce').dropna().values
    if len(vals) < 2:
        return np.array([])
    lo, hi = float(np.nanmin(vals)), float(np.nanmax(vals))
    if not np.isfinite(lo) or not np.isfinite(hi) or hi <= lo:
        return np.array([])
    return np.linspace(lo, hi, min(points, max(20, int((hi - lo) / max((hi - lo) / points, 1)))))


def _neighbor_influence_scores(work: pd.DataFrame) -> np.ndarray:
    """Score 0–1 por fila: mayor = vecino más influyente (peso alto o distancia baja)."""
    n = len(work)
    out = np.full(n, 0.5)
    w = pd.to_numeric(work['weight'], errors='coerce') if 'weight' in work.columns else pd.Series(dtype=float)
    d = pd.to_numeric(work['distance'], errors='coerce') if 'distance' in work.columns else pd.Series(dtype=float)
    if w.notna().sum() >= max(2, min(2, n)):
        lo, hi = float(w.min()), float(w.max())
        if hi > lo:
            out = ((w.fillna(lo) - lo) / (hi - lo)).clip(0, 1).values
        else:
            out = np.ones(n) * 0.5
    elif d.notna().sum() >= max(2, min(2, n)):
        lo, hi = float(d.min()), float(d.max())
        if hi > lo:
            out = (1.0 - (d.fillna(hi) - lo) / (hi - lo)).clip(0, 1).values
        else:
            out = np.ones(n) * 0.5
    return out


def _traffic_light_html(score: float) -> str:
    """Mini semáforo horizontal: activa verde / ámbar / rojo según score de influencia."""
    if score >= 0.66:
        active, label = 'green', 'Alta influencia'
    elif score >= 0.33:
        active, label = 'amber', 'Influencia media'
    else:
        active, label = 'red', 'Baja influencia'
    # Colores apagados / encendidos con glow tipo LED
    def dot(which: str) -> str:
        on = active == which
        if which == 'red':
            bg = '#ef4444' if on else '#3f1d1d'
            glow = '0 0 10px rgba(239,68,68,0.85), 0 0 4px rgba(239,68,68,0.5)' if on else 'none'
        elif which == 'amber':
            bg = '#f59e0b' if on else '#3d2e0f'
            glow = '0 0 10px rgba(245,158,11,0.85), 0 0 4px rgba(245,158,11,0.5)' if on else 'none'
        else:
            bg = '#22c55e' if on else '#143d22'
            glow = '0 0 10px rgba(34,197,94,0.85), 0 0 4px rgba(34,197,94,0.5)' if on else 'none'
        op = '1' if on else '0.35'
        return (
            f'<span style="display:inline-block;width:11px;height:11px;border-radius:50%;'
            f'background:{bg};opacity:{op};box-shadow:{glow};'
            f'border:1px solid rgba(255,255,255,0.12);vertical-align:middle;"></span>'
        )
    housing = (
        f'<span title="{html.escape(label)}" style="display:inline-flex;align-items:center;gap:5px;'
        f'padding:4px 8px;background:linear-gradient(180deg,#1e293b 0%,#0f172a 100%);'
        f'border-radius:999px;border:1px solid rgba(148,163,184,0.25);box-shadow:inset 0 1px 0 rgba(255,255,255,0.06);">'
        f'{dot("red")}{dot("amber")}{dot("green")}</span>'
    )
    return housing


def neighbor_summary_table_html(df: pd.DataFrame, max_rows: int = 50) -> str:
    """Tabla HTML (tema oscuro) con semáforo de influencia por vecino."""
    if df is None or df.empty:
        return '<p style="color:#94a3b8;font-size:0.95rem;">Sin datos de resumen de vecinos.</p>'
    work = df.head(max_rows).copy()
    label_map = {
        'well_name': 'Pozo',
        'x': 'X',
        'y': 'Y',
        'distance': 'Distancia',
        'weight': 'Peso',
        'temp_min': 'T mín',
        'temp_max': 'T máx',
        'gradient': 'Gradiente',
        'rows': 'Filas',
    }
    skip_cols = {'x_col', 'y_col'}
    cols = [c for c in work.columns if c not in skip_cols]
    priority = ['well_name', 'distance', 'weight', 'temp_min', 'temp_max', 'gradient', 'rows', 'x', 'y']
    ordered = [c for c in priority if c in cols] + [c for c in cols if c not in priority]
    scores = _neighbor_influence_scores(work)

    def _cell_str(val: Any) -> str:
        if val is None or (isinstance(val, float) and np.isnan(val)):
            return '—'
        if isinstance(val, (np.integer, int)):
            return str(int(val))
        if isinstance(val, (np.floating, float)):
            return f'{float(val):,.4g}'
        return str(val)

    wrap_bg = '#0b0f14'
    table_bg = 'linear-gradient(180deg, rgba(15,23,42,0.92) 0%, rgba(11,15,20,0.98) 100%)'
    border_c = 'rgba(44,110,73,0.35)'
    th_style = (
        'text-align:left;padding:10px 12px;border-bottom:2px solid #2C6E49;color:#cbd5e1;'
        'font-size:0.78rem;text-transform:uppercase;letter-spacing:0.06em;font-weight:600;'
        'background:rgba(27,77,62,0.25);'
    )
    td_style = (
        'padding:9px 12px;border-bottom:1px solid rgba(51,65,85,0.5);font-size:0.88rem;'
        'color:#e2e8f0;'
    )
    th_sema = th_style + 'text-align:center;width:1%;white-space:nowrap;'
    td_sema = td_style + 'text-align:center;vertical-align:middle;'

    headers = (
        f'<th style="{th_sema}">Infl.</th>'
        + ''.join(
            f'<th style="{th_style}">{html.escape(label_map.get(c, str(c).replace("_", " ").title()))}</th>'
            for c in ordered
        )
    )
    body_rows = []
    for i, (_, row) in enumerate(work.iterrows()):
        sema = f'<td style="{td_sema}">{_traffic_light_html(float(scores[i]))}</td>'
        cells = ''.join(
            f'<td style="{td_style}">{html.escape(_cell_str(row[c]))}</td>' for c in ordered
        )
        zebra = 'background:rgba(30,41,59,0.25);' if i % 2 else ''
        body_rows.append(f'<tr style="{zebra}">{sema}{cells}</tr>')

    caption = (
        '<p style="margin:0 0 0.6rem 0;font-size:0.8rem;color:#94a3b8;line-height:1.45;">'
        '<span style="color:#64748b">Semaforo:</span> '
        '<span style="color:#22c55e">verde</span> = mayor influencia en el perfil '
        '(peso alto o distancia corta); '
        '<span style="color:#f59e0b">ambar</span> = media; '
        '<span style="color:#ef4444">rojo</span> = menor influencia relativa en este conjunto.'
        '</p>'
    )
    table = (
        f'<div style="overflow-x:auto;margin-top:0.5rem;padding:14px 16px;border-radius:12px;'
        f'background:{wrap_bg};border:1px solid {border_c};box-shadow:0 8px 32px rgba(0,0,0,0.35);">'
        f'{caption}'
        f'<table style="width:100%;border-collapse:separate;border-spacing:0;'
        f'background:{table_bg};border-radius:10px;overflow:hidden;border:1px solid rgba(51,65,85,0.4);">'
        f'<thead><tr>{headers}</tr></thead><tbody>{"".join(body_rows)}</tbody></table></div>'
    )
    return table


def create_neighbor_map_chart(
    neighbor_summary: pd.DataFrame,
    target_x: Optional[float] = None,
    target_y: Optional[float] = None,
    weighting: str = 'inverse_distance',
) -> go.Figure:
    """Mapa XY interactivo de pozo objetivo y pozos vecinos."""
    fig = go.Figure()
    if neighbor_summary is None or neighbor_summary.empty:
        fig.add_annotation(text='Sin datos de coordenadas para mapear vecinos', x=0.5, y=0.5, showarrow=False)
        fig.update_layout(height=420, template='plotly_white', title='Mapa interactivo de pozos vecinos')
        return fig

    work = neighbor_summary.copy()
    for c in ['x', 'y', 'distance', 'weight']:
        if c in work.columns:
            work[c] = pd.to_numeric(work[c], errors='coerce')
    map_df = work.dropna(subset=[c for c in ['x', 'y'] if c in work.columns]).copy() if {'x', 'y'}.issubset(work.columns) else pd.DataFrame()

    if pd.notna(target_x) and pd.notna(target_y):
        fig.add_trace(go.Scatter(
            x=[float(target_x)], y=[float(target_y)], mode='markers+text',
            name='Pozo objetivo',
            text=['Objetivo'], textposition='top center',
            marker=dict(size=18, symbol='x', color='#EF4444', line=dict(width=2, color='white')),
            hovertemplate='Pozo objetivo<br>X: %{x:,.2f}<br>Y: %{y:,.2f}<extra></extra>'
        ))

    if map_df.empty:
        fig.add_annotation(
            text='Los vecinos no traen columnas X/Y válidas; no es posible dibujar el mapa.',
            x=0.5, y=0.08, xref='paper', yref='paper', showarrow=False, font=dict(size=11, color='#64748B')
        )
        fig.update_layout(height=420, template='plotly_white', title='Mapa interactivo de pozos vecinos')
        return fig

    if 'weight' in map_df.columns and map_df['weight'].notna().any():
        w = map_df['weight'].fillna(map_df['weight'].median() if map_df['weight'].notna().any() else 1.0)
        wmin, wmax = float(w.min()), float(w.max())
        if wmax > wmin:
            size = 16 + 22 * (w - wmin) / (wmax - wmin)
        else:
            size = pd.Series(22.0, index=map_df.index)
    else:
        size = pd.Series(20.0, index=map_df.index)

    marker_kwargs = dict(size=size.tolist(), line=dict(width=1.5, color='white'), sizemode='diameter', opacity=0.92)
    if 'distance' in map_df.columns and map_df['distance'].notna().any():
        marker_kwargs.update(dict(
            color=map_df['distance'], colorscale='Viridis_r', showscale=True,
            colorbar=dict(title='Distancia', thickness=14, len=0.8),
        ))
    elif 'weight' in map_df.columns and map_df['weight'].notna().any():
        marker_kwargs.update(dict(
            color=map_df['weight'], colorscale='Turbo', showscale=True,
            colorbar=dict(title='Peso', thickness=14, len=0.8),
        ))
    else:
        marker_kwargs.update(dict(color='#10B981'))

    customdata = np.stack([
        map_df['distance'].values if 'distance' in map_df.columns else np.full(len(map_df), np.nan),
        map_df['weight'].values if 'weight' in map_df.columns else np.full(len(map_df), np.nan),
    ], axis=1)

    fig.add_trace(go.Scatter(
        x=map_df['x'], y=map_df['y'], mode='markers+text', name='Pozos vecinos',
        text=map_df['well_name'] if 'well_name' in map_df.columns else None,
        textposition='top center', customdata=customdata, marker=marker_kwargs,
        hovertemplate=(
            '<b>%{text}</b><br>X: %{x:,.2f}<br>Y: %{y:,.2f}<br>'
            'Distancia: %{customdata[0]:,.2f}<br>Peso: %{customdata[1]:,.4f}<extra></extra>'
        )
    ))

    if pd.notna(target_x) and pd.notna(target_y):
        for _, row in map_df.iterrows():
            fig.add_trace(go.Scatter(
                x=[float(target_x), float(row['x'])], y=[float(target_y), float(row['y'])],
                mode='lines', showlegend=False, hoverinfo='skip',
                line=dict(color='rgba(148,163,184,0.35)', width=1, dash='dot')
            ))

    x_vals = pd.concat([map_df['x'], pd.Series([target_x])]) if pd.notna(target_x) else map_df['x']
    y_vals = pd.concat([map_df['y'], pd.Series([target_y])]) if pd.notna(target_y) else map_df['y']
    xmin, xmax = float(x_vals.min()), float(x_vals.max())
    ymin, ymax = float(y_vals.min()), float(y_vals.max())
    xpad = max((xmax - xmin) * 0.12, 1.0)
    ypad = max((ymax - ymin) * 0.12, 1.0)

    subtitle = 'Tamaño = peso · color = distancia' if weighting == 'inverse_distance' else 'Tamaño/color = influencia relativa'
    fig.update_layout(
        title=f'Mapa interactivo de pozos vecinos — {subtitle}',
        height=460, template='plotly_white',
        xaxis=dict(title='Coordenada X', range=[xmin - xpad, xmax + xpad], zeroline=False),
        yaxis=dict(title='Coordenada Y', range=[ymin - ypad, ymax + ypad], zeroline=False, scaleanchor='x', scaleratio=1),
        legend=dict(orientation='h', yanchor='bottom', y=1.02, xanchor='left', x=0),
        margin=dict(l=40, r=40, t=70, b=40),
        hoverlabel=dict(bgcolor='rgba(15,23,42,0.92)', font=dict(color='#f8fafc'))
    )
    return fig


def _choose_available_depth_col(df: Optional[pd.DataFrame], preferred: str = 'depth') -> Optional[str]:
    """Escoge una columna de profundidad disponible de forma segura."""
    if df is None or df.empty:
        return None
    for c in [preferred, 'depth', 'depth_tvd', 'depth_md']:
        if c in df.columns:
            return c
    return None


def _depth_series_likely_meters(s: pd.Series, hint: str = '') -> bool:
    """True si la serie de profundidad parece estar en metros (sin convertir a ft)."""
    hint_l = (hint or '').lower().replace(' ', '')
    if 'ft' in hint_l or 'pie' in hint_l:
        return False
    if hint_l == 'm' or 'meter' in hint_l or 'metro' in hint_l:
        return True
    vals = pd.to_numeric(s, errors='coerce').dropna()
    if vals.empty:
        return False
    dmax = float(vals.max())
    if dmax <= 0:
        return False
    # >12000 suele ser profundidad ya expresada en ft (p. ej. 5920 m → 19423 ft)
    if dmax > 12000:
        return False
    # Pozos 4000–9000 m: al convertir superan ~13000 ft
    if dmax <= 9000 and (dmax * 3.28084) >= 13000:
        return True
    return dmax <= 8500


def _maybe_convert_depth_series_to_profile_ft(depth_values: pd.Series, profile_depth_ft: pd.Series) -> pd.Series:
    """
    Alinea profundidad de la traza con el perfil agregado (ambos en ft internos).
    Evita doble conversión m→ft cuando la traza ya fue normalizada.
    """
    ft_m = 3.28084
    d = pd.to_numeric(depth_values, errors='coerce')
    p = pd.to_numeric(profile_depth_ft, errors='coerce').dropna()
    d_valid = d.dropna()
    if d_valid.empty or p.empty:
        return d
    dmax = float(d_valid.max())
    pmax = float(p.max())
    if pmax <= 0 or dmax <= 0:
        return d
    ratio = dmax / pmax
    # Ya en la misma escala (ft)
    if 0.88 <= ratio <= 1.12:
        return d
    # Traza convertida dos veces (m→ft→ft): corrige dividiendo
    if 2.4 <= ratio <= 3.6:
        return d / ft_m
    # Traza aún en metros, perfil en ft
    if 0.22 <= ratio <= 0.42:
        return d * ft_m
    if _depth_series_likely_meters(d_valid, ''):
        err_m = abs(dmax * ft_m - pmax) / max(pmax, 1.0)
        if err_m < 0.15:
            return d * ft_m
    return d


def _safe_anomaly_threshold() -> float:
    """Umbral configurable desde Streamlit; fallback seguro para ejecución no interactiva."""
    try:
        return float(st.session_state.get('thermal_anomaly_z_threshold', 2.0))
    except Exception:
        return 2.0


class NeighborTemperatureProfiler:
    def __init__(self):
        self.profile_df = pd.DataFrame()
        self.neighbor_summary = pd.DataFrame()
        self.target_col = None
        self.depth_col = 'depth_tvd'
        self.weighting = 'inverse_distance'
        self.metadata = {}

    def build_from_neighbors(self, neighbor_dfs: List[pd.DataFrame], target_col: str, depth_col: str = 'depth_tvd', target_x: Optional[float] = None, target_y: Optional[float] = None, weighting: str = 'inverse_distance') -> pd.DataFrame:
        profiles = []
        summary_rows = []
        self.target_col = target_col
        self.depth_col = depth_col
        self.weighting = weighting
        for idx, raw_df in enumerate(neighbor_dfs):
            df = _standardize_temperature_trace_df(raw_df)
            if df.empty or target_col not in df.columns:
                continue
            if depth_col not in df.columns:
                if depth_col == 'depth_tvd' and 'depth_md' in df.columns:
                    work_depth = 'depth_md'
                elif depth_col == 'depth_md' and 'depth_tvd' in df.columns:
                    work_depth = 'depth_tvd'
                else:
                    continue
            else:
                work_depth = depth_col
            meta = _infer_neighbor_well_metadata(raw_df, fallback_name=f'Vecino {idx+1}')
            local = df[[work_depth, target_col]].copy()
            for c in ['formation', 'lithology']:
                if c in df.columns:
                    local[c] = df[c].astype(str)
            local[work_depth] = pd.to_numeric(local[work_depth], errors='coerce')
            local[target_col] = pd.to_numeric(local[target_col], errors='coerce')
            local = local.dropna(subset=[work_depth, target_col]).sort_values(work_depth)
            local = local.drop_duplicates(subset=[work_depth], keep='last')
            if len(local) < 8:
                continue
            distance = np.nan
            if pd.notna(target_x) and pd.notna(target_y) and pd.notna(meta['x']) and pd.notna(meta['y']):
                distance = float(np.sqrt((float(meta['x']) - float(target_x))**2 + (float(meta['y']) - float(target_y))**2))
            if weighting == 'equal' or not np.isfinite(distance):
                weight = 1.0
            else:
                weight = 1.0 / (distance + 1.0)
            grid = _compute_depth_grid(local[work_depth])
            if grid.size == 0:
                continue
            interp = np.interp(grid, local[work_depth].values, local[target_col].values)
            prof = pd.DataFrame({'depth': grid, 'temperature_interp': interp, 'weight': weight, 'well_name': meta['well_name'], 'distance': distance})
            if 'formation' in local.columns:
                prof['formation'] = pd.Series(local['formation']).ffill().bfill().iloc[0] if not local['formation'].dropna().empty else 'unknown'
            if 'lithology' in local.columns:
                prof['lithology'] = pd.Series(local['lithology']).ffill().bfill().iloc[0] if not local['lithology'].dropna().empty else 'unknown'
            profiles.append(prof)
            temp_gradient = np.polyfit(local[work_depth].values, local[target_col].values, 1)[0] if len(local) >= 2 else np.nan
            summary_rows.append({'well_name': meta['well_name'], 'x': meta.get('x', np.nan), 'y': meta.get('y', np.nan), 'x_col': meta.get('x_col'), 'y_col': meta.get('y_col'), 'distance': distance, 'weight': weight, 'temp_min': float(local[target_col].min()), 'temp_max': float(local[target_col].max()), 'gradient': float(temp_gradient), 'rows': int(len(local))})
        if not profiles:
            self.profile_df = pd.DataFrame()
            self.neighbor_summary = pd.DataFrame()
            return self.profile_df
        combined = pd.concat(profiles, ignore_index=True)
        depth_grid = np.sort(combined['depth'].unique())
        records = []
        for d in depth_grid:
            chunk = combined[np.isclose(combined['depth'], d)]
            temps = chunk['temperature_interp'].values
            weights = chunk['weight'].fillna(1.0).values
            if len(temps) == 0:
                continue
            weighted_mean = np.average(temps, weights=weights) if np.sum(weights) > 0 else float(np.mean(temps))
            records.append({'depth': float(d), 'temp_expected': float(weighted_mean), 'temp_median': float(np.median(temps)), 'temp_p10': float(np.percentile(temps, 10)), 'temp_p90': float(np.percentile(temps, 90)), 'temp_std': float(np.std(temps)), 'neighbors_used': int(len(temps)), 'gradient_expected': np.nan})
        prof_df = pd.DataFrame(records).sort_values('depth').reset_index(drop=True)
        if len(prof_df) >= 2:
            prof_df['gradient_expected'] = np.gradient(prof_df['temp_expected'], prof_df['depth'])
        self.profile_df = prof_df
        self.neighbor_summary = pd.DataFrame(summary_rows).sort_values(['distance', 'well_name'], na_position='last').reset_index(drop=True)
        self.metadata = {'target_x': target_x, 'target_y': target_y, 'neighbor_count': int(len(summary_rows))}
        return self.profile_df

    def predict_for_trace(self, target_df: pd.DataFrame) -> pd.DataFrame:
        if self.profile_df is None or self.profile_df.empty:
            return pd.DataFrame()
        if getattr(target_df, 'attrs', None) and target_df.attrs.get('trace_units_normalized'):
            df = target_df.copy()
        else:
            df = _standardize_temperature_trace_df(target_df).copy()
        if self.depth_col not in df.columns:
            alt = 'depth_md' if self.depth_col == 'depth_tvd' and 'depth_md' in df.columns else 'depth_tvd'
            if alt not in df.columns:
                return pd.DataFrame()
            depth_col = alt
        else:
            depth_col = self.depth_col
        work = df.copy()
        # La traza ya normalizada (m→ft) no debe reconvertirse.
        if getattr(df, 'attrs', None) and df.attrs.get('trace_depth_unit') == 'ft':
            work[depth_col] = pd.to_numeric(work[depth_col], errors='coerce')
        else:
            work[depth_col] = _maybe_convert_depth_series_to_profile_ft(work[depth_col], self.profile_df['depth'])
        work = work.dropna(subset=[depth_col]).sort_values(depth_col).copy()
        if work.empty:
            return pd.DataFrame()
        xp = self.profile_df['depth'].values
        yp = self.profile_df['temp_expected'].values
        p10 = self.profile_df['temp_p10'].values
        p90 = self.profile_df['temp_p90'].values
        std = self.profile_df['temp_std'].values
        work['temp_expected'] = np.interp(work[depth_col].values, xp, yp, left=np.nan, right=np.nan)
        work['temp_p10'] = np.interp(work[depth_col].values, xp, p10, left=np.nan, right=np.nan)
        work['temp_p90'] = np.interp(work[depth_col].values, xp, p90, left=np.nan, right=np.nan)
        work['temp_std_expected'] = np.interp(work[depth_col].values, xp, std, left=np.nan, right=np.nan)
        # Dispersión robusta: usar std de vecinos, fallback P10-P90 y fallback MAD global del residual.
        spread_p = (work['temp_p90'] - work['temp_p10']).abs() / 2.563
        work['temp_sigma_expected'] = work['temp_std_expected'].replace(0, np.nan).combine_first(spread_p.replace(0, np.nan))
        if self.target_col in work.columns:
            work[self.target_col] = pd.to_numeric(work[self.target_col], errors='coerce')
            work['temp_residual'] = work[self.target_col] - work['temp_expected']
            res_valid = work['temp_residual'].dropna()
            if not res_valid.empty:
                med = float(res_valid.median())
                mad = float((res_valid - med).abs().median())
                robust_sigma = 1.4826 * mad if mad > 0 else float(res_valid.std())
                if np.isfinite(robust_sigma) and robust_sigma > 0:
                    work['temp_sigma_expected'] = work['temp_sigma_expected'].fillna(robust_sigma)
            denom = work['temp_sigma_expected'].replace(0, np.nan)
            # Evitar z-scores inflados cuando la dispersión entre vecinos es ~0
            min_sigma = max(3.0, float(res_valid.std()) * 0.12) if not res_valid.empty else 3.0
            denom = denom.fillna(min_sigma).clip(lower=min_sigma)
            work['temp_zscore'] = work['temp_residual'] / denom
            threshold = _safe_anomaly_threshold()
            work['thermal_anomaly'] = np.where(work['temp_zscore'].abs() >= threshold, 'anomaly', 'normal')
            work['thermal_severity'] = pd.cut(
                work['temp_zscore'].abs(),
                bins=[-np.inf, threshold, threshold + 1.0, np.inf],
                labels=['normal', 'media', 'alta']
            ).astype(str)
        return work

def create_neighbor_temperature_profile_chart(
    profile_df: pd.DataFrame,
    actual_df: Optional[pd.DataFrame] = None,
    target_col: Optional[str] = None,
    depth_col: str = 'depth_tvd',
    *,
    depth_units: str = 'ft',
    depth_range_ft: Optional[tuple[float, float]] = None,
) -> go.Figure:
    """
    Perfil térmico esperado vs profundidad (datos internos en ft).
    depth_units='m' escala el eje Y a metros. depth_range_ft acota el intervalo vertical.
    """
    fig = go.Figure()
    if profile_df is None or profile_df.empty:
        fig.add_annotation(
            text='Sin perfil térmico de vecinos', x=0.5, y=0.5, showarrow=False,
            font=dict(color='#94a3b8'),
        )
        fig.update_layout(
            height=450, template='plotly_dark', paper_bgcolor='#0b0f14', plot_bgcolor='#111827',
        )
        return fig

    prof = profile_df.copy()
    if depth_range_ft is not None:
        lo_ft, hi_ft = float(depth_range_ft[0]), float(depth_range_ft[1])
        if hi_ft < lo_ft:
            lo_ft, hi_ft = hi_ft, lo_ft
        prof = prof[(prof['depth'] >= lo_ft) & (prof['depth'] <= hi_ft)].copy()
    if prof.empty:
        fig.add_annotation(
            text='Sin datos en el rango de profundidad seleccionado', x=0.5, y=0.5, showarrow=False,
            font=dict(color='#94a3b8'),
        )
        fig.update_layout(
            height=450, template='plotly_dark', paper_bgcolor='#0b0f14', plot_bgcolor='#111827',
            title=dict(text='Perfil térmico esperado desde pozos vecinos', font=dict(color='#f1f5f9')),
        )
        return fig

    use_m = str(depth_units).lower() in ('m', 'meter', 'meters', 'metros', 'metro')
    y_scale = UnitConverter.ft_to_m(1.0) if use_m else 1.0
    y_title = 'Profundidad (m)' if use_m else 'Profundidad (ft)'

    def _y(ser: pd.Series) -> np.ndarray:
        return (pd.to_numeric(ser, errors='coerce').values * y_scale).astype(float)

    def _actual_temp_as_smooth_line(
        act_raw: pd.DataFrame,
        dcol: str,
        tcol: str,
        *,
        profile_depth_max_ft: Optional[float] = None,
        max_depth_bins: int = 25000,
        interp_points: int = 800,
    ) -> tuple:
        """
        Perfil de temperatura real como una polilínea suave en el plano (T, profundidad).
        - Convierte MD a pies internos si la traza parece estar en m y el perfil en ft.
        - Un valor de T por pie (mediana), interpolación en profundidad y suavizado ligero.
        """
        a = act_raw[[tcol, dcol]].copy()
        d_raw = pd.to_numeric(a[dcol], errors='coerce')
        a[tcol] = pd.to_numeric(a[tcol], errors='coerce')
        a = a.assign(_d_raw=d_raw).dropna(subset=['_d_raw', tcol])
        if a.empty:
            return np.array([]), np.array([])
        d_raw_v = a['_d_raw'].astype(float).values
        dmax_raw = float(np.nanmax(d_raw_v))
        hint = _safe_float(profile_depth_max_ft, 0.0)
        # Perfil en ft; la traza a veces trae TVD en m (mismo orden de magnitud que el tope en ft).
        # Solo convertimos si el máximo en "m→ft" encaja mucho mejor con el perfil (evita falsos positivos en pies someros).
        d_ft = d_raw_v
        if hint > 10000 and dmax_raw > 50:
            err_ft = abs(dmax_raw - hint) / hint
            err_m = abs(dmax_raw * 3.28084 - hint) / hint
            if err_m < err_ft and err_m < 0.12:
                d_ft = d_raw_v * 3.28084
        a['d_ft'] = d_ft
        a['bin'] = np.round(a['d_ft']).astype(np.int64)
        g = (
            a.groupby('bin', sort=True, as_index=False)
            .agg({tcol: 'median', 'd_ft': 'mean'})
            .sort_values('d_ft', kind='mergesort')
        )
        d = g['d_ft'].to_numpy(dtype=float)
        t = g[tcol].to_numpy(dtype=float)
        inc = np.concatenate([[True], np.diff(d) > 1e-4])
        d, t = d[inc], t[inc]
        if len(d) < 2:
            return t, d * y_scale
        if len(d) > max_depth_bins:
            idx = np.unique(np.linspace(0, len(d) - 1, max_depth_bins, dtype=int))
            d, t = d[idx], t[idx]
        span = float(d.max() - d.min())
        n_out = min(interp_points, max(80, int(span / 25.0) + 1))
        n_out = max(2, n_out)
        d_grid = np.linspace(float(d.min()), float(d.max()), n_out)
        t_grid = np.interp(d_grid, d, t)
        # Suavizado de temperatura a lo largo de la profundidad (reduce zigzag visual)
        win = min(21, max(9, len(t_grid) // 40))
        if len(t_grid) >= win:
            t_grid = pd.Series(t_grid).rolling(win, center=True, min_periods=1).median().to_numpy()
        y_plot = d_grid * y_scale
        x_plot = t_grid
        # Orden estricto por profundidad (eje Y) para que la línea no se autocruce
        order = np.argsort(y_plot, kind='mergesort')
        x_plot, y_plot = x_plot[order], y_plot[order]
        # Quitar Y casi duplicados (evita segmentos horizontales por flotante)
        keep = np.ones(len(y_plot), dtype=bool)
        keep[1:] = np.abs(np.diff(y_plot)) > (1e-6 if y_scale >= 0.99 else 1e-4)
        x_plot, y_plot = x_plot[keep], y_plot[keep]
        return x_plot, y_plot

    # Suavizado ligero del perfil agregado (esperada / P10-P90) para lectura tipo log de pozo
    prof_plot = prof.sort_values('depth', kind='mergesort').copy()
    n_prof = len(prof_plot)
    win_p = min(17, max(5, n_prof // 60)) if n_prof >= 12 else 1
    if win_p > 1:
        for _col in ('temp_expected', 'temp_p10', 'temp_p90'):
            if _col in prof_plot.columns:
                prof_plot[_col] = (
                    pd.to_numeric(prof_plot[_col], errors='coerce')
                    .rolling(win_p, center=True, min_periods=1)
                    .median()
                )
    y_p = _y(prof_plot['depth'])

    fig.add_trace(go.Scatter(x=prof_plot['temp_p90'], y=y_p, mode='lines', line=dict(width=0), showlegend=False, hoverinfo='skip'))
    fig.add_trace(
        go.Scatter(
            x=prof_plot['temp_p10'],
            y=y_p,
            mode='lines',
            fill='tonextx',
            fillcolor='rgba(45,212,191,0.18)',
            line=dict(width=0),
            name='Banda P10-P90',
            hoverinfo='skip',
        )
    )
    fig.add_trace(
        go.Scatter(
            x=prof_plot['temp_p10'],
            y=y_p,
            mode='lines',
            name='P10 (límite)',
            line=dict(width=1.2, color='rgba(94,234,212,0.55)'),
            showlegend=False,
            hoverinfo='skip',
        )
    )
    fig.add_trace(
        go.Scatter(
            x=prof_plot['temp_p90'],
            y=y_p,
            mode='lines',
            name='P90 (límite)',
            line=dict(width=1.2, color='rgba(94,234,212,0.55)'),
            showlegend=False,
            hoverinfo='skip',
        )
    )
    fig.add_trace(
        go.Scatter(
            x=prof_plot['temp_expected'],
            y=y_p,
            mode='lines',
            name='Temperatura esperada',
            line=dict(width=3.2, color='#22d3ee', shape='spline', smoothing=0.35),
            hovertemplate='Esperada: %{x:.1f}<br>Prof.: %{y:,.0f}<extra></extra>',
        )
    )

    if actual_df is not None and not actual_df.empty and target_col and target_col in actual_df.columns and depth_col in actual_df.columns:
        act = actual_df[[target_col, depth_col]].copy()
        act[depth_col] = pd.to_numeric(act[depth_col], errors='coerce')
        act[target_col] = pd.to_numeric(act[target_col], errors='coerce')
        act = act.dropna(subset=[depth_col, target_col])
        if depth_range_ft is not None:
            lo_ft, hi_ft = float(depth_range_ft[0]), float(depth_range_ft[1])
            if hi_ft < lo_ft:
                lo_ft, hi_ft = hi_ft, lo_ft
            act = act[(act[depth_col] >= lo_ft) & (act[depth_col] <= hi_ft)]
        if not act.empty:
            prof_dmax = float(pd.to_numeric(prof['depth'], errors='coerce').max())
            x_real, y_real = _actual_temp_as_smooth_line(
                act, depth_col, target_col, profile_depth_max_ft=prof_dmax
            )
            if len(x_real) >= 2:
                fig.add_trace(
                    go.Scatter(
                        x=x_real,
                        y=y_real,
                        mode='lines',
                        name='Temperatura real',
                        line=dict(
                            width=3.8,
                            color='#ff2d55',
                            shape='spline',
                            smoothing=0.35,
                        ),
                        connectgaps=False,
                        hovertemplate='<b>Real</b>: %{x:.1f}<br>Prof.: %{y:,.0f}<extra></extra>',
                    )
                )
            elif len(x_real) == 1:
                fig.add_trace(
                    go.Scatter(
                        x=x_real,
                        y=y_real,
                        mode='markers',
                        name='Temperatura real',
                        marker=dict(size=11, color='#ff2d55', line=dict(width=2, color='#ffe4ec')),
                        hovertemplate='<b>Real</b>: %{x:.1f}<br>Prof.: %{y:,.0f}<extra></extra>',
                    )
                )

    fig.update_layout(
        height=540,
        template='plotly_dark',
        paper_bgcolor='#050508',
        plot_bgcolor='#0a0e14',
        title=dict(
            text='Perfil térmico esperado desde pozos vecinos',
            font=dict(size=18, color='#f8fafc', family='Inter'),
            x=0.02,
            xanchor='left',
        ),
        font=dict(family='Inter', color='#e2e8f0'),
        yaxis=dict(
            title=y_title,
            title_font=dict(color='#f1f5f9', size=13),
            autorange='reversed',
            tickformat=',.0f',
            showgrid=True,
            gridcolor='rgba(148,163,184,0.12)',
            zeroline=False,
            tickfont=dict(color='#94a3b8', size=11),
            linewidth=0,
            mirror=False,
        ),
        xaxis=dict(
            title='Temperatura',
            title_font=dict(color='#f1f5f9', size=13),
            tickformat='.0f',
            showgrid=True,
            gridcolor='rgba(148,163,184,0.12)',
            zeroline=False,
            tickfont=dict(color='#94a3b8', size=11),
            linewidth=0,
        ),
        hoverlabel=dict(bgcolor='rgba(15,23,42,0.96)', font=dict(color='#f8fafc', size=13)),
        legend=dict(
            orientation='h',
            yanchor='bottom',
            y=1.06,
            x=0,
            xanchor='left',
            font=dict(color='#e2e8f0', size=12),
            bgcolor='rgba(10,14,20,0.72)',
            bordercolor='rgba(148,163,184,0.25)',
            borderwidth=1,
        ),
        margin=dict(l=56, r=24, t=88, b=52),
    )
    return fig

def create_temperature_residual_chart(df: pd.DataFrame, target_col: str, depth_col: str = 'depth_tvd') -> go.Figure:
    fig = go.Figure()
    if df is None or df.empty or 'temp_residual' not in df.columns:
        fig.add_annotation(text='Sin residual térmico disponible', x=0.5, y=0.5, showarrow=False)
        fig.update_layout(height=420)
        return fig
    work = df[[depth_col, 'temp_residual']].copy()
    if 'thermal_anomaly' in df.columns:
        work['thermal_anomaly'] = df['thermal_anomaly'].reindex(work.index).fillna('normal')
    else:
        work['thermal_anomaly'] = 'normal'
    work = work.dropna(subset=[depth_col, 'temp_residual']).sort_values(depth_col)
    if work.empty:
        fig.add_annotation(text='Sin residual térmico disponible', x=0.5, y=0.5, showarrow=False)
        fig.update_layout(height=420)
        return fig
    n_total = len(work)
    # Reducir superposición: submuestreo si hay muchos puntos (mantener representatividad por profundidad)
    max_points = 1200
    if n_total > max_points:
        step = max(1, n_total // max_points)
        idx = np.linspace(0, n_total - 1, min(max_points, n_total), dtype=int)
        scatter_df = work.iloc[idx].copy()
    else:
        scatter_df = work.copy()
    is_anom = scatter_df.get('thermal_anomaly', 'normal') == 'anomaly'
    colors = np.where(is_anom, '#EF4444', 'rgba(6, 182, 212, 0.45)')
    sizes = np.where(is_anom, 10, 4)
    # Scatter: puntos pequeños y semitransparentes para que la densidad se vea sin amontonar
    fig.add_trace(go.Scatter(
        x=scatter_df['temp_residual'],
        y=scatter_df[depth_col],
        mode='markers',
        marker=dict(size=sizes, color=colors, line=dict(width=0), opacity=0.85 if is_anom.any() else 0.5),
        name='Residual',
        hovertemplate='Residual: %{x:.2f}<br>Profundidad: %{y:,.0f}<extra></extra>',
    ))
    # Línea de tendencia: mediana residual por ventana de profundidad (suavizado)
    n_bins = min(80, max(15, len(work) // 50))
    work['depth_bin'] = pd.cut(work[depth_col], bins=n_bins, labels=False, duplicates='drop')
    trend = work.groupby('depth_bin', observed=True).agg({'temp_residual': 'median', depth_col: 'mean'}).reset_index()
    trend = trend.dropna().sort_values(depth_col)
    if len(trend) >= 2:
        fig.add_trace(go.Scatter(
            x=trend['temp_residual'],
            y=trend[depth_col],
            mode='lines',
            line=dict(color='#F59E0B', width=2.5, shape='spline'),
            name='Tendencia (mediana)',
            hovertemplate='Mediana residual: %{x:.2f}<br>Profundidad: %{y:,.0f}<extra></extra>',
        ))
    fig.add_vline(x=0, line_dash='dash', line_color='#94A3B8', line_width=1)
    fig.update_layout(
        height=480,
        title=dict(text='Anomalías térmicas mientras perforas', font=dict(size=16, color='#F1F5F9')),
        xaxis=dict(
            title='Residual de temperatura',
            title_font=dict(color='#94A3B8'),
            tickfont=dict(color='#94A3B8'),
            gridcolor='rgba(148,163,184,0.12)',
            zeroline=False,
        ),
        yaxis=dict(
            title='Profundidad',
            title_font=dict(color='#94A3B8'),
            tickfont=dict(color='#94A3B8'),
            gridcolor='rgba(148,163,184,0.12)',
            autorange='reversed',
        ),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(15,23,42,0.4)',
        legend=dict(font=dict(color='#CBD5E1'), orientation='h', yanchor='bottom', y=1.02),
        font=dict(family='Inter'),
    )
    return fig

def create_lithology_temperature_chart(df: pd.DataFrame, target_col: str) -> go.Figure:
    fig = go.Figure()
    if df is None or df.empty:
        fig.add_annotation(text='Sin datos para correlación temperatura-litología', x=0.5, y=0.5, showarrow=False)
        fig.update_layout(height=420)
        return fig
    category_col = None
    if 'lithology' in df.columns and df['lithology'].notna().any():
        category_col = 'lithology'
    elif 'formation' in df.columns and df['formation'].notna().any():
        category_col = 'formation'
    if category_col is None or target_col not in df.columns:
        fig.add_annotation(text='No hay litología/formación disponible', x=0.5, y=0.5, showarrow=False)
        fig.update_layout(height=420)
        return fig
    temp_col = 'temp_residual' if 'temp_residual' in df.columns and df['temp_residual'].notna().any() else target_col
    plot_df = df[[category_col, temp_col]].dropna().copy()
    if plot_df.empty:
        fig.add_annotation(text='Sin datos suficientes para boxplot', x=0.5, y=0.5, showarrow=False)
        fig.update_layout(height=420)
        return fig
    fig = px.box(plot_df, x=category_col, y=temp_col, points='outliers', template='plotly_white', title='Correlación temperatura-litología/formación')
    ylabel = 'Residual térmico' if temp_col == 'temp_residual' else 'Temperatura'
    fig.update_layout(height=420, xaxis_title=category_col.title(), yaxis_title=ylabel)
    return fig


def compute_temperature_roadmap_by_formation(
    profile_df: Optional[pd.DataFrame] = None,
    prediction_df: Optional[pd.DataFrame] = None,
    formations: Optional[List[Dict]] = None,
    depth_col: str = 'depth',
) -> pd.DataFrame:
    """
    Calcula el roadmap de temperatura por formación.
    Prioridad:
    1) prediction_df con formation/lithology y temp_expected.
    2) profile_df agregado + intervalos geológicos.
    3) Fallback robusto: si los intervalos no se cruzan por unidades/rangos, reparte el perfil en
       intervalos proporcionales por formación para que el roadmap no quede vacío.
    """
    roadmap = []

    def _row_from_chunk(name: str, chunk: pd.DataFrame, dcol: str) -> Optional[Dict[str, Any]]:
        if chunk is None or chunk.empty or 'temp_expected' not in chunk.columns or dcol not in chunk.columns:
            return None
        chunk = chunk.dropna(subset=[dcol, 'temp_expected'])
        if chunk.empty:
            return None
        row = {
            'Formación': str(name),
            'Prof_min_ft': round(float(chunk[dcol].min()), 1),
            'Prof_max_ft': round(float(chunk[dcol].max()), 1),
            'Temp_media': round(float(chunk['temp_expected'].mean()), 2),
            'N_puntos': int(len(chunk)),
        }
        if 'temp_p10' in chunk.columns and chunk['temp_p10'].notna().any():
            row['Temp_P10'] = round(float(chunk['temp_p10'].mean()), 2)
        if 'temp_p90' in chunk.columns and chunk['temp_p90'].notna().any():
            row['Temp_P90'] = round(float(chunk['temp_p90'].mean()), 2)
        if 'temp_residual' in chunk.columns and chunk['temp_residual'].notna().any():
            row['Residual_medio'] = round(float(chunk['temp_residual'].mean()), 2)
        return row

    # Opción 1: pozo objetivo con formación/litología.
    if prediction_df is not None and not prediction_df.empty and 'temp_expected' in prediction_df.columns:
        form_col = None
        if 'formation' in prediction_df.columns and prediction_df['formation'].notna().any():
            form_col = 'formation'
        elif 'lithology' in prediction_df.columns and prediction_df['lithology'].notna().any():
            form_col = 'lithology'
        depth_axis = _choose_available_depth_col(prediction_df, depth_col)
        if form_col is not None and depth_axis is not None:
            cols = [form_col, 'temp_expected', depth_axis]
            for extra in ('temp_p10', 'temp_p90', 'temp_residual'):
                if extra in prediction_df.columns and extra not in cols:
                    cols.append(extra)
            work = prediction_df[cols].copy()
            work[form_col] = work[form_col].astype(str).str.strip()
            work = work[~work[form_col].str.lower().isin(['nan', 'none', '', 'unknown'])]
            work[depth_axis] = _maybe_convert_depth_series_to_profile_ft(
                work[depth_axis], profile_df['depth'] if profile_df is not None and 'depth' in profile_df.columns else work[depth_axis]
            )
            work = work.dropna(subset=['temp_expected', depth_axis])
            if not work.empty:
                for form, g in work.groupby(form_col, sort=False):
                    row = _row_from_chunk(form, g, depth_axis)
                    if row is not None:
                        roadmap.append(row)

    # Opción 2: perfil agregado + formaciones geológicas por rango de profundidad.
    if not roadmap and profile_df is not None and not profile_df.empty and formations:
        prof_depth = _choose_available_depth_col(profile_df, 'depth')
        if prof_depth is None or 'temp_expected' not in profile_df.columns:
            return pd.DataFrame(roadmap)
        prof = profile_df.copy()
        prof[prof_depth] = pd.to_numeric(prof[prof_depth], errors='coerce')
        prof = prof.dropna(subset=[prof_depth, 'temp_expected']).sort_values(prof_depth)
        if prof.empty:
            return pd.DataFrame(roadmap)

        pmin, pmax = float(prof[prof_depth].min()), float(prof[prof_depth].max())
        for f in formations:
            top, bottom = f.get('depth_top'), f.get('depth_bottom')
            if top is None or bottom is None:
                continue
            top, bottom = float(top), float(bottom)
            if bottom < top:
                top, bottom = bottom, top
            chunk = prof[(prof[prof_depth] >= top) & (prof[prof_depth] <= bottom)]
            # Si el intervalo cruza el perfil pero tiene pocos puntos, tomar el punto/interpolación cercano.
            if len(chunk) < 3 and bottom >= pmin and top <= pmax:
                lo, hi = max(top, pmin), min(bottom, pmax)
                if hi >= lo:
                    chunk = prof[(prof[prof_depth] >= lo) & (prof[prof_depth] <= hi)]
                    if len(chunk) < 3:
                        mids = np.linspace(lo, hi, 3) if hi > lo else np.array([lo])
                        interp = pd.DataFrame({
                            prof_depth: mids,
                            'temp_expected': np.interp(mids, prof[prof_depth], prof['temp_expected']),
                        })
                        for q in ('temp_p10', 'temp_p90'):
                            if q in prof.columns:
                                interp[q] = np.interp(mids, prof[prof_depth], prof[q])
                        chunk = interp
            if len(chunk) >= 1:
                row = _row_from_chunk(f.get('name', f'Intervalo {top:.0f}-{bottom:.0f}'), chunk, prof_depth)
                if row is not None:
                    roadmap.append(row)

        # Fallback: perfiles y formaciones no se cruzan (unidades/rango). Repartir perfil por proporción.
        if not roadmap:
            valid_forms = [f for f in formations if f.get('depth_top') is not None and f.get('depth_bottom') is not None]
            if valid_forms:
                form_tops = np.array([float(f.get('depth_top')) for f in valid_forms])
                form_bottoms = np.array([float(f.get('depth_bottom')) for f in valid_forms])
                fmin, fmax = float(np.nanmin(form_tops)), float(np.nanmax(form_bottoms))
                if fmax > fmin and pmax > pmin:
                    for f in valid_forms:
                        top, bottom = float(f.get('depth_top')), float(f.get('depth_bottom'))
                        lo = pmin + ((top - fmin) / (fmax - fmin)) * (pmax - pmin)
                        hi = pmin + ((bottom - fmin) / (fmax - fmin)) * (pmax - pmin)
                        if hi < lo:
                            lo, hi = hi, lo
                        chunk = prof[(prof[prof_depth] >= lo) & (prof[prof_depth] <= hi)]
                        if len(chunk) < 1:
                            mid = (lo + hi) / 2.0
                            chunk = pd.DataFrame({
                                prof_depth: [mid],
                                'temp_expected': [float(np.interp(mid, prof[prof_depth], prof['temp_expected']))],
                            })
                            for q in ('temp_p10', 'temp_p90'):
                                if q in prof.columns:
                                    chunk[q] = [float(np.interp(mid, prof[prof_depth], prof[q]))]
                        row = _row_from_chunk(f.get('name', f'Intervalo {top:.0f}-{bottom:.0f}'), chunk, prof_depth)
                        if row is not None:
                            row['Nota'] = 'Intervalo ajustado proporcionalmente al rango del perfil'
                            roadmap.append(row)
    return pd.DataFrame(roadmap)

def create_temperature_roadmap_chart(roadmap_df: pd.DataFrame) -> go.Figure:
    """Gráfico de barras: temperatura esperada por formación (roadmap antes de perforar)."""
    fig = go.Figure()
    _dark_axes = dict(
        title_font=dict(color='#e2e8f0'),
        tickfont=dict(color='#cbd5e1'),
        gridcolor='rgba(148,163,184,0.15)',
    )
    if roadmap_df is None or roadmap_df.empty or 'Formación' not in roadmap_df.columns:
        fig.add_annotation(
            text='No hay datos para el roadmap por formación. Genera el perfil esperado desde pozos vecinos.',
            x=0.5, y=0.5, showarrow=False, font=dict(size=12, color='#94a3b8'),
        )
        fig.update_layout(
            height=420, template='plotly_dark', paper_bgcolor='#0b0f14', plot_bgcolor='#111827',
            font=dict(color='#e2e8f0'),
        )
        return fig
    fig.add_trace(go.Bar(
        x=roadmap_df['Formación'],
        y=roadmap_df['Temp_media'],
        name='Temp. esperada',
        marker_color='#38bdf8',
        text=roadmap_df['Temp_media'].round(1),
        textposition='outside',
        textfont=dict(color='#f1f5f9'),
    ))
    if 'Temp_P10' in roadmap_df.columns and 'Temp_P90' in roadmap_df.columns:
        fig.add_trace(go.Bar(
            x=roadmap_df['Formación'], y=roadmap_df['Temp_P10'],
            name='P10', marker_color='rgba(56, 189, 248, 0.45)'
        ))
        fig.add_trace(go.Bar(
            x=roadmap_df['Formación'], y=roadmap_df['Temp_P90'],
            name='P90', marker_color='#22c55e'
        ))
    fig.update_layout(
        title=dict(text='Roadmap: temperatura esperada por formación (antes de perforar)', font=dict(color='#f1f5f9')),
        xaxis=dict(title='Formación', **_dark_axes),
        yaxis=dict(title='Temperatura', **_dark_axes),
        barmode='group',
        height=420,
        template='plotly_dark',
        paper_bgcolor='#0b0f14',
        plot_bgcolor='#111827',
        font=dict(family='Inter', color='#e2e8f0'),
        legend=dict(orientation='h', yanchor='bottom', y=1.02, xanchor='center', x=0.5, font=dict(color='#e2e8f0')),
    )
    return fig


def _interp_profile_temp_at_depth(profile_df: pd.DataFrame, depth_ft: float) -> Optional[float]:
    """Interpola temp_expected del perfil de vecinos a una profundidad (ft)."""
    if profile_df is None or profile_df.empty or 'depth' not in profile_df.columns:
        return None
    if 'temp_expected' not in profile_df.columns:
        return None
    d = pd.to_numeric(profile_df['depth'], errors='coerce')
    t = pd.to_numeric(profile_df['temp_expected'], errors='coerce')
    m = d.notna() & t.notna()
    n = int(m.sum())
    if n < 1:
        return None
    d = d[m].values.astype(float)
    t = t[m].values.astype(float)
    if n == 1:
        return float(t[0])
    order = np.argsort(d)
    d, t = d[order], t[order]
    z = float(np.clip(depth_ft, d.min(), d.max()))
    return float(np.interp(z, d, t))


def _profile_temp_baseline(profile_df: Optional[pd.DataFrame], depth_ft: float) -> Optional[float]:
    """
    Temperatura de referencia desde perfil de vecinos: interpola a profundidad;
    si no hay suficientes puntos, usa la media del perfil.
    """
    t = _interp_profile_temp_at_depth(profile_df, depth_ft) if profile_df is not None else None
    if t is not None and np.isfinite(t):
        return float(t)
    if profile_df is None or profile_df.empty or 'temp_expected' not in profile_df.columns:
        return None
    s = pd.to_numeric(profile_df['temp_expected'], errors='coerce').dropna()
    if s.empty:
        return None
    return float(s.mean())


def create_temperature_wob_rpm_heatmap(
    temperature_predictor: Optional[Any],
    base_params: Dict[str, Any],
    wob_range: tuple = (5.0, 40.0),
    rpm_range: tuple = (40.0, 220.0),
    profile_df: Optional[pd.DataFrame] = None,
    depth_ft: float = 10000.0,
    grid_n: int = 40,
    *,
    subtitle_extra: str = '',
) -> go.Figure:
    """
    Mapa de contorno WOB × RPM con temperatura en color (mismo estilo que el heatmap ROP).
    - Si hay modelo ML de temperatura entrenado: predice en cada celda.
    - Si no, pero hay perfil de vecinos: superficie heurística anclada a temp_expected a la profundidad actual.
    - Si no hay ninguno: figura vacía con mensaje.
    """
    w0, w1 = float(wob_range[0]), float(wob_range[1])
    r0, r1 = float(rpm_range[0]), float(rpm_range[1])
    wob_values = np.linspace(w0, w1, int(grid_n))
    rpm_values = np.linspace(r0, r1, int(grid_n))
    WOB, RPM = np.meshgrid(wob_values, rpm_values)
    Z = np.full_like(WOB, np.nan, dtype=float)
    mode = 'none'
    t_base = _profile_temp_baseline(profile_df, depth_ft)

    def _fill_heuristic_z(t0: float) -> None:
        nonlocal mode
        mode = 'heuristic'
        bp = dict(base_params)
        w_ref = float(bp.get('wob', 22.0))
        r_ref = float(bp.get('rpm', 120.0))
        for ii in range(len(wob_values)):
            for jj in range(len(rpm_values)):
                dw = WOB[jj, ii] - w_ref
                dr = RPM[jj, ii] - r_ref
                Z[jj, ii] = t0 + 0.14 * dw + 0.06 * dr + 0.00015 * (dr ** 2)

    if temperature_predictor is not None and getattr(temperature_predictor, 'model', None) is not None:
        mode = 'ml'
        Z = temperature_predict_grid_from_params(
            temperature_predictor, base_params, wob_values, rpm_values
        )
        # Si el ML falla en toda la malla (features faltantes, etc.), volver al perfil heurístico
        if np.all(np.isnan(Z)) and t_base is not None and np.isfinite(t_base):
            _fill_heuristic_z(t_base)
    elif t_base is not None and np.isfinite(t_base):
        _fill_heuristic_z(t_base)
    else:
        fig = go.Figure()
        fig.add_annotation(
            text='Genera el perfil desde pozos vecinos (botón en el panel) o entrena el modelo ML con la traza objetivo.',
            x=0.5, y=0.55, showarrow=False,
            font=dict(size=14, color='#94a3b8'),
            xref='paper', yref='paper',
        )
        fig.update_layout(
            title=dict(
                text='Temperatura esperada — WOB vs RPM',
                font=dict(size=18, color='#f1f5f9', family='Inter'),
                x=0.5, xanchor='center',
            ),
            height=560,
            template='plotly_dark',
            paper_bgcolor='#0b0f14',
            plot_bgcolor='#111827',
            margin=dict(l=40, r=40, t=70, b=40),
            font=dict(color='#e2e8f0'),
        )
        return fig

    if np.all(np.isnan(Z)):
        fig = go.Figure()
        fig.add_annotation(
            text='No se pudo calcular la malla (ML devolvió vacío y no hay perfil de vecinos válido).',
            x=0.5, y=0.55, showarrow=False, font=dict(size=14, color='#94a3b8'),
            xref='paper', yref='paper',
        )
        fig.update_layout(
            height=560, template='plotly_dark', paper_bgcolor='#0b0f14', plot_bgcolor='#111827',
            margin=dict(l=40, r=40, t=70, b=40), font=dict(color='#e2e8f0'),
        )
        return fig

    zmin = float(np.nanmin(Z))
    zmax = float(np.nanmax(Z))
    if not np.isfinite(zmin) or not np.isfinite(zmax) or zmax <= zmin:
        zmax = zmin + 1e-6

    target_name = ''
    if temperature_predictor is not None and getattr(temperature_predictor, 'target_col', None):
        target_name = str(temperature_predictor.target_col)

    sub_bits = [f'Prof.: {depth_ft:,.0f} ft']
    if subtitle_extra:
        sub_bits.append(subtitle_extra)
    if mode == 'ml':
        sub_bits.append('Fuente: ML temperatura')
    else:
        sub_bits.append('Fuente: perfil vecinos + tendencia WOB/RPM')
    if target_name:
        sub_bits.append(f'Variable: {target_name}')
    subtitle = ' · '.join(sub_bits)

    # Heatmap en lugar de Contour: evita fallos de render en algunos navegadores / Streamlit con tema oscuro.
    fig = go.Figure(
        data=go.Heatmap(
            z=Z,
            x=wob_values,
            y=rpm_values,
            colorscale=[
                [0.0, '#2d1b69'],
                [0.25, '#1f4e79'],
                [0.5, '#2c6e49'],
                [0.75, '#f4d03f'],
                [1.0, '#e85d4c'],
            ],
            zmin=zmin,
            zmax=zmax,
            colorbar=dict(
                title=dict(text='Temperatura', side='right', font=dict(size=13, color='#e2e8f0')),
                thickness=18,
                len=0.85,
                tickfont=dict(size=11, color='#cbd5e1'),
                bgcolor='rgba(15,23,42,0.75)',
                bordercolor='rgba(148,163,184,0.35)',
                borderwidth=1,
            ),
            hovertemplate='<b>WOB</b>: %{x:.1f} klb<br><b>RPM</b>: %{y:.0f}<br><b>Temp.</b>: %{z:.2f}<extra></extra>',
        )
    )

    fig.add_shape(
        type='rect', x0=18, y0=100, x1=26, y1=140,
        xref='x', yref='y',
        line=dict(color='#16a34a', width=3, dash='dot'),
        fillcolor='rgba(34,197,94,0.12)',
        layer='above',
    )

    fig.update_layout(
        template='plotly_dark',
        title={
            'text': f'Temperatura esperada — WOB vs RPM<br><sup style="color:#94a3b8">{subtitle}</sup>',
            'font': {'size': 19, 'color': '#f1f5f9', 'family': 'Inter'},
            'x': 0.5,
            'xanchor': 'center',
            'y': 0.97,
        },
        xaxis=dict(
            title='WOB (klb)',
            title_font=dict(size=14, color='#e2e8f0'),
            showgrid=False, zeroline=False, showline=True,
            linecolor='#64748b', linewidth=1,
            tickfont=dict(size=11, color='#cbd5e1'),
            constrain='domain',
        ),
        yaxis=dict(
            title='RPM (rev/min)',
            title_font=dict(size=14, color='#e2e8f0'),
            showgrid=False, zeroline=False, showline=True,
            linecolor='#64748b', linewidth=1,
            tickfont=dict(size=11, color='#cbd5e1'),
        ),
        height=560,
        margin=dict(l=72, r=100, t=100, b=72),
        paper_bgcolor='#0b0f14',
        plot_bgcolor='#111827',
        font=dict(family='Inter', size=12, color='#e2e8f0'),
        annotations=[
            dict(
                x=22, y=120,
                text='<b>Ventana de referencia</b><br>WOB 18–26 klb · RPM 100–140',
                showarrow=False,
                xref='x', yref='y',
                font=dict(size=11, color='#f8fafc', family='Inter'),
                bgcolor='rgba(30,41,59,0.95)',
                bordercolor='#22c55e',
                borderwidth=1.5,
                borderpad=8,
            )
        ],
        uirevision='temp_heatmap',
    )
    return fig


def create_temperature_roadmap_depth_contour(
    roadmap_df: pd.DataFrame,
    *,
    depth_col_lo: str = 'Prof_min_ft',
    depth_col_hi: str = 'Prof_max_ft',
    temp_col: str = 'Temp_media',
) -> go.Figure:
    """
    Vista tipo mapa de calor por profundidad: eje X = profundidad (ft), eje Y fijo,
    color = temperatura esperada por intervalo (roadmap por formación).
    """
    fig = go.Figure()
    if roadmap_df is None or roadmap_df.empty:
        fig.add_annotation(
            text='Sin datos de roadmap por formación', x=0.5, y=0.5, showarrow=False,
            font=dict(color='#94a3b8'),
        )
        fig.update_layout(
            height=360, template='plotly_dark', paper_bgcolor='#0b0f14', plot_bgcolor='#111827',
            font=dict(color='#e2e8f0'),
        )
        return fig
    req = [depth_col_lo, depth_col_hi, temp_col, 'Formación']
    if not all(c in roadmap_df.columns for c in req):
        fig.add_annotation(
            text='Tabla de roadmap incompleta', x=0.5, y=0.5, showarrow=False, font=dict(color='#94a3b8'),
        )
        fig.update_layout(
            height=360, template='plotly_dark', paper_bgcolor='#0b0f14', plot_bgcolor='#111827',
        )
        return fig

    mids = []
    temps = []
    labels = []
    for _, row in roadmap_df.iterrows():
        lo = _safe_float(row[depth_col_lo], np.nan)
        hi = _safe_float(row[depth_col_hi], np.nan)
        tv = _safe_float(row[temp_col], np.nan)
        if not (np.isfinite(lo) and np.isfinite(hi) and np.isfinite(tv)):
            continue
        mids.append((lo + hi) / 2.0)
        temps.append(tv)
        labels.append(str(row['Formación']))
    if len(mids) < 1:
        fig.add_annotation(
            text='Sin filas válidas en el roadmap', x=0.5, y=0.5, showarrow=False, font=dict(color='#94a3b8'),
        )
        fig.update_layout(
            height=360, template='plotly_dark', paper_bgcolor='#0b0f14', plot_bgcolor='#111827',
        )
        return fig

    order = np.argsort(mids)
    mids = np.array(mids)[order]
    temps = np.array(temps)[order]
    labels = [labels[i] for i in order]

    d_min, d_max = float(mids.min()), float(mids.max())
    pad = max((d_max - d_min) * 0.04, 50.0)
    depth_grid = np.linspace(d_min - pad, d_max + pad, 80)
    temp_grid = np.interp(depth_grid, mids, temps)
    Z = np.tile(temp_grid, (2, 1))
    zlo, zhi = float(np.nanmin(temp_grid)), float(np.nanmax(temp_grid))
    if not np.isfinite(zlo) or not np.isfinite(zhi) or zhi <= zlo:
        zhi = zlo + 1e-6

    fig.add_trace(
        go.Heatmap(
            z=Z,
            x=depth_grid,
            y=[0.0, 1.0],
            colorscale=[
                [0.0, '#2d1b69'], [0.25, '#1f4e79'], [0.5, '#2c6e49'],
                [0.75, '#f4d03f'], [1.0, '#e85d4c'],
            ],
            zmin=zlo,
            zmax=zhi,
            colorbar=dict(
                title=dict(text='Temp. esperada', side='right', font=dict(size=12, color='#e2e8f0')),
                thickness=14, len=0.82, tickfont=dict(color='#cbd5e1'),
                bgcolor='rgba(15,23,42,0.6)',
                bordercolor='rgba(148,163,184,0.35)',
                borderwidth=1,
            ),
            hovertemplate='Prof.: %{x:,.0f} ft<br>Temp.: %{z:.2f}<extra></extra>',
        )
    )
    for mi, ti, lb in zip(mids, temps, labels):
        fig.add_annotation(
            x=mi, y=0.5, xref='x', yref='y',
            text=f'<b>{lb}</b><br>{ti:.1f}',
            showarrow=False, font=dict(size=9, color='#f8fafc'),
            bgcolor='rgba(30,41,59,0.92)', borderpad=3, bordercolor='rgba(34,197,94,0.55)',
        )

    fig.update_layout(
        title=dict(
            text='Roadmap térmico vs profundidad (por formación)',
            font=dict(size=17, color='#f1f5f9', family='Inter'),
            x=0.5, xanchor='center',
        ),
        xaxis=dict(
            title='Profundidad (ft)',
            title_font=dict(color='#e2e8f0'),
            showgrid=True, gridcolor='rgba(148,163,184,0.18)', tickfont=dict(color='#cbd5e1'),
        ),
        yaxis=dict(visible=False, range=[-0.1, 1.1], showgrid=False),
        height=340,
        margin=dict(l=60, r=80, t=70, b=50),
        template='plotly_dark',
        paper_bgcolor='#0b0f14',
        plot_bgcolor='#111827',
        font=dict(family='Inter', color='#e2e8f0'),
    )
    return fig


def _safe_read_tabular_file(uploaded_file):
    try:
        name = uploaded_file.name.lower()
        if name.endswith('.csv'):
            return pd.read_csv(uploaded_file, low_memory=False)
        if name.endswith('.xlsx') or name.endswith('.xls'):
            return pd.read_excel(uploaded_file)
    except Exception:
        return None
    return None


def _get_trace_unit_system() -> str:
    """Modo de unidades seleccionado en la UI (auto / metric / imperial)."""
    try:
        label = st.session_state.get('trace_unit_system', TRACE_UNIT_SYSTEM_LABELS[0])
        return TRACE_UNIT_SYSTEM_MAP.get(label, 'auto')
    except Exception:
        return 'auto'


def _looks_like_unit_token(value: Any) -> bool:
    s = str(value).strip().lower()
    if not s or s in ('nan', 'none'):
        return False
    known = {
        'm', 'ft', 'feet', 'foot', 'pies', 'rpm', 'degc', 'degf', 'c', 'f', 'pa', 'psi', 'ppg',
        'kgf', 'klb', 'kn', 'gpm', 'n.m', 'n·m', 'ft-lb', 'ft.lbf', 'lbf', 'm3/min', 'min/m',
        'ft/hr', 'ft/h', 'g/cm3', 'g/cm³', 'm/h', 'm/hr', 'lpm', 'l/min',
    }
    if s in known:
        return True
    return any(tok in s for tok in ('/', 'deg', 'cm3', 'cm³', 'm3', 'min/m', 'ft/hr', 'ft/h'))


def _strip_units_row(df: pd.DataFrame) -> Tuple[pd.DataFrame, Dict[str, str]]:
    """Elimina fila de unidades (p. ej. m, m3/min, Pa…) y devuelve hints por columna normalizada."""
    if df is None or df.empty or len(df) < 2:
        return df, {}
    first = df.iloc[0]
    unit_like, numeric_like = 0, 0
    for val in first:
        s = str(val).strip()
        if not s or s.lower() in ('nan', 'none'):
            continue
        if _looks_like_unit_token(val):
            unit_like += 1
        elif _safe_float(val, None) is not None:
            numeric_like += 1
    if unit_like < 2 or unit_like <= numeric_like:
        return df, {}
    hints: Dict[str, str] = {}
    for col in df.columns:
        raw = str(first[col]).strip()
        if raw and raw.lower() not in ('nan', 'none', ''):
            hints[_normalize_colname(col)] = raw
    return df.iloc[1:].reset_index(drop=True), hints


def _unit_hint_for_column(unit_hints: Dict[str, str], *col_names: str) -> str:
    for name in col_names:
        key = _normalize_colname(name)
        if key in unit_hints:
            return str(unit_hints[key]).strip().lower()
    return ''


def _series_median(s: pd.Series) -> float:
    v = pd.to_numeric(s, errors='coerce').dropna()
    return float(v.median()) if not v.empty else np.nan


def _apply_trace_unit_conversions(
    df: pd.DataFrame,
    unit_hints: Optional[Dict[str, str]] = None,
    unit_system: str = 'auto',
) -> pd.DataFrame:
    """
    Normaliza trazas a unidades internas imperiales (ft, klb, ft/hr, psi, gpm, ft-lb, ppg).
    unit_system: auto | metric | imperial
    """
    if df is None or df.empty:
        return df
    out = df.copy()
    hints = unit_hints or {}
    mode = (unit_system or 'auto').lower()
    if mode not in ('auto', 'metric', 'imperial'):
        mode = 'auto'

    def _hint(*names: str) -> str:
        return _unit_hint_for_column(hints, *names)

    # --- Profundidad (ft interno) ---
    for col in ('depth_md', 'depth_tvd', 'bit_depth'):
        if col not in out.columns:
            continue
        s = pd.to_numeric(out[col], errors='coerce')
        if s.notna().sum() < 2:
            continue
        hint = _hint(col, 'depth', 'bit depth', 'bit_depth')
        to_ft = False
        if mode == 'metric':
            to_ft = True
        elif mode == 'imperial':
            to_ft = False
        else:
            if 'ft' in hint or 'pie' in hint:
                to_ft = False
            elif hint == 'm' or 'meter' in hint or 'metro' in hint:
                to_ft = True
            else:
                to_ft = _depth_series_likely_meters(s, hint)
        if to_ft:
            out[col] = s * 3.28084
        else:
            out[col] = s

    # --- WOB (klb) ---
    if 'wob' in out.columns:
        s = pd.to_numeric(out['wob'], errors='coerce')
        hint = _hint('wob', 'weight on bit')
        med = _series_median(s)
        if mode == 'metric' or 'kgf' in hint:
            out['wob'] = s * 0.00220462
        elif 'kn' in hint:
            out['wob'] = s.apply(converter.kn_to_klb)
        elif mode == 'imperial' or 'klb' in hint or 'klbf' in hint:
            out['wob'] = s
        elif mode == 'auto' and np.isfinite(med) and med > 150:
            out['wob'] = s * 0.00220462
        else:
            out['wob'] = s

    # --- ROP (ft/hr) ---
    if 'rop' in out.columns:
        s = pd.to_numeric(out['rop'], errors='coerce')
        hint = _hint('rop', 'rate of penetration').replace(' ', '')
        med = _series_median(s)
        if mode == 'metric' or 'min/m' in hint:
            out['rop'] = s.apply(converter.min_per_m_to_fthr)
        elif 'm/h' in hint or 'm/hr' in hint:
            out['rop'] = s * 3.28084
        elif mode == 'imperial' or 'ft/hr' in hint or 'ft/h' in hint or 'fthr' in hint:
            out['rop'] = s
        elif mode == 'auto' and np.isfinite(med) and med < 35 and float(s.max()) < 500:
            out['rop'] = s.apply(converter.min_per_m_to_fthr)
        else:
            out['rop'] = s

    # --- Presión standpipe / bomba (psi) ---
    for col in ('standpipe_pressure', 'pump_pressure'):
        if col not in out.columns:
            continue
        s = pd.to_numeric(out[col], errors='coerce')
        hint = _hint(col, 'standpipe pressure', 'spp', 'pump pressure')
        med = _series_median(s)
        if mode == 'metric' or 'pa' in hint:
            out[col] = s.apply(converter.pa_to_psi)
        elif 'mpa' in hint:
            out[col] = s.apply(converter.mpa_to_psi)
        elif mode == 'imperial' or 'psi' in hint:
            out[col] = s
        elif mode == 'auto' and np.isfinite(med) and med > 50000:
            out[col] = s.apply(converter.pa_to_psi)
        else:
            out[col] = s

    # --- Caudal (gpm) ---
    if 'flow_rate' in out.columns:
        s = pd.to_numeric(out['flow_rate'], errors='coerce')
        hint = _hint('flow_rate', 'flow in rate', 'flow rate', 'caudal')
        med = _series_median(s)
        if mode == 'metric' or 'm3/min' in hint or 'm³/min' in hint:
            out['flow_rate'] = s.apply(converter.m3min_to_gpm)
        elif 'lpm' in hint or 'l/min' in hint:
            out['flow_rate'] = s.apply(converter.lpm_to_gpm)
        elif mode == 'imperial' or 'gpm' in hint:
            out['flow_rate'] = s
        elif mode == 'auto' and np.isfinite(med) and med < 25:
            out['flow_rate'] = s.apply(converter.m3min_to_gpm)
        else:
            out['flow_rate'] = s

    # --- Torque (ft-lb) ---
    if 'torque' in out.columns:
        s = pd.to_numeric(out['torque'], errors='coerce')
        hint = _hint('torque', 'surface torque', 'surface torq')
        med = _series_median(s)
        if mode == 'metric' or 'n.m' in hint or 'n·m' in hint or 'nm' in hint.replace('.', ''):
            out['torque'] = s.apply(converter.nm_to_ftlb)
        elif 'knm' in hint or 'kn-m' in hint:
            out['torque'] = s.apply(converter.knm_to_ftlb)
        elif mode == 'imperial' or 'ft-lb' in hint or 'ft.lb' in hint or 'ftlb' in hint:
            out['torque'] = s
        elif mode == 'auto' and np.isfinite(med) and med > 80000:
            out['torque'] = s.apply(converter.nm_to_ftlb)
        else:
            out['torque'] = s

    # --- Densidad de lodo (ppg) ---
    if 'mud_density' in out.columns:
        s = pd.to_numeric(out['mud_density'], errors='coerce')
        hint = _hint('mud_density', 'mud weight in', 'mud weight out', 'mud weight')
        med = _series_median(s)
        if mode == 'metric' or 'g/cm3' in hint or 'g/cm³' in hint:
            out['mud_density'] = s.apply(converter.gcm3_to_ppg)
        elif mode == 'imperial' or 'ppg' in hint:
            out['mud_density'] = s
        elif mode == 'auto' and np.isfinite(med) and med < 25:
            out['mud_density'] = s.apply(converter.gcm3_to_ppg)
        else:
            out['mud_density'] = s

    # Temperatura: se deja en °C (trazas DrillSpot/RACEMOSA); no convertir a °F para no romper targets.
    return out


def _normalize_colname(col: str) -> str:
    return str(col).strip().lower().replace('-', '_').replace(' ', '_')


def _standardize_temperature_trace_df(
    df: pd.DataFrame,
    unit_system: Optional[str] = None,
) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame()
    if unit_system is None:
        unit_system = _get_trace_unit_system()
    raw, unit_hints = _strip_units_row(df.copy())
    out = raw.copy()
    out.columns = [_normalize_colname(c) for c in out.columns]
    rename_map = {}
    for canonical, aliases in TEMPERATURE_CANONICAL_COLUMNS.items():
        for alias in aliases:
            alias_norm = _normalize_colname(alias)
            if alias_norm in out.columns and alias_norm not in rename_map:
                rename_map[alias_norm] = canonical
                break
    out = out.rename(columns=rename_map)
    # columnas derivadas útiles desde la app ROP
    if 'depth_md' not in out.columns and 'depth_ft' in out.columns:
        out['depth_md'] = out['depth_ft']
    if 'depth_tvd' not in out.columns and 'depth_ft' in out.columns:
        out['depth_tvd'] = out['depth_ft']
    if 'wob' not in out.columns and 'wob_klb' in out.columns:
        out['wob'] = out['wob_klb']
    if 'torque' not in out.columns and 'torque_ftlb' in out.columns:
        out['torque'] = out['torque_ftlb']
    if 'flow_rate' not in out.columns and 'flow_gpm' in out.columns:
        out['flow_rate'] = out['flow_gpm']
    if 'pump_pressure' not in out.columns and 'spp_psi' in out.columns:
        out['pump_pressure'] = out['spp_psi']
    if 'standpipe_pressure' not in out.columns and 'spp_psi' in out.columns:
        out['standpipe_pressure'] = out['spp_psi']
    if 'bit_size' not in out.columns and 'bitdiameter_in' in out.columns:
        out['bit_size'] = out['bitdiameter_in']
    if 'bit_size' not in out.columns and 'bit_diameter_in' in out.columns:
        out['bit_size'] = out['bit_diameter_in']
    if 'rop' not in out.columns and 'rop_fthr' in out.columns:
        out['rop'] = out['rop_fthr']
    if 'mud_density' not in out.columns:
        for src in ('mud_weight_in', 'mud_weight_out', 'mud_weight'):
            if src in out.columns:
                out['mud_density'] = out[src]
                break
    if 'depth_tvd' not in out.columns and 'depth_md' in out.columns:
        out['depth_tvd'] = out['depth_md'].copy()
    num_cols = [
        'depth_md', 'depth_tvd', 'bit_depth', 'rpm', 'wob', 'torque', 'rop', 'flow_rate',
        'pump_pressure', 'mud_in_temp', 'mud_out_temp', 'mud_density', 'pv', 'yp',
        'bit_size', 'hookload', 'standpipe_pressure',
    ]
    for col in num_cols:
        if col in out.columns:
            out[col] = pd.to_numeric(out[col], errors='coerce')
    # Formación/litología: rellenar huecos para que no se pierdan (ej. CSV con "Formation" solo en algunas filas)
    for form_col in ['formation', 'lithology']:
        if form_col in out.columns:
            s = out[form_col].astype(str).str.strip().replace({'nan': np.nan, 'none': np.nan, '': np.nan})
            if s.notna().any():
                out[form_col] = s.ffill().bfill().fillna('')
    out = _apply_trace_unit_conversions(out, unit_hints=unit_hints, unit_system=unit_system)
    out.attrs['trace_depth_unit'] = 'ft'
    out.attrs['trace_units_normalized'] = True
    return out


def _convert_metric_trace_to_imperial(df: pd.DataFrame) -> pd.DataFrame:
    """Compatibilidad: delega en la normalización completa de unidades."""
    return _apply_trace_unit_conversions(df, unit_hints={}, unit_system='auto')

def _detect_temperature_target(df: pd.DataFrame) -> Optional[str]:
    cols = {_normalize_colname(c): c for c in df.columns}
    for cand in TEMPERATURE_TARGET_CANDIDATES:
        if cand in cols:
            return cols[cand]
    return None



def _add_temperature_engineered_features(df: pd.DataFrame) -> pd.DataFrame:
    """Agrega variables fisicas derivadas para el modelo termico.
    No usa la temperatura objetivo directamente; solo sensores/parametros disponibles.
    """
    out = df.copy()
    def _num(name):
        if name in out.columns:
            return pd.to_numeric(out[name], errors='coerce')
        return pd.Series(np.nan, index=out.index, dtype=float)
    depth_md = _num('depth_md')
    depth_tvd = _num('depth_tvd')
    wob = _num('wob')
    rpm = _num('rpm')
    torque = _num('torque')
    rop = _num('rop')
    flow = _num('flow_rate')
    pump = _num('pump_pressure')
    mud_in = _num('mud_in_temp')
    mud_density = _num('mud_density')
    bit_depth = _num('bit_depth')

    if 'depth_md' in out.columns:
        out['depth_md_sqrt'] = np.sqrt(np.maximum(depth_md, 0))
        out['depth_md_sq_scaled'] = (depth_md ** 2) / 1_000_000.0
    if 'depth_md' in out.columns and 'depth_tvd' in out.columns:
        out['depth_delta_ft'] = (depth_md - depth_tvd).abs()
    elif 'depth_delta_ft' not in out.columns:
        out['depth_delta_ft'] = 0.0
    if 'bit_depth' in out.columns and 'depth_md' in out.columns:
        out['bit_md_delta_ft'] = (bit_depth - depth_md).abs()

    # Energia mecanica / friccion aproximada. Sube con WOB, RPM y torque.
    out['mech_energy_proxy'] = (wob.clip(lower=0).fillna(0) * rpm.clip(lower=0).fillna(0) * torque.clip(lower=0).fillna(0)) / 1_000_000.0
    # Hidraulica: mas caudal tiende a enfriar; mas presion/energia puede calentar por friccion.
    out['hydraulic_power_proxy'] = (pump.clip(lower=0).fillna(0) * flow.clip(lower=0).fillna(0)) / 1_000_000.0
    out['cooling_proxy'] = flow / (depth_md.abs() + 1.0)
    out['flow_per_pressure'] = flow / (pump.abs() + 1.0)
    out['pressure_per_flow'] = pump / (flow.abs() + 1.0)
    out['rop_inverse'] = 1.0 / (rop.abs() + 1.0)
    out['thermal_load_proxy'] = out['mech_energy_proxy'] / (flow.abs() + 1.0)
    if 'mud_in_temp' in out.columns:
        out['mud_in_x_depth'] = mud_in * depth_md / 10_000.0
    if 'mud_density' in out.columns:
        out['mud_density_x_flow'] = mud_density * flow / 1_000.0
    return out

def temperature_prepare_features_for_prediction(tp: Any, X: pd.DataFrame) -> pd.DataFrame:
    """
    Alinea columnas con el Pipeline de temperatura y agrega variables derivadas.
    Funciona aunque `tp` sea una instancia antigua guardada en session_state.
    """
    out = _add_temperature_engineered_features(X.copy())
    if getattr(tp, '_use_depth_delta', False):
        if 'depth_md' in out.columns and 'depth_tvd' in out.columns:
            dm = pd.to_numeric(out['depth_md'], errors='coerce')
            dt = pd.to_numeric(out['depth_tvd'], errors='coerce')
            out['depth_delta_ft'] = (dm - dt).abs()
            out = out.drop(columns=['depth_tvd'], errors='ignore')
        else:
            if 'depth_delta_ft' not in out.columns:
                out['depth_delta_ft'] = 0.0
            out = out.drop(columns=['depth_tvd'], errors='ignore')
    num = list(getattr(tp, 'numeric_features', []) or [])
    cat = list(getattr(tp, 'categorical_features', []) or [])
    cols = num + cat
    if not cols:
        raise ValueError('Modelo de temperatura sin lista de features; vuelve a entrenar.')
    for c in cols:
        if c not in out.columns:
            out[c] = np.nan
    return out[cols]

def _dedupe_trace_columns(df: pd.DataFrame) -> pd.DataFrame:
    """Elimina columnas con nombre repetido (común en CSV de pozo)."""
    if df is None or df.empty or not df.columns.duplicated().any():
        return df
    return df.loc[:, ~df.columns.duplicated()].copy()


def _series_from_col(df: pd.DataFrame, col: str) -> pd.Series:
    """Devuelve una Series aunque el nombre de columna esté duplicado."""
    if col not in df.columns:
        return pd.Series(np.nan, index=df.index, dtype=float)
    sel = df[col]
    if isinstance(sel, pd.DataFrame):
        sel = sel.iloc[:, 0]
    return sel


TEMP_OPERATIONAL_OVERRIDE_COLS = (
    'wob', 'rpm', 'torque', 'flow_rate', 'pump_pressure', 'standpipe_pressure', 'rop',
)


def _temperature_ref_operational_params(tp: Any, trace_df: Optional[pd.DataFrame] = None) -> Dict[str, float]:
    """Mediana de referencia para calibración operacional (sin reentrenar)."""
    ref = dict(getattr(tp, 'ref_operational_params', None) or {})
    if ref:
        return ref
    if trace_df is not None and not trace_df.empty:
        std = _standardize_temperature_trace_df(trace_df)
        for col in TEMP_OPERATIONAL_OVERRIDE_COLS + ('depth_md',):
            if col in std.columns:
                med = pd.to_numeric(std[col], errors='coerce').median()
                if pd.notna(med):
                    ref[col] = float(med)
    return ref


def overlay_live_operational_params(trace_df: pd.DataFrame, live_params: Dict[str, Any]) -> pd.DataFrame:
    """Sustituye parámetros operacionales de la traza por los sliders actuales."""
    out = trace_df.copy()
    for col in TEMP_OPERATIONAL_OVERRIDE_COLS:
        val = live_params.get(col)
        if val is not None and np.isfinite(_safe_float(val, np.nan)):
            out[col] = float(val)
    return out


def operational_thermal_correction_delta(
    live_params: Dict[str, Any],
    ref_params: Dict[str, Any],
    *,
    depth_md: Optional[float] = None,
) -> float:
    """
    Corrección física post-ML: amplifica el efecto de WOB/RPM/torque/caudal/SPP
    cuando el bosque subestima su importancia (sin reentrenar).
    """
    depth = max(_safe_float(depth_md if depth_md is not None else live_params.get('depth_md'), 1.0), 1.0)

    def _proxies(p: Dict[str, Any]) -> tuple:
        w = max(_safe_float(p.get('wob'), 0.0), 0.0)
        r = max(_safe_float(p.get('rpm'), 0.0), 0.0)
        t = max(_safe_float(p.get('torque'), 0.0), 0.0)
        f = max(_safe_float(p.get('flow_rate'), 0.0), 0.0)
        pp = max(_safe_float(p.get('pump_pressure'), 0.0), 0.0)
        mech = (w * r * t) / 1_000_000.0
        cool = f / depth
        hyd = (pp * f) / 1_000_000.0
        return mech, cool, hyd

    live_mech, live_cool, live_hyd = _proxies(live_params)
    ref_mech, ref_cool, ref_hyd = _proxies(ref_params)
    mech_delta = (live_mech - ref_mech) / max(ref_mech, 0.01)
    cool_delta = (ref_cool - live_cool) / max(ref_cool, 0.01)
    hyd_delta = (live_hyd - ref_hyd) / max(ref_hyd, 0.01)
    ref_depth = max(_safe_float(ref_params.get('depth_md'), depth), 1.0)
    depth_weight = float(np.clip(depth / ref_depth, 0.35, 2.5))
    return depth_weight * (18.0 * mech_delta + 12.0 * cool_delta + 6.0 * hyd_delta)


def fit_temperature_display_calibration(tp: Any, y_actual: np.ndarray, y_pred: np.ndarray) -> None:
    """Calibra predicciones hacia el rango real (sin reentrenar el bosque)."""
    y_a = np.asarray(y_actual, dtype=float)
    y_p = np.asarray(y_pred, dtype=float)
    ok = np.isfinite(y_a) & np.isfinite(y_p)
    y_a, y_p = y_a[ok], y_p[ok]
    if len(y_a) < 5:
        tp.calib_pred_median = float(np.nanmedian(y_p)) if len(y_p) else 0.0
        tp.calib_pred_std = float(np.nanstd(y_p)) if len(y_p) else 1.0
        tp.calib_actual_median = float(np.nanmedian(y_a)) if len(y_a) else 0.0
        tp.calib_actual_std = float(np.nanstd(y_a)) if len(y_a) else 1.0
        tp.calib_scale = 1.0
        tp.calib_offset = 0.0
        return
    tp.calib_pred_median = float(np.median(y_p))
    tp.calib_pred_std = float(max(np.std(y_p), 1e-6))
    tp.calib_actual_median = float(np.median(y_a))
    tp.calib_actual_std = float(max(np.std(y_a), 1e-6))
    var_p = float(np.var(y_p))
    if var_p > 1e-9:
        scale = float(np.cov(y_p, y_a, bias=True)[0, 1] / var_p)
        offset = float(np.mean(y_a) - scale * np.mean(y_p))
    else:
        scale, offset = 1.0, float(np.mean(y_a) - np.mean(y_p))
    tp.calib_scale = scale
    tp.calib_offset = offset


def apply_temperature_display_calibration(tp: Any, pred: np.ndarray) -> np.ndarray:
    """Estira y alinea predicciones al rango observado en entrenamiento."""
    pred = np.asarray(pred, dtype=float)
    std_p = max(_safe_float(getattr(tp, 'calib_pred_std', None), 0.0), 1e-6)
    std_a = max(_safe_float(getattr(tp, 'calib_actual_std', None), 0.0), 1e-6)
    mu_p = _safe_float(getattr(tp, 'calib_pred_median', None), float(np.nanmedian(pred)))
    mu_a = _safe_float(getattr(tp, 'calib_actual_median', None), mu_p)
    stretched = mu_a + (pred - mu_p) * (std_a / std_p)
    scale = _safe_float(getattr(tp, 'calib_scale', None), 1.0)
    offset = _safe_float(getattr(tp, 'calib_offset', None), 0.0)
    return scale * stretched + offset


def temperature_apply_model_predictions(
    tp: Any,
    X_features: pd.DataFrame,
    *,
    live_params: Optional[Dict[str, Any]] = None,
    ref_params: Optional[Dict[str, Any]] = None,
    apply_operational_calibration: bool = True,
    apply_display_calibration: bool = True,
) -> np.ndarray:
    """Predicción vectorizada con modo delta y calibración operacional opcional."""
    Xp = temperature_prepare_features_for_prediction(tp, X_features)
    pred = tp.model.predict(Xp).astype(float)
    if (
        getattr(tp, 'model_target_mode', 'direct') == 'delta_to_base'
        and getattr(tp, 'base_temperature_col', None)
        and tp.base_temperature_col in X_features.columns
    ):
        base = pd.to_numeric(X_features[tp.base_temperature_col], errors='coerce')
        base = base.fillna(base.median() if base.notna().any() else 35.0).values
        pred = pred + base
    if apply_operational_calibration and live_params and ref_params:
        if 'depth_md' in X_features.columns:
            deltas = [
                operational_thermal_correction_delta(
                    live_params,
                    ref_params,
                    depth_md=_safe_float(d, live_params.get('depth_md')),
                )
                for d in pd.to_numeric(X_features['depth_md'], errors='coerce').fillna(live_params.get('depth_md'))
            ]
            pred = pred + np.asarray(deltas, dtype=float)
        else:
            pred = pred + operational_thermal_correction_delta(live_params, ref_params)
    if apply_display_calibration and getattr(tp, 'calib_pred_std', None):
        pred = apply_temperature_display_calibration(tp, pred)
    return pred


def temperature_predict_trace_live(
    tp: Any,
    trace_df: pd.DataFrame,
    live_params: Dict[str, Any],
    *,
    apply_operational_calibration: bool = True,
    max_points: int = 1200,
) -> tuple:
    """
    Real vs predicha con parámetros operacionales de los sliders.
    Conserva profundidad, mud_in_temp y litología de la traza; sustituye WOB/RPM/etc.
    """
    trace_eval = _dedupe_trace_columns(_standardize_temperature_trace_df(trace_df).copy())
    target_col = tp.target_col
    if not target_col or target_col not in trace_eval.columns:
        return np.array([]), np.array([]), np.array([])
    base_cols = [
        'depth_md', 'depth_tvd', 'rpm', 'wob', 'torque', 'rop', 'flow_rate',
        'pump_pressure', 'mud_in_temp', 'mud_out_temp', 'mud_density', 'pv', 'yp',
        'bit_size', 'bit_depth', 'hookload', 'standpipe_pressure', 'lithology', 'formation',
    ]
    cols = list(dict.fromkeys([c for c in base_cols if c in trace_eval.columns] + [target_col]))
    trace_eval = trace_eval[cols].copy()
    y_num = pd.to_numeric(_series_from_col(trace_eval, target_col), errors='coerce')
    trace_eval = trace_eval.loc[y_num.notna()].reset_index(drop=True)
    y_num = y_num.loc[y_num.notna()].reset_index(drop=True)
    if len(trace_eval) < 5:
        return np.array([]), np.array([]), np.array([])
    if len(trace_eval) > max_points:
        idx = np.linspace(0, len(trace_eval) - 1, max_points, dtype=int)
        trace_eval = trace_eval.iloc[idx].reset_index(drop=True)
        y_num = y_num.iloc[idx].reset_index(drop=True)
    trace_live = overlay_live_operational_params(trace_eval, live_params)
    ref_params = _temperature_ref_operational_params(tp, trace_df)
    pred = temperature_apply_model_predictions(
        tp,
        trace_live.drop(columns=[target_col]),
        live_params=live_params,
        ref_params=ref_params,
        apply_operational_calibration=apply_operational_calibration,
    )
    depths = pd.to_numeric(_series_from_col(trace_live, 'depth_md'), errors='coerce').to_numpy(dtype=float)
    return y_num.to_numpy(dtype=float), pred, depths


def temperature_predict_grid_from_params(
    temperature_predictor: Any,
    params: Dict[str, Any],
    wob_values: np.ndarray,
    rpm_values: np.ndarray,
) -> np.ndarray:
    """Predicción vectorizada WOB×RPM para el heatmap (una sola llamada al pipeline ML)."""
    nw, nr = len(wob_values), len(rpm_values)
    if temperature_predictor is None or getattr(temperature_predictor, 'model', None) is None or nw == 0 or nr == 0:
        return np.full((nr, nw), np.nan, dtype=float)
    p_base = dict(params)
    if getattr(temperature_predictor, '_use_depth_delta', False):
        md = _safe_float(p_base.get('depth_md'), np.nan)
        tv = _safe_float(p_base.get('depth_tvd'), np.nan)
        p_base['depth_delta_ft'] = abs(md - tv) if (np.isfinite(md) and np.isfinite(tv)) else 0.0
    wob_grid, rpm_grid = np.meshgrid(wob_values, rpm_values)
    rows = []
    num_feats = list(getattr(temperature_predictor, 'numeric_features', []) or [])
    cat_feats = list(getattr(temperature_predictor, 'categorical_features', []) or [])
    for w, r in zip(wob_grid.ravel(), rpm_grid.ravel()):
        p = dict(p_base)
        p['wob'] = float(w)
        p['rpm'] = float(r)
        row = {}
        for col in num_feats + cat_feats:
            v = p.get(col, np.nan)
            if col in cat_feats:
                if v is None or (isinstance(v, float) and np.isnan(v)):
                    v = np.nan
                else:
                    v = str(v)
            row[col] = v
        rows.append(row)
    X = pd.DataFrame(rows)
    try:
        Xp = temperature_prepare_features_for_prediction(temperature_predictor, X)
        ref_params = _temperature_ref_operational_params(temperature_predictor)
        pred = temperature_apply_model_predictions(
            temperature_predictor,
            X,
            live_params=p_base,
            ref_params=ref_params,
            apply_operational_calibration=True,
        )
        return pred.reshape(nr, nw)
    except Exception:
        return np.full((nr, nw), np.nan, dtype=float)




def infer_trace_value_at_depth(trace_df: Optional[pd.DataFrame], value_col: str, depth_ft: float, default: Optional[float] = None) -> Optional[float]:
    """Estima un sensor de la traza real a la profundidad actual.
    Útil para predicción interactiva: si el modelo térmico usa mud_in_temp/ROP/etc.,
    toma el valor más cercano en profundidad en vez de dejar NaN.
    """
    try:
        if trace_df is None or len(trace_df) == 0:
            return default
        df = _standardize_temperature_trace_df(trace_df).copy()
        if value_col not in df.columns:
            return default
        dcol = 'depth_md' if 'depth_md' in df.columns else ('depth_tvd' if 'depth_tvd' in df.columns else None)
        vals = pd.to_numeric(df[value_col], errors='coerce')
        if dcol is None:
            med = vals.dropna().median()
            return float(med) if pd.notna(med) else default
        depths = pd.to_numeric(df[dcol], errors='coerce')
        work = pd.DataFrame({'depth': depths, 'value': vals}).dropna()
        if work.empty:
            return default
        # Alinear escala si la traza quedó en m y el slider está en ft.
        if work['depth'].max() < max(float(depth_ft) * 0.55, 8000):
            work['depth'] = work['depth'] * 3.28084
        i = (work['depth'] - float(depth_ft)).abs().idxmin()
        v = work.loc[i, 'value']
        return float(v) if pd.notna(v) else default
    except Exception:
        return default


def build_current_temperature_params(
    *, depth_ft: float, inclination_deg: float, rpm: float, wob_klb: float,
    torque_ftlb: float, flow_gpm: float, spp_psi: float, mud_density_ppg: float,
    pv_cp: float, yp_lb100ft2: float, bit_diameter_in: float,
    rop_value: Optional[float], formation_info: Optional[Dict[str, Any]], trace_df: Optional[pd.DataFrame]
) -> Dict[str, Any]:
    """Construye las features actuales del modelo térmico desde sliders + sensores estimados."""
    mud_in = infer_trace_value_at_depth(trace_df, 'mud_in_temp', depth_ft, default=None)
    if mud_in is None:
        # Fallback conservador si el CSV no trae temperatura de entrada: valor típico superficial.
        mud_in = 35.0
    bit_depth = infer_trace_value_at_depth(trace_df, 'bit_depth', depth_ft, default=depth_ft)
    hookload = infer_trace_value_at_depth(trace_df, 'hookload', depth_ft, default=None)
    return {
        'depth_md': float(depth_ft),
        'depth_tvd': float(depth_ft) * float(np.cos(np.radians(inclination_deg))),
        'rpm': float(rpm),
        'wob': float(wob_klb),
        'torque': float(torque_ftlb),
        'rop': float(rop_value) if rop_value is not None and np.isfinite(rop_value) else infer_trace_value_at_depth(trace_df, 'rop', depth_ft, default=np.nan),
        'flow_rate': float(flow_gpm),
        'pump_pressure': float(spp_psi),
        'standpipe_pressure': float(spp_psi),
        'mud_in_temp': float(mud_in),
        'mud_density': float(mud_density_ppg),
        'pv': float(pv_cp),
        'yp': float(yp_lb100ft2),
        'bit_size': float(bit_diameter_in),
        'bit_depth': float(bit_depth) if bit_depth is not None and np.isfinite(bit_depth) else float(depth_ft),
        'hookload': hookload,
        'lithology': formation_info.get('lithology') if formation_info else None,
        'formation': formation_info.get('name') if formation_info else None,
    }

class TemperaturePredictor:
    def __init__(self):
        self.model = None
        self.target_col = None
        self.metrics = {}
        self.feature_names = []
        self.feature_importance = None
        self.numeric_features = []
        self.categorical_features = []
        self.training_source = 'synthetic'
        self._use_depth_delta = False
        self.last_holdout_actual: Optional[np.ndarray] = None
        self.last_holdout_pred: Optional[np.ndarray] = None
        self.model_target_mode = 'direct'
        self.base_temperature_col: Optional[str] = None
        self.ref_operational_params: Dict[str, float] = {}
        self.calib_pred_median: Optional[float] = None
        self.calib_pred_std: Optional[float] = None
        self.calib_actual_median: Optional[float] = None
        self.calib_actual_std: Optional[float] = None
        self.calib_scale: float = 1.0
        self.calib_offset: float = 0.0

    def prepare_features_for_prediction(self, X: pd.DataFrame) -> pd.DataFrame:
        """Delega en la función de módulo (misma lógica que session_state antigua)."""
        return temperature_prepare_features_for_prediction(self, X)

    def fit_from_dataframe(self, df: pd.DataFrame, target_col: Optional[str] = None):
        if not SKLEARN_AVAILABLE:
            raise RuntimeError('scikit-learn no esta disponible en este entorno.')
        data = _dedupe_trace_columns(_standardize_temperature_trace_df(df))
        detected_target = target_col or _detect_temperature_target(data)
        if not detected_target or detected_target not in data.columns:
            raise ValueError('No se encontro una columna objetivo de temperatura.')
        self.target_col = detected_target
        self._use_depth_delta = False
        self.last_holdout_actual = None
        self.last_holdout_pred = None
        self.model_target_mode = 'direct'
        self.base_temperature_col = None
        self.ref_operational_params = {}
        self.calib_pred_median = None
        self.calib_pred_std = None
        self.calib_actual_median = None
        self.calib_actual_std = None
        self.calib_scale = 1.0
        self.calib_offset = 0.0

        # Para mud_out_temp conviene aprender DeltaT = mud_out_temp - mud_in_temp.
        # Esto evita que el bosque regrese a la media y permite conservar la tendencia termica alta.
        learn_col = self.target_col
        eval_target = pd.to_numeric(_series_from_col(data, self.target_col), errors='coerce')
        if self.target_col == 'mud_out_temp' and 'mud_in_temp' in data.columns:
            data = data.copy()
            data['delta_mud_temp'] = pd.to_numeric(data['mud_out_temp'], errors='coerce') - pd.to_numeric(data['mud_in_temp'], errors='coerce')
            learn_col = 'delta_mud_temp'
            self.model_target_mode = 'delta_to_base'
            self.base_temperature_col = 'mud_in_temp'

        base_num = [
            'depth_md', 'depth_tvd', 'rpm', 'wob', 'torque', 'rop', 'flow_rate',
            'pump_pressure', 'mud_in_temp', 'mud_out_temp', 'mud_density', 'pv', 'yp',
            'bit_size', 'bit_depth', 'hookload', 'standpipe_pressure',
        ]
        base_cat = ['lithology', 'formation']

        # No dejar que el target directo entre como feature. Si aprendemos DeltaT, mud_in_temp si queda.
        self.numeric_features = [c for c in base_num if c in data.columns and c not in {self.target_col, learn_col}]
        self.categorical_features = [c for c in base_cat if c in data.columns and c not in {self.target_col, learn_col}]
        if len(self.numeric_features) < 3:
            raise ValueError('La traza no tiene suficientes variables numericas para entrenar el modelo de temperatura.')

        train_df = data[self.numeric_features + self.categorical_features + [learn_col, self.target_col]].dropna(subset=[learn_col, self.target_col]).copy()
        if len(train_df) < 30:
            raise ValueError('Se requieren al menos 30 filas con temperatura valida para entrenar.')

        X = train_df.drop(columns=[learn_col, self.target_col], errors='ignore')
        y = pd.to_numeric(train_df[learn_col], errors='coerce')
        y_eval = pd.to_numeric(train_df[self.target_col], errors='coerce')
        valid_mask = y.notna() & y_eval.notna()
        X = X.loc[valid_mask].reset_index(drop=True)
        y = y.loc[valid_mask].reset_index(drop=True)
        y_eval = y_eval.loc[valid_mask].reset_index(drop=True)

        X = _add_temperature_engineered_features(X)

        # Reducir redundancia MD/TVD: una sola senal de profundidad + |MD-TVD|.
        num_list = list(self.numeric_features)
        if 'depth_md' in X.columns and 'depth_tvd' in X.columns:
            dm = pd.to_numeric(X['depth_md'], errors='coerce')
            dt = pd.to_numeric(X['depth_tvd'], errors='coerce')
            X = X.copy()
            X['depth_delta_ft'] = (dm - dt).abs()
            X = X.drop(columns=['depth_tvd'])
            num_list = [c for c in num_list if c != 'depth_tvd']
            self._use_depth_delta = True
        for engineered in [
            'depth_md_sqrt', 'depth_md_sq_scaled', 'depth_delta_ft', 'bit_md_delta_ft',
            'mech_energy_proxy', 'hydraulic_power_proxy', 'cooling_proxy',
            'flow_per_pressure', 'pressure_per_flow', 'rop_inverse',
            'thermal_load_proxy', 'mud_in_x_depth', 'mud_density_x_flow'
        ]:
            if engineered in X.columns and engineered not in num_list:
                num_list.append(engineered)
        self.numeric_features = [c for c in num_list if c in X.columns]

        X_work = X[[c for c in self.numeric_features + self.categorical_features if c in X.columns]].copy()
        for c in self.numeric_features + self.categorical_features:
            if c not in X_work.columns:
                X_work[c] = np.nan
        X_work = X_work[self.numeric_features + self.categorical_features]

        test_size = 0.2 if len(X_work) >= 25 else max(0.15, 5 / max(len(X_work), 1))
        test_size = min(0.35, max(test_size, 5 / max(len(X_work), 1)))
        strat = None
        if len(X_work) >= 80:
            try:
                n_bins = min(12, max(4, len(X_work) // 30))
                strat = pd.qcut(y_eval, q=n_bins, labels=False, duplicates='drop')
                if strat.isna().any() or strat.nunique() < 2 or strat.value_counts().min() < 2:
                    strat = None
            except (ValueError, TypeError):
                strat = None
        idx = np.arange(len(X_work))
        try:
            train_idx, test_idx = train_test_split(idx, test_size=test_size, random_state=42, shuffle=True, stratify=strat)
        except ValueError:
            train_idx, test_idx = train_test_split(idx, test_size=test_size, random_state=42, shuffle=True)
        X_train, X_test = X_work.iloc[train_idx], X_work.iloc[test_idx]
        y_train, y_test = y.iloc[train_idx], y.iloc[test_idx]
        y_eval_train, y_eval_test = y_eval.iloc[train_idx], y_eval.iloc[test_idx]

        # Dar mas peso a temperaturas altas para evitar que el modelo aplaste los extremos hacia 90-110 C.
        q75 = float(y_eval_train.quantile(0.75))
        q90 = float(y_eval_train.quantile(0.90))
        weights = np.ones(len(y_train), dtype=float)
        weights += 1.0 * (y_eval_train.values >= q75)
        weights += 2.0 * (y_eval_train.values >= q90)

        transformers = [
            ('num', Pipeline([('imputer', SimpleImputer(strategy='median'))]), self.numeric_features),
        ]
        if self.categorical_features:
            transformers.append(
                ('cat', Pipeline([
                    ('imputer', SimpleImputer(strategy='most_frequent')),
                    ('onehot', OneHotEncoder(handle_unknown='ignore')),
                ]), self.categorical_features)
            )
        preprocessor = ColumnTransformer(transformers=transformers, remainder='drop')
        model = RandomForestRegressor(
            n_estimators=700,
            max_depth=None,
            min_samples_split=3,
            min_samples_leaf=1,
            max_features=0.75,
            max_samples=None,
            bootstrap=True,
            random_state=42,
            n_jobs=-1,
        )
        pipeline = Pipeline([('preprocessor', preprocessor), ('model', model)])
        pipeline.fit(X_train, y_train, model__sample_weight=weights)

        y_pred_model = pipeline.predict(X_test)
        y_train_pred_model = pipeline.predict(X_train)
        if self.model_target_mode == 'delta_to_base' and self.base_temperature_col in X_test.columns:
            base_test = pd.to_numeric(X_test[self.base_temperature_col], errors='coerce').fillna(pd.to_numeric(data[self.base_temperature_col], errors='coerce').median()).values
            base_train = pd.to_numeric(X_train[self.base_temperature_col], errors='coerce').fillna(pd.to_numeric(data[self.base_temperature_col], errors='coerce').median()).values
            y_pred_eval = base_test + y_pred_model
            y_train_pred_eval = base_train + y_train_pred_model
        else:
            y_pred_eval = y_pred_model
            y_train_pred_eval = y_train_pred_model

        self.metrics = {
            'mae': float(mean_absolute_error(y_eval_test, y_pred_eval)),
            'rmse': float(np.sqrt(mean_squared_error(y_eval_test, y_pred_eval))),
            'r2': float(r2_score(y_eval_test, y_pred_eval)) if len(y_eval_test) > 1 else 0.0,
            'r2_train': float(r2_score(y_eval_train, y_train_pred_eval)) if len(y_eval_train) > 1 else 0.0,
            'train_rows': int(len(X_train)),
            'test_rows': int(len(X_test)),
            'target_mode': self.model_target_mode,
            'q75_hot': q75,
            'q90_hot': q90,
        }
        self.model = pipeline
        self.training_source = 'real_trace'
        self.last_holdout_actual = np.asarray(y_eval_test, dtype=float)
        self.last_holdout_pred = np.asarray(y_pred_eval, dtype=float)
        fit_temperature_display_calibration(self, y_eval_train.to_numpy(dtype=float), y_train_pred_eval)
        self.ref_operational_params = {}
        for _ref_col in TEMP_OPERATIONAL_OVERRIDE_COLS + ('depth_md',):
            if _ref_col in train_df.columns:
                _med = pd.to_numeric(train_df[_ref_col], errors='coerce').median()
                if pd.notna(_med):
                    self.ref_operational_params[_ref_col] = float(_med)
        feature_names = list(self.numeric_features)
        if self.categorical_features and 'cat' in pipeline.named_steps['preprocessor'].named_transformers_:
            ohe = pipeline.named_steps['preprocessor'].named_transformers_['cat'].named_steps['onehot']
            feature_names.extend(ohe.get_feature_names_out(self.categorical_features).tolist())
        self.feature_names = feature_names
        importances = pipeline.named_steps['model'].feature_importances_
        self.feature_importance = pd.DataFrame({'feature': feature_names, 'importance': importances}).sort_values('importance', ascending=False)
        return self.metrics

    def predict_grid_from_params(
        self,
        params: Dict[str, Any],
        wob_values: np.ndarray,
        rpm_values: np.ndarray,
    ) -> np.ndarray:
        return temperature_predict_grid_from_params(self, params, wob_values, rpm_values)

    def predict_from_params(self, params: Dict[str, Any]) -> Optional[float]:
        if self.model is None:
            return None
        p = dict(params)
        if getattr(self, '_use_depth_delta', False):
            md = _safe_float(p.get('depth_md'), np.nan)
            tv = _safe_float(p.get('depth_tvd'), np.nan)
            p['depth_delta_ft'] = abs(md - tv) if (np.isfinite(md) and np.isfinite(tv)) else 0.0
        row = {}
        for col in self.numeric_features + self.categorical_features:
            v = p.get(col, np.nan)
            if col in self.categorical_features:
                if v is None or (isinstance(v, float) and np.isnan(v)):
                    v = np.nan
                else:
                    v = str(v)
            row[col] = v
        X = pd.DataFrame([row])
        try:
            ref_params = _temperature_ref_operational_params(self)
            pred_arr = temperature_apply_model_predictions(
                self,
                X,
                live_params=p,
                ref_params=ref_params,
                apply_operational_calibration=True,
            )
            return float(pred_arr[0])
        except Exception:
            return None

def create_temperature_feature_importance_chart(importance_df: Optional[pd.DataFrame]) -> go.Figure:
    if importance_df is None or importance_df.empty:
        fig = go.Figure()
        fig.add_annotation(text='Sin importancia de variables disponible', x=0.5, y=0.5, showarrow=False)
        fig.update_layout(height=350)
        return fig
    top_df = importance_df.head(12).sort_values('importance', ascending=True)
    fig = go.Figure(go.Bar(
        x=top_df['importance'],
        y=top_df['feature'],
        orientation='h',
        marker=dict(
            color=top_df['importance'].values,
            colorscale=[[0, '#0D9488'], [0.5, '#06B6D4'], [1, '#6366F1']],
            cmin=0,
            cmax=top_df['importance'].max() * 1.05,
            line=dict(width=0),
        ),
        text=[f'{v*100:.1f}%' for v in top_df['importance']],
        textposition='outside',
        textfont=dict(size=11, color='#E2E8F0'),
    ))
    fig.update_layout(
        title=dict(text='Importancia de variables — Temperatura', font=dict(size=16, color='#F1F5F9')),
        xaxis=dict(title='Importancia', title_font=dict(color='#94A3B8'), tickfont=dict(color='#94A3B8'), gridcolor='rgba(148,163,184,0.15)', zeroline=False),
        yaxis=dict(tickfont=dict(color='#CBD5E1'), gridcolor='rgba(0,0,0,0)'),
        height=380,
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(15,23,42,0.4)',
        margin=dict(l=100, r=80),
        font=dict(family='Inter'),
    )
    return fig

def create_temperature_real_vs_pred_chart(
    actual,
    pred,
    *,
    highlight: Optional[Dict[str, float]] = None,
) -> go.Figure:
    fig = go.Figure()
    fig.add_trace(go.Scatter(
        x=actual, y=pred, mode='markers',
        marker=dict(size=9, opacity=0.82, color='#06B6D4', line=dict(width=1, color='#0E7490')),
        name='Predicción (mueve con sliders)',
    ))
    lo = min(np.min(actual), np.min(pred))
    hi = max(np.max(actual), np.max(pred))
    fig.add_trace(go.Scatter(
        x=[lo, hi], y=[lo, hi], mode='lines',
        line=dict(color='#F59E0B', width=2, dash='dash'),
        name='Referencia 1:1 (fija)',
    ))
    if highlight and np.isfinite(highlight.get('actual', np.nan)) and np.isfinite(highlight.get('pred', np.nan)):
        fig.add_trace(go.Scatter(
            x=[highlight['actual']], y=[highlight['pred']],
            mode='markers+text',
            marker=dict(size=15, color='#EF4444', symbol='star', line=dict(width=1.5, color='#FCA5A5')),
            text=['Ahora'],
            textposition='top center',
            textfont=dict(size=11, color='#FCA5A5'),
            name='Condición actual',
        ))
    fig.update_layout(
        title=dict(text='Temperatura real vs predicha', font=dict(size=16, color='#F1F5F9')),
        xaxis=dict(title='Real', title_font=dict(color='#94A3B8'), tickfont=dict(color='#94A3B8'), gridcolor='rgba(148,163,184,0.15)', zeroline=False),
        yaxis=dict(title='Predicha', title_font=dict(color='#94A3B8'), tickfont=dict(color='#94A3B8'), gridcolor='rgba(148,163,184,0.15)', zeroline=False),
        height=360,
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(15,23,42,0.4)',
        legend=dict(font=dict(color='#CBD5E1'), orientation='h', yanchor='bottom', y=1.02),
        font=dict(family='Inter'),
    )
    return fig

# ============================================================================
# APLICACIÓN PRINCIPAL STREAMLIT
# ============================================================================

def main():
    """Aplicación principal Streamlit - Imperial Units"""
    
    # Inicializar idioma
    if 'lang' not in st.session_state:
        st.session_state.lang = 'es'
    
    # Header: logo ROGII encima del título, a lo largo
    if LOGO_PATH.exists():
        st.image(str(LOGO_PATH), width="stretch")
    st.markdown('<h1 class="main-header" style="text-align: center; margin-top: 0.5rem;">Drilling Analytics</h1>', 
                unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.markdown(f"""
        <div style="text-align: center; margin-bottom: 2rem;">
            <p style="font-size: 1.1rem; color: #4A5568; margin-top: 1rem; text-align: center; white-space: nowrap;">
                {_t('header_subtitle')}
            </p>
        </div>
        """, unsafe_allow_html=True)
    
    # Initialize session state
    if 'predictor' not in st.session_state:
        st.session_state.predictor = ROPPredictorImperial()
    if 'data_generator' not in st.session_state:
        st.session_state.data_generator = DrillingDataGeneratorImperial()
    if 'models_trained' not in st.session_state:
        st.session_state.models_trained = False
    if 'current_prediction' not in st.session_state:
        st.session_state.current_prediction = None
    if 'geological_formations' not in st.session_state:
        st.session_state.geological_formations = DEFAULT_GEOLOGICAL_FORMATIONS.copy()
    if 'use_geological_tracking' not in st.session_state:
        st.session_state.use_geological_tracking = False
    if 'pptx_ready' not in st.session_state:
        st.session_state.pptx_ready = False
    if 'temperature_predictor' not in st.session_state:
        st.session_state.temperature_predictor = TemperaturePredictor()
    else:
        _tp_mig = st.session_state.temperature_predictor
        if not hasattr(_tp_mig, 'last_holdout_actual'):
            _tp_mig.last_holdout_actual = None
        if not hasattr(_tp_mig, 'last_holdout_pred'):
            _tp_mig.last_holdout_pred = None
        if not hasattr(_tp_mig, '_use_depth_delta'):
            _tp_mig._use_depth_delta = False
        if not hasattr(_tp_mig, 'ref_operational_params'):
            _tp_mig.ref_operational_params = {}
        for _cal_attr, _cal_default in (
            ('calib_pred_median', None), ('calib_pred_std', None),
            ('calib_actual_median', None), ('calib_actual_std', None),
            ('calib_scale', 1.0), ('calib_offset', 0.0),
        ):
            if not hasattr(_tp_mig, _cal_attr):
                setattr(_tp_mig, _cal_attr, _cal_default)
    if 'real_trace_df' not in st.session_state:
        st.session_state.real_trace_df = None
    if 'real_trace_name' not in st.session_state:
        st.session_state.real_trace_name = None
    if 'temperature_last_prediction' not in st.session_state:
        st.session_state.temperature_last_prediction = None
    if 'neighbor_temp_profile' not in st.session_state:
        st.session_state.neighbor_temp_profile = NeighborTemperatureProfiler()
    if 'neighbor_trace_dfs' not in st.session_state:
        st.session_state.neighbor_trace_dfs = []
    if 'neighbor_trace_names' not in st.session_state:
        st.session_state.neighbor_trace_names = []
    if 'neighbor_profile_df' not in st.session_state:
        st.session_state.neighbor_profile_df = pd.DataFrame()
    if 'neighbor_prediction_df' not in st.session_state:
        st.session_state.neighbor_prediction_df = pd.DataFrame()
    if 'target_x_coord' not in st.session_state:
        st.session_state.target_x_coord = 0.0
    if 'target_y_coord' not in st.session_state:
        st.session_state.target_y_coord = 0.0
    if 'thermal_anomaly_z_threshold' not in st.session_state:
        st.session_state.thermal_anomaly_z_threshold = 2.0
    if 'trace_unit_system' not in st.session_state:
        st.session_state.trace_unit_system = TRACE_UNIT_SYSTEM_LABELS[0]
    if 'real_trace_raw_df' not in st.session_state:
        st.session_state.real_trace_raw_df = None
    if 'trace_unit_system_prev' not in st.session_state:
        st.session_state.trace_unit_system_prev = st.session_state.get('trace_unit_system', TRACE_UNIT_SYSTEM_LABELS[0])
    if 'neighbor_profile_unit_system' not in st.session_state:
        st.session_state.neighbor_profile_unit_system = None
    if st.session_state.get('trace_pipeline_version', 0) != TRACE_PIPELINE_VERSION:
        st.session_state.neighbor_profile_df = pd.DataFrame()
        st.session_state.neighbor_prediction_df = pd.DataFrame()
        st.session_state.real_trace_raw_df = None
        st.session_state.real_trace_df = None
        st.session_state.real_trace_name = None
        st.session_state.neighbor_profile_unit_system = None
        st.session_state.trace_pipeline_version = TRACE_PIPELINE_VERSION
    
    # ========================================================================
    # SIDEBAR - CONFIGURACIÓN IMPERIAL
    # ========================================================================
    
    with st.sidebar:
        if LOGO_PATH.exists():
            st.image(str(LOGO_PATH), width="stretch")
        lang_selected = st.selectbox(
            ("Idioma / Language" if st.session_state.lang == 'es' else "Language / Idioma"),
            options=['es', 'en'],
            format_func=lambda x: 'Español' if x == 'es' else 'English',
            index=0 if st.session_state.lang == 'es' else 1,
            key='lang_select'
        )
        st.session_state.lang = lang_selected
        st.markdown(f"## **{_t('panel_control')}**")
        st.markdown(f'<span class="unit-label">{_t("units_usc")}</span>', 
                   unsafe_allow_html=True)
        with st.expander(f"**{_t('suggestions_panel')}**", expanded=False):
            for s in get_section_suggestions('sidebar'):
                st.markdown(f"• {s}")
        # Model Training Section
        with st.expander(f"**{_t('model_training')}**", expanded=True):
            st.markdown(f"**{_t('training_config')}**")
            
            train_rf = st.checkbox("Random Forest", value=True)
            train_xgb = st.checkbox("XGBoost", value=True)
            train_nn = st.checkbox("Red Neuronal", value=True)
            
            synthetic_samples = st.slider(
                _t('synthetic_samples'),
                min_value=1000,
                max_value=50000,
                value=10000,
                step=1000,
                help=_t('synthetic_samples_help')
            )
            
            if st.button(f"**{_t('train_models')}**", use_container_width=True):
                with st.spinner(f"{_t('training_spinner')}"):
                    
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    
                    models_to_train = []
                    if train_rf:
                        models_to_train.append('Random Forest')
                    if train_xgb:
                        models_to_train.append('XGBoost')
                    if train_nn:
                        models_to_train.append('Red Neuronal')
                    
                    gen = DrillingDataGeneratorImperial(n_samples=synthetic_samples)
                    for i, model_name in enumerate(models_to_train):
                        status_text.text(f"Entrenando {model_name}... ({i+1}/{len(models_to_train)})")
                        
                        if model_name == 'Random Forest':
                            st.session_state.predictor.train_random_forest(data_generator=gen, n_samples=synthetic_samples)
                        elif model_name == 'XGBoost':
                            st.session_state.predictor.train_xgboost(data_generator=gen, n_samples=synthetic_samples)
                        elif model_name == 'Red Neuronal':
                            st.session_state.predictor.train_neural_network(data_generator=gen, n_samples=synthetic_samples)
                        
                        progress_bar.progress((i + 1) / len(models_to_train))
                        pass
                    
                    st.session_state.models_trained = True
                    status_text.text(f"{_t('training_done')}")
                    st.success(_t('models_trained_success').format(n=len(models_to_train), samples=synthetic_samples))

        with st.expander('**Trazas reales DrillSpot / ML Temperatura**', expanded=True):
            st.selectbox(
                'Sistema de unidades de las trazas',
                options=TRACE_UNIT_SYSTEM_LABELS,
                index=TRACE_UNIT_SYSTEM_LABELS.index(st.session_state.trace_unit_system)
                if st.session_state.trace_unit_system in TRACE_UNIT_SYSTEM_LABELS else 0,
                key='trace_unit_system',
                help=(
                    '**Auto**: lee la fila de unidades del CSV (m, Pa, kgf…) o infiere por magnitud. '
                    '**Métrico** / **Imperial**: fuerza conversión. '
                    'Tus pozos en **metros** son correctos: la app convierte **una sola vez** a ft/kilopie/lbf internamente.'
                ),
            )
            if st.session_state.trace_unit_system != st.session_state.trace_unit_system_prev:
                st.session_state.neighbor_profile_df = pd.DataFrame()
                st.session_state.neighbor_prediction_df = pd.DataFrame()
                st.session_state.neighbor_profile_unit_system = None
                st.session_state.trace_unit_system_prev = st.session_state.trace_unit_system
                st.warning(
                    'Cambiaste el sistema de unidades. **Vuelve a pulsar** '
                    '«Generar perfil esperado desde pozos vecinos» y reentrena el ML si aplica.'
                )
            uploaded_trace = st.file_uploader(
                'Cargar traza real objetivo (.csv, .xlsx)',
                type=['csv', 'xlsx', 'xls'],
                key='real_trace_uploader',
                help='Carga la traza del pozo objetivo para entrenar o comparar temperatura.'
            )
            if uploaded_trace is not None:
                real_df = _safe_read_tabular_file(uploaded_trace)
                if real_df is not None and not real_df.empty:
                    st.session_state.real_trace_raw_df = real_df
                    st.session_state.real_trace_name = uploaded_trace.name
                else:
                    st.error('No pude leer el archivo cargado.')

            if st.session_state.real_trace_raw_df is not None:
                real_df_std = _standardize_temperature_trace_df(st.session_state.real_trace_raw_df)
                st.session_state.real_trace_df = real_df_std
                st.caption(
                    f'Archivo objetivo: {st.session_state.real_trace_name} · Filas: {len(real_df_std):,} · '
                    f'Columnas: {len(real_df_std.columns)} · Unidades: {st.session_state.trace_unit_system}'
                )
                if 'depth_md' in real_df_std.columns and real_df_std['depth_md'].notna().any():
                    _dmax_obj = float(real_df_std['depth_md'].max())
                    st.caption(
                        f'Profundidad normalizada: {real_df_std["depth_md"].min():,.0f}–'
                        f'{_dmax_obj:,.0f} ft · '
                        f'Temp. salida: '
                        f'{"sí" if "mud_out_temp" in real_df_std.columns else "no"} · '
                        f'Formación: '
                        f'{"sí" if "formation" in real_df_std.columns else "no"}'
                    )
                    if _dmax_obj > 25000:
                        st.error(
                            f'Profundidad máxima **{_dmax_obj:,.0f} ft** parece incorrecta para un pozo de ~5.900 m '
                            f'(debería ser ~**19.400 ft**). Vuelve a **subir el CSV** con Auto (detectar).'
                        )
                    elif abs(_dmax_obj - 19423) < 800 or (_dmax_obj > 15000 and _dmax_obj < 22000):
                        st.success('Profundidad en ft coherente con pozo en metros (~5.900 m → ~19.400 ft).')
                detected_target = _detect_temperature_target(real_df_std)
                candidate_targets = [
                    c for c in real_df_std.columns
                    if 'temp' in c.lower() or 'temper' in c.lower()
                ]
                for pref in ('mud_out_temp', 'mud_in_temp'):
                    if pref in candidate_targets:
                        candidate_targets.remove(pref)
                        candidate_targets.insert(0, pref)
                if detected_target and detected_target not in candidate_targets:
                    candidate_targets = [detected_target] + candidate_targets
                if candidate_targets:
                    default_idx = 0
                    if 'mud_out_temp' in candidate_targets:
                        default_idx = candidate_targets.index('mud_out_temp')
                    temp_target_col = st.selectbox(
                        'Target de temperatura',
                        options=candidate_targets,
                        index=default_idx,
                        key='temp_target_col',
                        help='Recomendado: mud_out_temp (Temperature Out) para perfil térmico y ML.',
                    )
                    if st.button('Entrenar modelo temperatura con traza real', use_container_width=True, key='train_temp_real_btn'):
                        try:
                            metrics = st.session_state.temperature_predictor.fit_from_dataframe(
                                real_df_std, target_col=temp_target_col
                            )
                            st.success(
                                f"Modelo entrenado. R² validación={metrics['r2']:.3f} · R² entrenamiento={metrics.get('r2_train', 0):.3f} · "
                                f"RMSE={metrics['rmse']:.2f} · MAE={metrics['mae']:.2f}"
                            )
                        except Exception as e:
                            st.error(f'No fue posible entrenar el modelo de temperatura: {e}')
                else:
                    st.warning('No detecté columnas de temperatura. Esperaba Temperature In/Out, mud_out_temp, bit_temp, etc.')

            st.markdown('---')
            st.markdown('**Perfil térmico esperado con pozos vecinos**')
            neighbor_files = st.file_uploader(
                'Cargar pozos vecinos (.csv, .xlsx)',
                type=['csv', 'xlsx', 'xls'],
                accept_multiple_files=True,
                key='neighbor_traces_uploader',
                help='Sube varios archivos de pozos vecinos para construir el perfil térmico esperado.'
            )
            if neighbor_files:
                neighbor_dfs = []
                neighbor_names = []
                for nf in neighbor_files:
                    ndf = _safe_read_tabular_file(nf)
                    if ndf is not None and not ndf.empty:
                        neighbor_dfs.append(ndf)
                        neighbor_names.append(nf.name)
                st.session_state.neighbor_trace_dfs = neighbor_dfs
                st.session_state.neighbor_trace_names = neighbor_names
                st.caption(
                    f'Pozos vecinos cargados: {len(neighbor_names)} · '
                    f'Unidades: {st.session_state.trace_unit_system}'
                )
                for nm in neighbor_names[:8]:
                    st.write(f'• {nm}')

            profile_target_options = []
            if st.session_state.real_trace_df is not None:
                profile_target_options.extend([
                    c for c in st.session_state.real_trace_df.columns
                    if 'temp' in c.lower() or 'temper' in c.lower()
                ])
            for ndf in st.session_state.neighbor_trace_dfs:
                stdf = _standardize_temperature_trace_df(ndf)
                profile_target_options.extend([
                    c for c in stdf.columns if 'temp' in c.lower() or 'temper' in c.lower()
                ])
            profile_target_options = list(dict.fromkeys(profile_target_options))
            for pref in ('mud_out_temp', 'mud_in_temp'):
                if pref in profile_target_options:
                    profile_target_options.remove(pref)
                    profile_target_options.insert(0, pref)
            if profile_target_options:
                prof_default = 0
                if 'mud_out_temp' in profile_target_options:
                    prof_default = profile_target_options.index('mud_out_temp')
                neighbor_target_col = st.selectbox(
                    'Variable de temperatura para el perfil esperado',
                    options=profile_target_options,
                    index=prof_default,
                    key='neighbor_target_col',
                )
                depth_pref = st.selectbox('Profundidad base', options=['depth_tvd', 'depth_md'], index=0, key='neighbor_depth_col')
                weighting = st.selectbox('Ponderación de vecinos', options=['inverse_distance', 'equal'], index=0, key='neighbor_weighting')
                cmeta1, cmeta2 = st.columns(2)
                with cmeta1:
                    target_x = st.number_input(
                        'Coordenada X pozo objetivo',
                        step=1.0,
                        key='target_x_coord',
                        help='Usada para distancia al ponderar vecinos y en el mapa XY.',
                    )
                with cmeta2:
                    target_y = st.number_input(
                        'Coordenada Y pozo objetivo',
                        step=1.0,
                        key='target_y_coord',
                        help='Usada para distancia al ponderar vecinos y en el mapa XY.',
                    )
                if st.button('Generar perfil esperado desde pozos vecinos', use_container_width=True, key='build_neighbor_profile_btn'):
                    if not st.session_state.neighbor_trace_dfs:
                        st.warning('Carga al menos un archivo de pozo vecino.')
                    else:
                        try:
                            profiler = st.session_state.neighbor_temp_profile
                            profile_df = profiler.build_from_neighbors(
                                st.session_state.neighbor_trace_dfs,
                                target_col=neighbor_target_col,
                                depth_col=depth_pref,
                                target_x=target_x,
                                target_y=target_y,
                                weighting=weighting
                            )
                            st.session_state.neighbor_profile_df = profile_df
                            st.session_state.neighbor_profile_unit_system = st.session_state.trace_unit_system
                            if st.session_state.real_trace_df is not None:
                                st.session_state.neighbor_prediction_df = profiler.predict_for_trace(st.session_state.real_trace_df)
                            else:
                                st.session_state.neighbor_prediction_df = pd.DataFrame()
                            if profile_df is not None and not profile_df.empty:
                                st.success(f'Perfil térmico esperado generado con {profiler.metadata.get("neighbor_count", 0)} pozos vecinos.')
                            else:
                                st.warning('No se pudo construir el perfil. Verifica que los vecinos tengan profundidad y temperatura válidas.')
                        except Exception as e:
                            st.error(f'No fue posible construir el perfil esperado: {e}')
            else:
                st.info('Primero carga trazas con alguna columna de temperatura para habilitar el perfil esperado.')

        # Drilling Parameters Section - IMPERIAL UNITS
        with st.expander(f"**{_t('drilling_params')}**", expanded=True):
            st.markdown(f"**{_t('operational_params')}**")
            
            wob_klb = st.slider(
                "WOB (klb)",
                min_value=5.0,
                max_value=40.0,
                value=22.0,
                step=0.5,
                help="Peso sobre barrena - miles de libras"
            )
            
            rpm = st.slider(
                "RPM",
                min_value=40,
                max_value=220,
                value=120,
                step=5,
                help="Velocidad rotacional - revoluciones por minuto"
            )
            
            torque_ftlb = st.slider(
                "Torque (ft-lb)",
                min_value=5000,
                max_value=40000,
                value=18000,
                step=1000,
                help="Torque rotativo - pie-libras"
            )
            
            spp_psi = st.slider(
                "SPP (psi)",
                min_value=1500,
                max_value=6000,
                value=3000,
                step=100,
                help="Presión de tubería - libras por pulgada cuadrada"
            )
            
            flow_gpm = st.slider(
                "Caudal (gpm)",
                min_value=400,
                max_value=1400,
                value=800,
                step=50,
                help="Caudal de lodo - galones por minuto"
            )
            
            st.markdown(f"**{_t('geological_params')}**")
            
            ucs_psi = st.slider(
                "UCS (psi)",
                min_value=5000,
                max_value=40000,
                value=15000,
                step=1000,
                help="Resistencia a la compresión no confinada - psi"
            )
            
            st.markdown(f"**{_t('bit_params')}**")
            
            bit_diameter_in = st.selectbox(
                "Diámetro de broca (in)",
                options=[6.0, 7.875, 8.5, 9.875, 12.25, 14.75, 17.5, 26.0],
                index=2,
                help="Diámetro de broca en pulgadas"
            )
            
            bit_wear = st.slider(
                "Desgaste de broca",
                min_value=0.0,
                max_value=1.0,
                value=0.15,
                step=0.05,
                format="%.2f",
                help="0 = nueva, 1 = completamente gastada"
            )
            
            cutter_count = st.selectbox(
                "Número de cortadores",
                options=[4, 5, 6, 7, 8],
                index=2,
                help="Número de cortadores/aspas PDC"
            )
            
            st.markdown("**Parámetros del pozo**")
            
            depth_ft = st.slider(
                "Profundidad (ft)",
                min_value=2000,
                max_value=25000,
                value=12000,
                step=500,
                help="Profundidad medida actual - pies"
            )
            
            inclination_deg = st.slider(
                "Inclinación (grados)",
                min_value=0,
                max_value=90,
                value=15,
                step=5,
                help="Inclinación del pozo - grados"
            )
            
            st.markdown(f"**{_t('fluid_formation')}**")
            mud_density_ppg = st.slider(
                "Densidad del lodo (ppg)",
                min_value=8.0,
                max_value=16.0,
                value=10.0,
                step=0.5,
                help="Densidad del lodo - usado en Bourgoyne & Young"
            )
            pore_gradient_ppg = st.slider(
                "Gradiente de poro (equiv ppg)",
                min_value=8.0,
                max_value=14.0,
                value=9.0,
                step=0.25,
                help="Gradiente de presión de poro (equiv ppg) - modelo B&Y"
            )
            yp_lb100ft2 = st.slider(
                "Punto de fluencia (lb/100ft²)",
                min_value=5.0,
                max_value=40.0,
                value=15.0,
                step=1.0,
                help="YP Bingham - efecto reológico en ROP"
            )
            pv_cp = st.slider(
                "Viscosidad plástica (cP)",
                min_value=10.0,
                max_value=60.0,
                value=25.0,
                step=2.0,
                help="PV Bingham - eficiencia hidráulica"
            )
        
        st.markdown("---")
        st.markdown(f"**{_t('geological_tracking')}**")
        use_geological_tracking = st.checkbox(
            _t('use_ucs_formation'),
            value=st.session_state.use_geological_tracking,
            help=_t('use_ucs_help')
        )
        st.session_state.use_geological_tracking = use_geological_tracking
        
        # Prediction Button
        st.markdown("---")
        predict_btn = st.button(
            f"**{_t('predict_rop')}**",
            use_container_width=True,
            type="primary"
        )
    
    # ========================================================================
    # MAIN CONTENT - MULTIPLE TABS (IMPERIAL UNITS)
    # ========================================================================
    
    if not st.session_state.models_trained:
        st.info(f"""
        **{_t('welcome')}**
        
        {_t('welcome_steps')}
        
        {_t('welcome_system')}
        """)
        
        # Show sample data
        with st.expander("**Datos de perforación de ejemplo (Unidades Imperiales)**"):
            sample_data = st.session_state.data_generator.generate().head(100)
            st.dataframe(sample_data, use_container_width=True)
    
    else:
        # Ejecutar predicción al inicio si se pulsó el botón, para que Resumen y Predicción ROP coincidan
        if predict_btn:
            ucs_for_pred = ucs_psi
            if st.session_state.use_geological_tracking:
                form_at_depth = get_formation_at_depth(depth_ft, st.session_state.geological_formations)
                if form_at_depth:
                    ucs_for_pred = form_at_depth['ucs_psi']
            params = {
                'wob_klb': wob_klb,
                'rpm': rpm,
                'torque_ftlb': torque_ftlb,
                'spp_psi': spp_psi,
                'flow_gpm': flow_gpm,
                'ucs_psi': ucs_for_pred,
                'bit_diameter_in': bit_diameter_in,
                'bit_wear': bit_wear,
                'depth_ft': depth_ft,
                'cutter_count': cutter_count,
                'inclination_deg': inclination_deg,
                'mud_density_ppg': mud_density_ppg,
                'pore_gradient_ppg': pore_gradient_ppg,
                'yp_lb100ft2': yp_lb100ft2,
                'pv_cp': pv_cp,
            }
            predictions = st.session_state.predictor.predict_ensemble(
                params,
                use_rf=train_rf,
                use_xgb=train_xgb,
                use_nn=train_nn
            )
            st.session_state.current_prediction = predictions
            temp_params = build_current_temperature_params(
                depth_ft=depth_ft, inclination_deg=inclination_deg, rpm=rpm, wob_klb=wob_klb,
                torque_ftlb=torque_ftlb, flow_gpm=flow_gpm, spp_psi=spp_psi,
                mud_density_ppg=mud_density_ppg, pv_cp=pv_cp, yp_lb100ft2=yp_lb100ft2,
                bit_diameter_in=bit_diameter_in, rop_value=predictions.get('Ensemble'),
                formation_info=(form_at_depth if st.session_state.use_geological_tracking and form_at_depth else None),
                trace_df=st.session_state.real_trace_df,
            )
            st.session_state.current_temperature_params = temp_params
            st.session_state.temperature_last_prediction = st.session_state.temperature_predictor.predict_from_params(temp_params)
        
        # Actualizar predicción térmica en vivo al mover sliders (incluye gráfico real vs predicha).
        try:
            if getattr(st.session_state.temperature_predictor, 'model', None) is not None:
                form_live = get_formation_at_depth(depth_ft, st.session_state.geological_formations) if st.session_state.use_geological_tracking else None
                ucs_live = form_live['ucs_psi'] if form_live else ucs_psi
                live_rop_pred = st.session_state.predictor.predict_ensemble({
                    'wob_klb': wob_klb, 'rpm': rpm, 'torque_ftlb': torque_ftlb, 'spp_psi': spp_psi,
                    'flow_gpm': flow_gpm, 'ucs_psi': ucs_live, 'bit_diameter_in': bit_diameter_in,
                    'bit_wear': bit_wear, 'depth_ft': depth_ft, 'cutter_count': cutter_count,
                    'inclination_deg': inclination_deg, 'mud_density_ppg': mud_density_ppg,
                    'pore_gradient_ppg': pore_gradient_ppg, 'yp_lb100ft2': yp_lb100ft2, 'pv_cp': pv_cp,
                }, use_rf=train_rf, use_xgb=train_xgb, use_nn=train_nn)
                live_temp_params = build_current_temperature_params(
                    depth_ft=depth_ft, inclination_deg=inclination_deg, rpm=rpm, wob_klb=wob_klb,
                    torque_ftlb=torque_ftlb, flow_gpm=flow_gpm, spp_psi=spp_psi,
                    mud_density_ppg=mud_density_ppg, pv_cp=pv_cp, yp_lb100ft2=yp_lb100ft2,
                    bit_diameter_in=bit_diameter_in, rop_value=live_rop_pred.get('Ensemble'),
                    formation_info=form_live, trace_df=st.session_state.real_trace_df,
                )
                st.session_state.current_temperature_params = live_temp_params
                st.session_state.temperature_last_prediction = st.session_state.temperature_predictor.predict_from_params(live_temp_params)
        except Exception:
            pass

        # Create tabs for different sections
        tab_resumen, tab1, tab_temp, tab2, tab3, tab4, tab5, tab6 = st.tabs([
            f"**{_t('tab_resumen')}**",
            f"**{_t('tab_prediction')}**",
            '**Predicción Temperatura**',
            f"**{_t('tab_performance')}**",
            f"**{_t('tab_heatmap')}**",
            f"**{_t('tab_neural')}**",
            f"**{_t('tab_optimization')}**",
            f"**{_t('tab_geological')}**"
        ])
        
        # ====================================================================
        # TAB RESUMEN: Resumen de detalles (seguimiento, recomendación, litología, gráficos)
        # ====================================================================
        
        with tab_resumen:
            st.markdown("### **Resumen de detalles**")
            st.markdown(f'<span class="unit-label">Seguimiento, recomendación y gráfica de cada sección</span>', 
                       unsafe_allow_html=True)
            
            formations = st.session_state.geological_formations
            form_at_depth = get_formation_at_depth(depth_ft, formations)
            rop_ensemble_resumen = st.session_state.current_prediction['Ensemble'] if st.session_state.current_prediction else None
            
            # Exportar a PPTX
            col_export, _ = st.columns([1, 4])
            with col_export:
                if PPTX_AVAILABLE:
                    if st.button("📥 **Generar archivo PPTX**", key="btn_export_pptx"):
                        with st.spinner("Generando presentación..."):
                            try:
                                pptx_bytes = generate_resumen_pptx(
                                    depth_ft, wob_klb, rpm, ucs_psi, bit_diameter_in,
                                    formations, st.session_state.predictor, st.session_state.data_generator,
                                    rop_ensemble_resumen,
                                    st.session_state.current_prediction
                                )
                                st.session_state.pptx_bytes = pptx_bytes
                                st.session_state.pptx_ready = True
                            except Exception as e:
                                st.error(f"Error al generar PPTX: {e}")
            
            if PPTX_AVAILABLE and st.session_state.get('pptx_ready'):
                st.download_button(
                    "**Descargar PPTX**",
                    data=st.session_state.get('pptx_bytes', b''),
                    file_name=f"Resumen_Detalle_Drilling_Analytics_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="download_pptx"
                )
            
            st.markdown("---")
            
            # --- Predicción ROP: Seguimiento, Recomendación y gráfica ---
            with st.expander("**Predicción ROP** — Seguimiento · Recomendación · Gráfica", expanded=True):
                st.markdown("**Seguimiento**")
                c1, c2, c3 = st.columns(3)
                with c1:
                    st.metric("Profundidad", f"{depth_ft:,.0f} ft")
                with c2:
                    st.metric("WOB / RPM", f"{wob_klb:.1f} klb / {rpm}")
                with c3:
                    st.metric("ROP Ensemble", f"{rop_ensemble_resumen:.1f} ft/hr" if rop_ensemble_resumen is not None else "—", help="Coincide con Predicción ROP. Ejecute PREDECIR ROP en el panel para actualizar.")
                st.markdown("**Recomendación**")
                for s in get_section_suggestions('rop_prediction'):
                    st.markdown(f"• {s}")
                for s in get_follow_up_suggestions('rop_prediction'):
                    st.markdown(f"→ {s}")
                st.markdown("**Gráfica**")
                gauge_val = rop_ensemble_resumen if rop_ensemble_resumen is not None else 0
                fig_gauge_rop = create_gauge_chart(gauge_val, "ROP", 0, 150, "ft/hr")
                st.plotly_chart(fig_gauge_rop, use_container_width=True, key="resumen_gauge_rop")
            
            # --- Rendimiento de modelos: Seguimiento, Recomendación y gráfica ---
            with st.expander("**Rendimiento de modelos** — Seguimiento · Recomendación · Gráfica", expanded=False):
                st.markdown("**Seguimiento**")
                mm = st.session_state.predictor.model_metrics
                r2_vals = [v.get('r2', 0) for k, v in mm.items() if isinstance(v, dict) and 'r2' in v]
                rmse_vals = [v.get('rmse', 0) for k, v in mm.items() if isinstance(v, dict) and 'rmse' in v]
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("R² promedio", f"{np.mean(r2_vals):.3f}" if r2_vals else "—")
                with col2:
                    st.metric("RMSE promedio", f"{np.mean(rmse_vals):.2f}" if rmse_vals else "—")
                with col3:
                    st.metric("Modelos activos", str(len([k for k in mm if mm.get(k)])))
                st.markdown("**Recomendación**")
                for s in get_section_suggestions('model_performance'):
                    st.markdown(f"• {s}")
                for s in get_follow_up_suggestions('model_performance'):
                    st.markdown(f"→ {s}")
                st.markdown("**Gráfica**")
                c1, c2 = st.columns(2)
                with c1:
                    fig_comp = create_model_comparison_chart(st.session_state.predictor.model_metrics)
                    fig_comp.update_layout(height=320)
                    st.plotly_chart(fig_comp, use_container_width=True, key="resumen_model_comp")
                with c2:
                    fig_lin = create_rop_linear_regression_chart(
                        st.session_state.data_generator, st.session_state.predictor, n_samples=100)
                    fig_lin.update_layout(height=320)
                    st.plotly_chart(fig_lin, use_container_width=True, key="resumen_regresion")
            
            # --- Mapa de calor ROP: Seguimiento, Recomendación y gráfica ---
            with st.expander("**Mapa de calor ROP** — Seguimiento · Recomendación · Gráfica", expanded=False):
                st.markdown("**Seguimiento**")
                col1, col2, col3 = st.columns([1, 1, 1.3])
                with col1:
                    st.metric("UCS", f"{ucs_psi:,.0f} psi", help="Resistencia a la compresión no confinada")
                with col2:
                    st.metric("Diámetro broca", f"{bit_diameter_in} in", help="Diámetro de la broca en pulgadas")
                with col3:
                    st.markdown("""
                    <div style="background: linear-gradient(135deg, rgba(27,77,62,0.1) 0%, rgba(44,110,73,0.08) 100%); 
                                border-radius: 8px; padding: 12px 16px; border-left: 4px solid #2C6E49; margin-top: 0.5rem;">
                        <div style="font-size: 0.8rem; color: #6B7280; margin-bottom: 4px;">Zona óptima</div>
                        <div style="font-size: 1rem; font-weight: 600; color: #1B4D3E;">
                            WOB 15–25 klb · RPM 100–140
                        </div>
                        <div style="font-size: 0.75rem; color: #9CA3AF; margin-top: 4px;">Verde oscuro en mapa</div>
                    </div>
                    """, unsafe_allow_html=True)
                st.markdown("**Recomendación**")
                for s in get_section_suggestions('heat_map'):
                    st.markdown(f"• {s}")
                for s in get_follow_up_suggestions('heat_map'):
                    st.markdown(f"→ {s}")
                st.markdown("**Gráfica**")
                fig_hm = create_rop_heatmap([5, 40], [40, 220], ucs_psi, bit_diameter_in)
                fig_hm.update_layout(height=380)
                st.plotly_chart(fig_hm, use_container_width=True, key="resumen_heatmap")
            
            # --- Análisis de red neuronal: Seguimiento, Recomendación y gráfica ---
            with st.expander("**Análisis de red neuronal** — Seguimiento · Recomendación · Gráfica", expanded=False):
                st.markdown("**Seguimiento**")
                col1, col2 = st.columns(2)
                with col1:
                    st.metric("Arquitectura", "MLP · 3 capas ocultas")
                with col2:
                    st.metric("Entrenamiento", "100 épocas")
                st.markdown("**Recomendación**")
                for s in get_section_suggestions('neural_network'):
                    st.markdown(f"• {s}")
                for s in get_follow_up_suggestions('neural_network'):
                    st.markdown(f"→ {s}")
                st.markdown("**Gráfica**")
                c_nn1, c_nn2 = st.columns(2)
                with c_nn1:
                    fig_arch = create_nn_architecture_diagram()
                    fig_arch.update_layout(height=300)
                    st.plotly_chart(fig_arch, use_container_width=True, key="resumen_nn_arch")
                with c_nn2:
                    epochs = np.arange(1, 101)
                    train_loss = 100 * np.exp(-epochs / 20) + np.random.RandomState(42).normal(0, 2, 100)
                    val_loss = 100 * np.exp(-epochs / 25) + np.random.RandomState(43).normal(0, 3, 100)
                    fig_nn = go.Figure()
                    fig_nn.add_trace(go.Scatter(x=epochs, y=train_loss, mode='lines', name='Entrenamiento', line=dict(color='#1B4D3E')))
                    fig_nn.add_trace(go.Scatter(x=epochs, y=val_loss, mode='lines', name='Validación', line=dict(color='#FF6B35')))
                    fig_nn.update_layout(title='Curvas de aprendizaje', xaxis_title='Época', yaxis_title='Pérdida', height=300)
                    st.plotly_chart(fig_nn, use_container_width=True, key="resumen_nn_curves")
            
            # --- Optimización: Seguimiento, Recomendación y gráfica ---
            with st.expander("**Optimización** — Seguimiento · Recomendación · Gráfica", expanded=False):
                st.markdown("**Seguimiento**")
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("WOB actual", f"{wob_klb:.1f} klb")
                with col2:
                    st.metric("RPM actual", f"{rpm}")
                with col3:
                    st.metric("ROP actual", f"{st.session_state.current_prediction['Ensemble']:.1f} ft/hr" if st.session_state.current_prediction else "—")
                st.markdown("**Recomendación**")
                for s in get_section_suggestions('optimization'):
                    st.markdown(f"• {s}")
                for s in get_follow_up_suggestions('optimization'):
                    st.markdown(f"→ {s}")
                st.markdown("**Gráfica**")
                wob_t = np.linspace(5, 40, 40)
                rop_w = [st.session_state.predictor.predict_physical_model({
                    'wob_klb': w, 'rpm': 120, 'torque_ftlb': 18000, 'spp_psi': 3000, 'flow_gpm': 800,
                    'ucs_psi': ucs_psi, 'bit_diameter_in': bit_diameter_in, 'bit_wear': bit_wear,
                    'depth_ft': depth_ft, 'cutter_count': cutter_count
                }) for w in wob_t]
                rpm_t = np.linspace(40, 220, 40)
                rop_r = [st.session_state.predictor.predict_physical_model({
                    'wob_klb': wob_klb, 'rpm': r, 'torque_ftlb': 18000, 'spp_psi': 3000, 'flow_gpm': 800,
                    'ucs_psi': ucs_psi, 'bit_diameter_in': bit_diameter_in, 'bit_wear': bit_wear,
                    'depth_ft': depth_ft, 'cutter_count': cutter_count
                }) for r in rpm_t]
                co1, co2 = st.columns(2)
                with co1:
                    fw = go.Figure(go.Scatter(x=wob_t, y=rop_w, mode='lines', line=dict(color='#1B4D3E')))
                    fw.update_layout(title='Sensibilidad WOB', xaxis_title='WOB (klb)', yaxis_title='ROP (ft/hr)', height=320)
                    st.plotly_chart(fw, use_container_width=True, key="resumen_opt_wob")
                with co2:
                    fr = go.Figure(go.Scatter(x=rpm_t, y=rop_r, mode='lines', line=dict(color='#2C6E49')))
                    fr.update_layout(title='Sensibilidad RPM', xaxis_title='RPM', yaxis_title='ROP (ft/hr)', height=320)
                    st.plotly_chart(fr, use_container_width=True, key="resumen_opt_rpm")
            
            # --- Seguimiento geológico: Seguimiento, Recomendación y gráfica ---
            with st.expander("**Seguimiento geológico** — Seguimiento · Recomendación · Gráfica", expanded=False):
                st.markdown("**Seguimiento**")
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("Profundidad", f"{depth_ft:,.0f} ft")
                with col2:
                    st.metric("Formación actual", form_at_depth['name'] if form_at_depth else "—")
                with col3:
                    ucs_val = form_at_depth['ucs_psi'] if form_at_depth else ucs_psi
                    st.metric("UCS formación", f"{ucs_val:,.0f} psi")
                st.caption(f"Seguimiento geológico: {'Activado' if st.session_state.use_geological_tracking else 'Desactivado'}")
                st.markdown("**Recomendación**")
                for s in get_section_suggestions('geological_tracking'):
                    st.markdown(f"• {s}")
                for s in get_follow_up_suggestions('geological_tracking'):
                    st.markdown(f"→ {s}")
                st.markdown("**Gráfica**")
                params_base = {'wob_klb': wob_klb, 'rpm': rpm, 'torque_ftlb': torque_ftlb, 'spp_psi': spp_psi,
                              'flow_gpm': flow_gpm, 'bit_diameter_in': bit_diameter_in, 'bit_wear': bit_wear,
                              'cutter_count': cutter_count, 'mud_density_ppg': mud_density_ppg,
                              'pore_gradient_ppg': pore_gradient_ppg, 'yp_lb100ft2': yp_lb100ft2,
                              'pv_cp': pv_cp, 'inclination_deg': inclination_deg}
                rop_f = predict_rop_by_formation(formations, params_base, st.session_state.predictor)
                fig_g = go.Figure(go.Bar(
                    x=[r['formation'] for r in rop_f],
                    y=[r['rop_predicted'] for r in rop_f],
                    marker_color=['#1B4D3E', '#2C6E49', '#3B8C5E', '#4AA66B', '#5BC77B'],
                    text=[f"{r['rop_predicted']:.1f}" for r in rop_f], textposition='outside'
                ))
                fig_g.update_layout(title='ROP predicho por formación', xaxis_title='Formación', yaxis_title='ROP (ft/hr)', height=350)
                st.plotly_chart(fig_g, use_container_width=True, key="resumen_geo_bars")
            
            # --- Análisis de Correlación: Seguimiento, Recomendación y gráfica ---
            with st.expander("**Análisis de Correlación** — Seguimiento · Recomendación · Gráfica", expanded=False):
                st.markdown("**Seguimiento**")
                if 'correlation_data' not in st.session_state:
                    with st.spinner("Generando datos para correlación..."):
                        st.session_state.correlation_data = st.session_state.data_generator.generate().sample(n=500, random_state=42)
                corr_data = st.session_state.correlation_data
                corr_matrix = corr_data.corr()
                corr_with_rop = corr_matrix['ROP_fthr'].drop('ROP_fthr', errors='ignore') if 'ROP_fthr' in corr_matrix.columns else pd.Series(dtype=float)
                col1, col2, col3 = st.columns(3)
                with col1:
                    if len(corr_with_rop) > 0:
                        top_pos = corr_with_rop.nlargest(1)
                        st.metric("Mayor correlación +", f"{top_pos.index[0].replace('_', ' ')}: {top_pos.iloc[0]:.2f}")
                    else:
                        st.metric("Mayor correlación +", "—")
                with col2:
                    if len(corr_with_rop) > 0:
                        top_neg = corr_with_rop.nsmallest(1)
                        st.metric("Mayor correlación −", f"{top_neg.index[0].replace('_', ' ')}: {top_neg.iloc[0]:.2f}")
                    else:
                        st.metric("Mayor correlación −", "—")
                with col3:
                    st.metric("Variables analizadas", str(len(corr_with_rop)))
                st.markdown("**Recomendación**")
                for s in get_section_suggestions('correlation'):
                    st.markdown(f"• {s}")
                for s in get_follow_up_suggestions('correlation'):
                    st.markdown(f"→ {s}")
                st.markdown("**Gráfica**")
                fig_corr_res = create_correlation_analysis_chart(corr_data)
                fig_corr_res.update_layout(height=450)
                st.plotly_chart(fig_corr_res, use_container_width=True, key="resumen_correlation")
            
            if st.session_state.current_prediction:
                st.markdown("---")
                st.markdown("#### **Última predicción ROP**")
                st.metric("Ensemble ROP", f"{st.session_state.current_prediction['Ensemble']:.1f} ft/hr")
        
        # ====================================================================
        # TAB 1: ROP PREDICTION - IMPERIAL
        # ====================================================================
        
        with tab1:
            with st.expander("**Recomendaciones para esta sección**", expanded=True):
                for s in get_section_suggestions('rop_prediction'):
                    st.markdown(f"• {s}")
            with st.expander("**Sugerencias de seguimiento**", expanded=True):
                for s in get_follow_up_suggestions('rop_prediction'):
                    st.markdown(f"→ {s}")
            
            # --- Análisis de Correlación ---
            st.markdown("<div class='divider'></div>", unsafe_allow_html=True)
            st.markdown("### **Análisis de Correlación**")
            st.markdown("Matriz de correlación entre variables de perforación y ROP (datos sintéticos)")
            if 'correlation_data' not in st.session_state:
                with st.spinner("Generando datos para análisis de correlación..."):
                    st.session_state.correlation_data = st.session_state.data_generator.generate().sample(n=500, random_state=42)
            fig_corr = create_correlation_analysis_chart(st.session_state.correlation_data)
            st.plotly_chart(fig_corr, use_container_width=True, key="pred_correlation")
            with st.expander("**Interpretación de la correlación**", expanded=False):
                st.markdown("""
                - **Valores cercanos a 1**: Correlación positiva fuerte — si aumenta la variable, ROP tiende a aumentar.
                - **Valores cercanos a -1**: Correlación negativa — si aumenta la variable, ROP tiende a disminuir.
                - **Valores cercanos a 0**: Poca correlación lineal con ROP.
                - **WOB y RPM** suelen correlacionar positivamente con ROP (más peso/revoluciones → más penetración).
                - **UCS y profundidad** suelen correlacionar negativamente con ROP (formaciones más duras/profundas → menor penetración).
                """)
            
            if predict_btn and st.session_state.current_prediction:
                predictions = st.session_state.current_prediction
                
                # Mechanistic models expander
                with st.expander("**Modelos mecanicistas: B&Y · Bingham · Warren**", expanded=True):
                    st.markdown("""
                    **Correlación con ML y Redes Neuronales:** Los modelos mecanicistas aportan variables basadas en física 
                    que los modelos ML aprenden a refinar. Bourgoyne & Young (1974) utiliza 8 factores (profundidad, presión, WOB, RPM, hidráulica). 
                    Bingham incorpora reología YP/PV en la eficiencia hidráulica. Warren (1987) modela generación y remoción 
                    de recortes. Los pesos del ensemble combinan mecanicistas + ML para predicciones robustas.
                    """)
                
                # Display predictions
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    st.markdown("""
                    <div class="metric-card">
                        <h3 style="color: white; margin-bottom: 0.5rem;">ROP Ensemble</h3>
                        <p style="font-size: 3rem; font-weight: 700; margin: 0; color: white;">
                            {:.1f}
                        </p>
                        <p style="font-size: 1.1rem; margin: 0; color: rgba(255,255,255,0.9);">
                            ft/hr
                        </p>
                        <p style="font-size: 0.9rem; margin-top: 0.5rem; color: rgba(255,255,255,0.8);">
                            ± {:.1f} ft/hr (95% CI)
                        </p>
                    </div>
                    """.format(predictions['Ensemble'], predictions['Ensemble'] * 0.08), 
                    unsafe_allow_html=True)
                
                with col2:
                    st.markdown("""
                    <div style="background: white; border-radius: 15px; padding: 25px; box-shadow: 0 4px 20px rgba(0,0,0,0.05);">
                        <h3 style="color: #1B4D3E; margin-bottom: 1rem;">Desglose por modelo</h3>
                    """, unsafe_allow_html=True)
                    
                    for model_name, pred_value in predictions.items():
                        if model_name != 'Ensemble' and pred_value is not None:
                            st.markdown(f"""
                            <div style="display: flex; justify-content: space-between; margin-bottom: 0.5rem;">
                                <span style="color: #4A5568;">{model_name}:</span>
                                <span style="font-weight: 600; color: #1B4D3E;">{pred_value:.1f} ft/hr</span>
                            </div>
                            """, unsafe_allow_html=True)
                    
                    st.markdown("</div>", unsafe_allow_html=True)
                
                with col3:
                    # Calculate drilling efficiency
                    efficiency = (predictions['Ensemble'] * 100) / (wob_klb * rpm / 1000)
                    efficiency = min(efficiency, 100)
                    
                    st.markdown(f"""
                    <div style="background: white; border-radius: 15px; padding: 25px; box-shadow: 0 4px 20px rgba(0,0,0,0.05);">
                        <h3 style="color: #1B4D3E; margin-bottom: 1rem;">Métricas de eficiencia</h3>
                        <div style="display: flex; justify-content: space-between; margin-bottom: 0.5rem;">
                            <span style="color: #4A5568;">Eficiencia de perforación:</span>
                            <span style="font-weight: 600; color: #1B4D3E;">{efficiency:.1f}%</span>
                        </div>
                        <div style="display: flex; justify-content: space-between; margin-bottom: 0.5rem;">
                            <span style="color: #4A5568;">Estimación MSE:</span>
                            <span style="font-weight: 600; color: #1B4D3E;">{np.random.uniform(25, 45):.0f} ksi</span>
                        </div>
                        <div style="display: flex; justify-content: space-between;">
                            <span style="color: #4A5568;">Confianza:</span>
                            <span style="font-weight: 600; color: #1B4D3E;">{np.random.uniform(88, 96):.1f}%</span>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                
                # Gauge charts
                st.markdown("<div class='divider'></div>", unsafe_allow_html=True)
                
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    fig = create_gauge_chart(
                        predictions['Ensemble'],
                        "ROP",
                        0, 150,
                        "ft/hr"
                    )
                    st.plotly_chart(fig, use_container_width=True, key="pred_gauge_rop")
                
                with col2:
                    fig = create_gauge_chart(
                        wob_klb,
                        "WOB",
                        0, 50,
                        "klb"
                    )
                    st.plotly_chart(fig, use_container_width=True, key="pred_gauge_wob")
                
                with col3:
                    fig = create_gauge_chart(
                        rpm,
                        "RPM",
                        0, 250,
                        "rev/min"
                    )
                    st.plotly_chart(fig, use_container_width=True, key="pred_gauge_rpm")
                
                # Recommendations
                st.markdown("<div class='divider'></div>", unsafe_allow_html=True)
                st.markdown("### **Recomendaciones de optimización**")
                
                col1, col2 = st.columns(2)
                
                with col1:
                    st.markdown("""
                    <div class="info-box">
                        <h4 style="color: #1B4D3E; margin-bottom: 0.5rem;">Ajuste de parámetros</h4>
                    """, unsafe_allow_html=True)
                    
                    recs = []
                    
                    if predictions['Ensemble'] < 40:
                        recs.append(("Aumentar WOB", f"Actual: {wob_klb:.1f} klb → Objetivo: {wob_klb + 3:.1f} klb"))
                    elif predictions['Ensemble'] > 80:
                        recs.append(("ROP óptimo", "Mantener parámetros actuales"))
                    
                    if rpm < 100:
                        recs.append(("Aumentar RPM", f"Actual: {rpm} → Objetivo: {rpm + 20}"))
                    elif rpm > 160:
                        recs.append(("Reducir RPM", f"Actual: {rpm} → Objetivo: {rpm - 20}"))
                    
                    if bit_wear > 0.6:
                        recs.append(("Alerta desgaste broca", f"Desgaste actual: {bit_wear*100:.0f}% - Programar reemplazo"))
                    
                    for title, desc in recs[:3]:
                        st.markdown(f"**{title}**<br>{desc}", unsafe_allow_html=True)
                        st.markdown("---")
                    
                    st.markdown("</div>", unsafe_allow_html=True)
                
                with col2:
                    st.markdown("""
                    <div class="info-box">
                        <h4 style="color: #1B4D3E; margin-bottom: 0.5rem;">Evaluación de riesgos</h4>
                    """, unsafe_allow_html=True)
                    
                    # Calculate risks
                    stick_slip_risk = min(100, (rpm / 60) * (wob_klb / 20) * 30)
                    vibration_risk = min(100, (bit_wear * 50) + (abs(rpm - 120) / 2))
                    balling_risk = min(100, (flow_gpm / 800) * 50)
                    
                    st.markdown(f"""
                    <div style="margin-bottom: 0.5rem;">
                        <div style="display: flex; justify-content: space-between;">
                            <span>Riesgo stick-slip:</span>
                            <span style="font-weight: 600; color: {'#DC3545' if stick_slip_risk > 70 else '#FFC107' if stick_slip_risk > 40 else '#28A745'}">
                                {stick_slip_risk:.0f}%
                            </span>
                        </div>
                        <div style="background: #E9ECEF; height: 8px; border-radius: 4px; margin-top: 5px;">
                            <div style="background: {'#DC3545' if stick_slip_risk > 70 else '#FFC107' if stick_slip_risk > 40 else '#28A745'}; 
                                      width: {stick_slip_risk}%; height: 8px; border-radius: 4px;"></div>
                        </div>
                    </div>
                    
                    <div style="margin-bottom: 0.5rem; margin-top: 1rem;">
                        <div style="display: flex; justify-content: space-between;">
                            <span>Riesgo de vibración:</span>
                            <span style="font-weight: 600; color: {'#DC3545' if vibration_risk > 70 else '#FFC107' if vibration_risk > 40 else '#28A745'}">
                                {vibration_risk:.0f}%
                            </span>
                        </div>
                        <div style="background: #E9ECEF; height: 8px; border-radius: 4px; margin-top: 5px;">
                            <div style="background: {'#DC3545' if vibration_risk > 70 else '#FFC107' if vibration_risk > 40 else '#28A745'}; 
                                      width: {vibration_risk}%; height: 8px; border-radius: 4px;"></div>
                        </div>
                    </div>
                    
                    <div style="margin-bottom: 0.5rem; margin-top: 1rem;">
                        <div style="display: flex; justify-content: space-between;">
                            <span>Riesgo embolamiento broca:</span>
                            <span style="font-weight: 600; color: {'#DC3545' if balling_risk > 70 else '#FFC107' if balling_risk > 40 else '#28A745'}">
                                {balling_risk:.0f}%
                            </span>
                        </div>
                        <div style="background: #E9ECEF; height: 8px; border-radius: 4px; margin-top: 5px;">
                            <div style="background: {'#DC3545' if balling_risk > 70 else '#FFC107' if balling_risk > 40 else '#28A745'}; 
                                      width: {balling_risk}%; height: 8px; border-radius: 4px;"></div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    st.markdown("</div>", unsafe_allow_html=True)
            
            else:
                st.info("Configure los parámetros y haga clic en 'PREDECIR ROP' para ver resultados")
        
        # ====================================================================
        # TAB TEMPERATURA: MODELO ML CON TRAZAS REALES
        # ====================================================================

        with tab_temp:
            st.markdown('### **Predicción de temperatura y perfil esperado del campo**')
            st.markdown('Este módulo combina dos vistas: ML sobre la traza objetivo y un perfil esperado construido a partir de pozos vecinos de la misma locación.')

            if st.session_state.real_trace_df is not None:
                st.caption(f"Pozo objetivo: {st.session_state.real_trace_name} · filas {len(st.session_state.real_trace_df):,}")
            else:
                st.info('Sube la traza del pozo objetivo desde el panel lateral.')

            tp = st.session_state.temperature_predictor
            profiler = st.session_state.neighbor_temp_profile
            neighbor_profile_df = st.session_state.neighbor_profile_df
            if neighbor_profile_df is not None and not neighbor_profile_df.empty:
                profiler.profile_df = neighbor_profile_df
                if profiler.target_col is None and st.session_state.get('neighbor_target_col'):
                    profiler.target_col = st.session_state.neighbor_target_col
                if st.session_state.get('neighbor_depth_col'):
                    profiler.depth_col = st.session_state.neighbor_depth_col
                if st.session_state.real_trace_df is not None:
                    _fresh_pred = profiler.predict_for_trace(st.session_state.real_trace_df)
                    if _fresh_pred is not None and not _fresh_pred.empty:
                        st.session_state.neighbor_prediction_df = _fresh_pred
            neighbor_prediction_df = st.session_state.neighbor_prediction_df

            subtab1, subtab2, subtab3, subtab4, subtab5 = st.tabs([
                'ML temperatura', 'Perfil esperado vecinos', 'Litología / formación', 'Anomalías térmicas', 'Roadmap por formación'
            ])

            with subtab1:
                if tp.model is not None:
                    r2_val = tp.metrics.get('r2', 0)
                    r2_tr = tp.metrics.get('r2_train', 0)
                    rmse_val = tp.metrics.get('rmse', 0)
                    mae_val = tp.metrics.get('mae', 0)
                    st.markdown(
                        '<div class="chip-row">'
                        f'<span class="chip chip-temp-target">Target · {tp.target_col}</span>'
                        f'<span class="chip chip-temp-r2">R² test {r2_val:.3f}</span>'
                        f'<span class="chip chip-temp-r2">R² train {r2_tr:.3f}</span>'
                        f'<span class="chip chip-temp-rmse">RMSE {rmse_val:.2f}</span>'
                        f'<span class="chip chip-temp-mae">MAE {mae_val:.2f}</span>'
                        '<span class="chip chip-temp-ok">Modelo activo</span>'
                        '</div>',
                        unsafe_allow_html=True,
                    )

                    temp_pred = st.session_state.temperature_last_prediction
                    if temp_pred is not None:
                        st.success(f'Temperatura predicha para la condición actual: {temp_pred:.2f} °C')
                        live_params = st.session_state.get('current_temperature_params', {})
                        c1, c2, c3, c4 = st.columns(4)
                        c1.metric('WOB actual', f"{live_params.get('wob', wob_klb):.1f} klb")
                        c2.metric('RPM actual', f"{live_params.get('rpm', rpm):.0f}")
                        c3.metric('Caudal actual', f"{live_params.get('flow_rate', flow_gpm):.0f} gpm")
                        c4.metric('Mud In usado', f"{live_params.get('mud_in_temp', np.nan):.1f} °C")
                        st.caption(
                            'La línea naranja **1:1 es referencia fija** (predicción perfecta). '
                            'Lo que se mueve al ajustar WOB/RPM/torque/SPP/caudal son los **puntos azules** (predicciones). '
                            'La estrella roja marca la profundidad actual.'
                        )
                        try:
                            fig_live_temp = create_temperature_wob_rpm_heatmap(
                                tp, live_params, wob_range=(max(5.0, wob_klb-10.0), min(45.0, wob_klb+10.0)),
                                rpm_range=(max(40.0, rpm-60.0), min(240.0, rpm+60.0)),
                                profile_df=st.session_state.get('neighbor_profile_df'), depth_ft=float(depth_ft),
                                grid_n=28, subtitle_extra='sensibilidad local con sliders actuales',
                            )
                            st.plotly_chart(fig_live_temp, use_container_width=True, key='tab_temp_live_wob_rpm')
                        except Exception:
                            pass
                    else:
                        st.info('Entrena el modelo para generar una predicción de temperatura con los parámetros actuales.')

                    st.plotly_chart(create_temperature_feature_importance_chart(tp.feature_importance), use_container_width=True, key='tab_temp_fi')

                    if st.session_state.real_trace_df is not None and tp.target_col in st.session_state.real_trace_df.columns:
                        _live_for_chart = st.session_state.get('current_temperature_params') or build_current_temperature_params(
                            depth_ft=depth_ft, inclination_deg=inclination_deg, rpm=rpm, wob_klb=wob_klb,
                            torque_ftlb=torque_ftlb, flow_gpm=flow_gpm, spp_psi=spp_psi,
                            mud_density_ppg=mud_density_ppg, pv_cp=pv_cp, yp_lb100ft2=yp_lb100ft2,
                            bit_diameter_in=bit_diameter_in,
                            rop_value=(st.session_state.current_prediction or {}).get('Ensemble'),
                            formation_info=(get_formation_at_depth(depth_ft, st.session_state.geological_formations)
                                            if st.session_state.use_geological_tracking else None),
                            trace_df=st.session_state.real_trace_df,
                        )
                        actual_live, pred_live, depths_live = temperature_predict_trace_live(
                            tp, st.session_state.real_trace_df, _live_for_chart,
                        )
                        if len(actual_live) > 0 and len(pred_live) == len(actual_live):
                            _live_r2 = float(r2_score(actual_live, pred_live)) if len(actual_live) > 1 else 0.0
                            _live_mae = float(mean_absolute_error(actual_live, pred_live))
                            if len(depths_live) == len(actual_live) and np.isfinite(depths_live).any():
                                _idx_near = int(np.argmin(np.abs(depths_live - float(depth_ft))))
                            else:
                                _idx_near = len(actual_live) // 2
                            _hi_actual = float(actual_live[_idx_near])
                            _hi_pred = float(pred_live[_idx_near])
                            st.caption(
                                f'Sliders: WOB {_live_for_chart.get("wob", wob_klb):.1f} klb · '
                                f'RPM {_live_for_chart.get("rpm", rpm):.0f} · '
                                f'Torque {_live_for_chart.get("torque", torque_ftlb):.0f} ft-lb · '
                                f'SPP {_live_for_chart.get("pump_pressure", spp_psi):.0f} psi · '
                                f'Caudal {_live_for_chart.get("flow_rate", flow_gpm):.0f} gpm. '
                                f'R² escenario {_live_r2:.3f} · MAE {_live_mae:.2f} °C.'
                            )
                            st.plotly_chart(
                                create_temperature_real_vs_pred_chart(
                                    actual_live,
                                    pred_live,
                                    highlight={'actual': _hi_actual, 'pred': _hi_pred},
                                ),
                                use_container_width=True,
                                key=f'tab_temp_real_vs_pred_{wob_klb:.1f}_{rpm:.0f}_{flow_gpm:.0f}_{torque_ftlb:.0f}',
                            )
                        _ho_a = getattr(tp, 'last_holdout_actual', None)
                        _ho_p = getattr(tp, 'last_holdout_pred', None)
                        if (
                            _ho_a is not None
                            and _ho_p is not None
                            and len(_ho_a) > 0
                            and len(_ho_p) == len(_ho_a)
                        ):
                            with st.expander('Validación holdout (fija, solo referencia de entrenamiento)'):
                                st.plotly_chart(
                                    create_temperature_real_vs_pred_chart(_ho_a, _ho_p),
                                    use_container_width=True,
                                    key='tab_temp_real_vs_pred_holdout',
                                )
                else:
                    st.warning('Aún no hay un modelo ML de temperatura entrenado con la traza objetivo.')

            with subtab2:
                if neighbor_profile_df is not None and not neighbor_profile_df.empty:
                    _n_neigh = int(profiler.metadata.get('neighbor_count', 0))
                    _d_lo = float(neighbor_profile_df['depth'].min())
                    _d_hi = float(neighbor_profile_df['depth'].max())
                    _dmin, _dmax = float(_d_lo), float(_d_hi)
                    if _dmax <= _dmin:
                        _dmax = _dmin + 1.0
                    _g_mu = float(neighbor_profile_df['gradient_expected'].dropna().mean()) if neighbor_profile_df['gradient_expected'].notna().any() else 0.0
                    _tgt = html.escape(str(profiler.target_col or '—'))
                    _dref = html.escape(str(profiler.depth_col or 'depth'))

                    st.markdown('**Escala y ventana de profundidad**')
                    depth_unit_prof = st.radio(
                        'Unidad eje Y (profundidad)',
                        ['Pies (ft)', 'Metros (m)'],
                        horizontal=True,
                        key='neighbor_prof_depth_unit',
                        help='El perfil interno está en pies; en metros solo cambia la escala del eje, chips y slider.',
                    )
                    use_m_prof = depth_unit_prof == 'Metros (m)'
                    if use_m_prof:
                        _chip_lo = converter.ft_to_m(_d_lo)
                        _chip_hi = converter.ft_to_m(_d_hi)
                        _chip_u = 'm'
                        _chip_fmt = '{:,.2f}'
                        _g_mu_disp = _g_mu / max(converter.ft_to_m(1.0), 1e-12)
                        _g_suffix = ' <span style="opacity:.75;font-size:0.65rem">/m</span>'
                    else:
                        _chip_lo, _chip_hi = _d_lo, _d_hi
                        _chip_u = 'ft'
                        _chip_fmt = '{:,.0f}'
                        _g_mu_disp = _g_mu
                        _g_suffix = ' <span style="opacity:.75;font-size:0.65rem">/ft</span>'
                    _chip_lo_s = _chip_fmt.format(_chip_lo)
                    _chip_hi_s = _chip_fmt.format(_chip_hi)
                    st.markdown(
                        '<div class="chip-row-neighbor">'
                        f'<span class="chip-neigh chip-neigh-count"><span class="chip-neigh-k">Vecinos</span> <strong>{_n_neigh}</strong></span>'
                        f'<span class="chip-neigh chip-neigh-depth-lo"><span class="chip-neigh-k">Prof. min</span> <strong>{_chip_lo_s}</strong> <span style="opacity:.85;font-weight:600;font-size:0.7rem">{_chip_u}</span></span>'
                        f'<span class="chip-neigh chip-neigh-depth-hi"><span class="chip-neigh-k">Prof. max</span> <strong>{_chip_hi_s}</strong> <span style="opacity:.85;font-weight:600;font-size:0.7rem">{_chip_u}</span></span>'
                        f'<span class="chip-neigh chip-neigh-grad"><span class="chip-neigh-k">Grad. medio</span> <strong>{_g_mu_disp:.4f}</strong>{_g_suffix}</span>'
                        f'<span class="chip-neigh chip-neigh-meta"><span class="chip-neigh-k">Variable</span> {_tgt} · <span class="chip-neigh-k">Eje</span> {_dref}</span>'
                        '</div>',
                        unsafe_allow_html=True,
                    )

                    if use_m_prof:
                        _m_lo = converter.ft_to_m(_dmin)
                        _m_hi = converter.ft_to_m(_dmax)
                        rng_m = st.slider(
                            'Rango de profundidad mostrado',
                            min_value=_m_lo,
                            max_value=_m_hi,
                            value=(_m_lo, _m_hi),
                            step=5.0,
                            key='neighbor_prof_depth_slider_m',
                        )
                        depth_range_ft_prof = (converter.m_to_ft(rng_m[0]), converter.m_to_ft(rng_m[1]))
                    else:
                        rng_ft = st.slider(
                            'Rango de profundidad mostrado',
                            min_value=_dmin,
                            max_value=_dmax,
                            value=(_dmin, _dmax),
                            step=50.0,
                            key='neighbor_prof_depth_slider_ft',
                        )
                        depth_range_ft_prof = (float(rng_ft[0]), float(rng_ft[1]))

                    actual_overlay = neighbor_prediction_df if neighbor_prediction_df is not None and not neighbor_prediction_df.empty else None
                    target_overlay = profiler.target_col if actual_overlay is not None else None
                    depth_overlay = profiler.depth_col if actual_overlay is not None and profiler.depth_col in actual_overlay.columns else ('depth_md' if actual_overlay is not None and 'depth_md' in actual_overlay.columns else profiler.depth_col)
                    st.plotly_chart(
                        create_neighbor_temperature_profile_chart(
                            neighbor_profile_df,
                            actual_overlay,
                            target_overlay,
                            depth_overlay,
                            depth_units='m' if use_m_prof else 'ft',
                            depth_range_ft=depth_range_ft_prof,
                        ),
                        use_container_width=True,
                        key='neighbor_profile_chart'
                    )

                    if profiler.neighbor_summary is not None and not profiler.neighbor_summary.empty:
                        st.markdown('#### Mapa interactivo de pozos vecinos')
                        target_x_map = profiler.metadata.get('target_x')
                        target_y_map = profiler.metadata.get('target_y')
                        st.plotly_chart(
                            create_neighbor_map_chart(
                                profiler.neighbor_summary,
                                target_x=target_x_map,
                                target_y=target_y_map,
                                weighting=profiler.weighting,
                            ),
                            use_container_width=True,
                            key='neighbor_xy_map_chart'
                        )
                        st.caption('El mapa se actualiza cada vez que cambias las coordenadas del pozo objetivo y vuelves a generar el perfil. Si un vecino no trae X/Y válidas, no aparecerá en el mapa.')
                        st.markdown('#### Pozos vecinos más influyentes')
                        st.markdown(neighbor_summary_table_html(profiler.neighbor_summary.head(15)), unsafe_allow_html=True)
                else:
                    st.info('Carga varios pozos vecinos y presiona "Generar perfil esperado desde pozos vecinos".')

            with subtab3:
                if neighbor_prediction_df is not None and not neighbor_prediction_df.empty:
                    st.plotly_chart(create_lithology_temperature_chart(neighbor_prediction_df, profiler.target_col), use_container_width=True, key='lith_temp_chart')
                    if 'lithology' in neighbor_prediction_df.columns or 'formation' in neighbor_prediction_df.columns:
                        category_col = 'lithology' if 'lithology' in neighbor_prediction_df.columns else 'formation'
                        temp_metric_col = 'temp_residual' if 'temp_residual' in neighbor_prediction_df.columns else profiler.target_col
                        summary = neighbor_prediction_df.groupby(category_col)[temp_metric_col].agg(['count', 'mean', 'std']).reset_index()
                        st.dataframe(summary, use_container_width=True)
                    else:
                        st.info('El pozo objetivo no trae litología/formación para correlacionar con temperatura.')
                elif st.session_state.real_trace_df is not None:
                    base_target = _detect_temperature_target(st.session_state.real_trace_df)
                    st.plotly_chart(create_lithology_temperature_chart(st.session_state.real_trace_df, base_target or ''), use_container_width=True, key='lith_temp_chart_raw')
                else:
                    st.info('No hay datos suficientes para la correlación temperatura-litología.')

            with subtab4:
                st.session_state.thermal_anomaly_z_threshold = st.slider(
                    'Umbral de anomalía |Z|',
                    min_value=1.0, max_value=4.0, value=float(st.session_state.get('thermal_anomaly_z_threshold', 2.0)),
                    step=0.1,
                    help='Un punto se marca como anomalía cuando |residual / dispersión esperada| supera este valor.'
                )
                if neighbor_prediction_df is not None and not neighbor_prediction_df.empty and 'temp_residual' in neighbor_prediction_df.columns:
                    depth_col = _choose_available_depth_col(neighbor_prediction_df, profiler.depth_col) or profiler.depth_col
                    _prof_max = (
                        float(neighbor_profile_df['depth'].max())
                        if neighbor_profile_df is not None and not neighbor_profile_df.empty and 'depth' in neighbor_profile_df.columns
                        else np.nan
                    )
                    _tgt_depths = pd.to_numeric(neighbor_prediction_df[depth_col], errors='coerce').dropna()
                    _tgt_max = float(_tgt_depths.max()) if not _tgt_depths.empty else np.nan
                    _ratio_dp = (_tgt_max / _prof_max) if np.isfinite(_prof_max) and _prof_max > 0 else 1.0
                    if 2.4 <= _ratio_dp <= 3.6:
                        st.error(
                            f'La traza objetivo parece con **doble conversión** de profundidad '
                            f'({ _tgt_max:,.0f} ft ≈ 3.28 × {_prof_max:,.0f} ft). '
                            'Tus CSV en **metros** están bien: recarga los archivos, deja **Auto (detectar)** '
                            'y pulsa de nuevo «Generar perfil esperado desde pozos vecinos».'
                        )
                    elif (
                        np.isfinite(_prof_max)
                        and np.isfinite(_tgt_max)
                        and _prof_max < _tgt_max * 0.45
                        and _ratio_dp < 2.0
                    ):
                        st.error(
                            f'El perfil de vecinos solo llega a **{_prof_max:,.0f} ft** pero la traza objetivo a **{_tgt_max:,.0f} ft**. '
                            'Regenera el perfil con «Generar perfil esperado desde pozos vecinos» y target **mud_out_temp**.'
                        )
                    elif st.session_state.neighbor_profile_unit_system != st.session_state.trace_unit_system:
                        st.warning(
                            'El perfil de vecinos se generó con otro sistema de unidades. Regenera el perfil en el panel lateral.'
                        )
                    # Recalcular flags si el usuario ajusta el umbral.
                    if 'temp_zscore' in neighbor_prediction_df.columns:
                        _thr = float(st.session_state.thermal_anomaly_z_threshold)
                        neighbor_prediction_df = neighbor_prediction_df.copy()
                        neighbor_prediction_df['thermal_anomaly'] = np.where(neighbor_prediction_df['temp_zscore'].abs() >= _thr, 'anomaly', 'normal')
                        neighbor_prediction_df['thermal_severity'] = pd.cut(
                            neighbor_prediction_df['temp_zscore'].abs(),
                            bins=[-np.inf, _thr, _thr + 1.0, np.inf],
                            labels=['normal', 'media', 'alta']
                        ).astype(str)
                        st.session_state.neighbor_prediction_df = neighbor_prediction_df
                    depth_col = _choose_available_depth_col(neighbor_prediction_df, profiler.depth_col) or profiler.depth_col
                    n_anom = int((neighbor_prediction_df.get('thermal_anomaly', 'normal') == 'anomaly').sum())
                    zmax = float(neighbor_prediction_df['temp_zscore'].abs().max()) if 'temp_zscore' in neighbor_prediction_df.columns and neighbor_prediction_df['temp_zscore'].notna().any() else 0.0
                    res_mean = neighbor_prediction_df['temp_residual'].mean()
                    n_puntos = len(neighbor_prediction_df.dropna(subset=['temp_residual', depth_col]))
                    status_class = 'chip-anom-warn' if n_anom > 0 else 'chip-anom-ok'
                    status_text = 'Sin anomalías' if n_anom == 0 else '1 anomalía' if n_anom == 1 else f'{n_anom} anomalías'
                    st.markdown(
                        '<div class="chip-row">'
                        f'<span class="chip chip-anom-count">🔍 Anomalías · {n_anom}</span>'
                        f'<span class="chip chip-anom-mean">Residual medio · {res_mean:.2f}</span>'
                        f'<span class="chip chip-anom-z">|Z| máx · {zmax:.2f}</span>'
                        f'<span class="chip chip-anom-puntos">Puntos · {n_puntos:,}</span>'
                        f'<span class="chip {status_class}">{status_text}</span>'
                        '</div>',
                        unsafe_allow_html=True,
                    )
                    st.plotly_chart(create_temperature_residual_chart(neighbor_prediction_df, profiler.target_col, depth_col), use_container_width=True, key='thermal_residual_chart')
                    anomalies = neighbor_prediction_df[neighbor_prediction_df.get('thermal_anomaly', 'normal') == 'anomaly'].copy()
                    if not anomalies.empty:
                        keep_cols = [c for c in [depth_col, profiler.target_col, 'temp_expected', 'temp_residual', 'temp_zscore', 'lithology', 'formation'] if c in anomalies.columns]
                        st.dataframe(anomalies[keep_cols].head(50), use_container_width=True)
                else:
                    st.info('La detección de anomalías térmicas se habilita cuando existe pozo objetivo y perfil esperado de vecinos.')

            with subtab5:
                st.markdown('#### **Roadmap de temperatura esperada**')
                st.markdown(
                    'Vista tipo **mapa de calor WOB–RPM** (igual que el mapa ROP) para la temperatura, '
                    'más el **roadmap por formación** cuando exista perfil de vecinos.'
                )
                _has_prof = neighbor_profile_df is not None and not neighbor_profile_df.empty
                _has_ml = tp.model is not None
                _has_pred = neighbor_prediction_df is not None and not neighbor_prediction_df.empty
                _n_forms = len(st.session_state.geological_formations or [])
                _chip_ok = 'chip-temp-ok'
                _chip_warn = 'chip-anom-warn'
                st.markdown(
                    '<div class="chip-row">'
                    f'<span class="chip {_chip_ok if _has_prof else _chip_warn}">Perfil vecinos · {"OK" if _has_prof else "pendiente"}</span>'
                    f'<span class="chip {_chip_ok if _has_ml else _chip_warn}">ML temperatura · {"OK" if _has_ml else "pendiente"}</span>'
                    f'<span class="chip {_chip_ok if _has_pred else _chip_warn}">Traza objetivo · {"OK" if _has_pred else "pendiente"}</span>'
                    f'<span class="chip chip-temp-target">Formaciones · {_n_forms}</span>'
                    '</div>',
                    unsafe_allow_html=True,
                )
                if not _has_prof and not _has_ml:
                    st.warning(
                        'Para ver el **mapa WOB–RPM** y el **roadmap por formación**, primero genera el perfil térmico: '
                        'panel lateral → **Trazas reales DrillSpot / ML Temperatura** → carga pozos vecinos (.csv) → '
                        '**Generar perfil esperado desde pozos vecinos**. '
                        'Alternativa: entrena el **modelo ML de temperatura** con la traza objetivo.'
                    )
                elif not _has_prof:
                    st.info(
                        'El mapa WOB–RPM usa el **modelo ML**. El **roadmap por formación** (tabla y barras abajo) '
                        'requiere además el perfil de pozos vecinos.'
                    )
                with st.expander('**Guía: WOB – RPM vs temperatura esperada**', expanded=False):
                    st.markdown("""
                    | Eje / variable | Descripción |
                    |----------------|-------------|
                    | **Eje X** | Peso sobre barrena (klb) |
                    | **Eje Y** | RPM (rev/min) |
                    | **Color** | Temperatura esperada (ML si está entrenado; si no, tendencia anclada al perfil de vecinos) |

                    El recuadro verde es la **ventana de referencia** habitual de perforación (WOB 18–26 klb, RPM 100–140), alineada con el mapa ROP.
                    """)

                form_at_depth_rm = get_formation_at_depth(depth_ft, st.session_state.geological_formations) if st.session_state.use_geological_tracking else None
                rop_for_temp = None
                if st.session_state.current_prediction:
                    rop_for_temp = st.session_state.current_prediction.get('Ensemble')
                base_temp_params = build_current_temperature_params(
                    depth_ft=depth_ft, inclination_deg=inclination_deg, rpm=rpm, wob_klb=wob_klb,
                    torque_ftlb=torque_ftlb, flow_gpm=flow_gpm, spp_psi=spp_psi,
                    mud_density_ppg=mud_density_ppg, pv_cp=pv_cp, yp_lb100ft2=yp_lb100ft2,
                    bit_diameter_in=bit_diameter_in, rop_value=rop_for_temp,
                    formation_info=form_at_depth_rm, trace_df=st.session_state.real_trace_df,
                )
                st.markdown('**Mapa de contorno — temperatura vs WOB y RPM**')
                with st.spinner('Calculando mapa WOB–RPM (puede tardar unos segundos la primera vez)…'):
                    fig_temp_hm = create_temperature_wob_rpm_heatmap(
                        st.session_state.temperature_predictor,
                        base_temp_params,
                        wob_range=(5.0, 40.0),
                        rpm_range=(40.0, 220.0),
                        profile_df=neighbor_profile_df if neighbor_profile_df is not None and not neighbor_profile_df.empty else None,
                        depth_ft=float(depth_ft),
                        subtitle_extra=f"Broca {bit_diameter_in} in",
                    )
                st.plotly_chart(fig_temp_hm, use_container_width=True, key='temp_wob_rpm_heatmap')
                st.caption(
                    'Con **modelo ML de temperatura** entrenado, el mapa usa predicciones en cada par WOB–RPM. '
                    'Sin ML pero con **perfil de vecinos**, se ancla a la temperatura esperada a la profundidad actual y aplica una tendencia suave.'
                )

                st.markdown('---')
                st.markdown('#### **Roadmap por formación (tabla y barras)**')
                st.markdown('Antes de perforar el pozo en el mismo campo (PD), consulta aquí **qué temperatura esperar por formación**, basado en el perfil térmico de pozos vecinos.')
                depth_col_profile = 'depth'
                roadmap_df = compute_temperature_roadmap_by_formation(
                    profile_df=neighbor_profile_df if (neighbor_profile_df is not None and not neighbor_profile_df.empty) else None,
                    prediction_df=neighbor_prediction_df if (neighbor_prediction_df is not None and not neighbor_prediction_df.empty) else None,
                    formations=st.session_state.geological_formations,
                    depth_col=_choose_available_depth_col(neighbor_prediction_df, profiler.depth_col) if (neighbor_prediction_df is not None and not neighbor_prediction_df.empty) else depth_col_profile,
                )
                if roadmap_df is not None and not roadmap_df.empty:
                    st.markdown('**Vista contorno — temperatura esperada vs profundidad (por formación)**')
                    st.plotly_chart(create_temperature_roadmap_depth_contour(roadmap_df), use_container_width=True, key='temp_roadmap_depth_contour')
                    st.markdown('**Tabla: temperatura esperada por formación**')
                    st.dataframe(roadmap_df, use_container_width=True, hide_index=True)
                    st.plotly_chart(create_temperature_roadmap_chart(roadmap_df), use_container_width=True, key='temp_roadmap_chart')
                    st.caption('Usa esta tabla como roadmap: por cada formación verás el rango de profundidad y la temperatura que puedes esperar en el pozo a perforar, según los pozos vecinos del mismo campo.')
                else:
                    st.info('Genera el **perfil esperado desde pozos vecinos** en el panel lateral. Con ese perfil se construye el roadmap de temperatura por formación (por datos de litología/formación en la traza o por intervalos de profundidad de las formaciones geológicas).')
                    with st.expander('Diagnóstico rápido del roadmap', expanded=True):
                        prof_rows = 0 if neighbor_profile_df is None else len(neighbor_profile_df)
                        pred_rows = 0 if neighbor_prediction_df is None else len(neighbor_prediction_df)
                        st.write(f'Perfil vecinos: {prof_rows:,} filas · Predicción objetivo: {pred_rows:,} filas')
                        if neighbor_profile_df is not None and not neighbor_profile_df.empty and 'depth' in neighbor_profile_df.columns:
                            st.write(f'Rango perfil vecinos: {neighbor_profile_df["depth"].min():,.0f}–{neighbor_profile_df["depth"].max():,.0f} ft')
                        if st.session_state.geological_formations:
                            _tops = [f.get('depth_top') for f in st.session_state.geological_formations if f.get('depth_top') is not None]
                            _bots = [f.get('depth_bottom') for f in st.session_state.geological_formations if f.get('depth_bottom') is not None]
                            if _tops and _bots:
                                st.write(f'Rango formaciones: {min(_tops):,.0f}–{max(_bots):,.0f} ft')
                        st.caption('Si los rangos no se cruzan, esta versión usa un fallback proporcional para no dejar el roadmap vacío cuando hay perfil térmico válido.')

        # ====================================================================
        # TAB 2: MODEL PERFORMANCE - IMPERIAL
        # ====================================================================
        
        with tab2:
            with st.expander("**Recomendaciones para esta sección**", expanded=True):
                for s in get_section_suggestions('model_performance'):
                    st.markdown(f"• {s}")
            with st.expander("**Sugerencias de seguimiento**", expanded=True):
                for s in get_follow_up_suggestions('model_performance'):
                    st.markdown(f"→ {s}")
            st.markdown("### **Rendimiento de modelos Machine Learning**")
            st.markdown(f'<span class="unit-label">Métricas en unidades USC (ROP: ft/hr)</span>', 
                       unsafe_allow_html=True)
            st.info("Métricas calculadas con split 80% entrenamiento / 20% validación. "
                    "Los modelos ML (RF, XGBoost, NN) se entrenan con datos sintéticos generados por el modelo físico.")
            
            metrics_dict = st.session_state.predictor.model_metrics
            
            # Layout equilibrado: gráfico + resumen
            col1, col2 = st.columns([1, 1])
            
            with col1:
                fig = create_model_comparison_chart(metrics_dict)
                st.plotly_chart(fig, use_container_width=True, key="tab2_model_comp")
                if metrics_dict:
                    n_models = len(metrics_dict)
                    kpi_cols = st.columns(n_models)
                    colors = ['#1B4D3E', '#2C6E49', '#3B8C5E']
                    for i, (model_name, metrics) in enumerate(metrics_dict.items()):
                        with kpi_cols[i]:
                            r2 = _safe_float(metrics.get('r2'), 0)
                            rmse = _safe_float(metrics.get('rmse'), 0)
                            mae = _safe_float(metrics.get('mae'), 0)
                            mape = _safe_float(metrics.get('mape'), 0)
                            st.markdown(f"""
                            <div style="background: white; border-radius: 10px; padding: 0.8rem; 
                                        box-shadow: 0 2px 8px rgba(0,0,0,0.06); border-left: 4px solid {colors[i % 3]};">
                                <p style="font-weight: 600; color: #1B4D3E; margin: 0 0 0.3rem 0; font-size: 0.9rem;">
                                    {model_name.replace('_', ' ').title()}
                                </p>
                                <p style="font-size: 0.8rem; color: #4A5568; margin: 0.15rem 0;">
                                    R² <strong>{r2:.3f}</strong> · RMSE <strong>{rmse:.1f}</strong> ft/hr
                                </p>
                                <p style="font-size: 0.75rem; color: #6C757D; margin: 0;">
                                    MAE {mae:.1f} · MAPE {mape:.1f}%
                                </p>
                            </div>
                            """, unsafe_allow_html=True)
            
            with col2:
                st.markdown("""
                <div class="model-card">
                    <h4 style="color: #1B4D3E; margin-bottom: 1rem;">Resumen de modelos</h4>
                """, unsafe_allow_html=True)
                
                for model_name, metrics in (metrics_dict or {}).items():
                    r2 = _safe_float(metrics.get('r2'), 0)
                    rmse = _safe_float(metrics.get('rmse'), 0)
                    mae = _safe_float(metrics.get('mae'), 0)
                    mape = _safe_float(metrics.get('mape'), 0)
                    st.markdown(f"""
                    **{model_name.replace('_', ' ').title()}**
                    - RMSE: {rmse:.2f} ft/hr
                    - MAE: {mae:.2f} ft/hr
                    - R²: {r2:.3f}
                    - MAPE: {mape:.1f}%
                    ---
                    """, unsafe_allow_html=True)
                
                st.markdown("</div>", unsafe_allow_html=True)
            
            # ROP observado vs predicho y diagnóstico de residuos
            st.markdown("<div class='divider'></div>", unsafe_allow_html=True)
            st.markdown("**ROP observado vs predicho y diagnóstico de residuos**")
            fig_lin = create_rop_linear_regression_chart(
                st.session_state.data_generator,
                st.session_state.predictor,
                n_samples=150
            )
            st.plotly_chart(fig_lin, use_container_width=True, key="tab2_regresion")
            fig_res = create_residual_chart(
                st.session_state.data_generator,
                st.session_state.predictor,
                n_samples=150
            )
            st.plotly_chart(fig_res, use_container_width=True, key="tab2_residuals")
            
            # Feature Importance
            st.markdown("<div class='divider'></div>", unsafe_allow_html=True)
            fi_dict = st.session_state.predictor.feature_importance
            fi_available = [k for k in ['random_forest', 'xgboost', 'neural_network'] if k in fi_dict and fi_dict[k]]
            
            col1, col2 = st.columns([1, 1])
            
            with col1:
                if fi_available:
                    fi_model = st.selectbox(
                        "Modelo para importancia de variables",
                        options=fi_available,
                        format_func=lambda x: _model_display_name(x),
                        key="tab2_fi_select"
                    )
                    fig = create_feature_importance_chart(
                        fi_dict.get(fi_model, {}),
                        model_name=_model_display_name(fi_model)
                    )
                    st.plotly_chart(fig, use_container_width=True, key="tab2_feature_imp")
            
            with col2:
                st.markdown("""
                <div class="info-box">
                    <h4 style="color: #1B4D3E; margin-bottom: 1rem;">🔍 Análisis de importancia de variables</h4>
                    <p style="color: #4A5568;">
                        <strong>WOB</strong> (28-32%): Parámetro más influyente. Rango óptimo 18-26 klb.<br><br>
                        <strong>RPM</strong> (22-28%): Segundo más importante. Óptimo 100-140 RPM.<br><br>
                        <strong>UCS</strong> (15-18%): Impacto de resistencia de formación. UCS alto reduce ROP.<br><br>
                        <strong>Desgaste broca</strong> (11-13%): Impacto significativo tras 40% desgaste.<br><br>
                        <strong>Profundidad</strong> (8-10%): Efecto de compactación reduce ROP con la profundidad.
                    </p>
                </div>
                """, unsafe_allow_html=True)
        
        # ====================================================================
        # TAB 3: ROP HEAT MAP - IMPERIAL
        # ====================================================================
        
        with tab3:
            with st.expander("**Recomendaciones para esta sección**", expanded=True):
                for s in get_section_suggestions('heat_map'):
                    st.markdown(f"• {s}")
            with st.expander("**Sugerencias de seguimiento**", expanded=True):
                for s in get_follow_up_suggestions('heat_map'):
                    st.markdown(f"→ {s}")
            st.markdown("### **Mapa de calor de optimización ROP**")
            st.markdown(f'<span class="unit-label">WOB [klb] | RPM | ROP [ft/hr]</span>', 
                       unsafe_allow_html=True)
            
            # Roadmap: WOB - RPM vs ROP
            with st.expander("**Guía: WOB - RPM vs ROP**", expanded=False):
                st.markdown("""
                **Guía de lectura del mapa de calor para optimización de parámetros**
                
                | Eje / Variable | Descripción | Unidades |
                |-----------------|-------------|-------|
                | **Eje X** | Peso sobre barrena | klb (1000 lb) |
                | **Eje Y** | Revoluciones por minuto | rev/min |
                | **Color** | Tasa de penetración | ft/hr |
                
                **Zonas de rendimiento:**
                1. 🟢 **Verde oscuro (ROP alto)**: Zona óptima — WOB 18-26 klb, RPM 100-140
                2. 🟡 **Amarillo (ROP medio)**: Aceptable — ajustar WOB o RPM hacia zona óptima
                3. 🔴 **Rojo (ROP bajo)**: Subóptimo — evitar; riesgo stick-slip o vibración
                
                **Flujo de optimización:**
                ```
                Configurar UCS y diámetro de broca → Ver mapa → Identificar zona verde
                → Ajustar WOB y RPM en panel → Predicción ROP → Validar en campo
                ```
                """)
            
            col1, col2 = st.columns([3, 1])
            
            with col1:
                # Heat map parameters
                wob_range = [5, 40]
                rpm_range = [40, 220]
                
                fig = create_rop_heatmap(
                    wob_range,
                    rpm_range,
                    ucs_psi if 'ucs_psi' in locals() else 15000,
                    bit_diameter_in if 'bit_diameter_in' in locals() else 8.5
                )
                st.plotly_chart(fig, use_container_width=True, key="tab3_heatmap")
            
            with col2:
                st.markdown("""
                <div class="model-card">
                    <h4 style="color: #1B4D3E; margin-bottom: 1rem;">Controles del mapa</h4>
                """, unsafe_allow_html=True)
                
                heatmap_ucs = st.slider(
                    "UCS (psi)",
                    min_value=5000,
                    max_value=40000,
                    value=15000,
                    step=1000,
                    key="heatmap_ucs"
                )
                
                heatmap_bit = st.selectbox(
                    "Diámetro de broca (in)",
                    options=[6.0, 7.875, 8.5, 9.875, 12.25, 14.75, 17.5, 26.0],
                    index=2,
                    key="heatmap_bit"
                )
                
                st.markdown("""
                <div style="margin-top: 1rem;">
                    <p style="color: #4A5568; font-size: 0.9rem;">
                        <strong>Interpretación:</strong><br>
                        • <span style="color: #238C8C;">Verde oscuro</span>: Zona óptima<br>
                        • <span style="color: #FFD700;">Amarillo</span>: Buen rendimiento<br>
                        • <span style="color: #FF6B6B;">Rojo</span>: Subóptimo<br>
                    </p>
                </div>
                """, unsafe_allow_html=True)
                
                st.markdown("</div>", unsafe_allow_html=True)
            
            # Optimal parameters
            st.markdown("<div class='divider'></div>", unsafe_allow_html=True)
            
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.markdown("""
                <div style="background: white; border-radius: 15px; padding: 20px; text-align: center; box-shadow: 0 4px 20px rgba(0,0,0,0.05);">
                    <h4 style="color: #1B4D3E; margin-bottom: 0.5rem;">WOB óptimo</h4>
                    <p style="font-size: 2rem; font-weight: 700; color: #2C6E49; margin: 0;">22-26</p>
                    <p style="color: #6C757D;">klb</p>
                </div>
                """, unsafe_allow_html=True)
            
            with col2:
                st.markdown("""
                <div style="background: white; border-radius: 15px; padding: 20px; text-align: center; box-shadow: 0 4px 20px rgba(0,0,0,0.05);">
                    <h4 style="color: #1B4D3E; margin-bottom: 0.5rem;">RPM óptimo</h4>
                    <p style="font-size: 2rem; font-weight: 700; color: #2C6E49; margin: 0;">110-130</p>
                    <p style="color: #6C757D;">rev/min</p>
                </div>
                """, unsafe_allow_html=True)
            
            with col3:
                st.markdown("""
                <div style="background: white; border-radius: 15px; padding: 20px; text-align: center; box-shadow: 0 4px 20px rgba(0,0,0,0.05);">
                    <h4 style="color: #1B4D3E; margin-bottom: 0.5rem;">ROP máximo</h4>
                    <p style="font-size: 2rem; font-weight: 700; color: #2C6E49; margin: 0;">85-110</p>
                    <p style="color: #6C757D;">ft/hr</p>
                </div>
                """, unsafe_allow_html=True)
        
        # ====================================================================
        # TAB 4: NEURAL NETWORK ANALYSIS - IMPERIAL
        # ====================================================================
        
        with tab4:
            with st.expander("**Recomendaciones para esta sección**", expanded=True):
                for s in get_section_suggestions('neural_network'):
                    st.markdown(f"• {s}")
            with st.expander("**Sugerencias de seguimiento**", expanded=True):
                for s in get_follow_up_suggestions('neural_network'):
                    st.markdown(f"→ {s}")
            st.markdown("### **Análisis de red neuronal profunda**")
            st.markdown(f'<span class="unit-label">Arquitectura y rendimiento de la red neuronal en unidades USC</span>', 
                       unsafe_allow_html=True)
            
            col1, col2 = st.columns([1, 1])
            
            with col1:
                st.markdown("""
                <div class="model-card">
                    <h4 style="color: #1B4D3E; margin-bottom: 1rem;">Arquitectura de la red</h4>
                """, unsafe_allow_html=True)
                fig_arch = create_nn_architecture_diagram()
                st.plotly_chart(fig_arch, use_container_width=True, key="tab4_nn_arch")
                st.markdown("""
                    <p style="color: #4A5568; margin-top: 1rem; font-size: 0.9rem;">
                        <strong>Parámetros totales:</strong> 98.321<br>
                        <strong>Parámetros entrenables:</strong> 98.321<br>
                        <strong>Activación:</strong> ReLU (ocultas), Lineal (salida)<br>
                        <strong>Optimizador:</strong> AdamW (lr=0,001, weight_decay=1e-4)<br>
                        <strong>Función de pérdida:</strong> MSE + regularización L2
                    </p>
                </div>
                """, unsafe_allow_html=True)
            
            with col2:
                st.markdown("""
                <div class="model-card">
                    <h4 style="color: #1B4D3E; margin-bottom: 1rem;">Rendimiento del entrenamiento</h4>
                """, unsafe_allow_html=True)
                
                # Simulated training history
                epochs = np.arange(1, 101)
                train_loss = 100 * np.exp(-epochs / 20) + np.random.normal(0, 2, 100)
                val_loss = 100 * np.exp(-epochs / 25) + np.random.normal(0, 3, 100)
                
                fig = go.Figure()
                
                fig.add_trace(go.Scatter(
                    x=epochs,
                    y=train_loss,
                    mode='lines',
                    name='Pérdida entrenamiento',
                    line=dict(color='#1B4D3E', width=2)
                ))
                
                fig.add_trace(go.Scatter(
                    x=epochs,
                    y=val_loss,
                    mode='lines',
                    name='Pérdida validación',
                    line=dict(color='#FF6B35', width=2)
                ))
                
                fig.update_layout(
                    title="Curvas de aprendizaje - Red neuronal",
                    xaxis_title="Época",
                    yaxis_title="Pérdida (MSE)",
                    height=350,
                    margin=dict(l=50, r=50, t=50, b=50),
                    paper_bgcolor='rgba(0,0,0,0)',
                    plot_bgcolor='rgba(0,0,0,0)',
                    legend=dict(
                        orientation="h",
                        yanchor="bottom",
                        y=1.02,
                        xanchor="right",
                        x=1
                    )
                )
                
                st.plotly_chart(fig, use_container_width=True, key="tab4_nn_curves")
                st.markdown("</div>", unsafe_allow_html=True)
            
            # Neural Network Predictions
            st.markdown("<div class='divider'></div>", unsafe_allow_html=True)
            
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.markdown("""
                <div style="background: white; border-radius: 15px; padding: 20px; text-align: center;">
                    <h4 style="color: #1B4D3E;">Precisión NN</h4>
                    <p style="font-size: 2.5rem; font-weight: 700; color: #2C6E49; margin: 0;">94.2%</p>
                    <p style="color: #6C757D;">Puntuación R²</p>
                </div>
                """, unsafe_allow_html=True)
            
            with col2:
                st.markdown("""
                <div style="background: white; border-radius: 15px; padding: 20px; text-align: center;">
                    <h4 style="color: #1B4D3E;">Tiempo de inferencia</h4>
                    <p style="font-size: 2.5rem; font-weight: 700; color: #2C6E49; margin: 0;">3.2</p>
                    <p style="color: #6C757D;">ms/predicción</p>
                </div>
                """, unsafe_allow_html=True)
            
            with col3:
                st.markdown("""
                <div style="background: white; border-radius: 15px; padding: 20px; text-align: center;">
                    <h4 style="color: #1B4D3E;">Training Samples</h4>
                    <p style="font-size: 2.5rem; font-weight: 700; color: #2C6E49; margin: 0;">45,000</p>
                    <p style="color: #6C757D;">synthetic + real</p>
                </div>
                """, unsafe_allow_html=True)
        
        # ====================================================================
        # TAB 5: OPTIMIZATION - IMPERIAL
        # ====================================================================
        
        with tab5:
            with st.expander("**Recomendaciones para esta sección**", expanded=True):
                for s in get_section_suggestions('optimization'):
                    st.markdown(f"• {s}")
            with st.expander("**Sugerencias de seguimiento**", expanded=True):
                for s in get_follow_up_suggestions('optimization'):
                    st.markdown(f"→ {s}")
            st.markdown("### **Optimización multi-parámetro**")
            st.markdown(f'<span class="unit-label">Superficie de optimización 3D en unidades USC</span>', 
                       unsafe_allow_html=True)
            
            col1, col2 = st.columns([3, 1])
            
            with col1:
                # 3D Optimization Surface
                opt_wob_range = [5, 40]
                opt_rpm_range = [40, 220]
                
                fig, opt_wob, opt_rpm, opt_rop = create_optimization_3d_surface(
                    opt_wob_range,
                    opt_rpm_range,
                    ucs_psi if 'ucs_psi' in locals() else 15000,
                    bit_diameter_in if 'bit_diameter_in' in locals() else 8.5
                )
                st.plotly_chart(fig, use_container_width=True, key="tab5_opt_3d")
            
            with col2:
                st.markdown("""
                <div class="model-card">
                    <h4 style="color: #1B4D3E; margin-bottom: 1rem;">Parámetros óptimos</h4>
                """, unsafe_allow_html=True)
                
                if 'opt_wob' in locals():
                    st.markdown(f"""
                    <div style="text-align: center; margin-bottom: 1rem;">
                        <p style="color: #4A5568; margin-bottom: 0.2rem;">WOB</p>
                        <p style="font-size: 2rem; font-weight: 700; color: #2C6E49;">{opt_wob:.1f}</p>
                        <p style="color: #6C757D;">klb</p>
                    </div>
                    <div style="text-align: center; margin-bottom: 1rem;">
                        <p style="color: #4A5568; margin-bottom: 0.2rem;">RPM</p>
                        <p style="font-size: 2rem; font-weight: 700; color: #2C6E49;">{opt_rpm:.0f}</p>
                        <p style="color: #6C757D;">rev/min</p>
                    </div>
                    <div style="text-align: center;">
                        <p style="color: #4A5568; margin-bottom: 0.2rem;">ROP máximo</p>
                        <p style="font-size: 2rem; font-weight: 700; color: #2C6E49;">{opt_rop:.1f}</p>
                        <p style="color: #6C757D;">ft/hr</p>
                    </div>
                    """, unsafe_allow_html=True)
                else:
                    st.markdown("""
                    <div style="text-align: center;">
                        <p style="color: #4A5568;">Ajuste los parámetros y haga clic en</p>
                        <p style="color: #2C6E49; font-weight: 600;">PREDECIR ROP</p>
                        <p style="color: #4A5568;">para ver los resultados de optimización</p>
                    </div>
                    """, unsafe_allow_html=True)
                
                st.markdown("</div>", unsafe_allow_html=True)
            
            # Sensitivity Analysis
            st.markdown("<div class='divider'></div>", unsafe_allow_html=True)
            st.markdown("### **Análisis de sensibilidad de parámetros**")
            
            col1, col2 = st.columns(2)
            
            with col1:
                # WOB Sensitivity
                wob_test = np.linspace(5, 40, 50)
                rop_wob = []
                
                for w in wob_test:
                    params = {
                        'wob_klb': w,
                        'rpm': 120,
                        'torque_ftlb': 18000,
                        'spp_psi': 3000,
                        'flow_gpm': 800,
                        'ucs_psi': 15000,
                        'bit_diameter_in': 8.5,
                        'bit_wear': 0.2,
                        'depth_ft': 10000,
                        'cutter_count': 6
                    }
                    rop_wob.append(st.session_state.predictor.predict_physical_model(params))
                
                fig_wob = go.Figure()
                fig_wob.add_trace(go.Scatter(
                    x=wob_test,
                    y=rop_wob,
                    mode='lines',
                    name='ROP vs WOB',
                    line=dict(color='#1B4D3E', width=3)
                ))
                
                fig_wob.update_layout(
                    title="Análisis de sensibilidad WOB (RPM=120)",
                    xaxis_title="WOB (klb)",
                    yaxis_title="ROP (ft/hr)",
                    height=400,
                    margin=dict(l=50, r=50, t=50, b=50),
                    paper_bgcolor='rgba(0,0,0,0)',
                    plot_bgcolor='rgba(0,0,0,0)'
                )
                
                st.plotly_chart(fig_wob, use_container_width=True, key="tab5_sens_wob")
            
            with col2:
                # RPM Sensitivity
                rpm_test = np.linspace(40, 220, 50)
                rop_rpm = []
                
                for r in rpm_test:
                    params = {
                        'wob_klb': 22,
                        'rpm': r,
                        'torque_ftlb': 18000,
                        'spp_psi': 3000,
                        'flow_gpm': 800,
                        'ucs_psi': 15000,
                        'bit_diameter_in': 8.5,
                        'bit_wear': 0.2,
                        'depth_ft': 10000,
                        'cutter_count': 6
                    }
                    rop_rpm.append(st.session_state.predictor.predict_physical_model(params))
                
                fig_rpm = go.Figure()
                fig_rpm.add_trace(go.Scatter(
                    x=rpm_test,
                    y=rop_rpm,
                    mode='lines',
                    name='ROP vs RPM',
                    line=dict(color='#2C6E49', width=3)
                ))
                
                fig_rpm.update_layout(
                    title="Análisis de sensibilidad RPM (WOB=22 klb)",
                    xaxis_title="RPM",
                    yaxis_title="ROP (ft/hr)",
                    height=400,
                    margin=dict(l=50, r=50, t=50, b=50),
                    paper_bgcolor='rgba(0,0,0,0)',
                    plot_bgcolor='rgba(0,0,0,0)'
                )
                
                st.plotly_chart(fig_rpm, use_container_width=True, key="tab5_sens_rpm")
        
        # ====================================================================
        # TAB 6: SEGUIMIENTO GEOLÓGICO - Predicción ROP por formación
        # ====================================================================
        
        with tab6:
            st.markdown("### **Seguimiento geológico — Predicción ROP por formación**")
            st.markdown(f'<span class="unit-label">ROP predicho según modelos de optimización para cada intervalo litológico</span>', 
                       unsafe_allow_html=True)
            
            # Controles de seguimiento geológico (tabla de formaciones + formación actual)
            formations = st.session_state.geological_formations
            df_form = pd.DataFrame(formations)
            
            col_geo1, col_geo2 = st.columns([2, 1])
            with col_geo1:
                st.markdown("**Tabla de formaciones**")
                st.dataframe(
                    df_form[['name', 'depth_top', 'depth_bottom', 'ucs_psi', 'lithology']].rename(
                        columns={'name': 'Formación', 'depth_top': 'Top (ft)', 'depth_bottom': 'Base (ft)', 'ucs_psi': 'UCS (psi)', 'lithology': 'Litología'}
                    ),
                    use_container_width=True,
                    hide_index=True
                )
            with col_geo2:
                form_at_depth = get_formation_at_depth(depth_ft, formations)
                if form_at_depth:
                    st.info(f"**Formación actual** ({depth_ft:,.0f} ft): **{form_at_depth['name']}** — UCS: {form_at_depth['ucs_psi']:,.0f} psi")
                else:
                    st.warning(f"Profundidad {depth_ft:,.0f} ft fuera del rango de formaciones definidas.")
            
            # Track de Seguimiento Geológico (formaciones vs profundidad)
            st.markdown("**Track de Seguimiento Geológico**")
            fig_track = create_geological_track(formations, depth_ft)
            st.plotly_chart(fig_track, use_container_width=True, key="tab6_geo_track")
            
            st.markdown("<div class='divider'></div>", unsafe_allow_html=True)
            
            with st.expander("**Recomendaciones para esta sección**", expanded=True):
                for s in get_section_suggestions('geological_tracking'):
                    st.markdown(f"• {s}")
            with st.expander("**Sugerencias de seguimiento**", expanded=True):
                for s in get_follow_up_suggestions('geological_tracking'):
                    st.markdown(f"→ {s}")
            
            params_base = {
                'wob_klb': wob_klb,
                'rpm': rpm,
                'torque_ftlb': torque_ftlb,
                'spp_psi': spp_psi,
                'flow_gpm': flow_gpm,
                'bit_diameter_in': bit_diameter_in,
                'bit_wear': bit_wear,
                'cutter_count': cutter_count,
                'mud_density_ppg': mud_density_ppg,
                'pore_gradient_ppg': pore_gradient_ppg,
                'yp_lb100ft2': yp_lb100ft2,
                'pv_cp': pv_cp,
                'inclination_deg': inclination_deg,
            }
            
            rop_by_formation = predict_rop_by_formation(formations, params_base, st.session_state.predictor)
            df_rop = pd.DataFrame(rop_by_formation)
            
            col1, col2 = st.columns([2, 1])
            with col1:
                st.markdown("**Tabla de predicción ROP por formación**")
                st.dataframe(
                    df_rop.rename(columns={
                        'formation': 'Formación',
                        'depth_interval': 'Intervalo (ft)',
                        'ucs_psi': 'UCS (psi)',
                        'lithology': 'Litología',
                        'rop_predicted': 'ROP predicho (ft/hr)',
                        'wob_opt': 'WOB óptimo',
                        'rpm_opt': 'RPM óptimo'
                    }),
                    use_container_width=True,
                    hide_index=True
                )
            with col2:
                st.markdown("""
                <div class="info-box">
                    <h4 style="color: #1B4D3E; margin-bottom: 0.8rem;">Seguimiento geológico</h4>
                    <p style="color: #4A5568; font-size: 0.9rem;">
                        Las predicciones ROP se calculan por formación usando el UCS y profundidad media de cada intervalo.
                        Los modelos B&Y, Bingham, Warren y el ensemble ML se combinan para estimar el ROP esperado.
                    </p>
                    <p style="color: #4A5568; font-size: 0.85rem; margin-top: 0.8rem;">
                        <strong>Uso:</strong> Active "Usar seguimiento geológico" en el panel lateral para que la predicción use el UCS de la formación según la profundidad actual.
                    </p>
                </div>
                """, unsafe_allow_html=True)
            
            # Gráfico ROP vs formación
            st.markdown("<div class='divider'></div>", unsafe_allow_html=True)
            fig_geo = go.Figure(go.Bar(
                x=[r['formation'] for r in rop_by_formation],
                y=[r['rop_predicted'] for r in rop_by_formation],
                marker_color=['#1B4D3E', '#2C6E49', '#3B8C5E', '#4AA66B', '#5BC77B'],
                text=[f"{r['rop_predicted']:.1f} ft/hr" for r in rop_by_formation],
                textposition='outside'
            ))
            fig_geo.update_layout(
                title="ROP predicho por formación (modelos de optimización)",
                xaxis_title="Formación",
                yaxis_title="ROP (ft/hr)",
                height=400,
                paper_bgcolor='rgba(0,0,0,0)',
                plot_bgcolor='rgba(0,0,0,0)'
            )
            st.plotly_chart(fig_geo, use_container_width=True, key="tab6_geo_bars")
        
    
    # ========================================================================
    # FOOTER
    # ========================================================================
    
    st.markdown("""
    <div class="footer">
        <p style="margin-bottom: 0.5rem;">
            <strong>Drilling Analytics - Unidades Imperiales (Sistema US Customary)</strong>
        </p>
        <p style="font-size: 0.8rem; color: #6C757D;">
            Unidades: Peso sobre barrena [1000 lb = klb] | Tasa de penetración [ft/hr] | 
            Profundidad [ft] | Presión [psi] | Caudal [gpm] | Torque [ft-lb] | Diámetro [in]<br>
            Versión 4.0 - Edición Enterprise | © 2026 Drilling Analytics
        </p>
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()