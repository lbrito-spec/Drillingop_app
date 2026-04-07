# Density @ °C unificada en una sola columna (valor @ temperatura)
"""Drilling KPI & Mechanical Efficiency Report (Streamlit).

Versión .py con mejoras profesionales:
- funciones separadas y reutilizables
- validaciones de columnas y datos
- cache para carga de Excel
- manejo consistente de session_state
- estilos y layout más limpios
"""

from __future__ import annotations

import copy
import io
import json
import math
import os
import re
import shutil
import tempfile
import time
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path


def _combine_value_temp(value, temp):
    if value is None and temp is None:
        return None
    v = "" if value is None else str(value)
    t = "" if temp is None else str(temp)
    return f"{v} @ {t}" if t != "" else v
from typing import Iterable, List, Tuple
from urllib.parse import urlencode
from textwrap import wrap
import smtplib
from email.message import EmailMessage

import numpy as np
import pandas as pd

# ===============================
# Rig Activity: códigos -> etiquetas
# ===============================
RIG_ACTIVITY_CODE_MAP = {
    0: "In Slips",
    11: "In Slips - Pump",
    21: "Drilling",
    22: "Slide Drill",
    23: "Slide Oscillate Drill",
    31: "Reaming",
    32: "Back Reaming",
    50: "Static",
    51: "Static - Rotate&Pump",
    52: "Static - Pump",
    53: "Static - Rotate",
    54: "Surface Operations",
    61: "Run In - Tripping",
    62: "Run In - Pump",
    63: "Run In - Rotate",
    64: "Pull Up - Pump",
    65: "Pull Up - Rotate",
    66: "Pull Up - Trip Out",
    98: "Unknown",
    99: "Missing Input",
}

DEFAULT_DRILLING_ACTIVITY_CODES = {21, 22, 23}


def normalize_rig_activity_series(s: pd.Series) -> pd.Series:
    """Convierte una serie de Rig Activity (códigos numéricos o texto) a etiquetas legibles."""
    if s is None:
        return s
    # Try numeric conversion
    s_num = pd.to_numeric(s, errors="coerce")
    if s_num.notna().any():
        return s_num.map(lambda x: RIG_ACTIVITY_CODE_MAP.get(int(x), f"Code {int(x)}") if pd.notna(x) else None)
    # Fallback string normalization
    return s.astype(str)


# =====================
# Filtros de perforación
# =====================
def filter_drilling_physical(df: pd.DataFrame, depth_col: str | None, min_delta_m: float = 0.05) -> pd.DataFrame:
    """Filtra perforación real usando avance de Hole/Bit Depth (delta profundidad > umbral)."""
    if df is None or df.empty or not depth_col or depth_col not in df.columns:
        return df
    d = df.copy()
    d[depth_col] = pd.to_numeric(d[depth_col], errors="coerce")
    d["_d_depth"] = d[depth_col].diff()
    return d[d["_d_depth"] > float(min_delta_m)].drop(columns=["_d_depth"], errors="ignore")


def filter_by_rig_activity(df: pd.DataFrame, activity_col: str | None, drilling_values: list[str]) -> pd.DataFrame:
    """Filtra por estados de Rig Activity considerados perforación."""
    if df is None or df.empty or not activity_col or activity_col not in df.columns or not drilling_values:
        return df
    allowed = {str(v).upper().strip() for v in drilling_values}
    s = df[activity_col].astype(str).str.upper().str.strip()
    return df[s.isin(allowed)]

import plotly.express as px
import plotly.graph_objects as go
import requests
from scipy.signal import find_peaks, savgol_filter
from scipy.stats import binned_statistic_2d
import streamlit as st
from dotenv import load_dotenv
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


# ---- Anotaciones sobre evidencia (Alertas) ----
_ANNOTATION_COLORS = {
    "red": (239, 68, 68),
    "yellow": (234, 179, 8),
    "blue": (37, 99, 235),
    "green": (34, 197, 94),
}


def _resolve_annotation_color(c) -> tuple[int, int, int]:
    """Acepta nombre (red), hex (#RRGGBB) o tupla RGB."""
    if isinstance(c, (tuple, list)) and len(c) >= 3:
        return int(c[0]), int(c[1]), int(c[2])
    if not isinstance(c, str):
        return _ANNOTATION_COLORS["red"]
    s = c.strip()
    if s.lower() in _ANNOTATION_COLORS:
        return _ANNOTATION_COLORS[s.lower()]
    if s.startswith("#"):
        h = s[1:]
        if len(h) == 6:
            try:
                return tuple(int(h[i : i + 2], 16) for i in (0, 2, 4))  # type: ignore[return-value]
            except ValueError:
                pass
        if len(h) == 3:
            try:
                return tuple(int(h[i] + h[i], 16) for i in range(3))  # type: ignore[return-value]
            except (ValueError, TypeError):
                pass
    return _ANNOTATION_COLORS["red"]


def _draw_arrow_line(draw, x1: int, y1: int, x2: int, y2: int, color, width: int = 4, head: int = 18) -> None:
    """Segmento + punta triangular en (x2, y2)."""
    dx = float(x2 - x1)
    dy = float(y2 - y1)
    length = math.hypot(dx, dy)
    if length < 1e-6:
        return
    ux, uy = dx / length, dy / length
    # Recortar el trazo un poco para que la punta no quede debajo del triángulo
    trim = min(head * 0.85, length * 0.95)
    bx = x2 - ux * trim
    by = y2 - uy * trim
    draw.line((x1, y1, bx, by), fill=color, width=width)
    spread = 0.45  # radianes
    lx = x2 - head * math.cos(math.atan2(dy, dx) - spread)
    ly = y2 - head * math.sin(math.atan2(dy, dx) - spread)
    rx = x2 - head * math.cos(math.atan2(dy, dx) + spread)
    ry = y2 - head * math.sin(math.atan2(dy, dx) + spread)
    draw.polygon([(x2, y2), (lx, ly), (rx, ry)], fill=color)


def draw_annotations(img, annotations):
    if img is None:
        return img
    img = img.copy()
    draw = ImageDraw.Draw(img)
    for a in annotations:
        fill_rgb = _resolve_annotation_color(a.get("color", "red"))
        if a["type"] == "arrow":
            _draw_arrow_line(
                draw,
                int(a["x1"]),
                int(a["y1"]),
                int(a["x2"]),
                int(a["y2"]),
                fill_rgb,
                width=4,
                head=18,
            )
        elif a["type"] == "box":
            x1, y1, x2, y2 = int(a["x1"]), int(a["y1"]), int(a["x2"]), int(a["y2"])
            draw.rectangle((x1, y1, x2, y2), outline=fill_rgb, width=4)
        elif a["type"] == "text":
            x, y = int(a["x1"]), int(a["y1"])
            # Pins del lienzo (círculo Fabric): siempre dibujar el contorno aunque el texto esté vacío
            rx = a.get("rx")
            ry = a.get("ry")
            irx = iry = 0
            if rx is not None and ry is not None:
                irx = max(int(rx), 3)
                iry = max(int(ry), 3)
                pen = max(2, min(irx, iry) // 12 + 3)
                pen = min(pen, 6)
                draw.ellipse((x - irx, y - iry, x + irx, y + iry), outline=fill_rgb, width=pen)
            txt = str(a.get("text") or "")
            if txt:
                ty = y + (iry + 6 if iry else 0)
                for ox, oy in ((-1, 0), (1, 0), (0, -1), (0, 1)):
                    draw.text((x + ox, ty + oy), txt, fill=(0, 0, 0))
                draw.text((x, ty), txt, fill=fill_rgb)
    return img


def _prepare_evidence_for_canvas(
    img: Image.Image, max_side: int = 960
) -> tuple[Image.Image, float, float, int, int]:
    """Escala la imagen para el lienzo y devuelve factores sx,sy hacia píxeles originales."""
    img = img.convert("RGB")
    w0, h0 = img.size
    if w0 < 1 or h0 < 1:
        return img, 1.0, 1.0, w0, h0
    s = min(1.0, float(max_side) / float(max(w0, h0)))
    cw = max(1, int(round(w0 * s)))
    ch = max(1, int(round(h0 * s)))
    bg = img.resize((cw, ch), Image.Resampling.LANCZOS)
    sx = w0 / cw
    sy = h0 / ch
    return bg, sx, sy, w0, h0


def _fabric_line_global_xy(obj: dict) -> tuple[float, float, float, float]:
    """Extremos de una línea Fabric.js en coordenadas del lienzo."""
    left = float(obj.get("left") or 0)
    top = float(obj.get("top") or 0)
    x1 = float(obj.get("x1") or 0)
    y1 = float(obj.get("y1") or 0)
    x2 = float(obj.get("x2") or 0)
    y2 = float(obj.get("y2") or 0)
    ang = math.radians(float(obj.get("angle") or 0))

    def to_world(px: float, py: float) -> tuple[float, float]:
        rx = px * math.cos(ang) - py * math.sin(ang)
        ry = px * math.sin(ang) + py * math.cos(ang)
        return left + rx, top + ry

    ax1, ay1 = to_world(x1, y1)
    ax2, ay2 = to_world(x2, y2)
    return ax1, ay1, ax2, ay2


def _fabric_rect_global_xy(obj: dict) -> tuple[float, float, float, float]:
    """Esquinas de rectángulo Fabric en coordenadas del lienzo (eje alineado)."""
    left = float(obj.get("left") or 0)
    top = float(obj.get("top") or 0)
    w = float(obj.get("width") or 0) * float(obj.get("scaleX") or 1)
    h = float(obj.get("height") or 0) * float(obj.get("scaleY") or 1)
    ang = math.radians(float(obj.get("angle") or 0))
    # Vértices locales desde esquina superior izquierda lógica
    corners = [(0, 0), (w, 0), (w, h), (0, h)]
    wx: list[float] = []
    wy: list[float] = []
    for px, py in corners:
        rx = px * math.cos(ang) - py * math.sin(ang)
        ry = px * math.sin(ang) + py * math.cos(ang)
        wx.append(left + rx)
        wy.append(top + ry)
    return min(wx), min(wy), max(wx), max(wy)


def _fabric_circle_center_xy(obj: dict) -> tuple[float, float]:
    """Centro aproximado de círculo/elipse Fabric."""
    left = float(obj.get("left") or 0)
    top = float(obj.get("top") or 0)
    return left, top


def _fabric_pin_radii_canvas(obj: dict) -> tuple[float, float]:
    """Radios del círculo/elipse en coords del lienzo (antes de sx/sy → imagen evidencia)."""
    typ = str(obj.get("type") or "").lower()
    scx = abs(float(obj.get("scaleX") or 1))
    scy = abs(float(obj.get("scaleY") or 1))
    if typ == "circle":
        r = float(obj.get("radius") or 0)
        if r <= 0:
            w = float(obj.get("width") or 0)
            r = (w / 2) if w > 0 else 24.0
        return max(r * scx, 4.0), max(r * scy, 4.0)
    rw = float(obj.get("width") or 0)
    rh = float(obj.get("height") or 0)
    rx = ((rw / 2) if rw > 0 else 24.0) * scx
    ry = ((rh / 2) if rh > 0 else 24.0) * scy
    return max(rx, 4.0), max(ry, 4.0)


def _fabric_pin_center_canvas(obj: dict, rx: float, ry: float) -> tuple[float, float]:
    """Centro en coords del lienzo según originX/originY (Fabric); sin origen → left/top = centro (canvas drawable)."""
    left = float(obj.get("left") or 0)
    top = float(obj.get("top") or 0)
    if "originX" not in obj and "originY" not in obj:
        return left, top
    ox = str(obj.get("originX") or "left").lower()
    oy = str(obj.get("originY") or "top").lower()
    if ox == "center":
        cx = left
    elif ox == "right":
        cx = left - rx
    else:
        cx = left + rx
    if oy == "center":
        cy = top
    elif oy == "bottom":
        cy = top - ry
    else:
        cy = top + ry
    return cx, cy


def _fabric_make_selectable_for_transform(d: dict | None) -> dict | None:
    """Fabric: asegura selección y movimiento en modo transform (selectable, locks, controles)."""
    if not isinstance(d, dict):
        return d
    out = copy.deepcopy(d)

    def _unlock_movement(o: dict) -> None:
        o["selectable"] = True
        o["evented"] = True
        o["hasControls"] = True
        o["hasBorders"] = True
        o["lockMovementX"] = False
        o["lockMovementY"] = False
        o["lockScalingX"] = False
        o["lockScalingY"] = False
        o["lockRotation"] = False

    def fix_obj(o: dict) -> None:
        if not isinstance(o, dict):
            return
        t = str(o.get("type") or "").lower()
        if t == "group" and isinstance(o.get("objects"), list):
            _unlock_movement(o)
            for c in o["objects"]:
                fix_obj(c)
            return
        _unlock_movement(o)

    for o in out.get("objects") or []:
        fix_obj(o)
    return out


def _fabric_flat_objects(json_data) -> list[dict]:
    """Aplana objetos Fabric (incluye círculos dentro de `group`)."""
    if json_data is None:
        return []
    if isinstance(json_data, str):
        try:
            json_data = json.loads(json_data)
        except Exception:
            return []
    root = json_data.get("objects") if isinstance(json_data, dict) else None
    if not root:
        return []

    out: list[dict] = []

    def walk(objs) -> None:
        for o in objs or []:
            if not isinstance(o, dict):
                continue
            typ = str(o.get("type") or "").lower()
            if typ == "group" and isinstance(o.get("objects"), list):
                walk(o["objects"])
            else:
                out.append(o)

    walk(root)
    return out


def _stroke_to_store(stroke: str | None, fallback: str) -> str:
    if stroke and isinstance(stroke, str) and stroke.strip().startswith("#"):
        return stroke.strip()
    if stroke and isinstance(stroke, str) and "rgb" in stroke.lower():
        m = re.search(r"rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", stroke)
        if m:
            r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
            return f"#{r:02x}{g:02x}{b:02x}"
    return fallback


def _count_fabric_pins(json_data) -> int:
    """Cuenta círculos/elipses (marcas de comentario) en el JSON del canvas."""
    return sum(
        1
        for o in _fabric_flat_objects(json_data)
        if str(o.get("type") or "").lower() in ("circle", "ellipse")
    )


def canvas_fabric_json_to_annotations(
    json_data,
    sx: float,
    sy: float,
    pin_texts: list[str],
    default_stroke: str,
) -> list[dict]:
    """Convierte salida de streamlit-drawable-canvas (Fabric) en lista para draw_annotations."""
    if json_data is None:
        return []
    if isinstance(json_data, str):
        try:
            json_data = json.loads(json_data)
        except Exception:
            return []
    objs = _fabric_flat_objects(json_data)
    if not objs:
        return []

    ann: list[dict] = []
    circle_i = 0
    for obj in objs:
        if not isinstance(obj, dict):
            continue
        typ = str(obj.get("type") or "").lower()
        stroke = _stroke_to_store(obj.get("stroke"), default_stroke)

        if typ == "line":
            ax1, ay1, ax2, ay2 = _fabric_line_global_xy(obj)
            x1 = int(round(ax1 * sx))
            y1 = int(round(ay1 * sy))
            x2 = int(round(ax2 * sx))
            y2 = int(round(ay2 * sy))
            ann.append(
                {"type": "arrow", "color": stroke, "x1": x1, "y1": y1, "x2": x2, "y2": y2, "text": ""}
            )
        elif typ == "rect":
            x1f, y1f, x2f, y2f = _fabric_rect_global_xy(obj)
            ann.append(
                {
                    "type": "box",
                    "color": stroke,
                    "x1": int(round(x1f * sx)),
                    "y1": int(round(y1f * sy)),
                    "x2": int(round(x2f * sx)),
                    "y2": int(round(y2f * sy)),
                    "text": "",
                }
            )
        elif typ in ("circle", "ellipse"):
            rx_c, ry_c = _fabric_pin_radii_canvas(obj)
            cx_c, cy_c = _fabric_pin_center_canvas(obj, rx_c, ry_c)
            txt = ""
            if circle_i < len(pin_texts):
                txt = str(pin_texts[circle_i] or "")
            circle_i += 1
            ann.append(
                {
                    "type": "text",
                    "color": stroke,
                    "x1": int(round(cx_c * sx)),
                    "y1": int(round(cy_c * sy)),
                    "x2": int(round(cx_c * sx)),
                    "y2": int(round(cy_c * sy)),
                    "text": txt,
                    "rx": max(4, int(round(rx_c * sx))),
                    "ry": max(4, int(round(ry_c * sy))),
                }
            )
    return ann


from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from streamlit_echarts import st_echarts

try:
    import streamlit_drawable_canvas_compat  # noqa: F401 — image_to_url (Streamlit nuevo)
except ImportError:
    pass
try:
    from streamlit_drawable_canvas import st_canvas as _st_canvas_draw

    _HAS_DRAWABLE_CANVAS = True
except ImportError:
    _st_canvas_draw = None
    _HAS_DRAWABLE_CANVAS = False

try:
    import winsound  # type: ignore
except Exception:
    winsound = None


# =========================
# Configuración
# =========================
APP_TITLE = "Drilling KPI & Mechanical Efficiency Report"
BASE_DIR = Path(__file__).resolve().parent
def _resolve_logo_path() -> Path:
    candidates = [
        BASE_DIR / "assets" / "LogoDS.png",
        BASE_DIR / "LogoDS.png",
        BASE_DIR.parent / "assets" / "LogoDS.png",
        BASE_DIR.parent / "LogoDS.png",
        Path.cwd() / "assets" / "LogoDS.png",
        Path.cwd() / "LogoDS.png",
        Path.cwd() / "apps" / "DO_app_rogii" / "assets" / "LogoDS.png",
    ]
    for p in candidates:
        try:
            if p.exists():
                return p
        except Exception:
            continue
    return BASE_DIR / "assets" / "LogoDS.png"

LOGO_PATH = _resolve_logo_path()
SHEET_NAME = "worksheet"
PLOTLY_TEMPLATE = "plotly_white"
COLOR_SEQ = ["#2563EB", "#10B981", "#F59E0B", "#EF4444", "#8B5CF6", "#14B8A6"]
# Sin barra de herramientas (iconos de zoom, etc.) en gráficas Plotly
PLOTLY_CONFIG = {"displayModeBar": False, "displaylogo": False}
# Heatmap ROP dashboard: barra al pasar el mouse + export PNG a mayor escala
PLOTLY_CONFIG_ROP_DASH = {
    **PLOTLY_CONFIG,
    "displayModeBar": "hover",
    "modeBarButtonsToRemove": ["lasso2d", "select2d"],
    "toImageButtonOptions": {
        "format": "png",
        "filename": "heatmap_rop_mejor_zona",
        "scale": 3,
    },
}
# Degradado continuo tipo panel pro: azul profundo → teal/cian → amarillo → naranja (sin rojo duro)
ROP_HEATMAP_COLORSCALE: list[list] = [
    [0.0, "rgb(8, 16, 40)"],
    [0.17, "rgb(37, 72, 170)"],
    [0.38, "rgb(8, 131, 168)"],
    [0.55, "rgb(45, 185, 168)"],
    [0.74, "rgb(234, 179, 8)"],
    [1.0, "rgb(251, 146, 60)"],
]
ROP_HEATMAP_LABEL_TOP_FRACTION = 0.15
# Conversión MPa -> ksi (kpsi): 1 MPa = 0.145038 ksi
MPA_TO_KSI = 0.145038
MODE_NORMALIZATION = {
    "OSCILLATION_SLIDE": "SLIDE",
    "OSCILLATION": "SLIDE",
    "SLIDING": "SLIDE",
    "ROTARY_DRILLING": "ROTARY",
    "ROTATE": "ROTARY",
    "RSS": "ROTARY",
}


def _load_dotenv_files() -> None:
    """
    Carga variables desde archivos .env.

    Nota: el token NO va dentro de la carpeta `.venv` (eso es solo el entorno virtual de Python).
    Coloca un archivo llamado exactamente `.env` (no `.env.txt`) en:
      - la misma carpeta que este script (recomendado), y/o
      - la carpeta desde la que ejecutas `streamlit run` (cwd).

    Formato: SOLO_ACCESS_TOKEN=tu_token_aqui (sin comillas salvo que las necesites)

    Se prueba cwd → carpeta padre del script → carpeta del script; el último .env encontrado
    gana (override=True) para que el .env junto al .py tenga prioridad sobre uno genérico en cwd.
    """
    candidates = [
        Path.cwd() / ".env",
        BASE_DIR.parent / ".env",
        BASE_DIR / ".env",
    ]
    seen: set[Path] = set()
    for p in candidates:
        try:
            resolved = p.resolve()
        except OSError:
            continue
        if resolved in seen:
            continue
        if not resolved.is_file():
            continue
        seen.add(resolved)
        load_dotenv(resolved, override=True)


_load_dotenv_files()

API_DEFAULT_BASE_URL = os.getenv("SOLO_BASE_URL", "https://solo.cloud").rstrip("/")
API_DEFAULT_TOKEN = os.getenv("SOLO_ACCESS_TOKEN")

# SMTP para envío de bitácora de lodo (usando st.secrets)
def _secret(name, default=""):
    try:
        if name in st.secrets:
            return st.secrets[name]
    except Exception:
        pass
    return os.getenv(name, default)

MUD_SMTP_SERVER = _secret("MUD_SMTP_SERVER", "smtp.gmail.com")
MUD_SMTP_PORT = int(_secret("MUD_SMTP_PORT", "587"))
MUD_SMTP_USER = _secret("MUD_SMTP_USER", "")
MUD_SMTP_PASS = _secret("MUD_SMTP_PASS", "")
MUD_SMTP_FROM = _secret("MUD_SMTP_FROM", MUD_SMTP_USER)
MUD_SMTP_TO = _secret("MUD_SMTP_TO", "solobox+pemex@rogii.com")


def get_solo_credentials(prefix: str = "solo") -> tuple[str, str]:
    """
    Carga base_url y token desde:
      1) st.session_state (para no volver a pedirlo en la misma sesión)
      2) variables de entorno/.env (SOLO_BASE_URL, SOLO_ACCESS_TOKEN)
    """
    base_url_key = f"{prefix}_base_url"
    token_key = f"{prefix}_token"

    base_url = st.session_state.get(base_url_key)
    token = st.session_state.get(token_key)

    if not base_url:
        base_url = os.getenv("SOLO_BASE_URL", API_DEFAULT_BASE_URL).rstrip("/")
        st.session_state[base_url_key] = base_url

    if not token:
        token = (os.getenv("SOLO_ACCESS_TOKEN") or "").strip()
        if token:
            st.session_state[token_key] = token

    return st.session_state.get(base_url_key, API_DEFAULT_BASE_URL), st.session_state.get(token_key, "")


def render_solo_connection_ui(prefix: str = "solo", label: str = "Conexión SOLO") -> tuple[str, str]:
    """
    UI mínima: si el token está en .env, no lo pide.
    Permite cambiarlo manualmente si hace falta.
    """
    base_url_key = f"{prefix}_base_url"
    token_key = f"{prefix}_token"

    base_url, token = get_solo_credentials(prefix=prefix)

    with st.expander(label, expanded=False):
        st.text_input(tr("api_base_url"), value=base_url, key=base_url_key)

        if token:
            st.success(tr("solo_token_loaded"))
            if st.button(tr("solo_change_token"), key=f"{prefix}_change_token"):
                st.session_state[token_key] = ""
                st.rerun()
        else:
            st.text_input(
                tr("api_token"),
                value="",
                type="password",
                key=token_key,
                help=tr("solo_token_help"),
            )

    return st.session_state.get(base_url_key, base_url), st.session_state.get(token_key, "")


@dataclass(frozen=True)
class RunInfo:
    name: str
    start_depth: float
    end_depth: float
    start_time: pd.Timestamp
    end_time: pd.Timestamp


REQUIRED_COLUMNS = [
    "Mode",
    "Start Depth",
    "End Depth",
    "Survey MD",
    "Inclination",
    "Azimuth",
    "Distance",
    "ROP",
    "DLS",
    "WOB",
    "RPM",
    "Start",
    "End",
]

PLANNED_COL_CANDIDATES = {
    "ROP": [
        "ROP Plan",
        "ROP Planned",
        "ROP Programmed",
        "ROP Target",
        "ROP (Planned)",
        "ROP_prog",
        "ROP_Prog",
    ],
    "WOB": [
        "WOB Plan",
        "WOB Planned",
        "WOB Programmed",
        "WOB Target",
        "WOB (Planned)",
        "WOB_prog",
        "WOB_Prog",
    ],
    "RPM": [
        "RPM Plan",
        "RPM Planned",
        "RPM Programmed",
        "RPM Target",
        "RPM (Planned)",
        "RPM_prog",
        "RPM_Prog",
    ],
}

TORQUE_COL_CANDIDATES = [
    "Surface Torque",
    "Surface_Torque",
    "SurfaceTorque",
    "Torque",
    "Torque (Surface)",
]

# Columnas candidatas para Shocks & Vibs (WOB-RPM heatmap)
SHOCKS_VIB_COL_CANDIDATES = [
    "Shocks",
    "Vibs",
    "Shock",
    "Vibration",
    "Vibrations",
    "Shock Count",
    "Vibration Severity",
    "Shock & Vib",
    "ShockAndVib",
]


# =========================
# API helpers
# =========================
def api_get(
    path: str,
    params: dict | None = None,
    base_url: str | None = None,
    token: str | None = None,
):
    if not token:
        raise RuntimeError("Falta token de acceso a la API (SOLO_ACCESS_TOKEN).")
    if path.startswith("http://") or path.startswith("https://"):
        url = path
    else:
        base = (base_url or API_DEFAULT_BASE_URL).rstrip("/")
        url = f"{base}{path}"
    headers = {
        "Authorization": f"Bearer {token}",
        "Accept": "application/json",
    }
    r = requests.get(url, headers=headers, params=params, timeout=30)
    if r.status_code >= 400:
        raise RuntimeError(f"GET {url} -> {r.status_code}\n{r.text}")
    try:
        return r.json()
    except Exception:
        content_type = r.headers.get("content-type", "")
        body_preview = (r.text or "").strip()
        if len(body_preview) > 500:
            body_preview = body_preview[:500] + "..."
        raise RuntimeError(
            f"Respuesta no JSON desde {url} (content-type: {content_type}).\n{body_preview}"
        )


def get_trace_time_with_fallback(path, params, base_url, token):
    """Intenta /data/calculated/time y hace fallback a /data/time cuando no existe mapping."""
    res = api_get(path, params, base_url, token)
    if isinstance(res, dict) and res.get("type") == "NOT_FOUND" and "calculated trace" in str(res.get("message", "")):
        return api_get(path.replace("/data/calculated/time", "/data/time"), params, base_url, token)
    return res

def get_trace_depth_with_fallback(path, params, base_url, token):
    """Intenta /data/calculated/depth y hace fallback a /data/depth cuando no existe mapping."""
    res = api_get(path, params, base_url, token)
    if isinstance(res, dict) and res.get("type") == "NOT_FOUND" and "calculated trace" in str(res.get("message", "")):
        return api_get(path.replace("/data/calculated/depth", "/data/depth"), params, base_url, token)
    return res

def build_api_url(path: str, base_url: str | None = None) -> str:
    if path.startswith("http://") or path.startswith("https://"):
        return path
    base = (base_url or API_DEFAULT_BASE_URL).rstrip("/")
    return f"{base}{path}"


def api_get_raw(
    path: str,
    params: dict | None = None,
    base_url: str | None = None,
    token: str | None = None,
):
    if not token:
        raise RuntimeError("Falta token de acceso a la API (SOLO_ACCESS_TOKEN).")
    url = build_api_url(path, base_url)
    headers = {
        "Authorization": f"Bearer {token}",
        "Accept": "application/json",
    }
    return requests.get(url, headers=headers, params=params, timeout=30)


def normalize_list_response(resp) -> list:
    if isinstance(resp, dict):
        for key in ("content", "items", "data", "results"):
            val = resp.get(key)
            if isinstance(val, list):
                return val
    if isinstance(resp, list):
        return resp
    return []


def parse_params_input(text: str) -> dict | None:
    raw = (text or "").strip()
    if not raw:
        return None
    if raw.startswith("{") and raw.endswith("}"):
        try:
            data = json.loads(raw)
            return data if isinstance(data, dict) else None
        except Exception:
            return None
    params: dict[str, str] = {}
    for part in raw.split("&"):
        if not part:
            continue
        if "=" in part:
            k, v = part.split("=", 1)
            params[k.strip()] = v.strip()
        else:
            params[part.strip()] = ""
    return params or None


@st.cache_data(show_spinner=False)
def api_list_projects(base_url: str, token: str):
    return api_get("/public/api/v1/projects", {"offset": 0, "limit": 200}, base_url, token)


@st.cache_data(show_spinner=False)
def api_list_wells(base_url: str, token: str, project_uuid: str):
    return api_get(
        f"/public/api/v1/projects/{project_uuid}/wells",
        {"offset": 0, "limit": 200},
        base_url,
        token,
    )


@st.cache_data(show_spinner=False)
def api_list_laterals(
    base_url: str,
    token: str,
    project_uuid: str | None = None,
    well_uuid: str | None = None,
    custom_path: str | None = None,
    extra_params: dict | None = None,
):
    base_params = {"offset": 0, "limit": 200}
    if extra_params:
        base_params.update(extra_params)
    candidates: list[tuple[str, dict]] = []
    if custom_path:
        path = custom_path
        if well_uuid:
            path = path.replace("{well_uuid}", well_uuid)
        if project_uuid:
            path = path.replace("{project_uuid}", project_uuid)
        candidates.append((path, dict(base_params)))
    if well_uuid:
        candidates.extend(
            [
                (f"/public/api/v1/wells/{well_uuid}/laterals", dict(base_params)),
                (f"/api/v1/wells/{well_uuid}/laterals", dict(base_params)),
            ]
        )
    if project_uuid:
        candidates.extend(
            [
                (
                    "/public/api/v1/laterals",
                    {"projectUuid": project_uuid, **dict(base_params)},
                ),
                ("/api/v1/laterals", {"projectUuid": project_uuid, **dict(base_params)}),
            ]
        )

    errors: list[str] = []
    for path, params in candidates:
        try:
            return api_get_raw(path, params, base_url, token)
        except RuntimeError as e:
            errors.append(str(e))

    raise RuntimeError("No pude listar laterales. Errores:\n" + "\n".join(errors))


@st.cache_data(show_spinner=False)
def api_list_trace_definitions(base_url: str, token: str, custom_path: str | None = None):
    candidates: list[tuple[str, dict]] = []
    if custom_path:
        candidates.append((custom_path, {"offset": 0, "limit": 500}))
    candidates.extend(
        [
            ("/public/api/v1/traces", {"offset": 0, "limit": 500}),
            ("/api/v1/traces", {"offset": 0, "limit": 500}),
        ]
    )
    errors: list[str] = []
    for path, params in candidates:
        try:
            return api_get_raw(path, params, base_url, token)
        except RuntimeError as e:
            errors.append(str(e))
    raise RuntimeError("No pude listar tipos de traza. Errores:\n" + "\n".join(errors))


@st.cache_data(show_spinner=False)
def api_get_slide_sheet_intervals(
    base_url: str,
    token: str,
    well_uuid: str,
    depth_from: float | None = None,
    depth_to: float | None = None,
    custom_path: str | None = None,
):
    params = {}
    if depth_from is not None:
        params["from"] = depth_from
    if depth_to is not None:
        params["to"] = depth_to
    path = custom_path or f"/public/api/v1/wells/{well_uuid}/intervals/slide-sheet"
    path = path.replace("{well_uuid}", well_uuid)
    return api_get(path, params or None, base_url, token)


def slide_sheet_to_df(resp) -> pd.DataFrame:
    rows = normalize_list_response(resp)
    if not rows:
        return pd.DataFrame()
    df = pd.DataFrame(rows)
    rename_map = {
        "mode": "Mode",
        "depth_from": "Start Depth",
        "depth_to": "End Depth",
        "survey_md": "Survey MD",
        "incl": "Inclination",
        "azimuth": "Azimuth",
        "distance": "Distance",
        "rop_avg": "ROP",
        "dls": "DLS",
        "wob_avg": "WOB",
        "rpm_avg": "RPM",
        "from": "Start",
        "to": "End",
    }
    df = df.rename(columns=rename_map)
    for col in REQUIRED_COLUMNS:
        if col not in df.columns:
            df[col] = np.nan
    return df[REQUIRED_COLUMNS]

@st.cache_data(show_spinner=False)
def api_list_drilling_traces(
    base_url: str,
    token: str,
    well_uuid: str,
    custom_path: str | None = None,
    extra_params: dict | None = None,
):
    base_params = {"offset": 0, "limit": 200}
    if extra_params:
        base_params.update(extra_params)
    candidates: list[tuple[str, dict]] = []
    if custom_path:
        candidates.append(
            (
                custom_path.replace("{well_uuid}", well_uuid),
                dict(base_params),
            )
        )
    candidates.extend(
        [
            (f"/public/api/v1/wells/{well_uuid}/drilling-traces", dict(base_params)),
            (f"/api/v1/wells/{well_uuid}/drilling-traces", dict(base_params)),
            (f"/public/api/v1/wells/{well_uuid}/traces", dict(base_params)),
            (f"/api/v1/wells/{well_uuid}/traces", dict(base_params)),
            (
                "/public/api/v1/drilling-traces",
                {"wellUuid": well_uuid, **dict(base_params)},
            ),
            ("/api/v1/drilling-traces", {"wellUuid": well_uuid, **dict(base_params)}),
        ]
    )

    errors: list[str] = []
    for path, params in candidates:
        try:
            return get_trace_time_with_fallback(path, params, base_url, token)
        except RuntimeError as e:
            errors.append(str(e))

    raise RuntimeError("No pude listar trazas. Errores:\n" + "\n".join(errors))


@st.cache_data(show_spinner=False)
def api_list_mapped_traces(
    base_url: str,
    token: str,
    scope_uuid: str,
    trace_type: str,
    scope_kind: str = "lateral",
    custom_path: str | None = None,
    extra_params: dict | None = None,
):
    base_params = {"offset": 0, "limit": 200, "type": trace_type}
    type_only_params = {"type": trace_type}
    if extra_params:
        base_params.update(extra_params)
        type_only_params.update(extra_params)
    candidates: list[tuple[str, dict]] = []
    if custom_path:
        path = (
            custom_path.replace("{lateral_uuid}", scope_uuid).replace("{well_uuid}", scope_uuid)
        )
        candidates.append((path, dict(base_params)))
    if scope_kind == "well":
        candidates.extend(
            [
                (f"/public/api/v1/wells/{scope_uuid}/traces/mapped", dict(type_only_params)),
                (f"/public/api/v1/wells/{scope_uuid}/traces/mapped", dict(base_params)),
                (f"/api/v1/wells/{scope_uuid}/traces/mapped", dict(type_only_params)),
                (f"/api/v1/wells/{scope_uuid}/traces/mapped", dict(base_params)),
            ]
        )
    elif trace_type.upper() == "DEPTH":
        candidates.extend(
            [
                (f"/public/api/v1/wells/{scope_uuid}/traces/mapped", dict(type_only_params)),
                (f"/public/api/v1/wells/{scope_uuid}/traces/mapped", dict(base_params)),
                (
                    f"/public/api/v1/laterals/{scope_uuid}/mapped-depth-traces",
                    dict(base_params),
                ),
                (
                    f"/public/api/v1/laterals/{scope_uuid}/depth-traces",
                    dict(base_params),
                ),
                (f"/api/v1/wells/{scope_uuid}/traces/mapped", dict(type_only_params)),
                (f"/api/v1/wells/{scope_uuid}/traces/mapped", dict(base_params)),
            ]
        )
    else:
        candidates.extend(
            [
                (f"/public/api/v1/wells/{scope_uuid}/traces/mapped", dict(type_only_params)),
                (f"/public/api/v1/wells/{scope_uuid}/traces/mapped", dict(base_params)),
                (
                    f"/public/api/v1/laterals/{scope_uuid}/mapped-time-traces",
                    dict(base_params),
                ),
                (
                    f"/public/api/v1/laterals/{scope_uuid}/time-traces",
                    dict(base_params),
                ),
                (f"/api/v1/wells/{scope_uuid}/traces/mapped", dict(type_only_params)),
                (f"/api/v1/wells/{scope_uuid}/traces/mapped", dict(base_params)),
            ]
        )
    errors: list[str] = []
    for path, params in candidates:
        try:
            return get_trace_time_with_fallback(path, params, base_url, token)
        except RuntimeError as e:
            errors.append(str(e))
    raise RuntimeError("No pude listar trazas mapeadas. Errores:\n" + "\n".join(errors))


@st.cache_data(show_spinner=False)
def api_get_drilling_trace(
    base_url: str,
    token: str,
    trace_uuid: str,
    custom_path: str | None = None,
    params: dict | None = None,
):
    """
    Obtiene datos de una drilling-trace.

    Nota: dependiendo del tenant, los datos pueden estar en:
      - /drilling-traces/{uuid}/data
      - /drilling-traces/{uuid}/data/time
      - /drilling-traces/{uuid}/data/depth
      - /drilling-traces/{uuid}
    Por eso probamos varios candidatos.
    """
    candidates: list[tuple[str, dict | None]] = []
    if custom_path:
        candidates.append((custom_path.replace("{trace_uuid}", trace_uuid), params))
    # Endpoints más comunes (public primero)
    candidates.extend(
        [
            (f"/public/api/v1/drilling-traces/{trace_uuid}/data", params),
            (f"/public/api/v1/drilling-traces/{trace_uuid}/data/time", params),
            (f"/public/api/v1/drilling-traces/{trace_uuid}/data/depth", params),
            (f"/public/api/v1/drilling-traces/{trace_uuid}", params),
            (f"/api/v1/drilling-traces/{trace_uuid}/data", params),
            (f"/api/v1/drilling-traces/{trace_uuid}/data/time", params),
            (f"/api/v1/drilling-traces/{trace_uuid}/data/depth", params),
            (f"/api/v1/drilling-traces/{trace_uuid}", params),
        ]
    )

    errors: list[str] = []
    for path, p in candidates:
        try:
            return api_get(path, p, base_url, token)
        except RuntimeError as e:
            errors.append(str(e))

    raise RuntimeError("No pude obtener detalle de drilling-trace. Errores:\n" + "\n".join(errors))




def api_get_mapped_trace(
    base_url: str,
    token: str,
    scope_uuid: str,
    trace_uuid: str,
    trace_type: str,
    scope_kind: str = "lateral",
    custom_path: str | None = None,
    extra_params: dict | None = None,
    force_data_endpoint: bool = False,
    well_uuid: str | None = None,
):
    base_params = {"traceUuid": trace_uuid}
    if extra_params:
        base_params.update(extra_params)
    candidates: list[tuple[str, dict]] = []
    data_suffix = "time" if trace_type.upper() == "TIME" else "depth"
    data_scope_uuid = well_uuid or scope_uuid
    if force_data_endpoint:
        candidates.extend(
            [
                (
                    f"/public/api/v1/wells/{data_scope_uuid}/traces/{trace_uuid}/data/{data_suffix}",
                    dict(extra_params or {}),
                ),
                (
                    f"/api/v1/wells/{data_scope_uuid}/traces/{trace_uuid}/data/{data_suffix}",
                    dict(extra_params or {}),
                ),
            ]
        )
    else:
        if custom_path:
            path = (
                custom_path.replace("{lateral_uuid}", scope_uuid)
                .replace("{well_uuid}", scope_uuid)
                .replace("{trace_uuid}", trace_uuid)
            )
            candidates.append((path, dict(base_params)))
        if scope_kind == "well":
            candidates.extend(
                [
                    (
                        f"/public/api/v1/wells/{scope_uuid}/traces/{trace_uuid}/data/{data_suffix}",
                        dict(extra_params or {}),
                    ),
                    (
                        f"/api/v1/wells/{scope_uuid}/traces/{trace_uuid}/data/{data_suffix}",
                        dict(extra_params or {}),
                    ),
                ]
            )
    if not force_data_endpoint:
        well_params = {"trace_id": trace_uuid, "type": trace_type}
        if extra_params:
            well_params.update(extra_params)

        candidates.extend(
            [
                (f"/public/api/v1/wells/{scope_uuid}/traces/mapped", dict(well_params)),
                (f"/api/v1/wells/{scope_uuid}/traces/mapped", dict(well_params)),
            ]
        )
        if scope_kind != "well":
            if trace_type.upper() == "DEPTH":
                candidates.extend(
                    [
                        (
                            f"/public/api/v1/laterals/{scope_uuid}/mapped-depth-traces",
                            dict(base_params),
                        ),
                        (
                            f"/public/api/v1/laterals/{scope_uuid}/depth-traces",
                            dict(base_params),
                        ),
                        (f"/api/v1/laterals/{scope_uuid}/mapped-depth-traces", dict(base_params)),
                    ]
                )
            else:
                candidates.extend(
                    [
                        (
                            f"/public/api/v1/laterals/{scope_uuid}/mapped-time-traces",
                            dict(base_params),
                        ),
                        (
                            f"/public/api/v1/laterals/{scope_uuid}/time-traces",
                            dict(base_params),
                        ),
                        (f"/api/v1/laterals/{scope_uuid}/mapped-time-traces", dict(base_params)),
                    ]
                )
    errors: list[str] = []
    for path, params in candidates:
        try:
            return get_trace_time_with_fallback(path, params, base_url, token)
        except RuntimeError as e:
            errors.append(str(e))
    raise RuntimeError("No pude obtener datos de traza mapeada. Errores:\n" + "\n".join(errors))

@st.cache_data(show_spinner=False)
def api_list_traces_catalog(base_url: str, token: str, custom_path: str | None = None, params: dict | None = None):
    """
    Lista el catálogo global de trazas (predefinidas) según la doc de SOLO.
    Endpoint típico: GET /public/api/v1/traces
    """
    path = custom_path or "/public/api/v1/traces"
    return get_trace_time_with_fallback(path, params, base_url, token)


@st.cache_data(show_spinner=False)
def api_get_well_trace_data(
    base_url: str,
    token: str,
    well_uuid: str,
    trace_uuid: str,
    trace_type: str,
    params: dict | None = None,
    calculated: bool = False,
):
    """
    Obtiene series de datos de una traza por pozo:
      GET /public/api/v1/wells/{well_uuid}/traces/{trace_uuid}/data/{time|depth}
      o (si calculated=True):
      GET /public/api/v1/wells/{well_uuid}/traces/{trace_uuid}/data/calculated/{time|depth}

    Nota: este endpoint suele estar paginado (content + page/size). Si no mandas size,
    muchas veces SOLO devuelve ~10 filas por default.
    """
    suffix = "time" if str(trace_type).upper() == "TIME" else "depth"
    if calculated:
        path = f"/public/api/v1/wells/{well_uuid}/traces/{trace_uuid}/data/calculated/{suffix}"
    else:
        path = f"/public/api/v1/wells/{well_uuid}/traces/{trace_uuid}/data/{suffix}"

    p = dict(params or {})
    # Defaults de paginación (no pisan valores del usuario)
    p.setdefault("page", 0)
    p.setdefault("size", 5000)

    # Compatibilidad: algunas instalaciones usan offset/limit en vez de page/size
    try:
        _sz = int(p.get("size", 0))
        _pg = int(p.get("page", 0))
        p.setdefault("limit", _sz)
        p.setdefault("offset", _pg * _sz)
    except Exception:
        pass

    return api_get(path, p, base_url, token)


def probe_well_trace_data(
    base_url: str,
    token: str,
    well_uuid: str,
    trace_uuid: str,
    prefer_type: str,
    user_params: dict | None = None,
) -> tuple[pd.DataFrame, str, dict]:
    """
    Estrategia robusta basada en endpoints oficiales:
      1) Intento directo con preferencia (TIME/DEPTH) y params del usuario
      2) Si TIME vacío: ventanas hacia atrás (6h, 24h, 7d, 30d, 180d, 365d, 3a)
      3) Si sigue vacío: DEPTH con rangos crecientes
    """
    from datetime import datetime, timedelta, timezone

    prefer = str(prefer_type).upper()
    user_params = dict(user_params or {})

    def _get_df(ttype: str, params: dict) -> pd.DataFrame:
        # Intento 1: data/{time|depth}
        try:
            detail = api_get_well_trace_data(
                base_url=base_url,
                token=token,
                well_uuid=well_uuid,
                trace_uuid=str(trace_uuid),
                trace_type=ttype,
                params=params,
                calculated=False,
            )
            df = trace_detail_to_df(detail)
            if df is not None and not df.empty:
                return df
        except Exception:
            pass

        # Intento 2 (fallback): data/calculated/{time|depth}
        detail2 = api_get_well_trace_data(
            base_url=base_url,
            token=token,
            well_uuid=well_uuid,
            trace_uuid=str(trace_uuid),
            trace_type=ttype,
            params=params,
            calculated=True,
        )
        return trace_detail_to_df(detail2)

    # 1) intento directo
    try:
        df0 = _get_df(prefer, user_params)
        if not df0.empty:
            return df0, prefer, user_params
    except Exception:
        pass

    # FIX: If user explicitly specified a range (from/to), do NOT override it with probing.
    # In that case, return empty so the UI can reflect "no data for that range".
    _has_user_range = ("from" in user_params and "to" in user_params and str(user_params.get("from")).strip() != "" and str(user_params.get("to")).strip() != "")
    if _has_user_range:
        return pd.DataFrame(), prefer, user_params

    # 2) probing TIME
    if prefer == "TIME":
        now = datetime.now(timezone.utc)
        windows = [
            timedelta(hours=6),
            timedelta(hours=24),
            timedelta(days=7),
            timedelta(days=30),
            timedelta(days=180),
            timedelta(days=365),
            timedelta(days=365 * 3),
        ]
        for w in windows:
            p = dict(user_params)
            p["to"] = now.strftime("%Y-%m-%dT%H:%M:%SZ")
            p["from"] = (now - w).strftime("%Y-%m-%dT%H:%M:%SZ")
            try:
                df = _get_df("TIME", p)
                if not df.empty:
                    return df, "TIME", p
            except Exception:
                continue

    # 3) probing DEPTH
    depth_ranges = [(0, 2000), (0, 5000), (0, 10000), (0, 20000), (0, 50000)]
    for a, b in depth_ranges:
        p = dict(user_params)
        p["from"] = a
        p["to"] = b
        try:
            df = _get_df("DEPTH", p)
            if not df.empty:
                return df, "DEPTH", p
        except Exception:
            continue

    return pd.DataFrame(), prefer, user_params



def trace_detail_to_df(detail) -> pd.DataFrame:
    if detail is None:
        return pd.DataFrame()
    if isinstance(detail, list):
        if detail and isinstance(detail[0], dict):
            return pd.DataFrame(detail)
        return pd.DataFrame(detail)
    if not isinstance(detail, dict):
        return pd.DataFrame()

    for key in ("content", "points", "samples", "data", "rows", "values"):
        val = detail.get(key)
        if isinstance(val, list):
            if val and isinstance(val[0], dict):
                return pd.DataFrame(val)
            columns = detail.get("columns") or detail.get("fields")
            if isinstance(columns, list):
                return pd.DataFrame(val, columns=columns)
            return pd.DataFrame(val)

    data = detail.get("data")
    if isinstance(data, dict):
        rows = data.get("rows") or data.get("values")
        cols = data.get("columns") or data.get("fields")
        if isinstance(rows, list):
            if rows and isinstance(rows[0], dict):
                return pd.DataFrame(rows)
            if isinstance(cols, list):
                return pd.DataFrame(rows, columns=cols)
            return pd.DataFrame(rows)

    return pd.DataFrame()



def probe_mapped_trace_data(
    base_url: str,
    token: str,
    well_uuid: str,
    mapped_scope_uuid: str,
    trace_uuid: str,
    mapped_scope_kind: str,
    mapped_trace_path: str | None,
    force_data_endpoint: bool,
    prefer_type: str,
    user_params: dict | None = None,
    max_rows_preview: int = 5,
) -> tuple[pd.DataFrame, str, dict]:
    """
    Intenta obtener datos para una traza mapeada con una estrategia robusta:
      1) Preferencia (TIME o DEPTH) con params del usuario (si hay)
      2) Si TIME está vacío, prueba ventanas relativas hacia atrás (6h, 24h, 7d, 30d, 180d, 365d)
      3) Si sigue vacío, prueba DEPTH con rangos crecientes

    Retorna: (df, tipo_usado, params_usados)
    """
    prefer = str(prefer_type).upper()
    user_params = dict(user_params or {})

    def _get_df(ttype: str, params: dict) -> pd.DataFrame:
        detail = api_get_mapped_trace(
            base_url,
            token,
            mapped_scope_uuid,
            str(trace_uuid),
            ttype,
            mapped_scope_kind,
            mapped_trace_path or None,
            params,
            force_data_endpoint,
            well_uuid,
        )
        return trace_detail_to_df(detail)

    # 1) intento directo con lo que pidió el usuario
    try:
        df0 = _get_df(prefer, user_params)
        if not df0.empty:
            return df0, prefer, user_params
    except Exception:
        pass

    # 2) probing TIME con ventanas hacia atrás (si la preferencia o el modo actual es TIME)
    if prefer == "TIME":
        now = datetime.now(timezone.utc)
        windows = [
            timedelta(hours=6),
            timedelta(hours=24),
            timedelta(days=7),
            timedelta(days=30),
            timedelta(days=180),
            timedelta(days=365),
            timedelta(days=365 * 3),
        ]
        for w in windows:
            p = dict(user_params)
            p["to"] = now.strftime("%Y-%m-%dT%H:%M:%SZ")
            p["from"] = (now - w).strftime("%Y-%m-%dT%H:%M:%SZ")
            try:
                df = _get_df("TIME", p)
                if not df.empty:
                    return df, "TIME", p
            except Exception:
                continue

    # 3) probing DEPTH con rangos crecientes
    depth_ranges = [(0, 2000), (0, 5000), (0, 10000), (0, 20000), (0, 50000)]
    for a, b in depth_ranges:
        p = dict(user_params)
        p["from"] = a
        p["to"] = b
        try:
            df = _get_df("DEPTH", p)
            if not df.empty:
                return df, "DEPTH", p
        except Exception:
            continue

    return pd.DataFrame(), prefer, user_params



def merge_trace_frames(frames: list[tuple[str, pd.DataFrame]]) -> pd.DataFrame:
    merged = pd.DataFrame()
    x_col_name = None
    x_candidates = {
        "depth",
        "md",
        "measured_depth",
        "survey_md",
        "time",
        "timestamp",
        "datetime",
        "date",
    }

    for label, df in frames:
        if df.empty:
            continue
        df = df.copy()
        cols = list(df.columns)
        if not cols:
            continue

        x_col = None
        for c in cols:
            if str(c).strip().lower() in x_candidates:
                x_col = c
                break
        if x_col is None:
            x_col = cols[0]

        value_cols = [c for c in cols if c != x_col]
        if not value_cols:
            df = df.reset_index().rename(columns={"index": "index"})
            x_col = "index"
            value_cols = [c for c in df.columns if c != x_col]
        if not value_cols:
            continue

        if len(value_cols) == 1:
            df = df[[x_col, value_cols[0]]].rename(columns={value_cols[0]: label})
        else:
            rename_map = {c: f"{label}_{c}" for c in value_cols}
            df = df[[x_col] + value_cols].rename(columns=rename_map)

        if merged.empty:
            merged = df
            x_col_name = x_col
        else:
            if x_col_name and x_col != x_col_name:
                df = df.rename(columns={x_col: x_col_name})
            merged = pd.merge(merged, df, on=x_col_name, how="outer")

    if x_col_name and x_col_name in merged.columns:
        merged = merged.sort_values(by=x_col_name, ignore_index=True)
    return merged


def pick_default_column(columns: Iterable[str], candidates: List[str]) -> str | None:
    lowered = {c.lower(): c for c in columns}
    for cand in candidates:
        key = cand.lower()
        if key in lowered:
            return lowered[key]
    return None


# =========================
# UI & Estilos
# =========================
st.set_page_config(page_title=APP_TITLE, layout="wide")
px.defaults.template = PLOTLY_TEMPLATE
px.defaults.color_discrete_sequence = COLOR_SEQ

st.markdown(
    """
<style>
div[data-testid="stHorizontalBlock"]{
    gap: 2rem !important;
    margin-top: -8px !important;
    padding-top: 0 !important;
}
div[data-testid="stPlotlyChart"]{
    margin-top: 0px !important;
    margin-bottom: 26px !important;
}
section.main > div{
    padding-top: 0.5rem !important;
}
</style>
""",
    unsafe_allow_html=True,
)

def vspace(px: int = 16) -> None:
    st.markdown(f"<div style='height:{px}px'></div>", unsafe_allow_html=True)


def init_session_state() -> None:
    st.session_state.setdefault("report_ready", False)
    st.session_state.setdefault("captures_done", False)
    st.session_state.setdefault("pptx_path", None)
    st.session_state.setdefault("tmp_dir", None)
    st.session_state.setdefault("pdf_path", None)
    st.session_state.setdefault("ui_lang", "es")
    st.session_state.setdefault("alert_pdf_path", None)
    st.session_state.setdefault("alert_pptx_path", None)
    st.session_state.setdefault("alert_captures", [])
    st.session_state.setdefault("alert_last_capture", None)


# --- Internacionalización (Español por defecto + English + Русский) ---
I18N: dict[str, dict[str, str]] = {
    "es": {
        "intro_p1": "Sube tu **Slide Sheet Export Excel**, define corridas y genera el PPTX completo.",
        "intro_p2": "En la app verás **tablas + gráficas Plotly (BI Pro)**.",
        "intro_p3": "Opcional: **capturas al final (Windows/pyautogui)** con botón de inicio.",
        "tab_kpi": "KPI Report",
        "tab_bha": "Ingeniería BHA",
        "tab_roadmap": "Insights Ingeniería",
        "tab_trip": "Tripping Analysis",
        "tab_mud": "Mud Report",
        "data_source": "Fuente de datos",
        "src_excel": "Excel",
        "src_api": "API",
        "src_csv": "CSV",
        "upload_excel": "Upload Excel (.xlsx)",
        "sidebar_options": "Opciones",
        "show_plots": "Mostrar gráficos en la app",
        "region_captures": "Region Captures (Windows)",
        "enable_region_captures": "Enable region captures",
        "n_captures": "Number of captures",
        "interval_captures": "Interval between captures (s)",
        "region_x": "Region X",
        "region_y": "Region Y",
        "region_w": "Region Width",
        "region_h": "Region Height",
        "capture_prefix": "Capture title prefix",
        "runs_header": "Runs / Corridas",
        "generate_pptx": "Generate PPTX",
        "logo_missing": "No encontré el logo en:",
        "bha_subheader": "🔧 Módulo de Ingeniería: BHA Resonance & Operational Window",
        "bha_caption": "Este módulo es independiente del reporte KPI. Sube aquí el CSV con trazas de **RPM, Torque y WOB**.",
        "roadmap_subheader": "📊 Insights de Ingeniería (Enterprise)",
        "roadmap_caption": "Heatmap y gráficas avanzadas basadas en el CSV de ingeniería.",
        "roadmap_need_bha": "Primero carga el CSV en la pestaña **Ingeniería BHA**.",
        "roadmap_rt": "**Actualización en tiempo real**",
        "roadmap_auto": "Actualizar automáticamente cada 30 s",
        "roadmap_auto_help": "Vuelve a cargar esta vista cada X segundos (usa los datos ya cargados en BHA).",
        "interval_seconds": "Intervalo (segundos)",
        "interval_help_roadmap": "Cada cuántos segundos se actualiza la vista (10–300 s).",
        "mse_bit_diam": "Diámetro de barrena para MSE (in)",
        "kpis_exec": "**KPIs Ejecutivos**",
        "metric_valid_points": "Puntos válidos",
        "metric_resonant": "Bandas resonantes",
        "metric_safe": "Ventanas seguras",
        "metric_tol_hz": "Tolerancia (Hz)",
        "heatmap_eng_title": "🔥 Heatmap de correlación – Ingeniería",
        "trip_caption": "Analiza la traza de Hookload durante un periodo de viaje (Trip In / Trip Out) y calcula envolventes por profundidad y eventos de sobre-tensión (overpull).",
        "trip_mode_label": "Tipo de viaje",
        "trip_out": "Trip Out",
        "trip_in": "Trip In",
        "trip_dir_filter": "Filtrar por dirección del viaje (derivada de profundidad)",
        "trip_dir_help": "Trip Out = solo puntos con profundidad bajando; Trip In = solo subiendo. Si no hay datos, desmarca para usar todo el rango.",
        "trip_bin": "Bin de profundidad (m)",
        "trip_baseline": "Baseline (percentil)",
        "trip_thr": "Umbral Overpull (mismas unidades de Hookload)",
        "trip_rolling": "Ventana rolling (muestras) para vista pro",
        "trip_envelope": "Envelope por profundidad",
        "trip_env_exact": "Profundidad exacta (max Hookload por depth)",
        "trip_env_bin": "Bin (rangos de profundidad)",
        "trip_env_help": "Exacta: una fila por cada Bit depth (max Hookload). Bin: agrupar por rangos para baseline y overpull.",
        "trip_range_hdr": "#### Rango de análisis (opcional)",
        "trip_range_cap": "Restringe el análisis a un intervalo de tiempo y/o de profundidad. Si no defines rango, se usa todo el dato cargado.",
        "trip_use_time": "Aplicar rango de tiempo",
        "trip_from": "Desde (fecha/hora)",
        "mud_caption": "Carga reportes de lodo en PDF, Excel o CSV (subiendo archivos o desde correo). Se genera una bitácora unificada por día.",
        "mud_src_files": "Subir archivos",
        "mud_src_email": "Correo electrónico",
        "mud_chip_upload": "📁 Subir archivos",
        "mud_chip_email": "Correo electrónico",
        "api_config_slide": "Configuración API – Slide Sheet",
        "api_base_url": "Base URL API",
        "api_token": "Access Token",
        "api_slide_path": "Ruta slide sheet (opcional)",
        "api_slide_help": "Usa {well_uuid}.",
        "use_depth_range": "Usar rango de profundidad",
        "depth_from": "Depth from",
        "depth_to": "Depth to",
        "api_note_html": "Nota: evita /api/v1 si responde HTML; usa /public/api/v1.",
        "enter_token": "Ingresa un token válido para consultar la API.",
        "list_projects_err": "No pude listar proyectos:",
        "no_projects": "No hay proyectos disponibles para este token.",
        "unnamed": "Sin nombre",
        "project": "Proyecto",
        "project_no_uuid": "El proyecto seleccionado no tiene UUID.",
        "list_wells_err": "No pude listar pozos:",
        "no_wells": "No hay pozos disponibles en este proyecto.",
        "well": "Pozo",
        "well_no_uuid": "El pozo seleccionado no tiene UUID.",
        "kpi_rt": "**Actualización en tiempo real (API)**",
        "kpi_auto_refresh": "Actualizar automáticamente cada 30 s",
        "kpi_auto_help": "Vuelve a cargar el Slide Sheet desde la API cada X segundos.",
        "kpi_interval_help": "Cada cuántos segundos se vuelve a cargar (10–300 s).",
        "load_slide_sheet": "Cargar Slide Sheet",
        "kpi_active_source": "Fuente activa: Slide Sheet (API).",
        "show_rows": "Mostrar filas",
        "all_rows": "Todos",
        "next_refresh": "🔄 Próxima actualización en **{i}** s… (desmarca «Actualizar automáticamente» para detener)",
        "solo_expander": "Conexión SOLO (API)",
        "bha_api_cfg": "Configuración API",
        "bha_missing_token": "Falta SOLO_ACCESS_TOKEN. Configúralo en tu .env o pega el token aquí.",
        "solo_token_loaded": "Token cargado desde .env (no necesitas pegarlo).",
        "solo_change_token": "Cambiar token",
        "solo_token_help": "Configura SOLO_ACCESS_TOKEN en tu .env para no pegarlo aquí.",
        "kpi_api_no_intervals": "La API no devolvió intervalos para ese rango.",
        "kpi_slide_loaded": "Slide Sheet cargado: {n:,} intervalos.",
        "kpi_load_err": "No pude cargar slide sheet:",
        "trip_to": "Hasta (fecha/hora)",
        "trip_help_dt": "ISO 8601 o formato reconocible por pandas.",
        "trip_use_depth": "Aplicar rango de profundidad",
        "trip_depth_min_l": "Profundidad mínima (m)",
        "trip_depth_max_l": "Profundidad máxima (m)",
        "trip_depth_help_min": "Solo se analizan puntos con Bit depth ≥ este valor. La columna debe estar en metros.",
        "trip_depth_help_max": "Solo se analizan puntos con Bit depth ≤ este valor. Si no sale data, revisa el rango real del CSV (mín/máx de la columna Bit depth).",
        "trip_badge_hookload": "Hookload",
        "trip_badge_trip_io": "Trip In / Out",
        "trip_badge_overpull": "Overpull",
        "trip_gap_caption": "Si en el rango de fechas hay huecos sin datos, puedes rellenarlos con interpolación lineal.",
        "trip_interp_chk": "Rellenar huecos con interpolación lineal",
        "trip_interp_help": "Crea una malla regular en el tiempo (desde el primer hasta el último dato) e interpola Hookload y Bit depth en los puntos faltantes.",
        "trip_interp_interval_l": "Intervalo de interpolación (segundos)",
        "trip_interp_interval_help": "Cada cuántos segundos se genera un punto en la malla temporal. Valores bajos (1–5 s) dan más puntos.",
        "mud_chip_bitacora": "Bitácora",
        "mud_chip_formats": "PDF / Excel / CSV",
        "mud_chip_mail_short": "Correo",
        "mud_imap_expander": "Configuración de correo (IMAP)",
        "mud_imap_caption": "Usa variables de entorno **MUD_IMAP_SERVER**, **MUD_IMAP_USER** y **MUD_IMAP_PASS** en tu .env para no escribir la contraseña aquí.",
        "mud_imap_server": "Servidor IMAP",
        "mud_imap_server_help": "Ej: imap.gmail.com",
        "mud_imap_user": "Usuario (correo)",
        "mud_imap_pass": "Contraseña (App Password en Gmail)",
        "mud_imap_pass_help": "En Gmail usa una contraseña de aplicación, no la de la cuenta.",
        "mud_imap_filter": "Filtrar por nombre de archivo (opcional)",
        "mud_mark_read": "Marcar correos como leídos al descargar",
        "mud_auto_hdr": "**Revisión automática**",
        "mud_auto_chk": "Revisar correo automáticamente cada 60 s",
        "mud_auto_help": "Cada X segundos se consulta el correo y se actualiza la bitácora. Desmarca para detener.",
        "mud_interval_imap": "Intervalo (segundos)",
        "mud_interval_imap_help": "Cada cuántos segundos se revisa el correo (30–300 s).",
        "mud_fetch_btn": "🔥 Rogii – Revisar correo y cargar reportes",
        "mud_fetch_help": "Consulta IMAP y descarga adjuntos PDF/Excel/CSV de correos no leídos.",
        "mud_err_imap": "Completa servidor IMAP, usuario y contraseña (o configúralos en .env).",
        "mud_spinner_imap": "Conectando al correo y descargando adjuntos...",
        "mud_err_dl": "No se pudo conectar o descargar:",
        "mud_no_attach": "No se encontraron adjuntos PDF/Excel/CSV en correos no leídos (o no coinciden con el filtro).",
        "mud_success_attach": "Se descargaron **{n}** adjunto(s). Procesando...",
        "tab_trip_env": "Envelope / Overpull",
        "tab_trip_broom": "Broomstick (FF)",
        "mud_email_countdown": "🔥 **Rogii** – Próxima revisión de correo en **{i}** s… (desmarca «Revisar correo automáticamente» para detener)",
        "mud_upload_reports": "Subir reportes de lodo (PDF, Excel, CSV)",
    },
    "en": {
        "intro_p1": "Upload your **Slide Sheet Excel export**, define runs and generate the full PPTX.",
        "intro_p2": "In the app you will see **tables + Plotly charts (BI Pro)**.",
        "intro_p3": "Optional: **end-of-run captures (Windows/pyautogui)** with a start button.",
        "tab_kpi": "KPI Report",
        "tab_bha": "BHA Engineering",
        "tab_roadmap": "Engineering Insights",
        "tab_trip": "Tripping Analysis",
        "tab_mud": "Mud Report",
        "data_source": "Data source",
        "src_excel": "Excel",
        "src_api": "API",
        "src_csv": "CSV",
        "upload_excel": "Upload Excel (.xlsx)",
        "sidebar_options": "Options",
        "show_plots": "Show charts in the app",
        "region_captures": "Region captures (Windows)",
        "enable_region_captures": "Enable region captures",
        "n_captures": "Number of captures",
        "interval_captures": "Interval between captures (s)",
        "region_x": "Region X",
        "region_y": "Region Y",
        "region_w": "Region width",
        "region_h": "Region height",
        "capture_prefix": "Capture title prefix",
        "runs_header": "Runs",
        "generate_pptx": "Generate PPTX",
        "logo_missing": "Logo not found at:",
        "bha_subheader": "🔧 Engineering module: BHA resonance & operational window",
        "bha_caption": "This module is independent of the KPI report. Upload the CSV with **RPM, torque and WOB** traces.",
        "roadmap_subheader": "📊 Engineering insights (Enterprise)",
        "roadmap_caption": "Heatmap and advanced charts from the engineering CSV.",
        "roadmap_need_bha": "First load the CSV in the **BHA Engineering** tab.",
        "roadmap_rt": "**Live refresh**",
        "roadmap_auto": "Auto-refresh every 30 s",
        "roadmap_auto_help": "Reload this view every X seconds (uses data already loaded in BHA).",
        "interval_seconds": "Interval (seconds)",
        "interval_help_roadmap": "How often the view refreshes (10–300 s).",
        "mse_bit_diam": "Bit diameter for MSE (in)",
        "kpis_exec": "**Executive KPIs**",
        "metric_valid_points": "Valid points",
        "metric_resonant": "Resonant bands",
        "metric_safe": "Safe windows",
        "metric_tol_hz": "Tolerance (Hz)",
        "heatmap_eng_title": "🔥 Engineering correlation heatmap",
        "trip_caption": "Analyze Hookload during a trip (Trip In / Trip Out) and compute depth envelopes and overpull events.",
        "trip_mode_label": "Trip type",
        "trip_out": "Trip out",
        "trip_in": "Trip in",
        "trip_dir_filter": "Filter by trip direction (depth derivative)",
        "trip_dir_help": "Trip out = depth decreasing; Trip in = increasing. If data is missing, uncheck to use the full range.",
        "trip_bin": "Depth bin (m)",
        "trip_baseline": "Baseline (percentile)",
        "trip_thr": "Overpull threshold (same units as Hookload)",
        "trip_rolling": "Rolling window (samples) for pro view",
        "trip_envelope": "Depth envelope",
        "trip_env_exact": "Exact depth (max Hookload per depth)",
        "trip_env_bin": "Bin (depth ranges)",
        "trip_env_help": "Exact: one row per bit depth (max Hookload). Bin: group depth ranges for baseline and overpull.",
        "trip_range_hdr": "#### Analysis range (optional)",
        "trip_range_cap": "Restrict analysis to a time and/or depth interval. If empty, the full loaded dataset is used.",
        "trip_use_time": "Apply time range",
        "trip_from": "From (date/time)",
        "mud_caption": "Load mud reports from PDF, Excel or CSV (upload or email). A unified daily log is generated.",
        "mud_src_files": "Upload files",
        "mud_src_email": "Email",
        "mud_chip_upload": "📁 Upload files",
        "mud_chip_email": "Email",
        "api_config_slide": "API settings – Slide Sheet",
        "api_base_url": "API base URL",
        "api_token": "Access token",
        "api_slide_path": "Slide sheet path (optional)",
        "api_slide_help": "Use {well_uuid}.",
        "use_depth_range": "Use depth range",
        "depth_from": "Depth from",
        "depth_to": "Depth to",
        "api_note_html": "Note: avoid /api/v1 if it returns HTML; use /public/api/v1.",
        "enter_token": "Enter a valid token to query the API.",
        "list_projects_err": "Could not list projects:",
        "no_projects": "No projects available for this token.",
        "unnamed": "Unnamed",
        "project": "Project",
        "project_no_uuid": "The selected project has no UUID.",
        "list_wells_err": "Could not list wells:",
        "no_wells": "No wells available in this project.",
        "well": "Well",
        "well_no_uuid": "The selected well has no UUID.",
        "kpi_rt": "**Live refresh (API)**",
        "kpi_auto_refresh": "Auto-refresh every 30 s",
        "kpi_auto_help": "Reload the Slide Sheet from the API every X seconds.",
        "kpi_interval_help": "How often to reload (10–300 s).",
        "load_slide_sheet": "Load Slide Sheet",
        "kpi_active_source": "Active source: Slide Sheet (API).",
        "show_rows": "Show rows",
        "all_rows": "All",
        "next_refresh": "🔄 Next refresh in **{i}** s… (uncheck auto-refresh to stop)",
        "solo_expander": "SOLO connection (API)",
        "bha_api_cfg": "API settings",
        "bha_missing_token": "SOLO_ACCESS_TOKEN is missing. Set it in .env or paste the token here.",
        "solo_token_loaded": "Token loaded from .env (paste not required).",
        "solo_change_token": "Change token",
        "solo_token_help": "Set SOLO_ACCESS_TOKEN in .env to skip pasting here.",
        "kpi_api_no_intervals": "The API returned no intervals for that range.",
        "kpi_slide_loaded": "Slide Sheet loaded: {n:,} intervals.",
        "kpi_load_err": "Could not load Slide Sheet:",
        "trip_to": "To (date/time)",
        "trip_help_dt": "ISO 8601 or a format pandas can parse.",
        "trip_use_depth": "Apply depth range",
        "trip_depth_min_l": "Minimum depth (m)",
        "trip_depth_max_l": "Maximum depth (m)",
        "trip_depth_help_min": "Only points with Bit depth ≥ this value. Column must be in metres.",
        "trip_depth_help_max": "Only points with Bit depth ≤ this value. If empty, check the CSV depth range.",
        "trip_badge_hookload": "Hookload",
        "trip_badge_trip_io": "Trip In / Out",
        "trip_badge_overpull": "Overpull",
        "trip_gap_caption": "If the date range has gaps, you can fill them with linear interpolation.",
        "trip_interp_chk": "Fill gaps with linear interpolation",
        "trip_interp_help": "Builds a regular time grid and interpolates Hookload and Bit depth at missing times.",
        "trip_interp_interval_l": "Interpolation interval (seconds)",
        "trip_interp_interval_help": "Seconds between generated points. Lower (1–5 s) yields more points.",
        "mud_chip_bitacora": "Log",
        "mud_chip_formats": "PDF / Excel / CSV",
        "mud_chip_mail_short": "Email",
        "mud_imap_expander": "Email settings (IMAP)",
        "mud_imap_caption": "Use **MUD_IMAP_SERVER**, **MUD_IMAP_USER** and **MUD_IMAP_PASS** in .env to avoid typing the password here.",
        "mud_imap_server": "IMAP server",
        "mud_imap_server_help": "E.g. imap.gmail.com",
        "mud_imap_user": "User (email)",
        "mud_imap_pass": "Password (Gmail app password)",
        "mud_imap_pass_help": "In Gmail use an app password, not your account password.",
        "mud_imap_filter": "Filter by file name (optional)",
        "mud_mark_read": "Mark messages as read when downloaded",
        "mud_auto_hdr": "**Automatic polling**",
        "mud_auto_chk": "Check email automatically every 60 s",
        "mud_auto_help": "Every X seconds the inbox is checked and the log updates. Uncheck to stop.",
        "mud_interval_imap": "Interval (seconds)",
        "mud_interval_imap_help": "How often to poll email (30–300 s).",
        "mud_fetch_btn": "🔥 Rogii – Check email and load reports",
        "mud_fetch_help": "Query IMAP and download PDF/Excel/CSV attachments from unread mail.",
        "mud_err_imap": "Enter IMAP server, user and password (or set them in .env).",
        "mud_spinner_imap": "Connecting to mail and downloading attachments...",
        "mud_err_dl": "Could not connect or download:",
        "mud_no_attach": "No PDF/Excel/CSV attachments in unread mail (or filter did not match).",
        "mud_success_attach": "Downloaded **{n}** attachment(s). Processing...",
        "tab_trip_env": "Envelope / Overpull",
        "tab_trip_broom": "Broomstick (FF)",
        "mud_email_countdown": "🔥 **Rogii** – Next mail check in **{i}** s… (uncheck auto mail refresh to stop)",
        "mud_upload_reports": "Upload mud reports (PDF, Excel, CSV)",
    },
    "ru": {
        "intro_p1": "Загрузите **Excel экспорт Slide Sheet**, задайте рейсы и сформируйте полный PPTX.",
        "intro_p2": "В приложении — **таблицы и графики Plotly (BI Pro)**.",
        "intro_p3": "Опционально: **скриншоты в конце (Windows/pyautogui)** по кнопке.",
        "tab_kpi": "Отчёт KPI",
        "tab_bha": "Инженерия BHA",
        "tab_roadmap": "Инженерная аналитика",
        "tab_trip": "Анализ спуска/подъёма",
        "tab_mud": "Отчёт по раствору",
        "data_source": "Источник данных",
        "src_excel": "Excel",
        "src_api": "API",
        "src_csv": "CSV",
        "upload_excel": "Загрузить Excel (.xlsx)",
        "sidebar_options": "Параметры",
        "show_plots": "Показывать графики в приложении",
        "region_captures": "Захват области (Windows)",
        "enable_region_captures": "Включить захват области",
        "n_captures": "Число снимков",
        "interval_captures": "Интервал между снимками (с)",
        "region_x": "Область X",
        "region_y": "Область Y",
        "region_w": "Ширина области",
        "region_h": "Высота области",
        "capture_prefix": "Префикс заголовка снимка",
        "runs_header": "Рейсы",
        "generate_pptx": "Создать PPTX",
        "logo_missing": "Логотип не найден:",
        "bha_subheader": "🔧 Модуль инженерии: резонанс BHA и рабочее окно",
        "bha_caption": "Модуль независим от KPI. Загрузите CSV с трассами **RPM, момент и WOB**.",
        "roadmap_subheader": "📊 Инженерные инсайты (Enterprise)",
        "roadmap_caption": "Теплокарта и расширенные графики по CSV инженерии.",
        "roadmap_need_bha": "Сначала загрузите CSV на вкладке **Инженерия BHA**.",
        "roadmap_rt": "**Обновление в реальном времени**",
        "roadmap_auto": "Автообновление каждые 30 с",
        "roadmap_auto_help": "Перезагружать вид каждые X с (данные из BHA).",
        "interval_seconds": "Интервал (секунды)",
        "interval_help_roadmap": "Период обновления вида (10–300 с).",
        "mse_bit_diam": "Диаметр долота для MSE (дюйм)",
        "kpis_exec": "**Ключевые KPI**",
        "metric_valid_points": "Допустимые точки",
        "metric_resonant": "Резонансные полосы",
        "metric_safe": "Безопасные окна",
        "metric_tol_hz": "Допуск (Гц)",
        "heatmap_eng_title": "🔥 Теплокарта корреляций — инженерия",
        "trip_caption": "Анализ Hookload за период рейса (Trip In / Trip Out), огибающие по глубине и перетяжки (overpull).",
        "trip_mode_label": "Тип рейса",
        "trip_out": "Trip Out",
        "trip_in": "Trip In",
        "trip_dir_filter": "Фильтр по направлению рейса (производная глубины)",
        "trip_dir_help": "Trip Out — глубина снижается; Trip In — растёт. Нет данных — снимите флажок для всего диапазона.",
        "trip_bin": "Бин глубины (м)",
        "trip_baseline": "Базовый уровень (перцентиль)",
        "trip_thr": "Порог overpull (те же единицы, что Hookload)",
        "trip_rolling": "Окно скользящего среднего (выборки), pro-вид",
        "trip_envelope": "Огибающая по глубине",
        "trip_env_exact": "Точная глубина (max Hookload на глубину)",
        "trip_env_bin": "Бин (диапазоны глубины)",
        "trip_env_help": "Точная: строка на глубину долота (max Hookload). Бин: группировка для базы и overpull.",
        "trip_range_hdr": "#### Диапазон анализа (опционально)",
        "trip_range_cap": "Ограничьте анализ интервалом времени и/или глубины. Иначе — весь загруженный набор.",
        "trip_use_time": "Задать интервал времени",
        "trip_from": "С (дата/время)",
        "mud_caption": "Загрузка отчётов по раствору: PDF, Excel или CSV (файлы или почта). Единый дневной журнал.",
        "mud_src_files": "Загрузить файлы",
        "mud_src_email": "Электронная почта",
        "mud_chip_upload": "📁 Загрузить файлы",
        "mud_chip_email": "Электронная почта",
        "api_config_slide": "Настройки API — Slide Sheet",
        "api_base_url": "Базовый URL API",
        "api_token": "Токен доступа",
        "api_slide_path": "Путь slide sheet (опционально)",
        "api_slide_help": "Используйте {well_uuid}.",
        "use_depth_range": "Диапазон глубины",
        "depth_from": "Глубина от",
        "depth_to": "Глубина до",
        "api_note_html": "Не используйте /api/v1, если ответ HTML; берите /public/api/v1.",
        "enter_token": "Введите действительный токен для API.",
        "list_projects_err": "Не удалось получить список проектов:",
        "no_projects": "Нет проектов для этого токена.",
        "unnamed": "Без имени",
        "project": "Проект",
        "project_no_uuid": "У проекта нет UUID.",
        "list_wells_err": "Не удалось получить список скважин:",
        "no_wells": "В проекте нет скважин.",
        "well": "Скважина",
        "well_no_uuid": "У скважины нет UUID.",
        "kpi_rt": "**Обновление в реальном времени (API)**",
        "kpi_auto_refresh": "Автообновление каждые 30 с",
        "kpi_auto_help": "Перезагружать Slide Sheet из API каждые X с.",
        "kpi_interval_help": "Период перезагрузки (10–300 с).",
        "load_slide_sheet": "Загрузить Slide Sheet",
        "kpi_active_source": "Источник: Slide Sheet (API).",
        "show_rows": "Показать строки",
        "all_rows": "Все",
        "next_refresh": "🔄 Следующее обновление через **{i}** с… (снимите автообновление, чтобы остановить)",
        "solo_expander": "Подключение SOLO (API)",
        "bha_api_cfg": "Настройки API",
        "bha_missing_token": "Нет SOLO_ACCESS_TOKEN. Укажите в .env или вставьте токен.",
        "solo_token_loaded": "Токен загружен из .env (вставлять не нужно).",
        "solo_change_token": "Сменить токен",
        "solo_token_help": "Укажите SOLO_ACCESS_TOKEN в .env, чтобы не вставлять сюда.",
        "kpi_api_no_intervals": "API не вернул интервалы для этого диапазона.",
        "kpi_slide_loaded": "Slide Sheet загружен: {n:,} интервалов.",
        "kpi_load_err": "Не удалось загрузить Slide Sheet:",
        "trip_to": "По (дата/время)",
        "trip_help_dt": "ISO 8601 или формат, который понимает pandas.",
        "trip_use_depth": "Диапазон глубины",
        "trip_depth_min_l": "Мин. глубина (м)",
        "trip_depth_max_l": "Макс. глубина (м)",
        "trip_depth_help_min": "Только точки с Bit depth ≥ значения. Колонка в метрах.",
        "trip_depth_help_max": "Только точки с Bit depth ≤ значения. Нет данных — проверьте диапазон в CSV.",
        "trip_badge_hookload": "Hookload",
        "trip_badge_trip_io": "Trip In / Out",
        "trip_badge_overpull": "Overpull",
        "trip_gap_caption": "Если в диапазоне дат есть пропуски, их можно заполнить линейной интерполяцией.",
        "trip_interp_chk": "Заполнять пропуски линейной интерполяцией",
        "trip_interp_help": "Создаёт регулярную сетку по времени и интерполирует Hookload и Bit depth.",
        "trip_interp_interval_l": "Шаг интерполяции (сек)",
        "trip_interp_interval_help": "Интервал между точками. Меньше (1–5 с) — больше точек.",
        "mud_chip_bitacora": "Журнал",
        "mud_chip_formats": "PDF / Excel / CSV",
        "mud_chip_mail_short": "Почта",
        "mud_imap_expander": "Настройка почты (IMAP)",
        "mud_imap_caption": "Используйте **MUD_IMAP_SERVER**, **MUD_IMAP_USER** и **MUD_IMAP_PASS** в .env, чтобы не вводить пароль здесь.",
        "mud_imap_server": "Сервер IMAP",
        "mud_imap_server_help": "Напр. imap.gmail.com",
        "mud_imap_user": "Пользователь (email)",
        "mud_imap_pass": "Пароль (пароль приложения Gmail)",
        "mud_imap_pass_help": "В Gmail — пароль приложения, не от аккаунта.",
        "mud_imap_filter": "Фильтр по имени файла (опционально)",
        "mud_mark_read": "Помечать письма прочитанными после загрузки",
        "mud_auto_hdr": "**Автоопрос**",
        "mud_auto_chk": "Проверять почту каждые 60 с",
        "mud_auto_help": "Каждые X с опрашивается почта и обновляется журнал. Снимите флажок, чтобы остановить.",
        "mud_interval_imap": "Интервал (секунды)",
        "mud_interval_imap_help": "Как часто опрашивать почту (30–300 с).",
        "mud_fetch_btn": "🔥 Rogii — проверить почту и загрузить отчёты",
        "mud_fetch_help": "IMAP: скачать вложения PDF/Excel/CSV из непрочитанных писем.",
        "mud_err_imap": "Укажите сервер IMAP, пользователя и пароль (или задайте в .env).",
        "mud_spinner_imap": "Подключение к почте и загрузка вложений...",
        "mud_err_dl": "Не удалось подключиться или скачать:",
        "mud_no_attach": "Нет вложений PDF/Excel/CSV в непрочитанных (или фильтр не совпал).",
        "mud_success_attach": "Скачано вложений: **{n}**. Обработка...",
        "tab_trip_env": "Огибающая / Overpull",
        "tab_trip_broom": "Broomstick (FF)",
        "mud_email_countdown": "🔥 **Rogii** – Следующая проверка почты через **{i}** с… (снимите автоопрос почты)",
        "mud_upload_reports": "Загрузить отчёты по раствору (PDF, Excel, CSV)",
    },
}

# Valores internos estables (traducción solo en la etiqueta)
TRIP_ENV_EXACT = "exact"
TRIP_ENV_BIN = "bin"
MUD_SRC_FILES = "files"
MUD_SRC_EMAIL = "email"


def tr(key: str) -> str:
    """Texto de UI según `st.session_state['ui_lang']` (es/en/ru)."""
    lang = st.session_state.get("ui_lang", "es")
    bundle = I18N.get(lang) or I18N["es"]
    return bundle.get(key) or I18N["es"].get(key) or key


def render_language_selector_sidebar() -> None:
    """Selector visible arriba de la barra lateral (siempre)."""
    with st.sidebar:
        st.selectbox(
            "Idioma / Language / Язык",
            options=["es", "en", "ru"],
            format_func=lambda c: {"es": "Español", "en": "English", "ru": "Русский"}[c],
            key="ui_lang",
        )
        st.divider()


# =========================
# PPTX Helpers
# =========================
def add_table_slide(prs: Presentation, title: str, data: List[List], headers: List[str]):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    rows, cols = len(data) + 1, len(headers)

    table = slide.shapes.add_table(
        rows, cols, Inches(0.6), Inches(1.3), Inches(9.0), Inches(5.0)
    ).table

    for j, header in enumerate(headers):
        table.cell(0, j).text = str(header)

    for i, row in enumerate(data, start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = str(val)

    for cell in table.iter_cells():
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(40, 40, 40)

    return slide


def add_title_slide(prs: Presentation) -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = APP_TITLE
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"Generado: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(8.1), Inches(0.2), height=Inches(0.9))


def add_image_slide(prs: Presentation, title: str, buf: io.BytesIO) -> None:
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.add_picture(buf, Inches(0.8), Inches(1.3), Inches(8.0), Inches(4.5))


def add_text_slide(prs: Presentation, title: str, body: str) -> None:
    """Diapositiva con título + cuerpo multilínea (para resúmenes copiables a WhatsApp)."""
    try:
        slide_layout = prs.slide_layouts[6]
    except Exception:
        slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9.0), Inches(0.65))
    title_box.text_frame.text = title[:250]
    for p in title_box.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(16)
            r.font.bold = True
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9.0), Inches(6.0))
    tf = body_box.text_frame
    tf.word_wrap = True
    parts = [s.strip() for s in (body or "").split("\n\n") if s.strip()]
    if not parts:
        tf.text = ""
        return
    tf.text = parts[0]
    p0 = tf.paragraphs[0]
    for r in p0.runs:
        r.font.size = Pt(11)
    for para_text in parts[1:]:
        p = tf.add_paragraph()
        p.text = para_text
        p.space_after = Pt(6)
        for r in p.runs:
            r.font.size = Pt(11)


def export_pptx_to_pdf(pptx_path: Path) -> Path | None:
    try:
        import comtypes.client  # type: ignore
    except ImportError:
        return None

    pdf_path = pptx_path.with_suffix(".pdf")
    powerpoint = None
    presentation = None
    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
        presentation.ExportAsFixedFormat(str(pdf_path), 32)
        return pdf_path
    except Exception:
        return None
    finally:
        try:
            if presentation is not None:
                presentation.Close()
        finally:
            if powerpoint is not None:
                powerpoint.Quit()


def _copy_file_to_downloads(src: Path, dest: Path) -> Path | None:
    """
    Copia src -> dest leyendo en memoria (evita algunos bloqueos de copy2 en Windows).
    Si dest está en uso (WinError 32), intenta un nombre con sufijo de hora.
    """
    if not src.is_file():
        return None
    downloads_dir = dest.parent
    try:
        downloads_dir.mkdir(parents=True, exist_ok=True)
    except OSError:
        return None
    data = src.read_bytes()
    candidates = [dest]
    stem, suf = dest.stem, dest.suffix
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    candidates.append(downloads_dir / f"{stem}_{ts}{suf}")
    for i in range(1, 6):
        candidates.append(downloads_dir / f"{stem}_{ts}_{i}{suf}")

    for target in candidates:
        try:
            target.write_bytes(data)
            return target
        except PermissionError:
            continue
        except OSError:
            continue
    return None


def copy_report_to_downloads(
    pptx_path: Path | None, pdf_path: Path | None, base_name: str
) -> None:
    downloads_dir = Path.home() / "Downloads"
    if not downloads_dir.exists():
        return

    if pptx_path is not None and pptx_path.exists():
        pptx_target = downloads_dir / f"{base_name}.pptx"
        _copy_file_to_downloads(pptx_path, pptx_target)
    if pdf_path is not None and pdf_path.exists():
        pdf_target = downloads_dir / f"{base_name}.pdf"
        _copy_file_to_downloads(pdf_path, pdf_target)


def apply_pro_theme(fig, h: int = 420):
    fig.update_layout(
        template=PLOTLY_TEMPLATE,
        height=h,
        margin=dict(l=50, r=30, t=40, b=55),
        title=dict(x=0.02, xanchor="left"),
        title_pad=dict(t=4, b=4),
        font=dict(family="Segoe UI", size=12, color="#2A2A2A"),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
        uniformtext_minsize=10,
        uniformtext_mode="hide",
    )
    fig.update_xaxes(showgrid=True, gridcolor="rgba(0,0,0,0.06)", zeroline=False)
    fig.update_yaxes(showgrid=True, gridcolor="rgba(0,0,0,0.06)", zeroline=False)
    return fig


def apply_pro_theme_dark(fig, h: int = 420):
    """Tema oscuro consistente con Streamlit dark mode (fondo transparente + texto claro)."""
    fig.update_layout(
        template="plotly_dark",
        height=h,
        margin=dict(l=50, r=30, t=40, b=55),
        title=dict(x=0.02, xanchor="left"),
        title_pad=dict(t=4, b=4),
        font=dict(family="Segoe UI", size=12, color="#E5E7EB"),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
        uniformtext_minsize=10,
        uniformtext_mode="hide",
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(17,24,39,0.92)",
    )
    fig.update_xaxes(showgrid=True, gridcolor="rgba(255,255,255,0.08)", zeroline=False)
    fig.update_yaxes(showgrid=True, gridcolor="rgba(255,255,255,0.08)", zeroline=False)
    return fig


def is_streamlit_dark_mode() -> bool:
    """Devuelve True si Streamlit está en theme base 'dark'."""
    try:
        return str(st.get_option("theme.base")).lower() == "dark"
    except Exception:
        return False


def prettify_auto(fig, h: int = 420):
    """Prettify que respeta el tema (dark/light) automáticamente."""
    return apply_pro_theme_dark(fig, h=h) if is_streamlit_dark_mode() else apply_pro_theme(fig, h=h)




def prettify(fig, h: int = 420):
    return apply_pro_theme(fig, h=h)


def prettify_hist(fig, h: int = 420):
    fig.update_layout(
        template=PLOTLY_TEMPLATE,
        height=h,
        margin=dict(l=50, r=30, t=18, b=55),
        title=dict(x=0.02, xanchor="left", y=0.98),
        title_pad=dict(t=0, b=0),
        bargap=0.05,
        font=dict(family="Segoe UI", size=12, color="#2A2A2A"),
    )
    fig.update_xaxes(showgrid=True, gridcolor="rgba(0,0,0,0.06)")
    fig.update_yaxes(showgrid=True, gridcolor="rgba(0,0,0,0.06)")
    return fig


def save_and_show_plotly(
    prs: Presentation,
    title: str,
    fig,
    show_plots: bool = True,
    is_hist: bool = False,
    is_heatmap: bool = False,
) -> None:
    if is_heatmap:
        fig = prettify_heatmap(fig)
    else:
        fig = prettify_hist(fig) if is_hist else prettify(fig)
    png_bytes = fig.to_image(format="png", scale=2)
    buf = io.BytesIO(png_bytes)
    buf.seek(0)
    add_image_slide(prs, title, buf)

    if show_plots:
        st.plotly_chart(
            fig,
            use_container_width=True,
            config=PLOTLY_CONFIG,
        )


# Colores pro para medidor de eficiencia (alineados con MSE y tendencias)
GAUGE_TRACK_COLOR = "#334155"
GAUGE_FILL_COLOR = "#0ea5e9"
GAUGE_NEEDLE_COLOR = "#c2410c"
GAUGE_LABEL_COLOR = "#94a3b8"


def _efficiency_status(rotary_pct: float) -> tuple[str, str]:
    """Devuelve (etiqueta, color) para chip de estado según % rotary."""
    if rotary_pct >= 70:
        return "Alta eficiencia", "green"
    if rotary_pct >= 40:
        return "Eficiencia media", "orange"
    return "Baja eficiencia", "red"


def _render_chips_row(items: list[tuple[str, str]]) -> None:
    """Muestra una fila de chips (badges) pro. items = [(label, color), ...]."""
    if not items:
        return
    try:
        cols = st.columns(len(items))
        for i, (label, color) in enumerate(items):
            with cols[i]:
                st.badge(label, color=color, width="content")
    except Exception:
        parts = " ".join(f":{c}-badge[{l}]" for l, c in items)
        st.markdown(parts)


def render_efficiency_chips(rotary_pct: float, slide_pct: float, run_name: str) -> None:
    """Muestra una fila de chips (badges) pro para Rotary / Slide y estado."""
    status_label, status_color = _efficiency_status(rotary_pct)
    rp = round(rotary_pct, 1)
    sp = round(slide_pct, 1)
    try:
        c1, c2, c3, c4 = st.columns(4)
        with c1:
            st.badge(run_name, color="gray", width="content")
        with c2:
            st.badge(f"Rotary {rp}%", color="blue", width="content")
        with c3:
            st.badge(f"Slide {sp}%", color="gray", width="content")
        with c4:
            st.badge(status_label, color=status_color, width="content")
    except Exception:
        st.markdown(
            f":gray-badge[{run_name}] :blue-badge[Rotary {rp}%] :gray-badge[Slide {sp}%] :{status_color}-badge[{status_label}]"
        )


def gauge_efficiency(rotary_pct: float, run_name: str) -> None:
    value = round(min(100, max(0, float(rotary_pct))), 1)
    option = {
        "series": [
            {
                "type": "gauge",
                "min": 0,
                "max": 100,
                "startAngle": 210,
                "endAngle": -30,
                "progress": {
                    "show": True,
                    "width": 20,
                    "roundCap": True,
                    "itemStyle": {"color": GAUGE_FILL_COLOR},
                },
                "axisLine": {
                    "lineStyle": {
                        "width": 20,
                        "color": [[1, GAUGE_TRACK_COLOR]],
                    }
                },
                "axisTick": {"show": False},
                "splitLine": {"show": False},
                "axisLabel": {
                    "distance": 22,
                    "color": GAUGE_LABEL_COLOR,
                    "fontSize": 11,
                },
                "pointer": {
                    "width": 4,
                    "length": "65%",
                    "itemStyle": {"color": GAUGE_NEEDLE_COLOR},
                },
                "anchor": {"show": True, "size": 14, "itemStyle": {"borderColor": GAUGE_NEEDLE_COLOR, "borderWidth": 2}},
                "title": {
                    "show": True,
                    "offsetCenter": [0, "58%"],
                    "fontSize": 12,
                    "color": GAUGE_LABEL_COLOR,
                },
                "detail": {
                    "valueAnimation": True,
                    "formatter": "{value}%",
                    "fontSize": 26,
                    "fontWeight": "bold",
                    "offsetCenter": [0, "28%"],
                    "color": "#f8fafc",
                },
                "data": [{"value": value, "name": f"Rotary % · {run_name}"}],
            }
        ]
    }
    st_echarts(option, height="340px")


def prettify_heatmap(fig, h: int = 520):
    fig.update_layout(
        template=PLOTLY_TEMPLATE,
        height=h,
        margin=dict(l=60, r=30, t=48, b=60),
        title=dict(x=0.02, xanchor="left", font=dict(size=15)),
        font=dict(family="Segoe UI", size=12, color="#2A2A2A"),
        plot_bgcolor="rgba(248,250,252,0.65)",
    )
    fig.update_xaxes(showgrid=False, tickangle=-35)
    fig.update_yaxes(showgrid=False)
    return fig


def prettify_heatmap_auto(fig, h: int = 520):
    """Heatmap con tema claro/oscuro alineado al dashboard."""
    if is_streamlit_dark_mode():
        fig.update_layout(
            template="plotly_dark",
            height=h,
            margin=dict(l=60, r=30, t=48, b=60),
            title=dict(x=0.02, xanchor="left", font=dict(size=15, color="#F1F5F9")),
            font=dict(family="Segoe UI", size=12, color="#E2E8F0"),
            paper_bgcolor="rgba(0,0,0,0)",
            plot_bgcolor="rgba(15,23,42,0.5)",
        )
        fig.update_xaxes(showgrid=False, tickangle=-35, tickfont=dict(color="#CBD5E1"))
        fig.update_yaxes(showgrid=False, tickfont=dict(color="#CBD5E1"))
        return fig
    return prettify_heatmap(fig, h=h)


def heatmap_numeric_stats(df: pd.DataFrame, cols: list) -> pd.DataFrame:
    """Min / media / max / N por columna (datos fuente del heatmap de correlación)."""
    cols = [c for c in cols if c is not None and str(c).strip() and c in df.columns]
    if not cols:
        return pd.DataFrame()
    d = df[cols].copy()
    for c in cols:
        d[c] = pd.to_numeric(d[c], errors="coerce")
    rows: list[dict] = []
    for c in cols:
        s = d[c].dropna()
        if s.empty:
            rows.append(
                {
                    "Parámetro": str(c),
                    "Mínimo": np.nan,
                    "Promedio": np.nan,
                    "Máximo": np.nan,
                    "N": 0,
                }
            )
        else:
            rows.append(
                {
                    "Parámetro": str(c),
                    "Mínimo": float(s.min()),
                    "Promedio": float(s.mean()),
                    "Máximo": float(s.max()),
                    "N": int(len(s)),
                }
            )
    return pd.DataFrame(rows)


def build_minmax_mean_spine_figure(
    stats_df: pd.DataFrame,
    title: str = "Rango por parámetro (0=min, 1=max, ●=media)",
) -> go.Figure | None:
    """
    Por cada variable: segmento vertical 0→1 en escala normalizada al rango observado;
    marcador azul = posición de la media en ese rango.
    """
    if stats_df is None or stats_df.empty or "Parámetro" not in stats_df.columns:
        return None
    fig = go.Figure()
    for _, r in stats_df.iterrows():
        lo, mid, hi = r.get("Mínimo"), r.get("Promedio"), r.get("Máximo")
        p = str(r["Parámetro"])
        if pd.isna(lo) or pd.isna(mid) or pd.isna(hi):
            continue
        lo_f, mid_f, hi_f = float(lo), float(mid), float(hi)
        span = hi_f - lo_f
        if span <= 0 or not np.isfinite(span):
            ym = 0.5
        else:
            ym = (mid_f - lo_f) / span
            ym = float(min(1.0, max(0.0, ym)))
        fig.add_trace(
            go.Scatter(
                x=[p, p],
                y=[0.0, 1.0],
                mode="lines",
                line=dict(width=3, color="rgba(148,163,184,0.9)"),
                showlegend=False,
                hovertemplate=(
                    f"<b>{p}</b><br>min: {lo_f:.6g}<br>max: {hi_f:.6g}<br>n: {int(r.get('N', 0))}<extra></extra>"
                ),
            )
        )
        fig.add_trace(
            go.Scatter(
                x=[p],
                y=[ym],
                mode="markers",
                marker=dict(size=11, color="#0ea5e9", line=dict(width=2, color="white")),
                showlegend=False,
                hovertemplate=f"<b>{p}</b><br>media: {mid_f:.6g}<br>posición en rango: {ym:.2f}<extra></extra>",
            )
        )
    fig.update_layout(
        title=dict(text=title, x=0.02, xanchor="left"),
        xaxis_title="Parámetro",
        yaxis_title="Normalizado (min→max)",
        yaxis=dict(range=[-0.08, 1.08], tickvals=[0, 0.5, 1], ticktext=["Min", "0.5", "Max"]),
        height=400,
        template=PLOTLY_TEMPLATE,
        margin=dict(l=52, r=28, t=52, b=96),
        font=dict(family="Segoe UI", size=11, color="#334155"),
    )
    if is_streamlit_dark_mode():
        fig.update_layout(
            template="plotly_dark",
            paper_bgcolor="#0b0d14",
            plot_bgcolor="#0b0d14",
            font=dict(family="Segoe UI", size=11, color="#E2E8F0"),
            title=dict(font=dict(color="#F8FAFC")),
            xaxis=dict(tickfont=dict(color="#CBD5E1")),
            yaxis=dict(
                range=[-0.08, 1.08],
                tickvals=[0, 0.5, 1],
                ticktext=["Min", "0.5", "Max"],
                tickfont=dict(color="#CBD5E1"),
            ),
        )
    return fig


def build_rop_top_zones_bar_figure(
    zone_stats: dict,
    top_n: int = 8,
    title: str = "Top zonas operativas por ROP",
) -> go.Figure | None:
    """Barras horizontales: mejores celdas WOB×RPM por ROP medio (estilo panel oscuro)."""
    stat = np.asarray(zone_stats.get("stat"), dtype=float)
    x_edges = zone_stats.get("x_edges")
    y_edges = zone_stats.get("y_edges")
    if stat.size == 0 or x_edges is None or y_edges is None:
        return None
    rows: list[tuple[float, str]] = []
    nwx, nwy = stat.shape
    for i in range(nwx):
        for j in range(nwy):
            v = stat[i, j]
            if not np.isfinite(v):
                continue
            w0, w1 = float(x_edges[i]), float(x_edges[i + 1])
            r0, r1 = float(y_edges[j]), float(y_edges[j + 1])
            label = f"WOB {w0:.0f}-{w1:.0f} | RPM {r0:.0f}-{r1:.0f}"
            rows.append((float(v), label))
    if not rows:
        return None
    rows.sort(key=lambda t: -t[0])
    rows = rows[: max(1, int(top_n))]
    rows.reverse()
    rops = [r[0] for r in rows]
    labels = [r[1] for r in rows]
    # En barras horizontales, la 1ª categoría en ``y`` va **abajo** y la última **arriba**.
    # Tras reverse(), la mejor ROP queda al **final** → el naranja debe ir en el **último** ítem.
    n_b = len(rows)
    colors = ["#2dd4bf"] * n_b
    if n_b:
        colors[-1] = "#f97316"
    bar_text = [f"{v:.1f}" for v in rops]
    if n_b:
        bar_text[-1] = f"{rops[-1]:.1f} · mejor"
    fig = go.Figure(
        go.Bar(
            x=rops,
            y=labels,
            orientation="h",
            text=bar_text,
            textposition="outside",
            textfont=dict(color="#e2e8f0", size=11),
            hovertemplate="%{y}<br>ROP medio: %{x:.2f}<extra></extra>",
            marker=dict(
                color=colors,
                line=dict(color="rgba(255,255,255,0.22)", width=1),
            ),
        )
    )
    fig.update_layout(
        title=dict(text=title, x=0.02, xanchor="left", font=dict(size=15, color="#f1f5f9")),
        xaxis_title="ROP medio",
        yaxis_title="Zona WOB-RPM",
        template="plotly_dark",
        paper_bgcolor="#0b0d14",
        plot_bgcolor="#0b0d14",
        height=max(380, min(520, 44 * len(rows) + 140)),
        margin=dict(l=210, r=72, t=56, b=52),
        font=dict(family="Segoe UI", size=11, color="#e2e8f0"),
        xaxis=dict(gridcolor="rgba(255,255,255,0.08)", zeroline=False),
        yaxis=dict(gridcolor="rgba(255,255,255,0.04)", automargin=True),
    )
    # plotly_dark a veces aplana el color de las barras; reforzar por barra
    fig.update_traces(
        marker=dict(color=colors, line=dict(color="rgba(255,255,255,0.22)", width=1)),
    )
    return fig


def _x_contiguous_segments_where_true(x_sorted: np.ndarray, mask: np.ndarray) -> list[tuple[float, float]]:
    """Tramos contiguos en x donde ``mask`` es True (mismos índices que x ordenado)."""
    if x_sorted.size == 0 or mask.size != x_sorted.size:
        return []
    out: list[tuple[float, float]] = []
    n = int(mask.size)
    i = 0
    while i < n:
        if not bool(mask[i]):
            i += 1
            continue
        j = i + 1
        while j < n and bool(mask[j]):
            j += 1
        out.append((float(x_sorted[i]), float(x_sorted[j - 1])))
        i = j
    return out


def kpi_depth_optimal_zone_chips(
    xv: np.ndarray,
    y_r: np.ndarray,
    in_zone: np.ndarray,
    x_title: str,
    zone_stats: dict,
    n_segments: int,
) -> list[tuple[str, str]]:
    """Chips pro para la franja de profundidad/índice en la celda óptima WOB×RPM."""
    items: list[tuple[str, str]] = []
    br = float(zone_stats.get("best_rop", 0.0))
    items.append((f"Celda óptima heatmap · ROP bin {br:.1f}", "orange"))

    if not np.any(in_zone):
        items.append(("Sin puntos en franja (WOB×RPM fuera de celda)", "gray"))
        return items

    ntot = max(int(in_zone.size), 1)
    pct = 100.0 * float(np.count_nonzero(in_zone)) / ntot
    items.append((f"Muestras en franja: {pct:.1f}% del tramo", "blue"))
    rop_z = y_r[in_zone]
    items.append((f"ROP medio en franja: {float(np.mean(rop_z)):.1f}", "green"))
    items.append((f"ROP máx. en franja: {float(np.max(rop_z)):.1f}", "green"))
    xv_z = xv[in_zone]
    if "Profundidad" in x_title:
        items.append(
            (f"Prof. franja: {float(np.min(xv_z)):.0f} – {float(np.max(xv_z)):.0f}", "blue"),
        )
    else:
        items.append(
            (f"Tramos contiguos en celda: {n_segments}", "blue"),
        )
    w0, w1 = float(zone_stats["best_wob_low"]), float(zone_stats["best_wob_high"])
    r0, r1 = float(zone_stats["best_rpm_low"]), float(zone_stats["best_rpm_high"])
    items.append((f"WOB {w0:.0f}-{w1:.0f} · RPM {r0:.0f}-{r1:.0f}", "gray"))
    return items


def build_kpi_depth_curves_figure(
    df: pd.DataFrame,
    depth_col: str | None,
    rop_col: str,
    wob_col: str,
    rpm_col: str,
    title: str = "Curvas suavizadas de ROP, WOB y RPM",
    zone_stats: dict | None = None,
) -> tuple[go.Figure | None, list[tuple[str, str]]]:
    """
    Series en **unidades reales** vs profundidad (o índice): ROP (eje Y izq.), WOB y RPM (ejes Y derecha).

    Si ``zone_stats`` tiene la celda óptima, se sombrean **franjas verticales** donde WOB y RPM
    caen dentro de esos rangos (tramos contiguos en X).
    """
    need = [rop_col, wob_col, rpm_col]
    if not all(c in df.columns for c in need):
        return None, []
    y_r = pd.to_numeric(df[rop_col], errors="coerce")
    y_w = pd.to_numeric(df[wob_col], errors="coerce")
    y_m = pd.to_numeric(df[rpm_col], errors="coerce")
    if depth_col and str(depth_col).strip() and depth_col in df.columns:
        x = pd.to_numeric(df[depth_col], errors="coerce")
        x_title = "Profundidad"
    else:
        x = pd.Series(np.arange(len(df), dtype=float), index=df.index)
        x_title = "Índice de muestra (sin columna de profundidad)"
    m = x.notna() & y_r.notna() & y_w.notna() & y_m.notna()
    if int(m.sum()) < 2:
        return None, []
    xv = x[m].to_numpy(dtype=float)
    oi = np.argsort(xv, kind="mergesort")
    xv = xv[oi]
    y_r = y_r[m].to_numpy(dtype=float)[oi]
    y_w = y_w[m].to_numpy(dtype=float)[oi]
    y_m = y_m[m].to_numpy(dtype=float)[oi]

    chips: list[tuple[str, str]] = []
    segments: list[tuple[float, float]] = []
    in_zone = np.zeros(len(xv), dtype=bool)
    if zone_stats is not None:
        wlo = float(zone_stats["best_wob_low"])
        whi = float(zone_stats["best_wob_high"])
        rlo = float(zone_stats["best_rpm_low"])
        rhi = float(zone_stats["best_rpm_high"])
        in_zone = (y_w >= wlo) & (y_w <= whi) & (y_m >= rlo) & (y_m <= rhi)
        segments = _x_contiguous_segments_where_true(xv, in_zone)
        chips = kpi_depth_optimal_zone_chips(xv, y_r, in_zone, x_title, zone_stats, len(segments))

    fig = go.Figure()
    for x0, x1 in segments:
        span = max(x1 - x0, 1e-9)
        pad = min(span * 0.008, (float(np.nanmax(xv)) - float(np.nanmin(xv))) * 0.002 + 1e-9)
        fig.add_vrect(
            x0=x0 - pad,
            x1=x1 + pad,
            fillcolor="rgba(249,115,22,0.16)",
            line=dict(color="rgba(251,146,60,0.65)", width=1.2),
            layer="below",
        )
    fig.add_trace(
        go.Scatter(
            x=xv,
            y=y_r,
            name="ROP",
            mode="lines",
            line=dict(color="#f97316", width=2.4),
            yaxis="y",
        )
    )
    fig.add_trace(
        go.Scatter(
            x=xv,
            y=y_w,
            name="WOB",
            mode="lines",
            line=dict(color="#2dd4bf", width=2),
            yaxis="y2",
        )
    )
    fig.add_trace(
        go.Scatter(
            x=xv,
            y=y_m,
            name="RPM",
            mode="lines",
            line=dict(color="#64748b", width=2),
            yaxis="y3",
        )
    )
    fig.update_layout(
        template="plotly_dark",
        paper_bgcolor="#0b0d14",
        plot_bgcolor="#0b0d14",
        title=dict(text=title, x=0.02, font=dict(size=15, color="#f1f5f9")),
        xaxis=dict(title=x_title, gridcolor="rgba(255,255,255,0.08)", zeroline=False),
        yaxis=dict(
            title=dict(text="ROP", font=dict(color="#fdba74")),
            side="left",
            gridcolor="rgba(255,255,255,0.06)",
            zeroline=False,
            tickfont=dict(color="#fdba74"),
        ),
        yaxis2=dict(
            title=dict(text="WOB", font=dict(color="#5eead4")),
            overlaying="y",
            side="right",
            showgrid=False,
            zeroline=False,
            tickfont=dict(color="#5eead4"),
        ),
        yaxis3=dict(
            title=dict(text="RPM", font=dict(color="#94a3b8")),
            overlaying="y",
            side="right",
            anchor="free",
            position=0.97,
            showgrid=False,
            zeroline=False,
            tickfont=dict(color="#cbd5e1"),
        ),
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.02,
            xanchor="left",
            x=0,
            font=dict(size=11, color="#e2e8f0"),
            bgcolor="rgba(11,13,20,0.7)",
        ),
        margin=dict(l=58, r=92, t=72, b=48),
        height=460,
        font=dict(family="Segoe UI", color="#e2e8f0"),
        hovermode="x unified",
    )
    if segments:
        fig.add_trace(
            go.Scatter(
                x=[np.nan],
                y=[np.nan],
                mode="markers",
                marker=dict(size=14, color="rgba(249,115,22,0.55)", symbol="square"),
                name="Franja zona óptima (WOB×RPM)",
            )
        )
        fig.update_layout(
            title=dict(
                text=title + " · franjas = celda óptima",
                x=0.02,
                font=dict(size=15, color="#f1f5f9"),
            )
        )
    return fig, chips


def stats_df_to_heatmap_chips(stats_df: pd.DataFrame, max_chips: int = 12) -> list[tuple[str, str]]:
    """Etiquetas cortas para fila de chips (nombre + min–max + media)."""
    items: list[tuple[str, str]] = []
    for _, r in stats_df.iterrows():
        if int(r.get("N", 0) or 0) < 1:
            continue
        name = str(r["Parámetro"])
        if len(name) > 18:
            name = name[:16] + "…"
        lo, mid, hi = r["Mínimo"], r["Promedio"], r["Máximo"]
        sub = f"{format_num(lo, 1)}–{format_num(hi, 1)} · μ{format_num(mid, 1)}"
        items.append((f"{name}: {sub}", "blue"))
        if len(items) >= max_chips:
            break
    return items


def build_heatmap_marginal_max_curves(
    zone_stats: dict,
    x_label: str = "WOB (centro de bin)",
    y_label: str = "RPM (centro de bin)",
    z_label: str = "ROP máx. en bin",
) -> go.Figure | None:
    """
    Curvas adicionales: máximo de ROP a lo largo de cada eje del heatmap 2D (por bin).
    """
    stat = zone_stats.get("stat")
    x_c = zone_stats.get("x_centers")
    y_c = zone_stats.get("y_centers")
    if stat is None or x_c is None or y_c is None:
        return None
    stat = np.asarray(stat, dtype=float)
    if stat.size == 0:
        return None
    # stat[i,j] = valor en bin x_i, y_j
    max_along_rpm = np.nanmax(stat, axis=1)
    max_along_wob = np.nanmax(stat, axis=0)
    from plotly.subplots import make_subplots

    fig = make_subplots(
        rows=2,
        cols=1,
        row_heights=[0.45, 0.55],
        vertical_spacing=0.14,
        subplot_titles=(
            f"{z_label} vs {x_label.split('(')[0].strip()} (máx. por bin WOB)",
            f"{z_label} vs {y_label.split('(')[0].strip()} (máx. por bin RPM)",
        ),
    )
    fig.add_trace(
        go.Scatter(
            x=x_c,
            y=max_along_rpm,
            mode="lines+markers",
            name="Máx por WOB",
            line=dict(color="#2563EB", width=2.5),
            marker=dict(size=6),
        ),
        row=1,
        col=1,
    )
    fig.add_trace(
        go.Scatter(
            x=y_c,
            y=max_along_wob,
            mode="lines+markers",
            name="Máx por RPM",
            line=dict(color="#F59E0B", width=2.5),
            marker=dict(size=6),
        ),
        row=2,
        col=1,
    )
    fig.update_xaxes(title_text=x_label, row=1, col=1)
    fig.update_yaxes(title_text=z_label, row=1, col=1)
    fig.update_xaxes(title_text=y_label, row=2, col=1)
    fig.update_yaxes(title_text=z_label, row=2, col=1)
    fig.update_layout(
        height=520,
        showlegend=False,
        margin=dict(l=55, r=25, t=56, b=48),
        font=dict(family="Segoe UI", size=11),
    )
    if is_streamlit_dark_mode():
        fig.update_layout(
            template="plotly_dark",
            paper_bgcolor="rgba(0,0,0,0)",
            plot_bgcolor="rgba(17,24,39,0.85)",
        )
    else:
        fig.update_layout(template=PLOTLY_TEMPLATE, plot_bgcolor="rgba(248,250,252,0.9)")
    return fig


def rop_zone_dashboard_chips(zone_stats: dict) -> list[tuple[str, str]]:
    """Chips de KPI para el dashboard ROP (mejor celda + extremos en la grilla 2D)."""
    stat = np.asarray(zone_stats.get("stat"), dtype=float)
    items: list[tuple[str, str]] = []
    br = zone_stats.get("best_rop")
    if br is not None and np.isfinite(float(br)):
        items.append((f"Mejor ROP (celda): {float(br):.1f}", "green"))
    items.append(
        (
            f"WOB óptimo: {zone_stats.get('best_wob_center', 0):.1f}",
            "blue",
        )
    )
    items.append(
        (
            f"RPM óptimo: {zone_stats.get('best_rpm_center', 0):.1f}",
            "blue",
        )
    )
    if stat.size:
        gmax = np.nanmax(stat)
        gmin = np.nanmin(stat)
        if np.isfinite(gmax):
            items.append((f"ROP máx (grilla): {float(gmax):.1f}", "orange"))
        if np.isfinite(gmin):
            items.append((f"ROP mín (grilla): {float(gmin):.1f}", "gray"))
    n = zone_stats.get("best_count")
    if n is not None:
        items.append((f"Puntos en mejor bin: {int(n)}", "gray"))
    return items


def _rop_heatmap_label_matrix_top_fraction(
    z_t: np.ndarray,
    top_fraction: float = ROP_HEATMAP_LABEL_TOP_FRACTION,
) -> np.ndarray:
    """Texto por celda: solo valores en el top ``top_fraction`` (p. ej. 15 %) del ROP en la grilla."""
    z = np.asarray(z_t, dtype=float)
    flat = z[np.isfinite(z)]
    if flat.size == 0:
        return np.full(z.shape, "", dtype=object)
    pct = 100.0 * (1.0 - float(np.clip(top_fraction, 0.01, 0.5)))
    thr = float(np.nanpercentile(flat, pct))

    def _cell(v: object) -> str:
        if not np.isfinite(float(v)):
            return ""
        if float(v) < thr:
            return ""
        return f"{float(v):.1f}"

    return np.vectorize(_cell, otypes=[object])(z)


def build_optimal_rop_heatmap_with_marginals(
    zone_stats: dict,
    title: str = "Heatmap ROP vs WOB-RPM (con marginales)",
) -> go.Figure | None:
    """
    Vista tipo dashboard: heatmap central + curvas arriba (WOB vs ROP máx/mín por bin)
    + curvas a la derecha (RPM vs ROP máx/mín por bin), alineadas con los ejes del mapa.

    Siempre usa fondo oscuro tipo heatmap pro (alineado a prettify_heatmap_auto en dark)
    y canvas grande para mejor nitidez al escalar en Streamlit.
    """
    # Canvas grande para nitidez al escalar en pantalla / export PNG
    _dash_w, _dash_h = 1680, 1050
    # Un solo fondo (paper + plot + accesorios): evita el rectángulo más claro vs Streamlit dark
    _dash_bg = "#0b0d14"
    _fg = "#E2E8F0"
    _tick = "#CBD5E1"
    _grid_marg = "rgba(255,255,255,0.1)"
    stat = zone_stats.get("stat")
    x_c = zone_stats.get("x_centers")
    y_c = zone_stats.get("y_centers")
    if stat is None or x_c is None or y_c is None:
        return None
    stat = np.asarray(stat, dtype=float)
    if stat.size == 0:
        return None

    x_c = np.asarray(x_c, dtype=float)
    y_c = np.asarray(y_c, dtype=float)
    max_w = np.nanmax(stat, axis=1)
    min_w = np.nanmin(stat, axis=1)
    max_r = np.nanmax(stat, axis=0)
    min_r = np.nanmin(stat, axis=0)

    i, j = zone_stats["best_bin"]
    x0 = zone_stats["best_wob_low"]
    x1 = zone_stats["best_wob_high"]
    y0 = zone_stats["best_rpm_low"]
    y1 = zone_stats["best_rpm_high"]

    from plotly.subplots import make_subplots

    fig = make_subplots(
        rows=2,
        cols=2,
        row_heights=[0.24, 0.76],
        column_widths=[0.76, 0.24],
        specs=[[{"type": "scatter"}, None], [{"type": "heatmap"}, {"type": "scatter"}]],
        horizontal_spacing=0.04,
        vertical_spacing=0.07,
        shared_xaxes=True,
    )

    # (1,1) Marginal superior: WOB → ROP máx / mín al variar RPM en cada columna
    fig.add_trace(
        go.Scatter(
            x=x_c,
            y=max_w,
            mode="lines+markers",
            name="Máx ROP",
            line=dict(color="#F97316", width=3),
            marker=dict(size=7, line=dict(width=0)),
            legendgroup="top",
            showlegend=True,
            hovertemplate="WOB: %{x:.2f}<br>ROP máx. en bin: %{y:.2f}<extra></extra>",
        ),
        row=1,
        col=1,
    )
    fig.add_trace(
        go.Scatter(
            x=x_c,
            y=min_w,
            mode="lines+markers",
            name="Mín ROP",
            line=dict(color="#38BDF8", width=2.5, dash="dot"),
            marker=dict(size=6, line=dict(width=0)),
            legendgroup="top",
            showlegend=True,
            hovertemplate="WOB: %{x:.2f}<br>ROP mín. en bin: %{y:.2f}<extra></extra>",
        ),
        row=1,
        col=1,
    )

    z_t = stat.T
    z_ok = stat[np.isfinite(stat)]
    if z_ok.size:
        _zlo = float(np.nanmin(z_ok))
        _zhi = float(np.nanmax(z_ok))
        _zpad = max(0.15, (_zhi - _zlo) * 0.04)
        zmin_hm, zmax_hm = _zlo - _zpad, _zhi + _zpad
    else:
        zmin_hm, zmax_hm = None, None

    counts = zone_stats.get("counts")
    custom_cd = None
    if counts is not None:
        c_arr = np.asarray(counts, dtype=float)
        if c_arr.shape == stat.shape:
            custom_cd = c_arr.T

    hm_text = _rop_heatmap_label_matrix_top_fraction(z_t)
    hm_extras: dict = {
        "text": hm_text,
        "texttemplate": "%{text}",
        "textfont": dict(
            family="Segoe UI", size=13, color="rgba(248,250,252,0.94)"
        ),
    }

    # (2,1) Heatmap principal
    fig.add_trace(
        go.Heatmap(
            x=x_c,
            y=y_c,
            z=z_t,
            zmin=zmin_hm,
            zmax=zmax_hm,
            colorscale=ROP_HEATMAP_COLORSCALE,
            zsmooth="best",
            xgap=2,
            ygap=2,
            customdata=custom_cd,
            **hm_extras,
            colorbar=dict(
                title=dict(text="ROP<br>medio", font=dict(size=14, color=_fg)),
                tickfont=dict(size=13, color=_tick),
                tickformat=".1f",
                len=0.6,
                y=0.36,
                yanchor="middle",
                thickness=22,
                outlinewidth=0,
                bgcolor=_dash_bg,
                bordercolor="rgba(255,255,255,0.12)",
                borderwidth=1,
            ),
            hovertemplate=(
                "WOB (centro bin): %{x:.2f}<br>"
                "RPM (centro bin): %{y:.2f}<br>"
                "ROP medio: %{z:.2f}"
                + ("<br>Puntos en celda: %{customdata:.0f}" if custom_cd is not None else "")
                + "<extra></extra>"
            ),
            hoverongaps=False,
            showscale=True,
        ),
        row=2,
        col=1,
    )
    _bw = float(zone_stats["best_wob_center"])
    _br = float(zone_stats["best_rpm_center"])
    _brop = float(zone_stats["best_rop"])
    fig.add_trace(
        go.Scatter(
            x=[_bw],
            y=[_br],
            mode="markers",
            marker=dict(
                symbol="star",
                size=20,
                color="#FDE047",
                line=dict(color="rgb(15,23,42)", width=1.5),
            ),
            name="Pico",
            showlegend=False,
            hovertemplate=(
                "Pico operacional<br>WOB: %{x:.2f}<br>RPM: %{y:.2f}<br>"
                f"ROP: {_brop:.1f}<extra></extra>"
            ),
        ),
        row=2,
        col=1,
    )
    fig.add_shape(
        type="rect",
        x0=x0,
        x1=x1,
        y0=y0,
        y1=y1,
        line=dict(color="#ffffff", width=3.5),
        fillcolor="rgba(255,255,255,0)",
        row=2,
        col=1,
    )
    _yr_span = float(np.nanmax(y_c) - np.nanmin(y_c)) if y_c.size else 1.0
    _dy = (
        float(np.median(np.abs(np.diff(np.sort(y_c)))))
        if len(y_c) > 1
        else max(_yr_span * 0.04, 1.0)
    )
    fig.add_annotation(
        x=_bw,
        y=_br + 1.35 * _dy,
        text="Pico",
        showarrow=True,
        arrowhead=2,
        arrowsize=1,
        arrowwidth=1.5,
        arrowcolor="rgba(248,250,252,0.85)",
        ax=0,
        ay=-28,
        font=dict(color="#F8FAFC", size=12, family="Segoe UI"),
        bgcolor="rgba(0,0,0,0.55)",
        borderpad=4,
        row=2,
        col=1,
    )
    fig.add_annotation(
        x=(x0 + x1) / 2.0,
        y=(y0 + y1) / 2.0,
        text=f"Mejor zona<br>ROP {_brop:.1f}",
        showarrow=False,
        font=dict(color="#ffffff", size=12, family="Segoe UI"),
        bgcolor="rgba(0,0,0,0.55)",
        borderpad=5,
        row=2,
        col=1,
    )

    # (2,2) Marginal derecha: ROP en X, RPM en Y (mismo eje Y que el heatmap vía matches)
    fig.add_trace(
        go.Scatter(
            x=max_r,
            y=y_c,
            mode="lines+markers",
            name="Máx ROP",
            line=dict(color="#F97316", width=3),
            marker=dict(size=7, line=dict(width=0)),
            legendgroup="right",
            showlegend=False,
            hovertemplate="RPM: %{y:.2f}<br>ROP máx. en bin: %{x:.2f}<extra></extra>",
        ),
        row=2,
        col=2,
    )
    fig.add_trace(
        go.Scatter(
            x=min_r,
            y=y_c,
            mode="lines+markers",
            name="Mín ROP",
            line=dict(color="#38BDF8", width=2.5, dash="dot"),
            marker=dict(size=6, line=dict(width=0)),
            legendgroup="right",
            showlegend=False,
            hovertemplate="RPM: %{y:.2f}<br>ROP mín. en bin: %{x:.2f}<extra></extra>",
        ),
        row=2,
        col=2,
    )

    xr = float(np.nanmin(x_c)), float(np.nanmax(x_c))
    yr = float(np.nanmin(y_c)), float(np.nanmax(y_c))
    zmax = float(np.nanmax(stat)) if np.isfinite(np.nanmax(stat)) else 1.0
    zmin_marg = float(np.nanmin(np.concatenate([max_r, min_r, max_w, min_w]))) if max_r.size else 0.0
    zmax_marg = float(np.nanmax(np.concatenate([max_r, min_r, max_w, min_w]))) if max_r.size else zmax
    pad = max(1e-6, (zmax_marg - zmin_marg) * 0.06)
    zr_marg = (zmin_marg - pad, zmax_marg + pad)

    fig.update_xaxes(
        title_text="WOB",
        title_font=dict(size=14, color=_fg),
        range=list(xr),
        showgrid=False,
        tickangle=-35,
        tickfont=dict(size=12, color=_tick),
        row=2,
        col=1,
    )
    fig.update_yaxes(
        title_text="RPM",
        title_font=dict(size=14, color=_fg),
        range=list(yr),
        showgrid=False,
        tickfont=dict(size=12, color=_tick),
        row=2,
        col=1,
    )
    fig.update_xaxes(
        title_text="",
        showticklabels=False,
        range=list(xr),
        showgrid=True,
        gridcolor=_grid_marg,
        row=1,
        col=1,
    )
    fig.update_yaxes(
        title_text="ROP (máx/mín por WOB)",
        title_font=dict(size=12, color=_fg),
        range=list(zr_marg),
        showgrid=True,
        gridcolor=_grid_marg,
        tickfont=dict(size=11, color=_tick),
        row=1,
        col=1,
    )
    fig.update_xaxes(
        title_text="ROP (máx/mín por RPM)",
        title_font=dict(size=12, color=_fg),
        range=list(zr_marg),
        showgrid=True,
        gridcolor=_grid_marg,
        tickfont=dict(size=11, color=_tick),
        row=2,
        col=2,
    )
    fig.update_yaxes(
        title_text="",
        range=list(yr),
        showticklabels=True,
        showgrid=False,
        tickfont=dict(size=12, color=_tick),
        row=2,
        col=2,
    )

    fig.update_layout(
        template="plotly_dark",
        width=_dash_w,
        height=_dash_h,
        paper_bgcolor=_dash_bg,
        plot_bgcolor=_dash_bg,
        title=dict(
            text=title,
            x=0.02,
            xanchor="left",
            font=dict(size=18, color="#F1F5F9"),
        ),
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.02,
            xanchor="right",
            x=1,
            font=dict(size=13, color=_fg),
            bgcolor="rgba(11,13,20,0.92)",
            bordercolor="rgba(255,255,255,0.06)",
            borderwidth=1,
        ),
        margin=dict(l=72, r=64, t=92, b=72),
        font=dict(family="Segoe UI", size=13, color=_fg),
        hoverlabel=dict(
            bgcolor="rgba(15,23,42,0.96)",
            bordercolor="rgba(148,163,184,0.35)",
            font_size=14,
            font_family="Segoe UI",
        ),
    )

    # plotly_dark puede dejar ejes con otro tono; reforzar fondo único en todos los paneles
    fig.update_xaxes(
        showline=False,
        zeroline=False,
        linecolor="rgba(255,255,255,0.12)",
    )
    fig.update_yaxes(
        showline=False,
        zeroline=False,
        linecolor="rgba(255,255,255,0.12)",
    )

    return fig


def build_corr_heatmap(df_run: pd.DataFrame, run_name: str):
    corr_cols = ["ROP", "WOB", "RPM", "DLS", "Inclination", "Azimuth", "Distance"]
    corr_cols = [c for c in corr_cols if c in df_run.columns]
    corr_df = df_run[corr_cols].corr()
    if corr_df.isna().all().all():
        return None

    # Porcentaje para texto; NaN no se puede convertir a int → texto manual
    corr_pct = (corr_df * 100).round(0)
    text_arr = np.where(
        np.isnan(corr_pct.values),
        "",
        (np.nan_to_num(corr_pct.values, nan=0.0).astype(int)).astype(str) + "%",
    )
    fig = px.imshow(
        corr_df,
        text_auto=".2f",
        color_continuous_scale="RdBu",
        zmin=-1,
        zmax=1,
        title=f"{run_name} – Correlation Heatmap",
    )
    fig.update_traces(
        text=text_arr,
        texttemplate="%{text}",
        textfont=dict(size=11),
        xgap=1,
        ygap=1,
    )
    fig.update_layout(coloraxis_colorbar=dict(title="Corr (-1 a 1)"))
    return fig


def build_dls_vs_md_figure(df_run: pd.DataFrame, run_name: str):
    df_plot = df_run.dropna(subset=["Survey MD", "DLS"])
    if df_plot.empty:
        return None

    fig = px.line(
        df_plot,
        x="Survey MD",
        y="DLS",
        title=f"{run_name} – DLS vs MD",
        labels={"Survey MD": "Measured Depth (m)", "DLS": "DLS (°/30m)"},
    )
    fig.update_traces(mode="lines+markers", marker=dict(size=4, opacity=0.75))
    return fig


def compute_effective_time_hours(df_run: pd.DataFrame) -> float:
    if "Start" not in df_run.columns or "End" not in df_run.columns:
        return 0.0
    df_time = df_run.dropna(subset=["Start", "End"]).copy()
    if df_time.empty:
        return 0.0
    df_time["Duration_h"] = (df_time["End"] - df_time["Start"]).dt.total_seconds() / 3600.0
    df_time = df_time[df_time["Duration_h"] > 0]
    effective_mask = (df_time["Distance"] > 0) & (
        df_time["Mode_norm"].isin(["ROTARY", "SLIDE"])
    )
    return float(df_time.loc[effective_mask, "Duration_h"].sum())


def build_boxplot_by_mode(df_run: pd.DataFrame, col: str, title: str):
    df_plot = df_run.dropna(subset=[col, "Mode_norm"])
    if df_plot.empty:
        return None
    fig = px.box(
        df_plot,
        x="Mode_norm",
        y=col,
        points="outliers",
        title=title,
        labels={"Mode_norm": "Mode", col: col},
        color="Mode_norm",
    )
    fig.update_traces(marker=dict(size=4, opacity=0.65))
    return fig


def build_control_chart(df_run: pd.DataFrame, col: str, run_name: str):
    df_plot = df_run.dropna(subset=[col, "Depth_X"])
    if df_plot.empty:
        return None
    mean_val = df_plot[col].mean()
    std_val = df_plot[col].std()
    upper = mean_val + 3 * std_val
    lower = mean_val - 3 * std_val

    fig = go.Figure()
    fig.add_trace(
        go.Scatter(
            x=df_plot["Depth_X"],
            y=df_plot[col],
            mode="markers",
            name=col,
            marker=dict(size=5, color="#2563EB", opacity=0.7),
        )
    )
    fig.add_hline(y=mean_val, line_dash="dash", line_color="#10B981")
    fig.add_hline(y=upper, line_dash="dot", line_color="#EF4444")
    fig.add_hline(y=lower, line_dash="dot", line_color="#EF4444")
    fig.update_layout(
        title=f"{run_name} – Control chart {col}",
        xaxis_title="Depth (m)",
        yaxis_title=col,
    )
    return fig


def build_crossplot_dls_rop(df_run: pd.DataFrame, run_name: str):
    df_plot = df_run.dropna(subset=["DLS", "ROP", "WOB", "RPM"])
    if df_plot.empty:
        return None
    fig = px.scatter(
        df_plot,
        x="DLS",
        y="ROP",
        color="RPM",
        size="WOB",
        title=f"{run_name} – DLS vs ROP (color RPM, size WOB)",
        labels={"DLS": "DLS (°/30m)", "ROP": "ROP (m/hr)", "RPM": "RPM", "WOB": "WOB"},
        color_continuous_scale="Turbo",
    )
    fig.update_traces(marker=dict(opacity=0.75))
    return fig


def build_scatter_matrix(df_run: pd.DataFrame, run_name: str):
    cols = ["ROP", "WOB", "RPM", "DLS"]
    cols = [c for c in cols if c in df_run.columns]
    df_plot = df_run[cols].dropna()
    if df_plot.empty or len(cols) < 2:
        return None
    fig = px.scatter_matrix(
        df_plot,
        dimensions=cols,
        title=f"{run_name} – Scatter Matrix",
        color="ROP" if "ROP" in cols else None,
        color_continuous_scale="Turbo",
    )
    fig.update_traces(diagonal_visible=False, marker=dict(size=4, opacity=0.7))
    return fig


def build_depth_heatmap(df_run: pd.DataFrame, run_name: str):
    metrics = ["ROP", "WOB", "RPM"]
    df_plot = df_run.dropna(subset=["Depth_X"]).copy()
    if df_plot.empty:
        return None
    try:
        df_plot["DepthBin"] = pd.cut(df_plot["Depth_X"], bins=20)
    except Exception:
        return None
    agg = df_plot.groupby("DepthBin")[metrics].mean().T
    if agg.empty:
        return None
    agg.columns = agg.columns.astype(str)
    fig = px.imshow(
        agg,
        aspect="auto",
        color_continuous_scale="Turbo",
        title=f"{run_name} – Heatmap por profundidad (binning)",
        labels={"x": "Depth Bin", "y": "Metric", "color": "Mean"},
    )
    fig.update_traces(xgap=1, ygap=1)
    return fig


def build_cumulative_meters(df_run: pd.DataFrame, run_name: str):
    if "End" not in df_run.columns:
        return None
    df_plot = df_run.dropna(subset=["End", "Distance"]).copy()
    if df_plot.empty:
        return None
    df_plot = df_plot.sort_values("End")
    df_plot["Meters_cum"] = df_plot["Distance"].fillna(0).cumsum()
    fig = px.line(
        df_plot,
        x="End",
        y="Meters_cum",
        title=f"{run_name} – Metros acumulados vs tiempo",
        labels={"End": "Time", "Meters_cum": "Meters"},
    )
    return fig


def compute_mse(
    df_run: pd.DataFrame,
    wob_col: str,
    rpm_col: str,
    rop_col: str,
    torque_col: str,
    bit_diameter_in: float,
    depth_col: str | None = None,
):
    required_cols = [wob_col, rpm_col, rop_col, torque_col]
    if bit_diameter_in <= 0 or any(c not in df_run.columns for c in required_cols):
        return None
    df_plot = df_run.copy()
    for col in required_cols:
        df_plot[col] = pd.to_numeric(df_plot[col], errors="coerce")
    subset_cols = required_cols[:]
    if depth_col and depth_col in df_plot.columns:
        subset_cols.append(depth_col)
    df_plot = df_plot.dropna(subset=subset_cols)
    if df_plot.empty:
        return None

    wob_n = df_plot[wob_col] * 9.80665  # kgf -> N
    torque_nm = df_plot[torque_col]
    rpm = df_plot[rpm_col]
    rop_mpm = df_plot[rop_col] / 60.0  # m/hr -> m/min
    bit_diameter_m = bit_diameter_in * 0.0254
    area_m2 = np.pi * (bit_diameter_m / 2.0) ** 2

    valid = (rop_mpm > 0) & (area_m2 > 0)
    df_plot = df_plot[valid].copy()
    if df_plot.empty:
        return None

    mse_pa = (wob_n.loc[df_plot.index] / area_m2) + (
        2 * np.pi * torque_nm.loc[df_plot.index] * rpm.loc[df_plot.index]
    ) / (area_m2 * rop_mpm.loc[df_plot.index])
    df_plot["MSE_MPa"] = mse_pa / 1_000_000.0
    df_plot["MSE_ksi"] = df_plot["MSE_MPa"] * MPA_TO_KSI
    if depth_col and depth_col in df_plot.columns:
        df_plot["Depth_MSE"] = df_plot[depth_col]
    else:
        df_plot["Depth_MSE"] = np.arange(1, len(df_plot) + 1)
    return df_plot


def build_mse_vs_depth(df_plot: pd.DataFrame, run_name: str, x_col: str, x_label: str):
    if "MSE_ksi" not in df_plot.columns and "MSE_MPa" not in df_plot.columns:
        return None
    y_col = "MSE_ksi" if "MSE_ksi" in df_plot.columns else "MSE_MPa"
    y_label = "MSE (ksi)" if y_col == "MSE_ksi" else "MSE (MPa)"
    fig = px.line(
        df_plot,
        x=x_col,
        y=y_col,
        title=f"{run_name} – MSE vs Profundidad",
        labels={x_col: x_label, y_col: y_label},
        color_discrete_sequence=["#0ea5e9"],
    ).update_layout(xaxis_autorange="reversed")
    fig.update_traces(mode="lines", line=dict(width=2.2))
    return fig


def build_mse_hist(df_plot: pd.DataFrame, run_name: str):
    if "MSE_ksi" not in df_plot.columns and "MSE_MPa" not in df_plot.columns:
        return None
    x_col = "MSE_ksi" if "MSE_ksi" in df_plot.columns else "MSE_MPa"
    x_label = "MSE (ksi)" if x_col == "MSE_ksi" else "MSE (MPa)"
    vals = df_plot[x_col].dropna()
    if vals.empty:
        return None
    counts, bin_edges = np.histogram(vals, bins=30)
    bin_centers = (bin_edges[:-1] + bin_edges[1:]) / 2
    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=bin_centers,
            y=counts,
            name="Frecuencia",
            marker_color="#0ea5e9",
            marker_line_width=0,
            opacity=0.85,
        )
    )
    fig.add_trace(
        go.Scatter(
            x=bin_centers,
            y=counts,
            mode="lines+markers",
            name="Tendencia",
            line=dict(color="#c2410c", width=2.5, dash="solid"),
            marker=dict(size=6, color="#c2410c", symbol="circle"),
        )
    )
    fig.update_layout(
        title=f"{run_name} – Distribución de MSE",
        xaxis_title=x_label,
        yaxis_title="count",
        template=PLOTLY_TEMPLATE,
        barmode="overlay",
        showlegend=True,
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
    )
    fig.update_xaxes(showgrid=True, gridcolor="rgba(0,0,0,0.06)")
    fig.update_yaxes(showgrid=True, gridcolor="rgba(0,0,0,0.06)")
    return fig


def build_hist_with_trend(
    values,
    title: str,
    x_label: str,
    nbins: int = 30,
    bar_color: str = "#0ea5e9",
    line_color: str = "#c2410c",
) -> go.Figure:
    """Histograma con línea de tendencia que une los tops de las barras."""
    vals = pd.Series(values).dropna()
    if vals.empty:
        return go.Figure()
    counts, bin_edges = np.histogram(vals, bins=nbins)
    bin_centers = (bin_edges[:-1] + bin_edges[1:]) / 2
    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=bin_centers,
            y=counts,
            name="Frecuencia",
            marker_color=bar_color,
            marker_line_width=0,
            opacity=0.85,
        )
    )
    fig.add_trace(
        go.Scatter(
            x=bin_centers,
            y=counts,
            mode="lines+markers",
            name="Tendencia",
            line=dict(color=line_color, width=2.5, dash="solid"),
            marker=dict(size=6, color=line_color, symbol="circle"),
        )
    )
    fig.update_layout(
        title=title,
        xaxis_title=x_label,
        yaxis_title="count",
        template=PLOTLY_TEMPLATE,
        barmode="overlay",
        showlegend=True,
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
    )
    fig.update_xaxes(showgrid=True, gridcolor="rgba(0,0,0,0.06)")
    fig.update_yaxes(showgrid=True, gridcolor="rgba(0,0,0,0.06)")
    return fig


def format_num(val: float | int | None, digits: int = 2) -> str:
    if val is None or pd.isna(val):
        return "—"
    return f"{val:.{digits}f}"


TRENDLINE_COLOR = "#c2410c"


def set_trendline_color(fig, color: str = TRENDLINE_COLOR):
    """Pone en color la línea de tendencia OLS en figuras px.scatter(..., trendline='ols')."""
    for trace in fig.data:
        mode = getattr(trace, "mode", "")
        name = (trace.name or "").lower()
        if mode == "lines" or "trend" in name or "ols" in name or "fit" in name:
            if hasattr(trace, "line") and trace.line is not None:
                trace.line.color = color
    return fig


def safe_corr(df: pd.DataFrame, x: str, y: str) -> str:
    df_plot = df.dropna(subset=[x, y])
    if len(df_plot) < 2:
        return "n/a"
    return format_num(df_plot[x].corr(df_plot[y]), 2)


def series_summary(series: pd.Series) -> str:
    return (
        f"min {format_num(series.min())}, max {format_num(series.max())}, "
        f"avg {format_num(series.mean())}"
    )


def chart_notes(summary: str, legend: str) -> None:
    st.caption(f"**Resumen:** {summary}")
    st.caption(f"**Leyenda:** {legend}")


def heatmap_legend_text(scope: str) -> str:
    return (
        f"**Cómo leer el heatmap ({scope}):** "
        "Rojo = correlación positiva, azul = negativa, blanco ≈ 0. "
        "Valores de -1 a +1: +1 relación directa perfecta, -1 inversa perfecta."
    )


def summarize_heatmap(corr_df: pd.DataFrame) -> str:
    corr_abs = corr_df.abs().copy()
    np.fill_diagonal(corr_abs.values, 0)
    if (corr_abs.values == 0).all():
        return "No se observa correlación lineal fuerte entre variables."

    max_idx = np.unravel_index(np.argmax(corr_abs.values), corr_abs.shape)
    var_a = corr_abs.index[max_idx[0]]
    var_b = corr_abs.columns[max_idx[1]]
    strength = corr_df.loc[var_a, var_b]
    direction = "positiva" if strength >= 0 else "negativa"
    return (
        f"Correlación más fuerte entre **{var_a}** y **{var_b}**: "
        f"{strength:.2f} ({direction})."
    )


def summarize_heatmap_engineering_pct(corr_df: pd.DataFrame) -> str:
    """
    Resumen de la matriz de correlación en porcentaje y con comentarios,
    para el heatmap de ingeniería (WOB, RPM, Torque, Freq_Hz, Proximity_norm).
    """
    if corr_df is None or corr_df.empty:
        return "No hay datos de correlación."
    corr_abs = corr_df.abs().copy()
    np.fill_diagonal(corr_abs.values, 0)
    # Pares (|r|, r, var_a, var_b) sin diagonal, ordenados por |r| desc
    pairs: list[tuple[float, float, str, str]] = []
    for i in range(len(corr_df.index)):
        for j in range(len(corr_df.columns)):
            if i >= j:
                continue
            a, b = corr_df.index[i], corr_df.columns[j]
            r = float(corr_df.iloc[i, j])
            if pd.isna(r):
                continue
            pairs.append((abs(r), r, a, b))
    pairs.sort(key=lambda x: -x[0])

    def pct(r: float) -> str:
        return f"{int(round(r * 100))}%"

    lines: list[str] = []
    lines.append("**Resumen en % (fuerza de relación lineal)**")
    lines.append("")
    # Fuertes positivas
    strong_pos = [(r, a, b) for abs_r, r, a, b in pairs if r >= 0.4]
    if strong_pos:
        lines.append("*Correlaciones positivas fuertes:*")
        for r, a, b in strong_pos[:5]:
            comment = ""
            if "Freq_Hz" in (a, b) and ("RPM" in (a, b) or "Surface RPM" in (a, b)):
                comment = " — esperado (Freq_Hz = RPM/60)."
            elif "Torque" in (a, b) and ("RPM" in (a, b) or "Freq" in (a, b)):
                comment = " — torque y velocidad rotacional suelen moverse juntos."
            lines.append(f"  • **{a}** ↔ **{b}**: {pct(r)}{comment}")
        lines.append("")
    # Fuertes negativas
    strong_neg = [(r, a, b) for abs_r, r, a, b in pairs if r <= -0.2]
    if strong_neg:
        lines.append("*Correlaciones negativas:*")
        for r, a, b in strong_neg[:3]:
            lines.append(f"  • **{a}** ↔ **{b}**: {pct(r)}")
        lines.append("")
    # Débiles / cercanas a cero
    weak = [(r, a, b) for abs_r, r, a, b in pairs if -0.2 < r < 0.2 and abs_r > 1e-6]
    if weak:
        lines.append("*Sin correlación lineal apreciable (≈0%):*")
        for r, a, b in weak[:4]:
            lines.append(f"  • **{a}** ↔ **{b}**: {pct(r)}")
        lines.append("")
    # Comentario final
    if pairs:
        _, r_max, a_max, b_max = pairs[0]
        lines.append("**Conclusión:** La relación más marcada es **" + str(a_max) + "** con **" + str(b_max) + "** (" + pct(r_max) + "). "
                     "Valores cercanos a 0% indican que las variables no varían de forma lineal entre sí.")
    return "\n".join(lines)


def analyze_bha_resonance(
    df: pd.DataFrame,
    wob_col: str,
    rpm_col: str,
    torque_col: str | None = None,
    max_modes: int = 8,
    tolerance_hz: float | None = None,
):
    df = df.copy()
    cols_to_num = [wob_col, rpm_col]
    if torque_col:
        cols_to_num.append(torque_col)
    for c in cols_to_num:
        df[c] = pd.to_numeric(df[c], errors="coerce")

    df = df.dropna(subset=[wob_col, rpm_col])
    df = df[(df[wob_col] > 0) & (df[rpm_col] > 0)]
    if df.empty:
        return None

    df["Freq_Hz"] = df[rpm_col] / 60.0

    hist, bin_edges = np.histogram(df["Freq_Hz"], bins=200)
    bin_centers = (bin_edges[:-1] + bin_edges[1:]) / 2
    peak_height = max(5, np.percentile(hist, 90))
    peaks, props = find_peaks(hist, height=peak_height)
    if len(peaks) == 0:
        peak_height = np.percentile(hist, 80)
        peaks, props = find_peaks(hist, height=peak_height)

    if len(peaks) > 0:
        order = np.argsort(props["peak_heights"])[::-1]
        peaks = peaks[order][:max_modes]
        natural_modes = np.sort(bin_centers[peaks])
    else:
        natural_modes = np.array([])

    if tolerance_hz is not None and tolerance_hz >= 0:
        tolerance = float(tolerance_hz)
    elif len(natural_modes) >= 2:
        freq_spacing = np.median(np.diff(natural_modes))
        tolerance = min(0.20, freq_spacing / 4)
    else:
        tolerance = 0.10  # más estrecha por defecto → más ventanas seguras visibles

    resonant_bands = [((f - tolerance) * 60, (f + tolerance) * 60) for f in natural_modes]

    def proximity_to_resonance(rpm: float) -> float:
        if not resonant_bands:
            return 1.0
        distances = []
        for band_min, band_max in resonant_bands:
            if band_min <= rpm <= band_max:
                return 0.0
            distances.append(min(abs(rpm - band_min), abs(rpm - band_max)))
        return min(distances)

    df["Proximity"] = df[rpm_col].apply(proximity_to_resonance)
    max_distance = max(df["Proximity"].max(), 1e-6)
    df["Proximity_norm"] = df["Proximity"] / max_distance

    rpm_total_min, rpm_total_max = df[rpm_col].min(), df[rpm_col].max()
    safe_segments = []
    current_min = rpm_total_min
    for band_min, band_max in sorted(resonant_bands):
        if band_min > current_min:
            safe_segments.append((current_min, band_min))
        current_min = max(current_min, band_max)
    if current_min < rpm_total_max:
        safe_segments.append((current_min, rpm_total_max))

    report = None
    if safe_segments:
        report = pd.DataFrame(
            {
                "Window #": np.arange(1, len(safe_segments) + 1),
                "RPM Min": [a for a, b in safe_segments],
                "RPM Max": [b for a, b in safe_segments],
                "Width (RPM)": [b - a for a, b in safe_segments],
                "Freq Min (Hz)": [a / 60.0 for a, b in safe_segments],
                "Freq Max (Hz)": [b / 60.0 for a, b in safe_segments],
                "Center RPM": [(a + b) / 2.0 for a, b in safe_segments],
                "Center Hz": [((a + b) / 2.0) / 60.0 for a, b in safe_segments],
            }
        )

    return {
        "df": df,
        "natural_modes": natural_modes,
        "tolerance": tolerance,
        "resonant_bands": resonant_bands,
        "safe_segments": safe_segments,
        "report": report,
    }


def build_proximity_figure(
    df: pd.DataFrame,
    wob_col: str,
    rpm_col: str,
    analysis,
    torque_col: str | None = None,
    xaxis_range: tuple[float, float] | None = None,
    yaxis_range: tuple[float, float] | None = None,
) -> go.Figure:
    fig = go.Figure()
    x_max = (xaxis_range[1] if xaxis_range else df[wob_col].max()) * 1.02

    for seg_min, seg_max in analysis["safe_segments"]:
        fig.add_hrect(
            y0=seg_min,
            y1=seg_max,
            fillcolor="rgba(16,185,129,0.08)",
            line_width=0,
            layer="below",
        )
    for f in analysis["natural_modes"]:
        fmin, fmax = (f - analysis["tolerance"]) * 60, (f + analysis["tolerance"]) * 60
        fig.add_hrect(
            y0=fmin,
            y1=fmax,
            fillcolor="rgba(239,68,68,0.18)",
            line_width=0,
            layer="below",
        )
        fig.add_annotation(
            x=x_max,
            y=(fmin + fmax) / 2,
            text=f"{f:.2f} Hz",
            showarrow=False,
            xanchor="left",
            font=dict(color="#7F1D1D", size=11),
        )

    customdata = None
    if torque_col and torque_col in df.columns:
        customdata = df[[torque_col]].to_numpy()
    marker_size = 4 if len(df) > 600 else 6

    fig.add_trace(
        go.Scatter(
            x=df[wob_col],
            y=df[rpm_col],
            mode="markers",
            customdata=customdata,
            marker=dict(
                size=marker_size,
                color=df["Proximity_norm"],
                colorscale="RdYlGn",
                cmin=0,
                cmax=1,
                opacity=0.9,
                line=dict(width=0),
                colorbar=dict(title="Proximity (0= peligro, 1= seguro)"),
            ),
            hovertemplate=(
                "WOB: %{x:.2f}<br>RPM: %{y:.2f}"
                + ("<br>Torque: %{customdata[0]:.2f}" if customdata is not None else "")
                + "<extra></extra>"
            ),
        )
    )

    fig.update_layout(
        title="Operational Window & BHA Resonant Frequencies<br>Proximity to Resonance",
        xaxis_title="WOB [kgf]",
        yaxis_title="Surface RPM [rev/min]",
    )
    yr = yaxis_range if yaxis_range else [0, df[rpm_col].max() * 1.1]
    xr = xaxis_range if xaxis_range else None
    fig.update_yaxes(range=yr, showgrid=True, gridcolor="rgba(0,0,0,0.08)")
    if xr is not None:
        fig.update_xaxes(range=xr)
    return prettify(fig, h=520)


def build_frequency_figure(
    df: pd.DataFrame,
    wob_col: str,
    rpm_col: str,
    analysis,
    torque_col: str | None = None,
    xaxis_range: tuple[float, float] | None = None,
    yaxis_range: tuple[float, float] | None = None,
) -> go.Figure:
    fig = go.Figure()
    x_max = (xaxis_range[1] if xaxis_range else df[wob_col].max()) * 1.02

    for seg_min, seg_max in analysis["safe_segments"]:
        fig.add_hrect(
            y0=seg_min,
            y1=seg_max,
            fillcolor="rgba(16,185,129,0.08)",
            line_width=0,
            layer="below",
        )
    for f in analysis["natural_modes"]:
        fmin, fmax = (f - analysis["tolerance"]) * 60, (f + analysis["tolerance"]) * 60
        fig.add_hrect(
            y0=fmin,
            y1=fmax,
            fillcolor="rgba(239,68,68,0.18)",
            line_width=0,
            layer="below",
        )
        fig.add_annotation(
            x=x_max,
            y=(fmin + fmax) / 2,
            text=f"{f:.2f} Hz",
            showarrow=False,
            xanchor="left",
            font=dict(color="#7F1D1D", size=11),
        )

    customdata = None
    if torque_col and torque_col in df.columns:
        customdata = df[[torque_col]].to_numpy()
    marker_size = 4 if len(df) > 600 else 6

    fig.add_trace(
        go.Scatter(
            x=df[wob_col],
            y=df[rpm_col],
            mode="markers",
            customdata=customdata,
            marker=dict(
                size=marker_size,
                color=df["Freq_Hz"],
                colorscale="Turbo",
                opacity=0.85,
                line=dict(width=0),
                colorbar=dict(title="Rotational Frequency [Hz]"),
            ),
            hovertemplate=(
                "WOB: %{x:.2f}<br>RPM: %{y:.2f}<br>Freq: %{marker.color:.2f} Hz"
                + ("<br>Torque: %{customdata[0]:.2f}" if customdata is not None else "")
                + "<extra></extra>"
            ),
        )
    )

    fig.update_layout(
        title="Operational Window & Rotational Frequency Mapping (Resonance Zones)",
        xaxis_title="WOB [kgf]",
        yaxis_title="Surface RPM [rev/min]",
    )
    yr = yaxis_range if yaxis_range else [0, df[rpm_col].max() * 1.1]
    xr = xaxis_range if xaxis_range else None
    fig.update_yaxes(range=yr, showgrid=True, gridcolor="rgba(0,0,0,0.08)")
    if xr is not None:
        fig.update_xaxes(range=xr)
    return prettify(fig, h=520)


def find_shocks_vibs_column(df: pd.DataFrame) -> str | None:
    """Devuelve la primera columna que coincida con Shocks/Vibs (case-insensitive)."""
    cols_lower = {c.lower(): c for c in df.columns}
    for cand in SHOCKS_VIB_COL_CANDIDATES:
        if cand.lower() in cols_lower:
            return cols_lower[cand.lower()]
    return None


def build_wob_rpm_binned_heatmap(
    df: pd.DataFrame,
    wob_col: str,
    rpm_col: str,
    value_col: str,
    value_label: str,
    n_bins: int = 30,
    xaxis_range: tuple[float, float] | None = None,
    yaxis_range: tuple[float, float] | None = None,
    title: str | None = None,
) -> go.Figure | None:
    """
    Heatmap WOB vs RPM con bins cuadrados (n_bins x n_bins).
    Cada celda se colorea según la media de value_col en ese bin.
    """
    df_plot = df[[wob_col, rpm_col, value_col]].dropna()
    if df_plot.empty:
        return None
    x = df_plot[wob_col].to_numpy()
    y = df_plot[rpm_col].to_numpy()
    v = df_plot[value_col].to_numpy()
    # Rango según datos + margen 5% para que el heatmap se vea amplio (no forzar 0–30k)
    if xaxis_range:
        x_min, x_max = xaxis_range[0], xaxis_range[1]
    else:
        x_min_d, x_max_d = float(x.min()), float(x.max())
        span_x = max(x_max_d - x_min_d, 1.0)
        x_min = x_min_d - 0.05 * span_x
        x_max = x_max_d + 0.05 * span_x
    if yaxis_range:
        y_min, y_max = yaxis_range[0], yaxis_range[1]
    else:
        y_min_d, y_max_d = float(y.min()), float(y.max())
        span_y = max(y_max_d - y_min_d, 1.0)
        y_min = y_min_d - 0.05 * span_y
        y_max = y_max_d + 0.05 * span_y
    stat, x_edges, y_edges, _ = binned_statistic_2d(
        x, y, v, statistic="mean", bins=n_bins, range=[[x_min, x_max], [y_min, y_max]]
    )
    # Celdas vacías (sin datos) -> NaN para que el heatmap las muestre vacías o con color neutro
    x_centers = (x_edges[:-1] + x_edges[1:]) / 2
    y_centers = (y_edges[:-1] + y_edges[1:]) / 2
    # Plotly Heatmap: z en filas = y, columnas = x; xgap/ygap = borde visible entre bins
    fig = go.Figure(
        data=go.Heatmap(
            x=x_centers,
            y=y_centers,
            z=stat.T,
            colorscale="Viridis",
            colorbar=dict(title=value_label),
            xgap=1,
            ygap=1,
            hovertemplate="WOB: %{x:.0f}<br>RPM: %{y:.1f}<br>" + value_label + ": %{z:.2f}<extra></extra>",
        )
    )
    fig.update_layout(
        title=title or f"WOB–RPM Heatmap (squared bins) · {value_label}",
        xaxis_title="WOB [kgf]",
        yaxis_title="Surface RPM [rev/min]",
    )
    # Fijar rango de ejes para que se vea amplio (con el margen ya aplicado en x_min/x_max/y_min/y_max)
    fig.update_xaxes(range=[x_min, x_max])
    fig.update_yaxes(range=[y_min, y_max])
    return prettify(fig, h=520)


def render_bha_module() -> None:
    st.subheader(tr("bha_subheader"))
    st.caption(tr("bha_caption"))

    data_source = st.radio(
        tr("data_source"),
        ["CSV", "API"],
        horizontal=True,
        format_func=lambda x: tr(f"src_{x.lower()}"),
    )
    bha_df = None
    bha_wob_col = "WOB"
    bha_rpm_col = "Surface RPM"
    bha_torque_col = "Torque"
    bha_rop_col = "ROP"
    bha_depth_col = "Depth"

    if data_source == "API":
        with st.expander(tr("bha_api_cfg"), expanded=False):
            # Carga automática desde .env (SOLO_BASE_URL / SOLO_ACCESS_TOKEN) y session_state
            base_url, token = render_solo_connection_ui(prefix="bha_api", label=tr("solo_expander"))
            if not token:
                st.warning(tr("bha_missing_token"))
            traces_list_path = st.text_input(
                "Ruta listar trazas (opcional)",
                value="",
                help="Usa {well_uuid}. Ej: /public/api/v1/wells/{well_uuid}/drilling-traces",
                key="bha_api_traces_list_path",
            )
            traces_list_params = st.text_input(
                "Params listar trazas (opcional)",
                value="",
                help="Formato: key=value&key2=value2 o JSON",
                key="bha_api_traces_list_params",
            )
            laterals_list_path = st.text_input(
                "Ruta listar laterales (opcional)",
                value="",
                help="Usa {well_uuid} o {project_uuid}",
                key="bha_api_laterals_list_path",
            )
            laterals_list_params = st.text_input(
                "Params listar laterales (opcional)",
                value="",
                help="Formato: key=value&key2=value2 o JSON",
                key="bha_api_laterals_list_params",
            )
            trace_types_path = st.text_input(
                "Ruta listar tipos de traza (opcional)",
                value="",
                help="Ej: /public/api/v1/traces",
                key="bha_api_trace_types_path",
            )
            trace_detail_path = st.text_input(
                "Ruta detalle traza (opcional)",
                value="",
                help="Usa {trace_uuid}. Ej: /public/api/v1/drilling-traces/{trace_uuid}",
                key="bha_api_trace_detail_path",
            )
            mapped_trace_path = st.text_input(
                "Ruta datos traza por lateral (opcional)",
                value="",
                help="Usa {lateral_uuid} y {trace_uuid}",
                key="bha_api_mapped_trace_path",
            )
            st.session_state["api_token"] = token
            st.caption(tr("api_note_html"))

        if not token:
            st.info(tr("enter_token"))
            return

        try:
            projects_resp = api_list_projects(base_url, token)
            projects = normalize_list_response(projects_resp)
        except Exception as e:
            st.error(f"{tr('list_projects_err')} {e}")
            return

        if not projects:
            st.info(tr("no_projects"))
            return

        project_map = {
            f"{p.get('name', tr('unnamed'))} ({p.get('uuid', 'n/a')})": p for p in projects
        }
        project_label = st.selectbox(tr("project"), list(project_map.keys()), key="bha_project_label_api")
        project_uuid = project_map[project_label].get("uuid")

        if not project_uuid:
            st.error(tr("project_no_uuid"))
            return

        try:
            wells_resp = api_list_wells(base_url, token, project_uuid)
            wells = normalize_list_response(wells_resp)
        except Exception as e:
            st.error(f"{tr('list_wells_err')} {e}")
            return

        if not wells:
            st.info(tr("no_wells"))
            return

        well_map = {
            f"{w.get('name', tr('unnamed'))} ({w.get('uuid', 'n/a')})": w for w in wells
        }
        well_label = st.selectbox(tr("well"), list(well_map.keys()), key="bha_well_label_api")
        well_uuid = well_map[well_label].get("uuid")

        if not well_uuid:
            st.error(tr("well_no_uuid"))
            return

        laterals = []
        no_laterals_found = False
        try:
            extra_lateral_params = parse_params_input(laterals_list_params)
            laterals_resp = api_list_laterals(
                base_url,
                token,
                project_uuid=project_uuid,
                well_uuid=well_uuid,
                custom_path=laterals_list_path or None,
                extra_params=extra_lateral_params,
            )
            laterals = normalize_list_response(laterals_resp)
        except Exception:
            laterals = []

        lateral_uuid = None
        if laterals:
            lateral_map = {
                f"{l.get('name', 'Sin nombre')} ({l.get('uuid', 'n/a')})": l for l in laterals
            }
            lateral_label = st.selectbox("Lateral", list(lateral_map.keys()), key="bha_lateral_label")
            lateral_uuid = lateral_map[lateral_label].get("uuid")
            if not lateral_uuid:
                st.error("La lateral seleccionada no tiene UUID.")
                return
        else:
            # Algunos pozos no tienen laterales listados; en SOLO el well_uuid puede usarse como wellbore/lateral para trazas.
            lateral_uuid = well_uuid
            no_laterals_found = True
            st.info("No se encontraron laterales para este pozo. Usaré el Pozo como wellbore para consultar trazas.")

        trace_type = st.radio(
            "Dominio de mapeo",
            options=["TIME", "DEPTH"],
            horizontal=True,
            help="Selecciona el dominio para filtrar las trazas mapeadas (time o depth).",
        )

        # ---------- Parámetros para datos de traza (UI unificada, sin duplicar size/page con texto) ----------
        st.markdown("**Parámetros para datos de traza**")
        with st.expander("Rango y paginación", expanded=True):
            use_custom_range = st.checkbox(
                "Especificar rango (from / to)",
                value=st.session_state.get("bha_use_custom_range", True),
                key="bha_use_custom_range",
                help="Si no marcas esto, la API puede usar sus valores por defecto. Con 'Auto-detectar rango' la app puede probar otros rangos si no hay datos.",
            )
            range_from = ""
            range_to = ""
            if use_custom_range:
                if trace_type.upper() == "TIME":
                    col_from, col_to = st.columns(2)
                    with col_from:
                        range_from = st.text_input(
                            "Desde (ISO 8601, ej. 2020-12-29T08:00:00Z)",
                            value=st.session_state.get("bha_trace_from", "2020-12-29T08:00:00Z"),
                            key="bha_trace_from",
                        )
                    with col_to:
                        range_to = st.text_input(
                            "Hasta (ISO 8601, ej. 2020-12-29T08:50:00Z)",
                            value=st.session_state.get("bha_trace_to", "2020-12-29T08:50:00Z"),
                            key="bha_trace_to",
                        )
                else:
                    col_from, col_to = st.columns(2)
                    with col_from:
                        range_from = str(st.number_input("Desde (profundidad m)", value=1.0, step=100.0, key="bha_trace_from_num"))
                    with col_to:
                        range_to = str(st.number_input("Hasta (profundidad m)", value=5000.0, step=100.0, key="bha_trace_to_num"))
            else:
                range_from = ""
                range_to = ""

            st.caption("Paginación. Para **más datos**: sube «Tamaño de página» y/o «Páginas a cargar» (concatena página 0, 1, 2…).")
            col_size, col_page, col_multi = st.columns(3)
            with col_size:
                page_size = st.number_input(
                    "Tamaño de página (size)",
                    min_value=100,
                    max_value=50000,
                    value=5000,
                    step=500,
                    key="bha_api_page_size",
                    help="Registros por página. Más = más datos por petición.",
                )
            with col_page:
                page_number = st.number_input(
                    "Página (page)",
                    min_value=0,
                    max_value=100000,
                    value=0,
                    step=1,
                    key="bha_api_page_number",
                    help="Solo se usa si «Páginas a cargar» = 1. 0 = primera página.",
                )
            with col_multi:
                num_pages_to_load = st.number_input(
                    "Páginas a cargar",
                    min_value=1,
                    max_value=20,
                    value=1,
                    step=1,
                    key="bha_num_pages_load",
                    help="1 = solo una página. 2 o más = se piden páginas 0,1,… y se unen (más datos).",
                )
            if page_number > 0 and num_pages_to_load == 1:
                st.warning(
                    "⚠️ Estás pidiendo la **página " + str(page_number) + "**. "
                    "Para más datos usa **Página = 0** y sube **Páginas a cargar** a 2 o 3."
                )

        with st.expander("Params adicionales (avanzado)", expanded=False):
            extra_params_raw = st.text_input(
                "Query string extra (opcional)",
                value=st.session_state.get("bha_extra_trace_params", ""),
                placeholder="Ej: sort=time&order=asc",
                help="Solo parámetros adicionales. No incluyas from, to, size ni page (ya se usan arriba). Formato: key=value&key2=value2",
                key="bha_extra_trace_params",
            )

        with st.expander("Opciones de API y auto-detección", expanded=True):
            force_data_endpoint = st.checkbox(
                "Forzar endpoint de datos por pozo (time/depth)",
                value=st.session_state.get("bha_force_data_endpoint", True),
                help="Usa solo /public/api/v1/wells/{well_uuid}/traces/{trace_uuid}/data/{time|depth}.",
                key="bha_force_data_endpoint",
            )
            trip_auto_probe_data = st.checkbox(
                "Auto-detectar rango con datos (si la traza viene vacía)",
                value=True,
                help="Si no hay datos, la app probará ventanas de tiempo hacia atrás y luego DEPTH. Solo se aplica si no especificaste 'Desde' y 'Hasta' arriba.",
                key="bha_trip_auto_probe_data",
            )
            if use_custom_range:
                trip_auto_probe_data = False  # FIX: never auto-probe when user specified from/to

        # Construcción única de params: rango + paginación + adicionales (sin duplicar)
        def build_trace_params(page_override: int | None = None) -> dict:
            p: dict = {}
            if use_custom_range and range_from.strip() and range_to.strip():
                p["from"] = range_from.strip()
                p["to"] = range_to.strip()
            p["size"] = int(page_size)
            p["page"] = page_override if page_override is not None else int(page_number)
            extra = parse_params_input(extra_params_raw)
            if extra:
                for k, v in extra.items():
                    if k not in ("from", "to", "size", "page", "limit", "offset"):
                        p[k] = v
            return p

        # Variables de alcance para trazas mapeadas (se usan más adelante)
        mapped_scope_uuid = lateral_uuid
        mapped_scope_kind = "lateral"
        use_well_mapping = False
        if no_laterals_found:
            use_well_mapping = st.checkbox(
                "Usar mapeo por pozo (si aplica)",
                value=True,
                help="Algunas cuentas exponen trazas mapeadas a nivel de pozo.",
            )
            if use_well_mapping:
                mapped_scope_uuid = well_uuid
                mapped_scope_kind = "well"

        def apply_paging_params(params: dict | None) -> dict:
            """Mantiene compatibilidad: si reciben params externos, se asegura size/page."""
            base = build_trace_params()
            if params:
                for k, v in params.items():
                    if k not in ("size", "page", "limit", "offset"):
                        base[k] = v
            return base



                # ------------------------------------------------------------
        # Listado de trazas (catálogo global):
        #   GET /public/api/v1/traces  -> devuelve UUID + name
        # Luego, los datos se consultan por pozo:
        #   /public/api/v1/wells/{well_uuid}/traces/{trace_uuid}/data/{time|depth}
        # ------------------------------------------------------------
        traces: list[dict] = []
        trace_source: str | None = None
        drilling_error = None
        mapped_error = None
        catalog_error = None
        extra_params = parse_params_input(traces_list_params)

        if mapped_scope_uuid:
            try:
                traces_resp = api_list_mapped_traces(
                    base_url=base_url,
                    token=token,
                    scope_uuid=mapped_scope_uuid,
                    trace_type=trace_type,
                    scope_kind=mapped_scope_kind or "lateral",
                    custom_path=None,
                    extra_params=extra_params,
                )
                traces = normalize_list_response(traces_resp)
                if traces:
                    trace_source = "mapped"
            except Exception as e:
                mapped_error = str(e)
                traces = []

        if not traces:
            try:
                traces_resp = api_list_drilling_traces(
                    base_url=base_url,
                    token=token,
                    well_uuid=well_uuid,
                    custom_path=traces_list_path or None,
                    extra_params=extra_params,
                )
                traces = normalize_list_response(traces_resp)
                if traces:
                    trace_source = "drilling"
            except Exception as e:
                drilling_error = str(e)
                traces = []

        if not traces:
            try:
                traces_resp = api_list_traces_catalog(
                    base_url,
                    token,
                    trace_types_path or None,
                    extra_params,
                )
                traces = normalize_list_response(traces_resp)
                if traces:
                    trace_source = "catalog"
            except Exception as e:
                catalog_error = str(e)
                traces = []
        if not traces:
            st.info(
                "No se encontraron trazas mapeadas, drilling-traces ni catálogo. "
                "Selecciona otra lateral o revisa permisos."
            )
            with st.expander("Detalles de error"):
                if "mapped_error" in locals():
                    st.write(f"Mapeadas: {mapped_error}")
                if "drilling_error" in locals():
                    st.write(f"Drilling: {drilling_error}")
                if "catalog_error" in locals():
                    st.write(f"Catálogo: {catalog_error}")
            return

        if not traces:
            st.info("No se encontraron trazas para este pozo.")
            return

        scope_key = f"{project_uuid}|{well_uuid}|{lateral_uuid}|{trace_type}|{trace_source}"
        prev_scope_key = st.session_state.get("bha_trace_scope")
        if scope_key != prev_scope_key:
            st.session_state["selected_trace_uuids"] = []
            st.session_state["bha_auto_select_applied"] = False
            st.session_state["bha_trace_scope"] = scope_key

        trace_map = {}
        trace_map_by_uuid = {}
        for t in traces:
            label = t.get("name") or t.get("label") or t.get("type") or t.get("uuid") or "trace"
            trace_map[label] = t
            trace_uuid = t.get("uuid") or t.get("id")
            if trace_uuid:
                trace_map_by_uuid[str(trace_uuid)] = {"label": label, "trace": t}

        trace_map_by_uuid_full = dict(trace_map_by_uuid)

        trace_filter = st.text_input(
            "Filtro de trazas (nombre)",
            value="",
            help="Escribe parte del nombre, ej: Hookload, RPM, Torque.",
            key="trace_filter",
        ).strip()
        if trace_filter:
            trace_map = {
                k: v for k, v in trace_map.items() if trace_filter.lower() in k.lower()
            }

        if trace_filter:
            allowed_labels = set(trace_map.keys())
            trace_map_by_uuid = {
                k: v for k, v in trace_map_by_uuid.items() if v["label"] in allowed_labels
            }

        if not trace_map_by_uuid:
            st.info("No hay trazas que coincidan con el filtro.")
            return

        def suggest_trace_uuids(
            trace_map: dict[str, dict],
        ) -> tuple[dict[str, str], list[str]]:
            required = {
                "WOB": ["wob", "weight on bit"],
                "RPM": ["rpm", "rotary speed"],
                "Torque": ["torque", "tq"],
                "ROP": ["rop", "rate of penetration"],
                "Depth": ["depth", "md", "measured depth", "survey md"],
            }
            suggestions: dict[str, str] = {}
            used: set[str] = set()
            for req, keywords in required.items():
                candidates: list[tuple[int, int, str, str]] = []
                for uuid, data in trace_map.items():
                    label = str(data.get("label", "")).lower()
                    for idx, kw in enumerate(keywords):
                        if kw in label:
                            candidates.append((idx, len(label), uuid, data.get("label", "")))
                            break
                if candidates:
                    candidates.sort()
                    for _, _, uuid, _ in candidates:
                        if uuid not in used:
                            suggestions[req] = uuid
                            used.add(uuid)
                            break
            missing = [req for req in required if req not in suggestions]
            return suggestions, missing

        st.session_state.pop("selected_traces", None)

        prev_filter = st.session_state.get("trace_filter_prev", "")
        if trace_filter != prev_filter:
            st.session_state["selected_trace_uuids"] = []
            st.session_state["bha_auto_select_applied"] = False
            st.session_state["trace_filter_prev"] = trace_filter

        prev_selected_uuids = st.session_state.get("selected_trace_uuids", [])
        prev_selected_uuids = [
            u for u in prev_selected_uuids if u in trace_map_by_uuid
        ]

        auto_select = st.checkbox(
            "Auto-seleccionar trazas sugeridas (WOB, RPM, Torque, ROP, Depth)",
            value=st.session_state.get("bha_auto_select", True),
            key="bha_auto_select",
        )
        suggestions, missing_required = suggest_trace_uuids(trace_map_by_uuid_full)
        suggested_labels = []
        if suggestions:
            for req, uuid in suggestions.items():
                label = trace_map_by_uuid_full.get(uuid, {}).get("label", uuid)
                suggested_labels.append(f"{req} → {label}")
        if auto_select and not st.session_state.get("bha_auto_select_applied", False):
            suggested_uuids = [
                u for u in suggestions.values() if u in trace_map_by_uuid
            ]
            if suggested_uuids:
                prev_selected_uuids = list(dict.fromkeys(prev_selected_uuids + suggested_uuids))
            st.session_state["bha_auto_select_applied"] = True

        with st.expander("Trazas sugeridas", expanded=False):
            st.write("Sugerencias para estadísticas y MSE.")
            if suggested_labels:
                st.write("Detectadas: " + "; ".join(suggested_labels))
            if missing_required:
                st.warning(
                    "No encontré estas trazas por nombre: "
                    + ", ".join(missing_required)
                )

        st.session_state["selected_trace_uuids"] = prev_selected_uuids

        trace_rows = [
            {
                "Seleccionar": uuid in prev_selected_uuids,
                "Traza": data["label"],
                "UUID": uuid,
            }
            for uuid, data in trace_map_by_uuid.items()
        ]
        traces_df = pd.DataFrame(trace_rows)
        traces_df = traces_df.sort_values("Traza", ignore_index=True)
        st.markdown("**Selecciona trazas a mapear**")
        traces_df = st.data_editor(
            traces_df,
            use_container_width=True,
            hide_index=True,
            key="trace_selector",
            column_config={
                "Seleccionar": st.column_config.CheckboxColumn(
                    "Seleccionar", help="Marca las trazas a mapear."
                )
            },
            disabled=["Traza", "UUID"],
        )
        selected_trace_uuids = traces_df.loc[
            traces_df["Seleccionar"], "UUID"
        ].tolist()
        st.session_state["selected_trace_uuids"] = selected_trace_uuids

        if not selected_trace_uuids:
            st.info("Selecciona al menos una traza.")
            return

        st.markdown("**Actualización en tiempo real (API)**")
        bha_auto_refresh = st.checkbox(
            "Actualizar automáticamente cada 30 s",
            value=st.session_state.get("bha_auto_refresh", False),
            key="bha_auto_refresh",
            help="Vuelve a consultar la API y refrescar gráficas cada X segundos.",
        )
        if bha_auto_refresh:
            st.number_input(
                "Intervalo (segundos)",
                min_value=10,
                max_value=300,
                value=30,
                step=10,
                key="bha_auto_refresh_interval",
                help="Cada cuántos segundos se vuelve a ejecutar (10–300 s).",
            )

        with st.expander("Debug API de trazas", expanded=False):
            st.write("Prueba un GET y revisa URL, status y respuesta.")
            if not selected_trace_uuids:
                st.info("Selecciona al menos una traza para habilitar la prueba.")
            else:
                trace_label_map = {
                    data["label"]: uuid for uuid, data in trace_map_by_uuid_full.items()
                }
                default_label = trace_map_by_uuid_full.get(
                    selected_trace_uuids[0], {}
                ).get("label")
                trace_label = st.selectbox(
                    "Traza para prueba",
                    list(trace_label_map.keys()),
                    index=list(trace_label_map.keys()).index(default_label)
                    if default_label in trace_label_map
                    else 0,
                    key="bha_debug_trace_label",
                )
                debug_trace_uuid = trace_label_map.get(trace_label)
                debug_params = build_trace_params()
                data_suffix = "time" if trace_type.upper() == "TIME" else "depth"
                debug_path = (
                    f"/public/api/v1/wells/{well_uuid}/traces/{debug_trace_uuid}/data/{data_suffix}"
                )
                debug_url = build_api_url(debug_path, base_url)
                debug_url_full = (
                    f"{debug_url}?{urlencode(debug_params)}" if debug_params else debug_url
                )
                st.code(debug_url_full)
                if st.button("Probar traza (GET)", key="bha_debug_trace_btn"):
                    try:
                        r = api_get_raw(debug_path, debug_params, base_url, token)
                        st.write(f"Status: {r.status_code}")
                        st.write(f"Content-Type: {r.headers.get('content-type', '')}")
                        st.write(f"URL final: {r.url}")
                        if r.status_code >= 400:
                            st.error((r.text or "")[:2000])
                        else:
                            try:
                                data = r.json()
                                if isinstance(data, list) and len(data) > 200:
                                    st.json(data[:200])
                                    st.caption(
                                        f"Mostrando 200 de {len(data)} elementos."
                                    )
                                else:
                                    st.json(data)
                            except Exception:
                                st.text((r.text or "")[:2000])
                    except Exception as e:
                        st.error(f"No pude probar la traza: {e}")

        frames: list[tuple[str, pd.DataFrame]] = []
        used_params_rows: list[dict] = []
        page_indices = list(range(num_pages_to_load)) if num_pages_to_load > 1 else [int(page_number)]

        for trace_uuid in selected_trace_uuids:
            trace_entry = trace_map_by_uuid.get(str(trace_uuid))
            if not trace_entry:
                continue
            label = trace_entry["label"]
            try:
                used_type = str(trace_type).upper()
                used_probe = False
                list_df_pages: list[pd.DataFrame] = []

                for page_idx in page_indices:
                    user_params = build_trace_params(page_override=page_idx)
                    detail = None
                    if trace_source in ("mapped",):
                        detail = api_get_mapped_trace(
                            base_url,
                            token,
                            mapped_scope_uuid,
                            str(trace_uuid),
                            trace_type,
                            mapped_scope_kind,
                            mapped_trace_path or None,
                            user_params or None,
                            force_data_endpoint,
                            well_uuid,
                        )
                    elif trace_source == "catalog":
                        try:
                            detail = api_get_well_trace_data(
                                base_url=base_url,
                                token=token,
                                well_uuid=well_uuid,
                                trace_uuid=str(trace_uuid),
                                trace_type=trace_type,
                                params=user_params or None,
                                calculated=False,
                            )
                            df_tmp = trace_detail_to_df(detail)
                        except Exception:
                            df_tmp = pd.DataFrame()
                            detail = None
                        if df_tmp is None or df_tmp.empty:
                            detail = api_get_well_trace_data(
                                base_url=base_url,
                                token=token,
                                well_uuid=well_uuid,
                                trace_uuid=str(trace_uuid),
                                trace_type=trace_type,
                                params=user_params or None,
                                calculated=True,
                            )
                    else:
                        detail = api_get_drilling_trace(
                            base_url,
                            token,
                            str(trace_uuid),
                            trace_detail_path or None,
                            params=user_params or None,
                        )

                    df_page = trace_detail_to_df(detail)
                    if df_page.empty:
                        break
                    list_df_pages.append(df_page)

                used_params = dict(build_trace_params(page_override=page_indices[0]))
                if num_pages_to_load > 1:
                    used_params["_pages_loaded"] = str(len(list_df_pages))
                df_trace = pd.concat(list_df_pages, ignore_index=True) if list_df_pages else pd.DataFrame()

                if not df_trace.empty and len(list_df_pages) > 1:
                    x_candidates = {"time", "timestamp", "datetime", "date", "depth", "md", "measured_depth", "survey_md", "index"}
                    for c in df_trace.columns:
                        if str(c).strip().lower() in x_candidates:
                            df_trace = df_trace.drop_duplicates(subset=[c], keep="first")
                            break

                user_has_explicit_range = bool(used_params.get("from") and used_params.get("to"))
                if df_trace.empty and trace_source in ("mapped", "catalog") and trip_auto_probe_data and not user_has_explicit_range:
                    df_probe, used_type, used_params = probe_well_trace_data(
                        base_url=base_url,
                        token=token,
                        well_uuid=well_uuid,
                        trace_uuid=str(trace_uuid),
                        prefer_type=trace_type,
                        user_params=used_params,
                    )
                    if not df_probe.empty:
                        st.info(
                            f"Traza {label}: sin datos con params actuales; usando {used_type} con params detectados."
                        )
                        df_trace = df_probe
                        used_probe = True

            except Exception as e:
                msg = str(e)
                if trace_source == "mapped" and "missing required request parameters" in msg:
                    st.warning(
                        f"No pude leer la traza {label}: {e}\n"
                        "Sugerencia: verifica que tengas una lateral seleccionada "
                        "y el dominio (TIME/DEPTH) correcto."
                    )
                    continue
                st.warning(f"No pude leer la traza {label}: {e}")
                continue
            if df_trace.empty:
                st.warning(f"La traza {label} no tiene datos.")
                continue
            used_params_rows.append(
                {
                    "Traza": label,
                    "Fuente": trace_source or "n/a",
                    "Tipo usado": used_type,
                    "Params usados": json.dumps(used_params, ensure_ascii=False) if used_params else "",
                    "Auto-probe": "Sí" if used_probe else "No",
                }
            )
            frames.append((label, df_trace))

        bha_df = merge_trace_frames(frames)
        if bha_df.empty:
            st.error("No pude construir un DataFrame con las trazas seleccionadas.")
            return

        with st.expander("Params efectivos usados", expanded=False):
            if used_params_rows:
                st.dataframe(pd.DataFrame(used_params_rows), use_container_width=True, hide_index=True)
            else:
                st.info("No hay params efectivos para mostrar.")
            st.caption("Nota: el campo 'index' en la tabla proviene del timestamp real de cada muestra.")

        cols = list(bha_df.columns)
        time_col_default = pick_default_column(
            cols, ["time", "timestamp", "datetime", "date", "index"]
        )
        with st.expander("Filtro por tiempo (exacto)", expanded=False):
            apply_time_filter = st.checkbox(
                "Aplicar filtro por rango de tiempo", value=False, key="bha_apply_time_filter"
            )
            offset_hours = st.number_input(
                "Desfase horario (horas)",
                min_value=-12.0,
                max_value=14.0,
                value=0.0,
                step=1.0,
                key="bha_time_offset_hours",
                help="Ej: usa -6 para filtrar con hora local de México si la API está en UTC.",
            )
            use_local_input = st.checkbox(
                "Interpretar fechas como hora local (aplicar desfase)",
                value=False,
                key="bha_time_local_input",
            )
            time_col = st.selectbox(
                "Columna de tiempo",
                ["(ninguna)"] + cols,
                index=(["(ninguna)"] + cols).index(time_col_default)
                if time_col_default in cols
                else 0,
                key="bha_time_col",
            )
            time_from = st.text_input(
                "Desde (ISO 8601, UTC con Z)",
                value="",
                key="bha_time_from",
                help="Ej: 2025-10-21T01:19:00Z",
            )
            time_to = st.text_input(
                "Hasta (ISO 8601, UTC con Z)",
                value="",
                key="bha_time_to",
                help="Ej: 2025-10-21T03:30:00Z",
            )

            if apply_time_filter:
                if time_col == "(ninguna)":
                    st.warning("Selecciona una columna de tiempo válida.")
                else:
                    df_time = bha_df.copy()
                    df_time[time_col] = pd.to_datetime(
                        df_time[time_col], errors="coerce", utc=True
                    )
                    def _parse_time_input(val: str | None):
                        if not val:
                            return None
                        if use_local_input:
                            local_dt = pd.to_datetime(val, errors="coerce")
                            if local_dt is pd.NaT:
                                return pd.NaT
                            if local_dt.tzinfo is None:
                                return (local_dt - pd.to_timedelta(offset_hours, unit="h")).tz_localize("UTC")
                            return local_dt.tz_convert("UTC")
                        return pd.to_datetime(val, errors="coerce", utc=True)

                    from_dt = _parse_time_input(time_from)
                    to_dt = _parse_time_input(time_to)

                    if time_from and from_dt is pd.NaT:
                        st.warning("No pude parsear la fecha 'Desde'. Usa ISO 8601 con Z.")
                    if time_to and to_dt is pd.NaT:
                        st.warning("No pude parsear la fecha 'Hasta'. Usa ISO 8601 con Z.")

                    mask = pd.Series(True, index=df_time.index)
                    if from_dt is not None and from_dt is not pd.NaT:
                        mask &= df_time[time_col] >= from_dt
                    if to_dt is not None and to_dt is not pd.NaT:
                        mask &= df_time[time_col] <= to_dt

                    before = len(df_time)
                    df_time = df_time[mask].copy()
                    after = len(df_time)
                    bha_df = df_time
                    st.info(f"Filtrado por tiempo: {after} de {before} filas.")

        st.dataframe(bha_df.head(200), use_container_width=True, hide_index=True)
        if time_col_default in bha_df.columns:
            tmp_time = pd.to_datetime(bha_df[time_col_default], errors="coerce", utc=True)
            if tmp_time.notna().any():
                st.caption(
                    f"Rango del dataset ({time_col_default}): "
                    f"{tmp_time.min()} → {tmp_time.max()}"
                )
        st.caption(f"Columnas disponibles: {', '.join(map(str, bha_df.columns))}")

        cols = list(bha_df.columns)
        cols_options = [c for c in cols if str(c) != "index"]
        wob_default = pick_default_column(cols, ["WOB", "Weight on Bit", "Surface WOB"]) or cols[0]
        rpm_default = pick_default_column(cols, ["RPM", "Surface RPM"]) or cols[0]
        torque_default = pick_default_column(cols, ["Torque", "Surface Torque"]) or cols[0]
        rop_default = pick_default_column(cols, ["ROP"]) or cols[0]
        depth_default = pick_default_column(cols, ["Depth", "MD", "Measured Depth", "Survey MD"]) or cols[0]

        with st.expander("Configurar columnas", expanded=False):
            bha_wob_col = st.selectbox("Columna WOB", cols_options, index=cols_options.index(wob_default) if wob_default in cols_options else 0, key="bha_col_wob")
            rpm_options = ["(ninguna)"] + cols_options
            bha_rpm_col = st.selectbox(
                "Columna RPM",
                rpm_options,
                index=rpm_options.index(rpm_default) if rpm_default in rpm_options else 0,
                key="bha_col_rpm",
            )
            torque_options = ["(ninguna)"] + cols_options
            bha_torque_col = st.selectbox(
                "Columna Torque",
                torque_options,
                index=torque_options.index(torque_default) if torque_default in torque_options else 0,
                key="bha_col_torque",
            )
            rop_options = ["(ninguna)"] + cols
            depth_options = ["(ninguna)"] + cols
            bha_rop_col = st.selectbox(
                "Columna ROP",
                rop_options,
                index=rop_options.index(rop_default) if rop_default in rop_options else 0,
                key="bha_col_rop",
            )
            bha_depth_col = st.selectbox(
                "Columna Profundidad (MD)",
                depth_options,
                index=depth_options.index(depth_default)
                if depth_default in depth_options
                else 0,
                key="bha_col_depth",
            )
            bha_rpm_col = None if bha_rpm_col == "(ninguna)" else bha_rpm_col
            bha_torque_col = None if bha_torque_col == "(ninguna)" else bha_torque_col
            bha_rop_col = None if bha_rop_col == "(ninguna)" else bha_rop_col
            bha_depth_col = None if bha_depth_col == "(ninguna)" else bha_depth_col

            # Validación: evitar seleccionar la misma columna más de una vez (Plotly requiere nombres únicos)
            selected_cols = [bha_wob_col, bha_rpm_col, bha_torque_col, bha_rop_col, bha_depth_col]
            selected_cols = [c for c in selected_cols if c]
            if len(selected_cols) != len(set(selected_cols)):
                st.error(
                    "Tienes columnas repetidas seleccionadas. "
                    "Elige una columna distinta para cada traza. (RPM/Torque pueden ser '(ninguna)')"
                )
                return

            if not bha_rop_col or not bha_depth_col:
                st.info(
                    "Para MSE necesitas ROP y Profundidad (MD). "
                    "Si no están disponibles, MSE se omitirá."
                )

    elif data_source == "CSV":
        bha_file = st.file_uploader("Upload Engineering Traces (.csv)", type=["csv"])
        with st.expander("Configurar columnas", expanded=False):
            bha_wob_col = st.text_input("Columna WOB", value="WOB")
            bha_rpm_col = st.text_input("Columna RPM", value="Surface RPM")
            bha_torque_col = st.text_input("Columna Torque", value="Torque")
            bha_rop_col = st.text_input("Columna ROP", value="ROP")
            bha_depth_col = st.text_input("Columna Profundidad (MD)", value="Depth")

        if bha_file is None:
            return

        try:
            bha_df = pd.read_csv(bha_file, sep=",", low_memory=False)
        except Exception as e:
            st.error(f"No pude leer el CSV: {e}")
            return

    # Columnas mínimas: WOB siempre, RPM solo si está configurada (puede ser None)
    required_cols = [bha_wob_col] + ([bha_rpm_col] if bha_rpm_col else [])
    required_cols = [c for c in required_cols if c]  # seguridad

    missing_cols = [c for c in required_cols if c not in bha_df.columns]
    if missing_cols:
        st.error("No se encontraron estas columnas: " + ", ".join([str(c) for c in missing_cols if c is not None]))
        return

    for col in [bha_wob_col, bha_rpm_col, bha_torque_col, bha_rop_col, bha_depth_col]:
        if col and col in bha_df.columns:
            bha_df[col] = pd.to_numeric(bha_df[col], errors="coerce")

    # ===============================
    # Filtros de perforación
    # ===============================
    st.subheader("Filtros de perforación")

    use_physical_filter = st.checkbox(
        "Filtrar solo perforación (avance de Hole/Bit Depth)",
        value=True,
        key="bha_use_physical_filter",
    )
    min_delta_m = st.number_input(
        "Umbral mínimo de avance de profundidad (m) entre muestras",
        min_value=0.0,
        value=0.05,
        step=0.01,
        key="bha_min_delta_depth_m",
    )

    use_activity_filter = st.checkbox(
        "Filtrar por Rig Activity",
        value=False,
        key="bha_use_activity_filter",
    )

    activity_col = None
    activity_values: list[str] = []
    if use_activity_filter:
        activity_candidates = [c for c in bha_df.columns if "activity" in str(c).lower()]
        if activity_candidates:
            activity_col = st.selectbox(
                "Columna Rig Activity",
                activity_candidates,
                key="bha_activity_col",
            )
            activity_values = st.multiselect(
                "Valores considerados como perforación",
                options=sorted(bha_df[activity_col].astype(str).unique()),
                default=[v for v in ["DRILLING", "ROTARY DRILLING", "SLIDE DRILLING"] if v in set(map(str, bha_df[activity_col].unique()))],
                key="bha_activity_values",
            )
        else:
            st.info("No se encontró columna de Rig Activity en las trazas.")

    bha_df_filtered = bha_df.copy()
    if use_physical_filter:
        # Preferir Hole Depth en TIME; si no existe, usar la columna seleccionada por el usuario
        preferred_depth_col = None
        if "Hole Depth" in bha_df_filtered.columns:
            preferred_depth_col = "Hole Depth"
        elif bha_depth_col and bha_depth_col in bha_df_filtered.columns:
            preferred_depth_col = bha_depth_col

        if not preferred_depth_col:
            st.warning("Para filtrar perforación por avance de profundidad, configura la columna Hole Depth o Bit depth.")
        else:
            bha_df_filtered = filter_drilling_physical(
                bha_df_filtered,
                depth_col=preferred_depth_col,
                min_delta_m=float(min_delta_m),
            )

    if use_activity_filter and activity_col and activity_values:
        bha_df_filtered = filter_by_rig_activity(bha_df_filtered, activity_col=activity_col, drilling_values=activity_values)

    if bha_df_filtered.empty:
        st.warning(
            "En el rango seleccionado no hay perforación activa "
            "(sin avance de Hole/Bit Depth o sin estados de perforación)."
        )
        # Aún así mostramos una previsualización sin filtros para validar que llegaron trazas
        plot_cols_raw = [c for c in [bha_wob_col, bha_rpm_col, bha_torque_col, bha_rop_col, bha_depth_col] if c and c in bha_df.columns]
        if plot_cols_raw:
            st.subheader("📈 Previsualización de trazas (sin filtros)")
            df_plot = bha_df[plot_cols_raw].copy()
            df_reset = df_plot.reset_index()
            if "index" in df_reset.columns and any(c == "index" for c in plot_cols_raw):
                df_reset = df_reset.rename(columns={"index": "time_index"})
            x_col = "time" if "time" in df_reset.columns else ("time_index" if "time_index" in df_reset.columns else df_reset.columns[0])
            try:
                fig_preview = px.line(df_reset, x=x_col, y=plot_cols_raw, title="Trazas seleccionadas (sin filtros)")
                st.plotly_chart(
                    prettify(fig_preview, h=420),
                    use_container_width=True,
                    config=PLOTLY_CONFIG,
                )
            except Exception as e:
                st.info(f"No pude graficar la previsualización: {e}")
        return

    # Desde aquí en adelante, usa el DF filtrado
    bha_df = bha_df_filtered

    with st.expander("⚙️ Análisis de resonancia (tolerancia)", expanded=False):
        st.caption(
            "La banda de resonancia se define como modo ± tolerancia (Hz). "
            "Si ves «0 ventanas seguras», prueba bajar la tolerancia (ej. 0.08–0.10) para estrechar la banda roja."
        )
        use_custom_tolerance = st.checkbox(
            "Usar tolerancia manual (± Hz)",
            value=False,
            key="bha_use_custom_tolerance",
        )
        tolerance_hz_param: float | None = None
        if use_custom_tolerance:
            tolerance_hz_param = st.number_input(
                "Tolerancia (± Hz)",
                min_value=0.02,
                max_value=0.50,
                value=0.10,
                step=0.01,
                format="%.2f",
                key="bha_tolerance_hz",
                help="Menor valor = banda de resonancia más estrecha = más probabilidad de ventanas seguras.",
            )

    if (bha_rpm_col is not None and bha_wob_col is not None):
        analysis = analyze_bha_resonance(
            bha_df, bha_wob_col, bha_rpm_col, torque_col=bha_torque_col, tolerance_hz=tolerance_hz_param
        )
    else:
        st.warning("RPM o WOB no están configuradas. Se mostrará solo previsualización (sin resonancia).")
        analysis = None
    if analysis is None:
        st.warning("No hay datos válidos para procesar (para resonancia).")
        # Aún así mostramos una previsualización para confirmar que las trazas se cargaron
        plot_cols = [c for c in [bha_wob_col, bha_rpm_col, bha_torque_col, bha_rop_col, bha_depth_col] if c and c in bha_df.columns]
        if plot_cols:
            st.subheader("📈 Previsualización de trazas (sin análisis de resonancia)")
            df_plot = bha_df[plot_cols].copy()
            df_reset = df_plot.reset_index()
            if "index" in df_reset.columns and any(c == "index" for c in plot_cols):
                df_reset = df_reset.rename(columns={"index": "time_index"})
            x_col = "time" if "time" in df_reset.columns else ("time_index" if "time_index" in df_reset.columns else df_reset.columns[0])
            try:
                fig_preview = px.line(df_reset, x=x_col, y=plot_cols, title="Trazas seleccionadas")
                st.plotly_chart(
                    prettify(fig_preview, h=420),
                    use_container_width=True,
                    config=PLOTLY_CONFIG,
                )
            except Exception as e:
                st.info(f"No pude graficar la previsualización: {e}")
        st.info("Tip: para resonancia necesitas tramos con RPM y WOB > 0 (perforación activa).")
        return

    st.success(
        f"✅ Datos limpios: {len(analysis['df']):,} puntos válidos. "
        f"Modos detectados: {np.round(analysis['natural_modes'], 2).tolist()}"
    )
    st.session_state["bha_analysis"] = analysis
    st.session_state["bha_cols"] = {
        "wob": bha_wob_col,
        "rpm": bha_rpm_col,
        "torque": bha_torque_col,
        "rop": bha_rop_col,
        "depth": bha_depth_col,
    }
    if st.button("📄 Generar PPTX de Ingeniería"):
        st.session_state["bha_cols"]["bit_diameter_in"] = st.session_state.get(
            "mse_bit_diameter_in_eng", 8.5
        )
        pptx_str, pptx_path = export_engineering_pptx(analysis, st.session_state["bha_cols"])
        st.success("✅ PPTX de Ingeniería generado.")
        st.download_button(
            "Download PPTX – Ingeniería",
            data=Path(pptx_str).read_bytes(),
            file_name="Engineering_Insights_Report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        # Copiar PPTX a Downloads ANTES de abrirlo con PowerPoint (evita WinError 32)
        copy_report_to_downloads(Path(pptx_str), None, "Engineering_Insights_Report")
        pdf_path = export_pptx_to_pdf(Path(pptx_str))
        copy_report_to_downloads(None, pdf_path, "Engineering_Insights_Report")
        if pdf_path is not None and pdf_path.exists():
            st.download_button(
                "Download PDF – Ingeniería",
                data=pdf_path.read_bytes(),
                file_name="Engineering_Insights_Report.pdf",
                mime="application/pdf",
            )
        else:
            st.info(
                "No pude convertir a PDF. Revisa que PowerPoint esté instalado "
                "y el paquete comtypes disponible."
            )

    st.caption(
        "**Leyenda:** Verde = zona segura, rojo = zona de resonancia. "
        "Color de puntos indica proximidad (0 peligro → 1 seguro) o frecuencia (Hz)."
    )

    safe_count = len(analysis["safe_segments"])
    largest_window = 0
    if analysis["report"] is not None and not analysis["report"].empty:
        largest_window = analysis["report"]["Width (RPM)"].max()
    st.markdown(
        f"**Resumen:** {safe_count} ventanas seguras detectadas. "
        f"Mayor ventana: {largest_window:.1f} RPM. "
        f"Tolerancia: ±{analysis['tolerance']:.2f} Hz."
    )
    if safe_count == 0:
        st.info(
            "**0 ventanas seguras:** En el rango de datos cargado, todo el RPM operativo cae dentro de la banda de resonancia (zona roja). "
            "Puede ser un hallazgo real (se operó en zona resonante) o la banda es muy ancha. "
            "Prueba en «Análisis de resonancia (tolerancia)» bajar la tolerancia (ej. 0.10 Hz) para estrechar la banda roja y ver si aparecen ventanas verdes."
        )

    use_reference_scale = st.checkbox(
        "Usar escala de ejes como referencia (igual que CSV completo: WOB 0–30k, RPM 0–350)",
        value=st.session_state.get("bha_use_reference_scale", False),
        key="bha_use_reference_scale",
        help="Marca esto para que las gráficas se vean como cuando cargas el CSV completo, con los mismos ejes.",
    )
    x_range = (0.0, 30000.0) if use_reference_scale else None
    y_range = (0.0, 350.0) if use_reference_scale else None

    fig1 = build_proximity_figure(
        analysis["df"], bha_wob_col, bha_rpm_col, analysis, bha_torque_col,
        xaxis_range=x_range, yaxis_range=y_range,
    )
    st.plotly_chart(
        fig1,
        use_container_width=True,
        config=PLOTLY_CONFIG,
    )
    chart_notes(
        f"Ventanas seguras: {safe_count}, tolerancia ±{analysis['tolerance']:.2f} Hz.",
        "Puntos: WOB vs RPM; colores indican proximidad a resonancia.",
    )

    fig2 = build_frequency_figure(
        analysis["df"], bha_wob_col, bha_rpm_col, analysis, bha_torque_col,
        xaxis_range=x_range, yaxis_range=y_range,
    )
    st.plotly_chart(
        fig2,
        use_container_width=True,
        config=PLOTLY_CONFIG,
    )
    chart_notes(
        "Mapa de frecuencia rotacional por punto operativo.",
        "Colores representan Hz; bandas rojas = resonancia.",
    )

    # WOB–RPM Heatmap con bins cuadrados (color por ROP / MSE / Shocks & Vibs)
    st.subheader("WOB–RPM Heatmap (squared bins)")
    cols = st.session_state.get("bha_cols") or {}
    rop_col = cols.get("rop")
    depth_col = cols.get("depth")
    torque_col = cols.get("torque") or bha_torque_col
    bit_diameter_in = cols.get("bit_diameter_in", 8.5)
    df_hm = analysis["df"]
    has_rop = rop_col and rop_col in df_hm.columns
    df_mse_for_hm = None
    if has_rop and torque_col and torque_col in df_hm.columns:
        df_mse_for_hm = compute_mse(
            df_hm,
            wob_col=bha_wob_col,
            rpm_col=bha_rpm_col,
            rop_col=rop_col,
            torque_col=torque_col,
            bit_diameter_in=bit_diameter_in,
            depth_col=depth_col if depth_col in df_hm.columns else None,
        )
    has_mse = df_mse_for_hm is not None and not df_mse_for_hm.empty
    shocks_col = find_shocks_vibs_column(df_hm)
    color_options = []
    if has_rop:
        color_options.append("ROP")
    if has_mse:
        color_options.append("MSE")
    if shocks_col:
        color_options.append("Shocks & Vibs")
    if not color_options:
        color_options.append("ROP")
    color_by = st.selectbox(
        "Color bins by",
        options=color_options,
        index=0,
        key="bha_heatmap_color_by",
        help="Cada celda del heatmap muestra la media del indicador en ese bin WOB–RPM.",
    )
    if color_by == "ROP" and not has_rop:
        st.info("No hay columna ROP en los datos. Añade ROP o elige otra variable.")
    elif color_by == "MSE" and not has_mse:
        st.info("MSE no disponible (requiere ROP, WOB, RPM, Torque y profundidad).")
    elif color_by == "Shocks & Vibs" and not shocks_col:
        st.warning("No se encontró columna de Shocks/Vibs. Usa una columna con nombre tipo Shocks, Vibs, Vibration.")
    n_bins_hm = st.number_input("Number of bins (each axis)", min_value=10, max_value=60, value=30, key="bha_heatmap_bins")
    # Heatmap siempre con rango según datos (más amplio); no usar escala de referencia 0–30k
    if color_by == "ROP" and has_rop:
        fig_hm = build_wob_rpm_binned_heatmap(
            df_hm, bha_wob_col, bha_rpm_col, rop_col, "ROP (m/hr)", n_bins=n_bins_hm,
        )
    elif color_by == "MSE" and df_mse_for_hm is not None:
        fig_hm = build_wob_rpm_binned_heatmap(
            df_mse_for_hm, bha_wob_col, bha_rpm_col, "MSE_ksi", "MSE (ksi)", n_bins=n_bins_hm,
        )
    elif color_by == "Shocks & Vibs" and shocks_col:
        fig_hm = build_wob_rpm_binned_heatmap(
            df_hm, bha_wob_col, bha_rpm_col, shocks_col, shocks_col, n_bins=n_bins_hm,
        )
    else:
        fig_hm = None
    if fig_hm is not None:
        st.plotly_chart(fig_hm, use_container_width=True, config=PLOTLY_CONFIG)
        _df_bha_chip = df_mse_for_hm if color_by == "MSE" and df_mse_for_hm is not None else df_hm
        _vcol_bha = (
            rop_col
            if color_by == "ROP"
            else ("MSE_ksi" if color_by == "MSE" else shocks_col)
        )
        _cols_bha_hm = [bha_wob_col, bha_rpm_col]
        if _vcol_bha and _vcol_bha in _df_bha_chip.columns:
            _cols_bha_hm.append(_vcol_bha)
        _st_bha = heatmap_numeric_stats(_df_bha_chip, _cols_bha_hm)
        _ch_bha = stats_df_to_heatmap_chips(_st_bha, max_chips=8)
        if _ch_bha:
            st.caption("**Rangos de los datos usados en el mapa**")
            _render_chips_row([("Bins " + f"{int(n_bins_hm)}×{int(n_bins_hm)}", "gray")] + _ch_bha)
        _zs_bha = None
        if color_by == "ROP" and has_rop:
            _zs_bha = compute_rop_zone_stats(
                df_hm,
                bha_wob_col,
                bha_rpm_col,
                str(rop_col),
                bins=int(n_bins_hm),
                min_points_per_bin=1,
            )
        elif color_by == "MSE" and df_mse_for_hm is not None:
            _zs_bha = compute_rop_zone_stats(
                df_mse_for_hm,
                bha_wob_col,
                bha_rpm_col,
                "MSE_ksi",
                bins=int(n_bins_hm),
                min_points_per_bin=1,
            )
        elif color_by == "Shocks & Vibs" and shocks_col:
            _zs_bha = compute_rop_zone_stats(
                df_hm,
                bha_wob_col,
                bha_rpm_col,
                str(shocks_col),
                bins=int(n_bins_hm),
                min_points_per_bin=1,
            )
        _marg_bha = (
            build_heatmap_marginal_max_curves(
                _zs_bha,
                x_label="WOB (centro de bin)",
                y_label="RPM (centro de bin)",
                z_label=f"max {color_by} en bin",
            )
            if _zs_bha is not None
            else None
        )
        if _marg_bha is not None:
            st.caption("**Curvas de máximo por eje:** pico de la variable coloreada a lo largo de WOB y de RPM.")
            st.plotly_chart(_marg_bha, use_container_width=True, config=PLOTLY_CONFIG)
        _sp_bha = build_minmax_mean_spine_figure(
            _st_bha,
            title="Min · media · max (variables del heatmap WOB–RPM)",
        )
        if _sp_bha is not None:
            st.plotly_chart(_sp_bha, use_container_width=True, config=PLOTLY_CONFIG)
        chart_notes(
            f"Bins cuadrados {n_bins_hm}×{n_bins_hm}. Color = media de {color_by} en cada celda.",
            "Útil para ver zonas WOB–RPM con mejor ROP, menor MSE o mayor/menor shocks.",
        )

    if analysis["report"] is not None:
        st.subheader("🟢 Ventanas operativas seguras")
        st.dataframe(analysis["report"], use_container_width=True, hide_index=True)
    else:
        st.info("⚠️ No se detectaron ventanas seguras fuera de las bandas resonantes.")

    if data_source == "API" and bha_df is not None and st.session_state.get("bha_auto_refresh"):
        interval = int(st.session_state.get("bha_auto_refresh_interval", 30))
        interval = max(10, min(300, interval))
        countdown_placeholder = st.empty()
        for i in range(interval, 0, -1):
            countdown_placeholder.info(f"🔄 Próxima actualización en **{i}** s… (desmarca «Actualizar automáticamente» para detener)")
            time.sleep(1)
        countdown_placeholder.empty()
        st.rerun()


def render_roadmap() -> None:
    st.subheader(tr("roadmap_subheader"))
    st.caption(tr("roadmap_caption"))

    analysis = st.session_state.get("bha_analysis")
    cols = st.session_state.get("bha_cols")
    if analysis is None or cols is None:
        st.info(tr("roadmap_need_bha"))
        return

    st.markdown(tr("roadmap_rt"))
    roadmap_auto_refresh = st.checkbox(
        tr("roadmap_auto"),
        value=st.session_state.get("roadmap_auto_refresh", False),
        key="roadmap_auto_refresh",
        help=tr("roadmap_auto_help"),
    )
    if roadmap_auto_refresh:
        st.number_input(
            tr("interval_seconds"),
            min_value=10,
            max_value=300,
            value=30,
            step=10,
            key="roadmap_auto_refresh_interval",
            help=tr("interval_help_roadmap"),
        )

    df = analysis["df"]
    wob_col = cols["wob"]
    rpm_col = cols["rpm"]
    torque_col = cols["torque"]
    rop_col = cols.get("rop")
    depth_col = cols.get("depth")

    mse_bit_diameter_in = st.number_input(
        tr("mse_bit_diam"), 1.0, 36.0, 8.5, key="mse_bit_diameter_in_eng"
    )

    st.markdown(tr("kpis_exec"))
    c1, c2, c3, c4 = st.columns(4)
    with c1:
        st.metric(tr("metric_valid_points"), f"{len(df):,}")
    with c2:
        st.metric(tr("metric_resonant"), f"{len(analysis['resonant_bands']):,}")
    with c3:
        st.metric(tr("metric_safe"), f"{len(analysis['safe_segments']):,}")
    with c4:
        tol_hz = analysis["tolerance"]
        st.metric(tr("metric_tol_hz"), f"±{tol_hz:.2f}")

    st.subheader(tr("heatmap_eng_title"))
    corr_cols = [wob_col, rpm_col, torque_col, "Freq_Hz", "Proximity_norm"]
    corr_cols = [c for c in corr_cols if c in df.columns]
    corr_df = df[corr_cols].corr()
    _hm_stats_eng = heatmap_numeric_stats(df, corr_cols)
    st.caption("**KPI de variables** (misma base que la matriz de correlación)")
    _chip_eng = stats_df_to_heatmap_chips(_hm_stats_eng, max_chips=12)
    if _chip_eng:
        _render_chips_row([("n=" + f"{len(df):,}", "gray")] + _chip_eng[:11])
    if not _hm_stats_eng.empty:
        with st.expander("Tabla min / media / max por parámetro", expanded=False):
            st.dataframe(_hm_stats_eng, use_container_width=True, hide_index=True)
    # Opción: mostrar valores en % en la celda (escala -100 a 100)
    corr_pct = (corr_df * 100).round(0).astype(int)
    fig_corr = px.imshow(
        corr_df,
        text_auto=".2f",
        color_continuous_scale="RdBu",
        zmin=-1,
        zmax=1,
        title="Engineering Correlation Heatmap",
    )
    fig_corr.update_traces(
        text=corr_pct.values,
        texttemplate="%{text}%",
        textfont={"size": 11},
        xgap=1,
        ygap=1,
    )
    fig_corr.update_layout(coloraxis_colorbar_title="Corr (-1 a 1)")
    st.caption(
        "Colores: azul = correlación positiva, rojo = negativa, blanco ≈ 0. "
        "Números en celdas = **porcentaje** (ej. 60 = 60% de relación lineal)."
    )
    st.plotly_chart(
        prettify_heatmap_auto(fig_corr),
        use_container_width=True,
        config=PLOTLY_CONFIG,
    )
    _spine_eng = build_minmax_mean_spine_figure(
        _hm_stats_eng,
        title="Rango observado por parámetro (● = media en min–max)",
    )
    if _spine_eng is not None:
        st.caption("**Curvas de contexto:** cada barra gris es el rango min→max; el punto azul es la media.")
        st.plotly_chart(_spine_eng, use_container_width=True, config=PLOTLY_CONFIG)
    st.markdown("---")
    st.markdown("**Resumen de lo observado (en % y comentarios)**")
    st.markdown(summarize_heatmap_engineering_pct(corr_df))
    chart_notes(
        summarize_heatmap(corr_df),
        "Heatmap de correlaciones en el set de ingeniería.",
    )

    st.subheader("MSE (Energía Mecánica Específica)")
    if not rop_col or rop_col not in df.columns or not depth_col or depth_col not in df.columns:
        st.info(
            "Para MSE se requieren columnas de ROP y Profundidad (MD) en el CSV de ingeniería."
        )
    else:
        df_mse = compute_mse(
            df,
            wob_col=wob_col,
            rpm_col=rpm_col,
            rop_col=rop_col,
            torque_col=torque_col,
            bit_diameter_in=mse_bit_diameter_in,
            depth_col=depth_col,
        )
        if df_mse is None:
            st.info("No hay datos suficientes para MSE. Verifica torque, ROP y WOB.")
        else:
            c1, c2 = st.columns(2, gap="large")
            fig_mse_depth = build_mse_vs_depth(
                df_mse, "Ingeniería", "Depth_MSE", "Depth (m)"
            )
            fig_mse_hist = build_mse_hist(df_mse, "Ingeniería")
            with c1:
                st.plotly_chart(
                    prettify(fig_mse_depth),
                    use_container_width=True,
                    config=PLOTLY_CONFIG,
                )
                chart_notes(
                    f"{series_summary(df_mse['MSE_ksi'].dropna() if 'MSE_ksi' in df_mse.columns else df_mse['MSE_MPa'].dropna())}.",
                    "MSE vs profundidad; valores en ksi.",
                )
            with c2:
                st.plotly_chart(
                    prettify_hist(fig_mse_hist),
                    use_container_width=True,
                    config=PLOTLY_CONFIG,
                )
                chart_notes(
                    f"{series_summary(df_mse['MSE_ksi'].dropna() if 'MSE_ksi' in df_mse.columns else df_mse['MSE_MPa'].dropna())}.",
                    "Histograma de MSE en ksi.",
                )

    st.subheader("Gráficos cruzados (X-plot)")
    c1, c2 = st.columns(2, gap="large")
    with c1:
        fig_scatter = px.scatter(
            df,
            x=wob_col,
            y=torque_col,
            color=rpm_col,
            color_continuous_scale="Turbo",
            title="WOB vs Torque (color: RPM)",
            labels={wob_col: "WOB", torque_col: "Torque", rpm_col: "RPM"},
        )
        fig_scatter.update_traces(marker=dict(size=6, opacity=0.75))
        st.plotly_chart(
            prettify(fig_scatter),
            use_container_width=True,
            config=PLOTLY_CONFIG,
        )
        chart_notes(
            f"Correlación WOB vs Torque: {safe_corr(df, wob_col, torque_col)}.",
            "Color=RPM; puntos representan mediciones.",
        )
    with c2:
        fig_hist = build_hist_with_trend(
            df[rpm_col],
            title="Distribución de RPM",
            x_label="RPM",
            nbins=40,
        )
        st.plotly_chart(
            prettify_hist(fig_hist),
            use_container_width=True,
            config=PLOTLY_CONFIG,
        )
        chart_notes(
            f"{series_summary(df[rpm_col].dropna())}.",
            "Histograma de RPM (ingeniería).",
        )

    if analysis["report"] is not None:
        st.subheader("🟢 Ventanas seguras (ancho)")
        fig_windows = px.bar(
            analysis["report"],
            x="Window #",
            y="Width (RPM)",
            title="Safe Window Widths (RPM)",
            labels={"Window #": "Window", "Width (RPM)": "Width (RPM)"},
        )
        fig_windows.update_traces(marker_line_width=0)
        st.plotly_chart(
            prettify(fig_windows),
            use_container_width=True,
            config=PLOTLY_CONFIG,
        )
        chart_notes(
            f"Mayor ventana: {format_num(analysis['report']['Width (RPM)'].max())} RPM.",
            "Barras muestran ancho de cada ventana segura.",
        )

    if st.session_state.get("roadmap_auto_refresh"):
        interval = int(st.session_state.get("roadmap_auto_refresh_interval", 30))
        interval = max(10, min(300, interval))
        countdown_placeholder = st.empty()
        for i in range(interval, 0, -1):
            countdown_placeholder.info(f"🔄 Próxima actualización en **{i}** s… (desmarca «Actualizar automáticamente» para detener)")
            time.sleep(1)
        countdown_placeholder.empty()
        st.rerun()


def export_engineering_pptx(analysis, cols) -> tuple[str, Path]:
    prs = Presentation()
    add_title_slide(prs)

    df = analysis["df"]
    wob_col = cols["wob"]
    rpm_col = cols["rpm"]
    torque_col = cols["torque"]
    rop_col = cols.get("rop")
    depth_col = cols.get("depth")
    bit_diameter_in = cols.get("bit_diameter_in", 8.5)

    # Proximity figure
    fig1 = build_proximity_figure(df, wob_col, rpm_col, analysis, torque_col)
    buf1 = io.BytesIO(fig1.to_image(format="png", scale=2))
    add_image_slide(prs, "BHA Proximity to Resonance", buf1)

    # Frequency figure
    fig2 = build_frequency_figure(df, wob_col, rpm_col, analysis, torque_col)
    buf2 = io.BytesIO(fig2.to_image(format="png", scale=2))
    add_image_slide(prs, "BHA Rotational Frequency Mapping", buf2)

    # WOB–RPM binned heatmap (squared bins); color by ROP if available, else MSE, else Shocks & Vibs
    n_bins_pptx = 30
    fig_hm_pptx = None
    if rop_col and rop_col in df.columns:
        fig_hm_pptx = build_wob_rpm_binned_heatmap(
            df, wob_col, rpm_col, rop_col, "ROP (m/hr)", n_bins=n_bins_pptx,
            title="WOB–RPM Heatmap (squared bins) · ROP",
        )
    if fig_hm_pptx is None and rop_col and depth_col and depth_col in df.columns and torque_col:
        df_mse_pptx = compute_mse(
            df, wob_col=wob_col, rpm_col=rpm_col, rop_col=rop_col, torque_col=torque_col,
            bit_diameter_in=bit_diameter_in, depth_col=depth_col,
        )
        if df_mse_pptx is not None:
            fig_hm_pptx = build_wob_rpm_binned_heatmap(
                df_mse_pptx, wob_col, rpm_col, "MSE_ksi", "MSE (ksi)", n_bins=n_bins_pptx,
                title="WOB–RPM Heatmap (squared bins) · MSE",
            )
    if fig_hm_pptx is None:
        shocks_col_pptx = find_shocks_vibs_column(df)
        if shocks_col_pptx:
            fig_hm_pptx = build_wob_rpm_binned_heatmap(
                df, wob_col, rpm_col, shocks_col_pptx, shocks_col_pptx, n_bins=n_bins_pptx,
                title=f"WOB–RPM Heatmap (squared bins) · {shocks_col_pptx}",
            )
    if fig_hm_pptx is not None:
        buf_hm = io.BytesIO(fig_hm_pptx.to_image(format="png", scale=2))
        add_image_slide(prs, "WOB–RPM Heatmap (squared bins)", buf_hm)

    # Heatmap (porcentual)
    corr_cols = [wob_col, rpm_col, torque_col, "Freq_Hz", "Proximity_norm"]
    corr_cols = [c for c in corr_cols if c in df.columns]
    corr_df = df[corr_cols].corr()
    corr_pct = (corr_df * 100).round(0).astype(int)
    fig_corr = px.imshow(
        corr_df,
        text_auto=".2f",
        color_continuous_scale="RdBu",
        zmin=-1,
        zmax=1,
        title="Engineering Correlation Heatmap",
    )
    fig_corr.update_traces(
        text=corr_pct.values,
        texttemplate="%{text}%",
        textfont=dict(size=11),
        xgap=1,
        ygap=1,
    )
    fig_corr.update_layout(coloraxis_colorbar_title="Corr (-1 a 1)")
    fig_corr = prettify_heatmap(fig_corr)
    buf3 = io.BytesIO(fig_corr.to_image(format="png", scale=2))
    add_image_slide(prs, "Engineering Correlation Heatmap", buf3)
    _stats_pptx = heatmap_numeric_stats(df, corr_cols)
    _sp_pptx = build_minmax_mean_spine_figure(
        _stats_pptx,
        title="Parámetros – rango min–media–max (normalizado)",
    )
    if _sp_pptx is not None:
        buf_spine = io.BytesIO(_sp_pptx.to_image(format="png", scale=2))
        add_image_slide(prs, "Contexto heatmap – min / media / max por parámetro", buf_spine)

    # Safe windows bar
    if analysis["report"] is not None:
        fig_windows = px.bar(
            analysis["report"],
            x="Window #",
            y="Width (RPM)",
            title="Safe Window Widths (RPM)",
            labels={"Window #": "Window", "Width (RPM)": "Width (RPM)"},
        )
        fig_windows.update_traces(marker_line_width=0)
        fig_windows = prettify(fig_windows)
        buf4 = io.BytesIO(fig_windows.to_image(format="png", scale=2))
        add_image_slide(prs, "Safe Window Widths", buf4)

    # MSE slides (if columns available)
    if rop_col in df.columns and depth_col in df.columns:
        df_mse = compute_mse(
            df,
            wob_col=wob_col,
            rpm_col=rpm_col,
            rop_col=rop_col,
            torque_col=torque_col,
            bit_diameter_in=bit_diameter_in,
            depth_col=depth_col,
        )
        if df_mse is not None:
            fig_mse_depth = build_mse_vs_depth(
                df_mse, "Ingeniería", "Depth_MSE", "Depth (m)"
            )
            fig_mse_hist = build_mse_hist(df_mse, "Ingeniería")
            buf5 = io.BytesIO(fig_mse_depth.to_image(format="png", scale=2))
            add_image_slide(prs, "MSE vs Profundidad", buf5)
            buf6 = io.BytesIO(fig_mse_hist.to_image(format="png", scale=2))
            add_image_slide(prs, "MSE Distribution", buf6)

    tmp_dir = Path(tempfile.mkdtemp())
    pptx_path = tmp_dir / "Engineering_Insights_Report.pptx"
    prs.save(str(pptx_path))
    return str(pptx_path), pptx_path


# =========================
# Tripping Analysis – helpers y render
# =========================
def _normalize_iso_year(s: str) -> str:
    """Corrige año en formato ISO (ej. 026 -> 2026, 26 -> 2026) para evitar rangos inválidos."""
    if not s or not isinstance(s, str):
        return s or ""
    s = s.strip()
    m = re.match(r"^(\d{1,3})(-\d{2}-\d{2})", s)
    if not m:
        return s
    yy = int(m.group(1))
    if yy >= 1000:  # ya es 4 dígitos
        return s
    if yy <= 99:
        year_4 = 2000 + yy
    else:
        year_4 = 2000 + (yy % 100)
    return str(year_4) + s[m.end(1):]


def _trip_iso_to_utc_z(s: str, default_offset: str = "-06:00") -> str:
    """Normaliza un ISO string a UTC (sufijo 'Z').

    - Si el string ya trae zona (Z o ±HH:MM), se convierte a UTC.
    - Si NO trae zona (naive), se asume `default_offset` (por default UTC-06:00) y se convierte a UTC.
    - También normaliza formato con espacio -> 'T' y corrige año corto vía `_normalize_iso_year`.
    """
    if not s or not isinstance(s, str):
        return s or ""
    s0 = _normalize_iso_year(s.strip())
    if not s0:
        return ""
    # Normaliza separador fecha-hora
    if " " in s0 and "T" not in s0:
        s0 = s0.replace(" ", "T")
    # Detecta timezone explícita al final
    has_tz = bool(re.search(r"(Z|[+-]\d{2}:\d{2})$", s0))
    try:
        if s0.endswith("Z"):
            dt = datetime.fromisoformat(s0.replace("Z", "+00:00"))
            return dt.astimezone(timezone.utc).isoformat().replace("+00:00", "Z")
        if has_tz:
            dt = datetime.fromisoformat(s0)
            return dt.astimezone(timezone.utc).isoformat().replace("+00:00", "Z")
        # Naive: asume offset local y convierte a UTC
        dt = datetime.fromisoformat(s0 + default_offset)
        return dt.astimezone(timezone.utc).isoformat().replace("+00:00", "Z")
    except Exception:
        # Si algo falla, retorna lo original para no romper la query
        return s0



def _trip_pick_time_col(df: pd.DataFrame) -> str | None:
    """Heurística: encuentra la columna temporal (timestamp, datetime, o formato ISO como YYYY-MM-DDTHH:MM:SS)."""
    for c in df.columns:
        cl = str(c).strip().lower()
        if cl in ("timestamp", "time", "datetime", "date"):
            return c
        if "time" in cl or "fecha" in cl:
            return c
        if "yyyy-mm-dd" in cl and "hh:mm" in cl:
            return c
        if cl == "yyyy-mm-ddthh:mm:ss":
            return c
        if "hh:mm" in cl or "mm:ss" in cl:
            return c
    # Muchas APIs devuelven la dimensión tiempo/profundidad en una columna llamada "index"
    for c in df.columns:
        if str(c).strip().lower() == "index":
            return c
    return None


def _trip_pick_value_col(df: pd.DataFrame, prefer_names: list[str] | None = None) -> str | None:
    """Heurística: encuentra la columna de valor (por nombre o primera numérica)."""
    prefer_names = prefer_names or []
    lowered = {str(c).strip().lower(): c for c in df.columns}
    for name in prefer_names:
        key = str(name).strip().lower()
        if key in lowered:
            return lowered[key]
    # APIs que devuelven solo "index" y "value": usar "value" como columna de valor
    if "value" in lowered:
        return lowered["value"]
    for c in df.columns:
        if pd.api.types.is_numeric_dtype(df[c]):
            return c
    return None


def _trip_to_numeric(df: pd.DataFrame, col: str) -> pd.Series:
    return pd.to_numeric(df[col], errors="coerce")


def _trip_normalize_and_clean(df: pd.DataFrame) -> pd.DataFrame:
    """
    Preprocesamiento estándar para Tripping Analysis:
    Limpia nombres, fuerza Bit depth y Hookload a numérico, elimina NaN, convierte Timestamp a datetime.
    """
    df = df.copy()
    df.columns = [str(c).strip() for c in df.columns]
    if "Bit depth" in df.columns:
        df["Bit depth"] = pd.to_numeric(df["Bit depth"], errors="coerce")
    if "Hookload" in df.columns:
        df["Hookload"] = pd.to_numeric(df["Hookload"], errors="coerce")
    df = df.dropna(subset=["Bit depth", "Hookload"])
    if "Timestamp" in df.columns:
        df["Timestamp"] = pd.to_datetime(df["Timestamp"], errors="coerce")
    if "Timestamp" in df.columns:
        df = df.sort_values("Timestamp").reset_index(drop=True)
    return df


def _trip_parse_limits_csv(raw: pd.DataFrame) -> tuple[pd.DataFrame | None, str]:
    """
    Parsea un CSV de límites simulados. Espera columna de profundidad y columnas
    Trip Out, Trip In, Rotating (o Rotando). Devuelve (df con columnas Depth, TripOut, TripIn, Rotating), error_msg.
    """
    raw = raw.copy()
    raw.columns = [str(c).strip() for c in raw.columns]
    depth_col = None
    for c in raw.columns:
        cl = str(c).strip().lower()
        if cl in ("depth", "md", "bit depth", "profundidad", "measured depth"):
            depth_col = c
            break
        if "depth" in cl or "profundidad" in cl:
            depth_col = c
            break
    if not depth_col:
        return None, "No se encontró columna de profundidad (Depth, MD, Bit depth, Profundidad)."
    tripout_col = tripin_col = rot_col = None
    for c in raw.columns:
        cl = str(c).strip().lower()
        if "trip" in cl and "out" in cl:
            tripout_col = c
        elif "trip" in cl and "in" in cl:
            tripin_col = c
        elif "rotat" in cl or cl == "rotando":
            rot_col = c
    if not (tripout_col or tripin_col or rot_col):
        return None, "No se encontraron columnas Trip Out, Trip In o Rotating (Rotando)."
    out = raw[[depth_col]].copy()
    out = out.rename(columns={depth_col: "Depth"})
    out["Depth"] = pd.to_numeric(out["Depth"], errors="coerce")
    out = out.dropna(subset=["Depth"]).sort_values("Depth")
    if tripout_col and tripout_col in raw.columns:
        out["TripOut"] = pd.to_numeric(raw.loc[out.index, tripout_col], errors="coerce").values
    else:
        out["TripOut"] = np.nan
    if tripin_col and tripin_col in raw.columns:
        out["TripIn"] = pd.to_numeric(raw.loc[out.index, tripin_col], errors="coerce").values
    else:
        out["TripIn"] = np.nan
    if rot_col and rot_col in raw.columns:
        out["Rotating"] = pd.to_numeric(raw.loc[out.index, rot_col], errors="coerce").values
    else:
        out["Rotating"] = np.nan
    return out.reset_index(drop=True), ""


def _trip_prepare_limits_for_continuous_line(
    limits_df: pd.DataFrame,
    rolling_window: int | None = 25,
) -> pd.DataFrame:
    """
    Prepara el DataFrame de límites para dibujar una sola curva continua por columna.
    Ordena por Depth, colapsa duplicados (un valor por profundidad) y opcionalmente suaviza.
    Así se evita el efecto 'escalera' o 'peine' (líneas horizontales segmentadas).
    """
    if limits_df is None or limits_df.empty or "Depth" not in limits_df.columns:
        return pd.DataFrame()
    lim = limits_df.copy()
    lim["Depth"] = pd.to_numeric(lim["Depth"], errors="coerce")
    for c in ["TripOut", "TripIn", "Rotating"]:
        if c in lim.columns:
            lim[c] = pd.to_numeric(lim[c], errors="coerce")
    lim = lim.dropna(subset=["Depth"]).sort_values("Depth")
    # Un solo valor por profundidad (evita segmentos por punto)
    agg_dict = {c: "mean" for c in ["TripOut", "TripIn", "Rotating"] if c in lim.columns}
    if agg_dict:
        lim = lim.groupby("Depth", as_index=False).agg(agg_dict)
    # Suavizado opcional para curva más profesional
    if rolling_window and rolling_window >= 3:
        for c in ["TripOut", "TripIn", "Rotating"]:
            if c in lim.columns and lim[c].notna().any():
                lim[c] = lim[c].rolling(rolling_window, center=True, min_periods=1).median()
    return lim


# ===============================
# Broomstick: modelo por familia de FF (PU/SO/ROT) + interpolación
# ===============================
_TRIP_FF_FAMILY_RE = re.compile(r"^(PU|SO|ROT)[_ ]?([0-9]*\.?[0-9]+)$", re.IGNORECASE)

def parse_ff_family_csv(raw: pd.DataFrame) -> tuple[pd.DataFrame | None, dict, str]:
    """Parsea un CSV con familia de curvas por FF.

    Formato esperado (mínimo):
      Depth, PU_0.10, PU_0.20, ..., SO_0.10, SO_0.20, ... (y opcional ROT_xx)

    Devuelve:
      - df_model con columna 'Depth' y columnas numéricas por curva
      - fam_map: {'PU': {0.1:'PU_0.10', ...}, 'SO': {...}, 'ROT': {...}}
      - error_msg ("" si ok)
    """
    raw = raw.copy()
    raw.columns = [str(c).strip() for c in raw.columns]

    # Detectar Depth
    depth_col = None
    for c in raw.columns:
        cl = c.lower().strip()
        if cl in ("depth", "md", "bit depth", "profundidad", "measured depth"):
            depth_col = c
            break
    if not depth_col:
        for c in raw.columns:
            cl = c.lower()
            if "depth" in cl or "profundidad" in cl:
                depth_col = c
                break
    if not depth_col:
        return None, {}, "No se encontró columna Depth/MD/Bit depth/Profundidad."

    df = raw.rename(columns={depth_col: "Depth"}).copy()
    df["Depth"] = pd.to_numeric(df["Depth"], errors="coerce")
    df = df.dropna(subset=["Depth"]).sort_values("Depth").reset_index(drop=True)

    fam = {"PU": {}, "SO": {}, "ROT": {}}
    for c in df.columns:
        if c == "Depth":
            continue
        # Normalizar nombre para matching
        cname = str(c).strip()
        cname_norm = (
            cname.replace("-", "_")
            .replace("FF", "")
            .replace("ff", "")
            .replace(" ", "_")
        )
        m = _TRIP_FF_FAMILY_RE.match(cname_norm)
        if not m:
            continue
        mode = m.group(1).upper()
        ff = float(m.group(2))
        fam[mode][ff] = c
        df[c] = pd.to_numeric(df[c], errors="coerce")

    if sum(len(v) for v in fam.values()) == 0:
        return None, {}, "No se detectaron columnas tipo PU_0.10 / SO_0.10 / ROT_0.10."
    return df, fam, ""

# Conversión kg → klb para que el modelo simplificado use la misma escala que datos típicos de Hookload (klb)
KG_TO_KLB = 0.00220462


def generate_broomstick_curves_simplified(
    md_min_m: float,
    md_max_m: float,
    step_m: float,
    weight_per_m_kg: float,
    buoyancy_factor: float,
    inclination_deg: float,
    ff_values: list[float],
    output_klb: bool = True,
    block_weight_klb: float = 0.0,
    inclination_surface_deg: float | None = None,
    inclination_td_deg: float | None = None,
    survey_df: pd.DataFrame | None = None,
) -> tuple[pd.DataFrame, dict]:
    """
    Genera familia de curvas Hookload vs Depth (PU, SO, ROT) para varios FF.
    Si survey_df tiene columnas 'Depth' e 'Inclination', se usa inclinación del survey (curvas reales).
    Si no, inclination_surface/td distintos dan perfil lineal; si no, inclinación constante.
    block_weight_klb desplaza todas las curvas.
    """
    import math
    if step_m <= 0 or md_max_m <= md_min_m or not ff_values:
        return pd.DataFrame(), {}
    depths = np.arange(md_min_m, md_max_m + step_m * 0.5, step_m)
    depths = np.round(depths, 2)
    w_eff = float(weight_per_m_kg) * float(buoyancy_factor)
    scale = KG_TO_KLB if output_klb else 1.0
    block_klb = float(block_weight_klb or 0.0)

    use_survey = (
        survey_df is not None
        and not survey_df.empty
        and "Depth" in survey_df.columns
        and "Inclination" in survey_df.columns
    )
    if use_survey:
        survey_df = survey_df.dropna(subset=["Depth", "Inclination"]).sort_values("Depth")
        use_survey = len(survey_df) >= 2
    if use_survey:
        survey_depth = survey_df["Depth"].values.astype(float)
        survey_inc = survey_df["Inclination"].values.astype(float)
    else:
        survey_depth = survey_inc = None

    use_linear_inc = (
        not use_survey
        and inclination_surface_deg is not None
        and inclination_td_deg is not None
        and (abs(inclination_surface_deg - inclination_td_deg) > 0.01)
    )
    depth_span = float(md_max_m - md_min_m) if md_max_m != md_min_m else 1.0

    out = pd.DataFrame({"Depth": depths})
    fam = {"PU": {}, "SO": {}, "ROT": {}}

    for ff in sorted(set(ff_values)):
        ff = round(ff, 2)
        if use_survey and survey_depth is not None and survey_inc is not None:
            inc_deg_at_depths = np.interp(depths, survey_depth, survey_inc)
            pu = np.zeros_like(depths, dtype=float)
            so = np.zeros_like(depths, dtype=float)
            rot = np.zeros_like(depths, dtype=float)
            pu[0] = block_klb
            so[0] = block_klb
            rot[0] = block_klb
            for i in range(1, len(depths)):
                inc_rad = math.radians(float(inc_deg_at_depths[i]))
                cos_i = math.cos(inc_rad)
                sin_i = math.sin(inc_rad)
                step = depths[i] - depths[i - 1]
                pu[i] = pu[i - 1] + step * w_eff * (cos_i + ff * sin_i) * scale
                so[i] = so[i - 1] + step * w_eff * (cos_i - ff * sin_i) * scale
                rot[i] = rot[i - 1] + step * w_eff * cos_i * scale
        elif use_linear_inc:
            inc_surf_rad = math.radians(float(inclination_surface_deg))
            inc_td_rad = math.radians(float(inclination_td_deg))
            pu = np.zeros_like(depths, dtype=float)
            so = np.zeros_like(depths, dtype=float)
            rot = np.zeros_like(depths, dtype=float)
            pu[0] = block_klb
            so[0] = block_klb
            rot[0] = block_klb
            for i in range(1, len(depths)):
                d = depths[i]
                t = (d - md_min_m) / depth_span
                inc_rad = inc_surf_rad + t * (inc_td_rad - inc_surf_rad)
                cos_i = math.cos(inc_rad)
                sin_i = math.sin(inc_rad)
                step = depths[i] - depths[i - 1]
                pu[i] = pu[i - 1] + step * w_eff * (cos_i + ff * sin_i) * scale
                so[i] = so[i - 1] + step * w_eff * (cos_i - ff * sin_i) * scale
                rot[i] = rot[i - 1] + step * w_eff * cos_i * scale
        else:
            inc_rad = math.radians(float(inclination_deg))
            cos_i = math.cos(inc_rad)
            sin_i = math.sin(inc_rad)
            pu = block_klb + depths * w_eff * (cos_i + ff * sin_i) * scale
            so = block_klb + depths * w_eff * (cos_i - ff * sin_i) * scale
            rot = block_klb + depths * w_eff * cos_i * scale

        col_pu = f"PU_{ff:.2f}"
        col_so = f"SO_{ff:.2f}"
        col_rot = f"ROT_{ff:.2f}"
        out[col_pu] = pu
        out[col_so] = so
        out[col_rot] = rot
        fam["PU"][ff] = col_pu
        fam["SO"][ff] = col_so
        fam["ROT"][ff] = col_rot
    return out, fam


def interp_ff_curve(df_model: pd.DataFrame, fam_map: dict, mode: str, ff: float) -> pd.Series:
    """Interpola una curva Hookload(Depth) para un FF arbitrario usando la familia disponible."""
    mode = mode.upper()
    ffs = sorted(fam_map.get(mode, {}).keys())
    if not ffs:
        return pd.Series(index=df_model.index, dtype=float)

    if ff <= ffs[0]:
        return df_model[fam_map[mode][ffs[0]]]
    if ff >= ffs[-1]:
        return df_model[fam_map[mode][ffs[-1]]]

    lo = max([x for x in ffs if x <= ff])
    hi = min([x for x in ffs if x >= ff])
    if hi == lo:
        return df_model[fam_map[mode][lo]]

    w = (ff - lo) / (hi - lo)
    return (1 - w) * df_model[fam_map[mode][lo]] + w * df_model[fam_map[mode][hi]]

def _trip_dir_labels_from_depth(df: pd.DataFrame, eps_m: float = 0.02) -> pd.Series:
    """Etiqueta puntos como PU (depth bajando) o SO (depth subiendo)."""
    d = pd.to_numeric(df["Bit depth"], errors="coerce").diff()
    lab = pd.Series(index=df.index, dtype="object")
    lab[d < -eps_m] = "PU"
    lab[d > eps_m] = "SO"
    return lab
def _trip_build_exact_depth_envelope(df: pd.DataFrame) -> pd.DataFrame:
    """
    Envelope por profundidad exacta: para cada Bit depth conserva la fila donde Hookload es máximo.
    En conexiones (Bit depth fijo, Hookload variable) queda un punto por profundidad con el mayor Hookload.
    """
    if df.empty or "Bit depth" not in df.columns or "Hookload" not in df.columns:
        return pd.DataFrame()
    filtered = df.loc[df.groupby("Bit depth")["Hookload"].idxmax()].sort_values("Bit depth").reset_index(drop=True)
    return filtered


TRIP_SEVERITY_COLORS = {"low": "#FACC15", "medium": "#F97316", "high": "#EF4444"}
TRIP_SEVERITY_LABELS = {
    "low": "Puntos apretados aislados",
    "medium": "Frecuencia en aumento",
    "high": "Restricción / Riesgo de pegadura",
}


def _trip_compute_time_overpull(
    df: pd.DataFrame,
    rolling_window: int = 60,
    overpull_thr: float = 0.0,
) -> tuple[pd.DataFrame, pd.Timestamp]:
    """
    Calcula overpull en dominio tiempo: baseline = rolling median, overpull = Hookload - baseline.
    Devuelve (df con columnas extra Overpull_t, Event_t, Severity, Hours), t0 (Timestamp mínimo).
    """
    d = df.sort_values("Timestamp").reset_index(drop=True).copy()
    d["Hookload"] = pd.to_numeric(d["Hookload"], errors="coerce")
    d["Baseline_roll"] = d["Hookload"].rolling(window=rolling_window, min_periods=max(1, rolling_window // 2), center=True).median()
    d["Overpull_t"] = d["Hookload"] - d["Baseline_roll"]
    d["Overpull_t"] = d["Overpull_t"].clip(lower=0)
    d["Event_t"] = d["Overpull_t"] >= overpull_thr if overpull_thr > 0 else d["Overpull_t"] > 0
    op_vals = d.loc[d["Event_t"], "Overpull_t"]
    if len(op_vals) >= 3:
        p33 = float(op_vals.quantile(0.33))
        p66 = float(op_vals.quantile(0.66))
        if p66 <= p33:
            p66 = max(p33 + 1e-6, op_vals.max())

        def _sev(x):
            if pd.isna(x) or x <= 0:
                return "low"
            if x <= p33:
                return "low"
            if x <= p66:
                return "medium"
            return "high"
        d["Severity"] = d.apply(lambda r: _sev(r["Overpull_t"]) if r["Event_t"] else None, axis=1)
    else:
        d["Severity"] = d.apply(
            lambda r: "high" if r["Event_t"] and r["Overpull_t"] > 0 else (None if not r["Event_t"] else "low"),
            axis=1,
        )
    t0 = d["Timestamp"].min()
    d["Hours"] = (d["Timestamp"] - t0).dt.total_seconds() / 3600.0
    return d, t0



def _trip_build_events_timeline_figure(
    df_events: pd.DataFrame,
    t0: pd.Timestamp,
    mode_label: str,
    dark: bool = False,
) -> go.Figure:
    """Timeline de overpull: barras por evento coloreadas por severidad (legible y sin 'bloques' gigantes)."""
    fig = go.Figure()

    # Ancho de barra: basado en la separación típica de muestras (en horas).
    try:
        diffs = df_events["Hours"].sort_values().diff().dropna()
        step = float(diffs.median()) if len(diffs) else 0.01
        # 80% del paso típico, con límites razonables
        bar_w = max(0.002, min(0.08, 0.8 * step))
    except Exception:
        bar_w = 0.01

    for sev in ("low", "medium", "high"):
        sub = df_events[(df_events["Event_t"]) & (df_events["Severity"] == sev)]
        if sub.empty:
            continue
        fig.add_trace(
            go.Bar(
                x=sub["Hours"],
                y=sub["Overpull_t"],
                name=TRIP_SEVERITY_LABELS[sev],
                marker_color=TRIP_SEVERITY_COLORS[sev],
                marker_line_width=0,
                opacity=0.9,
                width=bar_w,
                hovertemplate="Tiempo: %{x:.2f} h<br>Overpull: %{y:.2f} klb<extra></extra>",
            )
        )

    fig.update_layout(
        title=dict(
            text=f"{mode_label} – Eventos de Overpull en el tiempo<br><sub style='font-size:11px;color:rgba(148,163,184,0.95);'>Barras = magnitud del overpull por evento · Color = severidad (aislado → frecuencia en aumento → riesgo pegadura)</sub>",
            font=dict(size=16),
        ),
        xaxis_title="Tiempo (h desde inicio)",
        yaxis_title="Overpull (klb)",
        barmode="overlay",
        bargap=0.10,
        height=520,
        template="plotly_dark" if dark else "plotly_white",
        margin=dict(l=60, r=25, t=100, b=60),
        font=dict(size=12),
        legend=dict(orientation="h", yanchor="bottom", y=1.18, xanchor="right", x=1),
        hovermode="x unified",
        xaxis=dict(
            showgrid=True,
            gridcolor="rgba(255,255,255,0.08)" if dark else "rgba(0,0,0,0.06)",
            zeroline=False,
        ),
        yaxis=dict(
            showgrid=True,
            gridcolor="rgba(255,255,255,0.08)" if dark else "rgba(0,0,0,0.06)",
            zeroline=False,
        ),
        # Transparente para que se integre bien con el fondo (oscuro o claro)
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
    )
    return fig



def _trip_build_hookload_with_events_figure(
    df: pd.DataFrame,
    mode_label: str,
    dark: bool = False,
) -> go.Figure:
    """Gráfico Hookload vs tiempo: línea base + puntos donde hay evento, coloreados por severidad."""
    fig = go.Figure()
    fig.add_trace(
        go.Scatter(
            x=df["Hours"],
            y=df["Hookload"],
            mode="lines",
            name="Hookload",
            line=dict(color="#94A3B8" if dark else "#64748B", width=1.2),
        )
    )
    for sev in ("low", "medium", "high"):
        sub = df[(df["Event_t"]) & (df["Severity"] == sev)]
        if sub.empty:
            continue
        color = TRIP_SEVERITY_COLORS[sev]
        fig.add_trace(
            go.Scatter(
                x=sub["Hours"],
                y=sub["Hookload"],
                mode="markers",
                name=TRIP_SEVERITY_LABELS[sev],
                marker=dict(size=8, color=color, line=dict(width=0), symbol="diamond"),
            )
        )
    fig.update_layout(
        title=dict(
            text=f"Hookload – {mode_label}<br><sub style='font-size:10px;color:rgba(148,163,184,0.95);'>Línea = carga en gancho · ◆ = eventos de overpull coloreados por severidad</sub>",
            font=dict(size=16),
        ),
        xaxis_title="Tiempo (h desde inicio)",
        yaxis_title="Hookload (klb)",
        height=420,
        template="plotly_dark" if dark else "plotly_white",
        margin=dict(l=55, r=25, t=85, b=55),
        font=dict(size=11),
        legend=dict(orientation="h", yanchor="bottom", y=1.14, xanchor="right", x=1),
        hovermode="x unified",
        xaxis=dict(
            showgrid=True,
            gridcolor="rgba(255,255,255,0.08)" if dark else "rgba(0,0,0,0.06)",
            zeroline=False,
        ),
        yaxis=dict(
            showgrid=True,
            gridcolor="rgba(255,255,255,0.08)" if dark else "rgba(0,0,0,0.06)",
            zeroline=False,
        ),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
    )
    return fig


def _trip_build_events_vs_depth_figure(
    df_events: pd.DataFrame,
    mode_label: str,
    dark: bool = False,
) -> go.Figure:
    """Gráfico eventos de overpull vs profundidad, coloreado por severidad.

    Nota: este gráfico se ve "blanco" fácilmente en Streamlit cuando theme.base no viene como 'dark'.
    Para que quede consistente con dashboards oscuros, lo renderizamos con estilo pro-dark y fondo transparente.
    """
    fig = go.Figure()
    if "Bit depth" not in df_events.columns:
        return fig

    ev = df_events[df_events["Event_t"]].copy()
    if ev.empty:
        return fig

    # Scatter por severidad
    for sev in ("low", "medium", "high"):
        sub = ev[ev["Severity"] == sev]
        if sub.empty:
            continue
        color = TRIP_SEVERITY_COLORS[sev]
        fig.add_trace(
            go.Scatter(
                x=sub["Bit depth"],
                y=sub["Overpull_t"],
                mode="markers",
                name=TRIP_SEVERITY_LABELS[sev],
                marker=dict(size=9, color=color, line=dict(width=0), symbol="diamond"),
                hovertemplate="<b>%{fullData.name}</b><br>Profundidad: %{x:.0f} m<br>Overpull: %{y:.2f}<extra></extra>",
            )
        )

    # Layout pro (oscuro + transparente)
    fig.update_layout(
        title=dict(
            text=f"{mode_label} – Eventos de overpull vs profundidad<br><sub style='font-size:10px;color:rgba(235,235,240,0.85);'>◆ = eventos por severidad · Eje Y = magnitud del overpull</sub>",
            font=dict(size=16),
        ),
        xaxis_title="Profundidad MD (m)",
        yaxis_title="Overpull (klb)",
        height=380,
        template="plotly_dark",  # forzar consistencia visual
        margin=dict(l=60, r=25, t=60, b=55),
        font=dict(size=12, color="rgba(240,240,245,0.95)"),
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.08,
            xanchor="right",
            x=1,
            bgcolor="rgba(0,0,0,0)",
        ),
        # Fondo 100% transparente para que se integre con el tema del dashboard
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        hovermode="x unified",
    )

    fig.update_xaxes(
        autorange="reversed",
        showgrid=True,
        gridcolor="rgba(255,255,255,0.08)",
        zeroline=False,
        linecolor="rgba(255,255,255,0.18)",
        tickfont=dict(color="rgba(235,235,240,0.95)"),
        title=dict(font=dict(color="rgba(235,235,240,0.95)")),
    )
    fig.update_yaxes(
        showgrid=True,
        gridcolor="rgba(255,255,255,0.08)",
        zeroline=False,
        linecolor="rgba(255,255,255,0.18)",
        tickfont=dict(color="rgba(235,235,240,0.95)"),
        title=dict(font=dict(color="rgba(235,235,240,0.95)")),
    )

    return fig


def _trip_generate_insight(df: pd.DataFrame, env_events_count: int) -> str:
    """Genera una frase interpretativa según el patrón de eventos."""
    events = df[df["Event_t"]]
    n = len(events)
    if n == 0:
        return "No se detectaron eventos de overpull por encima del baseline en el rango analizado."
    total_hrs = (df["Timestamp"].max() - df["Timestamp"].min()).total_seconds() / 3600.0
    rate = n / max(0.1, total_hrs)
    high = (events["Severity"] == "high").sum()
    if high >= n * 0.4 and rate > 2:
        return "**Cada overpull parecía manejable. El patrón no.** Frecuencia alta y muchos eventos de alta severidad sugieren restricción o riesgo de pegadura."
    if rate > 5:
        return "Frecuencia de eventos en aumento a lo largo del tiempo. Revisar tendencia y condiciones de pozo."
    if high > 0:
        return "Eventos de alta severidad presentes. Recomendable revisar profundidades y carga."
    return "Eventos aislados; patrón estable. Seguir monitoreando."


def render_tripping_analysis() -> None:
    # Migración: valores antiguos del radio de envelope (español) → códigos estables
    _tev = st.session_state.get("trip_envelope_method")
    if _tev == "Profundidad exacta (max Hookload por depth)":
        st.session_state["trip_envelope_method"] = TRIP_ENV_EXACT
    elif _tev == "Bin (rangos de profundidad)":
        st.session_state["trip_envelope_method"] = TRIP_ENV_BIN

    st.markdown(
        f"""
        <div style="margin-bottom: 0.5rem;">
            <span style="
                font-size: 1.5rem; font-weight: 600; color: inherit;
            ">{tr("tab_trip")}</span>
            <span style="
                display: inline-flex; align-items: center; gap: 0.35rem; margin-left: 0.75rem;
                flex-wrap: wrap;
            ">
                <span style="
                    background: linear-gradient(135deg, #1e3a5f 0%, #2563eb 100%);
                    color: #fff; font-size: 0.7rem; font-weight: 600;
                    padding: 0.2rem 0.55rem; border-radius: 999px; letter-spacing: 0.02em;
                ">{tr("trip_badge_hookload")}</span>
                <span style="
                    background: linear-gradient(135deg, #422006 0%, #f59e0b 100%);
                    color: #fff; font-size: 0.7rem; font-weight: 600;
                    padding: 0.2rem 0.55rem; border-radius: 999px; letter-spacing: 0.02em;
                ">{tr("trip_badge_trip_io")}</span>
                <span style="
                    background: linear-gradient(135deg, #450a0a 0%, #ef4444 100%);
                    color: #fff; font-size: 0.7rem; font-weight: 600;
                    padding: 0.2rem 0.55rem; border-radius: 999px; letter-spacing: 0.02em;
                ">{tr("trip_badge_overpull")}</span>
            </span>
        </div>
        """,
        unsafe_allow_html=True,
    )
    st.caption(tr("trip_caption"))

    # --- Fuente de datos (CSV o API)
    data_source = st.radio(
        tr("data_source"),
        ["CSV", "API"],
        horizontal=True,
        key="trip_data_source",
        format_func=lambda x: tr(f"src_{x.lower()}"),
    )

    # --- Inputs comunes
    mode = st.radio(
        tr("trip_mode_label"),
        ["Trip Out", "Trip In"],
        horizontal=True,
        key="trip_mode",
        format_func=lambda m: tr("trip_out") if m == "Trip Out" else tr("trip_in"),
    )
    use_direction_filter = st.checkbox(
        tr("trip_dir_filter"),
        value=True,
        key="trip_dir_filter",
        help=tr("trip_dir_help"),
    )

    c1, c2, c3, c4 = st.columns(4)
    with c1:
        bin_m = st.number_input(tr("trip_bin"), 0.1, 20.0, 1.0, 0.1, key="trip_bin_m")
    with c2:
        baseline_q = st.slider(tr("trip_baseline"), 10, 90, 50, key="trip_baseline_q")
    with c3:
        thr = st.number_input(tr("trip_thr"), 0.0, 1e9, 0.0, key="trip_thr")
    with c4:
        rolling_window = st.number_input(tr("trip_rolling"), 5, 500, 60, 5, key="trip_rolling_window")

    envelope_method = st.radio(
        tr("trip_envelope"),
        [TRIP_ENV_EXACT, TRIP_ENV_BIN],
        horizontal=True,
        key="trip_envelope_method",
        format_func=lambda x: tr("trip_env_exact") if x == TRIP_ENV_EXACT else tr("trip_env_bin"),
        help=tr("trip_env_help"),
    )

    st.markdown(tr("trip_range_hdr"))
    st.caption(tr("trip_range_cap"))
    use_time_range = st.checkbox(tr("trip_use_time"), value=False, key="trip_use_time_range")
    col_from, col_to = st.columns(2)
    with col_from:
        trip_from_dt = st.text_input(
            tr("trip_from"),
            value="",
            placeholder="2025-10-20T14:00:00",
            key="trip_range_from",
            help=tr("trip_help_dt"),
        )
    with col_to:
        trip_to_dt = st.text_input(
            tr("trip_to"),
            value="",
            placeholder="2025-10-21T08:00:00",
            key="trip_range_to",
            help=tr("trip_help_dt"),
        )
    use_depth_range = st.checkbox(tr("trip_use_depth"), value=False, key="trip_use_depth_range")
    col_dmin, col_dmax = st.columns(2)
    with col_dmin:
        trip_depth_min = st.number_input(
            tr("trip_depth_min_l"),
            value=0.0,
            step=10.0,
            key="trip_depth_min",
            help=tr("trip_depth_help_min"),
        )
    with col_dmax:
        trip_depth_max = st.number_input(
            tr("trip_depth_max_l"),
            value=5000.0,
            step=10.0,
            key="trip_depth_max",
            help=tr("trip_depth_help_max"),
        )

    st.caption(tr("trip_gap_caption"))
    use_interpolation = st.checkbox(
        tr("trip_interp_chk"),
        value=False,
        key="trip_use_interpolation",
        help=tr("trip_interp_help"),
    )
    if use_interpolation:
        trip_interp_interval = st.number_input(
            tr("trip_interp_interval_l"),
            value=5,
            min_value=1,
            step=1,
            key="trip_interp_interval",
            help=tr("trip_interp_interval_help"),
        )

    st.markdown("#### Límites simulados (Trip In / Trip Out / Rotating)")
    st.caption("Opcional: carga un CSV con pesos/límites por profundidad para comparar el viaje real vs límite.")
    limits_file = st.file_uploader(
        "Upload CSV de límites (Depth, Trip Out, Trip In, Rotating)",
        type=["csv"],
        key="trip_limits_csv",
    )
    if limits_file is not None:
        try:
            raw_limits = pd.read_csv(limits_file, sep=",", low_memory=False)
            limits_df_parsed, err = _trip_parse_limits_csv(raw_limits)
            if err:
                st.warning(f"Límites: {err} Revisa las columnas del CSV.")
                st.session_state.pop("trip_limits_df", None)
            elif limits_df_parsed is not None and not limits_df_parsed.empty:
                st.session_state["trip_limits_df"] = limits_df_parsed
                st.success(f"Límites cargados: {len(limits_df_parsed):,} filas (Depth, Trip Out, Trip In, Rotating).")
        except Exception as e:
            st.warning(f"No se pudo leer el CSV de límites: {e}")
            st.session_state.pop("trip_limits_df", None)
    else:
        st.session_state.pop("trip_limits_df", None)


    st.markdown("#### Modelo Broomstick por FF (opcional)")
    st.caption(
        "Carga un CSV con familia de curvas por FF **o** genera curvas con el modelo simplificado (soft-string)."
    )
    ff_source = st.radio(
        "Origen del modelo Broomstick",
        ["CSV (simulador/Excel)", "Modelo simplificado (generar en app)"],
        horizontal=True,
        key="trip_ff_source",
    )

    if ff_source == "CSV (simulador/Excel)":
        ff_family_file = st.file_uploader(
            "Upload CSV de modelo por FF (familia PU/SO/ROT)",
            type=["csv"],
            key="trip_ff_family_csv",
        )
        if ff_family_file is not None:
            try:
                raw_model = pd.read_csv(ff_family_file, sep=",", low_memory=False)
                model_df, fam_map, err = parse_ff_family_csv(raw_model)
                if err:
                    st.warning(f"Modelo FF: {err} Revisa el formato/columnas del CSV.")
                    st.session_state.pop("trip_ff_family_df", None)
                    st.session_state.pop("trip_ff_family_map", None)
                else:
                    st.session_state["trip_ff_family_df"] = model_df
                    st.session_state["trip_ff_family_map"] = fam_map
                    all_ffs = sorted(set([*fam_map.get("PU", {}).keys(), *fam_map.get("SO", {}).keys(), *fam_map.get("ROT", {}).keys()]))
                    if all_ffs:
                        st.session_state["trip_ff_min"] = float(all_ffs[0])
                        st.session_state["trip_ff_max"] = float(all_ffs[-1])
                    st.success("Modelo broomstick por FF cargado desde CSV.")
            except Exception as e:
                st.warning(f"No se pudo leer el CSV de modelo por FF: {e}")
                st.session_state.pop("trip_ff_family_df", None)
                st.session_state.pop("trip_ff_family_map", None)
        else:
            st.session_state.pop("trip_ff_family_df", None)
            st.session_state.pop("trip_ff_family_map", None)
    else:
        # No borrar curvas generadas: así persisten al abrir la pestaña Broomstick y se dibujan
        with st.expander("Parámetros del modelo simplificado (soft-string)", expanded=True):
            st.caption(
                "Hookload = Depth × peso unitario flotado × (cos(incl.) ± FF×sin(incl.)). "
                "PU = Trip Out, SO = Trip In, ROT = solo peso axial."
            )
            c1, c2, c3 = st.columns(3)
            with c1:
                md_min_m = st.number_input("MD mínimo (m)", value=0.0, step=100.0, key="trip_model_md_min")
                md_max_m = st.number_input("MD máximo (m)", value=3000.0, step=100.0, key="trip_model_md_max")
                step_m = st.number_input("Paso de profundidad (m)", value=50.0, min_value=1.0, step=10.0, key="trip_model_step")
            with c2:
                weight_per_m_kg = st.number_input("Peso unitario (kg/m)", value=30.0, min_value=0.1, step=1.0, key="trip_model_wpm")
                buoyancy_factor = st.number_input("Factor de flotación (0–1)", value=0.85, min_value=0.5, max_value=1.0, step=0.01, key="trip_model_bf")
                inclination_deg = st.number_input("Inclinación promedio (°)", value=0.0, min_value=0.0, max_value=90.0, step=1.0, key="trip_model_inc")
                block_weight_klb = st.number_input(
                    "Peso del block (klb)",
                    value=0.0,
                    min_value=-500.0,
                    max_value=500.0,
                    step=5.0,
                    key="trip_model_block_klb",
                    help="Desplaza todas las curvas del modelo. Valor negativo baja las curvas (acercarlas a datos si quedan altas); positivo las sube.",
                )
            with c3:
                ff_min_in = st.number_input("FF mínimo", value=0.10, min_value=0.01, max_value=0.99, step=0.01, key="trip_model_ff_min")
                ff_max_in = st.number_input("FF máximo", value=0.50, min_value=0.01, max_value=0.99, step=0.01, key="trip_model_ff_max")
                ff_step_in = st.number_input("Paso de FF", value=0.05, min_value=0.01, step=0.01, key="trip_model_ff_step")
            with st.expander("Inclinación variable (opcional)", expanded=False):
                st.caption("Si rellenas ambos valores y son distintos, la inclinación varía linealmente con la profundidad y las curvas dejan de ser rectas.")
                inc_surf_deg = st.number_input("Incl. en superficie (°)", value=0.0, min_value=0.0, max_value=90.0, step=1.0, key="trip_model_inc_surf")
                inc_td_deg = st.number_input("Incl. en TD (°)", value=0.0, min_value=0.0, max_value=90.0, step=1.0, key="trip_model_inc_td")
            with st.expander("Cargar survey / wellplan (opcional)", expanded=False):
                st.caption(
                    "CSV o Excel con profundidad (MD) e inclinación. Si lo cargas, las curvas del modelo usan esta trayectoria "
                    "y dejan de ser rectas (inclinación real vs profundidad)."
                )
                survey_file = st.file_uploader(
                    "Survey o wellplan (CSV / Excel)",
                    type=["csv", "xlsx", "xls"],
                    key="trip_survey_csv",
                )
                trip_survey_df: pd.DataFrame | None = None
                if survey_file is not None:
                    try:
                        name = (survey_file.name or "").lower()
                        if name.endswith((".xlsx", ".xls")):
                            raw_survey = pd.read_excel(survey_file, sheet_name=0)
                        else:
                            raw_survey = pd.read_csv(survey_file, sep=",", low_memory=False)
                        raw_survey.columns = [str(c).strip() for c in raw_survey.columns]
                        cols_survey = raw_survey.columns.tolist()
                        depth_candidates = [c for c in cols_survey if any(x in c.lower() for x in ["depth", "md", "measured", "survey", "profundidad"])]
                        incl_candidates = [c for c in cols_survey if any(x in c.lower() for x in ["incl", "inclination", "inclinación"])]
                        idx_d = cols_survey.index(depth_candidates[0]) if depth_candidates else 0
                        idx_i = cols_survey.index(incl_candidates[0]) if incl_candidates else 0
                        survey_md_col = st.selectbox("Columna MD / Profundidad", cols_survey, index=idx_d, key="trip_survey_md_col")
                        survey_incl_col = st.selectbox("Columna Inclinación (°)", cols_survey, index=idx_i, key="trip_survey_incl_col")
                        if survey_md_col and survey_incl_col:
                            trip_survey_df = raw_survey[[survey_md_col, survey_incl_col]].copy()
                            trip_survey_df.columns = ["Depth", "Inclination"]
                            trip_survey_df["Depth"] = pd.to_numeric(trip_survey_df["Depth"], errors="coerce")
                            trip_survey_df["Inclination"] = pd.to_numeric(trip_survey_df["Inclination"], errors="coerce")
                            trip_survey_df = trip_survey_df.dropna().sort_values("Depth")
                            if len(trip_survey_df) >= 2:
                                st.session_state["trip_survey_df"] = trip_survey_df
                                st.success(f"Survey cargado: {len(trip_survey_df):,} estaciones (MD {trip_survey_df['Depth'].min():.0f}–{trip_survey_df['Depth'].max():.0f} m).")
                            else:
                                st.warning("Se necesitan al menos 2 filas válidas (Depth, Inclination).")
                                st.session_state.pop("trip_survey_df", None)
                    except Exception as e:
                        st.warning(f"No se pudo leer el survey: {e}")
                        st.session_state.pop("trip_survey_df", None)
                else:
                    st.session_state.pop("trip_survey_df", None)
            if st.button("Generar curvas por FF", type="primary", key="trip_model_generate"):
                ff_list = list(np.arange(ff_min_in, ff_max_in + ff_step_in * 0.5, ff_step_in))
                ff_list = [round(f, 2) for f in ff_list]
                if not ff_list:
                    ff_list = [ff_min_in]
                inc_surf = inc_surf_deg if abs(inc_surf_deg - inc_td_deg) > 0.01 else None
                inc_td = inc_td_deg if abs(inc_surf_deg - inc_td_deg) > 0.01 else None
                survey_for_model = st.session_state.get("trip_survey_df")
                if survey_for_model is not None and ("Depth" not in survey_for_model.columns or "Inclination" not in survey_for_model.columns):
                    survey_for_model = None
                model_df, fam_map = generate_broomstick_curves_simplified(
                    md_min_m=md_min_m,
                    md_max_m=md_max_m,
                    step_m=step_m,
                    weight_per_m_kg=weight_per_m_kg,
                    buoyancy_factor=buoyancy_factor,
                    inclination_deg=inclination_deg,
                    ff_values=ff_list,
                    block_weight_klb=block_weight_klb,
                    inclination_surface_deg=inc_surf,
                    inclination_td_deg=inc_td,
                    survey_df=survey_for_model,
                )
                if not model_df.empty:
                    st.session_state["trip_ff_family_df"] = model_df
                    st.session_state["trip_ff_family_map"] = fam_map
                    st.session_state["trip_ff_min"] = float(min(ff_list))
                    st.session_state["trip_ff_max"] = float(max(ff_list))
                    st.success(f"Curvas generadas: {len(model_df):,} puntos, FF desde {min(ff_list):.2f} hasta {max(ff_list):.2f}.")
                else:
                    st.error("Revisa MD min/máx y paso. No se generaron puntos.")

    df = pd.DataFrame()
    used_type_hk = ""
    used_type_dp = ""
    used_params_hk: dict = {}
    used_params_dp: dict = {}

    if data_source == "CSV":
        st.markdown("#### Cargar CSV")
        st.caption(
            "El CSV debe contener columnas de tiempo (timestamp/datetime o formato ISO), Hookload y profundidad (Bit depth/MD). "
            "Puedes asignar las columnas en «Configurar columnas»."
        )
        trip_file = st.file_uploader("Upload Tripping Traces (.csv)", type=["csv"], key="trip_csv_upload")
        if not trip_file:
            cached_df = st.session_state.get("trip_analysis_df")
            cached_src = st.session_state.get("trip_analysis_data_source")
            if cached_src == "CSV" and cached_df is not None and not cached_df.empty:
                df = cached_df.copy()
                used_type_hk = "CSV"
                used_type_dp = "CSV"
                used_params_hk = {}
                used_params_dp = {}
            else:
                st.info("Sube un CSV con columnas de tiempo, Hookload y profundidad, luego configura y presiona **Ejecutar análisis**.")
                return
        if trip_file:
            try:
                raw_df = pd.read_csv(trip_file, sep=",", low_memory=False)
            except Exception as e:
                st.error(f"No pude leer el CSV: {e}")
                return
            raw_df.columns = [str(c).strip() for c in raw_df.columns]
            cols = raw_df.columns.tolist()

            tcol_auto = _trip_pick_time_col(raw_df)
            hk_col_auto = _trip_pick_value_col(raw_df, prefer_names=["hookload", "hl", "load", "weight"])
            dp_col_auto = _trip_pick_value_col(raw_df, prefer_names=["bit depth", "depth", "md", "measured_depth", "block position"])

            with st.expander("Configurar columnas", expanded=True):
                idx_t = cols.index(tcol_auto) if tcol_auto in cols else 0
                idx_h = cols.index(hk_col_auto) if hk_col_auto in cols else 0
                idx_d = cols.index(dp_col_auto) if dp_col_auto in cols else 0
                tcol = st.selectbox(
                    "Columna Timestamp",
                    cols,
                    index=idx_t,
                    key="trip_col_timestamp",
                    help="Columna de fecha/hora (ej. Timestamp, YYYY-MM-DDTHH:MM:SS).",
                )
                hk_col = st.selectbox(
                    "Columna Hookload",
                    cols,
                    index=idx_h,
                    key="trip_col_hookload",
                    help="Columna de carga en gancho (kgf o klb).",
                )
                dp_col = st.selectbox(
                    "Columna Bit depth",
                    cols,
                    index=idx_d,
                    key="trip_col_bitdepth",
                    help="Columna de profundidad (m).",
                )

            if not st.button("▶️ Ejecutar análisis", key="trip_run"):
                return

            if tcol == dp_col or tcol == hk_col or dp_col == hk_col:
                st.error("Elija columnas distintas para Timestamp, Hookload y Bit depth.")
                return

            df = raw_df[[tcol, hk_col, dp_col]].copy()
            df = df.rename(columns={tcol: "Timestamp", hk_col: "Hookload", dp_col: "Bit depth"})
            df["Timestamp"] = pd.to_datetime(df["Timestamp"], errors="coerce")
            df["Hookload"] = _trip_to_numeric(df, "Hookload")
            df["Bit depth"] = _trip_to_numeric(df, "Bit depth")
            df = df.dropna(subset=["Timestamp", "Hookload", "Bit depth"]).sort_values("Timestamp")
            if df.empty:
                st.error(
                    "El CSV quedó vacío tras limpiar NaNs. Compruebe que Bit depth y Hookload tengan valores numéricos válidos."
                )
                return
            st.session_state["trip_analysis_df"] = df
            st.session_state["trip_analysis_data_source"] = "CSV"
            used_type_hk = "CSV"
            used_type_dp = "CSV"
            used_params_hk = {}
            used_params_dp = {}
    else:
        # --- API (flujo alineado con Ingeniería BHA: Proyecto → Pozo → Lateral → Trazas)
        with st.expander("Configuración API", expanded=False):
            base_url, token = render_solo_connection_ui(prefix="trip_api", label=tr("solo_expander"))
            if not token:
                st.warning("Falta SOLO_ACCESS_TOKEN. Configúralo en tu .env o pega el token aquí.")
            trip_traces_list_path = st.text_input(
                "Ruta listar trazas (opcional)",
                value="",
                help="Usa {well_uuid}. Ej: /public/api/v1/wells/{well_uuid}/drilling-traces",
                key="trip_api_traces_list_path",
            )
            trip_traces_list_params = st.text_input(
                "Params listar trazas (opcional)",
                value="",
                help="Formato: key=value&key2=value2 o JSON",
                key="trip_api_traces_list_params",
            )
            trip_laterals_list_path = st.text_input(
                "Ruta listar laterales (opcional)",
                value="",
                help="Usa {well_uuid} o {project_uuid}",
                key="trip_api_laterals_list_path",
            )
            trip_laterals_list_params = st.text_input(
                "Params listar laterales (opcional)",
                value="",
                key="trip_api_laterals_list_params",
            )
            trip_trace_types_path = st.text_input(
                "Ruta listar tipos de traza (opcional)",
                value="",
                key="trip_api_trace_types_path",
            )
            st.caption("Nota: evita /api/v1 si responde HTML; usa /public/api/v1.")

        if not base_url or not token:
            base_url, token = get_solo_credentials(prefix="trip_api")
        if not token:
            st.info("Ingresa un token válido para consultar la API.")
            return

        try:
            projects_resp = api_list_projects(base_url, token)
            projects = normalize_list_response(projects_resp)
        except Exception as e:
            st.error(f"No pude listar proyectos: {e}")
            return
        if not projects:
            st.info("No hay proyectos disponibles para este token.")
            return

        trip_project_map = {
            f"{p.get('name', 'Sin nombre')} ({p.get('uuid', 'n/a')})": p for p in projects
        }
        trip_project_label = st.selectbox("Proyecto", list(trip_project_map.keys()), key="trip_project_label_api")
        project_uuid = trip_project_map[trip_project_label].get("uuid")
        if not project_uuid:
            st.error("El proyecto seleccionado no tiene UUID.")
            return

        try:
            wells_resp = api_list_wells(base_url, token, project_uuid)
            wells = normalize_list_response(wells_resp)
        except Exception as e:
            st.error(f"No pude listar pozos: {e}")
            return
        if not wells:
            st.info("No hay pozos disponibles en este proyecto.")
            return

        trip_well_map = {
            f"{w.get('name', 'Sin nombre')} ({w.get('uuid', 'n/a')})": w for w in wells
        }
        trip_well_label = st.selectbox("Pozo", list(trip_well_map.keys()), key="trip_well_label_api")
        well_uuid = trip_well_map[trip_well_label].get("uuid")
        if not well_uuid:
            st.error("El pozo seleccionado no tiene UUID.")
            return

        laterals = []
        trip_no_laterals = False
        try:
            extra_lateral_params = parse_params_input(trip_laterals_list_params or "")
            laterals_resp = api_list_laterals(
                base_url,
                token,
                project_uuid=project_uuid,
                well_uuid=well_uuid,
                custom_path=trip_laterals_list_path or None,
                extra_params=extra_lateral_params,
            )
            laterals = normalize_list_response(laterals_resp)
        except Exception:
            laterals = []
        lateral_uuid = None
        if laterals:
            trip_lateral_map = {
                f"{l.get('name', 'Sin nombre')} ({l.get('uuid', 'n/a')})": l for l in laterals
            }
            trip_lateral_label = st.selectbox("Lateral", list(trip_lateral_map.keys()), key="trip_lateral_label")
            lateral_uuid = trip_lateral_map[trip_lateral_label].get("uuid")
            if not lateral_uuid:
                st.error("La lateral seleccionada no tiene UUID.")
                return
        else:
            lateral_uuid = well_uuid
            trip_no_laterals = True
            st.info("No se encontraron laterales para este pozo. Se usará el Pozo para consultar trazas.")

        trace_type = "TIME"
        st.caption("Dominio de trazas: **TIME** (tripping por tiempo).")

        st.markdown("**Rango y paginación para datos de traza**")
        with st.expander("Rango y paginación", expanded=True):
            trip_use_custom_range = st.checkbox(
                "Especificar rango (from / to)",
                value=st.session_state.get("trip_use_custom_range", True),
                key="trip_use_custom_range",
                help="Si no marcas esto, se puede usar auto-probing.",
            )
            trip_range_from = ""
            trip_range_to = ""
            if trip_use_custom_range:
                col_from, col_to = st.columns(2)
                with col_from:
                    trip_range_from = st.text_input(
                        "Desde (ISO 8601)",
                        value=st.session_state.get("trip_trace_from", "2020-12-29T08:00:00Z"),
                        key="trip_trace_from",
                        help="Rango corto (pocos días) suele dar pocos puntos. Para más datos usa un intervalo mayor (ej. una o varias semanas).",
                    )
                with col_to:
                    trip_range_to = st.text_input(
                        "Hasta (ISO 8601)",
                        value=st.session_state.get("trip_trace_to", "2020-12-29T08:50:00Z"),
                        key="trip_trace_to",
                        help="Ejemplo: 2026-02-15T23:59:59Z. Cuanto más amplio el rango Desde–Hasta, más puntos se cargarán.",
                    )
            trip_load_all_pages = st.checkbox(
                "Cargar todo el dato del rango (todas las páginas automáticamente)",
                value=True,
                key="trip_load_all_pages",
                help="La app pedirá página 0, 1, 2… hasta que no venga más dato. No hace falta elegir número de páginas.",
            )
            st.caption("Paginación: más «Tamaño de página» = más datos por petición. Si no usas «Cargar todo», indica cuántas páginas.")
            col_size, col_page, col_multi = st.columns(3)
            with col_size:
                trip_page_size = st.number_input(
                    "Tamaño de página (size)",
                    min_value=100,
                    max_value=50000,
                    value=10000,
                    step=500,
                    key="trip_api_page_size",
                    help="Más valor = más puntos por petición. Con «Cargar todo» se piden páginas 0, 1, 2… hasta agotar.",
                )
            with col_page:
                trip_page_number = st.number_input(
                    "Página (page)",
                    min_value=0,
                    max_value=100000,
                    value=0,
                    step=1,
                    key="trip_api_page_number",
                    disabled=trip_load_all_pages,
                    help="Solo si no usas «Cargar todo el dato del rango».",
                )
            with col_multi:
                trip_num_pages = st.number_input(
                    "Páginas a cargar",
                    min_value=1,
                    max_value=50,
                    value=3,
                    step=1,
                    key="trip_num_pages_load",
                    disabled=trip_load_all_pages,
                    help="Solo si no usas «Cargar todo». 1 = una página; 2+ = se unen.",
                )
        with st.expander("Params adicionales (avanzado)", expanded=False):
            trip_extra_params_raw = st.text_input(
                "Query string extra (opcional)",
                value=st.session_state.get("trip_extra_trace_params", ""),
                placeholder="Ej: sort=time&order=asc",
                key="trip_extra_trace_params",
            )
        trip_force_data_endpoint = st.checkbox(
            "Forzar endpoint de datos por pozo (time)",
            value=st.session_state.get("trip_force_data_endpoint", True),
            key="trip_force_data_endpoint",
        )
        trip_trip_auto_probe_data = st.checkbox(
            "Auto-detectar rango con datos (si la traza viene vacía)",
            value=True,
            key="trip_trip_auto_probe_data",
        )

        def build_trace_params_trip(page_override: int | None = None) -> dict:
            p: dict = {}
            if trip_use_custom_range and (trip_range_from or "").strip() and (trip_range_to or "").strip():
                raw_from = (trip_range_from or "").strip()
                raw_to = (trip_range_to or "").strip()
                p["from"] = _trip_iso_to_utc_z(raw_from)
                p["to"] = _trip_iso_to_utc_z(raw_to)
            p["size"] = int(trip_page_size)
            p["page"] = page_override if page_override is not None else int(trip_page_number)
            # Compatibilidad con APIs que usan offset/limit en lugar de page/size
            try:
                _sz = int(p.get("size", 0))
                _pg = int(p.get("page", 0))
                if "limit" not in p and _sz > 0:
                    p["limit"] = _sz
                if "offset" not in p and _sz > 0 and _pg >= 0:
                    p["offset"] = _pg * _sz
            except Exception:
                pass
            extra = parse_params_input(trip_extra_params_raw or "")
            if extra:
                for k, v in extra.items():
                    if k not in ("from", "to", "size", "page", "limit", "offset"):
                        p[k] = v
            return p

        mapped_scope_uuid = lateral_uuid
        mapped_scope_kind = "lateral"
        if trip_no_laterals:
            use_well_mapping_trip = st.checkbox(
                "Usar mapeo por pozo (si aplica)",
                value=True,
                key="trip_use_well_mapping",
            )
            if use_well_mapping_trip:
                mapped_scope_uuid = well_uuid
                mapped_scope_kind = "well"

        traces_trip: list[dict] = []
        trace_source_trip: str | None = None
        trip_mapped_err = trip_drill_err = trip_cat_err = None
        extra_trace_params = parse_params_input(trip_traces_list_params or "")

        if mapped_scope_uuid:
            try:
                traces_resp = api_list_mapped_traces(
                    base_url=base_url,
                    token=token,
                    scope_uuid=mapped_scope_uuid,
                    trace_type=trace_type,
                    scope_kind=mapped_scope_kind or "lateral",
                    custom_path=None,
                    extra_params=extra_trace_params,
                )
                traces_trip = normalize_list_response(traces_resp)
                if traces_trip:
                    trace_source_trip = "mapped"
            except Exception as e:
                trip_mapped_err = str(e)
                traces_trip = []
        if not traces_trip:
            try:
                traces_resp = api_list_drilling_traces(
                    base_url=base_url,
                    token=token,
                    well_uuid=well_uuid,
                    custom_path=trip_traces_list_path or None,
                    extra_params=extra_trace_params,
                )
                traces_trip = normalize_list_response(traces_resp)
                if traces_trip:
                    trace_source_trip = "drilling"
            except Exception as e:
                trip_drill_err = str(e)
                traces_trip = []
        if not traces_trip:
            try:
                traces_resp = api_list_traces_catalog(
                    base_url,
                    token,
                    trip_trace_types_path or None,
                    extra_trace_params,
                )
                traces_trip = normalize_list_response(traces_resp)
                if traces_trip:
                    trace_source_trip = "catalog"
            except Exception as e:
                trip_cat_err = str(e)
                traces_trip = []

        if not traces_trip:
            st.info(
                "No se encontraron trazas (mapeadas, drilling ni catálogo). "
                "Selecciona otra lateral/pozo o revisa permisos."
            )
            with st.expander("Detalles de error"):
                if trip_mapped_err:
                    st.write(f"Mapeadas: {trip_mapped_err}")
                if trip_drill_err:
                    st.write(f"Drilling: {trip_drill_err}")
                if trip_cat_err:
                    st.write(f"Catálogo: {trip_cat_err}")
            return

        trip_trace_map_by_uuid = {}
        trip_label_to_uuid = {}
        for t in traces_trip:
            label = t.get("name") or t.get("label") or t.get("type") or t.get("uuid") or "trace"
            trace_uuid = t.get("uuid") or t.get("id")
            if trace_uuid:
                trip_trace_map_by_uuid[str(trace_uuid)] = {"label": label, "trace": t}
                trip_label_to_uuid[label] = str(trace_uuid)

        trip_trace_filter = st.text_input(
            "Filtro de trazas (nombre)",
            value="",
            help="Ej: Hookload, Depth, Bit depth.",
            key="trip_trace_filter",
        ).strip()
        if trip_trace_filter:
            trip_trace_map_by_uuid = {
                k: v for k, v in trip_trace_map_by_uuid.items()
                if trip_trace_filter.lower() in (v.get("label") or "").lower()
            }
            trip_label_to_uuid = {v["label"]: k for k, v in trip_trace_map_by_uuid.items()}

        if not trip_trace_map_by_uuid:
            st.info("No hay trazas que coincidan con el filtro.")
            return

        def suggest_trip_traces(trace_map_by_uuid: dict) -> tuple[str | None, str | None]:
            hookload_kw = ["hookload", "hl", "load", "weight", "tension"]
            depth_kw = ["depth", "md", "bit depth", "block", "measured", "block position"]
            hookload_uuid = depth_uuid = None
            for uuid, data in trace_map_by_uuid.items():
                label = (data.get("label") or "").lower()
                if not hookload_uuid:
                    for kw in hookload_kw:
                        if kw in label:
                            hookload_uuid = uuid
                            break
                if not depth_uuid:
                    for kw in depth_kw:
                        if kw in label:
                            depth_uuid = uuid
                            break
                if hookload_uuid and depth_uuid:
                    break
            return hookload_uuid, depth_uuid

        suggested_hk, suggested_dp = suggest_trip_traces(trip_trace_map_by_uuid)
        labels_sorted = sorted(set(v["label"] for v in trip_trace_map_by_uuid.values()))
        idx_hk = 0
        idx_dp = 0
        if suggested_hk and suggested_dp:
            lab_hk = trip_trace_map_by_uuid.get(suggested_hk, {}).get("label")
            lab_dp = trip_trace_map_by_uuid.get(suggested_dp, {}).get("label")
            if lab_hk and lab_hk in labels_sorted:
                idx_hk = labels_sorted.index(lab_hk)
            if lab_dp and lab_dp in labels_sorted:
                idx_dp = labels_sorted.index(lab_dp)

        st.markdown("**Configurar trazas**")
        st.caption("Elige la traza que corresponde a Hookload y la que corresponde a Bit depth.")
        col_hl, col_dp = st.columns(2)
        with col_hl:
            selected_hookload_label = st.selectbox(
                "Traza Hookload",
                options=labels_sorted,
                index=min(idx_hk, len(labels_sorted) - 1) if labels_sorted else 0,
                key="trip_hookload_trace_label",
            )
        with col_dp:
            selected_depth_label = st.selectbox(
                "Traza Bit depth",
                options=labels_sorted,
                index=min(idx_dp, len(labels_sorted) - 1) if labels_sorted else 0,
                key="trip_depth_trace_label",
            )
        hookload_trace_uuid = trip_label_to_uuid.get(selected_hookload_label)
        depth_trace_uuid = trip_label_to_uuid.get(selected_depth_label)
        if not hookload_trace_uuid or not depth_trace_uuid:
            st.error("No se pudo resolver UUID de Hookload o Bit depth.")
            return
        if hookload_trace_uuid == depth_trace_uuid:
            st.warning("Elige trazas distintas para Hookload y Bit depth.")
            return

        st.markdown("**Actualización en tiempo real (API)**")
        trip_auto_refresh = st.checkbox(
            "Actualizar automáticamente cada 30 s",
            value=st.session_state.get("trip_auto_refresh", False),
            key="trip_auto_refresh",
            help="Vuelve a consultar la API y refrescar gráficas cada X segundos. Desmarca para detener.",
        )
        if trip_auto_refresh:
            st.number_input(
                "Intervalo (segundos)",
                min_value=10,
                max_value=300,
                value=30,
                step=10,
                key="trip_auto_refresh_interval",
                help="Cada cuántos segundos se vuelve a ejecutar el análisis (10–300 s).",
            )

        auto_rerun = st.session_state.pop("trip_auto_rerun_trigger", False)
        run_clicked = st.button("▶️ Ejecutar análisis", key="trip_run")
        use_cached_api = (
            not (run_clicked or (trip_auto_refresh and auto_rerun))
            and st.session_state.get("trip_analysis_data_source") == "API"
            and st.session_state.get("trip_analysis_df") is not None
            and not st.session_state.get("trip_analysis_df").empty
        )
        if not (run_clicked or (trip_auto_refresh and auto_rerun)) and not use_cached_api:
            st.info("Configura Proyecto, Pozo, trazas y presiona **Ejecutar análisis**.")
            return
        if use_cached_api:
            df = st.session_state["trip_analysis_df"].copy()
            used_type_hk = st.session_state.get("trip_used_type_hk", "TIME")
            used_type_dp = st.session_state.get("trip_used_type_dp", "TIME")
            used_params_hk = st.session_state.get("trip_used_params_hk", {})
            used_params_dp = st.session_state.get("trip_used_params_dp", {})
        else:
            if trip_use_custom_range and (trip_range_from or "").strip() and (trip_range_to or "").strip():
                try:
                    dt_from = pd.to_datetime((trip_range_from or "").strip(), errors="coerce")
                    dt_to = pd.to_datetime((trip_range_to or "").strip(), errors="coerce")
                    if pd.notna(dt_from) and pd.notna(dt_to) and dt_from > dt_to:
                        st.error(
                            "El rango no es válido: **Desde** es posterior a **Hasta**. "
                            "Revisa las fechas (ej: Desde 2026-02-07T08:00:00Z, Hasta 2026-09-07T08:50:00Z)."
                        )
                        return
                except Exception:
                    pass

            page_size_int = int(trip_page_size)
            max_auto_pages = 50
            if trip_load_all_pages:
                page_indices_trip = []  # se llenará en el bucle hasta que venga menos de size
            else:
                page_indices_trip = list(range(int(trip_num_pages))) if int(trip_num_pages) > 1 else [int(trip_page_number)]
            list_hk_pages: list[pd.DataFrame] = []
            list_dp_pages: list[pd.DataFrame] = []
            used_type_hk = used_type_dp = "TIME"
            used_params_hk = used_params_dp = {}
            try:
                page_idx = 0
                while True:
                    if not trip_load_all_pages and page_idx >= len(page_indices_trip):
                        break
                    if trip_load_all_pages and page_idx >= max_auto_pages:
                        break
                    current_page = page_indices_trip[page_idx] if not trip_load_all_pages else page_idx
                    user_params = build_trace_params_trip(page_override=current_page)
                    d_hk, used_type_hk, used_params_hk = probe_well_trace_data(
                        base_url=base_url,
                        token=token,
                        well_uuid=well_uuid,
                        trace_uuid=hookload_trace_uuid,
                        prefer_type="TIME",
                        user_params=user_params,
                    )
                    d_dp, used_type_dp, used_params_dp = probe_well_trace_data(
                        base_url=base_url,
                        token=token,
                        well_uuid=well_uuid,
                        trace_uuid=depth_trace_uuid,
                        prefer_type="TIME",
                        user_params=user_params,
                    )
                    if not d_hk.empty:
                        list_hk_pages.append(d_hk)
                    if not d_dp.empty:
                        list_dp_pages.append(d_dp)
                    if d_hk.empty and d_dp.empty:
                        break
                    if trip_load_all_pages:
                        n_hk = len(d_hk)
                        n_dp = len(d_dp)
                        # FIX: do NOT stop just because len < size (API may cap page size)
                        if n_hk == 0 and n_dp == 0:
                            break

                    page_idx += 1
                df_hk = pd.concat(list_hk_pages, ignore_index=True) if list_hk_pages else pd.DataFrame()
                df_dp = pd.concat(list_dp_pages, ignore_index=True) if list_dp_pages else pd.DataFrame()
                if len(list_hk_pages) > 1 and not df_hk.empty:
                    tcol = _trip_pick_time_col(df_hk)
                    if tcol:
                        df_hk = df_hk.drop_duplicates(subset=[tcol], keep="first").sort_values(tcol).reset_index(drop=True)
                if len(list_dp_pages) > 1 and not df_dp.empty:
                    tcol = _trip_pick_time_col(df_dp)
                    if tcol:
                        df_dp = df_dp.drop_duplicates(subset=[tcol], keep="first").sort_values(tcol).reset_index(drop=True)
            except Exception as e:
                st.error(f"No pude leer trazas: {e}")
                return

            if df_hk.empty or df_dp.empty:
                st.error(
                    "Hookload o Depth vienen vacíos con los params/rango actuales. "
                    "Si usaste un rango de fechas, revisa que **Hasta** y **Desde** tengan el año correcto (ej. **2026**, no 026)."
                )
                if not df_hk.empty:
                    st.caption(f"Hookload OK ({used_type_hk}) params: {used_params_hk}")
                if not df_dp.empty:
                    st.caption(f"Depth OK ({used_type_dp}) params: {used_params_dp}")
                return

            df_hk.columns = [str(c).strip() for c in df_hk.columns]
            df_dp.columns = [str(c).strip() for c in df_dp.columns]
            tcol_hk = _trip_pick_time_col(df_hk)
            vcol_hk = _trip_pick_value_col(df_hk, prefer_names=["hookload", "hl", "load", "value"])
            tcol_dp = _trip_pick_time_col(df_dp)
            vcol_dp = _trip_pick_value_col(df_dp, prefer_names=["bit depth", "depth", "md", "measured_depth", "value"])

            if not tcol_hk or not vcol_hk or not tcol_dp or not vcol_dp:
                st.error("No pude inferir columnas de tiempo/valor en las trazas.")
                st.write("Hookload cols:", df_hk.columns.tolist())
                st.write("Depth cols:", df_dp.columns.tolist())
                return

            hk = df_hk[[tcol_hk, vcol_hk]].rename(columns={tcol_hk: "Timestamp", vcol_hk: "Hookload"}).copy()
            dp = df_dp[[tcol_dp, vcol_dp]].rename(columns={tcol_dp: "Timestamp", vcol_dp: "Bit depth"}).copy()
            hk["Timestamp"] = pd.to_datetime(hk["Timestamp"], errors="coerce")
            dp["Timestamp"] = pd.to_datetime(dp["Timestamp"], errors="coerce")
            hk["Hookload"] = _trip_to_numeric(hk, "Hookload")
            dp["Bit depth"] = _trip_to_numeric(dp, "Bit depth")
            hk_before = len(hk)
            dp_before = len(dp)
            hk = hk.dropna(subset=["Timestamp", "Hookload"]).sort_values("Timestamp")
            dp = dp.dropna(subset=["Timestamp", "Bit depth"]).sort_values("Timestamp")

            if hk.empty or dp.empty:
                st.error("Los datos quedaron vacíos después de limpiar NaNs / parseo de tiempo.")
                # Diagnóstico para ver por qué se vació
                with st.expander("Diagnóstico (por qué quedaron vacíos)", expanded=True):
                    st.markdown("**Traza Hookload**")
                    st.caption(f"Filas recibidas de la API: {len(df_hk):,}. Columnas usadas: tiempo → «{tcol_hk}», valor → «{vcol_hk}».")
                    st.write(
                        f"Tras parsear tiempo y valor: {hk_before:,} filas. Tras quitar NaNs en Timestamp/Hookload: **{len(hk):,}** filas."
                    )
                    if hk_before > 0 and len(hk) == 0:
                        sample = df_hk[[tcol_hk, vcol_hk]].head(3)
                        st.write("Muestra de los primeros valores (antes de renombrar):", sample)
                        st.caption("Si todos los tiempos son NaT, revisa el formato de fecha (ISO 8601 o compatible). Si Hookload queda NaN, revisa que la columna sea numérica.")
                    st.markdown("**Traza Bit depth**")
                    st.caption(f"Filas recibidas de la API: {len(df_dp):,}. Columnas usadas: tiempo → «{tcol_dp}», valor → «{vcol_dp}».")
                    st.write(
                        f"Tras parsear tiempo y valor: {dp_before:,} filas. Tras quitar NaNs en Timestamp/Bit depth: **{len(dp):,}** filas."
                    )
                    if dp_before > 0 and len(dp) == 0:
                        sample_dp = df_dp[[tcol_dp, vcol_dp]].head(3)
                        st.write("Muestra de los primeros valores (antes de renombrar):", sample_dp)
                        st.caption("Si todos los tiempos son NaT o Bit depth queda NaN, revisa formato y que la columna de profundidad sea numérica.")
                return

            df = pd.merge_asof(
                hk,
                dp,
                on="Timestamp",
                direction="nearest",
                tolerance=pd.Timedelta("2s"),
            ).dropna(subset=["Hookload", "Bit depth"])

            if df.empty:
                st.error("No pude alinear Hookload y Depth por tiempo (merge_asof).")
                return

            if len(df) < 100:
                st.warning(
                    f"**API:** Tras unir Hookload y Bit depth por tiempo quedaron **{len(df):,}** filas "
                    f"(Hookload trajo {len(hk):,}, Bit depth trajo {len(dp):,}). "
                    "Si una traza trae muy pocas filas, la API está devolviendo poco dato para ese rango: marca **«Cargar todo el dato del rango»**, "
                    "amplía **Desde/Hasta** y vuelve a ejecutar."
                )

    # --- Aplicar rango de análisis, interpolación, dirección y preprocesamiento (común CSV y API)
    if use_time_range and (trip_from_dt or trip_to_dt):
        from_dt = pd.to_datetime(trip_from_dt, errors="coerce", utc=True) if trip_from_dt else None
        to_dt = pd.to_datetime(trip_to_dt, errors="coerce", utc=True) if trip_to_dt else None
        if from_dt is not None and pd.notna(from_dt):
            df = df[df["Timestamp"] >= from_dt]
        if to_dt is not None and pd.notna(to_dt):
            df = df[df["Timestamp"] <= to_dt]
        if df.empty:
            st.warning("No quedan datos tras aplicar el rango de tiempo. Revisa Desde/Hasta.")
            return
    if use_depth_range:
        depth_min_val = float(trip_depth_min)
        depth_max_val = float(trip_depth_max)
        actual_min = df["Bit depth"].min()
        actual_max = df["Bit depth"].max()
        df = df[(df["Bit depth"] >= depth_min_val) & (df["Bit depth"] <= depth_max_val)]
        if df.empty:
            st.warning(
                "No quedan datos tras aplicar el rango de profundidad. "
                f"Pediste **{depth_min_val:.0f}–{depth_max_val:.0f} m**, pero en los datos **Bit depth** va de **{actual_min:.1f}** a **{actual_max:.1f}**. "
                "Revisa el rango mín/máx (m)."
            )
            return
    if use_interpolation and len(df) >= 2:
        interval_sec = int(trip_interp_interval) if trip_interp_interval and trip_interp_interval >= 1 else 5
        df = df.sort_values("Timestamp").drop_duplicates(subset=["Timestamp"], keep="first")
        t_min = df["Timestamp"].min()
        t_max = df["Timestamp"].max()
        tz = getattr(t_min, "tzinfo", None)
        new_index = pd.date_range(
            start=t_min, end=t_max, freq=pd.Timedelta(seconds=interval_sec), inclusive="both", tz=tz
        )
        df_idx = df.set_index("Timestamp")
        df_reindexed = df_idx.reindex(new_index)
        df_reindexed = df_reindexed.interpolate(method="linear").ffill().bfill()
        df = df_reindexed.reset_index()
        if df.columns[0] != "Timestamp":
            df = df.rename(columns={df.columns[0]: "Timestamp"})
        df = df[["Timestamp", "Hookload", "Bit depth"]].dropna(subset=["Hookload", "Bit depth"])
    if use_direction_filter:
        df = df.sort_values("Timestamp")
        df["_d"] = df["Bit depth"].diff()
        if mode == "Trip Out":
            df = df[df["_d"] < 0]
        else:
            df = df[df["_d"] > 0]
        df = df.drop(columns=["_d"], errors="ignore")
        if df.empty:
            st.warning(
                "Sin datos después de filtrar por dirección del viaje. "
                "Desmarca «Filtrar por dirección del viaje» o amplía el rango."
            )
            return
    df = _trip_normalize_and_clean(df)
    if df.empty:
        st.warning("Sin datos tras preprocesamiento.")
        return
    if data_source == "API":
        st.session_state["trip_analysis_df"] = df
        st.session_state["trip_analysis_data_source"] = "API"
        st.session_state["trip_used_type_hk"] = used_type_hk
        st.session_state["trip_used_type_dp"] = used_type_dp
        st.session_state["trip_used_params_hk"] = used_params_hk
        st.session_state["trip_used_params_dp"] = used_params_dp

    # --- Envelope por profundidad exacta (max Hookload por Bit depth), si aplica
    envelope_exact = pd.DataFrame()
    if envelope_method == TRIP_ENV_EXACT:
        envelope_exact = _trip_build_exact_depth_envelope(df)

    # --- Envelope por bin de profundidad
    df["Depth_bin"] = (df["Bit depth"] / float(bin_m)).round().astype(int) * float(bin_m)

    grp = df.groupby("Depth_bin")["Hookload"]
    env = grp.max().to_frame("Hookload_max")
    env["Hookload_baseline"] = grp.quantile(float(baseline_q) / 100.0)
    env["Overpull"] = env["Hookload_max"] - env["Hookload_baseline"]
    env = env.reset_index().sort_values("Depth_bin")

    if thr > 0:
        env["Event"] = env["Overpull"] >= float(thr)
    else:
        env["Event"] = env["Overpull"] > 0

    tab_env, tab_broom = st.tabs([tr("tab_trip_env"), tr("tab_trip_broom")])

    with tab_env:
        st.markdown("### Resultados")
        n_raw = len(df) if isinstance(df, pd.DataFrame) else 0
        span_text = ""
        # Duración real de los datos cargados (para diagnóstico cuando "solo salen X min")
        if "Timestamp" in df.columns and n_raw > 0:
            t_min = df["Timestamp"].min()
            t_max = df["Timestamp"].max()
            duration_sec = (t_max - t_min).total_seconds()
            duration_min = duration_sec / 60.0
            duration_h = duration_sec / 3600.0
            if duration_h >= 1:
                span_text = f"duración **{duration_h:.1f} h**"
            else:
                span_text = f"duración **{duration_min:.0f} min**"
            if data_source == "API" and duration_min > 0 and duration_min < 60:
                st.warning(
                    f"Los datos cargados abarcan solo **{duration_min:.0f} min** (desde {t_min} hasta {t_max}). "
                    "Para cargar **todo el rango** que definiste en Desde/Hasta: marca **«Cargar todo el dato del rango (todas las páginas automáticamente)»** "
                    "en «Rango y paginación» y vuelve a ejecutar."
                )
        st.caption(
            f"**Puntos cargados (tras filtros):** {n_raw:,}. "
            + (f"Rango temporal: {span_text}. " if span_text else "")
            + (
                f"El envelope tiene **{len(envelope_exact):,}** filas (una por profundidad distinta: max Hookload por depth)."
                if not envelope_exact.empty
                else f"Envelope por bin: **{len(env):,}** rangos de profundidad."
            )
        )
        if data_source == "CSV":
            st.caption("Fuente: CSV. Columnas: tiempo, Hookload, profundidad.")
        else:
            st.caption(
                f"Hookload: {used_type_hk}; Depth: {used_type_dp}. "
                "Si el rango no fue especificado, el probing pudo ajustar from/to."
            )
        with st.expander("Preprocesamiento aplicado", expanded=False):
            range_lines = []
            if use_time_range and (trip_from_dt or trip_to_dt):
                range_lines.append(f"- **Rango de tiempo:** desde {trip_from_dt or '—'} hasta {trip_to_dt or '—'}\n")
            if use_depth_range:
                range_lines.append(f"- **Rango de profundidad:** {trip_depth_min} m – {trip_depth_max} m\n")
            if use_interpolation:
                range_lines.append(
                    f"- **Interpolación:** huecos rellenados con interpolación lineal (intervalo {trip_interp_interval or 5} s)\n"
                )
            st.markdown(
                "".join(range_lines)
                + "- **Columnas:** nombres limpiados (strip)\n"
                "- **Bit depth y Hookload:** forzados a numérico\n"
                "- **Filas:** eliminadas donde Bit depth o Hookload son NaN\n"
                "- **Timestamp:** convertido a datetime (si existe)\n"
                f"- **Dirección:** filtro por viaje **{mode}** (derivada de profundidad)\n"
                + (
                    "- **Envelope por profundidad exacta:** max Hookload por cada Bit depth (una fila por depth)\n"
                    if not envelope_exact.empty
                    else "- **Envelope por bin:** agrupación por rangos de profundidad\n"
                )
            )
            st.caption(f"Filas tras preprocesamiento: {len(df):,}. Columnas: {list(df.columns)}.")

        # --- Gráfico principal de envelope: exacta (por Bit depth) o por bin según método elegido
        fig_exact = None
        n_env = 0
        if not envelope_exact.empty:
            st.markdown("#### Envelope Hookload vs profundidad (exacta)")
            st.caption(
                "Una fila por Bit depth: la de **mayor Hookload** en esa profundidad (curva de carga máxima del viaje). "
                "En conexiones, donde la profundidad se mantiene fija y el Hookload varía, se toma el máximo Hookload en esa profundidad."
            )
            n_env = len(envelope_exact)
            n_distinct_depths = df["Bit depth"].nunique() if "Bit depth" in df.columns else n_env
            if n_env <= 2:
                if n_raw >= 200:
                    st.info(
                        f"La envolvente tiene **{n_env}** punto(s) porque en los datos hay solo **{n_distinct_depths}** profundidad(es) distinta(s) (Bit depth) "
                        f"con **{n_raw:,}** puntos crudos. "
                        "En **conexiones** (profundidad fija, Hookload variable) es normal: se toma el **máximo Hookload** en esa profundidad. "
                        "Para ver una **curva con muchas profundidades**: amplía el rango **Desde/Hasta** en «Rango y paginación» para incluir un tramo del viaje donde el Bit depth sí cambie, "
                        "o revisa que la traza **Bit depth** desde la API devuelva más de un valor en el rango."
                    )
                else:
                    st.info(
                        f"La envolvente tiene solo **{n_env}** punto(s) (una fila por profundidad distinta). "
                        "Para ver una **curva completa**: (1) En «Rango y paginación» marca **«Cargar todo el dato del rango»** para traer todas las páginas; "
                        "(2) amplía **Desde/Hasta** si hace falta y vuelve a ejecutar."
                    )
            marker_size = 14 if n_env <= 2 else 3
            fig_exact = px.line(
                envelope_exact,
                x="Bit depth",
                y="Hookload",
                title=f"Hookload envelope – {mode} (profundidad exacta)<br><sub style='font-size:10px;color:#64748b;'>Max Hookload por profundidad · Líneas punteadas = límites simulados (si aplica)</sub>",
                labels={"Bit depth": "Profundidad MD (m)", "Hookload": "Hookload (klb)"},
            ).update_layout(xaxis_autorange="reversed")
            fig_exact.update_traces(selector=dict(name="Hookload"), mode="lines+markers", marker=dict(size=marker_size))
        elif not env.empty:
            # Método por bin: no hay envelope_exact; usar agregado por Depth_bin (evita NameError en fig_exact / n_env)
            st.markdown("#### Envelope Hookload vs profundidad (por bin)")
            st.caption(
                f"Agregación por bin de **{bin_m:.2f} m**: se muestra el **máximo Hookload** por rango de profundidad."
            )
            n_env = len(env)
            marker_size = 14 if n_env <= 2 else 3
            fig_exact = px.line(
                env,
                x="Depth_bin",
                y="Hookload_max",
                title=f"Hookload envelope – {mode} (por bin)<br><sub style='font-size:10px;color:#64748b;'>Max Hookload por bin de profundidad · Límites simulados opcionales</sub>",
                labels={"Depth_bin": "Profundidad (m)", "Hookload_max": "Hookload max (klb)"},
            ).update_layout(xaxis_autorange="reversed")
            fig_exact.update_traces(selector=dict(name="Hookload_max"), mode="lines+markers", marker=dict(size=marker_size))
        else:
            st.warning("No hay datos para dibujar el envelope (exacto ni por bin).")

        # --- Overlay: límites simulados (una curva continua por límite, sin efecto "peine")
        limits_df_overlay = st.session_state.get("trip_limits_df")
        if fig_exact is not None and limits_df_overlay is not None and not limits_df_overlay.empty and "Depth" in limits_df_overlay.columns:
            try:
                lim = _trip_prepare_limits_for_continuous_line(limits_df_overlay, rolling_window=25)
                if not lim.empty:
                    # Una sola traza por curva: x=Depth, y=valor límite (Hookload)
                    if "TripOut" in lim.columns and lim["TripOut"].notna().any():
                        fig_exact.add_trace(
                            go.Scatter(
                                x=lim["Depth"].values,
                                y=lim["TripOut"].values,
                                mode="lines",
                                name="Límite simulado – Trip Out",
                                line=dict(dash="dash", width=2),
                                connectgaps=True,
                                hovertemplate="MD: %{x:.0f} m<br>Límite: %{y:.1f}<extra></extra>",
                            )
                        )
                    if "TripIn" in lim.columns and lim["TripIn"].notna().any():
                        fig_exact.add_trace(
                            go.Scatter(
                                x=lim["Depth"].values,
                                y=lim["TripIn"].values,
                                mode="lines",
                                name="Límite simulado – Trip In",
                                line=dict(dash="dash", width=2),
                                connectgaps=True,
                                hovertemplate="MD: %{x:.0f} m<br>Límite: %{y:.1f}<extra></extra>",
                            )
                        )
                    if "Rotating" in lim.columns and lim["Rotating"].notna().any():
                        fig_exact.add_trace(
                            go.Scatter(
                                x=lim["Depth"].values,
                                y=lim["Rotating"].values,
                                mode="lines",
                                name="Límite simulado – Rotating",
                                line=dict(width=2, shape="spline", smoothing=1.3),
                                connectgaps=True,
                                hovertemplate="MD: %{x:.0f} m<br>Límite: %{y:.1f}<extra></extra>",
                            )
                        )
            except Exception as _e_lim:
                st.warning(f"No pude superponer límites simulados en el gráfico: {_e_lim}")

        # --- Ajustes de legibilidad (tema + layout)
        if fig_exact is not None:
            use_dark_env = is_streamlit_dark_mode()
            # Si el envelope es grande, evita markers (se ve como "peine")
            if n_env > 400:
                try:
                    fig_exact.update_traces(mode="lines")
                except Exception:
                    pass
            fig_exact.update_layout(
                legend=dict(orientation="h", yanchor="bottom", y=1.10, xanchor="right", x=1),
                hovermode="x unified",
            )
            fig_exact.update_xaxes(autorange="reversed", title="Profundidad MD (m)")
            fig_exact.update_yaxes(title="Hookload (klb)")

            fig_exact = prettify_auto(fig_exact, h=480) if "prettify_auto" in globals() else (apply_pro_theme_dark(fig_exact, h=480) if use_dark_env else apply_pro_theme(fig_exact, h=480))
            st.plotly_chart(fig_exact, use_container_width=True, config=PLOTLY_CONFIG)

        if not envelope_exact.empty:
            st.dataframe(envelope_exact.head(100), use_container_width=True, hide_index=True)
            st.caption(
                f"Mostrando hasta 100 filas de {len(envelope_exact):,} (envelope = una fila por profundidad distinta). "
                f"Puntos crudos en el rango: **{n_raw:,}**."
                + (f" Para más puntos, amplía el rango **Desde/Hasta** en «Rango y paginación»." if n_raw < 500 and data_source == "API" else "")
            )
        elif not env.empty:
            st.dataframe(env.head(100), use_container_width=True, hide_index=True)
            st.caption(
                f"Mostrando hasta 100 filas de {len(env):,} (envelope por bin de {bin_m:.2f} m). "
                f"Puntos crudos en el rango: **{n_raw:,}**."
                + (f" Para más puntos, amplía el rango **Desde/Hasta** en «Rango y paginación»." if n_raw < 500 and data_source == "API" else "")
            )

        # --- Vista Pro: overpull en tiempo + severidad
        df_with_events, t0 = _trip_compute_time_overpull(
            df,
            rolling_window=max(5, int(rolling_window)),
            overpull_thr=float(thr),
        )
        mode_label = "Trip Out" if mode == "Trip Out" else "Trip In"
        dark_pro = st.checkbox("Tema oscuro (vista pro)", value=is_streamlit_dark_mode(), key="trip_dark_pro")

        # --- Panel de KPIs
        total_hrs = (df["Timestamp"].max() - df["Timestamp"].min()).total_seconds() / 3600.0
        max_hl = float(df["Hookload"].max())
        ev_t = df_with_events[df_with_events["Event_t"]]
        n_ev = len(ev_t)
        n_low = (ev_t["Severity"] == "low").sum()
        n_med = (ev_t["Severity"] == "medium").sum()
        n_high = (ev_t["Severity"] == "high").sum()
        pct_overpull = 100.0 * n_ev / len(df_with_events) if len(df_with_events) > 0 else 0.0
        max_overpull_t = float(ev_t["Overpull_t"].max()) if n_ev > 0 else 0.0
        max_overpull_env = float(env["Overpull"].max()) if not env.empty else 0.0

        st.markdown("#### KPIs del análisis")
        pct_bar = min(100, max(0, pct_overpull))
        st.markdown(
            f"""
            <div style="display: flex; flex-wrap: wrap; gap: 0.6rem; align-items: center; margin-bottom: 0.5rem;">
                <span style="
                    background: linear-gradient(135deg, #0f172a 0%, #1e293b 100%);
                    color: #e2e8f0; font-size: 0.8rem; font-weight: 600;
                    padding: 0.4rem 0.75rem; border-radius: 999px; letter-spacing: 0.02em;
                    box-shadow: 0 1px 2px rgba(0,0,0,0.2);
                "><span style="color: #94a3b8; font-weight: 500;">Tiempo</span> {total_hrs:.1f} h</span>
                <span style="
                    background: linear-gradient(135deg, #1e3a5f 0%, #2563eb 100%);
                    color: #fff; font-size: 0.8rem; font-weight: 600;
                    padding: 0.4rem 0.75rem; border-radius: 999px; letter-spacing: 0.02em;
                    box-shadow: 0 1px 2px rgba(0,0,0,0.2);
                "><span style="opacity: 0.9; font-weight: 500;">Hookload máx.</span> {max_hl:,.0f}</span>
                <span style="
                    background: linear-gradient(135deg, #422006 0%, #f59e0b 100%);
                    color: #fff; font-size: 0.8rem; font-weight: 600;
                    padding: 0.4rem 0.75rem; border-radius: 999px; letter-spacing: 0.02em;
                    box-shadow: 0 1px 2px rgba(0,0,0,0.2);
                "><span style="opacity: 0.9; font-weight: 500;">Overpull máx.</span> {max(max_overpull_t, max_overpull_env):,.1f}</span>
                <span style="
                    background: linear-gradient(135deg, #450a0a 0%, #b91c1c 50%, #ef4444 100%);
                    color: #fff; font-size: 0.8rem; font-weight: 600;
                    padding: 0.4rem 0.75rem; border-radius: 999px; letter-spacing: 0.02em;
                    box-shadow: 0 1px 2px rgba(0,0,0,0.2);
                "><span style="opacity: 0.9; font-weight: 500;">Eventos</span> {n_low} / {n_med} / {n_high}</span>
                <span style="
                    background: linear-gradient(135deg, #14532d 0%, #166534 50%, #22c55e 100%);
                    color: #fff; font-size: 0.8rem; font-weight: 600;
                    padding: 0.4rem 0.75rem; border-radius: 999px; letter-spacing: 0.02em;
                    box-shadow: 0 1px 2px rgba(0,0,0,0.2);
                    min-width: 140px;
                "><span style="opacity: 0.9; font-weight: 500;">% overpull</span> {pct_overpull:.1f}%
                    <div style="margin-top: 4px; height: 4px; background: rgba(255,255,255,0.3); border-radius: 999px; overflow: hidden;">
                        <div style="height: 100%; width: {pct_bar}%; background: #fff; border-radius: 999px; transition: width 0.3s;"></div>
                    </div>
                </span>
            </div>
            """,
            unsafe_allow_html=True,
        )

        # --- Comparación con límites simulados (Trip In / Trip Out / Rotating)
        limits_df = st.session_state.get("trip_limits_df")
        if limits_df is not None and not limits_df.empty:
            limit_col = "TripOut" if mode == "Trip Out" else "TripIn"
            if limit_col not in limits_df.columns or limits_df[limit_col].isna().all():
                limit_col = "Rotating" if "Rotating" in limits_df.columns else None
            if limit_col and limits_df[limit_col].notna().any():
                depth_lim = limits_df["Depth"].values
                val_lim = limits_df[limit_col].values
                df["Limit_sim"] = np.interp(df["Bit depth"].values, depth_lim, val_lim)
                df["Over_limit"] = df["Hookload"] > df["Limit_sim"]
                df["Exceedance"] = (df["Hookload"] - df["Limit_sim"]).clip(lower=0)
                n_over = df["Over_limit"].sum()
                pct_over = 100.0 * n_over / len(df) if len(df) > 0 else 0.0
                max_exc = float(df["Exceedance"].max()) if n_over > 0 else 0.0
                depth_max_exc = (
                    df.loc[df["Exceedance"].idxmax(), "Bit depth"] if n_over > 0 else None
                )

                st.markdown("#### Estadísticas vs límite simulado")
                st.caption(f"Límite usado: **{limit_col}** (según tipo de viaje). Cuando Hookload real > límite, se cuenta como excedencia.")
                pct_over_bar = min(100, max(0, pct_over))
                depth_str = f"{depth_max_exc:,.0f} m" if depth_max_exc is not None else "—"
                st.markdown(
                    f"""
                    <div style="display: flex; flex-wrap: wrap; gap: 0.6rem; align-items: center; margin-bottom: 0.5rem;">
                        <span style="
                            background: linear-gradient(135deg, #7c2d12 0%, #ea580c 100%);
                            color: #fff; font-size: 0.8rem; font-weight: 600;
                            padding: 0.4rem 0.75rem; border-radius: 999px; letter-spacing: 0.02em;
                            box-shadow: 0 1px 2px rgba(0,0,0,0.2);
                        "><span style="opacity: 0.9; font-weight: 500;">Puntos sobre límite</span> {int(n_over):,}</span>
                        <span style="
                            background: linear-gradient(135deg, #4c1d95 0%, #7c3aed 100%);
                            color: #fff; font-size: 0.8rem; font-weight: 600;
                            padding: 0.4rem 0.75rem; border-radius: 999px; letter-spacing: 0.02em;
                            box-shadow: 0 1px 2px rgba(0,0,0,0.2);
                            min-width: 150px;
                        "><span style="opacity: 0.9; font-weight: 500;">% viaje sobre límite</span> {pct_over:.1f}%
                            <div style="margin-top: 4px; height: 4px; background: rgba(255,255,255,0.3); border-radius: 999px; overflow: hidden;">
                                <div style="height: 100%; width: {pct_over_bar}%; background: #fff; border-radius: 999px;"></div>
                            </div>
                        </span>
                        <span style="
                            background: linear-gradient(135deg, #991b1b 0%, #ef4444 100%);
                            color: #fff; font-size: 0.8rem; font-weight: 600;
                            padding: 0.4rem 0.75rem; border-radius: 999px; letter-spacing: 0.02em;
                            box-shadow: 0 1px 2px rgba(0,0,0,0.2);
                        "><span style="opacity: 0.9; font-weight: 500;">Excedencia máx.</span> {max_exc:,.1f}</span>
                        <span style="
                            background: linear-gradient(135deg, #0f766e 0%, #14b8a6 100%);
                            color: #fff; font-size: 0.8rem; font-weight: 600;
                            padding: 0.4rem 0.75rem; border-radius: 999px; letter-spacing: 0.02em;
                            box-shadow: 0 1px 2px rgba(0,0,0,0.2);
                        "><span style="opacity: 0.9; font-weight: 500;">Prof. max exced.</span> {depth_str}</span>
                    </div>
                    """,
                    unsafe_allow_html=True,
                )

                over_df = df[df["Over_limit"]].copy()
                if not over_df.empty:
                    over_df = over_df.sort_values("Bit depth")
                    exc = over_df["Exceedance"]
                    p33 = float(exc.quantile(0.33)) if len(exc) >= 2 else 0.0
                    p66 = float(exc.quantile(0.66)) if len(exc) >= 2 else float(exc.max())
                    if p66 <= p33:
                        p66 = max(p33 + 1e-6, float(exc.max()))

                    def _estado_lim(x):
                        if pd.isna(x) or x <= 0:
                            return "🟢 Bajo"
                        if x <= p33:
                            return "🟢 Bajo"
                        if x <= p66:
                            return "🟡 Medio"
                        return "🔴 Alto"

                    over_df["Estado_lim"] = over_df["Exceedance"].apply(_estado_lim)
                    st.markdown("##### Puntos del viaje que superan el límite simulado")
                    st.dataframe(
                        over_df[["Timestamp", "Bit depth", "Hookload", "Limit_sim", "Exceedance", "Estado_lim"]].rename(columns={"Limit_sim": "Límite", "Estado_lim": "Estado"}),
                        use_container_width=True,
                        hide_index=True,
                    )
                    st.caption("Solo se muestran filas donde Hookload real > límite. Estado por excedencia (percentiles).")
                else:
                    st.success("Ningún punto del viaje supera el límite simulado.")

            else:
                st.caption("Límites cargados pero no hay columna válida para este tipo de viaje (Trip Out / Trip In / Rotating).")

        st.markdown("#### Vista Pro – Eventos de overpull en el tiempo")
        st.caption("Barras = eventos de overpull; color = severidad (amarillo aislado → naranja frecuencia en aumento → rojo restricción/riesgo de pegadura).")
        fig_events = _trip_build_events_timeline_figure(df_with_events, t0, mode_label, dark=dark_pro)
        st.plotly_chart(fig_events, use_container_width=True, config={"displayModeBar": "hover", "displaylogo": False})

        st.markdown("#### Hookload con eventos resaltados")
        st.caption("Línea = Hookload; diamantes = eventos de overpull coloreados por severidad.")
        fig_hookload_pro = _trip_build_hookload_with_events_figure(df_with_events, mode_label, dark=dark_pro)
        st.plotly_chart(fig_hookload_pro, use_container_width=True, config={"displayModeBar": "hover", "displaylogo": False})

        insight = _trip_generate_insight(df_with_events, env[env["Event"]].shape[0])
        st.info(f"💡 **Insight:** {insight}")

        st.markdown("#### Eventos de overpull vs profundidad")
        st.caption("Dónde aparecen los eventos a lo largo del pozo (eje X = profundidad). Color = severidad.")
        fig_events_depth = _trip_build_events_vs_depth_figure(df_with_events, mode_label, dark=dark_pro)
        if fig_events_depth.data:
            st.plotly_chart(fig_events_depth, use_container_width=True, config={"displayModeBar": "hover", "displaylogo": False})
        else:
            st.caption("No hay eventos de overpull para graficar vs profundidad.")

        st.markdown("---")
        st.markdown("#### Envolvente por profundidad (bin)")
        fig_env = px.line(
            env,
            x="Depth_bin",
            y=["Hookload_max", "Hookload_baseline", "Overpull"],
            title="Hookload envelope vs profundidad (bin)<br><sub style='font-size:10px;color:#64748b;'>Max, baseline y overpull por bin de profundidad (m)</sub>",
            labels={"Depth_bin": "Profundidad (m)", "value": "Hookload / Overpull (klb)", "variable": "Serie"},
        ).update_layout(xaxis_autorange="reversed")
        st.plotly_chart(prettify(fig_env), use_container_width=True, config={"displayModeBar": "hover", "displaylogo": False})

        st.markdown("#### Eventos de overpull (por profundidad)")
        env_events = env[env["Event"]].copy()
        if not env_events.empty:
            op = env_events["Overpull"]
            p33 = float(op.quantile(0.33))
            p66 = float(op.quantile(0.66))
            if p66 <= p33:
                p66 = max(p33 + 1e-6, float(op.max()))
            def _estado(val):
                if pd.isna(val) or val <= 0:
                    return "🟢 Bajo"
                if val <= p33:
                    return "🟢 Bajo"
                if val <= p66:
                    return "🟡 Medio"
                return "🔴 Alto"
            env_events["Estado"] = env_events["Overpull"].apply(_estado)
            env_events_display = env_events[["Depth_bin", "Hookload_max", "Hookload_baseline", "Overpull", "Estado"]]
            st.markdown(
                '<span style="font-size:0.85rem;">'
                '<span style="background:#dcfce7;color:#166534;padding:0.15rem 0.5rem;border-radius:999px;margin-right:0.5rem;">🟢 Bajo</span>'
                '<span style="background:#fef9c3;color:#854d0e;padding:0.15rem 0.5rem;border-radius:999px;margin-right:0.5rem;">🟡 Medio</span>'
                '<span style="background:#fee2e2;color:#991b1b;padding:0.15rem 0.5rem;border-radius:999px;">🔴 Alto</span>'
                '</span>',
                unsafe_allow_html=True,
            )
            st.dataframe(env_events_display, use_container_width=True, hide_index=True)
        else:
            st.caption("No hay eventos de overpull en el envelope por profundidad.")

        st.markdown("#### Muestra de datos alineados (tiempo)")
        st.dataframe(df.head(200), use_container_width=True, hide_index=True)

        # --- Actualización automática cada N segundos (solo API)
        if data_source == "API" and st.session_state.get("trip_auto_refresh"):
            interval = int(st.session_state.get("trip_auto_refresh_interval", 30))
            interval = max(10, min(300, interval))
            countdown_placeholder = st.empty()
            for i in range(interval, 0, -1):
                countdown_placeholder.info(f"🔄 Próxima actualización en **{i}** s… (desmarca «Actualizar automáticamente» para detener)")
                time.sleep(1)
            countdown_placeholder.empty()
            st.session_state["trip_auto_rerun_trigger"] = True
            st.rerun()

    with tab_broom:
        st.markdown("### Broomstick (Hookload vs Profundidad)")
        st.caption("Puntos medidos (PU/SO) + curvas del modelo por FF. **Para mover las curvas:** usa el **Offset de curvas (klb)** (negativo = bajar, positivo = subir). El **slider FF** cambia la forma (fricción).")

        # --- Configuración del broomstick
        colb1, colb2, colb3 = st.columns(3)
        with colb1:
            eps_dir = st.number_input("Epsilon dirección (m por muestra)", min_value=0.001, value=0.02, step=0.01, key="trip_broom_eps_dir")
        with colb2:
            max_pts = st.number_input("Máx puntos (scatter)", min_value=200, max_value=20000, value=2500, step=100, key="trip_broom_max_pts")
        with colb3:
            one_per_stand = st.checkbox("One per stand (downsample por profundidad)", value=False, key="trip_broom_one_per_stand",
                                        help="Reduce puntos agrupando por tramo de profundidad (aprox) para una vista más limpia.")

        df_rt = df.copy()
        if "Timestamp" in df_rt.columns:
            df_rt = df_rt.sort_values("Timestamp")
        df_rt["Dir"] = _trip_dir_labels_from_depth(df_rt, eps_m=float(eps_dir))
        df_rt = df_rt.dropna(subset=["Dir"])

        # Opcional: downsample tipo "one per stand"
        if one_per_stand and not df_rt.empty:
            stand_m = 30.0
            df_rt["_stand"] = (df_rt["Bit depth"] / stand_m).round().astype(int)
            df_rt = (
                df_rt.sort_values("Timestamp")
                    .groupby(["Dir", "_stand"], as_index=False)
                    .agg({"Bit depth": "mean", "Hookload": "median"})
            )

        # Limitar puntos para performance
        if not df_rt.empty and len(df_rt) > int(max_pts):
            df_rt = df_rt.sample(int(max_pts), random_state=7).sort_values("Bit depth")

        # --- Chips y etiquetas sobre la gráfica
        model_df = st.session_state.get("trip_ff_family_df")
        fam_map = st.session_state.get("trip_ff_family_map") or {}
        limits_df = st.session_state.get("trip_limits_df")
        has_limits = limits_df is not None and not limits_df.empty
        n_pu = len(df_rt[df_rt["Dir"] == "PU"]) if not df_rt.empty else 0
        n_so = len(df_rt[df_rt["Dir"] == "SO"]) if not df_rt.empty else 0
        ff_val = float(st.session_state.get("trip_broom_ff", 0.25)) if (model_df is not None and not model_df.empty) else None

        chip_items = []
        if n_pu > 0:
            chip_items.append(("Medido PU (Trip Out)", "blue", f"{n_pu:,} pts"))
        if n_so > 0:
            chip_items.append(("Medido SO (Trip In)", "orange", f"{n_so:,} pts"))
        if model_df is not None and not model_df.empty and ff_val is not None:
            chip_items.append((f"Modelo FF = {ff_val:.2f}", "green", "PU/SO/ROT"))
        if has_limits:
            chip_items.append(("Límites cargados", "red", "Trip Out/In/Rot"))

        if chip_items:
            color_map = {"blue": ("#1e3a5f", "#2563eb"), "orange": ("#422006", "#f59e0b"), "green": ("#14532d", "#22c55e"), "red": ("#450a0a", "#ef4444")}
            cols_chip = st.columns(min(len(chip_items), 5))
            for i, (label, color_key, sub) in enumerate(chip_items):
                if i < len(cols_chip):
                    c1, c2 = color_map.get(color_key, ("#334155", "#64748b"))
                    with cols_chip[i]:
                        style_outer = (
                            "display:inline-flex;align-items:center;gap:0.35rem;"
                            f"background:linear-gradient(135deg,{c1},{c2});"
                            "color:#fff;font-size:0.75rem;font-weight:600;padding:0.25rem 0.6rem;"
                            "border-radius:999px;letter-spacing:0.02em;box-shadow:0 1px 2px rgba(0,0,0,0.2);"
                        )
                        st.markdown(
                            f'<span style="{style_outer}">{label} <span style="opacity:0.9;font-weight:500">({sub})</span></span>',
                            unsafe_allow_html=True,
                        )

        # --- Alerta: puntos sobre límite (si hay límites y datos medidos)
        alert_lines = []
        if has_limits and not df_rt.empty and "Depth" in limits_df.columns:
            lim_prep = _trip_prepare_limits_for_continuous_line(limits_df, rolling_window=25)
            if not lim_prep.empty and "TripOut" in lim_prep.columns and "TripIn" in lim_prep.columns:
                depth_lim = lim_prep["Depth"].values
                tripout_lim = lim_prep["TripOut"].values
                tripin_lim = lim_prep["TripIn"].values
                pu_pts = df_rt[df_rt["Dir"] == "PU"]
                so_pts = df_rt[df_rt["Dir"] == "SO"]
                over_tripout = 0
                under_tripin = 0
                for _, row in pu_pts.iterrows():
                    d, hl = row["Bit depth"], row["Hookload"]
                    lim_val = np.interp(d, depth_lim, tripout_lim)
                    if hl > lim_val:
                        over_tripout += 1
                for _, row in so_pts.iterrows():
                    d, hl = row["Bit depth"], row["Hookload"]
                    lim_val = np.interp(d, depth_lim, tripin_lim)
                    if hl < lim_val and lim_val > 0:
                        under_tripin += 1
                if over_tripout > 0:
                    alert_lines.append(f"⚠️ **{over_tripout:,}** puntos PU (Trip Out) **sobre** el límite simulado.")
                if under_tripin > 0:
                    alert_lines.append(f"⚠️ **{under_tripin:,}** puntos SO (Trip In) **bajo** el límite simulado.")
        if alert_lines:
            st.warning(" | ".join(alert_lines))

        fig = go.Figure()

        # Scatter medido
        if not df_rt.empty:
            pu_pts = df_rt[df_rt["Dir"] == "PU"]
            so_pts = df_rt[df_rt["Dir"] == "SO"]
            if not pu_pts.empty:
                fig.add_trace(go.Scatter(
                    x=pu_pts["Hookload"], y=pu_pts["Bit depth"],
                    mode="markers",
                    name="Medido PU (Trip Out)",
                    marker=dict(size=5, opacity=0.6),
                ))
            if not so_pts.empty:
                fig.add_trace(go.Scatter(
                    x=so_pts["Hookload"], y=so_pts["Bit depth"],
                    mode="markers",
                    name="Medido SO (Trip In)",
                    marker=dict(size=5, opacity=0.6),
                ))

        # Curvas del modelo por FF (familia)
        model_df = st.session_state.get("trip_ff_family_df")
        fam_map = st.session_state.get("trip_ff_family_map") or {}
        if model_df is None or model_df.empty or not fam_map:
            st.info("Carga un CSV de **Modelo Broomstick por FF** para habilitar el slider que mueve las curvas.")
        else:
            ff_min = float(st.session_state.get("trip_ff_min", 0.1))
            ff_max = float(st.session_state.get("trip_ff_max", 0.5))
            ff_val = st.slider("Factor de fricción (FF)", min_value=ff_min, max_value=ff_max, value=float((ff_min+ff_max)/2), step=0.01, key="trip_broom_ff")
            st.caption("**Mover curvas del modelo:**")
            curve_offset_klb = st.number_input(
                "Offset de curvas (klb) — sube/baja las líneas PU, SO y ROT",
                value=0.0,
                min_value=-300.0,
                max_value=300.0,
                step=5.0,
                key="trip_broom_curve_offset",
                help="Negativo = bajar curvas (acercar a los puntos). Positivo = subir. Se aplica al instante.",
            )

            # Curvas interpoladas para el FF seleccionado
            pu_curve = interp_ff_curve(model_df, fam_map, "PU", float(ff_val))
            so_curve = interp_ff_curve(model_df, fam_map, "SO", float(ff_val))
            rot_curve = interp_ff_curve(model_df, fam_map, "ROT", float(ff_val))
            if pu_curve is not None and curve_offset_klb != 0:
                pu_curve = pu_curve + float(curve_offset_klb)
            if so_curve is not None and curve_offset_klb != 0:
                so_curve = so_curve + float(curve_offset_klb)
            if rot_curve is not None and curve_offset_klb != 0:
                rot_curve = rot_curve + float(curve_offset_klb)

            if pu_curve is not None and not pu_curve.empty:
                fig.add_trace(go.Scatter(
                    x=pu_curve, y=model_df["Depth"],
                    mode="lines",
                    name=f"Modelo PU @ FF={ff_val:.2f}",
                    line=dict(width=2, shape="spline", smoothing=1.3),
                ))
            if so_curve is not None and not so_curve.empty:
                fig.add_trace(go.Scatter(
                    x=so_curve, y=model_df["Depth"],
                    mode="lines",
                    name=f"Modelo SO @ FF={ff_val:.2f}",
                    line=dict(width=2, shape="spline", smoothing=1.3),
                ))
            # ROT opcional
            if rot_curve is not None and not rot_curve.empty and rot_curve.notna().any():
                fig.add_trace(go.Scatter(
                    x=rot_curve, y=model_df["Depth"],
                    mode="lines",
                    name=f"Modelo ROT @ FF={ff_val:.2f}",
                    line=dict(width=2, shape="spline", smoothing=1.3),
                ))

        # Límites simples (una curva continua por límite; x=Hookload, y=Depth)
        limits_df = st.session_state.get("trip_limits_df")
        if limits_df is not None and not limits_df.empty:
            lim_broom = _trip_prepare_limits_for_continuous_line(limits_df, rolling_window=25)
            if not lim_broom.empty:
                if "TripOut" in lim_broom.columns and lim_broom["TripOut"].notna().any():
                    fig.add_trace(go.Scatter(
                        x=lim_broom["TripOut"].values,
                        y=lim_broom["Depth"].values,
                        mode="lines",
                        name="Límite Trip Out",
                        line=dict(width=2, shape="spline", smoothing=1.3),
                        connectgaps=True,
                    ))
                if "TripIn" in lim_broom.columns and lim_broom["TripIn"].notna().any():
                    fig.add_trace(go.Scatter(
                        x=lim_broom["TripIn"].values,
                        y=lim_broom["Depth"].values,
                        mode="lines",
                        name="Límite Trip In",
                        line=dict(width=2, shape="spline", smoothing=1.3),
                        connectgaps=True,
                    ))
                if "Rotating" in lim_broom.columns and lim_broom["Rotating"].notna().any():
                    fig.add_trace(go.Scatter(
                        x=lim_broom["Rotating"].values,
                        y=lim_broom["Depth"].values,
                        mode="lines",
                        name="Límite Rotating",
                        line=dict(width=2, shape="spline", smoothing=1.3),
                        connectgaps=True,
                    ))

        fig.update_layout(
            xaxis_title="Hookload",
            yaxis_title="Profundidad (m)",
            template="plotly_dark",
            height=650,
            legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
        )
        fig.update_yaxes(autorange="reversed")
        st.plotly_chart(fig, use_container_width=True, config=PLOTLY_CONFIG)
        st.caption(
            "**Leyenda:** PU = Trip Out (sacar tubería), SO = Trip In (meter tubería), ROT = rotando. "
            "Puntos = mediciones reales; líneas = modelo por factor de fricción (FF); límites = curvas de diseño/simulador."
        )

        with st.expander("Datos usados (puntos medidos)", expanded=False):
            st.dataframe(df_rt.head(200), use_container_width=True, hide_index=True)


def find_planned_column(columns: Iterable[str], candidates: List[str]) -> str | None:
    lowered = {c.lower(): c for c in columns}
    for cand in candidates:
        key = cand.lower()
        if key in lowered:
            return lowered[key]
    return None


def find_torque_column(columns: Iterable[str], preferred: str | None) -> str | None:
    if preferred:
        col = find_planned_column(columns, [preferred])
        if col:
            return col
    return find_planned_column(columns, TORQUE_COL_CANDIDATES)


def add_real_vs_planned_section(
    prs: Presentation,
    df_run: pd.DataFrame,
    run_name: str,
    show_plots: bool,
) -> None:
    df_run_sorted = df_run.sort_values("Depth_X")
    planned_map = {
        metric: find_planned_column(df_run_sorted.columns, candidates)
        for metric, candidates in PLANNED_COL_CANDIDATES.items()
    }

    available = {m: c for m, c in planned_map.items() if c}
    if not available:
        st.info("No se encontraron columnas programadas para comparar.")
        return

    st.subheader(f"{run_name} – Real vs Programado")
    vspace(6)

    for metric, planned_col in available.items():
        real_col = metric
        if real_col not in df_run_sorted.columns:
            continue

        df_plot = df_run_sorted[["Depth_X", real_col, planned_col]].copy()
        df_plot = df_plot.dropna(subset=["Depth_X"])

        fig = px.line(
            df_plot,
            x="Depth_X",
            y=[real_col, planned_col],
            title=f"{run_name} – {metric} Real vs Programado",
            labels={"Depth_X": "Depth (m)", "value": metric, "variable": "Serie"},
        ).update_layout(xaxis_autorange="reversed")

        if len(fig.data) >= 2:
            fig.data[1].line.dash = "dash"

        if show_plots:
            st.plotly_chart(
                prettify(fig),
                use_container_width=True,
                config=PLOTLY_CONFIG,
            )
            mae = (df_plot[real_col] - df_plot[planned_col]).abs().mean()
            chart_notes(
                f"MAE {metric}: {format_num(mae)}.",
                "Línea sólida = real, línea discontinua = programado.",
            )

        save_and_show_plotly(prs, f"{run_name} – {metric} Real vs Programado", fig, False)


def capture_region_screenshots(
    region: Tuple[int, int, int, int],
    n_shots: int,
    interval_s: int,
    title_prefix: str,
    prs: Presentation,
    show_plots: bool,
    download_dir: str,
) -> None:
    try:
        import pyautogui  # type: ignore
    except Exception as e:
        st.error(f"❌ No se pudo importar pyautogui: {e}")
        return

    x, y, w, h = region
    st.info(
        f"🟡 Iniciando capturas: {n_shots} tomas cada {interval_s}s.\n"
        f"Región (x,y,w,h): {region}\n"
        "Asegúrate de que el área esté visible en pantalla."
    )

    prog = st.progress(0)

    for i in range(n_shots):
        if winsound:
            try:
                winsound.Beep(1000, 300)
            except Exception:
                pass

        st.write(f"📸 Captura {i+1}/{n_shots} a las {datetime.now().strftime('%H:%M:%S')}")

        img = pyautogui.screenshot(region=(x, y, w, h))
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        file_png = os.path.join(download_dir, f"region_capture_{ts}_{i+1}.png")
        img.save(file_png)

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)

        slide_title = f"{title_prefix} #{i+1}"
        add_image_slide(prs, slide_title, buf)

        if show_plots:
            st.image(img, caption=slide_title, use_container_width=True)

        prog.progress(int(((i + 1) / n_shots) * 100))
        if i < n_shots - 1:
            time.sleep(interval_s)

    st.success("✅ Capturas completadas.")


# =========================
# Carga y validación de datos
# =========================
@st.cache_data(show_spinner=False)
def _uploaded_file_bytes(file) -> bytes:
    if isinstance(file, (bytes, bytearray)):
        return bytes(file)
    if hasattr(file, "getvalue"):
        return bytes(file.getvalue())
    if hasattr(file, "read"):
        pos = None
        try:
            pos = file.tell()
        except Exception:
            pos = None
        data = file.read()
        try:
            if pos is not None:
                file.seek(pos)
        except Exception:
            pass
        return data if isinstance(data, bytes) else bytes(data)
    raise ValueError("No pude leer el archivo Excel subido.")


def _find_geopark_header_row(xl: pd.ExcelFile, sheet_name: str) -> int | None:
    probe = pd.read_excel(xl, sheet_name=sheet_name, header=None, nrows=80)
    for i in range(len(probe)):
        row = [str(v).strip().lower() for v in probe.iloc[i].tolist() if pd.notna(v)]
        if not row:
            continue
        joined = " | ".join(row)
        if (
            "actividad" in joined
            and "metros perforados" in joined
            and ("rpm superficie" in joined or "rpm totales" in joined)
        ):
            return int(i)
    return None


def _detect_excel_slide_sheet_format(file) -> tuple[str, str | None, int | None]:
    content = _uploaded_file_bytes(file)
    xl = pd.ExcelFile(io.BytesIO(content))
    sheet_names = list(xl.sheet_names)

    if SHEET_NAME in sheet_names:
        try:
            df_std = pd.read_excel(xl, sheet_name=SHEET_NAME, nrows=5)
            cols_std = {str(c).strip().lower() for c in df_std.columns}
            if {"mode", "start depth", "end depth", "rop", "wob", "rpm"}.issubset(cols_std):
                return "standard", SHEET_NAME, 0
        except Exception:
            pass

    for sh in sheet_names:
        hdr_row = _find_geopark_header_row(xl, sh)
        if hdr_row is not None:
            return "geopark", sh, hdr_row

    for sh in sheet_names:
        try:
            df_std = pd.read_excel(xl, sheet_name=sh, nrows=5)
            cols_std = {str(c).strip().lower() for c in df_std.columns}
            if {"mode", "start depth", "end depth", "rop", "wob", "rpm"}.issubset(cols_std):
                return "standard", sh, 0
        except Exception:
            continue

    return "unknown", sheet_names[0] if sheet_names else None, None


def _excel_serial_to_datetime(series: pd.Series) -> pd.Series:
    s_num = pd.to_numeric(series, errors="coerce")
    out = pd.to_datetime(s_num, unit="D", origin="1899-12-30", errors="coerce")
    if out.notna().any():
        return out
    return pd.to_datetime(series, errors="coerce")


def _combine_date_and_time_columns(date_series: pd.Series, time_series: pd.Series) -> pd.Series:
    date_part = pd.to_datetime(date_series, errors="coerce").dt.normalize()
    time_part = pd.to_datetime(time_series.astype(str), errors="coerce")
    out = pd.Series(pd.NaT, index=date_series.index, dtype="datetime64[ns]")
    mask = date_part.notna() & time_part.notna()
    if mask.any():
        out.loc[mask] = date_part.loc[mask] + (
            time_part.loc[mask] - time_part.loc[mask].dt.normalize()
        )
    return out


def _classify_geopark_mode(activity, rpm_surface, rpm_total, meters) -> str:
    act = str(activity or "").strip().lower()
    meters_v = pd.to_numeric(pd.Series([meters]), errors="coerce").iloc[0]
    rpm_s = pd.to_numeric(pd.Series([rpm_surface]), errors="coerce").iloc[0]
    rpm_t = pd.to_numeric(pd.Series([rpm_total]), errors="coerce").iloc[0]

    if "desliza" in act or "orienta tf" in act:
        return "SLIDE"
    if any(k in act for k in ["rota", "perfora", "repasa"]):
        return "ROTARY"
    if pd.notna(meters_v) and meters_v > 0:
        if pd.notna(rpm_s) and rpm_s > 10:
            return "ROTARY"
        if pd.notna(rpm_t) and rpm_t > 30:
            return "SLIDE"
        return "ROTARY"
    if "circula" in act:
        return "Static - Pump"
    if "survey" in act:
        return "Static"
    if any(k in act for k in ["arma bha", "desarma bha", "agrega", "otros"]):
        return "Surface Operations"
    if any(k in act for k in ["rih", "run in"]):
        return "Run In - Tripping"
    if any(k in act for k in ["pooh", "pull out"]):
        return "Pull Up - Trip Out"
    return str(activity or "Unknown")


def _load_geopark_slide_sheet(file, sheet_name: str, header_row: int) -> pd.DataFrame:
    content = _uploaded_file_bytes(file)
    xl = pd.ExcelFile(io.BytesIO(content))
    df = pd.read_excel(xl, sheet_name=sheet_name, header=header_row)
    df = df.dropna(how="all").copy()

    unnamed = [c for c in df.columns if str(c).lower().startswith("unnamed")]
    if unnamed:
        df = df.drop(columns=unnamed, errors="ignore")

    rename_map = {
        "Hasta ": "Hasta_time",
        "Desde.1": "Start_dt_raw",
        "Hasta": "End_dt_raw",
        "Desde.2": "Start Depth",
        "Hasta.1": "End Depth",
        "Metros Perforados": "Distance",
        "RPM Superficie": "RPM Surface",
        "RPM Totales": "RPM Total",
        "Tq. En Fondo": "Torque",
        "Tq. Fuera de Fondo": "Torque Off Bottom",
        "Presión Diferencial": "Differential Pressure",
        "Densidad del Lodo": "Mud Density",
        "BHA # (*)": "BHA",
        "INCL": "Inclination",
        "AZ": "Azimuth",
        "MD": "Survey MD",
        "Fecha": "Date",
        "Actividad": "Activity",
        "WOB": "WOB",
        "ROP": "ROP",
    }
    df = df.rename(columns=rename_map)

    if "Activity" not in df.columns:
        raise ValueError(
            "No encontré la columna 'Actividad' en el slide sheet Geopark."
        )

    unit_mask = df["Activity"].astype(str).str.contains("actividad", case=False, na=False)
    df = df.loc[~unit_mask].copy()

    if "Date" in df.columns:
        df["Date"] = pd.to_datetime(df["Date"], errors="coerce")

    start_dt = _excel_serial_to_datetime(df.get("Start_dt_raw", pd.Series(index=df.index, dtype="object")))
    end_dt = _excel_serial_to_datetime(df.get("End_dt_raw", pd.Series(index=df.index, dtype="object")))

    if "Desde" in df.columns and "Date" in df.columns:
        start_fallback = _combine_date_and_time_columns(df["Date"], df["Desde"])
        start_dt = start_dt.where(start_dt.notna(), start_fallback)
    if "Hasta_time" in df.columns and "Date" in df.columns:
        end_fallback = _combine_date_and_time_columns(df["Date"], df["Hasta_time"])
        end_dt = end_dt.where(end_dt.notna(), end_fallback)

    df["Start"] = start_dt
    df["End"] = end_dt

    numeric_cols = [
        "Start Depth",
        "End Depth",
        "Distance",
        "ROP",
        "RPM Surface",
        "RPM Total",
        "WOB",
        "Torque",
        "Torque Off Bottom",
        "Differential Pressure",
        "Mud Density",
        "Survey MD",
        "Inclination",
        "Azimuth",
        "BHA",
    ]
    for col in numeric_cols:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce")

    if "End Depth" in df.columns and "Start Depth" in df.columns and "Distance" in df.columns:
        df["Distance"] = df["Distance"].where(df["Distance"].notna(), df["End Depth"] - df["Start Depth"])
        df["Start Depth"] = df["Start Depth"].where(
            df["Start Depth"].notna(),
            df["End Depth"] - df["Distance"],
        )
        df["End Depth"] = df["End Depth"].where(
            df["End Depth"].notna(),
            df["Start Depth"] + df["Distance"],
        )

    if "Survey MD" not in df.columns:
        df["Survey MD"] = np.nan
    df["Survey MD"] = df["Survey MD"].where(df["Survey MD"].notna(), df.get("End Depth"))

    if "RPM Total" in df.columns or "RPM Surface" in df.columns:
        rpm_total = df["RPM Total"] if "RPM Total" in df.columns else pd.Series(np.nan, index=df.index)
        rpm_surface = df["RPM Surface"] if "RPM Surface" in df.columns else pd.Series(np.nan, index=df.index)
        df["RPM"] = rpm_total.where(rpm_total.notna(), rpm_surface)
    else:
        df["RPM"] = np.nan

    df["Mode"] = [
        _classify_geopark_mode(a, rs, rt, m)
        for a, rs, rt, m in zip(
            df.get("Activity", pd.Series(index=df.index)),
            df.get("RPM Surface", pd.Series(index=df.index)),
            df.get("RPM Total", pd.Series(index=df.index)),
            df.get("Distance", pd.Series(index=df.index)),
        )
    ]
    df["DLS"] = np.nan
    df["Source Format"] = "geopark"

    ordered_cols = list(dict.fromkeys(REQUIRED_COLUMNS + [
        "Activity",
        "RPM Surface",
        "RPM Total",
        "Torque",
        "Torque Off Bottom",
        "Differential Pressure",
        "Mud Density",
        "BHA",
        "Source Format",
        "Comments",
        "Comentarios",
        "Date",
    ]))
    for col in ordered_cols:
        if col not in df.columns:
            df[col] = np.nan

    df = df.dropna(subset=["Start", "End", "Start Depth", "End Depth"], how="all")
    return df[ordered_cols].reset_index(drop=True)


@st.cache_data(show_spinner=False)
def load_excel_data(file) -> pd.DataFrame:
    fmt, sheet_name, header_row = _detect_excel_slide_sheet_format(file)
    if fmt == "standard":
        content = _uploaded_file_bytes(file)
        df = pd.read_excel(io.BytesIO(content), sheet_name=sheet_name or SHEET_NAME)
        df = df.dropna(how="all").reset_index(drop=True)
        df["Source Format"] = "standard"
        return df
    if fmt == "geopark" and sheet_name is not None and header_row is not None:
        return _load_geopark_slide_sheet(file, sheet_name, header_row)
    raise ValueError(
        "No pude reconocer el formato del Slide Sheet. Soporto el formato clásico 'worksheet' y el formato Geopark."
    )


def suggest_runs_from_loaded_df(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame()
    d = df.copy()
    if "Start" in d.columns:
        d["Start"] = pd.to_datetime(d["Start"], errors="coerce")
    if "End" in d.columns:
        d["End"] = pd.to_datetime(d["End"], errors="coerce")
    for col in ["Start Depth", "End Depth", "Distance", "BHA"]:
        if col in d.columns:
            d[col] = pd.to_numeric(d[col], errors="coerce")

    if "BHA" in d.columns and d["BHA"].notna().any():
        rows = []
        for bha, g in d.groupby("BHA", dropna=True):
            g = g.copy()
            g["depth_lo"] = g[["Start Depth", "End Depth"]].min(axis=1)
            g["depth_hi"] = g[["Start Depth", "End Depth"]].max(axis=1)
            start_depth = g["depth_lo"].min()
            end_depth = g["depth_hi"].max()
            start_time = g["Start"].min()
            end_time = g["End"].max()
            if pd.notna(start_depth) and pd.notna(end_depth) and end_depth > start_depth and pd.notna(start_time) and pd.notna(end_time):
                bha_num = int(bha) if float(bha).is_integer() else bha
                rows.append(
                    {
                        "Run": f"BHA {bha_num}",
                        "start_depth": float(start_depth),
                        "end_depth": float(end_depth),
                        "start_time": pd.Timestamp(start_time).isoformat(),
                        "end_time": pd.Timestamp(end_time).isoformat(),
                    }
                )
        if rows:
            return pd.DataFrame(rows).sort_values(["start_depth", "start_time"]).reset_index(drop=True)

    if {"Start Depth", "End Depth", "Start", "End"}.issubset(d.columns):
        d["depth_lo"] = d[["Start Depth", "End Depth"]].min(axis=1)
        d["depth_hi"] = d[["Start Depth", "End Depth"]].max(axis=1)
        start_depth = d["depth_lo"].min()
        end_depth = d["depth_hi"].max()
        start_time = d["Start"].min()
        end_time = d["End"].max()
        if pd.notna(start_depth) and pd.notna(end_depth) and end_depth > start_depth and pd.notna(start_time) and pd.notna(end_time):
            return pd.DataFrame(
                [
                    {
                        "Run": "Run 1",
                        "start_depth": float(start_depth),
                        "end_depth": float(end_depth),
                        "start_time": pd.Timestamp(start_time).isoformat(),
                        "end_time": pd.Timestamp(end_time).isoformat(),
                    }
                ]
            )
    return pd.DataFrame()



def normalize_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    missing = [c for c in REQUIRED_COLUMNS if c not in df.columns]
    if missing:
        st.error(
            "Faltan columnas obligatorias en el Excel:\n"
            + ", ".join(missing)
        )
        st.stop()

    df["Mode"] = df["Mode"].astype(str).str.strip().str.upper()
    df["Mode"] = df["Mode"].replace(MODE_NORMALIZATION)

    num_cols = [
        "Start Depth",
        "End Depth",
        "Survey MD",
        "Inclination",
        "Azimuth",
        "Distance",
        "ROP",
        "DLS",
        "WOB",
        "RPM",
    ]
    for col in num_cols:
        df[col] = pd.to_numeric(df[col], errors="coerce")

    df["Start"] = pd.to_datetime(df["Start"], errors="coerce")
    df["End"] = pd.to_datetime(df["End"], errors="coerce")
    return df


def parse_runs(runs_df: pd.DataFrame) -> List[RunInfo]:
    runs_list = runs_df.to_dict(orient="records")
    bad = []
    parsed: List[RunInfo] = []

    for r in runs_list:
        try:
            start_depth = float(r["start_depth"])
            end_depth = float(r["end_depth"])
            if start_depth >= end_depth:
                bad.append(r["Run"])
                continue

            run = RunInfo(
                name=str(r["Run"]),
                start_depth=start_depth,
                end_depth=end_depth,
                start_time=pd.to_datetime(r["start_time"]),
                end_time=pd.to_datetime(r["end_time"]),
            )
            parsed.append(run)
        except Exception:
            bad.append(r.get("Run", "Unnamed"))

    if bad:
        st.error(f"Estas corridas tienen profundidades inválidas: {bad}")
        st.stop()

    return parsed


# =========================
# Cálculos por corrida
# =========================
def compute_run_stats(df_run: pd.DataFrame, run: RunInfo):
    df_run["Mode_norm"] = df_run["Mode"].replace(MODE_NORMALIZATION)
    df_run["Depth_X"] = df_run["Survey MD"].fillna(df_run["End Depth"])
    df_run = df_run.dropna(subset=["Depth_X", "ROP", "WOB", "RPM"])
    if df_run.empty:
        return None

    duration = run.end_time - run.start_time
    total_hours = duration.total_seconds() / 3600.0
    effective_hours = compute_effective_time_hours(df_run)
    npt_hours = max(total_hours - effective_hours, 0.0)
    start_inc, end_inc = df_run["Inclination"].min(), df_run["Inclination"].max()
    start_az, end_az = df_run["Azimuth"].min(), df_run["Azimuth"].max()
    max_dls, avg_dls = df_run["DLS"].max(), df_run["DLS"].mean()

    wob_min, wob_max, wob_avg = (
        df_run["WOB"].min(),
        df_run["WOB"].max(),
        df_run["WOB"].mean(),
    )
    rpm_min, rpm_max, rpm_avg = (
        df_run["RPM"].min(),
        df_run["RPM"].max(),
        df_run["RPM"].mean(),
    )

    total_rotary = df_run.loc[df_run["Mode_norm"] == "ROTARY", "Distance"].sum()
    total_slide = df_run.loc[df_run["Mode_norm"] == "SLIDE", "Distance"].sum()
    total_distance = df_run["Distance"].sum()

    rotary_pct = (total_rotary / total_distance) * 100 if total_distance > 0 else 0
    slide_pct = (total_slide / total_distance) * 100 if total_distance > 0 else 0

    slide_data = df_run[df_run["Mode_norm"] == "SLIDE"]
    slide_start, slide_end = slide_data["Start Depth"].min(), slide_data["End Depth"].max()
    slide_intervals = slide_data.shape[0]
    avg_dls_slide = slide_data["DLS"].mean()
    max_slide_len = slide_data["Distance"].max()
    max_rotary_len = df_run.loc[df_run["Mode_norm"] == "ROTARY", "Distance"].max()
    avg_rop_rotary = df_run.loc[df_run["Mode_norm"] == "ROTARY", "ROP"].mean()
    avg_rop_slide = df_run.loc[df_run["Mode_norm"] == "SLIDE", "ROP"].mean()

    general_table = [
        ["Start Time", run.start_time],
        ["End Time", run.end_time],
        ["Total Duration", duration],
        ["Depth Start", run.start_depth],
        ["Depth End", run.end_depth],
        ["Inclination Start", start_inc],
        ["Inclination End", end_inc],
        ["Azimuth Start", start_az],
        ["Azimuth End", end_az],
        ["Max DLS", max_dls],
        ["Avg DLS", avg_dls],
        ["WOB Min", f"{wob_min:.2f} kgf"],
        ["WOB Max", f"{wob_max:.2f} kgf"],
        ["WOB Avg", f"{wob_avg:.2f} kgf"],
        ["RPM Min", f"{rpm_min:.2f}"],
        ["RPM Max", f"{rpm_max:.2f}"],
        ["RPM Avg", f"{rpm_avg:.2f}"],
    ]
    drilling_table = [
        ["Rotary", total_rotary, rotary_pct],
        ["Slide (incl. Oscillation)", total_slide, slide_pct],
        ["Total", total_distance, 100],
    ]
    slide_table = [
        ["Slide Depth Range", f"{slide_start:.2f} – {slide_end:.2f} m"],
        ["Slide Intervals", slide_intervals],
        ["Avg DLS Slide", f"{avg_dls_slide:.2f}°/30m"],
        ["Max Slide Interval", max_slide_len],
        ["Max Rotary Interval", max_rotary_len],
        ["Avg ROP Rotary", f"{avg_rop_rotary:.2f} m/h"],
        ["Avg ROP Slide", f"{avg_rop_slide:.2f} m/h"],
    ]

    return {
        "df_run": df_run,
        "general_table": general_table,
        "drilling_table": drilling_table,
        "slide_table": slide_table,
        "rotary_pct": rotary_pct,
        "slide_pct": slide_pct,
        "total_hours": total_hours,
        "effective_hours": effective_hours,
        "npt_hours": npt_hours,
        "wob_avg": wob_avg,
        "wob_min": wob_min,
        "wob_max": wob_max,
        "rpm_avg": rpm_avg,
        "rpm_min": rpm_min,
        "rpm_max": rpm_max,
        "avg_rop_rotary": avg_rop_rotary,
        "avg_rop_slide": avg_rop_slide,
        "max_rotary_len": max_rotary_len,
        "start_inc": float(start_inc) if pd.notna(start_inc) else None,
        "end_inc": float(end_inc) if pd.notna(end_inc) else None,
        "start_az": float(start_az) if pd.notna(start_az) else None,
        "end_az": float(end_az) if pd.notna(end_az) else None,
        "max_dls": float(max_dls) if pd.notna(max_dls) else None,
        "avg_dls": float(avg_dls) if pd.notna(avg_dls) else None,
    }


def _format_duration_whatsapp_en(td) -> str:
    """Formatea timedelta como '1 day 7h 02m' o '7h 02m'."""
    if td is None:
        return "n/a"
    total_sec = int(abs(td.total_seconds()))
    days, rem = divmod(total_sec, 86400)
    hours, rem2 = divmod(rem, 3600)
    minutes, _ = divmod(rem2, 60)
    parts: list[str] = []
    if days:
        parts.append(f"{days} day{'s' if days != 1 else ''}")
    parts.append(f"{hours}h {minutes:02d}m")
    return " ".join(parts) if parts else "0h 00m"


def _format_duration_whatsapp_es(td) -> str:
    """Formatea timedelta como '1 día 7h 02m' o '7h 02m'."""
    if td is None:
        return "n/d"
    total_sec = int(abs(td.total_seconds()))
    days, rem = divmod(total_sec, 86400)
    hours, rem2 = divmod(rem, 3600)
    minutes, _ = divmod(rem2, 60)
    parts: list[str] = []
    if days:
        parts.append(f"{days} día{'s' if days != 1 else ''}")
    parts.append(f"{hours}h {minutes:02d}m")
    return " ".join(parts) if parts else "0h 00m"


_MONTHS_ES = (
    "enero",
    "febrero",
    "marzo",
    "abril",
    "mayo",
    "junio",
    "julio",
    "agosto",
    "septiembre",
    "octubre",
    "noviembre",
    "diciembre",
)


def _ts_fmt_es(ts) -> str:
    """Fecha legible en español: '14 de julio de 2025, 04:15'."""
    if ts is None or pd.isna(ts):
        return "n/d"
    t = pd.Timestamp(ts)
    return f"{t.day} de {_MONTHS_ES[t.month - 1]} de {t.year}, {t.hour:02d}:{t.minute:02d}"


def _ts_fmt_en(ts) -> str:
    if ts is None or pd.isna(ts):
        return "n/a"
    t = pd.Timestamp(ts)
    return t.strftime("%B %d, %Y, at %H:%M")


def _wob_to_tonf(wob_kgf: float | None) -> float | None:
    if wob_kgf is None or pd.isna(wob_kgf):
        return None
    return float(wob_kgf) / 1000.0


def build_kpi_whatsapp_summary_en(run: RunInfo, stats: dict, df_run: pd.DataFrame) -> str:
    """
    Resumen narrativo en inglés (estilo WhatsApp), alineado al ejemplo del usuario.
    WOB: asume unidades tipo kgf en datos → se muestra en tonf (÷1000).
    """
    st_t = _ts_fmt_en(run.start_time)
    en_t = _ts_fmt_en(run.end_time)
    dur = run.end_time - run.start_time if run.end_time is not None and run.start_time is not None else None
    dur_txt = _format_duration_whatsapp_en(dur)

    d0, d1 = float(run.start_depth), float(run.end_depth)
    si, ei = stats.get("start_inc"), stats.get("end_inc")
    sa, ea = stats.get("start_az"), stats.get("end_az")

    if si is not None and ei is not None and abs(si - ei) < 1e-6:
        inc_phrase = f"inclination constant at {si:.1f}°"
    elif si is not None and ei is not None:
        inc_phrase = f"inclination from {si:.1f}° to {ei:.1f}°"
    else:
        inc_phrase = "inclination n/a"

    if sa is not None and ea is not None and abs(sa - ea) < 1e-6:
        az_phrase = f"azimuth constant at {sa:.1f}°"
    elif sa is not None and ea is not None:
        az_phrase = f"azimuth from {sa:.1f}° to {ea:.1f}°"
    else:
        az_phrase = "azimuth n/a"

    max_dls = stats.get("max_dls")
    avg_dls = stats.get("avg_dls")
    dls_max_s = f"{max_dls:.1f}°/30m" if max_dls is not None and pd.notna(max_dls) else "n/a"
    dls_avg_s = f"{avg_dls:.1f}°/30m" if avg_dls is not None and pd.notna(avg_dls) else "n/a"

    total_rotary = stats["drilling_table"][0][1]
    total_slide = stats["drilling_table"][1][1]
    total_dist = stats["drilling_table"][2][1]
    rp = stats["rotary_pct"]
    sp = stats["slide_pct"]

    max_rot = stats.get("max_rotary_len")
    if max_rot is None or (isinstance(max_rot, float) and pd.isna(max_rot)):
        max_rot_s = "n/a"
    else:
        max_rot_s = f"{float(max_rot):.2f} m"

    avg_rop_r = stats.get("avg_rop_rotary")
    avg_rop_r_s = f"{float(avg_rop_r):.2f} m/h" if avg_rop_r is not None and pd.notna(avg_rop_r) else "n/a"

    wmn = _wob_to_tonf(stats.get("wob_min"))
    wmx = _wob_to_tonf(stats.get("wob_max"))
    wav = _wob_to_tonf(stats.get("wob_avg"))
    wob_line = (
        f"WOB ranged from {wmn:.2f} to {wmx:.2f} tonf, with an average of {wav:.2f} tonf."
        if wmn is not None and wmx is not None and wav is not None
        else "WOB: n/a."
    )

    rpm_min, rpm_max, rpm_avg = stats.get("rpm_min"), stats.get("rpm_max"), stats.get("rpm_avg")
    rpm_line = (
        f"RPM ranged from {rpm_min:.2f} to {rpm_max:.2f}, with an average of {rpm_avg:.2f}."
        if rpm_min is not None and rpm_max is not None and rpm_avg is not None
        else "RPM: n/a."
    )

    p1 = (
        f"Drilling started on {st_t} and ended on {en_t}, lasting {dur_txt}."
    )
    p2 = (
        f"Depth ranged from {d0:.2f} m to {d1:.2f} m, with {inc_phrase} and {az_phrase}."
    )
    p3 = f"Max DLS was {dls_max_s}. Average DLS was {dls_avg_s}."
    p4 = (
        f"Drilling Summary:\n"
        f"A total of {total_dist:.2f} m was drilled: Rotary {total_rotary:.2f} m ({rp:.1f}%), "
        f"Slide {total_slide:.2f} m ({sp:.1f}%). Rotary Details:\n"
        f"Max rotary interval: {max_rot_s}, average ROP: {avg_rop_r_s}."
    )
    p5 = f"{wob_line}\n{rpm_line}"

    return "\n\n".join([p1, p2, p3, p4, p5])


def build_kpi_whatsapp_global_en(
    summary_df: pd.DataFrame,
    df: pd.DataFrame,
    run_summaries: list[tuple[str, str]],
) -> str:
    """Resumen global en inglés para pegar en WhatsApp."""
    n_runs = len(summary_df) if summary_df is not None and not summary_df.empty else 0
    lines = [
        f"Global drilling KPI summary ({n_runs} run(s)).",
        "",
    ]
    if summary_df is not None and not summary_df.empty:
        r_mean = summary_df["Rotary %"].mean()
        s_mean = summary_df["Slide %"].mean()
        lines.append(
            f"Across runs: mean Rotary {r_mean:.1f}%, mean Slide {s_mean:.1f}% "
            f"(by distance-weighted run table)."
        )
    corr_cols = [c for c in ["ROP", "WOB", "RPM", "DLS"] if c in df.columns]
    if len(corr_cols) >= 2:
        csub = df[corr_cols].dropna()
        if len(csub) >= 5:
            r_rw = csub["ROP"].corr(csub["WOB"]) if "ROP" in csub and "WOB" in csub else None
            r_rr = csub["ROP"].corr(csub["RPM"]) if "ROP" in csub and "RPM" in csub else None
            extra = []
            if r_rw is not None and pd.notna(r_rw):
                extra.append(f"Global ROP–WOB correlation r={r_rw:.2f}")
            if r_rr is not None and pd.notna(r_rr):
                extra.append(f"ROP–RPM r={r_rr:.2f}")
            if extra:
                lines.append("Dataset: " + "; ".join(extra) + ".")
    lines.append("")
    lines.append("— Per-run narratives below —")
    for name, _txt in run_summaries:
        lines.append(f"• {name}")
    return "\n".join(lines)


def build_kpi_whatsapp_summary_es(run: RunInfo, stats: dict, df_run: pd.DataFrame) -> str:
    """
    Resumen narrativo en español (estilo WhatsApp).
    WOB: asume kgf en datos → tonf (÷1000).
    """
    st_t = _ts_fmt_es(run.start_time)
    en_t = _ts_fmt_es(run.end_time)
    dur = (
        run.end_time - run.start_time
        if run.end_time is not None and run.start_time is not None
        else None
    )
    dur_txt = _format_duration_whatsapp_es(dur)

    d0, d1 = float(run.start_depth), float(run.end_depth)
    si, ei = stats.get("start_inc"), stats.get("end_inc")
    sa, ea = stats.get("start_az"), stats.get("end_az")

    if si is not None and ei is not None and abs(si - ei) < 1e-6:
        inc_phrase = f"inclinación constante en {si:.1f}°"
    elif si is not None and ei is not None:
        inc_phrase = f"inclinación de {si:.1f}° a {ei:.1f}°"
    else:
        inc_phrase = "inclinación n/d"

    if sa is not None and ea is not None and abs(sa - ea) < 1e-6:
        az_phrase = f"acimut constante en {sa:.1f}°"
    elif sa is not None and ea is not None:
        az_phrase = f"acimut de {sa:.1f}° a {ea:.1f}°"
    else:
        az_phrase = "acimut n/d"

    max_dls = stats.get("max_dls")
    avg_dls = stats.get("avg_dls")
    dls_max_s = f"{max_dls:.1f}°/30m" if max_dls is not None and pd.notna(max_dls) else "n/d"
    dls_avg_s = f"{avg_dls:.1f}°/30m" if avg_dls is not None and pd.notna(avg_dls) else "n/d"

    total_rotary = stats["drilling_table"][0][1]
    total_slide = stats["drilling_table"][1][1]
    total_dist = stats["drilling_table"][2][1]
    rp = stats["rotary_pct"]
    sp = stats["slide_pct"]

    max_rot = stats.get("max_rotary_len")
    if max_rot is None or (isinstance(max_rot, float) and pd.isna(max_rot)):
        max_rot_s = "n/d"
    else:
        max_rot_s = f"{float(max_rot):.2f} m"

    avg_rop_r = stats.get("avg_rop_rotary")
    avg_rop_r_s = (
        f"{float(avg_rop_r):.2f} m/h"
        if avg_rop_r is not None and pd.notna(avg_rop_r)
        else "n/d"
    )

    wmn = _wob_to_tonf(stats.get("wob_min"))
    wmx = _wob_to_tonf(stats.get("wob_max"))
    wav = _wob_to_tonf(stats.get("wob_avg"))
    wob_line = (
        f"El WOB osciló entre {wmn:.2f} y {wmx:.2f} tonf, con un promedio de {wav:.2f} tonf."
        if wmn is not None and wmx is not None and wav is not None
        else "WOB: n/d."
    )

    rpm_min, rpm_max, rpm_avg = stats.get("rpm_min"), stats.get("rpm_max"), stats.get("rpm_avg")
    rpm_line = (
        f"El RPM osciló entre {rpm_min:.2f} y {rpm_max:.2f}, con un promedio de {rpm_avg:.2f}."
        if rpm_min is not None and rpm_max is not None and rpm_avg is not None
        else "RPM: n/d."
    )

    p1 = f"La perforación comenzó el {st_t} y terminó el {en_t}, con una duración de {dur_txt}."
    p2 = (
        f"La profundidad osciló entre {d0:.2f} m y {d1:.2f} m, con {inc_phrase} y {az_phrase}."
    )
    p3 = f"DLS máximo: {dls_max_s}. DLS promedio: {dls_avg_s}."
    p4 = (
        f"Resumen de perforación:\n"
        f"Se perforaron en total {total_dist:.2f} m: Rotary {total_rotary:.2f} m ({rp:.1f}%), "
        f"Slide {total_slide:.2f} m ({sp:.1f}%). Detalle Rotary:\n"
        f"Intervalo rotary máximo: {max_rot_s}, ROP promedio: {avg_rop_r_s}."
    )
    p5 = f"{wob_line}\n{rpm_line}"

    return "\n\n".join([p1, p2, p3, p4, p5])


def build_kpi_whatsapp_global_es(
    summary_df: pd.DataFrame,
    df: pd.DataFrame,
    run_summaries: list[tuple[str, str]],
) -> str:
    """Resumen global en español para pegar en WhatsApp."""
    n_runs = len(summary_df) if summary_df is not None and not summary_df.empty else 0
    if n_runs == 0 and run_summaries:
        n_runs = len(run_summaries)
    lines = [
        f"Resumen global de KPI de perforación ({n_runs} corrida(s)).",
        "",
    ]
    if summary_df is not None and not summary_df.empty:
        r_mean = summary_df["Rotary %"].mean()
        s_mean = summary_df["Slide %"].mean()
        lines.append(
            f"A través de las corridas: Rotary medio {r_mean:.1f}%, Slide medio {s_mean:.1f}% "
            f"(tabla de corridas por distancia)."
        )
    corr_cols = [c for c in ["ROP", "WOB", "RPM", "DLS"] if c in df.columns]
    if len(corr_cols) >= 2:
        csub = df[corr_cols].dropna()
        if len(csub) >= 5:
            r_rw = csub["ROP"].corr(csub["WOB"]) if "ROP" in csub and "WOB" in csub else None
            r_rr = csub["ROP"].corr(csub["RPM"]) if "ROP" in csub and "RPM" in csub else None
            extra = []
            if r_rw is not None and pd.notna(r_rw):
                extra.append(f"correlación global ROP–WOB r={r_rw:.2f}")
            if r_rr is not None and pd.notna(r_rr):
                extra.append(f"ROP–RPM r={r_rr:.2f}")
            if extra:
                lines.append("Conjunto de datos: " + "; ".join(extra) + ".")
    lines.append("")
    lines.append("— Narrativas por corrida —")
    for name, _txt in run_summaries:
        lines.append(f"• {name}")
    return "\n".join(lines)


def save_whatsapp_summaries_excel(
    path: Path,
    run_summaries: list[tuple[str, str]],
    global_text: str,
) -> None:
    """Excel con hoja por corridas + hoja global para copiar/pegar."""
    rows = [{"Run": name, "WhatsApp_summary_ES": text} for name, text in run_summaries]
    df_runs = pd.DataFrame(rows)
    with pd.ExcelWriter(path, engine="openpyxl") as writer:
        df_runs.to_excel(writer, sheet_name="By_run", index=False)
        pd.DataFrame([{"Section": "Global_summary_ES", "Text": global_text}]).to_excel(
            writer, sheet_name="Global", index=False
        )
        full = "\n\n---\n\n".join([f"## {n}\n\n{t}" for n, t in run_summaries])
        if global_text:
            full = global_text + "\n\n---\n\n" + full
        pd.DataFrame([{"All_text_combined": full}]).to_excel(writer, sheet_name="All_combined", index=False)


# =========================
# UI: Header
# =========================
init_session_state()
render_language_selector_sidebar()

col_logo, col_title = st.columns([1, 6], vertical_alignment="center")
with col_logo:
    if LOGO_PATH.exists():
        st.image(str(LOGO_PATH), width=90)
    else:
        st.warning(f"{tr('logo_missing')} {LOGO_PATH}")

with col_title:
    st.title(APP_TITLE)

if LOGO_PATH.exists():
    try:
        st.sidebar.image(str(LOGO_PATH), width=130)
    except Exception:
        pass

st.markdown(
    f"{tr('intro_p1')}\n\n{tr('intro_p2')}\n\n{tr('intro_p3')}"
)



def _best_matching_column(columns: Iterable[str], candidates: list[str]) -> str | None:
    lowered = {str(c).strip().lower(): c for c in columns}
    for cand in candidates:
        key = str(cand).strip().lower()
        if key in lowered:
            return lowered[key]
    for cand in candidates:
        key = str(cand).strip().lower()
        for low, original in lowered.items():
            if key in low:
                return original
    return None


@st.cache_data(show_spinner=False)
def load_continuous_csv_data(file_bytes: bytes) -> pd.DataFrame:
    df = pd.read_csv(io.BytesIO(file_bytes))
    df = df.dropna(how="all").copy()
    if df.empty:
        return df

    # Muchos exports traen una primera fila con unidades (m, RPM, klbf, etc.)
    first_row = df.iloc[0]
    first_numeric_ratio = pd.to_numeric(first_row, errors="coerce").notna().mean()
    if first_numeric_ratio < 0.5:
        df = df.iloc[1:].reset_index(drop=True)

    for col in df.columns:
        df[col] = pd.to_numeric(df[col], errors="coerce")

    return df.dropna(how="all").reset_index(drop=True)


def preprocess_kpi_csv_for_rop(
    df: pd.DataFrame,
    numeric_cols: list[str],
    interpolation_method: str = "linear",
    smoothing_method: str = "rolling",
    smoothing_window: int = 7,
    polyorder: int = 2,
) -> pd.DataFrame:
    out = df.copy()
    for col in numeric_cols:
        if col in out.columns:
            out[col] = pd.to_numeric(out[col], errors="coerce")

    interp_cols = [c for c in numeric_cols if c in out.columns]
    if interp_cols:
        out[interp_cols] = out[interp_cols].interpolate(
            method=interpolation_method,
            limit_direction="both",
        )

    smooth_cols = [c for c in numeric_cols if c in out.columns]
    if not smooth_cols:
        return out

    window = max(3, int(smoothing_window))
    if window % 2 == 0:
        window += 1

    if smoothing_method == "savitzky_golay":
        for col in smooth_cols:
            series = pd.to_numeric(out[col], errors="coerce")
            valid = series.dropna()
            if len(valid) < max(window, polyorder + 2):
                out[f"{col}_smooth"] = series.rolling(window=max(3, min(5, len(series))), min_periods=1, center=True).mean()
                continue
            safe_window = min(window, len(valid) if len(valid) % 2 == 1 else len(valid) - 1)
            safe_window = max(safe_window, polyorder + 2 + ((polyorder + 2) % 2 == 0))
            if safe_window % 2 == 0:
                safe_window += 1
            safe_window = min(safe_window, len(valid) if len(valid) % 2 == 1 else len(valid) - 1)
            if safe_window <= polyorder or safe_window < 3:
                out[f"{col}_smooth"] = series.rolling(window=max(3, min(5, len(series))), min_periods=1, center=True).mean()
                continue
            smoothed = pd.Series(index=valid.index, data=savgol_filter(valid.to_numpy(), safe_window, min(polyorder, safe_window - 1)))
            out[f"{col}_smooth"] = smoothed.reindex(series.index).interpolate(limit_direction="both")
    else:
        for col in smooth_cols:
            out[f"{col}_smooth"] = (
                pd.to_numeric(out[col], errors="coerce")
                .rolling(window=window, min_periods=1, center=True)
                .mean()
            )

    return out


def compute_rop_zone_stats(
    df: pd.DataFrame,
    wob_col: str,
    rpm_col: str,
    rop_col: str,
    bins: int = 20,
    min_points_per_bin: int = 3,
    density_percentile_trim: tuple[float, float] | None = None,
) -> dict | None:
    """
    Agrupa WOB×RPM en una grilla y calcula ROP medio por celda.

    Las celdas quedan **vacías** (NaN) si:
    - no cayó ningún punto en ese bin, o
    - hay menos de ``min_points_per_bin`` puntos (para no confiar en medias con N muy bajo).

    Por defecto los bordes de la grilla van del **mínimo al máximo** de WOB y RPM en el archivo;
    si casi todo el trabajo fue en una “nube” pequeña, verás **mucho negro** fuera de esa nube.
    Opcionalmente ``density_percentile_trim=(2, 98)`` recorta outliers en el plano WOB×RPM antes
    de binar, para que los bins se concentren donde hay más datos.
    """
    needed = [wob_col, rpm_col, rop_col]
    if any(c not in df.columns for c in needed):
        return None

    d = df[needed].copy()
    d = d.dropna()
    d = d[(d[wob_col] > 0) & (d[rpm_col] > 0)]
    if d.empty:
        return None

    if density_percentile_trim is not None:
        p_lo = float(np.clip(density_percentile_trim[0], 0.0, 49.0))
        p_hi = float(np.clip(density_percentile_trim[1], 51.0, 100.0))
        wa = d[wob_col].to_numpy(dtype=float)
        ra = d[rpm_col].to_numpy(dtype=float)
        wl, wh = np.percentile(wa, p_lo), np.percentile(wa, p_hi)
        rl, rh = np.percentile(ra, p_lo), np.percentile(ra, p_hi)
        mask = d[wob_col].between(wl, wh) & d[rpm_col].between(rl, rh)
        d_sub = d.loc[mask]
        if len(d_sub) >= max(30, int(bins) * 2):
            d = d_sub

    stat, x_edges, y_edges, binnumber = binned_statistic_2d(
        d[wob_col].to_numpy(),
        d[rpm_col].to_numpy(),
        d[rop_col].to_numpy(),
        statistic="mean",
        bins=bins,
    )
    counts, _, _, _ = binned_statistic_2d(
        d[wob_col].to_numpy(),
        d[rpm_col].to_numpy(),
        d[rop_col].to_numpy(),
        statistic="count",
        bins=[x_edges, y_edges],
    )
    stat = np.where(counts >= max(1, int(min_points_per_bin)), stat, np.nan)
    if np.isnan(stat).all():
        return None

    max_idx = np.nanargmax(stat)
    i, j = np.unravel_index(max_idx, stat.shape)
    best_wob_low, best_wob_high = x_edges[i], x_edges[i + 1]
    best_rpm_low, best_rpm_high = y_edges[j], y_edges[j + 1]
    center_wob = (best_wob_low + best_wob_high) / 2.0
    center_rpm = (best_rpm_low + best_rpm_high) / 2.0
    best_rop = float(stat[i, j])
    best_count = int(counts[i, j])

    x_centers = (x_edges[:-1] + x_edges[1:]) / 2.0
    y_centers = (y_edges[:-1] + y_edges[1:]) / 2.0

    return {
        "stat": stat,
        "counts": counts,
        "x_edges": x_edges,
        "y_edges": y_edges,
        "x_centers": x_centers,
        "y_centers": y_centers,
        "best_bin": (i, j),
        "best_rop": best_rop,
        "best_wob_low": float(best_wob_low),
        "best_wob_high": float(best_wob_high),
        "best_rpm_low": float(best_rpm_low),
        "best_rpm_high": float(best_rpm_high),
        "best_wob_center": float(center_wob),
        "best_rpm_center": float(center_rpm),
        "best_count": best_count,
    }


def build_optimal_rop_heatmap(zone_stats: dict, title: str = "Heatmap ROP vs WOB-RPM") -> go.Figure:
    i, j = zone_stats["best_bin"]
    x0 = zone_stats["best_wob_low"]
    x1 = zone_stats["best_wob_high"]
    y0 = zone_stats["best_rpm_low"]
    y1 = zone_stats["best_rpm_high"]

    stat = np.asarray(zone_stats["stat"], dtype=float)
    z_t = stat.T
    z_ok = stat[np.isfinite(stat)]
    if z_ok.size:
        zlo, zhi = float(np.nanmin(z_ok)), float(np.nanmax(z_ok))
        zpad = max(0.15, (zhi - zlo) * 0.04)
        zmin_hm, zmax_hm = zlo - zpad, zhi + zpad
    else:
        zmin_hm, zmax_hm = None, None

    counts = zone_stats.get("counts")
    custom_cd = None
    if counts is not None:
        c_arr = np.asarray(counts, dtype=float)
        if c_arr.shape == stat.shape:
            custom_cd = c_arr.T

    hm_text = _rop_heatmap_label_matrix_top_fraction(z_t)
    hm_extras: dict = {
        "text": hm_text,
        "texttemplate": "%{text}",
        "textfont": dict(size=12, color="rgba(248,250,252,0.92)"),
    }
    _bw = float(zone_stats["best_wob_center"])
    _br = float(zone_stats["best_rpm_center"])
    _brop = float(zone_stats["best_rop"])
    y_centers = np.asarray(zone_stats["y_centers"], dtype=float)
    _yr_span = float(np.nanmax(y_centers) - np.nanmin(y_centers)) if y_centers.size else 1.0
    _dy = (
        float(np.median(np.abs(np.diff(np.sort(y_centers)))))
        if len(y_centers) > 1
        else max(_yr_span * 0.04, 1.0)
    )

    fig = go.Figure(
        data=[
            go.Heatmap(
                x=zone_stats["x_centers"],
                y=zone_stats["y_centers"],
                z=z_t,
                zmin=zmin_hm,
                zmax=zmax_hm,
                colorscale=ROP_HEATMAP_COLORSCALE,
                zsmooth="best",
                xgap=2,
                ygap=2,
                colorbar=dict(title="ROP medio", tickformat=".1f", thickness=20),
                customdata=custom_cd,
                hovertemplate=(
                    "WOB: %{x:.2f}<br>RPM: %{y:.2f}<br>ROP medio: %{z:.2f}"
                    + ("<br>Puntos: %{customdata:.0f}" if custom_cd is not None else "")
                    + "<extra></extra>"
                ),
                hoverongaps=False,
                **hm_extras,
            ),
            go.Scatter(
                x=[_bw],
                y=[_br],
                mode="markers",
                marker=dict(
                    symbol="star",
                    size=18,
                    color="#FDE047",
                    line=dict(color="rgb(15,23,42)", width=1.2),
                ),
                name="Pico",
                showlegend=False,
                hovertemplate=(
                    "Pico operacional<br>WOB: %{x:.2f}<br>RPM: %{y:.2f}<br>"
                    f"ROP: {_brop:.1f}<extra></extra>"
                ),
            ),
        ]
    )
    fig.add_shape(
        type="rect",
        x0=x0,
        x1=x1,
        y0=y0,
        y1=y1,
        line=dict(color="#ffffff", width=3),
        fillcolor="rgba(255,255,255,0)",
    )
    fig.add_annotation(
        x=_bw,
        y=_br + 1.35 * _dy,
        text="Pico",
        showarrow=True,
        arrowhead=2,
        arrowsize=1,
        arrowwidth=1.5,
        arrowcolor="rgba(248,250,252,0.85)",
        ax=0,
        ay=-24,
        font=dict(color="#F8FAFC", size=11),
        bgcolor="rgba(0,0,0,0.5)",
        borderpad=3,
    )
    fig.add_annotation(
        x=(x0 + x1) / 2.0,
        y=(y0 + y1) / 2.0,
        text=f"Mejor zona<br>ROP {_brop:.1f}",
        showarrow=False,
        font=dict(color="#ffffff", size=12),
        bgcolor="rgba(0,0,0,0.4)",
    )
    fig.update_layout(title=title, xaxis_title="WOB", yaxis_title="RPM")
    return prettify_heatmap_auto(fig, h=640)


def render_kpi_csv_optimizer() -> None:
    st.markdown("### Optimizador de ROP desde CSV")
    st.caption(
        "Carga un CSV continuo para interpolar huecos, suavizar curvas y encontrar la mejor zona de ROP según WOB y RPM."
    )

    uploaded_csv = st.file_uploader(
        "Sube CSV de datos continuos",
        type=["csv"],
        key="kpi_csv_optimizer_upload",
        help="Ejemplo esperado: DEPTH, ROP, Surface Torque, WOB, Surface RPM.",
    )
    if uploaded_csv is None:
        return

    try:
        df_raw = load_continuous_csv_data(_uploaded_file_bytes(uploaded_csv))
    except Exception as e:
        st.error(f"No pude leer el CSV: {e}")
        return

    if df_raw.empty:
        st.warning("El CSV no contiene filas útiles después de limpiar encabezados/unidades.")
        return

    st.success(f"CSV cargado: {len(df_raw):,} filas útiles y {len(df_raw.columns)} columnas.")

    default_depth = _best_matching_column(df_raw.columns, ["DEPTH", "Depth", "Bit Depth", "Hole Depth", "MD"])
    default_rop = _best_matching_column(df_raw.columns, ["ROP", "ROP Avg", "Rate of Penetration"])
    default_wob = _best_matching_column(df_raw.columns, ["WOB", "Weight on Bit", "Hookload WOB"])
    default_rpm = _best_matching_column(df_raw.columns, ["Surface RPM", "RPM", "Rotary RPM"])
    default_torque = _best_matching_column(df_raw.columns, ["Surface Torque", "Torque", "Torque (Surface)"])

    cols1 = st.columns(5)
    with cols1[0]:
        depth_col = st.selectbox("Profundidad", options=["<ninguna>"] + list(df_raw.columns), index=(list(df_raw.columns).index(default_depth) + 1 if default_depth in df_raw.columns else 0), key="kpi_csv_depth_col")
    with cols1[1]:
        rop_col = st.selectbox("ROP", options=list(df_raw.columns), index=(list(df_raw.columns).index(default_rop) if default_rop in df_raw.columns else 0), key="kpi_csv_rop_col")
    with cols1[2]:
        wob_col = st.selectbox("WOB", options=list(df_raw.columns), index=(list(df_raw.columns).index(default_wob) if default_wob in df_raw.columns else 0), key="kpi_csv_wob_col")
    with cols1[3]:
        rpm_col = st.selectbox("RPM", options=list(df_raw.columns), index=(list(df_raw.columns).index(default_rpm) if default_rpm in df_raw.columns else 0), key="kpi_csv_rpm_col")
    with cols1[4]:
        torque_options = ["<ninguna>"] + list(df_raw.columns)
        torque_idx = (torque_options.index(default_torque) if default_torque in torque_options else 0)
        torque_col = st.selectbox("Torque (opcional)", options=torque_options, index=torque_idx, key="kpi_csv_torque_col")

    cols2 = st.columns(5)
    with cols2[0]:
        interpolation_method = st.selectbox("Interpolación", options=["linear", "nearest"], index=0, key="kpi_csv_interp")
    with cols2[1]:
        smoothing_method = st.selectbox("Suavizado", options=["rolling", "savitzky_golay"], index=0, key="kpi_csv_smooth_method")
    with cols2[2]:
        smoothing_window = st.number_input("Ventana suavizado", min_value=3, max_value=51, value=7, step=2, key="kpi_csv_smooth_window")
    with cols2[3]:
        bins = st.number_input("Bins WOB-RPM", min_value=8, max_value=60, value=20, step=1, key="kpi_csv_bins")
    with cols2[4]:
        min_points = st.number_input("Mín. puntos/bin", min_value=1, max_value=25, value=3, step=1, key="kpi_csv_min_points")

    focus_dense_cloud = st.checkbox(
        "Enfocar heatmap en núcleo de datos (P2–P98 WOB y RPM)",
        value=False,
        key="kpi_csv_focus_dense_hm",
        help="Quita outliers en el plano WOB×RPM antes de armar la grilla: el mapa usa el rango donde está la mayoría de puntos y se reduce el área negra vacía.",
    )

    selected_numeric = [c for c in [rop_col, wob_col, rpm_col] if c and c != "<ninguna>"]
    if depth_col != "<ninguna>":
        selected_numeric.append(depth_col)
    if torque_col != "<ninguna>":
        selected_numeric.append(torque_col)

    df_processed = preprocess_kpi_csv_for_rop(
        df_raw,
        numeric_cols=selected_numeric,
        interpolation_method=interpolation_method,
        smoothing_method=smoothing_method,
        smoothing_window=int(smoothing_window),
    )

    rop_work_col = f"{rop_col}_smooth" if f"{rop_col}_smooth" in df_processed.columns else rop_col
    wob_work_col = f"{wob_col}_smooth" if f"{wob_col}_smooth" in df_processed.columns else wob_col
    rpm_work_col = f"{rpm_col}_smooth" if f"{rpm_col}_smooth" in df_processed.columns else rpm_col

    zone_stats = compute_rop_zone_stats(
        df_processed,
        wob_col=wob_work_col,
        rpm_col=rpm_work_col,
        rop_col=rop_work_col,
        bins=int(bins),
        min_points_per_bin=int(min_points),
        density_percentile_trim=(2.0, 98.0) if focus_dense_cloud else None,
    )

    na_counts = df_raw[[c for c in [rop_col, wob_col, rpm_col] if c in df_raw.columns]].isna().sum()
    filled_points = int(na_counts.sum())

    m1, m2, m3, m4 = st.columns(4)
    m1.metric("Filas útiles", f"{len(df_processed):,}")
    m2.metric("Huecos detectados", f"{filled_points:,}")
    m3.metric("ROP promedio", format_num(pd.to_numeric(df_processed[rop_work_col], errors="coerce").mean()))
    m4.metric("ROP máximo", format_num(pd.to_numeric(df_processed[rop_work_col], errors="coerce").max()))

    stats_cols = [c for c in [depth_col if depth_col != "<ninguna>" else None, rop_col, wob_col, rpm_col, torque_col if torque_col != "<ninguna>" else None] if c]
    stats_df = df_processed[stats_cols].agg(["min", "mean", "max"]).T.reset_index().rename(columns={"index": "Parámetro", "min": "Mínimo", "mean": "Promedio", "max": "Máximo"})
    st.dataframe(stats_df, use_container_width=True, hide_index=True)
    _rop_opt_cols = [wob_work_col, rpm_work_col, rop_work_col]
    _st_rop = heatmap_numeric_stats(df_processed, _rop_opt_cols)
    _ch_rop = stats_df_to_heatmap_chips(_st_rop, max_chips=8)
    if _ch_rop:
        st.caption("**Chips — min–max y media (WOB, RPM, ROP trabajados)**")
        _render_chips_row(_ch_rop)

    if zone_stats is None:
        st.warning("No encontré suficiente data válida para calcular la mejor zona de ROP. Revisa columnas, bins o cantidad mínima por celda.")
    else:
        st.success(
            "Mejor zona de ROP encontrada: "
            f"WOB {zone_stats['best_wob_low']:.2f}–{zone_stats['best_wob_high']:.2f}, "
            f"RPM {zone_stats['best_rpm_low']:.2f}–{zone_stats['best_rpm_high']:.2f}, "
            f"ROP medio {zone_stats['best_rop']:.2f} ({zone_stats['best_count']} puntos)."
        )
        _dash_chips = rop_zone_dashboard_chips(zone_stats)
        if _dash_chips:
            st.caption("**Resumen rápido (mejor zona y grilla)**")
            _render_chips_row(_dash_chips)
        fig_dash = build_optimal_rop_heatmap_with_marginals(
            zone_stats,
            title="Heatmap de mejor zona de ROP",
        )
        if fig_dash is not None:
            st.plotly_chart(fig_dash, use_container_width=True, config=PLOTLY_CONFIG_ROP_DASH)
            st.caption(
                "**Heatmap:** degradado azul → teal → amarillo/naranja; marginales = ROP máx/mín por bin de WOB y de RPM."
            )
            st.caption(
                "**Zonas vacías (negro):** en ese bin WOB×RPM no hubo puntos **o** hay menos muestras que tu **mín. puntos/bin** "
                "(la media no se muestra para no inflar ruido). Además la grilla abarca el **rango completo** de WOB/RPM del archivo: "
                "si casi todo el trabajo fue en una “nube” (p. ej. alto RPM y WOB medio), el resto del rectángulo se ve vacío. "
                "Activa **Enfocar heatmap en núcleo P2–P98** arriba para recortar outliers y ampliar la zona útil. "
                "En el gráfico, **no hay tooltip** al pasar el mouse por celdas sin datos (evita el mensaje confuso ROP NaN / 0 puntos)."
            )
        else:
            fig_hm = build_optimal_rop_heatmap(zone_stats, title="Heatmap de mejor zona de ROP")
            st.plotly_chart(fig_hm, use_container_width=True, config=PLOTLY_CONFIG)
            st.caption(
                "**Zonas vacías:** igual que en el dashboard — pocos puntos por celda o combinación WOB×RPM no usada; opción **P2–P98** para enfocar el núcleo."
            )

    _depth_sel = depth_col if depth_col != "<ninguna>" else None
    _fig_top_zones = build_rop_top_zones_bar_figure(zone_stats, top_n=8) if zone_stats is not None else None
    _fig_depth_curves, _depth_franja_chips = build_kpi_depth_curves_figure(
        df_processed,
        _depth_sel,
        rop_work_col,
        wob_work_col,
        rpm_work_col,
        zone_stats=zone_stats,
    )
    if _fig_top_zones is not None or _fig_depth_curves is not None:
        st.markdown("#### Contexto operativo")
        _ctx_lines = []
        if zone_stats is not None:
            _ctx_lines.append(
                "En el heatmap, la **estrella** marca el pico operativo y el **recuadro blanco** la celda óptima; "
                f"solo se etiquetan celdas del **top {int(ROP_HEATMAP_LABEL_TOP_FRACTION * 100)} %** por ROP."
            )
            _ctx_lines.append(
                "En **curvas vs profundidad**, las **franjas naranjas** marcan tramos donde WOB y RPM caen en esa misma celda óptima."
            )
        _ctx_lines.append(
            "**Top zonas** y **curvas suavizadas** usan **valores reales** (ROP, WOB, RPM en sus unidades), sin normalizar el eje Y."
        )
        st.caption(" ".join(_ctx_lines))
        if _fig_top_zones is not None and _fig_depth_curves is not None:
            _z1, _z2 = st.columns(2)
            with _z1:
                st.plotly_chart(_fig_top_zones, use_container_width=True, config=PLOTLY_CONFIG_ROP_DASH)
            with _z2:
                if _depth_franja_chips:
                    st.caption("**Chips — franja de mejor ROP (celda WOB×RPM)**")
                    _render_chips_row(_depth_franja_chips)
                st.plotly_chart(_fig_depth_curves, use_container_width=True, config=PLOTLY_CONFIG_ROP_DASH)
        elif _fig_depth_curves is not None:
            if _depth_franja_chips:
                st.caption("**Chips — franja de mejor ROP (celda WOB×RPM)**")
                _render_chips_row(_depth_franja_chips)
            st.plotly_chart(_fig_depth_curves, use_container_width=True, config=PLOTLY_CONFIG_ROP_DASH)
        else:
            st.plotly_chart(_fig_top_zones, use_container_width=True, config=PLOTLY_CONFIG_ROP_DASH)

    with st.expander("Perfil min–media–max (comparación normalizada)", expanded=False):
        st.caption(
            "**Qué hace el normalizado:** para cada variable (WOB, RPM, ROP) se toma el **mínimo y máximo** del archivo y se dibuja una escala **0 → 1**. "
            "La **línea gris** va de min a max; el **punto azul** es la **media** en ese rango. Sirve para comparar *dónde cae la media* entre parámetros de distinta magnitud; "
            "**el eje Y no son unidades físicas**. Para magnitudes reales usa la tabla, el heatmap y **Curvas suavizadas** arriba."
        )
        _sp_rop = build_minmax_mean_spine_figure(
            _st_rop,
            title="Perfil min · media · max — WOB, RPM, ROP (normalizado 0–1 por variable)",
        )
        if _sp_rop is not None:
            st.plotly_chart(_sp_rop, use_container_width=True, config=PLOTLY_CONFIG)

    if zone_stats is not None:
        st.caption(
            "La celda óptima del heatmap es la combinación WOB–RPM con mayor ROP promedio por bin, tras interpolar y suavizar."
        )

    preview_cols = [c for c in [depth_col if depth_col != "<ninguna>" else None, rop_col, rop_work_col if rop_work_col != rop_col else None, wob_col, wob_work_col if wob_work_col != wob_col else None, rpm_col, rpm_work_col if rpm_work_col != rpm_col else None] if c]
    with st.expander("Vista previa de datos procesados", expanded=False):
        st.dataframe(df_processed[preview_cols].head(200), use_container_width=True, hide_index=True)


def render_kpi_module() -> None:
    # =========================
    # Inputs
    # =========================
    data_source = st.radio(
        tr("data_source"),
        ["Excel", "API"],
        horizontal=True,
        format_func=lambda x: tr(f"src_{x.lower()}"),
    )
    uploaded = None
    if data_source == "Excel":
        uploaded = st.file_uploader(tr("upload_excel"), type=["xlsx"])

    render_kpi_csv_optimizer()
    st.divider()

    with st.sidebar:
        st.header(tr("sidebar_options"))
        show_plots = st.checkbox(tr("show_plots"), True)

        st.subheader(tr("region_captures"))
        enable_region_captures = st.checkbox(tr("enable_region_captures"), False)
        n_captures = st.number_input(tr("n_captures"), 1, 10, 3)
        interval_seconds = st.number_input(tr("interval_captures"), 1, 60, 10)
        region_x = st.number_input(tr("region_x"), 0, 5000, 456)
        region_y = st.number_input(tr("region_y"), 0, 5000, 196)
        region_w = st.number_input(tr("region_w"), 10, 8000, 2088)
        region_h = st.number_input(tr("region_h"), 10, 8000, 1319)
        capture_title_prefix = st.text_input(
            tr("capture_prefix"), "Drill Spot – Regional Snapshot"
        )

    preloaded_excel_df = None
    detected_excel_format = None
    if data_source == "Excel" and uploaded is not None:
        try:
            preloaded_excel_df = load_excel_data(uploaded)
            detected_excel_format = str(preloaded_excel_df.get("Source Format", pd.Series(["standard"])).dropna().iloc[0]) if "Source Format" in preloaded_excel_df.columns else "standard"
            st.caption(f"Formato detectado automáticamente: **{detected_excel_format}**")
        except Exception as e:
            st.warning(f"No pude previsualizar el Slide Sheet: {e}")

    st.subheader(tr("runs_header"))
    default_runs = pd.DataFrame(
        [
            {
                "Run": "Run 100",
                "start_depth": 14,
                "end_depth": 66.32,
                "start_time": "2025-10-20T20:52:25-06:00",
                "end_time": "2025-10-25T13:52:10-06:00",
            },
            {
                "Run": "Run 200",
                "start_depth": 66.32,
                "end_depth": 130.94,
                "start_time": "2025-10-25T14:36:15-06:00",
                "end_time": "2025-10-25T23:54:25-06:00",
            },
            {
                "Run": "Run 300",
                "start_depth": 130,
                "end_depth": 921.11,
                "start_time": "2025-10-27T00:00:10-06:00",
                "end_time": "2025-10-31T06:42:45-06:00",
            },
            {
                "Run": "Run 400",
                "start_depth": 921,
                "end_depth": 1212.3,
                "start_time": "2025-11-02T06:26:25-06:00",
                "end_time": "2025-03-18T18:21:49-07:00",
            },
            {
                "Run": "Run 500",
                "start_depth": 1212,
                "end_depth": 2000.03,
                "start_time": "2025-11-15T13:47:40-06:00",
                "end_time": "2025-11-17T13:59:35-06:00",
            },
            {
                "Run": "Run 600",
                "start_depth": 2000,
                "end_depth": 3932.13,
                "start_time": "2025-11-21T09:45-06:00",
                "end_time": "2025-12-02T21:05:45-06:00",
            },
        ]
    )
    if preloaded_excel_df is not None:
        auto_runs = suggest_runs_from_loaded_df(preloaded_excel_df)
        if not auto_runs.empty:
            default_runs = auto_runs
            st.caption("Se detectaron corridas automáticamente desde el Excel. Puedes editarlas abajo si quieres.")

    runs_df = st.data_editor(default_runs, num_rows="dynamic", use_container_width=True)
    api_df_ready = st.session_state.get("kpi_df_api") is not None
    data_ready = uploaded is not None if data_source == "Excel" else api_df_ready
    generate = st.button(tr("generate_pptx"), type="primary", disabled=not data_ready)

    if data_source == "API":
        with st.expander(tr("api_config_slide"), expanded=False):
            base_url = st.text_input(tr("api_base_url"), value=API_DEFAULT_BASE_URL)
            token_default = st.session_state.get("api_token") or API_DEFAULT_TOKEN or ""
            token = st.text_input(tr("api_token"), value=token_default, type="password")
            slide_sheet_path = st.text_input(
                tr("api_slide_path"),
                value="/public/api/v1/wells/{well_uuid}/intervals/slide-sheet",
                help=tr("api_slide_help"),
            )
            use_depth_range = st.checkbox(tr("use_depth_range"), value=False)
            depth_from = st.number_input(
                tr("depth_from"),
                value=0.0,
                step=10.0,
                disabled=not use_depth_range,
            )
            depth_to = st.number_input(
                tr("depth_to"),
                value=0.0,
                step=10.0,
                disabled=not use_depth_range,
            )
            st.session_state["api_token"] = token
            st.caption(tr("api_note_html"))

        if not token:
            st.info(tr("enter_token"))
            return

        try:
            projects_resp = api_list_projects(base_url, token)
            projects = normalize_list_response(projects_resp)
        except Exception as e:
            st.error(f"{tr('list_projects_err')} {e}")
            return

        if not projects:
            st.info(tr("no_projects"))
            return

        project_map = {
            f"{p.get('name', tr('unnamed'))} ({p.get('uuid', 'n/a')})": p for p in projects
        }
        project_label = st.selectbox(tr("project"), list(project_map.keys()), key="bha_project_label_legacy")
        project_uuid = project_map[project_label].get("uuid")

        if not project_uuid:
            st.error(tr("project_no_uuid"))
            return

        try:
            wells_resp = api_list_wells(base_url, token, project_uuid)
            wells = normalize_list_response(wells_resp)
        except Exception as e:
            st.error(f"{tr('list_wells_err')} {e}")
            return

        if not wells:
            st.info(tr("no_wells"))
            return

        well_map = {
            f"{w.get('name', tr('unnamed'))} ({w.get('uuid', 'n/a')})": w for w in wells
        }
        well_label = st.selectbox(tr("well"), list(well_map.keys()), key="bha_well_label_legacy")
        well_uuid = well_map[well_label].get("uuid")

        if not well_uuid:
            st.error(tr("well_no_uuid"))
            return

        st.markdown(tr("kpi_rt"))
        kpi_auto_refresh = st.checkbox(
            tr("kpi_auto_refresh"),
            value=st.session_state.get("kpi_auto_refresh", False),
            key="kpi_auto_refresh",
            help=tr("kpi_auto_help"),
        )
        if kpi_auto_refresh:
            st.number_input(
                tr("interval_seconds"),
                min_value=10,
                max_value=300,
                value=30,
                step=10,
                key="kpi_auto_refresh_interval",
                help=tr("kpi_interval_help"),
            )

        kpi_auto_rerun = st.session_state.pop("kpi_auto_rerun_trigger", False)
        if kpi_auto_rerun and st.session_state.get("kpi_auto_refresh"):
            try:
                df_api = slide_sheet_to_df(
                    api_get_slide_sheet_intervals(
                        base_url,
                        token,
                        well_uuid,
                        depth_from if use_depth_range else None,
                        depth_to if use_depth_range else None,
                        slide_sheet_path or None,
                    )
                )
                if not df_api.empty:
                    st.session_state["kpi_df_api"] = df_api
            except Exception:
                pass

        if st.button(tr("load_slide_sheet"), type="secondary"):
            try:
                df_api = slide_sheet_to_df(
                    api_get_slide_sheet_intervals(
                        base_url,
                        token,
                        well_uuid,
                        depth_from if use_depth_range else None,
                        depth_to if use_depth_range else None,
                        slide_sheet_path or None,
                    )
                )
                if df_api.empty:
                    st.warning(tr("kpi_api_no_intervals"))
                else:
                    st.session_state["kpi_df_api"] = df_api
                    st.success(tr("kpi_slide_loaded").format(n=len(df_api)))
                    st.rerun()  # Re-evaluar data_ready y habilitar "Generate PPTX"
            except Exception as e:
                st.error(f"{tr('kpi_load_err')} {e}")

        if st.session_state.get("kpi_df_api") is not None:
            st.caption(tr("kpi_active_source"))
            show_options = [50, 100, 200, 500, "Todos"]
            if "kpi_rows_preview" not in st.session_state:
                st.session_state["kpi_rows_preview"] = 200
            show_choice = st.selectbox(
                tr("show_rows"),
                show_options,
                index=show_options.index(st.session_state["kpi_rows_preview"]),
                key="kpi_rows_preview",
                format_func=lambda x: tr("all_rows") if x == "Todos" else str(x),
            )
            df_api = st.session_state["kpi_df_api"]
            df_preview = df_api if show_choice == "Todos" else df_api.head(int(show_choice))
            st.dataframe(df_preview, use_container_width=True, hide_index=True)

            if st.session_state.get("kpi_auto_refresh") and not generate:
                interval = int(st.session_state.get("kpi_auto_refresh_interval", 30))
                interval = max(10, min(300, interval))
                countdown_placeholder = st.empty()
                for i in range(interval, 0, -1):
                    countdown_placeholder.info(tr("next_refresh").format(i=i))
                    time.sleep(1)
                countdown_placeholder.empty()
                st.session_state["kpi_auto_rerun_trigger"] = True
                st.rerun()

    # =========================
    # Generar reporte
    # =========================
    if generate and data_ready:
        runs = parse_runs(runs_df)

        with st.spinner("Leyendo datos y limpiando..."):
            try:
                if data_source == "Excel":
                    df = load_excel_data(uploaded)
                else:
                    df = st.session_state.get("kpi_df_api")
                    if df is None or df.empty:
                        st.error("Primero carga el Slide Sheet desde la API.")
                        st.stop()
            except Exception as e:
                st.error(f"No pude leer los datos: {e}")
                st.stop()
            df = normalize_dataframe(df)

        prs = Presentation()
        add_title_slide(prs)
        tmp_dir = Path(tempfile.mkdtemp())
        pptx_path = tmp_dir / "Final_Drilling_KPI_Report.pptx"
        whatsapp_excel_path = tmp_dir / "Drilling_KPI_WhatsApp_Summary.xlsx"

        summary_stats = []
        all_wob, all_rpm = [], []
        whatsapp_run_texts: list[tuple[str, str]] = []

        with st.spinner("Procesando corridas y generando BI Pro..."):
            for run in runs:
                df_run = df[
                    (df["End Depth"] >= run.start_depth)
                    & (df["Start Depth"] <= run.end_depth)
                ].copy()

                if df_run.empty:
                    st.warning(f"⚠️ {run.name} no tiene datos válidos.")
                    continue

                stats = compute_run_stats(df_run, run)
                if stats is None:
                    st.warning(f"⚠️ {run.name} no tiene datos suficientes tras limpieza.")
                    continue
                df_run = stats["df_run"]

                all_wob.extend(df_run["WOB"].dropna().tolist())
                all_rpm.extend(df_run["RPM"].dropna().tolist())

                try:
                    wa_txt = build_kpi_whatsapp_summary_es(run, stats, df_run)
                    whatsapp_run_texts.append((run.name, wa_txt))
                except Exception:
                    whatsapp_run_texts.append(
                        (run.name, "(No se pudo generar el resumen WhatsApp para esta corrida.)")
                    )

                # Tablas
                with st.expander(f"📋 Tablas – {run.name}", expanded=False):
                    st.markdown("**General Information**")
                    st.dataframe(
                        pd.DataFrame(stats["general_table"], columns=["Parameter", "Value"]),
                        use_container_width=True,
                        hide_index=True,
                    )
                    st.markdown("**Drilling Summary**")
                    st.dataframe(
                        pd.DataFrame(
                            stats["drilling_table"],
                            columns=["Mode", "Distance (m)", "Percentage (%)"],
                        ),
                        use_container_width=True,
                        hide_index=True,
                    )
                    st.markdown("**Slide Specifics**")
                    st.dataframe(
                        pd.DataFrame(stats["slide_table"], columns=["Parameter", "Value"]),
                        use_container_width=True,
                        hide_index=True,
                    )

                add_table_slide(
                    prs,
                    f"{run.name} – General Information",
                    stats["general_table"],
                    ["Parameter", "Value"],
                )
                add_table_slide(
                    prs,
                    f"{run.name} – Drilling Summary",
                    stats["drilling_table"],
                    ["Mode", "Distance (m)", "Percentage (%)"],
                )
                add_table_slide(
                    prs,
                    f"{run.name} – Slide Specifics",
                    stats["slide_table"],
                    ["Parameter", "Value"],
                )

                # Gauge + chips pro
                if show_plots:
                    st.subheader(f"Mechanical Efficiency – {run.name}")
                    render_efficiency_chips(
                        stats["rotary_pct"], stats["slide_pct"], run.name
                    )
                    vspace(4)
                    gauge_efficiency(stats["rotary_pct"], run.name)
                    chart_notes(
                        f"Rotary {format_num(stats['rotary_pct'])}% vs Slide {format_num(stats['slide_pct'])}%.",
                        "Indicador muestra porcentaje de metros en modo Rotary.",
                    )
                    vspace(12)

                # Rotary vs Slide
                fig_rot_slide = px.bar(
                    x=["Rotary", "Slide"],
                    y=[stats["drilling_table"][0][1], stats["drilling_table"][1][1]],
                    text=[
                        round(stats["drilling_table"][0][1], 1),
                        round(stats["drilling_table"][1][1], 1),
                    ],
                    labels={"x": "Mode", "y": "Meters Drilled"},
                    title=f"{run.name} – Rotary vs Slide",
                    color=["Rotary", "Slide"],
                    color_discrete_map={"Rotary": "#2563EB", "Slide": "#F59E0B"},
                )
                fig_rot_slide.update_traces(textposition="outside", marker_line_width=0)
                save_and_show_plotly(prs, f"{run.name} – Rotary vs Slide", fig_rot_slide, show_plots)
                chart_notes(
                    f"Rotary {format_num(stats['drilling_table'][0][1])} m "
                    f"({format_num(stats['rotary_pct'])}%) y Slide "
                    f"{format_num(stats['drilling_table'][1][1])} m "
                    f"({format_num(stats['slide_pct'])}%).",
                    "Barras comparan metros perforados por modo.",
                )
                vspace(14)

                # Scatter Pair
                st.subheader(f"{run.name} – ROP Relationships")
                _render_chips_row(
                    [
                        (run.name, "gray"),
                        (f"n={len(df_run):,}", "gray"),
                        (f"ROP–WOB r={safe_corr(df_run, 'ROP', 'WOB')}", "blue"),
                        (f"ROP–RPM r={safe_corr(df_run, 'ROP', 'RPM')}", "blue"),
                    ]
                )
                vspace(6)
                c1, c2 = st.columns(2, gap="large", vertical_alignment="top")

                fig_rop_wob = px.scatter(
                    df_run,
                    x="WOB",
                    y="ROP",
                    trendline="ols",
                    labels={"WOB": "WOB (kgf)", "ROP": "ROP (m/hr)"},
                    title=f"{run.name} – ROP vs WOB",
                )
                fig_rop_rpm = px.scatter(
                    df_run,
                    x="RPM",
                    y="ROP",
                    trendline="ols",
                    labels={"RPM": "RPM", "ROP": "ROP (m/hr)"},
                    title=f"{run.name} – ROP vs RPM",
                )
                fig_rop_wob.update_traces(marker=dict(size=6, opacity=0.7))
                fig_rop_rpm.update_traces(marker=dict(size=6, opacity=0.7))
                set_trendline_color(fig_rop_wob)
                set_trendline_color(fig_rop_rpm)

                if show_plots:
                    with c1:
                        st.plotly_chart(
                            prettify(fig_rop_wob),
                            use_container_width=True,
                            config=PLOTLY_CONFIG,
                        )
                        chart_notes(
                            f"Correlación ROP vs WOB: {safe_corr(df_run, 'ROP', 'WOB')}. n = {len(df_run):,} intervalos.",
                            "X=WOB (kgf), Y=ROP (m/hr). Línea de tendencia OLS.",
                        )
                    with c2:
                        st.plotly_chart(
                            prettify(fig_rop_rpm),
                            use_container_width=True,
                            config=PLOTLY_CONFIG,
                        )
                        chart_notes(
                            f"Correlación ROP vs RPM: {safe_corr(df_run, 'ROP', 'RPM')}. n = {len(df_run):,} intervalos.",
                            "X=RPM, Y=ROP (m/hr). Línea de tendencia OLS.",
                        )

                save_and_show_plotly(prs, f"{run.name} – ROP vs WOB", fig_rop_wob, False)
                save_and_show_plotly(prs, f"{run.name} – ROP vs RPM", fig_rop_rpm, False)
                vspace(18)

                # DLS vs MD
                st.subheader(f"{run.name} – DLS vs MD")
                dls_vals = df_run["DLS"].dropna() if "DLS" in df_run.columns else pd.Series(dtype=float)
                _render_chips_row(
                    [(run.name, "gray"), (f"n={len(dls_vals):,}", "gray"), (f"DLS: {series_summary(dls_vals)}", "blue")] if len(dls_vals) else [(run.name, "gray")]
                )
                vspace(6)
                fig_dls = build_dls_vs_md_figure(df_run, run.name)
                if fig_dls is None:
                    st.info("No hay datos suficientes para DLS vs MD.")
                else:
                    if show_plots:
                        st.plotly_chart(
                            prettify(fig_dls),
                            use_container_width=True,
                            config=PLOTLY_CONFIG,
                        )
                        dls_series = df_run["DLS"].dropna()
                        chart_notes(
                            f"{series_summary(dls_series)}.",
                            "X=MD (m), Y=DLS (°/30m), línea con marcadores.",
                        )
                    save_and_show_plotly(prs, f"{run.name} – DLS vs MD", fig_dls, False)
                vspace(18)

                # Boxplots por modo
                st.subheader(f"{run.name} – Boxplots por modo")
                vspace(6)
                box_metrics = [
                    ("ROP", "ROP"),
                    ("WOB", "WOB"),
                    ("RPM", "RPM"),
                    ("DLS", "DLS"),
                ]
                c1, c2 = st.columns(2, gap="large", vertical_alignment="top")
                for idx, (col, label) in enumerate(box_metrics):
                    fig_box = build_boxplot_by_mode(
                        df_run, col, f"{run.name} – {label} por modo"
                    )
                    if fig_box is None:
                        continue
                    if show_plots:
                        with (c1 if idx % 2 == 0 else c2):
                            st.plotly_chart(
                                prettify(fig_box),
                                use_container_width=True,
                                config=PLOTLY_CONFIG,
                            )
                            medians = (
                                df_run.dropna(subset=[col, "Mode_norm"])
                                .groupby("Mode_norm")[col]
                                .median()
                                .to_dict()
                            )
                            summary_parts = [
                                f"{k} mediana {format_num(v)}" for k, v in medians.items()
                            ]
                            chart_notes(
                                ", ".join(summary_parts) + ".",
                                "Caja=IQR, línea=mediana, puntos=outliers.",
                            )
                    save_and_show_plotly(prs, f"{run.name} – {label} por modo", fig_box, False)
                vspace(18)

                # Control charts
                st.subheader(f"{run.name} – Control charts")
                vspace(6)
                control_metrics = ["ROP", "WOB", "RPM"]
                c1, c2 = st.columns(2, gap="large", vertical_alignment="top")
                for idx, col in enumerate(control_metrics):
                    fig_ctrl = build_control_chart(df_run, col, run.name)
                    if fig_ctrl is None:
                        continue
                    if show_plots:
                        with (c1 if idx % 2 == 0 else c2):
                            st.plotly_chart(
                                prettify(fig_ctrl),
                                use_container_width=True,
                                config=PLOTLY_CONFIG,
                            )
                            series = df_run[col].dropna()
                            mean_val = series.mean()
                            std_val = series.std()
                            upper = mean_val + 3 * std_val
                            lower = mean_val - 3 * std_val
                            outliers = ((series > upper) | (series < lower)).sum()
                            chart_notes(
                                f"Media {format_num(mean_val)}, "
                                f"±3σ [{format_num(lower)}, {format_num(upper)}], "
                                f"outliers {outliers}.",
                                "Línea central=media, líneas punteadas=límites.",
                            )
                    save_and_show_plotly(prs, f"{run.name} – Control chart {col}", fig_ctrl, False)
                vspace(18)

                # Crossplot DLS vs ROP
                st.subheader(f"{run.name} – DLS vs ROP (trade-off)")
                vspace(6)
                fig_cross = build_crossplot_dls_rop(df_run, run.name)
                if fig_cross is None:
                    st.info("No hay datos suficientes para DLS vs ROP.")
                else:
                    if show_plots:
                        st.plotly_chart(
                            prettify(fig_cross),
                            use_container_width=True,
                            config=PLOTLY_CONFIG,
                        )
                        chart_notes(
                            f"Correlación DLS vs ROP: {safe_corr(df_run, 'DLS', 'ROP')}.",
                            "Color=RPM, tamaño=WOB, puntos=intervalos.",
                        )
                    save_and_show_plotly(prs, f"{run.name} – DLS vs ROP", fig_cross, False)
                vspace(18)

                # Scatter matrix
                st.subheader(f"{run.name} – Scatter matrix")
                vspace(6)
                fig_matrix = build_scatter_matrix(df_run, run.name)
                if fig_matrix is None:
                    st.info("No hay datos suficientes para scatter matrix.")
                else:
                    if show_plots:
                        st.plotly_chart(
                            prettify(fig_matrix, h=520),
                            use_container_width=True,
                            config=PLOTLY_CONFIG,
                        )
                        cols = ["ROP", "WOB", "RPM", "DLS"]
                        cols = [c for c in cols if c in df_run.columns]
                        corr_df = df_run[cols].corr()
                        corr_abs = corr_df.abs().copy()
                        np.fill_diagonal(corr_abs.values, 0)
                        if (corr_abs.values == 0).all():
                            summary = "Sin correlaciones fuertes visibles."
                        else:
                            max_idx = np.unravel_index(
                                np.argmax(corr_abs.values), corr_abs.shape
                            )
                            var_a = corr_abs.index[max_idx[0]]
                            var_b = corr_abs.columns[max_idx[1]]
                            summary = (
                                f"Mayor relación entre {var_a} y {var_b}: "
                                f"{format_num(corr_df.loc[var_a, var_b])}."
                            )
                        chart_notes(
                            summary,
                            "Matriz de dispersión entre variables clave.",
                        )
                    save_and_show_plotly(prs, f"{run.name} – Scatter Matrix", fig_matrix, False)
                vspace(18)

                # Metros acumulados vs tiempo
                st.subheader(f"{run.name} – Metros acumulados vs tiempo")
                vspace(6)
                fig_cum = build_cumulative_meters(df_run, run.name)
                if fig_cum is None:
                    st.info("No hay datos suficientes para metros acumulados.")
                else:
                    if show_plots:
                        st.plotly_chart(
                            prettify(fig_cum),
                            use_container_width=True,
                            config=PLOTLY_CONFIG,
                        )
                        df_cum = df_run.dropna(subset=["End", "Distance"]).copy()
                        df_cum = df_cum.sort_values("End")
                        total_m = df_cum["Distance"].fillna(0).sum()
                        t0 = df_cum["End"].min()
                        t1 = df_cum["End"].max()
                        chart_notes(
                            f"Total {format_num(total_m)} m entre {t0} y {t1}.",
                            "Curva acumulada de metros perforados en el tiempo.",
                        )
                    save_and_show_plotly(prs, f"{run.name} – Metros acumulados", fig_cum, False)
                vspace(18)

                # Real vs Programado (si hay columnas)
                add_real_vs_planned_section(prs, df_run, run.name, show_plots)
                vspace(12)

                # Distributions Pair
                st.subheader(f"{run.name} – Distributions")
                _render_chips_row(
                    [(run.name, "gray"), (f"WOB: {series_summary(df_run['WOB'].dropna())}", "blue"), (f"RPM: {series_summary(df_run['RPM'].dropna())}", "blue")],
                )
                vspace(6)
                c1, c2 = st.columns(2, gap="large", vertical_alignment="top")

                fig_wob_hist = px.histogram(
                    df_run,
                    x="WOB",
                    nbins=20,
                    title=f"{run.name} – WOB Distribution",
                    labels={"WOB": "WOB (kgf)"},
                )
                fig_rpm_hist = build_hist_with_trend(
                    df_run["RPM"],
                    title=f"{run.name} – RPM Distribution",
                    x_label="RPM",
                    nbins=20,
                )

                if show_plots:
                    with c1:
                        st.plotly_chart(
                            prettify_hist(fig_wob_hist),
                            use_container_width=True,
                            config=PLOTLY_CONFIG,
                        )
                        chart_notes(
                            f"{series_summary(df_run['WOB'].dropna())}.",
                            "Histograma de WOB (kgf).",
                        )
                    with c2:
                        st.plotly_chart(
                            prettify_hist(fig_rpm_hist),
                            use_container_width=True,
                            config=PLOTLY_CONFIG,
                        )
                        chart_notes(
                            f"{series_summary(df_run['RPM'].dropna())}.",
                            "Histograma de RPM.",
                        )

                save_and_show_plotly(
                    prs, f"{run.name} – WOB Distribution", fig_wob_hist, False, True
                )
                save_and_show_plotly(
                    prs, f"{run.name} – RPM Distribution", fig_rpm_hist, False, True
                )
                vspace(18)

                # Trends Pair
                st.subheader(f"{run.name} – Trends vs Depth")
                vspace(6)
                df_run_sorted = df_run.sort_values("Depth_X")
                wob_roll = df_run_sorted["WOB"].rolling(window=10, min_periods=1).mean()
                rpm_roll = df_run_sorted["RPM"].rolling(window=10, min_periods=1).mean()

                df_trend = pd.DataFrame(
                    {
                        "Depth": df_run_sorted["Depth_X"],
                        "WOB_roll": wob_roll,
                        "RPM_roll": rpm_roll,
                    }
                )

                fig_wob_trend = px.line(
                    df_trend,
                    x="Depth",
                    y="WOB_roll",
                    title=f"{run.name} – WOB Trend (rolling)",
                    labels={"Depth": "Depth (m)", "WOB_roll": "WOB (kgf)"},
                ).update_layout(xaxis_autorange="reversed")

                fig_rpm_trend = px.line(
                    df_trend,
                    x="Depth",
                    y="RPM_roll",
                    title=f"{run.name} – RPM Trend (rolling)",
                    labels={"Depth": "Depth (m)", "RPM_roll": "RPM"},
                ).update_layout(xaxis_autorange="reversed")

                if show_plots:
                    c1, c2 = st.columns(2, gap="large", vertical_alignment="top")
                    with c1:
                        st.plotly_chart(
                            prettify(fig_wob_trend),
                            use_container_width=True,
                            config=PLOTLY_CONFIG,
                        )
                        wob_delta = df_trend["WOB_roll"].iloc[-1] - df_trend["WOB_roll"].iloc[0]
                        chart_notes(
                            f"Cambio neto WOB: {format_num(wob_delta)}.",
                            "Línea = WOB promedio móvil vs profundidad.",
                        )
                    with c2:
                        st.plotly_chart(
                            prettify(fig_rpm_trend),
                            use_container_width=True,
                            config=PLOTLY_CONFIG,
                        )
                        rpm_delta = df_trend["RPM_roll"].iloc[-1] - df_trend["RPM_roll"].iloc[0]
                        chart_notes(
                            f"Cambio neto RPM: {format_num(rpm_delta)}.",
                            "Línea = RPM promedio móvil vs profundidad.",
                        )

                save_and_show_plotly(prs, f"{run.name} – WOB Trend", fig_wob_trend, False)
                save_and_show_plotly(prs, f"{run.name} – RPM Trend", fig_rpm_trend, False)
                vspace(22)

                summary_stats.append(
                    [
                        run.name,
                        stats["wob_avg"],
                        stats["wob_min"],
                        stats["wob_max"],
                        stats["rpm_avg"],
                        stats["rpm_min"],
                        stats["rpm_max"],
                        stats["avg_rop_rotary"],
                        stats["avg_rop_slide"],
                        stats["rotary_pct"],
                        stats["slide_pct"],
                    ]
                )

            # Summary Global
            summary_df = pd.DataFrame(
                summary_stats,
                columns=[
                    "Run",
                    "Avg WOB",
                    "Min WOB",
                    "Max WOB",
                    "Avg RPM",
                    "Min RPM",
                    "Max RPM",
                    "Avg ROP Rotary",
                    "Avg ROP Slide",
                    "Rotary %",
                    "Slide %",
                ],
            )

            st.subheader("Summary Across All Runs (App)")
            if not summary_df.empty:
                st.dataframe(summary_df, use_container_width=True, hide_index=True)
            else:
                st.info("No hay datos de summary para mostrar.")

            add_table_slide(
                prs,
                "Summary Across All Runs",
                summary_df.values.tolist(),
                list(summary_df.columns),
            )

            if not summary_df.empty:
                fig_wob_sum = px.line(
                    summary_df,
                    x="Run",
                    y=["Avg WOB", "Min WOB", "Max WOB"],
                    markers=True,
                    title="WOB Summary per Run (Avg/Min/Max)",
                    labels={"value": "WOB (kgf)", "variable": "Metric"},
                )
                save_and_show_plotly(
                    prs, "WOB Summary per Run (Avg/Min/Max)", fig_wob_sum, show_plots
                )
                if show_plots:
                    chart_notes(
                        "Comparación de promedio, mínimo y máximo de WOB por corrida.",
                        "Líneas por métrica a través de las corridas.",
                    )

                fig_rpm_sum = px.line(
                    summary_df,
                    x="Run",
                    y=["Avg RPM", "Min RPM", "Max RPM"],
                    markers=True,
                    title="RPM Summary per Run (Avg/Min/Max)",
                    labels={"value": "RPM", "variable": "Metric"},
                )
                save_and_show_plotly(
                    prs, "RPM Summary per Run (Avg/Min/Max)", fig_rpm_sum, show_plots
                )
                if show_plots:
                    chart_notes(
                        "Comparación de promedio, mínimo y máximo de RPM por corrida.",
                        "Líneas por métrica a través de las corridas.",
                    )

                rop_melt = summary_df.melt(
                    id_vars=["Run"],
                    value_vars=["Avg ROP Rotary", "Avg ROP Slide"],
                    var_name="Mode",
                    value_name="ROP (m/hr)",
                )
                fig_rop_sum = px.bar(
                    rop_melt,
                    x="Run",
                    y="ROP (m/hr)",
                    color="Mode",
                    barmode="group",
                    title="Average ROP by Run – Rotary vs Slide",
                    text="ROP (m/hr)",
                    color_discrete_map={
                        "Avg ROP Rotary": "#2563EB",
                        "Avg ROP Slide": "#F59E0B",
                    },
                )
                fig_rop_sum.update_traces(textposition="outside", marker_line_width=0)
                save_and_show_plotly(
                    prs, "Average ROP by Run – Rotary vs Slide", fig_rop_sum, show_plots
                )
                if show_plots:
                    chart_notes(
                        "ROP promedio comparado entre Rotary y Slide por corrida.",
                        "Barras agrupadas por modo de perforación.",
                    )
                vspace(10)

                # Global Distributions
                st.subheader("Global Distributions")
                vspace(6)
                c1, c2 = st.columns(2, gap="large", vertical_alignment="top")

                fig_wob_global = px.histogram(
                    x=all_wob,
                    nbins=25,
                    title="Global WOB Distribution",
                    labels={"x": "WOB (kgf)"},
                )
                fig_rpm_global = build_hist_with_trend(
                    all_rpm,
                    title="Global RPM Distribution",
                    x_label="RPM",
                    nbins=25,
                )

                if show_plots:
                    with c1:
                        st.plotly_chart(
                            prettify_hist(fig_wob_global),
                            use_container_width=True,
                            config=PLOTLY_CONFIG,
                        )
                        chart_notes(
                            f"{series_summary(pd.Series(all_wob))}.",
                            "Histograma global de WOB.",
                        )
                    with c2:
                        st.plotly_chart(
                            prettify_hist(fig_rpm_global),
                            use_container_width=True,
                            config=PLOTLY_CONFIG,
                        )
                        chart_notes(
                            f"{series_summary(pd.Series(all_rpm))}.",
                            "Histograma global de RPM.",
                        )

                save_and_show_plotly(prs, "Global WOB Distribution", fig_wob_global, False, True)
                save_and_show_plotly(prs, "Global RPM Distribution", fig_rpm_global, False, True)

            global_wa = build_kpi_whatsapp_global_es(
                summary_df if not summary_df.empty else pd.DataFrame(),
                df,
                whatsapp_run_texts,
            )
            try:
                save_whatsapp_summaries_excel(whatsapp_excel_path, whatsapp_run_texts, global_wa)
            except Exception:
                pass
            st.session_state["kpi_whatsapp_excel_path"] = str(whatsapp_excel_path)
            st.session_state["kpi_whatsapp_run_texts"] = whatsapp_run_texts
            st.session_state["kpi_whatsapp_global_text"] = global_wa

            for _name, _txt in whatsapp_run_texts:
                add_text_slide(prs, f"{_name} – Resumen WhatsApp (ES)", _txt)
            add_text_slide(prs, "Global – Resumen WhatsApp (ES)", global_wa)

            prs.save(str(pptx_path))
            pdf_path = export_pptx_to_pdf(pptx_path)
            st.session_state["pdf_path"] = str(pdf_path) if pdf_path else None
            copy_report_to_downloads(
                pptx_path,
                pdf_path,
                "Final_Drilling_KPI_Report",
            )
            try:
                if whatsapp_excel_path.exists():
                    _copy_file_to_downloads(
                        whatsapp_excel_path,
                        Path.home() / "Downloads" / "Drilling_KPI_WhatsApp_Summary.xlsx",
                    )
            except Exception:
                pass

        st.session_state["pptx_path"] = str(pptx_path)
        st.session_state["tmp_dir"] = str(tmp_dir)
        st.session_state["report_ready"] = True
        st.session_state["captures_done"] = False

        st.success(
            "✅ PPTX generado (sin capturas). Incluye diapositivas de **resumen WhatsApp (español)** al final. "
            "Excel de textos: `Drilling_KPI_WhatsApp_Summary.xlsx`. Si activaste capturas, presiona el botón para iniciarlas."
        )
        _wa_parts = []
        for _wn, _wt in st.session_state.get("kpi_whatsapp_run_texts") or []:
            _wa_parts.append(f"=== {_wn} ===\n\n{_wt}")
        _gw = st.session_state.get("kpi_whatsapp_global_text") or ""
        _wa_all = "\n\n\n".join(_wa_parts) + (
            "\n\n\n=== RESUMEN GLOBAL ===\n\n" + _gw if _gw else ""
        )
        st.session_state["kpi_whatsapp_all_clipboard"] = _wa_all
        with st.expander("📋 Resumen tipo WhatsApp (español) – copiar y pegar", expanded=False):
            st.caption(
                "Mismo texto que en las últimas diapositivas del PPTX/PDF y en el Excel "
                "(hojas By_run, Global, All_combined). WOB en **tonf** = datos asumidos en **kgf** ÷ 1000."
            )
            st.text_area(
                "Texto completo (todas las corridas + global)",
                value=_wa_all,
                height=420,
                key="kpi_whatsapp_textarea_all",
            )

    # =========================
    # Post-proceso: capturas
    # =========================
    if st.session_state.get("report_ready"):
        pptx_path = Path(st.session_state["pptx_path"])
        tmp_dir = Path(st.session_state["tmp_dir"])

        st.download_button(
            "Download PPTX",
            data=pptx_path.read_bytes(),
            file_name="Final_Drilling_KPI_Report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        _wa_xlsx = st.session_state.get("kpi_whatsapp_excel_path")
        if _wa_xlsx and Path(_wa_xlsx).exists():
            st.download_button(
                "Descargar resumen WhatsApp (Excel)",
                data=Path(_wa_xlsx).read_bytes(),
                file_name="Drilling_KPI_WhatsApp_Summary.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="kpi_dl_whatsapp_xlsx",
            )
        pdf_path = st.session_state.get("pdf_path")
        if pdf_path:
            pdf_path_obj = Path(pdf_path)
            if pdf_path_obj.exists():
                st.download_button(
                    "Download PDF",
                    data=pdf_path_obj.read_bytes(),
                    file_name="Final_Drilling_KPI_Report.pdf",
                    mime="application/pdf",
                )
        st.caption(f"PPTX temporal en: {pptx_path}")

        if enable_region_captures and not st.session_state.get("captures_done", False):
            st.warning("📸 Las capturas se añadirán al FINAL del PPTX.")
            st.info("Acomoda la región en pantalla y presiona el botón cuando estés listo.")

            start_caps = st.button("📸 Iniciar capturas de pantalla ahora")

            if start_caps:
                countdown_placeholder = st.empty()
                for t in range(5, 0, -1):
                    countdown_placeholder.warning(f"Tomando capturas en {t} segundos…")
                    time.sleep(1)
                countdown_placeholder.empty()

                prs2 = Presentation(str(pptx_path))

                capture_region_screenshots(
                    region=(int(region_x), int(region_y), int(region_w), int(region_h)),
                    n_shots=int(n_captures),
                    interval_s=int(interval_seconds),
                    title_prefix=capture_title_prefix,
                    prs=prs2,
                    show_plots=show_plots,
                    download_dir=str(tmp_dir),
                )

                prs2.save(str(pptx_path))
                st.session_state["captures_done"] = True
                pdf_path = export_pptx_to_pdf(pptx_path)
                st.session_state["pdf_path"] = str(pdf_path) if pdf_path else None
                copy_report_to_downloads(
                    pptx_path,
                    pdf_path,
                    "Final_Drilling_KPI_Report_with_captures",
                )

                st.success("✅ Capturas agregadas al PPTX.")
                st.download_button(
                    "Download PPTX (con capturas)",
                    data=pptx_path.read_bytes(),
                    file_name="Final_Drilling_KPI_Report_with_captures.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                if pdf_path is not None and Path(pdf_path).exists():
                    st.download_button(
                        "Download PDF (con capturas)",
                        data=Path(pdf_path).read_bytes(),
                        file_name="Final_Drilling_KPI_Report_with_captures.pdf",
                        mime="application/pdf",
                    )


# =========================
# Mud Report – bitácora de propiedades de fluidos
# =========================
# Aliases por propiedad canónica (nombre en reporte -> clave bitácora)
MUD_PROPERTY_ALIASES = {
    "Density": ["density", "densidad", "mw", "mw (g/l)", "density @ c", "density @", "mud weight", "peso lodo", "densidad sp.gr"],
    "Marsh": ["marsh", "visc. marsh", "viscosidad marsh"],
    "Temperature": ["temperatura salida", "temp. de salida", "temperatura", "temp. de analisis", "temp. de análisis", "temp de analisis"],
    "VA": ["va", "visc.aparente", "visc. aparente", "viscosidad aparente"],
    "FV": ["fv", "fv @ c", "fv @ °c", "funnel viscosity", "viscosidad embudo"],
    "PV": ["pv", "pv (cp)", "pv @ c", "pv @ °c", "plastic viscosity", "viscosidad plástica", "viscoplastic", "visc. plastica", "visc.plastica"],
    "YP": ["yp", "yv", "yield point", "punto de cedencia", "yp (lb/100ft2)", "lb/100ft²", "pc"],
    "Gel_10s": ["gel 10s", "gel 10s/10m/30m", "gels 10s", "10s", "gel (10s)", "gel 10s/10m", "geles"],
    "Gel_10min": ["gel 10m", "gel 10min", "10min", "10m"],
    "Gel_30min": ["gel 30m", "gel 30min", "30min", "30m"],
    "L600": ["lectura 600", "l600"],
    "L300": ["lectura 300", "l300"],
    "L200": ["lectura 200", "l200"],
    "L100": ["lectura 100", "l100"],
    "L6": ["lectura 6", "l6"],
    "L3": ["lectura 3", "l3"],
    "Filtrado": ["filtrado", "filtrate", "fl temp", "hthp", "api filtrate", "fluid loss", "cake (hthp)", "filtrado hpht", "filtrado apat"],
    "Enjarre": ["enjarre", "cake"],
    "LGS": ["lgs", "lgs/hgs", "low gravity solids", "lgs (%)"],
    "HGS": ["hgs", "high gravity solids", "hgs (%)"],
    "Chlorides": ["chlorides", "cloruros", "chlorides (ppm)", "chlorides / calcium"],
    "Solids": ["solids", "corr solid", "solids content %", "sand %", "% sólidos", "% solidos", "sólidos no corregidos", "solidos no corregidos", "solidos corregidos"],
    "Oil": ["% aceite", "aceite %vol", "%oil", "oil"],
    "Water": ["% agua", "agua no correg", "%water", "water"],
    "RAA": ["raa", "r. aceite / agua", "rel. aceite/agua", "aceite/agua"],
    "AgNO3": ["agno3"],
    "Salinity": ["salinidad"],
    "Electrical_Stability": ["est. electrica", "estabilidad", "est. elect", "elec. stability"],
    "Alkalinity": ["alcalinidad"],
    "Excess_Cal": ["exceso de cal", "exc.cal", "exc cal"],
}
MUD_CANONICAL_ORDER = [
    "Date", "DateTime", "Properties", "Depth (MD)", "Depth (TVD)", "Fluid set", "Source", "Time", "FL Temp",
    "Density @ °C", "FV", "FV Temp", "FV @ °C", "PV", "PV Temp", "PV @ °C", "YP",
    "Gel_10s", "Gel_10min", "Gel_30min", "tau0",
    "L600", "L300", "L200", "L100", "L6", "L3",
    "HTHP", "HTHP @ °C", "Corr Solid", "NAP", "Water", "NAP Ratio", "Water Ratio",
    "Sand", "Cake (HTHP)", "Chlorides", "Calcium", "CaCl2", "Water Phase Salinity",
    "NaCL (Sol/Insol)", "Excess Lime", "Electrical_Stability",
    "LGS (%)", "HGS (%)", "LGS (kg/m³)", "HGS (kg/m³)", "ASG",
    "Additional Properties", "n (HB)", "K (HB)", "Viscometer Sag Shoe Test", "(VSST)",
    "Marsh", "Temperature", "VA", "Filtrado", "Enjarre", "LGS", "HGS", "Solids", "Oil", "RAA",
    "AgNO3", "Salinity", "Alkalinity", "Excess_Cal",
]
MUD_METADATA_COLUMNS = {
    "Date", "DateTime", "Properties", "Depth (MD)", "Depth (TVD)", "Fluid set", "Source", "Time",
    "Additional Properties",
}
MUD_ANALYTIC_EXCLUDE = {
    "DateTime", "Properties", "Depth (MD)", "Depth (TVD)", "Fluid set", "Source", "Time",
    "Density @ °C", "FV @ °C", "PV @ °C",
}
MUD_EXPORT_HEADER_SPECS = [
    ("Depth (MD)", "Depth (MD)", "m"),
    ("Depth (TVD)", "Depth (TVD)", "m"),
    ("Properties", "Properties", "N°"),
    ("Fluid set", "Fluid set", "Fluid"),
    ("Source", "Source", "Source"),
    ("Time", "Time", "time"),
    ("DateTime", "DateTime", "YYYY-MM-DDTHH:MM:SS"),
    ("FL Temp", "FL Temp", "°C"),
    ("Density @ °C", "D @ °C", "kg/m³"),
    ("FV @ °C", "Fv @ °C", "s/qt"),
    ("PV @ °C", "PV @ °C", "cP"),
    ("YP", "YP", "lb/100ft²"),
    ("Gel_10s", "GELS 10s", "lb/100ft²"),
    ("Gel_10min", "GELS 10min", "lb/100ft²"),
    ("Gel_30min", "GELS 30min", "lb/100ft²"),
    ("tau0", "tau0", "lb/100ft²"),
    ("L600", "600", "600"),
    ("L300", "300", "300"),
    ("L200", "200", "200"),
    ("L100", "100", "100"),
    ("L6", "6", "6"),
    ("L3", "3", "3"),
    ("HTHP", "HTHP", "HTHP"),
    ("HTHP @ °C", "°C", "°C"),
    ("Corr Solid", "Corr Solid", "%"),
    ("NAP", "NAP", "%"),
    ("Water", "Water", "%"),
    ("NAP Ratio", "NAP", "%"),
    ("Water Ratio", "Water Ratio", "%"),
    ("Sand", "Sand", "%"),
    ("Cake (HTHP)", "Cake (HTHP)", "32nd"),
    ("Chlorides", "Chlorides", "mg/L"),
    ("Calcium", "Calcium", "mg/L"),
    ("CaCl2", "CaCl2", "mg/L"),
    ("Water Phase Salinity", "Water Phase Salinity", "ppm"),
    ("NaCL (Sol/Insol)", "NaCL (Sol/Insol)", "kg/m³"),
    ("Excess Lime", "Excess Lime", "kg/m³"),
    ("Electrical_Stability", "Elec. Stability", "V"),
    ("LGS (%)", "LGS", "%"),
    ("HGS (%)", "HGS", "%"),
    ("LGS (kg/m³)", "LGS", "kg/m³"),
    ("HGS (kg/m³)", "HGS", "kg/m³"),
    ("ASG", "ASG", "SG"),
    ("Additional Properties", "Additional Properties", "Properties"),
    ("n (HB)", "n (HB)", "dec"),
    ("K (HB)", "K (HB)", "lb*s^n'/100ft2"),
    ("Viscometer Sag Shoe Test", "Viscometer Sag Shoe Test", "lbm/gal"),
    ("(VSST)", "(VSST)", ""),
]


def _normalize_mud_property_name(label: str) -> str | None:
    """Mapea etiqueta de reporte a nombre canónico."""
    if not label or not isinstance(label, str):
        return None
    key = str(label).strip().lower()
    key = re.sub(r"\s+", " ", key)
    for canonical, aliases in MUD_PROPERTY_ALIASES.items():
        for a in aliases:
            if a in key or key in a:
                return canonical
    if "gel" in key and "10s" in key:
        return "Gel_10s"
    if "gel" in key and "10" in key and "30" not in key:
        return "Gel_10min"
    if "gel" in key and "30" in key:
        return "Gel_30min"
    return None


def _extract_numeric(val) -> float | None:
    """Extrae un número de una celda o string libre."""
    if val is None or (isinstance(val, float) and np.isnan(val)):
        return None
    if isinstance(val, (int, float)):
        return float(val)
    s = str(val).strip()
    if not s:
        return None
    s = s.replace(" ", " ")
    # quitar miles y normalizar decimales
    s = re.sub(r"(?<=\d),(?=\d{3}(?:\D|$))", "", s)
    s = s.replace(",", ".")
    m = re.search(r"[-+]?\d*\.?\d+", s)
    if m:
        try:
            return float(m.group(0))
        except ValueError:
            return None
    return None


def _extract_all_numbers(val) -> list[float]:
    if val is None or (isinstance(val, float) and np.isnan(val)):
        return []
    s = str(val).strip().replace(" ", " ")
    if not s:
        return []
    s = re.sub(r"(?<=\d),(?=\d{3}(?:\D|$))", "", s)
    s = s.replace(",", ".")
    nums = []
    for tok in re.findall(r"[-+]?\d*\.?\d+", s):
        try:
            nums.append(float(tok))
        except ValueError:
            pass
    return nums


def _parse_gel_triple(val) -> tuple[float | None, float | None, float | None]:
    """Parsea '10/15/17' -> (10, 15, 17) y variantes '8/16/22'."""
    nums = _extract_all_numbers(val)
    if len(nums) >= 3:
        return nums[0], nums[1], nums[2]
    if len(nums) == 2:
        return nums[0], nums[1], None
    if len(nums) == 1:
        return nums[0], None, None
    return None, None, None


def _extract_date_from_text(text: str) -> pd.Timestamp | None:
    if not text:
        return None
    month_map = {
        "ene": 1, "feb": 2, "mar": 3, "abr": 4, "may": 5, "jun": 6,
        "jul": 7, "ago": 8, "sep": 9, "oct": 10, "nov": 11, "dic": 12,
    }
    # yyyy-mm-dd / yyyy.mm.dd / yyyy/mm/dd
    m = re.search(r"((?:19|20)\d{2})[./-](\d{1,2})[./-](\d{1,2})", text)
    if m:
        y, mo, d = int(m.group(1)), int(m.group(2)), int(m.group(3))
        try:
            return pd.Timestamp(y, mo, d).normalize()
        except ValueError:
            pass
    # dd-mm-yyyy / dd/mm/yyyy
    m = re.search(r"(\d{1,2})[/-](\d{1,2})[/-](\d{2,4})", text)
    if m:
        d, mo, y = int(m.group(1)), int(m.group(2)), int(m.group(3))
        if y < 100:
            y += 2000
        try:
            return pd.Timestamp(y, mo, d).normalize()
        except ValueError:
            pass
    m = re.search(r"(\d{1,2})-([A-Za-z]{3})-(\d{2,4})", text)
    if m:
        d = int(m.group(1))
        mo = month_map.get(m.group(2).lower()[:3])
        y = int(m.group(3))
        if y < 100:
            y += 2000
        if mo:
            try:
                return pd.Timestamp(y, mo, d).normalize()
            except ValueError:
                pass
    return None


def _date_from_filename_or_today(name: str) -> pd.Timestamp:
    """Extrae fecha de nombre de archivo o usa hoy."""
    if not name:
        return pd.Timestamp.now().normalize()
    d = _extract_date_from_text(name)
    return d if d is not None else pd.Timestamp.now().normalize()


def _mud_num_to_text(val) -> str:
    num = _extract_numeric(val)
    if num is None:
        s = str(val).strip() if val is not None else ""
        return s
    if abs(num - round(num)) < 1e-9:
        return str(int(round(num)))
    return f"{num:.12g}"


def _mud_clean_cell_text(val) -> str:
    if val is None or (isinstance(val, float) and np.isnan(val)):
        return ""
    if isinstance(val, pd.Timestamp):
        return val.strftime("%Y-%m-%d %H:%M:%S")
    try:
        import datetime as _dt
        if isinstance(val, _dt.time):
            return val.strftime("%H:%M")
    except Exception:
        pass
    return str(val).strip().replace("\u00a0", " ")


def _mud_parse_time_value(val):
    s = _mud_clean_cell_text(val)
    if not s:
        return None
    try:
        import datetime as _dt
        if hasattr(val, "hour") and hasattr(val, "minute") and not isinstance(val, pd.Timestamp):
            return _dt.time(val.hour, val.minute, getattr(val, "second", 0))
    except Exception:
        pass
    for fmt in ("%H:%M:%S", "%H:%M"):
        try:
            return datetime.strptime(s, fmt).time()
        except Exception:
            pass
    dt = pd.to_datetime(s, errors="coerce")
    if pd.notna(dt):
        return dt.time()
    return None


def _mud_compose_datetime(date_value, time_value) -> pd.Timestamp | None:
    base_date = pd.to_datetime(date_value, errors="coerce")
    if pd.isna(base_date):
        return None
    t = _mud_parse_time_value(time_value)
    if t is None:
        return base_date
    return pd.Timestamp.combine(base_date.normalize().date(), t)


def _mud_isoformat_no_tz(ts) -> str:
    ts = pd.to_datetime(ts, errors="coerce")
    if pd.isna(ts):
        return ""
    return ts.strftime("%Y-%m-%dT%H:%M:%S")


def _mud_pair_string(raw_value) -> str:
    nums = _extract_all_numbers(raw_value)
    if len(nums) >= 2:
        return f"{_mud_num_to_text(nums[0])} @ {_mud_num_to_text(nums[1])}"
    if len(nums) == 1:
        return _mud_num_to_text(nums[0])
    return _mud_clean_cell_text(raw_value)


def _mud_apply_daily_property(row_record: dict, label: str, unit: str, raw_value) -> None:
    low = re.sub(r"\s+", " ", _mud_clean_cell_text(label).lower())
    unit_low = re.sub(r"\s+", " ", _mud_clean_cell_text(unit).lower())
    nums = _extract_all_numbers(raw_value)
    if not low:
        return
    if low.startswith("depth"):
        if len(nums) >= 1:
            row_record["Depth (MD)"] = nums[0]
        if len(nums) >= 2:
            row_record["Depth (TVD)"] = nums[1]
        return
    if low.startswith("fl temp"):
        if nums:
            row_record["FL Temp"] = nums[0]
        return
    if low.startswith("density"):
        row_record["Density @ °C"] = _mud_pair_string(raw_value)
        if nums:
            row_record["Density"] = nums[0]
        if len(nums) >= 2:
            row_record["Density Temp"] = nums[1]
        return
    if low.startswith("fv @"):
        if nums:
            row_record["FV"] = nums[0]
        if len(nums) >= 2:
            row_record["FV Temp"] = nums[1]
        row_record["FV @ °C"] = _mud_pair_string(raw_value)
        return
    if low.startswith("pv @"):
        if nums:
            row_record["PV"] = nums[0]
        if len(nums) >= 2:
            row_record["PV Temp"] = nums[1]
        row_record["PV @ °C"] = _mud_pair_string(raw_value)
        return
    if low == "yp" or low.startswith("yp "):
        if nums:
            row_record["YP"] = nums[0]
        return
    if low.startswith("gels"):
        g1, g2, g3 = _parse_gel_triple(raw_value)
        if g1 is not None:
            row_record["Gel_10s"] = g1
        if g2 is not None:
            row_record["Gel_10min"] = g2
        if g3 is not None:
            row_record["Gel_30min"] = g3
        return
    if low == "tau0":
        if nums:
            row_record["tau0"] = nums[0]
        return
    if low.startswith("600/300"):
        if len(nums) >= 1:
            row_record["L600"] = nums[0]
        if len(nums) >= 2:
            row_record["L300"] = nums[1]
        return
    if low.startswith("200/100"):
        if len(nums) >= 1:
            row_record["L200"] = nums[0]
        if len(nums) >= 2:
            row_record["L100"] = nums[1]
        return
    if low.startswith("6/3"):
        if len(nums) >= 1:
            row_record["L6"] = nums[0]
        if len(nums) >= 2:
            row_record["L3"] = nums[1]
        return
    if low.startswith("hthp"):
        if nums:
            row_record["HTHP"] = nums[0]
        if len(nums) >= 2:
            row_record["HTHP @ °C"] = nums[1]
        return
    if low.startswith("corr solid"):
        if nums:
            row_record["Corr Solid"] = nums[0]
        return
    if low.startswith("nap / water ratio"):
        if len(nums) >= 1:
            row_record["NAP Ratio"] = nums[0]
        if len(nums) >= 2:
            row_record["Water Ratio"] = nums[1]
        return
    if low.startswith("nap / water"):
        if len(nums) >= 1:
            row_record["NAP"] = nums[0]
        if len(nums) >= 2:
            row_record["Water"] = nums[1]
        return
    if low == "sand" or low.startswith("sand "):
        if nums:
            row_record["Sand"] = nums[0]
        return
    if low.startswith("cake"):
        if nums:
            row_record["Cake (HTHP)"] = nums[0]
        return
    if low.startswith("chlorides / calcium"):
        if len(nums) >= 1:
            row_record["Chlorides"] = nums[0]
        if len(nums) >= 2:
            row_record["Calcium"] = nums[1]
        return
    if low == "cacl2" or low.startswith("cacl2 "):
        if nums:
            row_record["CaCl2"] = nums[0]
        return
    if low.startswith("water phase salinity"):
        if nums:
            row_record["Water Phase Salinity"] = nums[0]
        return
    if low.startswith("nacl"):
        txt = _mud_clean_cell_text(raw_value)
        if txt and txt != "/":
            row_record["NaCL (Sol/Insol)"] = txt
        return
    if low.startswith("excess lime"):
        if nums:
            row_record["Excess Lime"] = nums[0]
        return
    if low.startswith("elec. stability"):
        if nums:
            row_record["Electrical_Stability"] = nums[0]
        return
    if low.startswith("lgs / hgs"):
        if len(nums) >= 1:
            target_lgs = "LGS (kg/m³)" if "kg/" in unit_low else "LGS (%)"
            row_record[target_lgs] = nums[0]
        if len(nums) >= 2:
            target_hgs = "HGS (kg/m³)" if "kg/" in unit_low else "HGS (%)"
            row_record[target_hgs] = nums[1]
        return
    if low == "asg":
        if nums:
            row_record["ASG"] = nums[0]
        return
    if low.startswith("n (hb)"):
        if nums:
            row_record["n (HB)"] = nums[0]
        return
    if low.startswith("k (hb)"):
        if nums:
            row_record["K (HB)"] = nums[0]
        return
    if low.startswith("viscometer sag shoe test"):
        if nums:
            row_record["Viscometer Sag Shoe Test"] = nums[0]
        return
    if low.startswith("(vsst)"):
        if nums:
            row_record["(VSST)"] = nums[0]
        else:
            txt = _mud_clean_cell_text(raw_value)
            if txt:
                row_record["(VSST)"] = txt
        return
    canonical = _normalize_mud_property_name(label)
    if canonical:
        _mud_apply_canonical_value(row_record, canonical, raw_value)


def _parse_mud_daily_report_sheet(df_raw: pd.DataFrame, source_name: str = "") -> list[dict]:
    df = df_raw.copy()
    if df.empty or df.shape[0] < 6 or df.shape[1] < 5:
        return []
    cell_a1 = _mud_clean_cell_text(df.iat[0, 0]) if df.shape[0] > 0 else ""
    cell_a2 = _mud_clean_cell_text(df.iat[1, 0]) if df.shape[0] > 1 else ""
    cell_a5 = _mud_clean_cell_text(df.iat[4, 0]) if df.shape[0] > 4 else ""
    if "daily fluid properties" not in cell_a1.lower() or "properties" not in cell_a2.lower() or "time" not in cell_a5.lower():
        return []

    report_date = _extract_date_from_text(cell_a1) or _date_from_filename_or_today(source_name)
    sample_cols = []
    for j in range(2, df.shape[1]):
        prop_txt = _mud_clean_cell_text(df.iat[1, j]) if df.shape[0] > 1 else ""
        fluid_txt = _mud_clean_cell_text(df.iat[2, j]) if df.shape[0] > 2 else ""
        src_txt = _mud_clean_cell_text(df.iat[3, j]) if df.shape[0] > 3 else ""
        time_txt = _mud_clean_cell_text(df.iat[4, j]) if df.shape[0] > 4 else ""
        if time_txt or fluid_txt or src_txt or (prop_txt and (fluid_txt or src_txt)):
            sample_cols.append(j)
    if not sample_cols:
        return []

    records: list[dict] = []
    for idx, j in enumerate(sample_cols, start=1):
        prop_id = _extract_numeric(df.iat[1, j]) if df.shape[0] > 1 else None
        time_raw = df.iat[4, j] if df.shape[0] > 4 else None
        ts = _mud_compose_datetime(report_date, time_raw)
        rec = {
            "Date": ts if ts is not None else report_date,
            "DateTime": _mud_isoformat_no_tz(ts if ts is not None else report_date),
            "Properties": int(prop_id) if prop_id is not None else idx,
            "Fluid set": _mud_clean_cell_text(df.iat[2, j]) if df.shape[0] > 2 else "",
            "Source": _mud_clean_cell_text(df.iat[3, j]) if df.shape[0] > 3 else source_name,
            "Time": _mud_parse_time_value(time_raw).strftime("%H:%M") if _mud_parse_time_value(time_raw) else _mud_clean_cell_text(time_raw),
            "Additional Properties": int(prop_id) if prop_id is not None else idx,
        }
        records.append(rec)

    row_texts = []
    for i in range(df.shape[0]):
        row_texts.append(" ".join(_mud_clean_cell_text(df.iat[i, c]) for c in range(df.shape[1]) if _mud_clean_cell_text(df.iat[i, c])))

    for i in range(5, df.shape[0]):
        label = _mud_clean_cell_text(df.iat[i, 0])
        if not label:
            continue
        unit = _mud_clean_cell_text(df.iat[i, 1]) if df.shape[1] > 1 else ""
        raw_vals = [df.iat[i, j] for j in sample_cols]
        whole_nums = _extract_all_numbers(row_texts[i])
        use_sequence = len(whole_nums) == len(sample_cols) and any(
            (not _mud_clean_cell_text(v)) or len(_extract_all_numbers(v)) != 1 for v in raw_vals
        )
        for idx, rec in enumerate(records):
            raw_value = whole_nums[idx] if use_sequence else raw_vals[idx]
            if not _mud_clean_cell_text(raw_value) and not isinstance(raw_value, (int, float)):
                continue
            _mud_apply_daily_property(rec, label, unit, raw_value)

    return [r for r in records if any(
        k not in MUD_METADATA_COLUMNS and pd.notna(v) and _mud_clean_cell_text(v) not in ("", "/")
        for k, v in r.items()
    )]


def _mud_apply_canonical_value(row_record: dict, canonical: str, raw_value) -> None:
    if canonical.startswith("Gel"):
        g1, g2, g3 = _parse_gel_triple(raw_value)
        if g1 is not None:
            row_record["Gel_10s"] = g1
        if g2 is not None:
            row_record["Gel_10min"] = g2
        if g3 is not None:
            row_record["Gel_30min"] = g3
        return
    if canonical == "RAA":
        nums = _extract_all_numbers(raw_value)
        if len(nums) >= 1:
            row_record[canonical] = nums[0]
        return
    num = _extract_numeric(raw_value)
    if num is not None:
        row_record[canonical] = num


def _parse_mud_text_block(text: str, row_record: dict) -> None:
    if not text:
        return
    text = text.replace("\u00a0", " ")
    if pd.isna(row_record.get("Date")) or row_record.get("Date") is None:
        d = _extract_date_from_text(text)
        if d is not None:
            row_record["Date"] = d

    patterns = [
        ("Density", [r"densidad[^\n\r:]*[: ]+([0-9.,]+)", r"density[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("Marsh", [r"visc\.? marsh[^\n\r:]*[: ]+([0-9.,]+)", r"viscosidad marsh[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("Temperature", [r"temperatura salida[^\n\r:]*[: ]+([0-9.,]+)", r"temp\. de salida[^\n\r:]*[: ]+([0-9.,]+)", r"temp\. de an[aá]lisis[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("VA", [r"visc\.?aparente(?:\(va\))?[^\n\r:]*[: ]+([0-9.,]+)", r"visc\.? aparente[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("PV", [r"visc\.?plastica(?:\(vp\))?[^\n\r:]*[: ]+([0-9.,]+)", r"visc\.? plastica[^\n\r:]*[: ]+([0-9.,]+)", r"\bPV\b[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("YP", [r"punto cedente(?:\(yp\))?[^\n\r:]*[: ]+([0-9.,]+)", r"\bPC\b[^\n\r:]*[: ]+([0-9.,]+)", r"\bYP\b[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("L600", [r"lectura 600[^\n\r:]*[: ]+([0-9.,]+)", r"l600[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("L300", [r"lectura 300[^\n\r:]*[: ]+([0-9.,]+)", r"l300[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("L200", [r"lectura 200[^\n\r:]*[: ]+([0-9.,]+)", r"l200[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("L100", [r"lectura 100[^\n\r:]*[: ]+([0-9.,]+)", r"l100[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("L6", [r"lectura 6[^\n\r:]*[: ]+([0-9.,]+)", r"l6[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("L3", [r"lectura 3[^\n\r:]*[: ]+([0-9.,]+)", r"l3[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("Filtrado", [r"filtrado hpht[^\n\r:]*[: ]+([0-9.,]+)", r"filtrado apat[^\n\r:]*[: ]+([0-9.,]+)", r"filtrado[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("Enjarre", [r"enjarre[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("Solids", [r"%\s*s[óo]lidos[^\n\r:]*[: ]+([0-9.,]+)", r"s[óo]lidos no corregidos[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("Oil", [r"%\s*aceite[^\n\r:]*[: ]+([0-9.,]+)", r"aceite\s*%vol[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("Water", [r"%\s*agua[^\n\r:]*[: ]+([0-9.,]+)", r"agua no correg[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("RAA", [r"raa[^\n\r:]*[: ]+([0-9.,]+)", r"rel\. aceite/agua[^\n\r:]*[: ]+([0-9.,]+)", r"aceite/agua[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("AgNO3", [r"agno3[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("Chlorides", [r"cloruros[^\n\r:]*[: ]+([0-9.,]+)", r"chlorides[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("Salinity", [r"salinidad[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("Electrical_Stability", [r"est\.? electrica[^\n\r:]*[: ]+([0-9.,]+)", r"estabilidad[^\n\r:]*[: ]+([0-9.,]+)", r"est\.? elect\.?[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("Alkalinity", [r"alcalinidad[^\n\r:]*[: ]+([0-9.,]+)"]),
        ("Excess_Cal", [r"exceso de cal[^\n\r:]*[: ]+([0-9.,]+)", r"exc\.cal[^\n\r:]*[: ]+([0-9.,]+)"]),
    ]
    for canonical, regexes in patterns:
        if canonical in row_record and pd.notna(row_record.get(canonical)):
            continue
        for pat in regexes:
            m = re.search(pat, text, flags=re.IGNORECASE)
            if m:
                _mud_apply_canonical_value(row_record, canonical, m.group(1))
                break

    m = re.search(r"geles?[^\n\r:]*[: ]+([0-9.,/ ]+)", text, flags=re.IGNORECASE)
    if m:
        _mud_apply_canonical_value(row_record, "Gel_10s", m.group(1))
    m = re.search(r"gel\s*10s/10m[^\n\r:]*[: ]+([0-9.,/ ]+)", text, flags=re.IGNORECASE)
    if m:
        _mud_apply_canonical_value(row_record, "Gel_10s", m.group(1))




def _parse_mud_lines(text: str, row_record: dict) -> None:
    if not text:
        return
    for raw_line in text.splitlines():
        line = (raw_line or "").strip()
        if not line:
            continue
        low = re.sub(r"\s+", " ", line.lower())
        nums = _extract_all_numbers(line)
        if not nums:
            continue

        if low.startswith("densidad") or low.startswith("density"):
            row_record["Density @ °C"] = _mud_pair_string(line)
            row_record["Density"] = nums[0] if nums else row_record.get("Density")
            if len(nums) >= 2:
                row_record["Density Temp"] = nums[1]
        elif low.startswith("visc. marsh") or low.startswith("viscosidad marsh"):
            row_record["Marsh"] = nums[-1]
        elif low.startswith("temperatura salida") or low.startswith("temp. de salida"):
            row_record["Temperature"] = nums[-1]
        elif low.startswith("temp. de analisis") or low.startswith("temp. de análisis"):
            row_record["Temperature"] = nums[-1]
        elif low.startswith("visc.aparente") or low.startswith("visc. aparente") or low.startswith("viscosidad aparente"):
            row_record["VA"] = nums[-1]
        elif low.startswith("visc.plastica") or low.startswith("visc. plastica") or low.startswith("pv ") or low == "pv":
            row_record["PV"] = nums[-1]
        elif low.startswith("punto cedente") or low.startswith("pc ") or low == "pc" or low.startswith("yp ") or low == "yp":
            row_record["YP"] = nums[-1]
        elif low.startswith("lectura 600"):
            row_record["L600"] = nums[0]
        elif low.startswith("lectura 300"):
            row_record["L300"] = nums[0]
        elif low.startswith("lectura 200"):
            row_record["L200"] = nums[0]
        elif low.startswith("lectura 100"):
            row_record["L100"] = nums[0]
        elif low.startswith("lectura 6"):
            row_record["L6"] = nums[0]
        elif low.startswith("lectura 3"):
            row_record["L3"] = nums[0]
        elif low.startswith("l600/l300") and len(nums) >= 2:
            row_record["L600"] = nums[0]
            row_record["L300"] = nums[1]
        elif low.startswith("l200/l100") and len(nums) >= 2:
            row_record["L200"] = nums[0]
            row_record["L100"] = nums[1]
        elif low.startswith("l6/l3") and len(nums) >= 2:
            row_record["L6"] = nums[0]
            row_record["L3"] = nums[1]
        elif low.startswith("filtrado hpht") or low.startswith("filtrado apat") or low.startswith("filtrado"):
            row_record["Filtrado"] = nums[-1]
        elif low.startswith("enjarre"):
            row_record["Enjarre"] = nums[-1]
        elif low.startswith("geles") or low.startswith("gel 10s/10m"):
            if len(nums) >= 1:
                row_record["Gel_10s"] = nums[0]
            if len(nums) >= 2:
                row_record["Gel_10min"] = nums[1]
            if len(nums) >= 3:
                row_record["Gel_30min"] = nums[2]
        elif low.startswith("% sólidos") or low.startswith("% solidos") or low.startswith("sólidos no corregidos") or low.startswith("solidos no corregidos"):
            row_record["Solids"] = nums[-1]
        elif low.startswith("% aceite") or low.startswith("aceite %vol"):
            row_record["Oil"] = nums[-1]
        elif low.startswith("% agua") or low.startswith("agua no correg"):
            row_record["Water"] = nums[-1]
        elif low.startswith("raa") or low.startswith("rel. aceite/agua") or low.startswith("aceite/agua"):
            row_record["RAA"] = nums[-2] if len(nums) >= 2 else nums[0]
        elif low.startswith("agno3"):
            row_record["AgNO3"] = nums[-1]
        elif low.startswith("cloruros"):
            row_record["Chlorides"] = nums[-1]
        elif low.startswith("salinidad"):
            row_record["Salinity"] = nums[-1]
        elif low.startswith("est. electrica") or low.startswith("estabilidad") or low.startswith("est. elect"):
            row_record["Electrical_Stability"] = nums[-1]
        elif low.startswith("alcalinidad"):
            row_record["Alkalinity"] = nums[-1]
        elif low.startswith("exceso de cal") or low.startswith("exc.cal"):
            row_record["Excess_Cal"] = nums[-1]
def _parse_mud_excel_sheet(df_raw: pd.DataFrame, source_name: str = "") -> list[dict]:
    """Parsea una hoja Excel de propiedades de lodo (formato filas propiedad/valor o tabla)."""
    daily_rows = _parse_mud_daily_report_sheet(df_raw, source_name)
    if daily_rows:
        return daily_rows

    out: list[dict] = []
    date = _date_from_filename_or_today(source_name)
    row_record: dict = {"Date": date, "Source": source_name}

    df = df_raw.copy()
    df.columns = [str(c).strip() for c in df.columns]
    prop_col = None
    for c in df.columns:
        cl = c.lower()
        if "propert" in cl or "parámetro" in cl or "parameter" in cl or c == "Unnamed: 0" or cl == "0":
            prop_col = c
            break
    if prop_col is None and len(df.columns) >= 1:
        prop_col = df.columns[0]

    if prop_col is not None:
        for _, r in df.iterrows():
            label = r.get(prop_col)
            if pd.isna(label):
                continue
            canonical = _normalize_mud_property_name(str(label))
            if canonical:
                for c in df.columns:
                    if c == prop_col:
                        continue
                    v = r.get(c)
                    if pd.isna(v):
                        continue
                    _mud_apply_canonical_value(row_record, canonical, v)
                    if canonical in row_record or canonical.startswith("Gel"):
                        break

    if not any(k for k in row_record if k not in ("Date", "Source")):
        for col in df.columns:
            canonical = _normalize_mud_property_name(col)
            if canonical:
                vals = df[col].dropna().tolist()
                if vals:
                    _mud_apply_canonical_value(row_record, canonical, vals[0])

    text_blob = "\n".join(
        " ".join(str(v) for v in row.tolist() if pd.notna(v))
        for _, row in df.iterrows()
    )
    _parse_mud_text_block(text_blob, row_record)
    _parse_mud_lines(text_blob, row_record)

    if any(k for k in row_record if k not in ("Date", "Source")):
        out.append(row_record)
    return out


def _parse_mud_csv(df_raw: pd.DataFrame, source_name: str = "") -> list[dict]:
    """Parsea CSV de propiedades de lodo (igual lógica que Excel)."""
    return _parse_mud_excel_sheet(df_raw, source_name)


def _parse_mud_pdf(file, source_name: str = "") -> list[dict]:
    """Extrae tablas/texto de PDF y parsea propiedades conocidas."""
    out: list[dict] = []
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        return out
    name = source_name or getattr(file, "name", "") or ""
    row_record: dict = {"Date": _date_from_filename_or_today(name), "Source": name}
    try:
        with pdfplumber.open(file) as pdf:
            full_text_parts = []
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                if page_text:
                    full_text_parts.append(page_text)
                tables = page.extract_tables() or []
                for table in tables:
                    for row in table or []:
                        if not row:
                            continue
                        for idx, cell in enumerate(row):
                            if cell is None:
                                continue
                            canonical = _normalize_mud_property_name(str(cell))
                            if not canonical:
                                continue
                            for other in row[idx + 1:]:
                                if other is None:
                                    continue
                                _mud_apply_canonical_value(row_record, canonical, other)
                                if canonical in row_record or canonical.startswith("Gel"):
                                    break
                        full_text = "\n".join(full_text_parts)
            dt_pdf = _extract_date_from_text(full_text)
            if dt_pdf is not None:
                row_record["Date"] = dt_pdf
            _parse_mud_text_block(full_text, row_record)
            _parse_mud_lines(full_text, row_record)
        if any(k for k in row_record if k not in ("Date", "Source")):
            out.append(row_record)
    except Exception:
        pass
    return out


def _fetch_mud_attachments_from_email(
    imap_server: str,
    imap_user: str,
    imap_pass: str,
    filename_contains: str | None = None,
    mark_read: bool = True,
) -> list[tuple[str, bytes]]:
    """
    Descarga adjuntos PDF/Excel/CSV de correos no leídos por IMAP.
    filename_contains: filtro opcional (ej. "Daily Full Report" o "LA-358").
    mark_read: si True, marca los correos como leídos tras descargar.
    Retorna lista de (nombre_archivo, contenido_bytes).
    """
    results: list[tuple[str, bytes]] = []
    try:
        import imaplib
        import email as email_module
    except ImportError:
        return results
    try:
        with imaplib.IMAP4_SSL(imap_server, timeout=30) as mail:
            mail.login(imap_user, imap_pass)
            mail.select("inbox")
            status, messages = mail.search(None, "(UNSEEN)")
            if status != "OK":
                return results
            for num in (messages[0] or b"").split():
                if not num:
                    continue
                status, data = mail.fetch(num, "(RFC822)")
                if status != "OK":
                    continue
                msg = email_module.message_from_bytes(data[0][1])
                for part in msg.walk():
                    if part.get_content_disposition() != "attachment":
                        continue
                    filename = part.get_filename()
                    if not filename:
                        continue
                    filename = str(filename).strip()
                    ext = (filename or "").lower()
                    if not (
                        ext.endswith(".pdf")
                        or ext.endswith(".xlsx")
                        or ext.endswith(".xls")
                        or ext.endswith(".csv")
                    ):
                        continue
                    if filename_contains and filename_contains.strip():
                        if filename_contains.strip().lower() not in filename.lower():
                            continue
                    payload = part.get_payload(decode=True)
                    if payload:
                        results.append((filename, bytes(payload)))
                if mark_read and results:
                    try:
                        mail.store(num, "+FLAGS", "\\Seen")
                    except Exception:
                        pass
    except Exception:
        raise
    return results


def _build_mud_bitacora(parsed_rows: list[dict]) -> pd.DataFrame:
    """Construye DataFrame bitácora con columnas canónicas."""
    if not parsed_rows:
        return pd.DataFrame()
    all_keys = set()
    for r in parsed_rows:
        all_keys.update(r.keys())
    cols = [c for c in MUD_CANONICAL_ORDER if c in all_keys]
    for c in sorted(all_keys):
        if c not in cols:
            cols.append(c)
    rows = []
    for r in parsed_rows:
        row = {}
        for k in cols:
            row[k] = r.get(k)
        rows.append(row)
    df = pd.DataFrame(rows)
    if "Date" in df.columns:
        df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
        df = df.dropna(subset=["Date"]).sort_values("Date").reset_index(drop=True)
    if "DateTime" in df.columns:
        df["DateTime"] = df["DateTime"].fillna("")
    return df


def _mud_numeric_property_columns(bitacora: pd.DataFrame) -> list[str]:
    preferred = [
        "FL Temp", "Density @ °C", "FV", "FV Temp", "PV", "PV Temp", "YP",
        "Gel_10s", "Gel_10min", "Gel_30min", "tau0",
        "L600", "L300", "L200", "L100", "L6", "L3",
        "HTHP", "HTHP @ °C", "Corr Solid", "NAP", "Water", "NAP Ratio", "Water Ratio",
        "Sand", "Cake (HTHP)", "Chlorides", "Calcium", "CaCl2", "Water Phase Salinity",
        "Excess Lime", "Electrical_Stability", "LGS (%)", "HGS (%)", "LGS (kg/m³)", "HGS (kg/m³)",
        "ASG", "n (HB)", "K (HB)", "Viscometer Sag Shoe Test", "(VSST)",
        "Marsh", "Temperature", "VA", "Filtrado", "Enjarre", "LGS", "HGS", "Solids", "Oil", "RAA",
        "AgNO3", "Salinity", "Alkalinity", "Excess_Cal",
    ]
    cols = []
    for c in preferred:
        if c in bitacora.columns and pd.api.types.is_numeric_dtype(bitacora[c]):
            cols.append(c)
    for c in bitacora.columns:
        if c in cols or c == "Date" or c in MUD_ANALYTIC_EXCLUDE:
            continue
        if pd.api.types.is_numeric_dtype(bitacora[c]):
            cols.append(c)
    return cols


def _mud_build_view_df(bitacora: pd.DataFrame) -> pd.DataFrame:
    if bitacora is None or bitacora.empty:
        return pd.DataFrame()
    view = bitacora.copy()
    if "DateTime" not in view.columns and "Date" in view.columns:
        view["DateTime"] = pd.to_datetime(view["Date"], errors="coerce").dt.strftime("%Y-%m-%dT%H:%M:%S")
    else:
        dt_series = pd.to_datetime(view.get("Date"), errors="coerce")
        mask = view["DateTime"].astype(str).str.strip().eq("")
        if mask.any():
            view.loc[mask, "DateTime"] = dt_series.loc[mask].dt.strftime("%Y-%m-%dT%H:%M:%S")
    if "Time" not in view.columns and "Date" in view.columns:
        view["Time"] = pd.to_datetime(view["Date"], errors="coerce").dt.strftime("%H:%M")
    if "FV @ °C" not in view.columns and "FV" in view.columns:
        if "FV Temp" in view.columns:
            view["FV @ °C"] = view.apply(lambda r: f"{_mud_num_to_text(r['FV'])} @ {_mud_num_to_text(r['FV Temp'])}" if pd.notna(r.get("FV")) and pd.notna(r.get("FV Temp")) else (_mud_num_to_text(r['FV']) if pd.notna(r.get("FV")) else ""), axis=1)
        else:
            view["FV @ °C"] = view["FV"].map(_mud_num_to_text)
    if "PV @ °C" not in view.columns and "PV" in view.columns:
        if "PV Temp" in view.columns:
            view["PV @ °C"] = view.apply(lambda r: f"{_mud_num_to_text(r['PV'])} @ {_mud_num_to_text(r['PV Temp'])}" if pd.notna(r.get("PV")) and pd.notna(r.get("PV Temp")) else (_mud_num_to_text(r['PV']) if pd.notna(r.get("PV")) else ""), axis=1)
        else:
            view["PV @ °C"] = view["PV"].map(_mud_num_to_text)
    if "Properties" not in view.columns:
        view["Properties"] = np.arange(1, len(view) + 1)
    if "Additional Properties" not in view.columns:
        view["Additional Properties"] = view["Properties"]
    export_cols = [c for c, _, _ in MUD_EXPORT_HEADER_SPECS if c in view.columns]
    for c in export_cols:
        if c in ("DateTime", "Time"):
            view[c] = view[c].fillna("")
    return view[export_cols]


def _export_mud_bitacora_excel(view_df: pd.DataFrame) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = "Table 1.1"

    headers = [(c, h1, h2) for c, h1, h2 in MUD_EXPORT_HEADER_SPECS if c in view_df.columns]
    last_col = len(headers)
    if last_col == 0:
        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)
        return buf.getvalue()

    title_date = ""
    if "Date" in view_df.columns:
        dt0 = pd.to_datetime(view_df["Date"], errors="coerce")
        if hasattr(dt0, "notna") and dt0.notna().any():
            title_date = dt0.dropna().min().strftime("%Y-%m-%d")
    elif "DateTime" in view_df.columns:
        dt0 = pd.to_datetime(view_df["DateTime"], errors="coerce")
        if hasattr(dt0, "notna") and dt0.notna().any():
            title_date = dt0.dropna().min().strftime("%Y-%m-%d")
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=last_col)
    ws.cell(1, 1).value = f"Daily Fluid Properties Daily Report\nReport: {title_date}" if title_date else "Daily Fluid Properties Daily Report"
    ws.cell(1, 1).alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    ws.cell(1, 1).font = Font(size=14, bold=True)
    ws.row_dimensions[1].height = 34

    fill_header = PatternFill("solid", fgColor="D9D9D9")
    fill_sub = PatternFill("solid", fgColor="EDEDED")
    thin_gray = Side(style="thin", color="BFBFBF")
    border = Border(top=thin_gray, bottom=thin_gray)

    for idx, (_, h1, h2) in enumerate(headers, start=1):
        c2 = ws.cell(2, idx, h1)
        c3 = ws.cell(3, idx, h2)
        for cell in (c2, c3):
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            cell.font = Font(bold=True)
            cell.border = border
        c2.fill = fill_header
        c3.fill = fill_sub

    widths = {
        "Depth (MD)": 11, "Depth (TVD)": 11, "Properties": 10, "Fluid set": 16, "Source": 14,
        "Time": 10, "DateTime": 22, "FL Temp": 10, "Density @ °C": 14,
        "FV @ °C": 12, "PV @ °C": 12, "YP": 10, "Gel_10s": 10, "Gel_10min": 11,
        "Gel_30min": 11, "tau0": 10, "L600": 9, "L300": 9, "L200": 9, "L100": 9, "L6": 9, "L3": 9,
        "HTHP": 10, "HTHP @ °C": 9, "Corr Solid": 11, "NAP": 9, "Water": 9, "NAP Ratio": 10,
        "Water Ratio": 12, "Sand": 9, "Cake (HTHP)": 12, "Chlorides": 12, "Calcium": 12,
        "CaCl2": 11, "Water Phase Salinity": 16, "NaCL (Sol/Insol)": 16, "Excess Lime": 12,
        "Electrical_Stability": 13, "LGS (%)": 10, "HGS (%)": 10, "LGS (kg/m³)": 12, "HGS (kg/m³)": 12,
        "ASG": 9, "Additional Properties": 16, "n (HB)": 12, "K (HB)": 12,
        "Viscometer Sag Shoe Test": 20, "(VSST)": 10,
    }
    num_format = "0.00"
    int_format = "0"
    row_start = 4
    for r_idx, (_, row) in enumerate(view_df.iterrows(), start=row_start):
        for c_idx, (col_name, _, _) in enumerate(headers, start=1):
            val = row.get(col_name)
            cell = ws.cell(r_idx, c_idx, val)
            if pd.isna(val):
                cell.value = None
            elif col_name == "DateTime" and str(val).strip():
                cell.number_format = "@"
            elif isinstance(val, (int, np.integer)):
                cell.number_format = int_format
            elif isinstance(val, (float, np.floating)) and np.isfinite(float(val)):
                cell.number_format = num_format if abs(float(val) - round(float(val))) > 1e-9 else int_format
            cell.alignment = Alignment(horizontal="center", vertical="center")
        ws.row_dimensions[r_idx].height = 21

    for idx, (col_name, _, _) in enumerate(headers, start=1):
        ws.column_dimensions[get_column_letter(idx)].width = widths.get(col_name, 12)
    ws.freeze_panes = "A4"

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _send_mud_bitacora_email(
    attachment_bytes: bytes,
    to_email: str,
    subject: str,
    body: str,
    filename: str = "mud_bitacora.xlsx",
    smtp_server: str = MUD_SMTP_SERVER,
    smtp_port: int = MUD_SMTP_PORT,
    smtp_user: str = MUD_SMTP_USER,
    smtp_pass: str = MUD_SMTP_PASS,
    from_email: str = MUD_SMTP_FROM,
) -> tuple[bool, str]:
    """Envía la bitácora Excel por correo como adjunto."""
    if not smtp_user or not smtp_pass:
        return False, "Faltan credenciales SMTP. Configura MUD_SMTP_USER y MUD_SMTP_PASS en secrets."
    try:
        msg = EmailMessage()
        msg["Subject"] = subject
        msg["From"] = from_email
        msg["To"] = to_email
        msg.set_content(body)
        msg.add_attachment(
            attachment_bytes,
            maintype="application",
            subtype="vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=filename,
        )
        with smtplib.SMTP(smtp_server, smtp_port, timeout=30) as server:
            server.starttls()
            server.login(smtp_user, smtp_pass)
            server.send_message(msg)
        return True, f"Bitácora enviada correctamente a {to_email}."
    except Exception as e:
        return False, str(e)


def render_mud_report() -> None:
    _ms = st.session_state.get("mud_data_source")
    if _ms == "Correo electrónico":
        st.session_state["mud_data_source"] = MUD_SRC_EMAIL
    elif _ms == "Subir archivos":
        st.session_state["mud_data_source"] = MUD_SRC_FILES

    st.markdown(
        f"""
        <div style="margin-bottom: 0.5rem;">
            <span style="font-size: 1.5rem; font-weight: 600;">{tr("tab_mud")}</span>
            <span style="display: inline-flex; align-items: center; gap: 0.35rem; margin-left: 0.75rem; flex-wrap: wrap;">
                <span style="background: linear-gradient(135deg, #b91c1c 0%, #ea580c 50%, #f59e0b 100%); color: #fff; font-size: 0.7rem; font-weight: 700; padding: 0.22rem 0.6rem; border-radius: 999px; letter-spacing: 0.03em; box-shadow: 0 1px 3px rgba(234,88,12,0.4);">🔥 Rogii</span>
                <span style="background: linear-gradient(135deg, #0f766e 0%, #14b8a6 100%); color: #fff; font-size: 0.7rem; font-weight: 600; padding: 0.2rem 0.55rem; border-radius: 999px;">{tr("mud_chip_bitacora")}</span>
                <span style="background: linear-gradient(135deg, #1e3a5f 0%, #2563eb 100%); color: #fff; font-size: 0.7rem; font-weight: 600; padding: 0.2rem 0.55rem; border-radius: 999px;">{tr("mud_chip_formats")}</span>
                <span style="background: linear-gradient(135deg, #7c2d12 0%, #ea580c 100%); color: #fff; font-size: 0.7rem; font-weight: 600; padding: 0.2rem 0.55rem; border-radius: 999px;">{tr("mud_chip_mail_short")}</span>
            </span>
        </div>
        """,
        unsafe_allow_html=True,
    )
    st.caption(tr("mud_caption"))

    mud_source = st.radio(
        tr("data_source"),
        [MUD_SRC_FILES, MUD_SRC_EMAIL],
        horizontal=True,
        key="mud_data_source",
        format_func=lambda x: tr("mud_src_files") if x == MUD_SRC_FILES else tr("mud_src_email"),
    )

    parsed: list[dict] = []

    # Chips de contexto (Rogii + fuente + Auto 60s si aplica)
    mud_chip_items = [
        ("🔥 Rogii", "#b91c1c", "#ea580c"),
        (tr("mud_chip_email"), "#1e3a5f", "#2563eb")
        if mud_source == MUD_SRC_EMAIL
        else (tr("mud_chip_upload"), "#0f766e", "#14b8a6"),
    ]
    if mud_source == MUD_SRC_EMAIL and st.session_state.get("mud_auto_refresh", False):
        mud_chip_items.append(("Auto 60s 🔥", "#7c2d12", "#ea580c"))
    mud_cols = st.columns(len(mud_chip_items))
    for i, (label, c1, c2) in enumerate(mud_chip_items):
        with mud_cols[i]:
            st.markdown(
                f'<span style="display:inline-flex;align-items:center;gap:0.25rem;'
                f"background:linear-gradient(135deg,{c1},{c2});color:#fff;font-size:0.75rem;font-weight:600;"
                f'padding:0.25rem 0.6rem;border-radius:999px;box-shadow:0 1px 2px rgba(0,0,0,0.2);">{label}</span>',
                unsafe_allow_html=True,
            )
    st.markdown("<div style='height:6px'></div>", unsafe_allow_html=True)

    if mud_source == MUD_SRC_EMAIL:
        with st.expander(tr("mud_imap_expander"), expanded=True):
            st.caption(tr("mud_imap_caption"))
            col_imap1, col_imap2 = st.columns(2)
            with col_imap1:
                imap_server = st.text_input(
                    tr("mud_imap_server"),
                    value=os.getenv("MUD_IMAP_SERVER", "imap.gmail.com"),
                    key="mud_imap_server",
                    help=tr("mud_imap_server_help"),
                )
                imap_user = st.text_input(
                    tr("mud_imap_user"),
                    value=os.getenv("MUD_IMAP_USER", ""),
                    key="mud_imap_user",
                )
            with col_imap2:
                imap_pass = st.text_input(
                    tr("mud_imap_pass"),
                    value=os.getenv("MUD_IMAP_PASS", ""),
                    type="password",
                    key="mud_imap_pass",
                    help=tr("mud_imap_pass_help"),
                )
                filename_filter = st.text_input(
                    tr("mud_imap_filter"),
                    value=os.getenv("MUD_IMAP_FILTER", ""),
                    placeholder='Ej: "Daily Full Report" o "LA-358"',
                    key="mud_imap_filter",
                )
            mark_read = st.checkbox(
                tr("mud_mark_read"),
                value=True,
                key="mud_imap_mark_read",
            )

        st.markdown(tr("mud_auto_hdr"))
        mud_auto_refresh = st.checkbox(
            tr("mud_auto_chk"),
            value=st.session_state.get("mud_auto_refresh", False),
            key="mud_auto_refresh",
            help=tr("mud_auto_help"),
        )
        if mud_auto_refresh:
            mud_refresh_interval = st.number_input(
                tr("mud_interval_imap"),
                min_value=30,
                max_value=300,
                value=60,
                step=15,
                key="mud_auto_refresh_interval",
                help=tr("mud_interval_imap_help"),
            )

        run_fetch = st.button(
            tr("mud_fetch_btn"),
            type="primary",
            key="mud_fetch_email_btn",
            help=tr("mud_fetch_help"),
        ) or (
            mud_auto_refresh
            and st.session_state.pop("mud_auto_rerun_trigger", False)
        )

        if run_fetch:
            if not imap_server or not imap_user or not imap_pass:
                st.error(tr("mud_err_imap"))
            else:
                with st.spinner(tr("mud_spinner_imap")):
                    try:
                        attachments = _fetch_mud_attachments_from_email(
                            imap_server.strip(),
                            imap_user.strip(),
                            imap_pass.strip(),
                            filename_contains=filename_filter.strip() or None,
                            mark_read=mark_read,
                        )
                    except Exception as e:
                        st.error(f"{tr('mud_err_dl')} {e}")
                        attachments = []
                if not attachments:
                    st.info(tr("mud_no_attach"))
                else:
                    st.success(tr("mud_success_attach").format(n=len(attachments)))
                    for name, data in attachments:
                        try:
                            if name.lower().endswith(".pdf"):
                                buf = io.BytesIO(data)
                                parsed.extend(_parse_mud_pdf(buf, source_name=name))
                            elif name.lower().endswith((".xlsx", ".xls")):
                                xl = pd.ExcelFile(io.BytesIO(data))
                                for sh in xl.sheet_names[:5]:
                                    df_raw = pd.read_excel(xl, sheet_name=sh, header=None)
                                    parsed.extend(_parse_mud_excel_sheet(df_raw, name))
                            else:
                                df_raw = pd.read_csv(io.BytesIO(data), sep=None, engine="python", low_memory=False)
                                parsed.extend(_parse_mud_csv(df_raw, name))
                        except Exception as e:
                            st.warning(f"No se pudo procesar **{name}**: {e}")
                    if parsed:
                        bitacora_new = _build_mud_bitacora(parsed)
                        existing = st.session_state.get("mud_bitacora")
                        if existing is not None and not existing.empty:
                            bitacora_combined = pd.concat([existing, bitacora_new], ignore_index=True)
                            bitacora_combined["Date"] = pd.to_datetime(bitacora_combined["Date"], errors="coerce")
                            bitacora_combined = bitacora_combined.dropna(subset=["Date"]).sort_values("Date").drop_duplicates().reset_index(drop=True)
                            st.session_state["mud_bitacora"] = bitacora_combined
                        else:
                            st.session_state["mud_bitacora"] = bitacora_new
                        st.success(f"Bitácora actualizada con **{len(parsed)}** registro(s) desde correo.")
                        st.rerun()
                    else:
                        st.warning("No se detectaron propiedades de lodo en los adjuntos.")

    else:
        uploaded = st.file_uploader(
            tr("mud_upload_reports"),
            type=["pdf", "xlsx", "xls", "csv"],
            accept_multiple_files=True,
            key="mud_upload",
        )

        if uploaded:
            for f in uploaded:
                name = getattr(f, "name", "") or ""
                try:
                    if name.lower().endswith(".pdf"):
                        parsed.extend(_parse_mud_pdf(f, source_name=name))
                    elif name.lower().endswith((".xlsx", ".xls")):
                        xl = pd.ExcelFile(f)
                        for sh in xl.sheet_names[:5]:
                            df_raw = pd.read_excel(xl, sheet_name=sh, header=None)
                            parsed.extend(_parse_mud_excel_sheet(df_raw, name))
                    else:
                        df_raw = pd.read_csv(f, sep=None, engine="python", low_memory=False)
                        parsed.extend(_parse_mud_csv(df_raw, name))
                except Exception as e:
                    st.warning(f"No se pudo procesar **{name}**: {e}")

            if parsed:
                bitacora = _build_mud_bitacora(parsed)
                st.session_state["mud_bitacora"] = bitacora
            else:
                st.warning("No se detectaron propiedades de lodo en los archivos. Revisa que contengan columnas o celdas como Density, MW, PV, YP, Gels, etc.")
                if "mud_bitacora" in st.session_state:
                    del st.session_state["mud_bitacora"]

    bitacora = st.session_state.get("mud_bitacora")
    if bitacora is None or bitacora.empty:
        st.info("Sube uno o más reportes (PDF, Excel o CSV) para generar la bitácora.")
        return
    bitacora_view = _mud_build_view_df(bitacora)

    # Chips pro Rogii sobre la bitácora
    n_reg = len(bitacora)
    bitacora_chips = [
        ("🔥 Rogii", "#b91c1c", "#ea580c"),
        (f"{n_reg:,} registros", "#0f766e", "#14b8a6"),
        ("Bitácora", "#1e3a5f", "#2563eb"),
    ]
    if bitacora["Date"].notna().any():
        d_min = bitacora["Date"].min()
        d_max = bitacora["Date"].max()
        if hasattr(d_min, "strftime"):
            bitacora_chips.append((f"{d_min.strftime('%d/%b')} – {d_max.strftime('%d/%b')}", "#334155", "#64748b"))
    bitacora_cols = st.columns(len(bitacora_chips))
    for i, (label, c1, c2) in enumerate(bitacora_chips):
        with bitacora_cols[i]:
            st.markdown(
                f'<span style="display:inline-flex;align-items:center;gap:0.25rem;'
                f"background:linear-gradient(135deg,{c1},{c2});color:#fff;font-size:0.75rem;font-weight:600;"
                f'padding:0.28rem 0.65rem;border-radius:999px;box-shadow:0 1px 3px rgba(0,0,0,0.15);">{label}</span>',
                unsafe_allow_html=True,
            )
    st.success(f"Bitácora: **{n_reg:,}** registros por fecha.")
    tab_bitacora, tab_graficas, tab_stats = st.tabs(["Bitácora", "Gráficas y evolución", "Estadísticas"])

    with tab_bitacora:
        st.subheader("Bitácora de propiedades de fluidos")
        st.dataframe(bitacora_view, use_container_width=True, hide_index=True)

        buf_csv = io.BytesIO()
        bitacora_view.to_csv(buf_csv, index=False, encoding="utf-8-sig")
        buf_csv.seek(0)
        xlsx_bytes = _export_mud_bitacora_excel(bitacora_view)

        col1, col2, col3 = st.columns(3)
        with col1:
            st.download_button(
                "Exportar bitácora (CSV)",
                data=buf_csv.getvalue(),
                file_name="mud_bitacora.csv",
                mime="text/csv",
                key="mud_export_csv",
            )
        with col2:
            st.download_button(
                "Exportar bitácora (Excel)",
                data=xlsx_bytes,
                file_name="mud_bitacora.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="mud_export_xlsx",
            )
        with col3:
            if st.button("Enviar bitácora por correo", key="mud_send_email_btn", type="secondary"):
                date_label = ""
                try:
                    if "Date" in bitacora.columns and bitacora["Date"].notna().any():
                        dmax = pd.to_datetime(bitacora["Date"], errors="coerce").dropna().max()
                        if pd.notna(dmax):
                            date_label = dmax.strftime("%Y-%m-%d")
                except Exception:
                    date_label = ""
                subject = f"Mud bitácora {date_label}".strip()
                body = (
                    "Hola,\n\n"
                    "Adjunto la bitácora de propiedades de fluidos generada desde la app.\n\n"
                    "Saludos."
                )
                ok, msg = _send_mud_bitacora_email(
                    attachment_bytes=xlsx_bytes,
                    to_email=MUD_SMTP_TO,
                    subject=subject,
                    body=body,
                    filename="mud_bitacora.xlsx",
                )
                if ok:
                    st.success(msg)
                else:
                    st.error(f"No se pudo enviar la bitácora por correo: {msg}")

        with st.expander("Configuración de envío por correo", expanded=False):
            st.caption("Estos valores se leen desde st.secrets o variables de entorno.")
            e1, e2 = st.columns(2)
            with e1:
                st.text_input("SMTP server", value=MUD_SMTP_SERVER, disabled=True, key="mud_smtp_server_view")
                st.text_input("SMTP user", value=MUD_SMTP_USER, disabled=True, key="mud_smtp_user_view")
                st.text_input("From", value=MUD_SMTP_FROM, disabled=True, key="mud_smtp_from_view")
            with e2:
                st.text_input("SMTP port", value=str(MUD_SMTP_PORT), disabled=True, key="mud_smtp_port_view")
                st.text_input("To", value=MUD_SMTP_TO, disabled=True, key="mud_smtp_to_view")
                st.text_input("SMTP password", value=("********" if MUD_SMTP_PASS else ""), type="password", disabled=True, key="mud_smtp_pass_view")

    with tab_graficas:
        st.subheader("Evolución de propiedades por día")
        props = _mud_numeric_property_columns(bitacora)
        if not props:
            st.info("No hay columnas numéricas para graficar.")
        else:
            st.caption("Selecciona las propiedades a graficar")
            default_first = min(4, len(props))
            n_cols = 4
            n_rows = (len(props) + n_cols - 1) // n_cols
            checkbox_state = {}
            for row in range(n_rows):
                cols = st.columns(n_cols)
                for col_idx in range(n_cols):
                    i = row * n_cols + col_idx
                    if i >= len(props):
                        break
                    p = props[i]
                    default = i < default_first
                    checkbox_state[p] = st.checkbox(
                        p,
                        value=st.session_state.get(f"mud_cb_{p}", default),
                        key=f"mud_cb_{p}",
                        label_visibility="visible",
                    )
            selected = [p for p in props if checkbox_state.get(p, False)]
            if selected:
                df_plot = bitacora[["Date"] + selected].copy()
                df_plot = df_plot.set_index("Date").sort_index()
                df_plot = df_plot.reset_index()
                fig = go.Figure()
                for p in selected:
                    fig.add_trace(
                        go.Scatter(
                            x=df_plot["Date"],
                            y=df_plot[p],
                            mode="lines+markers",
                            name=p,
                            line=dict(width=2),
                            marker=dict(size=8),
                        )
                    )
                fig.update_layout(
                    title="Evolución de propiedades de lodo",
                    xaxis_title="Fecha",
                    yaxis_title="Valor",
                    legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
                    hovermode="x unified",
                )
                fig = prettify_auto(fig, h=480)
                st.plotly_chart(fig, use_container_width=True, config=PLOTLY_CONFIG)
                # Chips por propiedad
                chip_items = []
                for p in selected:
                    s = bitacora[p].dropna()
                    if len(s):
                        chip_items.append((p, "blue", f"min {format_num(s.min())} · max {format_num(s.max())} · n={len(s):,}"))
                if chip_items:
                    nchips = len(chip_items)
                    cols_chip = st.columns(min(nchips, 4))
                    for i, (label, color, sub) in enumerate(chip_items):
                        with cols_chip[i % len(cols_chip)]:
                            st.caption(f"**{label}** — {sub}")

            st.markdown("---")
            st.subheader("Gráfica por propiedad (selección única)")
            single_prop = st.selectbox("Propiedad", ["(ninguna)"] + props, key="mud_single_prop")
            if single_prop and single_prop != "(ninguna)":
                s = bitacora[single_prop].dropna()
                if len(s):
                    fig1 = px.line(bitacora, x="Date", y=single_prop, title=f"Evolución – {single_prop}", markers=True)
                    fig1.update_traces(line=dict(width=2), marker=dict(size=10))
                    fig1 = prettify(fig1, h=420)
                    st.plotly_chart(fig1, use_container_width=True, config=PLOTLY_CONFIG)
                    _render_chips_row([(single_prop, "blue"), (f"min {format_num(s.min())}", "gray"), (f"max {format_num(s.max())}", "gray"), (f"promedio {format_num(s.mean())}", "green")])

            st.markdown("---")
            st.subheader("🔥 Gráficas adicionales pro")
            st.caption("Heatmap de correlación, gráfico de control y perfil radar para análisis de fluidos.")

            # 1) Heatmap de correlación entre propiedades de lodo
            if len(props) >= 2:
                st.markdown("**Correlación entre propiedades**")
                corr_mud = bitacora[props].corr()
                if not corr_mud.isna().all().all():
                    _mud_hm_stats = heatmap_numeric_stats(bitacora, props)
                    _mud_chips = stats_df_to_heatmap_chips(_mud_hm_stats, max_chips=10)
                    if _mud_chips:
                        st.caption("**Chips — min–max y media por propiedad (base del heatmap)**")
                        _render_chips_row(_mud_chips)
                    with st.expander("Min / media / max por propiedad (lodo)", expanded=False):
                        st.dataframe(_mud_hm_stats, use_container_width=True, hide_index=True)
                    corr_pct = (corr_mud * 100).round(0)
                    text_arr = np.where(
                        np.isnan(corr_pct.values),
                        "",
                        (np.nan_to_num(corr_pct.values, nan=0.0).astype(int)).astype(str) + "%",
                    )
                    fig_corr_mud = px.imshow(
                        corr_mud,
                        color_continuous_scale="RdBu",
                        zmin=-1,
                        zmax=1,
                        title="Mud Report – Correlación entre propiedades",
                    )
                    fig_corr_mud.update_traces(
                        text=text_arr,
                        texttemplate="%{text}",
                        textfont=dict(size=11),
                        xgap=1,
                        ygap=1,
                    )
                    fig_corr_mud.update_layout(coloraxis_colorbar=dict(title="Corr (-1 a 1)"))
                    fig_corr_mud = prettify_heatmap_auto(fig_corr_mud, h=420)
                    st.plotly_chart(fig_corr_mud, use_container_width=True, config=PLOTLY_CONFIG)
                    _sp_mud = build_minmax_mean_spine_figure(
                        _mud_hm_stats,
                        title="Perfil min · media · max — propiedades de lodo (normalizado)",
                    )
                    if _sp_mud is not None:
                        st.caption("**Curvas pro:** rango observado por variable (● = media en el rango min–max).")
                        st.plotly_chart(_sp_mud, use_container_width=True, config=PLOTLY_CONFIG)
                    st.caption("Rojo = correlación positiva, azul = negativa. Valores en % de fuerza lineal.")
                else:
                    st.info("No hay suficientes datos para calcular correlaciones.")
            else:
                st.caption("Se necesitan al menos 2 propiedades numéricas para el heatmap de correlación.")

            # 2) Gráfico de control (propiedad vs fecha, media ± 2σ)
            st.markdown("**Gráfico de control (media ± 2σ)**")
            ctrl_prop = st.selectbox(
                "Propiedad para control",
                props,
                key="mud_ctrl_prop",
                help="Puntos fuera de la banda se marcan en naranja.",
            )
            if ctrl_prop:
                df_ctrl = bitacora[["Date", ctrl_prop]].dropna()
                if len(df_ctrl) >= 2:
                    mean_val = float(df_ctrl[ctrl_prop].mean())
                    std_val = float(df_ctrl[ctrl_prop].std()) or 1e-6
                    upper = mean_val + 2 * std_val
                    lower = mean_val - 2 * std_val
                    df_ctrl = df_ctrl.copy()
                    df_ctrl["_out"] = (df_ctrl[ctrl_prop] > upper) | (df_ctrl[ctrl_prop] < lower)
                    fig_ctrl = go.Figure()
                    in_spec = df_ctrl[~df_ctrl["_out"]]
                    out_spec = df_ctrl[df_ctrl["_out"]]
                    if not in_spec.empty:
                        fig_ctrl.add_trace(
                            go.Scatter(
                                x=in_spec["Date"],
                                y=in_spec[ctrl_prop],
                                mode="markers",
                                name="Dentro de límites",
                                marker=dict(size=8, color="#2563EB", opacity=0.85),
                            )
                        )
                    if not out_spec.empty:
                        fig_ctrl.add_trace(
                            go.Scatter(
                                x=out_spec["Date"],
                                y=out_spec[ctrl_prop],
                                mode="markers",
                                name="Fuera de límites (±2σ)",
                                marker=dict(size=10, color="#EA580C", symbol="diamond"),
                            )
                        )
                    fig_ctrl.add_hline(y=mean_val, line_dash="dash", line_color="#10B981", annotation_text="Media")
                    fig_ctrl.add_hline(y=upper, line_dash="dot", line_color="#EF4444", annotation_text="UCL")
                    fig_ctrl.add_hline(y=lower, line_dash="dot", line_color="#EF4444", annotation_text="LCL")
                    fig_ctrl.update_layout(
                        title=f"Control chart – {ctrl_prop}",
                        xaxis_title="Fecha",
                        yaxis_title=ctrl_prop,
                        height=400,
                        showlegend=True,
                        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
                    )
                    fig_ctrl = prettify_auto(fig_ctrl, h=400)
                    st.plotly_chart(fig_ctrl, use_container_width=True, config=PLOTLY_CONFIG)
                    st.caption(f"Media = {format_num(mean_val)}, LCL = {format_num(lower)}, UCL = {format_num(upper)}. Puntos fuera de banda = {len(out_spec)}.")
                else:
                    st.info("Se necesitan al menos 2 puntos para el gráfico de control.")

            # 3) Radar / perfil del lodo (último registro o por fecha)
            st.markdown("**Perfil radar (comparación normalizada)**")
            radar_props = [p for p in props if bitacora[p].notna().any()]
            if len(radar_props) >= 3:
                dates_opt = bitacora["Date"].dropna().unique()
                if len(dates_opt) > 0:
                    dates_sorted = sorted(dates_opt, reverse=True)
                    default_idx = 0
                    radar_date = st.selectbox(
                        "Fecha del perfil",
                        options=dates_sorted,
                        format_func=lambda x: x.strftime("%Y-%m-%d %H:%M") if hasattr(x, "strftime") else str(x),
                        index=default_idx,
                        key="mud_radar_date",
                        help="Perfil normalizado 0–100% para esa fecha.",
                    )
                    row = bitacora[bitacora["Date"] == radar_date].iloc[-1]
                    r_vals = []
                    for p in radar_props:
                        v = row.get(p)
                        if pd.isna(v):
                            r_vals.append(0.0)
                        else:
                            s_col = bitacora[p].dropna()
                            if len(s_col) and s_col.max() != s_col.min():
                                norm = (float(v) - float(s_col.min())) / (float(s_col.max()) - float(s_col.min()))
                            else:
                                norm = 0.5
                            r_vals.append(round(norm * 100, 1))
                    fig_radar = go.Figure()
                    fig_radar.add_trace(
                        go.Scatterpolar(
                            r=r_vals + [r_vals[0]],
                            theta=radar_props + [radar_props[0]],
                            fill="toself",
                            name="Perfil normalizado",
                            line=dict(color="#EA580C", width=2),
                            fillcolor="rgba(234,88,12,0.25)",
                        )
                    )
                    fig_radar.update_layout(
                        polar=dict(radialaxis=dict(visible=True, range=[0, 100], tickfont=dict(size=10))),
                        title=f"Perfil de lodo – {radar_date.strftime('%Y-%m-%d') if hasattr(radar_date, 'strftime') else radar_date}",
                        height=460,
                        showlegend=False,
                    )
                    fig_radar = prettify_auto(fig_radar, h=460)
                    st.plotly_chart(fig_radar, use_container_width=True, config=PLOTLY_CONFIG)
                    st.caption("Cada eje = propiedad normalizada 0–100% respecto al min/max del histórico. Útil para comparar perfiles por fecha.")
                else:
                    st.info("No hay fechas válidas para el radar.")
            else:
                st.caption("Se necesitan al menos 3 propiedades numéricas para el perfil radar.")

    with tab_stats:
        st.subheader("Estadísticas por propiedad")
        props = _mud_numeric_property_columns(bitacora)
        if not props:
            st.info("No hay columnas numéricas.")
        else:
            stats_rows = []
            for p in props:
                s = bitacora[p].dropna()
                if len(s):
                    stats_rows.append({
                        "Propiedad": p,
                        "N": len(s),
                        "Min": s.min(),
                        "Max": s.max(),
                        "Media": s.mean(),
                        "Mediana": s.median(),
                        "Desv. est.": s.std(),
                    })
            if stats_rows:
                stats_df = pd.DataFrame(stats_rows)
                st.dataframe(stats_df, use_container_width=True, hide_index=True)
                st.markdown("---")
                st.subheader("Distribución (histograma)")
                prop_hist = st.selectbox("Propiedad para histograma", props, key="mud_hist_prop")
                if prop_hist:
                    vals = bitacora[prop_hist].dropna()
                    fig_hist = build_hist_with_trend(vals, title=f"Distribución – {prop_hist}", x_label=prop_hist, nbins=25)
                    st.plotly_chart(prettify_hist(fig_hist), use_container_width=True, config=PLOTLY_CONFIG)
                    st.caption(f"**Resumen:** {series_summary(vals)}.")

    # Auto-refresh correo cada N segundos (solo si fuente = Correo y hay bitácora)
    if (
        st.session_state.get("mud_data_source") == MUD_SRC_EMAIL
        and st.session_state.get("mud_auto_refresh")
    ):
        interval = int(st.session_state.get("mud_auto_refresh_interval", 60))
        interval = max(30, min(300, interval))
        countdown_placeholder = st.empty()
        for i in range(interval, 0, -1):
            countdown_placeholder.info(tr("mud_email_countdown").format(i=i))
            time.sleep(1)
        countdown_placeholder.empty()
        st.session_state["mud_auto_rerun_trigger"] = True
        st.rerun()



# =========================
# ALERTAS PRO (ROGII)
# =========================
PRIMARY = "#0ea5e9"
DARK = "#0B1220"
GRAY = "#94A3B8"

ALERT_TYPE_META = {
    "OPERATIVA": {"badge_color": "blue", "hex": "#2563EB"},
    "DESEMPEÑO": {"badge_color": "orange", "hex": "#F59E0B"},
    "NOTIFICACIÓN": {"badge_color": "gray", "hex": "#6B7280"},
    "CALIDAD_DE_DATOS": {"badge_color": "red", "hex": "#DC2626"},
}

ALERT_LEVEL_META = {
    "BAJA": {"badge_color": "green", "hex": "#10B981"},
    "MEDIA": {"badge_color": "orange", "hex": "#F59E0B"},
    "ALTA": {"badge_color": "red", "hex": "#EF4444"},
    "CRITICA": {"badge_color": "red", "hex": "#B91C1C"},
}


def _safe_font(size: int = 18, bold: bool = False):
    candidates = []
    if os.name == "nt":
        if bold:
            candidates = [
                "C:/Windows/Fonts/arialbd.ttf",
                "C:/Windows/Fonts/segoeuib.ttf",
            ]
        else:
            candidates = [
                "C:/Windows/Fonts/arial.ttf",
                "C:/Windows/Fonts/segoeui.ttf",
            ]
    for p in candidates:
        if os.path.exists(p):
            try:
                return ImageFont.truetype(p, size=size)
            except Exception:
                pass
    return ImageFont.load_default()


def _alert_chip_html(text: str, bg: str = "#0F172A", fg: str = "white") -> str:
    return (
        f"<span style='display:inline-block;padding:6px 12px;margin:4px;"
        f"border-radius:999px;background:{bg};color:{fg};font-size:13px;"
        f"font-weight:600;border:1px solid rgba(0,0,0,0.08);'>{text}</span>"
    )


def render_alert_chips(payload: dict) -> None:
    tmeta = ALERT_TYPE_META.get(payload["tipo_alerta"], ALERT_TYPE_META["OPERATIVA"])
    lmeta = ALERT_LEVEL_META.get(payload["nivel_alerta"], ALERT_LEVEL_META["MEDIA"])

    chips = [
        _alert_chip_html(payload["tipo_alerta"], tmeta["hex"], "white"),
        _alert_chip_html(f"Nivel {payload['nivel_alerta']}", lmeta["hex"], "white"),
        _alert_chip_html(payload["pozo"], "#E0F2FE", "#075985"),
        _alert_chip_html(payload["equipo"], "#ECFCCB", "#3F6212"),
        _alert_chip_html(payload["etapa"], "#F3E8FF", "#6B21A8"),
        _alert_chip_html(payload["actividad"], "#FEF3C7", "#92400E"),
    ]
    st.markdown("".join(chips), unsafe_allow_html=True)


def build_alert_title(payload: dict) -> str:
    fecha_txt = pd.to_datetime(payload["fecha_apertura"]).strftime("%d-%b-%Y")
    return (
        f"{payload['pozo']} | {payload['equipo']} | "
        f"{payload['actividad']} | {payload['tipo_alerta']} | {fecha_txt}"
    )


def build_alert_insight(payload: dict) -> str:
    if payload["tipo_alerta"] == "OPERATIVA":
        return "Condición operativa fuera de comportamiento esperado detectada en tiempo real."
    if payload["tipo_alerta"] == "DESEMPEÑO":
        return "Desviación de desempeño respecto a parámetros objetivo."
    if payload["tipo_alerta"] == "CALIDAD_DE_DATOS":
        return "Inconsistencia en datos detectada que puede afectar decisiones operativas."
    return "Evento informativo registrado."


def build_alert_recommendations(payload: dict) -> list[str]:
    recomendaciones = []
    txt = (
        f"{payload.get('resumen','')} "
        f"{payload.get('descripcion','')} "
        f"{payload.get('condiciones_operacion','')}"
    ).lower()

    if payload["tipo_alerta"] == "OPERATIVA":
        recomendaciones.extend([
            "Validar tendencia de presión, torque y arrastre antes de continuar la maniobra.",
            "Circular y estabilizar parámetros antes de retomar avance.",
            "Confirmar condición de sarta y herramienta con el equipo de perforación.",
        ])
    elif payload["tipo_alerta"] == "DESEMPEÑO":
        recomendaciones.extend([
            "Comparar ROP, WOB y RPM contra ventana objetivo del tramo.",
            "Revisar ineficiencias por deslizamiento, fricción o limpieza deficiente.",
            "Ajustar parámetros de perforación y validar respuesta en los próximos minutos.",
        ])
    elif payload["tipo_alerta"] == "NOTIFICACIÓN":
        recomendaciones.extend([
            "Registrar evento y comunicar estatus al equipo operativo.",
            "Mantener seguimiento hasta cierre o cambio de condición.",
        ])
    else:
        recomendaciones.extend([
            "Validar sensores, mapping y calidad de transmisión de datos.",
            "Corregir tags o unidades antes de usar la información para toma de decisiones.",
        ])

    if "atrap" in txt:
        recomendaciones.insert(0, "Reducir exposición al riesgo de tubería atrapada y monitorear arrastre/presión en tiempo real.")
    if "pres" in txt:
        recomendaciones.insert(0, "Monitorear incremento anormal de presión y confirmar integridad del sistema de circulación.")
    if "torque" in txt:
        recomendaciones.insert(0, "Verificar incremento de torque y evaluar limpieza del hoyo o interacción mecánica.")

    out = []
    seen = set()
    for r in recomendaciones:
        if r not in seen:
            out.append(r)
            seen.add(r)
    return out[:5]


def _to_pil_image(uploaded_file):
    if uploaded_file is None:
        return None
    try:
        return Image.open(uploaded_file).convert("RGB")
    except Exception:
        return None


def build_alert_canvas(payload: dict, recommendations: list[str], evidence_img=None):
    W, H = 1600, 900
    canvas = Image.new("RGB", (W, H), DARK)
    draw = ImageDraw.Draw(canvas)

    title_font = _safe_font(34, True)
    h1_font = _safe_font(22, True)
    body_font = _safe_font(18, False)
    small_font = _safe_font(15, False)

    tmeta = ALERT_TYPE_META.get(payload["tipo_alerta"], ALERT_TYPE_META["OPERATIVA"])
    level_hex = ALERT_LEVEL_META.get(payload["nivel_alerta"], ALERT_LEVEL_META["MEDIA"])["hex"]

    # Header Rogii
    draw.rounded_rectangle((30, 24, 1570, 120), radius=24, fill="#020617", outline="#1E293B", width=2)
    draw.text((55, 42), "ROGII ALERTA OPERATIVA", fill="white", font=title_font)
    draw.text((55, 82), build_alert_title(payload), fill=GRAY, font=small_font)

    # Type ribbons
    draw.rounded_rectangle((1200, 42, 1380, 76), radius=14, fill=tmeta["hex"])
    draw.text((1220, 50), payload["tipo_alerta"], fill="white", font=small_font)
    draw.rounded_rectangle((1395, 42, 1535, 76), radius=14, fill=level_hex)
    draw.text((1415, 50), payload["nivel_alerta"], fill="white", font=small_font)

    # Evidence panel (recuadro interior azul: 60..660 × 205..775 → 600×570 px)
    _ev_outer = (30, 145, 720, 860)
    _ev_inner = (60, 205, 660, 775)
    _bx0, _by0, _bx1, _by1 = _ev_inner
    _ev_pad = 10
    _ev_max_w = (_bx1 - _bx0) - 2 * _ev_pad
    _ev_max_h = (_by1 - _by0) - 2 * _ev_pad

    draw.rounded_rectangle(_ev_outer, radius=26, fill="#0F172A", outline="#334155", width=2)
    draw.text((52, 162), "Evidencia", fill="white", font=h1_font)
    if evidence_img is not None:
        img = evidence_img.copy().convert("RGB")
        # Contener la captura dentro del marco azul (antes thumbnail 620×560 > ancho útil 580)
        img.thumbnail((_ev_max_w, _ev_max_h), Image.Resampling.LANCZOS)
        ew, eh = img.size
        x = _bx0 + _ev_pad + (_ev_max_w - ew) // 2
        y = _by0 + _ev_pad + (_ev_max_h - eh) // 2
        canvas.paste(img, (x, y))
        draw.rounded_rectangle(_ev_inner, radius=20, outline=tmeta["hex"], width=3)
    else:
        draw.rounded_rectangle(_ev_inner, radius=20, outline="#475569", width=2)
        draw.text((235, 470), "Sin evidencia cargada", fill=GRAY, font=h1_font)

    # Executive summary
    draw.rounded_rectangle((750, 145, 1570, 355), radius=26, fill="#0F172A", outline="#334155", width=2)
    draw.text((772, 162), "Resumen ejecutivo", fill="white", font=h1_font)
    meta_lines = [
        f"Pozo: {payload['pozo']}",
        f"Equipo: {payload['equipo']}",
        f"Etapa: {payload['etapa']}",
        f"Actividad: {payload['actividad']}",
        f"Fecha apertura: {pd.to_datetime(payload['fecha_apertura']).strftime('%d-%b-%Y %H:%M')}",
        f"Responsable: {payload['responsable']}",
    ]
    y = 205
    for line in meta_lines:
        draw.text((775, y), line, fill="white", font=body_font)
        y += 26

    insight = build_alert_insight(payload)
    draw.text((775, y + 6), "Insight", fill=PRIMARY, font=body_font)
    y += 34
    for line in wrap(insight, width=63):
        draw.text((775, y), line, fill=GRAY, font=body_font)
        y += 24

    # Description
    draw.rounded_rectangle((750, 375, 1570, 590), radius=26, fill="#0F172A", outline="#334155", width=2)
    draw.text((772, 392), "Descripción / condiciones", fill="white", font=h1_font)
    desc = (payload["descripcion"].strip() + "\n\n" + payload["condiciones_operacion"].strip()).strip()
    desc_lines = wrap(desc if desc else "Sin descripción.", width=72)
    y = 430
    for line in desc_lines[:6]:
        draw.text((775, y), line, fill=GRAY, font=body_font)
        y += 24

    # Recommendations
    draw.rounded_rectangle((750, 610, 1570, 860), radius=26, fill="#0F172A", outline="#334155", width=2)
    draw.text((772, 627), "Recomendaciones", fill="white", font=h1_font)
    y = 670
    for idx, rec in enumerate(recommendations[:5], start=1):
        bullet = f"{idx}. {rec}"
        for line in wrap(bullet, width=70):
            draw.text((775, y), line, fill=GRAY, font=body_font)
            y += 24
        y += 6

    draw.text((35, 870), f"Generado por Rogii Streamlit · {datetime.now().strftime('%Y-%m-%d %H:%M')}", fill=GRAY, font=small_font)
    return canvas


def save_alert_outputs(payload: dict, recommendations: list[str], evidence_img=None):
    base_name = re.sub(r"[^A-Za-z0-9_-]+", "_", build_alert_title(payload))[:120]
    out_dir = Path(tempfile.mkdtemp(prefix="alerta_pro_"))

    board_img = build_alert_canvas(payload, recommendations, evidence_img)
    png_path = out_dir / f"{base_name}.png"
    pdf_path = out_dir / f"{base_name}.pdf"
    pptx_path = out_dir / f"{base_name}.pptx"

    board_img.save(png_path, format="PNG")
    board_img.save(pdf_path, format="PDF", resolution=150.0)
    try:
        board_img.close()
    except Exception:
        pass

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(png_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(pptx_path))

    copy_report_to_downloads(pptx_path, pdf_path, base_name)
    return png_path, pdf_path, pptx_path




def alert_capture_region(
    region: tuple[int, int, int, int],
    countdown_s: int,
    n_shots: int,
    interval_s: int,
) -> list:
    try:
        import pyautogui  # type: ignore
    except Exception as e:
        st.error(f"❌ No se pudo importar pyautogui: {e}")
        return []

    x, y, w, h = region
    status = st.empty()
    prog = st.progress(0)
    captured_images = []

    if countdown_s > 0:
        for rem in range(int(countdown_s), 0, -1):
            status.warning(f"⏱️ Captura en {rem}s · Región ({x}, {y}, {w}, {h})")
            time.sleep(1)

    for i in range(int(n_shots)):
        status.info(f"📸 Captura {i+1}/{n_shots} · Región ({x}, {y}, {w}, {h})")
        if winsound:
            try:
                winsound.Beep(1100, 180)
            except Exception:
                pass

        img = pyautogui.screenshot(region=(x, y, w, h))
        captured_images.append(img)
        st.session_state["alert_last_capture"] = img
        st.session_state["alert_captures"] = captured_images.copy()

        prog.progress(int(((i + 1) / max(int(n_shots), 1)) * 100))

        if i < int(n_shots) - 1 and int(interval_s) > 0:
            for rem in range(int(interval_s), 0, -1):
                status.info(f"🕒 Siguiente captura en {rem}s")
                time.sleep(1)

    status.success("✅ Captura completada.")
    return captured_images


def _normalize_capture_for_preview(img):
    if img is None:
        return None
    if isinstance(img, Image.Image):
        return img
    try:
        return Image.open(img).convert("RGB")
    except Exception:
        return img


def _merge_available_alert_captures() -> list:
    merged = []
    if st.session_state.get("alert_last_capture") is not None:
        merged.append(st.session_state.get("alert_last_capture"))
    for item in (st.session_state.get("alert_captures") or []):
        merged.append(item)
    for item in (st.session_state.get("captures") or []):
        merged.append(item)
    if st.session_state.get("last_capture") is not None:
        merged.append(st.session_state.get("last_capture"))

    out = []
    seen = set()
    for item in merged:
        marker = id(item)
        if item is not None and marker not in seen:
            out.append(item)
            seen.add(marker)
    return out


def render_alertas_module() -> None:
    st.subheader("🚨 Alertas")
    st.caption("Captura una región desde esta misma pestaña, arma la alerta y exporta PDF/PPTX en una sola hoja.")

    c1, c2, c3 = st.columns([1.15, 1.15, 0.9])
    with c1:
        tipo_alerta = st.selectbox("Tipo de alerta", list(ALERT_TYPE_META.keys()), key="alert_tipo_alerta")
    with c2:
        nivel_alerta = st.selectbox("Nivel de alerta", list(ALERT_LEVEL_META.keys()), index=1, key="alert_nivel_alerta")
    with c3:
        if hasattr(st, "datetime_input"):
            fecha_apertura = st.datetime_input(
                "Fecha apertura", value=datetime.now(), key="alert_fecha_apertura"
            )
        else:
            _fd, _ft = st.columns(2)
            with _fd:
                _d = st.date_input(
                    "Fecha",
                    value=datetime.now().date(),
                    key="alert_fecha_apertura_d",
                )
            with _ft:
                _t = st.time_input(
                    "Hora",
                    value=datetime.now().time(),
                    key="alert_fecha_apertura_t",
                )
            fecha_apertura = datetime.combine(_d, _t)

    col_a, col_b, col_c = st.columns(3)
    with col_a:
        pozo = st.text_input("Pozo", value="RACEMOSA 1001SON", key="alert_pozo")
        equipo = st.text_input("Equipo", value="PMX-305", key="alert_equipo")
    with col_b:
        etapa = st.text_input("Etapa", value="Standpipe", key="alert_etapa")
        actividad = st.text_input("Actividad", value="Levantando BHA de perforación", key="alert_actividad")
    with col_c:
        responsable = st.text_input("Responsable", value="Ingeniero (ADT)", key="alert_responsable")
        
    riesgos_pre = ["Tubería atrapada","Alta presión","Torque elevado","Pérdida de circulación","Falla de sensor"]
    riesgo = st.selectbox("Riesgo", riesgos_pre, key="alert_riesgo")
    

    st.markdown("#### Evidencia · Captura de región")
    cap_a, cap_b, cap_c, cap_d = st.columns(4)
    with cap_a:
        region_x = st.number_input("X", min_value=0, max_value=5000, value=456, step=1, key="alert_region_x")
        region_y = st.number_input("Y", min_value=0, max_value=5000, value=196, step=1, key="alert_region_y")
    with cap_b:
        region_w = st.number_input("Ancho", min_value=10, max_value=8000, value=2088, step=1, key="alert_region_w")
        region_h = st.number_input("Alto", min_value=10, max_value=8000, value=1319, step=1, key="alert_region_h")
    with cap_c:
        countdown_s = st.number_input("Temporizador (s)", min_value=0, max_value=30, value=5, step=1, key="alert_countdown_s")
        n_shots = st.number_input("Núm. capturas", min_value=1, max_value=10, value=1, step=1, key="alert_n_shots")
    with cap_d:
        interval_s = st.number_input("Intervalo entre capturas (s)", min_value=1, max_value=60, value=3, step=1, key="alert_interval_s")
        st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)
        if st.button("📸 Capturar región ahora", type="primary", use_container_width=True, key="alert_capture_now"):
            imgs = alert_capture_region(
                region=(int(region_x), int(region_y), int(region_w), int(region_h)),
                countdown_s=int(countdown_s),
                n_shots=int(n_shots),
                interval_s=int(interval_s),
            )
            if imgs:
                st.session_state["alert_last_capture"] = imgs[-1]
                st.session_state["alert_captures"] = imgs

    available_captures = _merge_available_alert_captures()
    evidence_img = None

    if available_captures:
        idx = st.selectbox(
            "Seleccionar evidencia",
            list(range(len(available_captures))),
            index=max(len(available_captures) - 1, 0),
            format_func=lambda x: f"Captura {x+1}",
            key="alert_capture_selector",
        )
        evidence_img = _normalize_capture_for_preview(available_captures[idx])
        st.image(evidence_img, caption="Captura seleccionada", use_container_width=True)
    else:
        st.info("No hay capturas todavía. Usa 'Capturar región ahora' en esta pestaña.")

    st.markdown("#### Contenido de la alerta")
    resumen = st.text_input("Resumen", value="Presión superior y valor de carga en gancho inverso.", key="alert_resumen")
    condiciones_operacion = st.text_area(
        "Condiciones de la operación",
        value=(
            "Con sarta estacionada y circulando a través de botella de circulación, "
            "se registra incremento de presión y comportamiento inverso del gancho. "
            "Se recomienda estabilizar parámetros antes de retomar viaje."
        ),
        height=120,
        key="alert_condiciones_operacion",
    )
    descripcion = st.text_area(
        "Descripción",
        value=(
            "Se observa presión por encima del comportamiento esperado y variación de carga "
            "en gancho durante la maniobra. Se requiere revisión operativa y seguimiento inmediato."
        ),
        height=140,
        key="alert_descripcion",
    )

    payload = {
        "tipo_alerta": tipo_alerta,
        "nivel_alerta": nivel_alerta,
        "fecha_apertura": fecha_apertura,
        "pozo": pozo,
        "equipo": equipo,
        "etapa": etapa,
        "actividad": actividad,
        "responsable": responsable,
        "riesgo": riesgo,
        "resumen": resumen,
        "descripcion": descripcion,
        "condiciones_operacion": condiciones_operacion,
    }

    render_alert_chips(payload)
    recommendations = build_alert_recommendations(payload)

    st.markdown("#### Recomendaciones sugeridas")
    for rec in recommendations:
        st.markdown(f"- {rec}")


    st.markdown("#### Anotaciones sobre la captura")

    if "annotations" not in st.session_state:
        st.session_state["annotations"] = []

    _iw = _ih = 1920
    if evidence_img is not None:
        _iw, _ih = evidence_img.size

    capture_idx = 0
    if available_captures:
        capture_idx = int(st.session_state.get("alert_capture_selector", 0))

    _MAX_PINS = 12
    _canvas_jd = None
    _sx = _sy = 1.0

    if evidence_img is not None and _HAS_DRAWABLE_CANVAS and _st_canvas_draw is not None:
        _sig = (capture_idx, evidence_img.size)
        _prev_sig = st.session_state.get(f"alert_canvas_ev_sig_{capture_idx}")
        if _prev_sig is not None and _prev_sig != _sig:
            st.session_state.pop(f"alert_canvas_json_{capture_idx}", None)
        st.session_state[f"alert_canvas_ev_sig_{capture_idx}"] = _sig

        bg_rgb, _sx, _sy, _w0, _h0 = _prepare_evidence_for_canvas(evidence_img, max_side=960)
        cw, ch = bg_rgb.size

        st.caption(
            f"Lienzo **profesional**: dibuja con el ratón sobre la captura (resolución de trabajo {cw}×{ch} px; "
            f"exportación a **{_w0}×{_h0} px**). Usa la barra del lienzo para **deshacer / rehacer / borrar todo**."
        )

        with st.expander("Instrucciones del lienzo", expanded=False):
            st.markdown(
                """
- **Flecha:** modo *Línea* — arrastra desde el origen hasta la punta (donde debe apuntar la flecha).
- **Recuadro:** modo *Rectángulo* — arrastra una esquina a la opuesta.
- **Comentario:** modo *Círculo* — coloca un círculo donde quieras el texto; luego escribe en los campos **Texto en pin** (orden = orden en que aparecen los círculos en el dibujo).
- **Mover / ajustar:** en el menú **Herramienta** elige *Mover / transformar*, haz **clic** en la figura (marco de selección) y **arrastra** o usa las esquinas para redimensionar.
                """
            )

        c_tool, c_col, c_w = st.columns([1.2, 0.8, 0.8])
        with c_tool:
            _mode_labels = {
                "line": "Flecha (línea)",
                "rect": "Caja (rectángulo)",
                "circle": "Pin de comentario (círculo)",
                "transform": "Mover / transformar",
            }
            _mode_keys = list(_mode_labels.keys())
            _dm = st.selectbox(
                "Herramienta",
                _mode_keys,
                index=0,
                format_func=lambda k: _mode_labels[k],
                key=f"alert_draw_mode_{capture_idx}",
            )
        with c_col:
            _stroke_hex = st.color_picker(
                "Color del trazo", "#ef4444", key=f"alert_stroke_{capture_idx}"
            )
        with c_w:
            _stroke_w = st.slider(
                "Grosor", 2, 14, 5, key=f"alert_stroke_w_{capture_idx}"
            )

        _init = st.session_state.get(f"alert_canvas_json_{capture_idx}")
        if isinstance(_init, str):
            try:
                _init = json.loads(_init)
            except Exception:
                _init = None
        _init_for_canvas = (
            _fabric_make_selectable_for_transform(_init) if isinstance(_init, dict) else None
        )

        _canvas_result = _st_canvas_draw(
            fill_color="rgba(255, 255, 255, 0)",
            stroke_width=int(_stroke_w),
            stroke_color=_stroke_hex,
            background_image=bg_rgb,
            update_streamlit=True,
            height=int(ch),
            width=int(cw),
            drawing_mode=_dm,
            initial_drawing=_init_for_canvas,
            display_toolbar=True,
            key=f"rogii_alert_canvas_{capture_idx}",
        )

        if _canvas_result.json_data is not None:
            _jd_raw = _canvas_result.json_data
            _jd_store = (
                _fabric_make_selectable_for_transform(_jd_raw)
                if isinstance(_jd_raw, dict)
                else _jd_raw
            )
            st.session_state[f"alert_canvas_json_{capture_idx}"] = _jd_store
            _canvas_jd = _jd_store
        else:
            _canvas_jd = st.session_state.get(f"alert_canvas_json_{capture_idx}")

        _n_pins = _count_fabric_pins(_canvas_jd)
        _pin_slots = min(_n_pins, _MAX_PINS)
        _pin_vals: list[str] = []
        if _n_pins > 0:
            st.markdown("##### Textos en los pins (círculos)")
            st.caption(
                "El **pin 1** es el primer círculo del dibujo. Si hay duda, usa *Mover / transformar* o redibuja."
            )
            if _n_pins > _MAX_PINS:
                st.warning(
                    f"Solo se pueden enlazar textos a los primeros **{_MAX_PINS}** círculos "
                    f"(hay **{_n_pins}** en el lienzo)."
                )
            for _pi in range(_pin_slots):
                _pin_vals.append(
                    st.text_input(
                        f"Texto en pin {_pi + 1}",
                        key=f"alert_pin_{capture_idx}_{_pi}",
                        label_visibility="visible",
                    )
                )

        _pin_use = _pin_vals
        st.session_state["annotations"] = canvas_fabric_json_to_annotations(
            _canvas_jd, _sx, _sy, _pin_use, _stroke_hex
        )

        b_clear = st.button("🗑️ Borrar dibujo del lienzo (solo esta captura)", key=f"alert_canvas_clear_{capture_idx}")
        if b_clear:
            st.session_state.pop(f"alert_canvas_json_{capture_idx}", None)
            for _pi in range(_MAX_PINS):
                k = f"alert_pin_{capture_idx}_{_pi}"
                if k in st.session_state:
                    del st.session_state[k]
            st.session_state["annotations"] = []
            st.rerun()

    elif evidence_img is not None and not _HAS_DRAWABLE_CANVAS:
        st.warning(
            "Para anotar con el ratón instala: `pip install streamlit-drawable-canvas` "
            "y reinicia la app."
        )
        with st.expander("Modo manual por coordenadas", expanded=False):
            st.caption(f"Imagen {_iw}×{_ih} px. Origen arriba-izquierda.")
            tipo_labels = {"arrow": "Flecha", "box": "Caja", "text": "Texto"}
            tipo_values = list(tipo_labels.keys())
            c1, c2, c3 = st.columns(3)
            with c1:
                atype = st.selectbox(
                    "Tipo",
                    tipo_values,
                    format_func=lambda k: tipo_labels[k],
                    key="alert_ann_type_fb",
                )
            with c2:
                color = st.selectbox(
                    "Color",
                    ["red", "yellow", "blue", "green"],
                    key="alert_ann_color_fb",
                )
            with c3:
                text_note = st.text_input("Texto", key="alert_ann_text_fb")
            xa, xb = st.columns(2)
            with xa:
                x1 = st.number_input("x1", 0, max(_iw, 1), 50, key="alert_ann_x1_fb")
                y1 = st.number_input("y1", 0, max(_ih, 1), 50, key="alert_ann_y1_fb")
            with xb:
                x2 = st.number_input("x2", 0, max(_iw, 1), 200, key="alert_ann_x2_fb")
                y2 = st.number_input("y2", 0, max(_ih, 1), 200, key="alert_ann_y2_fb")
            if st.button("Agregar", key="alert_ann_add_fb"):
                st.session_state["annotations"].append(
                    {
                        "type": atype,
                        "color": color,
                        "text": text_note,
                        "x1": int(x1),
                        "y1": int(y1),
                        "x2": int(x2),
                        "y2": int(y2),
                    }
                )
                st.rerun()
            if st.button("Limpiar lista", key="alert_ann_clear_fb"):
                st.session_state["annotations"] = []
                st.rerun()
    else:
        st.info("Carga una captura arriba para usar el lienzo de anotaciones.")

    _ann = st.session_state.get("annotations") or []
    if _ann:
        st.caption(
            f"**{len(_ann)}** marcas aplicadas a la evidencia (orden = orden en vista previa y PDF)."
        )
        with st.expander("Ver datos de marcas", expanded=False):
            st.dataframe(pd.DataFrame(_ann), use_container_width=True, hide_index=True)

    st.markdown("#### Vista previa ejecutiva")
    annotated = draw_annotations(evidence_img, st.session_state.get("annotations", []))
    preview_img = build_alert_canvas(payload, recommendations, annotated)
    st.image(preview_img, use_container_width=True)

    col_x, col_y = st.columns(2)
    with col_x:
        if st.button("📄 Generar PDF alerta", type="primary", use_container_width=True):
            _, pdf_path, pptx_path = save_alert_outputs(payload, recommendations, annotated)
            st.session_state["alert_pdf_path"] = str(pdf_path)
            st.session_state["alert_pptx_path"] = str(pptx_path)
            st.success("PDF generado correctamente.")
    with col_y:
        if st.button("📊 Generar PowerPoint alerta", use_container_width=True):
            _, pdf_path, pptx_path = save_alert_outputs(payload, recommendations, annotated)
            st.session_state["alert_pdf_path"] = str(pdf_path)
            st.session_state["alert_pptx_path"] = str(pptx_path)
            st.success("PowerPoint generado correctamente.")

    pdf_ready = st.session_state.get("alert_pdf_path")
    pptx_ready = st.session_state.get("alert_pptx_path")

    if pdf_ready and Path(pdf_ready).exists():
        with open(pdf_ready, "rb") as f:
            st.download_button(
                "Descargar PDF",
                data=f.read(),
                file_name=Path(pdf_ready).name,
                mime="application/pdf",
                use_container_width=True,
                key="alert_download_pdf",
            )

    if pptx_ready and Path(pptx_ready).exists():
        with open(pptx_ready, "rb") as f:
            st.download_button(
                "Descargar PowerPoint",
                data=f.read(),
                file_name=Path(pptx_ready).name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                key="alert_download_pptx",
            )

tab_kpi, tab_bha, tab_roadmap, tab_trip, tab_mud, tab_alertas = st.tabs(
    [
        tr("tab_kpi"),
        tr("tab_bha"),
        tr("tab_roadmap"),
        tr("tab_trip"),
        tr("tab_mud"),
        "Alertas",
    ]
)

with tab_kpi:
    render_kpi_module()

with tab_bha:
    render_bha_module()

with tab_roadmap:
    render_roadmap()

with tab_trip:
    render_tripping_analysis()

with tab_mud:
    render_mud_report()

with tab_alertas:
    render_alertas_module()
